from __future__ import annotations

import base64
import hashlib
import html
import importlib.util
import json
import os
import re
import sys
from urllib.parse import quote
from pathlib import Path
from typing import Any, Dict

import pandas as pd
import plotly.graph_objects as go
import plotly.io as pio


BASE_DIR = Path(__file__).resolve().parent
REPO_ROOT = next((parent for parent in BASE_DIR.parents if (parent / "package.json").is_file()), BASE_DIR.parents[-1])
for entry in (BASE_DIR, REPO_ROOT):
    entry_str = str(entry)
    if entry_str not in sys.path:
        sys.path.insert(0, entry_str)

try:
    from PyQt6.QtCore import QObject, pyqtSignal, pyqtSlot  # type: ignore
    from PyQt6.QtWidgets import QWidget  # type: ignore
except Exception:  # pragma: no cover
    QObject = object  # type: ignore
    QWidget = object  # type: ignore

    def pyqtSignal(*_args, **_kwargs):  # type: ignore
        return object()

    def pyqtSlot(*_args, **_kwargs):  # type: ignore
        def _decorator(fn):
            return fn

        return _decorator


def _ppt_dbg_on() -> bool:
    v = (os.environ.get("PPT_DEBUG") or "").strip().lower()
    return v in {"1", "true", "yes", "y", "on"}


def _ppt_dbg(msg: str) -> None:
    if _ppt_dbg_on():
        print(msg, flush=True)


_PPT_UI_WARMED = False
_PPT_UI_CACHE = {}

def _fig_to_jsonable(fig: go.Figure) -> dict:
    s = pio.to_json(fig, validate=False)
    return json.loads(s)


def _json_safe(x):
    """
    ###1. normalize JSON primitives (including NaN/Inf -> None)
    ###2. recurse list/tuple/dict
    ###3. normalize numpy/pandas scalars/arrays via module/type hints
    """
    import math

    if x is None:
        return None

    if x is pd.NA:
        return None

    if isinstance(x, bool):
        return x

    if isinstance(x, int):
        return x

    if isinstance(x, float):
        if math.isnan(x) or math.isinf(x):
            return None
        return x

    if isinstance(x, str):
        return x

    if isinstance(x, (list, tuple)):
        return [_json_safe(v) for v in x]

    if isinstance(x, dict):
        return {str(k): _json_safe(v) for k, v in x.items()}

    mod = type(x).__module__
    name = type(x).__name__

    if mod == "numpy":
        return _json_safe(x.tolist())

    if mod.startswith("pandas"):
        if name in {"Timestamp", "Timedelta"}:
            return str(x)
        if name in {"NAType"}:
            return None

    return str(x)


def _canonical_json(obj) -> str:
    safe = _json_safe(obj)
    return json.dumps(safe, sort_keys=True, separators=(",", ":"), ensure_ascii=False)


def _fig_json_md5(fig_json) -> str:
    if fig_json is None:
        return ""
    s = _canonical_json(fig_json)
    return hashlib.md5(s.encode("utf-8")).hexdigest()


def _deck_md5(slides: list[dict]) -> str:
    if not isinstance(slides, list) or not slides:
        return ""
    s = _canonical_json(slides)
    return hashlib.md5(s.encode("utf-8")).hexdigest()


def _ppts_cache_dir(collection_name: str) -> Path:
    name = str(collection_name or "").strip()
    d = (Path(MAIN_APP_CACHE_DIR) / "pages" / "ppts" / name).resolve()
    d.mkdir(parents=True, exist_ok=True)
    return d


def _normalize_img_value(value) -> str:
    if value is None:
        return ""
    s = str(value).strip()
    if not s:
        return ""
    if s.startswith("data:image/"):
        return s
    if s.startswith("file://"):
        return s
    p = Path(s)
    if p.exists():
        return _path_to_file_url(str(p))
    return ""


def _table_thumb_data_url() -> str:
    svg = (
        "<svg xmlns='http://www.w3.org/2000/svg' width='400' height='225' viewBox='0 0 400 225'>"
        "<rect width='400' height='225' fill='#1b202b'/>"
        "<rect x='28' y='28' width='344' height='169' rx='12' fill='#111521' stroke='#2b3345'/>"
        "<g fill='#2b3345'>"
        "<rect x='44' y='54' width='312' height='18' rx='6'/>"
        "<rect x='44' y='84' width='312' height='18' rx='6'/>"
        "<rect x='44' y='114' width='312' height='18' rx='6'/>"
        "<rect x='44' y='144' width='312' height='18' rx='6'/>"
        "</g>"
        "<text x='200' y='205' text-anchor='middle' font-family='Segoe UI, Arial, sans-serif' font-size='14' fill='#8aa0c6'>TABLE</text>"
        "</svg>"
    )
    return "data:image/svg+xml;utf8," + quote(svg)


def _extract_slide_img(slide: Dict[str, Any]) -> str:
    if not isinstance(slide, dict):
        return ""
    keys = ("img", "image", "png", "data_url", "image_data_url", "img_path", "image_path", "png_path")
    for k in keys:
        v = slide.get(k)
        if v:
            normalized = _normalize_img_value(v)
            return normalized or ""
    b64 = slide.get("png_b64") or slide.get("image_b64")
    if b64 and str(b64).strip():
        s = str(b64).strip()
        if s.startswith("data:image/"):
            return s
        return "data:image/png;base64," + s
    return ""

def _load_qwebchannel_js_text() -> str:
    """
    ###1. locate qwebchannel.js inside the PyQt6 installation
    ###2. return file contents as text, or "" if not found
    """
    import sys

    candidates = []

    for p in sys.path:
        if not p:
            continue
        root = Path(p)
        candidates.append(root / "PyQt6" / "Qt6" / "resources" / "qtwebchannel" / "qwebchannel.js")
        candidates.append(root / "PyQt6" / "Qt" / "resources" / "qtwebchannel" / "qwebchannel.js")
        candidates.append(root / "PyQt6" / "QtWebChannel" / "resources" / "qwebchannel.js")

    for fp in candidates:
        if fp.exists():
            return fp.read_text(encoding="utf-8", errors="replace")

    return ""


def _png_bytes_to_data_url(png_bytes: bytes) -> str:
    b = png_bytes or b""
    b64 = base64.b64encode(b).decode("ascii")
    _ppt_dbg(f"[ppt_debug] _png_bytes_to_data_url bytes={len(b)} b64={len(b64)}")
    return "data:image/png;base64," + b64


def _kaleido_available() -> bool:
    return importlib.util.find_spec("kaleido") is not None

def ppt_ui_warmup() -> dict:
    """
    ###1. read and cache template + local JS once
    ###2. invalidate cache if resources folder moved
    """
    global _PPT_UI_WARMED, _PPT_UI_CACHE

    if _PPT_UI_WARMED:
        cached_dir = Path(str(_PPT_UI_CACHE.get("res_dir") or "")).expanduser()
        if cached_dir.exists():
            return dict(_PPT_UI_CACHE)
        _PPT_UI_WARMED = False
        _PPT_UI_CACHE = {}

    res_dir = _resources_dir().resolve()

    template_path = res_dir / "visuals.html"
    template_text = _read_text(template_path) if template_path.exists() else ""

    js_dir = (res_dir / "js").resolve()
    plotly_fp = js_dir / "plotly-2.35.2.min.js"
    qweb_fp = js_dir / "qwebchannel.js"

    plotly_js = _read_text(plotly_fp) if plotly_fp.exists() else ""
    qweb_js = _read_text(qweb_fp) if qweb_fp.exists() else ""

    _PPT_UI_CACHE = {
        "res_dir": str(res_dir),
        "js_dir": str(js_dir),
        "template_path": str(template_path),
        "template_len": int(len(template_text)),
        "plotly_js_len": int(len(plotly_js)),
        "qweb_js_len": int(len(qweb_js)),
        "template_text": template_text,
        "png_cache": {},
    }

    _PPT_UI_WARMED = True
    return dict(_PPT_UI_CACHE)

def _coerce_fig_json(fig_spec) -> dict | None:
    if fig_spec is None:
        return None
    if isinstance(fig_spec, go.Figure):
        return fig_spec.to_plotly_json()
    if isinstance(fig_spec, str):
        s = fig_spec.strip()
        if not s:
            return None
        if s.startswith("{") or s.startswith("["):
            try:
                obj = json.loads(s)
            except json.JSONDecodeError:
                return None
            fig = go.Figure(obj)
            return fig.to_plotly_json()
        return None
    if isinstance(fig_spec, dict):
        return fig_spec
    try:
        fig = go.Figure(fig_spec)
    except Exception:
        return None
    return fig.to_plotly_json()


def _fig_spec_to_thumb_url(fig_spec, *, collection_name: str, fig_md5: str = "") -> str:
    if not _kaleido_available():
        _ppt_dbg("[ppt_debug] kaleido missing; cannot rasterize fig")
        return ""

    fig_json = _coerce_fig_json(fig_spec)
    if fig_json is None:
        return ""

    key = fig_md5 or _fig_json_md5(fig_json)
    if not key:
        return ""

    png_dir = (_ppts_cache_dir(collection_name) / "thumbs").resolve()
    png_dir.mkdir(parents=True, exist_ok=True)
    fp = (png_dir / (key + ".png")).resolve()

    if fp.exists():
        return _path_to_file_url(str(fp))

    fig = go.Figure(fig_json)
    png_bytes = pio.to_image(fig, format="png", width=400, height=225, scale=2)
    fp.write_bytes(png_bytes)
    return _path_to_file_url(str(fp))


def _normalize_slide_for_ui(slide: Dict[str, Any], *, want_png: bool = False, collection_name: str = "") -> Dict[str, Any]:
    """
    ###1. normalize keys to {section,title,notes,fig_json,fig_html,table_html,img}
    ###2. keep fig_json as a JSON object (dict), not a JSON-encoded string
    ###3. prefer existing 'img'; else rasterize fig_json to 'img' when want_png=True
    """
    out: Dict[str, Any] = {}

    sec = (
        slide.get("section")
        or slide.get("section_id")
        or slide.get("section_key")
        or slide.get("key")
        or slide.get("group")
    )
    if sec is not None and str(sec).strip():
        out["section"] = str(sec).strip()

    title = slide.get("title") or slide.get("name") or "Slide"
    out["title"] = str(title)

    for k in ("notes", "note"):
        v = slide.get(k)
        if v is not None and str(v).strip():
            out["notes"] = str(v)
            break

    fig_html = slide.get("fig_html")
    if fig_html is not None and str(fig_html).strip():
        out["fig_html"] = str(fig_html)

    thumb_img = _normalize_img_value(slide.get("thumb_img") or slide.get("thumb") or "")
    if thumb_img:
        out["thumb_img"] = thumb_img

    img = _extract_slide_img(slide)
    if img:
        out["img"] = img

    spec = (
            slide.get("fig_json")
            or slide.get("plotly_json")
            or slide.get("figure_json")
            or slide.get("fig")
            or slide.get("figure")
            or slide.get("plot")
            or slide.get("plotly_fig")
    )
    if spec is not None:
        if isinstance(spec, str):
            s = spec.strip()
            if s.startswith("{") or s.startswith("["):
                out["fig_json"] = json.loads(s)
            else:
                out["fig_json"] = {"data": [], "layout": {"annotations": [{"text": s, "showarrow": False}]}}
        elif isinstance(spec, dict):
            out["fig_json"] = spec
        elif isinstance(spec, go.Figure):
            out["fig_json"] = spec.to_plotly_json()
        else:
            out["fig_json"] = {"data": [], "layout": {"annotations": [{"text": str(spec), "showarrow": False}]}}

    fig_json = out.get("fig_json")
    if fig_json is not None:
        fig_md5 = _fig_json_md5(fig_json)
        if fig_md5:
            out["fig_md5"] = fig_md5
        if not out.get("thumb_img"):
            if img:
                out["thumb_img"] = img
            else:
                out["thumb_img"] = _fig_spec_to_thumb_url(
                    fig_json,
                    collection_name=str(collection_name),
                    fig_md5=fig_md5,
                )

    if "table_html" in slide and slide.get("table_html") is not None:
        out["table_html"] = str(slide.get("table_html"))
        if not out.get("thumb_img") and not out.get("fig_json"):
            out["thumb_img"] = _table_thumb_data_url()

    return _json_safe(out)


def _fig_json_meta(spec) -> dict:
    """
    ###1. summarize Plotly figure JSON shape and serializability
    ###2. compute stable hash + head sample for logs
    ###3. extract first-trace hints and layout keys
    """
    import hashlib

    out = {"present": False}

    if spec is None:
        return out

    out["present"] = True
    out["py_type"] = type(spec).__name__

    if isinstance(spec, str):
        s = spec.strip()
        out["is_str"] = True
        out["str_len"] = len(s)
        out["str_head"] = s[:220]
        if s.startswith("{") or s.startswith("["):
            out["looks_like_json"] = True
        return out

    if not isinstance(spec, dict):
        out["note"] = "non-dict figure spec"
        out["repr_head"] = str(spec)[:220]
        return out

    data = spec.get("data")
    layout = spec.get("layout")
    frames = spec.get("frames")

    out["data_is_list"] = isinstance(data, list)
    out["layout_is_dict"] = isinstance(layout, dict)
    out["frames_is_list"] = isinstance(frames, list)

    out["data_len"] = len(data) if isinstance(data, list) else 0
    out["layout_keys_n"] = len(layout.keys()) if isinstance(layout, dict) else 0
    out["layout_keys_head"] = sorted(list(layout.keys()))[:20] if isinstance(layout, dict) else []

    if isinstance(data, list) and len(data) > 0 and isinstance(data[0], dict):
        t0 = data[0]
        out["trace0_type"] = str(t0.get("type") or "")
        out["trace0_name"] = str(t0.get("name") or "")
        out["trace0_keys_head"] = sorted(list(t0.keys()))[:20]
    else:
        out["trace0_type"] = ""
        out["trace0_name"] = ""

    safe = _json_safe(spec)
    s = json.dumps(safe, ensure_ascii=False, allow_nan=False)
    out["json_len"] = len(s)
    out["json_head"] = s[:260]
    out["json_md5"] = hashlib.md5(s.encode("utf-8")).hexdigest()

    return out

def _ui_slide_debug(slide: dict, i: int) -> dict:
    """
    ###1. summarize normalized slide payload (types, lengths, heads)
    ###2. check fig_json shape and json serializability
    ###3. return stable dict for logs
    """
    s = slide if isinstance(slide, dict) else {}
    out = {"i": int(i), "keys": sorted(list(s.keys()))}

    title = s.get("title")
    out["title"] = str(title)[:120] if title is not None else ""

    sec = s.get("section")
    out["section"] = str(sec)[:80] if sec is not None else ""

    notes = s.get("notes")
    out["notes_len"] = len(str(notes)) if notes is not None else 0

    img = s.get("img")
    img_s = str(img) if isinstance(img, str) else ""
    out["has_img"] = bool(img_s.strip())
    out["img_head"] = img_s[:80]
    out["img_is_data_url"] = img_s.startswith("data:image/")

    tbl = s.get("table_html")
    tbl_s = str(tbl) if tbl is not None else ""
    out["has_table"] = bool(tbl_s.strip())
    out["table_len"] = len(tbl_s)
    out["table_head"] = tbl_s[:120]

    fig_html = s.get("fig_html")
    fh_s = str(fig_html) if fig_html is not None else ""
    out["has_fig_html"] = bool(fh_s.strip())
    out["fig_html_len"] = len(fh_s)
    out["fig_html_head"] = fh_s[:160]

    fig_json = s.get("fig_json")
    out["has_fig_json"] = fig_json is not None
    out["fig_json_type"] = type(fig_json).__name__ if fig_json is not None else "none"

    shape = "none"
    if fig_json is None:
        shape = "none"
    elif isinstance(fig_json, dict):
        ok = isinstance(fig_json.get("data"), list) and isinstance(fig_json.get("layout"), dict)
        shape = "dict_ok" if ok else "dict_bad"
    elif isinstance(fig_json, str):
        shape = "string"
    else:
        shape = "other"
    out["fig_shape"] = shape

    if fig_json is not None:
        out["fig_meta"] = _fig_json_meta(fig_json)

    safe = _json_safe(s)
    dumped = json.dumps(safe, ensure_ascii=False, allow_nan=False)
    out["json_len"] = len(dumped)
    out["json_head"] = dumped[:200]

    return out


def _ui_deck_debug(slides: list[dict]) -> dict:
    """
    ###1. compute deck coverage counts (fig_json/fig_html/img/table/text)
    ###2. track indices with suspect fig_json shapes or missing content
    ###3. return stable dict for logs
    """
    src = slides if isinstance(slides, list) else []
    n = len(src)

    n_fig_json = 0
    n_fig_html = 0
    n_img = 0
    n_table = 0
    n_text = 0

    bad_fig_shape = []
    empty_slide = []

    for i, s in enumerate(src):
        if not isinstance(s, dict):
            empty_slide.append(i)
            continue

        has_any = False

        fj = s.get("fig_json")
        if fj is not None:
            n_fig_json += 1
            has_any = True
            if isinstance(fj, dict):
                ok = isinstance(fj.get("data"), list) and isinstance(fj.get("layout"), dict)
                if not ok:
                    bad_fig_shape.append(i)
            elif not isinstance(fj, str):
                bad_fig_shape.append(i)

        fh = s.get("fig_html")
        if fh is not None and str(fh).strip():
            n_fig_html += 1
            has_any = True

        img = s.get("img")
        if isinstance(img, str) and img.strip():
            n_img += 1
            has_any = True

        th = s.get("table_html")
        if th is not None and str(th).strip():
            n_table += 1
            has_any = True

        notes = s.get("notes")
        if notes is not None and str(notes).strip():
            n_text += 1
            has_any = True

        if not has_any:
            empty_slide.append(i)

    return {
        "slides": n,
        "fig_json": n_fig_json,
        "fig_html": n_fig_html,
        "img": n_img,
        "table": n_table,
        "text": n_text,
        "bad_fig_shape_n": len(bad_fig_shape),
        "bad_fig_shape_head": bad_fig_shape[:12],
        "empty_slide_n": len(empty_slide),
        "empty_slide_head": empty_slide[:12],
    }

def _repoish_root(start: Path) -> Path | None:
    cur = start
    for _ in range(8):
        if (cur / "resources").exists():
            return cur
        if cur.parent == cur:
            break
        cur = cur.parent
    return None


def _resources_dir() -> Path:
    env = (os.environ.get("ANNOTARIUM_RESOURCES_DIR") or "").strip()
    if env:
        p = Path(env)
        if p.exists():
            return p

    here = Path(__file__).resolve()
    root = _repoish_root(here)
    if root is not None:
        return root / "resources"

    return here.parents[1] / "resources"


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _escape_cell(x) -> str:
    return html.escape("" if x is None else str(x))


def _df_to_html_table(df: pd.DataFrame, max_rows: int = 200) -> str:
    if not isinstance(df, pd.DataFrame) or df.empty:
        return ""

    df2 = df if len(df) <= max_rows else df.head(max_rows)
    cols = [str(c) for c in df2.columns]
    thead = "<tr>" + "".join(f"<th>{html.escape(c)}</th>" for c in cols) + "</tr>"

    rows = []
    for tup in df2.itertuples(index=False):
        tds = "".join(f"<td>{_escape_cell(v)}</td>" for v in tup)
        rows.append(f"<tr>{tds}</tr>")

    return f"<table><thead>{thead}</thead><tbody>{''.join(rows)}</tbody></table>"

def _print_ui_debug_dump(payload) -> None:
    """
    ###1. validate JS dump shape and print root-level context
    ###2. compute distributions across slide summaries (fig/img/table/text, fig_shape)
    ###3. infer likely failure modes (including script-not-run and zero-width Plotly DOM)
    ###4. print focused details for current slide and first few problematic slides
    """
    if not isinstance(payload, dict):
        print("[PptxExportWidget.ui_debug] no-payload", flush=True)
        return

    slides = payload.get("slides")
    if not isinstance(slides, list):
        slides = []

    deck_slides = len(slides)
    slides_len = payload.get("slides_len")
    index = payload.get("index")
    title = payload.get("title")
    current_keys = payload.get("current_keys")

    render_last = payload.get("render_last")

    has_fig_json = bool(payload.get("has_fig_json"))
    fig_type = payload.get("fig_json_type")

    fig_head = payload.get("fig_json_head")
    img_head = payload.get("img_head")

    fig_html_len = payload.get("fig_html_len")
    fig_html_head = payload.get("fig_html_head")
    fig_html_has_plotly_call = payload.get("fig_html_has_plotly_call")

    plot_host_rect = payload.get("plot_host_rect")
    gd_rect = payload.get("gd_rect")
    js_plot_rect = payload.get("js_plot_rect")
    js_plot_found = payload.get("js_plot_found")

    root_keys = sorted(list(payload.keys()))

    print(
        "[PptxExportWidget.ui_debug] "
        f"slides_len={slides_len} deck_slides={deck_slides} index={index} title={repr(title)}",
        flush=True,
    )
    print(
        "[PptxExportWidget.ui_debug] "
        f"root_keys_n={len(root_keys)} current_keys={current_keys} render_last={repr(render_last)}",
        flush=True,
    )
    print(
        "[PptxExportWidget.ui_debug] "
        f"current_has_fig_json={has_fig_json} fig_json_type={repr(fig_type)}",
        flush=True,
    )

    if isinstance(fig_html_len, int) and fig_html_len > 0:
        print(f"[PptxExportWidget.ui_debug] fig_html_len={fig_html_len}", flush=True)
    if fig_html_has_plotly_call is not None:
        print(f"[PptxExportWidget.ui_debug] fig_html_has_plotly_call={bool(fig_html_has_plotly_call)}", flush=True)
    if isinstance(fig_html_head, str) and fig_html_head.strip():
        print(f"[PptxExportWidget.ui_debug] fig_html_head={repr(fig_html_head[:520])}", flush=True)

    if isinstance(plot_host_rect, dict):
        print(f"[PptxExportWidget.ui_debug] plot_host_rect={json.dumps(plot_host_rect, ensure_ascii=False)}", flush=True)
    if isinstance(gd_rect, dict):
        print(f"[PptxExportWidget.ui_debug] gd_rect={json.dumps(gd_rect, ensure_ascii=False)}", flush=True)
    if js_plot_found is not None:
        print(f"[PptxExportWidget.ui_debug] js_plot_found={bool(js_plot_found)}", flush=True)
    if isinstance(js_plot_rect, dict):
        print(f"[PptxExportWidget.ui_debug] js_plot_rect={json.dumps(js_plot_rect, ensure_ascii=False)}", flush=True)

    if isinstance(fig_head, str) and fig_head.strip():
        print(f"[PptxExportWidget.ui_debug] fig_json_head={repr(fig_head[:260])}", flush=True)
    if isinstance(img_head, str) and img_head.strip():
        print(f"[PptxExportWidget.ui_debug] img_head={repr(img_head[:260])}", flush=True)

    n_has_fig = 0
    n_has_img = 0
    n_has_table = 0
    n_has_text = 0

    fig_shape_counts = {"none": 0, "string": 0, "dict_ok": 0, "dict_bad": 0, "other": 0}

    idx_missing_fig = []
    idx_bad_fig = []
    idx_string_fig = []
    idx_missing_any = []

    for s in slides:
        if not isinstance(s, dict):
            continue

        i = s.get("i")
        i_ok = isinstance(i, int)

        has_fig = bool(s.get("has_fig"))
        has_img = bool(s.get("has_img"))
        has_table = bool(s.get("has_table"))
        has_text = bool(s.get("has_text"))

        if has_fig:
            n_has_fig += 1
        if has_img:
            n_has_img += 1
        if has_table:
            n_has_table += 1
        if has_text:
            n_has_text += 1

        fig_shape = s.get("fig_shape")
        if fig_shape in fig_shape_counts:
            fig_shape_counts[fig_shape] += 1
        else:
            fig_shape_counts["other"] += 1

        if i_ok:
            if not (has_fig or has_img or has_table or has_text):
                idx_missing_any.append(i)

            if not has_fig:
                idx_missing_fig.append(i)
            else:
                if fig_shape == "dict_bad":
                    idx_bad_fig.append(i)
                if fig_shape == "string":
                    idx_string_fig.append(i)

    print(
        "[PptxExportWidget.ui_debug] "
        f"deck_counts has_fig={n_has_fig} has_img={n_has_img} has_table={n_has_table} has_text={n_has_text}",
        flush=True,
    )
    print(
        "[PptxExportWidget.ui_debug] "
        f"fig_shape_counts={json.dumps(fig_shape_counts, ensure_ascii=False)}",
        flush=True,
    )

    total = deck_slides if deck_slides > 0 else 1
    print(
        "[PptxExportWidget.ui_debug] "
        f"fig_coverage={n_has_fig}/{total} bad_fig={len(idx_bad_fig)} fig_as_string={len(idx_string_fig)} empty_slides={len(idx_missing_any)}",
        flush=True,
    )

    if idx_bad_fig:
        print(f"[PptxExportWidget.ui_debug] bad_fig_indices_head={idx_bad_fig[:12]}", flush=True)
    if idx_string_fig:
        print(f"[PptxExportWidget.ui_debug] fig_string_indices_head={idx_string_fig[:12]}", flush=True)
    if idx_missing_any:
        print(f"[PptxExportWidget.ui_debug] empty_slide_indices_head={idx_missing_any[:12]}", flush=True)

    def _rect_w(rect):
        if not isinstance(rect, dict):
            return None
        w = rect.get("w")
        return w if isinstance(w, (int, float)) else None

    inferred = []

    if deck_slides == 0:
        inferred.append("no_slides_received_by_js")
    if has_fig_json and (fig_type is None or str(fig_type).strip() == ""):
        inferred.append("current_fig_json_present_but_type_missing")
    if n_has_fig == 0 and n_has_img == 0 and n_has_table == 0 and n_has_text == 0 and deck_slides > 0:
        inferred.append("js_summary_flags_all_false_possible_schema_mismatch")
    if len(idx_string_fig) > 0:
        inferred.append("fig_json_sent_as_string_expect_object_dict_for_plotly")
    if len(idx_bad_fig) > 0:
        inferred.append("fig_json_missing_data_or_layout_expected_keys")

    if isinstance(fig_html_len, int) and fig_html_len > 0 and js_plot_found is False:
        inferred.append("fig_html_present_but_js_plot_node_missing_script_not_run_or_target_mismatch")

    ph_w = _rect_w(plot_host_rect)
    jp_w = _rect_w(js_plot_rect)
    gd_w = _rect_w(gd_rect)
    if isinstance(fig_html_len, int) and fig_html_len > 0 and isinstance(ph_w, (int, float)):
        if ph_w > 0 and ((isinstance(jp_w, (int, float)) and jp_w == 0) or (isinstance(gd_w, (int, float)) and gd_w == 0)):
            inferred.append("plotly_dom_zero_width_css_flex_sizing_issue")

    if inferred:
        print(
            "[PptxExportWidget.ui_debug] inferred_issues="
            + json.dumps(inferred, ensure_ascii=False),
            flush=True,
        )

    by_i = {}
    for s in slides:
        if isinstance(s, dict) and isinstance(s.get("i"), int):
            by_i[int(s.get("i"))] = s

    focus = []
    if isinstance(index, int):
        focus.append(index)
    for i in idx_bad_fig[:4]:
        focus.append(i)
    for i in idx_string_fig[:4]:
        focus.append(i)
    for i in idx_missing_any[:4]:
        focus.append(i)

    seen = set()
    ordered = []
    for i in focus:
        if not isinstance(i, int):
            continue
        if i in seen:
            continue
        seen.add(i)
        ordered.append(i)
        if len(ordered) >= 6:
            break

    for i in ordered:
        s = by_i.get(i)
        if not isinstance(s, dict):
            continue

        item_keys = sorted(list(s.keys()))
        fig_shape = s.get("fig_shape")
        has_fig = bool(s.get("has_fig"))
        has_img = bool(s.get("has_img"))
        has_table = bool(s.get("has_table"))
        has_text = bool(s.get("has_text"))

        sec = s.get("section")
        stitle = s.get("title")

        fig_head_s = s.get("fig_head")
        img_head_s = s.get("img_head")

        fig_html_len_s = s.get("fig_html_len")
        fig_html_head_s = s.get("fig_html_head")
        fig_html_has_plotly_call_s = s.get("fig_html_has_plotly_call")

        js_plot_found_s = s.get("js_plot_found")
        js_plot_rect_s = s.get("js_plot_rect")
        plot_host_rect_s = s.get("plot_host_rect")
        gd_rect_s = s.get("gd_rect")

        line = {
            "i": i,
            "keys_n": len(item_keys),
            "has_fig": has_fig,
            "fig_shape": fig_shape,
            "has_img": has_img,
            "has_table": has_table,
            "has_text": has_text,
            "section": sec,
            "title": stitle,
        }
        print("[PptxExportWidget.ui_debug.slide] " + json.dumps(_json_safe(line), ensure_ascii=False), flush=True)

        if isinstance(fig_html_len_s, int) and fig_html_len_s > 0:
            print(f"[PptxExportWidget.ui_debug.slide] i={i} fig_html_len={fig_html_len_s}", flush=True)
        if fig_html_has_plotly_call_s is not None:
            print(
                f"[PptxExportWidget.ui_debug.slide] i={i} fig_html_has_plotly_call={bool(fig_html_has_plotly_call_s)}",
                flush=True,
            )
        if isinstance(fig_html_head_s, str) and fig_html_head_s.strip():
            print(f"[PptxExportWidget.ui_debug.slide] i={i} fig_html_head={repr(fig_html_head_s[:520])}", flush=True)

        if js_plot_found_s is not None:
            print(f"[PptxExportWidget.ui_debug.slide] i={i} js_plot_found={bool(js_plot_found_s)}", flush=True)
        if isinstance(plot_host_rect_s, dict):
            print(
                f"[PptxExportWidget.ui_debug.slide] i={i} plot_host_rect={json.dumps(plot_host_rect_s, ensure_ascii=False)}",
                flush=True,
            )
        if isinstance(gd_rect_s, dict):
            print(f"[PptxExportWidget.ui_debug.slide] i={i} gd_rect={json.dumps(gd_rect_s, ensure_ascii=False)}", flush=True)
        if isinstance(js_plot_rect_s, dict):
            print(
                f"[PptxExportWidget.ui_debug.slide] i={i} js_plot_rect={json.dumps(js_plot_rect_s, ensure_ascii=False)}",
                flush=True,
            )

        if isinstance(fig_head_s, str) and fig_head_s.strip():
            print(f"[PptxExportWidget.ui_debug.slide] i={i} fig_head={repr(fig_head_s[:220])}", flush=True)
        if isinstance(img_head_s, str) and img_head_s.strip():
            print(f"[PptxExportWidget.ui_debug.slide] i={i} img_head={repr(img_head_s[:220])}", flush=True)

def _insert_into_host_div(template_html: str, host_id: str, inner_html: str) -> str:
    pat = re.compile(
        rf'(<[^>]+\bid="{re.escape(host_id)}"[^>]*>)(.*?)(</[^>]+>)',
        flags=re.IGNORECASE | re.DOTALL,
    )
    m = pat.search(template_html)
    if not m:
        return template_html
    return template_html[: m.start()] + m.group(1) + inner_html + m.group(3) + template_html[m.end() :]


def _build_options_html(schema: list[dict], params: dict) -> str:
    if not isinstance(schema, list) or not schema:
        return "<div class='muted'>No inputs.</div>"

    out = []
    for fld in schema:
        if not isinstance(fld, dict):
            continue

        typ = (fld.get("type") or "").strip().lower()
        key = (fld.get("key") or "").strip()
        label = (fld.get("label") or key).strip()
        if not key:
            continue

        cur = params.get(key, fld.get("default", ""))

        if typ == "select":
            opts = fld.get("options") or []
            rows = []
            for opt in opts:
                if not isinstance(opt, dict):
                    continue
                v = "" if opt.get("value") is None else str(opt.get("value"))
                t = str(opt.get("label", v))
                sel = " selected" if str(cur) == v else ""
                rows.append(f"<option value='{html.escape(v)}'{sel}>{html.escape(t)}</option>")

            out.append(
                "<div class='field'>"
                f"<label class='label'>{html.escape(label)}</label>"
                f"<select data-param='{html.escape(key)}' data-type='str'>"
                f"{''.join(rows)}"
                "</select>"
                "</div>"
            )
            continue

        if typ == "number":
            if isinstance(cur, (int, float)):
                v = int(cur)
            else:
                s = str(cur).strip()
                v = int(s) if s.isdigit() else int(fld.get("default", 0) or 0)

            mn = fld.get("min")
            mx = fld.get("max")
            mn_attr = f" min='{int(mn)}'" if isinstance(mn, (int, float)) else ""
            mx_attr = f" max='{int(mx)}'" if isinstance(mx, (int, float)) else ""
            out.append(
                "<div class='field'>"
                f"<label class='label'>{html.escape(label)}</label>"
                f"<input type='number' value='{v}' data-param='{html.escape(key)}' data-type='int'{mn_attr}{mx_attr}/>"
                "</div>"
            )
            continue

        if typ in ("textarea", "keywords"):
            v = "" if cur is None else str(cur)
            out.append(
                "<div class='field'>"
                f"<label class='label'>{html.escape(label)}</label>"
                f"<textarea data-param='{html.escape(key)}' data-type='csv_keywords'>{html.escape(v)}</textarea>"
                "</div>"
            )
            continue

        if typ in ("bool", "checkbox"):
            checked = " checked" if bool(cur) else ""
            hint = str(fld.get("hint") or "").strip()
            out.append(
                "<label class='check'>"
                f"<input type='checkbox' data-param='{html.escape(key)}' data-type='bool'{checked}/>"
                "<div class='meta'>"
                f"<div class='name'>{html.escape(label)}</div>"
                + (f"<div class='hint'>{html.escape(hint)}</div>" if hint else "")
                + "</div></label>"
            )
            continue

        v = "" if cur is None else str(cur)
        dtype = "opt" if typ == "text_opt" else "str"
        out.append(
            "<div class='field'>"
            f"<label class='label'>{html.escape(label)}</label>"
            f"<input type='text' value='{html.escape(v)}' data-param='{html.escape(key)}' data-type='{dtype}'/>"
            "</div>"
        )

    return "".join(out)


class _Bridge(QObject):
    pptx = pyqtSignal(dict)
    params_changed = pyqtSignal(dict)

    def _to_dict(self, payload):
        if isinstance(payload, dict):
            return dict(payload)

        if isinstance(payload, str):
            s = payload.strip()
            if not s:
                return {}
            obj = json.loads(s)
            return obj if isinstance(obj, dict) else {}

        d = dict(payload)
        return d if isinstance(d, dict) else {}

    @pyqtSlot(str, "QVariant")
    def emit(self, action: str, payload):
        a = (action or "").strip().lower()
        p = self._to_dict(payload)
        if a == "pptx":
            self.pptx.emit(p)
            return
        if a == "params":
            self.params_changed.emit(p)
            return

    @pyqtSlot(str)
    def params(self, payload_json: str):
        p = self._to_dict(payload_json)
        self.params_changed.emit(p)


class PptxExportWidget(QWidget):
    """
    ###1. render the HTML template (plot/table + right settings)
    ###2. wire QWebChannel bridge so JS -> Python signals work
    ###3. expose set_slide_deck / set_status for MainWindow updates
    """

    pptx_requested = pyqtSignal(dict)

    def __init__(self, *, schema: list[dict] | None = None, sections: list[dict] | None = None, collection_name: str = "", parent=None):
        super().__init__(parent)
        from PyQt6.QtWebEngineCore import QWebEngineSettings

        self._schema = schema or []
        self._sections = sections or []
        self._collection_name = str(collection_name or "").strip()

        if not self._sections:
            self._sections = [
                {
                    "id": "all",
                    "title": "All slides",
                    "hint": "Default section (fallback when no sections are declared).",
                    "default_checked": True,
                }
            ]

        self._params: dict = {}

        self._slide_deck: list[dict] = []
        self._deck_index: int = 0
        self._last_deck_hash: str = ""

        self._page_ready = False
        self._pending_slide_payload: dict | None = None
        self._push_attempts_left = 0


        template_path = PPT_HTML

        self._template_path = template_path
        self._template_text = _read_text(self._template_path)

        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)

        self.view = QWebEngineView()
        self.view.setZoomFactor(1.0)
        s = self.view.settings()

        self.view.setFocusPolicy(Qt.FocusPolicy.StrongFocus)
        self.view.setAttribute(Qt.WidgetAttribute.WA_Hover, True)
        self.view.setAttribute(Qt.WidgetAttribute.WA_AcceptTouchEvents, True)
        self.view.setMouseTracking(True)
        for w in self.view.findChildren(QWidget):
            w.setMouseTracking(True)
            w.setAttribute(Qt.WidgetAttribute.WA_Hover, True)

        s = self.view.settings()
        s.setAttribute(QWebEngineSettings.WebAttribute.JavascriptEnabled, True)
        s.setAttribute(QWebEngineSettings.WebAttribute.LocalContentCanAccessFileUrls, True)
        s.setAttribute(QWebEngineSettings.WebAttribute.LocalContentCanAccessRemoteUrls, True)
        s.setAttribute(QWebEngineSettings.WebAttribute.FullScreenSupportEnabled, True)

        layout.addWidget(self.view, 1)

        self._bridge = _Bridge()
        self._bridge.pptx.connect(self._on_pptx_from_js)
        self._bridge.params_changed.connect(self._on_params_from_js)

        self._channel = QWebChannel(self.view.page())
        self._channel.registerObject("bridge", self._bridge)
        self.view.page().setWebChannel(self._channel)

        self.view.loadFinished.connect(self._on_load_finished)

        self._load_page()

    def _dbg_enabled(self) -> bool:
        v = (os.environ.get("PPT_DEBUG") or "").strip().lower()
        return v in {"1", "true", "yes", "y", "on"}

    def _dbg(self, msg: str) -> None:
        if self._dbg_enabled():
            print(str(msg), flush=True)

    def _base_url(self) -> QUrl:
        base_dir = self._template_path.parent if hasattr(self, "_template_path") else PPT_HTML.parent
        return QUrl.fromLocalFile(str(base_dir.resolve()) + os.sep)

    def _render_template(self) -> str:
        """
        ###1. replace COMMON_MONITOR_CSS + OPTIONS_HTML placeholders
        ###2. inject window.__PPTX_SECTIONS
        """
        t = self._template_text or ""

        options_html = _build_options_html(self._schema, self._params)
        t = t.replace("{{OPTIONS_HTML}}", options_html)
        t = t.replace("{{COMMON_MONITOR_CSS}}", "")

        sections_js = json.dumps(self._sections or [], ensure_ascii=False)
        inject_sections = "<script>window.__PPTX_SECTIONS = " + sections_js + ";</script>\n"
        inject_defer = "<script>window.__DEFER_BRIDGE_INIT = true;</script>\n"

        if "</head>" in t:
            t = t.replace("</head>", inject_sections + inject_defer + "</head>")
        else:
            t = inject_sections + inject_defer + t

        return t

    def _load_page(self) -> None:
        page_html = self._render_template()
        self._page_ready = False
        self.view.setHtml(page_html, self._base_url())

    def _js_ready_state(self, cb) -> None:
        js = (
            "(function(){"
            "  var rs = (document && document.readyState) ? document.readyState : '';"
            "  return {"
            "    readyState: rs,"
            "    hasSetter: (typeof window.__setSlidesFromPy === 'function'),"
            "    hasRenderSlide: (typeof window.renderSlide === 'function'),"
            "    hasRenderThumbs: (typeof window.renderThumbnails === 'function'),"
            "    hasUpdateSummary: (typeof window.updateSummary === 'function')"
            "  };"
            "})()"
        )
        self.view.page().runJavaScript(js, cb)

    def _push_slide_payload_to_js(self, payload: dict) -> None:
        payload_json = json.dumps(payload, ensure_ascii=False, allow_nan=False)

        js = (
                "(function(){"
                "  var p = " + payload_json + ";"
                                              "  var slides = Array.isArray(p && p.slides) ? p.slides : [];"
                                              "  var idx = parseInt((p && p.index) ? p.index : 0, 10);"
                                              "  if (isNaN(idx) || idx < 0) idx = 0;"
                                              "  if (idx >= slides.length) idx = Math.max(0, slides.length - 1);"
                                              ""
                                              "  window.__slides = slides;"
                                              "  window.__slideIndex = idx;"
                                              ""
                                              "  if (typeof window.__setSlidesFromPy === 'function'){"
                                              "    window.__setSlidesFromPy({slides: slides, index: idx});"
                                              "  } else {"
                                              "    if (typeof window.renderSlide === 'function') window.renderSlide();"
                                              "    if (typeof window.renderThumbnails === 'function') window.renderThumbnails();"
                                              "    if (typeof window.updateSummary === 'function') window.updateSummary();"
                                              "  }"
                                              ""
                                              "  return {"
                                              "    ok: true,"
                                              "    hasSetter: (typeof window.__setSlidesFromPy === 'function'),"
                                              "    slides_len: (window.__slides && Array.isArray(window.__slides)) ? window.__slides.length : 0,"
                                              "    index: window.__slideIndex"
                                              "  };"
                                              "})()"
        )

        def _cb(res):
            self._dbg("[ppt_debug] push result=" + json.dumps(_json_safe(res), ensure_ascii=False))

            ok = isinstance(res, dict) and bool(res.get("ok"))
            slides_len = int(res.get("slides_len") or 0) if isinstance(res, dict) else 0

            if not ok or slides_len == 0:
                self._ensure_js_ready_and_push()

        self.view.page().runJavaScript(js, _cb)

        dbg_js = "(window.__debugDumpSlidesFromPy && window.__debugDumpSlidesFromPy()) || null;"
        self.view.page().runJavaScript(dbg_js, _print_ui_debug_dump)

    def _ensure_js_ready_and_push(self) -> None:
        from PyQt6.QtCore import QTimer

        payload = self._pending_slide_payload if isinstance(self._pending_slide_payload, dict) else None
        if not payload:
            return
        if not self._page_ready:
            return
        if self._push_attempts_left <= 0:
            self._dbg("[ppt_debug] push attempts exhausted; giving up")
            return

        self._push_attempts_left -= 1

        def _cb(state):
            d = state if isinstance(state, dict) else {}
            rs = str(d.get("readyState") or "")
            has_setter = bool(d.get("hasSetter"))
            has_fallback = bool(d.get("hasRenderSlide"))

            if rs in {"interactive", "complete"} and (has_setter or has_fallback):
                self._push_slide_payload_to_js(payload)
                return

            QTimer.singleShot(60, self._ensure_js_ready_and_push)

        self._js_ready_state(_cb)

    def _push_deck_to_js(self) -> None:
        payload = _json_safe({"slides": self._slide_deck, "index": self._deck_index})
        self._pending_slide_payload = payload

        if not self._page_ready:
            self._dbg("[ppt_debug] push_deck skipped; page_ready=False (payload cached)")
            return

        self._push_attempts_left = 40
        self._ensure_js_ready_and_push()

    def _push_index_to_js(self) -> None:
        if not self._page_ready:
            return

        idx = int(self._deck_index) if isinstance(self._deck_index, int) else 0
        js = (
            "(function(){"
            "  var idx = " + json.dumps(idx) + ";"
            "  var slides = (window.__slides && Array.isArray(window.__slides)) ? window.__slides : [];"
            "  var n = slides.length || 0;"
            "  if (n <= 0) return {ok:false, slides_len:n};"
            "  if (idx < 0) idx = 0;"
            "  if (idx >= n) idx = Math.max(0, n - 1);"
            "  window.__slideIndex = idx;"
            "  if (typeof window.renderSlide === 'function') window.renderSlide();"
            "  if (typeof window.renderThumbnails === 'function') window.renderThumbnails();"
            "  if (typeof window.updateSummary === 'function') window.updateSummary();"
            "  return {ok:true, slides_len:n, index: idx};"
            "})()"
        )

        def _cb(res):
            ok = isinstance(res, dict) and bool(res.get("ok"))
            slides_len = int(res.get("slides_len") or 0) if isinstance(res, dict) else 0
            if not ok or slides_len <= 0:
                self._push_deck_to_js()

        self.view.page().runJavaScript(js, _cb)

    def _on_load_finished(self, ok: bool) -> None:
        self._page_ready = bool(ok)
        if not ok:
            self.set_status("Error", "UI failed to load.")
            return

        self.view.page().runJavaScript(
            "if (typeof window.renderIncludeList==='function'){ window.renderIncludeList(); }"
        )

        self.set_status("Ready", "UI loaded.")

        self.view.page().runJavaScript(
            "(function(){"
            "  var hasPlotly = (typeof window.Plotly !== 'undefined');"
            "  if (window.__setPptxStatusFromPy){"
            "    window.__setPptxStatusFromPy({short:'Ready', detail:(hasPlotly ? 'UI loaded; Plotly loaded.' : 'UI loaded; Plotly missing.')});"
            "  }"
            "  if (typeof window.initWebChannel === 'function'){"
            "    window.initWebChannel();"
            "    return {called:true, plotly:hasPlotly};"
            "  }"
            "  return {called:false, plotly:hasPlotly};"
            "})()"
        )

        self._push_deck_to_js()

    def _on_pptx_from_js(self, payload: dict) -> None:
        if not isinstance(payload, dict):
            return

        mode = str(payload.get("mode") or "")
        include = payload.get("include")
        section = str(payload.get("section") or "")
        n_include = len(include) if isinstance(include, list) else 0

        self.set_status("Received", f"mode={mode} section={section} include={n_include}")
        self.pptx_requested.emit(payload)

    def _on_params_from_js(self, payload: dict) -> None:
        if not isinstance(payload, dict):
            return
        self._params.update(payload)

    def set_status(self, short: str, detail: str = "") -> None:
        if not self._page_ready:
            return
        p = {"short": str(short or ""), "detail": str(detail or "")}
        js = f"window.__setPptxStatusFromPy({json.dumps(p, ensure_ascii=False, allow_nan=False)});"
        self.view.page().runJavaScript(js)

    def set_slide_deck(self, slides: list[dict], index: int = 0, *, want_png: bool = False) -> None:
        raw = slides if isinstance(slides, list) else []
        normalized = [
            _normalize_slide_for_ui(s, want_png=True, collection_name=self._collection_name)
            for s in raw
        ]

        deck_hash = _deck_md5(normalized)
        same_deck = bool(deck_hash) and (deck_hash == self._last_deck_hash)
        self._last_deck_hash = deck_hash
        self._slide_deck = normalized

        self._deck_index = int(index) if isinstance(index, int) else 0

        if self._deck_index < 0:
            self._deck_index = 0
        if self._slide_deck and self._deck_index >= len(self._slide_deck):
            self._deck_index = len(self._slide_deck) - 1

        n = len(self._slide_deck)

        if self._dbg_enabled():
            deck_dbg = _ui_deck_debug(self._slide_deck)
            print(
                "[PptxExportWidget.set_slide_deck] "
                + json.dumps(_json_safe(deck_dbg), ensure_ascii=False),
                flush=True,
            )

            idxs = []
            bad_head = deck_dbg.get("bad_fig_shape_head")
            empty_head = deck_dbg.get("empty_slide_head")
            if isinstance(bad_head, list):
                idxs.extend([int(x) for x in bad_head if isinstance(x, int)])
            if isinstance(empty_head, list):
                idxs.extend([int(x) for x in empty_head if isinstance(x, int)])

            seen = set()
            pick = []
            for x in idxs:
                if x in seen:
                    continue
                if x < 0 or x >= n:
                    continue
                seen.add(x)
                pick.append(x)
                if len(pick) >= 6:
                    break

            for i in pick:
                d = _ui_slide_debug(self._slide_deck[i], i)
                print(
                    "[PptxExportWidget.set_slide_deck.slide_debug] "
                    + json.dumps(_json_safe(d), ensure_ascii=False),
                    flush=True,
                )

            if 0 <= self._deck_index < n:
                cur = _ui_slide_debug(self._slide_deck[self._deck_index], self._deck_index)
                print(
                    "[PptxExportWidget.set_slide_deck.current_debug] "
                    + json.dumps(_json_safe(cur), ensure_ascii=False),
                    flush=True,
                )

        if same_deck:
            self._pending_slide_payload = _json_safe({"slides": self._slide_deck, "index": self._deck_index})
            self._push_index_to_js()
        else:
            self._push_deck_to_js()

    def refresh_settings(self) -> None:
        self._load_page()

    def _build_ppt_page_ui(self):
        """
        ###1. build the PPT QWidget
        ###2. consume self._ppt_preloaded_data for preview
        ###3. return the QWidget
        """
        import re
        from pathlib import Path

        import pandas as pd
        from PyQt6.QtCore import QTimer
        from PyQt6.QtWidgets import QWidget, QVBoxLayout, QFileDialog
        from pptx import Presentation
        from pptx.util import Inches

        from .power_point_export import PptxExportWidget

        def _bool_param(p: dict, key: str, default: bool = False) -> bool:
            v = p.get(key) if isinstance(p, dict) else None
            if v is None:
                return default
            return str(v).strip().lower() in {"1", "true", "yes", "y", "on"}

        schema = [
            {"type": "number", "key": "top_n_authors", "label": "Top N authors", "default": 15, "min": 5, "max": 100},
            {"type": "number", "key": "production_top_n", "label": "Production top N", "default": 10, "min": 3,
             "max": 50},
            {"type": "number", "key": "top_ngram", "label": "Top n-gram", "default": 20, "min": 5, "max": 200},
            {
                "type": "select",
                "key": "slide_notes",
                "label": "slide_notes",
                "default": "false",
                "options": [{"label": "false", "value": "false"}, {"label": "true", "value": "true"}],
            },
        ]

        sections = [
            {"id": "Data_summary", "label": "Data summary", "hint": "Summary cards and key totals."},
            {"id": "Scope_and_shape", "label": "Scope and shape", "hint": "Volume, years, document types."},
            {"id": "Authors_overview", "label": "Authors overview", "hint": "Top authors and collaboration."},
            {"id": "Citations_overview", "label": "Citations overview", "hint": "Citation distribution and leaders."},
            {"id": "Words_and_topics", "label": "Words and topics", "hint": "Keywords, n-grams, topic signals."},
            {"id": "Affiliations_geo", "label": "Affiliations (geo)", "hint": "Institutions and geography."},
            {"id": "Temporal_analysis", "label": "Temporal analysis", "hint": "Trends over time."},
            {"id": "Research_design", "label": "Research design", "hint": "Design outputs and mix."},
            {"id": "Profiles", "label": "Profiles", "hint": "Feature profiles and summaries."},
            {"id": "Categorical_keywords", "label": "Categorical keywords", "hint": "Categorical keyword outputs."},
        ]

        w = QWidget()
        lay = QVBoxLayout(w)
        lay.setContentsMargins(0, 0, 0, 0)

        export = PptxExportWidget(schema=schema, sections=sections,collection_name="collection")
        lay.addWidget(export, 1)

        self._ppt_widget = export
        self._ppt_page = w

        def _preview_from_preloaded(section_id: str) -> list[dict]:
            from .power_point_export import _normalize_slide_for_ui

            ppt_data = self._ppt_preloaded_data or {}
            payload = ppt_data.get(section_id)
            print("[visualise][ppt_data][preview]", {"section": section_id, "has_payload": isinstance(payload, dict)}, flush=True)

            if not isinstance(payload, dict):
                return [{"title": section_id.replace("_", " "), "bullets": ["No preview data available."], "notes": ""}]

            candidates: list = []

            slides_list = payload.get("slides")
            if isinstance(slides_list, list) and slides_list:
                candidates.extend(slides_list)

            figs_list = payload.get("figs")
            if isinstance(figs_list, list) and figs_list:
                for i, fig in enumerate(figs_list):
                    candidates.append(
                        {
                            "title": section_id.replace("_",
                                                        " ") if i == 0 else f"{section_id.replace('_', ' ')} ({i + 1})",
                            "fig": fig,
                            "notes": "",
                        }
                    )

            if payload.get("fig") is not None or payload.get("plotly_json") is not None or payload.get(
                    "fig_json") is not None:
                candidates.append(payload)

            slides: list[dict] = []
            sec = str(section_id or "").strip()

            for item in candidates:
                if isinstance(item, dict):
                    already_norm = (
                            ("img" in item and str(item.get("img") or "").strip())
                            or ("table_html" in item and str(item.get("table_html") or "").strip())
                    )
                    s = dict(item) if already_norm else _normalize_slide_for_ui(item, want_png=False)
                else:
                    s = _normalize_slide_for_ui(
                        {"title": section_id.replace("_", " "), "fig": item, "notes": ""},
                        want_png=False,
                    )

                if isinstance(s, dict):
                    if not str(s.get("section") or "").strip():
                        s["section"] = sec
                    slides.append(s)

            if not slides:
                slides.append({"title": section_id.replace("_", " "), "bullets": ["No figures available for preview."],
                               "notes": ""})

            return slides

        def _run_section(section_id: str, params: dict) -> list[dict]:
            from pptx import Presentation

            df = self.loaded_dataframe
            if type(df) is not pd.DataFrame or df.empty:
                return [{"title": section_id.replace("_", " "), "bullets": ["No preview data available."], "notes": "", "section": section_id}]

            sec = str(section_id or "").strip()
            if not sec:
                return [{"title": "Preview", "bullets": ["No section selected."], "notes": ""}]

            if type(params) is not dict:
                params = {}

            slides: list[dict] = []

            if sec == "Data_summary":
                prs_tmp = Presentation()
                res_ds = add_data_summary_slides(
                    prs=prs_tmp,
                    df=df,
                    collection_name="Collection",
                    slide_notes=False,
                    return_payload=True,
                    export=False,
                )
                payload = res_ds if type(res_ds) is dict else {"slides": []}
                slides_list = payload.get("slides")
                if type(slides_list) is list:
                    for s in slides_list:
                        if type(s) is dict:
                            d = dict(s)
                            if not str(d.get("section") or "").strip():
                                d["section"] = sec
                            slides.append(d)
                return slides if slides else [{"title": sec.replace("_", " "), "bullets": ["No preview data available."], "notes": "", "section": sec}]

            if sec == "Scope_and_shape":
                prs_tmp2 = Presentation()
                res_sc = shape_scope(
                    prs=prs_tmp2,
                    df=df,
                    collection_name="Collection",
                    slide_notes=False,
                    export=False,
                    return_payload=True,
                )
                payload = res_sc if type(res_sc) is dict else {"slides": []}
                slides_list = payload.get("slides")
                if type(slides_list) is list:
                    for s in slides_list:
                        if type(s) is dict:
                            d = dict(s)
                            if not str(d.get("section") or "").strip():
                                d["section"] = sec
                            slides.append(d)
                return slides if slides else [{"title": sec.replace("_", " "), "bullets": ["No preview data available."], "notes": "", "section": sec}]

            if sec == "Authors_overview":
                prs_tmp3 = Presentation()
                res_au = authors_overview(
                    prs=prs_tmp3,
                    df=df,
                    collection_name="Collection",
                    slide_notes=False,
                    export=True,
                    return_payload=True,
                )
                payload = res_au if type(res_au) is dict else {"slides": []}
                slides_list = payload.get("slides")
                if type(slides_list) is list:
                    for s in slides_list:
                        if type(s) is dict:
                            d = dict(s)
                            if not str(d.get("section") or "").strip():
                                d["section"] = sec
                            slides.append(d)
                return slides if slides else [{"title": sec.replace("_", " "), "bullets": ["No preview data available."], "notes": "", "section": sec}]

            if sec == "Institutions_overview":
                prs_tmp4 = Presentation()
                res_inst = add_institution_section(
                    prs=prs_tmp4,
                    df=df,
                    collection_name="Collection",
                    slide_notes=False,
                    export=True,
                    return_payload=True,
                )
                payload = res_inst if type(res_inst) is dict else {"slides": []}
                slides_list = payload.get("slides")
                if type(slides_list) is list:
                    for s in slides_list:
                        if type(s) is dict:
                            d = dict(s)
                            if not str(d.get("section") or "").strip():
                                d["section"] = sec
                            slides.append(d)
                return slides if slides else [{"title": sec.replace("_", " "), "bullets": ["No preview data available."], "notes": "", "section": sec}]

            plan = self._ppt_section_plan()
            dispatch = self._analysis_dispatch_table()
            jobs = plan.get(sec, [])

            for analysis_id, job_params, title in jobs:
                if analysis_id not in dispatch:
                    continue
                p = dict(job_params) if type(job_params) is dict else {}
                if type(params) is dict:
                    p.update(params)
                _fn_name, df_result, fig_or_custom = self._run_analysis_core(
                    analysis_id,
                    p,
                    cache_name="Collection",
                    quiet=True,
                )
                payload0 = self._ppt_payload_from_result(
                    section_id=sec,
                    title=title,
                    df_result=df_result,
                    fig_or_custom=fig_or_custom,
                )
                built = payload0.get("slides") if type(payload0) is dict else None
                if type(built) is list and built:
                    for s in built:
                        if type(s) is dict:
                            slides.append(s)

            if not slides:
                slides.append({"title": sec.replace("_", " "), "bullets": ["No preview data available."], "notes": "", "section": sec})

            return slides

        def _build_pptx(include: list[str], params: dict, slides: list[dict] | None) -> Path | None:
            import base64
            import io
            from concurrent.futures import ThreadPoolExecutor

            df = self.loaded_dataframe
            if not isinstance(df, pd.DataFrame) or df.empty:
                return None

            collection_name = "Collection"
            slide_notes = _bool_param(params, "slide_notes", False)

            default_name = re.sub(r"[^\w\-_. ]", "_", f"{collection_name}_export.pptx").strip()
            out_path, _ = QFileDialog.getSaveFileName(self, "Save PowerPoint", default_name, "PowerPoint (*.pptx)")
            if not out_path:
                return None

            from .power_point_export import _normalize_slide_for_ui
            src_slides = slides if isinstance(slides, list) else []
            if not src_slides:
                ppt_data = self._ppt_preloaded_data or {}
                for sec in include:
                    payload = ppt_data.get(sec)
                    if isinstance(payload, dict) and isinstance(payload.get("slides"), list):
                        src_slides.extend(payload["slides"])

            if not src_slides:
                prs = Presentation()
                outp = Path(out_path)
                prs.save(str(outp))
                return outp

            def _decode_data_url_png(s: str) -> bytes | None:
                s = str(s or "").strip()
                if not s:
                    return None
                if s.startswith("data:image/png;base64,"):
                    b64 = s.split(",", 1)[1].strip()
                    return base64.b64decode(b64)
                if s.startswith("data:image/") and "base64," in s:
                    b64 = s.split("base64,", 1)[1].strip()
                    return base64.b64decode(b64)
                return None

            def _ensure_export_png(slide: dict) -> dict:
                if not isinstance(slide, dict):
                    return slide

                existing_img = None
                for key in ("img", "image", "png", "data_url", "image_data_url"):
                    v = slide.get(key)
                    if isinstance(v, str) and v.strip():
                        existing_img = v
                        break

                if existing_img:
                    return slide

                norm = _normalize_slide_for_ui(slide, want_png=True)
                img = norm.get("img")
                if not isinstance(img, str) or not img.strip():
                    return slide

                merged = dict(slide)
                merged["img"] = img

                if "fig_json" not in merged and norm.get("fig_json") is not None:
                    merged["fig_json"] = norm.get("fig_json")

                return merged

            export_slides = [_ensure_export_png(s) for s in src_slides if isinstance(s, dict)]

            def _slide_image_bytes(slide: dict) -> bytes | None:
                v = slide.get("img") or slide.get("image") or slide.get("png") or slide.get("data_url") or ""
                b = _decode_data_url_png(str(v))
                if b is not None:
                    return b

                b64 = slide.get("png_b64") or slide.get("image_b64") or ""
                if str(b64 or "").strip():
                    return base64.b64decode(str(b64).strip())

                return None

            def _prep_one(slide: dict) -> tuple[dict, bytes | None]:
                return slide, _slide_image_bytes(slide)

            with ThreadPoolExecutor(max_workers=4) as ex:
                prepared = list(ex.map(_prep_one, export_slides))
            prs = Presentation()
            blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

            for slide_dict, img_bytes in prepared:
                title = str(slide_dict.get("title") or "Slide")
                bullets = slide_dict.get("bullets")
                notes = str(slide_dict.get("notes") or slide_dict.get("note") or "")

                slide = prs.slides.add_slide(blank_layout)

                tx = slide.shapes.add_textbox(left=Inches(0.6), top=Inches(0.3), width=Inches(12.0), height=Inches(0.6))
                tf = tx.text_frame
                tf.clear()
                tf.text = title

                if img_bytes is not None:
                    bio = io.BytesIO(img_bytes)
                    slide.shapes.add_picture(bio, left=Inches(0.6), top=Inches(1.1), width=Inches(12.0))
                else:
                    body = slide.shapes.add_textbox(left=Inches(0.8), top=Inches(1.2), width=Inches(12.0),
                                                    height=Inches(5.5))
                    btf = body.text_frame
                    btf.clear()

                    if isinstance(bullets, list) and bullets:
                        btf.text = str(bullets[0] or "")
                        for b in bullets[1:]:
                            p = btf.add_paragraph()
                            p.text = str(b or "")
                            p.level = 0
                    elif notes.strip():
                        btf.text = notes
                    else:
                        btf.text = "No preview content."

                if slide_notes and notes.strip():
                    ns = slide.notes_slide
                    ns.notes_text_frame.text = notes

            outp = Path(out_path)
            prs.save(str(outp))
            return outp

        def _on_pptx_requested(payload: dict) -> None:
            mode = str(payload.get("mode") or "").lower()
            params = payload.get("params") or {}
            include = [str(x).strip() for x in payload.get("include", []) if str(x).strip()]

            if mode == "preview_section":
                sec = str(payload.get("section") or "")
                slides = _preview_from_preloaded(sec)
                export.set_slide_deck(slides, 0)
                export.set_status("Preview", f"{sec}: {len(slides)} slide(s)")
                return

            if mode == "preview":
                slides = []
                for sec in include:
                    slides.extend(_preview_from_preloaded(sec))
                export.set_slide_deck(slides, 0)
                export.set_status("Preview", f"{len(slides)} slide(s)")
                return

            if mode == "run_section":
                sec = str(payload.get("section") or "").strip()
                if not sec and include:
                    sec = include[0]
                slides = _run_section(sec, params)
                export.set_slide_deck(slides, 0)
                export.set_status("Preview", f"{sec}: {len(slides)} slide(s)")
                return

            if mode in {"build", "export_existing"}:
                slides_in = payload.get("slides")
                slides_list = slides_in if isinstance(slides_in, list) else None
                outp = _build_pptx(include, params, slides_list)
                if outp is None:
                    export.set_status("Cancelled", "No file generated.")
                    return
                export.set_status("Done", f"Saved: {str(outp)}")
                return

            export.set_status("Unknown", f"Unknown mode: {mode}")

        export.pptx_requested.connect(_on_pptx_requested)

        def _initial_preview():
            slides0 = []
            for sec in sections:
                sec_id = sec.get("id") if type(sec) is dict else None
                if type(sec_id) is not str or not sec_id.strip():
                    continue
                slides0.extend(_preview_from_preloaded(sec_id))
            export.set_slide_deck(slides0, 0)
            export.set_status("Preview", f"Loaded: {len(slides0)} slide(s)")

        QTimer.singleShot(0, _initial_preview)
        return w

    def _ppt_payload_from_result(self, *, section_id: str, title: str, df_result: Any, fig_or_custom: Any) -> dict:
        """
        ###1. convert (df, fig/custom) into PPT preview payload
        ###2. stamp slide["section"] so the UI can filter by key
        ###3. keep a stable {"slides":[...]} shape for the PPT UI
        """
        sec = str(section_id or "").strip()
        slides = []

        import pandas as pd
        import plotly.graph_objects as go

        fig_json = None

        if isinstance(fig_or_custom, go.Figure):
            fig_json = fig_or_custom.to_plotly_json()
        elif isinstance(fig_or_custom, dict) and ("data" in fig_or_custom and "layout" in fig_or_custom):
            fig_json = fig_or_custom

        if isinstance(fig_json, dict):
            slides.append({"title": title, "fig_json": fig_json, "notes": "", "section": sec})

        if isinstance(df_result, pd.DataFrame) and not df_result.empty:
            table_html = df_result.head(50).to_html(index=False, escape=True)
            slides.append({"title": f"{title} (table)", "table_html": table_html, "notes": "", "section": sec})

        if not slides:
            slides.append({"title": title, "bullets": ["No preview data available."], "notes": "", "section": sec})

        return {"slides": slides}

    def _analysis_dispatch_table(self) -> dict:
        """
        ###1. map analysis_id to a callable returning (fn_name, df_result, fig_or_custom)
        ###2. keep all analysis routing in one place
        ###3. allow reuse by run_requested_analysis and PPT preload
        """

        def _cb():
            return self.update_status_bar_message

        def _cb_quiet():
            def _noop(*args, **kwargs):
                return None

            return _noop

        def _run_author_impact(df, params, cb):
            return "analyze_author_impact", *analyze_author_impact(df, params, cb)

        def _run_author_collaboration(df, params, cb):
            return "analyze_author_collaboration", *analyze_author_collaboration(df, params, cb)

        def _run_author_trends(df, params, cb):
            return "analyze_author_trends", *analyze_author_trends(df, params, cb)

        def _run_lotka(df, params, cb):
            return "analyze_lotkas_law", *analyze_lotkas_law(df, params, cb)

        def _run_citations_overview(df, params, cb):
            return "analyze_citations_overview_plotly", *analyze_citations_overview_plotly(df, params, cb)

        def _run_temporal(df, params, cb):
            return "analyze_temporal_analysis", *analyze_temporal_analysis(df, params, cb)

        def _run_affiliations(df, params, cb):
            return "analyze_affiliations", *analyze_affiliations(df, params, cb)

        def _run_research_design(df, params, cb):
            return "analyze_research_design_suite", *analyze_research_design_suite(df, params, cb)

        def _run_profiles(df, params, cb):
            return "analyze_feature_profile", *analyze_feature_profile(df, params, cb)

        def _run_ngrams(df, params, cb, zot_client, cache_name):
            return "analyze_ngrams", *analyze_ngrams(df, params, cb, zot_client, cache_name)

        def _run_pdf_keywords(df, params, cb, zot_client):
            return "analyze_pdf_keywords_and_trends", *analyze_pdf_keywords_and_trends(
                df=df,
                params=params,
                progress_callback=cb,
                zotero_client=zot_client,
            )

        def _run_categorical_keywords(df, params, cb):
            return "analyze_categorical_keywords", *analyze_categorical_keywords(df, params, cb)

        def _run_word_analysis(df, params, cb, zot_client, cache_name):
            plot_type = str(params.get("plot_type") or "")
            if plot_type == "bar_chart":
                return "analyze_most_frequent_words", *analyze_most_frequent_words(df, params, cb, zot_client,
                                                                                   cache_name)
            if plot_type == "word_cloud":
                return "analyze_word_cloud", *analyze_word_cloud(df, params, cb, zot_client, cache_name)
            if plot_type == "treemap":
                return "analyze_word_treemap", *analyze_word_treemap(df, params, cb, zot_client, cache_name)
            if plot_type == "cooccurrence_network":
                params2 = dict(params)
                params2["plot_type"] = "network_graph"
                return "analyze_keyword_cooccurrence_network", *analyze_keyword_cooccurrence_network(
                    df, params2, cb, zot_client, cache_name
                )
            if plot_type == "heatmap":
                params2 = dict(params)
                params2["plot_type"] = "heatmap"
                return "analyze_keyword_heatmap", *analyze_keyword_heatmap(df, params2, cb, zot_client, cache_name)
            if plot_type == "words_over_time":
                return "analyze_words_over_time", *analyze_words_over_time(df, params, cb, zot_client, cache_name)
            raise ValueError(f"Unknown plot type '{plot_type}' for Word Analysis")

        table = {
            PAGE_ID_AUTHOR_IMPACT_ANALYSIS: lambda df, params, meta: _run_author_impact(df, params, meta["cb"]),
            PAGE_ID_AUTHOR_COLLABORATION: lambda df, params, meta: _run_author_collaboration(df, params, meta["cb"]),
            PAGE_ID_AUTHOR_TRENDS: lambda df, params, meta: _run_author_trends(df, params, meta["cb"]),
            PAGE_ID_LOTKA_LAW: lambda df, params, meta: _run_lotka(df, params, meta["cb"]),
            PAGE_ID_CITATIONS_OVERVIEW: lambda df, params, meta: _run_citations_overview(df, params, meta["cb"]),
            PAGE_ID_TEMPORAL_ANALYSIS: lambda df, params, meta: _run_temporal(df, params, meta["cb"]),
            PAGE_ID_AFFILIATION_GEO: lambda df, params, meta: _run_affiliations(df, params, meta["cb"]),
            PAGE_ID_RESEARCH_DESIGN_MAIN: lambda df, params, meta: _run_research_design(df, params, meta["cb"]),
            PAGE_ID_DESIGN_METHOD_MIX: lambda df, params, meta: _run_research_design(df, params, meta["cb"]),
            PAGE_ID_PROFILES_ANALYSIS: lambda df, params, meta: _run_profiles(df, params, meta["cb"]),
            PAGE_ID_WORD_ANALYSIS: lambda df, params, meta: _run_word_analysis(
                df, params, meta["cb"], meta["zot_client"], meta["cache_name"]
            ),
            PAGE_ID_NGRAM_ANALYSIS: lambda df, params, meta: _run_ngrams(
                df, params, meta["cb"], meta["zot_client"], meta["cache_name"]
            ),
            PAGE_ID_PDF_KEYWORD_SCAN: lambda df, params, meta: _run_pdf_keywords(
                df, params, meta["cb"], meta["zot_client"]
            ),
            PAGE_ID_CATEGORICAL_KEYWORD_ANALYSIS: lambda df, params, meta: _run_categorical_keywords(df, params,
                                                                                                     meta["cb"]),
        }
        return table
    def _run_analysis_core(self, analysis_id: str, params: dict, *, cache_name: str, quiet: bool) -> tuple[
        str, Any, Any]:
        """
        ###1. run an analysis by id using the central dispatcher
        ###2. return (fn_name, df_result, fig_or_custom)
        ###3. quiet=True uses a no-op callback
        """
        df = self.loaded_dataframe
        cb = self.update_status_bar_message
        if quiet:
            def _noop(*args, **kwargs):
                return None

            cb = _noop

        meta = {
            "cb": cb,
            "zot_client": zotero_client_instance,
            "cache_name": cache_name,
        }

        dispatch = self._analysis_dispatch_table()
        if analysis_id not in dispatch:
            raise ValueError(f"Analysis type '{analysis_id}' not recognized.")

        fn_name, df_result, fig_or_custom = dispatch[analysis_id](df, params, meta)
        return fn_name, df_result, fig_or_custom


    def _ppt_section_plan(self) -> dict:
        """
        ###1. map PPT section ids to analysis ids + default params
        ###2. keep defaults conservative so they work without extra UI inputs
        ###3. return {section_id: [(analysis_id, params, slide_title), ...]}
        """
        return {
            "Authors_overview": [
                (PAGE_ID_AUTHOR_IMPACT_ANALYSIS, {}, "Author impact"),
                (PAGE_ID_AUTHOR_COLLABORATION, {}, "Author collaboration"),
                (PAGE_ID_AUTHOR_TRENDS, {}, "Author trends"),
                (PAGE_ID_LOTKA_LAW, {}, "Lotka's law"),
            ],
            "Citations_overview": [
                (PAGE_ID_CITATIONS_OVERVIEW, {}, "Citations overview"),
            ],
            "Words_and_topics": [
                (PAGE_ID_WORD_ANALYSIS, {"plot_type": "bar_chart"}, "Most frequent words"),
                (PAGE_ID_NGRAM_ANALYSIS, {}, "N-grams"),
            ],
            "Affiliations_geo": [
                (PAGE_ID_AFFILIATION_GEO, {}, "Affiliations (geo)"),
            ],
            "Temporal_analysis": [
                (PAGE_ID_TEMPORAL_ANALYSIS, {}, "Temporal analysis"),
            ],
            "Research_design": [
                (PAGE_ID_RESEARCH_DESIGN_MAIN, {}, "Research design"),
            ],
            "Profiles": [
                (PAGE_ID_PROFILES_ANALYSIS, {}, "Profiles"),
            ],
            "Categorical_keywords": [
                (PAGE_ID_CATEGORICAL_KEYWORD_ANALYSIS, {}, "Categorical keywords"),
            ],
        }

    def _clear_ppt_cache(self) -> None:
        self._ppt_preloaded_data = {}
        self._ppt_payload_cache = {}
        self._save_ppt_cache_to_disk()

    def _ppt_cache_get(self, key: str):
        cache = self._ppt_payload_cache if isinstance(self._ppt_payload_cache, dict) else {}
        return cache.get(key)

    def _ppt_cache_set(self, key: str, payload) -> None:
        if not isinstance(self._ppt_payload_cache, dict):
            self._ppt_payload_cache = {}
        self._ppt_payload_cache[key] = payload

    def _ppt_cache_path(self) -> Path:
        base = (Path(MAIN_APP_CACHE_DIR) / "pages" / "ppts" / "ui_cache").resolve()
        base.mkdir(parents=True, exist_ok=True)
        return base / "ppt_payload_cache.json"

    def _load_ppt_cache_from_disk(self) -> None:
        path = self._ppt_cache_path()
        if not path.exists():
            return
        try:
            raw = path.read_text(encoding="utf-8")
            data = json.loads(raw) if raw.strip() else {}
        except Exception as exc:
            logging.warning("PPT cache load failed: %s", exc)
            return

        if isinstance(data, dict) and "entries" in data:
            entries = data.get("entries")
        else:
            entries = data

        if isinstance(entries, dict):
            self._ppt_payload_cache = entries

    def _save_ppt_cache_to_disk(self) -> None:
        path = self._ppt_cache_path()
        payload = {"version": 1, "entries": self._ppt_payload_cache}
        try:
            path.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")
        except Exception as exc:
            logging.warning("PPT cache save failed: %s", exc)

    def _refresh_ppt_preload_async(self) -> None:
        """
        ###1. refresh PPT preload using the existing BackgroundPageLoader only
        ###2. avoid calling load_ppt_data() from a second worker path
        ###3. keep behavior deterministic: one preload run per refresh request
        """
        df = getattr(self, "loaded_dataframe", None)
        if df is None or getattr(df, "empty", True):
            self._ppt_preloaded_data = {}
            return

        loader = getattr(self, "_ppt_loader", None)
        if loader is None:
            return

        state = str(getattr(loader, "state", "") or "")
        if state in {"loading", "waiting", "building"}:
            loader.kick()
            return

        loader.restart()

    def load_ppt_data(self, *, use_cache: bool = True) -> dict:
        """
        ###1. warm the PPT UI cache (template + local JS) and heavy deps early
        ###2. compute section payloads in a stable order
        ###3. normalize slides for UI, including optional PNG thumbnails (want_png=True)
        """
        import logging
        import re
        import json

        import pandas as pd
        from pptx import Presentation


        VERBOSE_PPT_PRELOAD = True
        sections_called = 0
        sections_cached = 0
        planned_jobs = 0
        called_jobs = 0
        skipped_missing_dispatch = 0

        def _stamp_section(payload, section_id: str) -> dict:
            sec = str(section_id or "").strip()
            if type(payload) is not dict:
                return {"slides": []}

            slides_list = payload.get("slides")
            figs_list = payload.get("figs")

            slides_out: list[dict] = []

            if type(slides_list) is list and slides_list:
                for s in slides_list:
                    if type(s) is dict:
                        d = dict(s)
                        if not str(d.get("section") or "").strip():
                            d["section"] = sec
                        slides_out.append(d)

            if not slides_out and type(figs_list) is list and figs_list:
                for i, fig in enumerate(figs_list):
                    title = sec.replace("_", " ") if i == 0 else f"{sec.replace('_', ' ')} ({i + 1})"
                    slides_out.append({"title": title, "fig_json": fig, "notes": "", "section": sec})

            outp = dict(payload)
            outp["slides"] = slides_out
            return outp

        def _normalize_deck(payload, section_id: str) -> dict:
            if type(payload) is not dict:
                return {"slides": [], "summary": {"slides": 0, "figs": 0, "imgs": 0, "tables": 0}}

            def _coerce_fig_json(v):
                """
                ###1. accept dict specs directly
                ###2. accept plotly.graph_objects.Figure and convert to plotly JSON dict
                ###3. accept JSON strings containing Plotly specs
                ###4. otherwise return None
                """
                if v is None:
                    return None

                if type(v) is dict:
                    return v

                import plotly.graph_objects as go

                if type(v) is go.Figure:
                    return v.to_plotly_json()

                if type(v) is str:
                    s = v.strip()
                    if s.startswith("{") or s.startswith("["):
                        parsed = json.loads(s)
                        if type(parsed) is dict:
                            return parsed
                    return None

                return None

            def _fig_meta(fig_json):
                if type(fig_json) is not dict:
                    return None
                data = fig_json.get("data")
                layout = fig_json.get("layout")
                data_len = len(data) if type(data) is list else None
                dragmode = ""
                x_fixed = None
                y_fixed = None

                if type(layout) is dict:
                    dragmode = str(layout.get("dragmode") or "")
                    xaxis = layout.get("xaxis")
                    yaxis = layout.get("yaxis")
                    if type(xaxis) is dict:
                        x_fixed = xaxis.get("fixedrange")
                    if type(yaxis) is dict:
                        y_fixed = yaxis.get("fixedrange")

                return {
                    "data_len": data_len,
                    "dragmode": dragmode,
                    "x_fixedrange": x_fixed,
                    "y_fixedrange": y_fixed,
                }

            slides_list = payload.get("slides")
            src = slides_list if type(slides_list) is list else []

            pre_norm = []
            for s in src:
                if type(s) is not dict:
                    continue

                d = dict(s)
                if d.get("fig_json") is not None:
                    d["fig_json"] = _coerce_fig_json(d.get("fig_json"))

                pre_norm.append(d)

            for i, s in enumerate(pre_norm):
                fig_json = s.get("fig_json")
                if fig_json is None:
                    continue

                meta = _fig_meta(fig_json)
                if meta is None:
                    continue

                logging.info("PPT preload fig_meta [%s] slide=%d %s", section_id, i, meta)

            norm = []
            for s in pre_norm:
                norm.append(_normalize_slide_for_ui(s, want_png=False))

            n_slides = len(norm)
            n_figs = 0
            n_imgs = 0
            n_tables = 0

            for s in norm:
                if type(s) is not dict:
                    continue

                fig_html = s.get("fig_html")
                if type(fig_html) is str and fig_html.strip():
                    n_figs += 1

                if type(s.get("img")) is str and s.get("img").strip():
                    n_imgs += 1

                if s.get("table_html") is not None and str(s.get("table_html") or "").strip():
                    n_tables += 1

            outp = dict(payload)
            outp["slides"] = norm
            outp["summary"] = {"slides": n_slides, "figs": n_figs, "imgs": n_imgs, "tables": n_tables}

            logging.info(
                "PPT preload [%s]: slides=%d figs=%d imgs=%d tables=%d",
                section_id,
                n_slides,
                n_figs,
                n_imgs,
                n_tables,
            )
            return outp

        def _cache_ready(payload) -> bool:
            if type(payload) is not dict:
                return False
            slides = payload.get("slides")
            return type(slides) is list and len(slides) > 0

        def _payload_stats(payload) -> tuple[int, int, int, int]:
            if type(payload) is not dict:
                return 0, 0, 0, 0
            summary = payload.get("summary")
            if type(summary) is dict:
                s = summary.get("slides")
                f = summary.get("figs")
                i = summary.get("imgs")
                t = summary.get("tables")
                if type(s) is int and type(f) is int and type(i) is int and type(t) is int:
                    return s, f, i, t

            slides = payload.get("slides")
            if type(slides) is not list:
                return 0, 0, 0, 0

            n_slides = len(slides)
            n_figs = 0
            n_imgs = 0
            n_tables = 0
            for s in slides:
                if type(s) is not dict:
                    continue
                fig_html = s.get("fig_html")
                if type(fig_html) is str and fig_html.strip():
                    n_figs += 1
                img_val = s.get("img")
                if type(img_val) is str and img_val.strip():
                    n_imgs += 1
                table_html = s.get("table_html")
                if table_html is not None and str(table_html or "").strip():
                    n_tables += 1
            return n_slides, n_figs, n_imgs, n_tables

        def _print_section_line(
                section_id: str,
                payload: dict,
                *,
                cached: bool,
                planned_jobs: int,
                called_jobs: int,
        ) -> None:
            s, f, i, t = _payload_stats(payload)
            cache_tag = "cached" if cached else "built"
            print(
                f"[PPT section] {section_id} {cache_tag} "
                f"planned_jobs={planned_jobs} called_jobs={called_jobs} "
                f"slides={s} figs={f} imgs={i} tables={t}",
                flush=True,
            )

        section_ids = [
            "Data_summary",
            "Scope_and_shape",
            "Authors_overview",
            "Institutions_overview",
            "Citations_overview",
            "Words_and_topics",
            "Affiliations_geo",
            "Temporal_analysis",
            "Research_design",
            "Profiles",
            "Categorical_keywords",
        ]

        out = {sid: {"slides": [], "summary": {"slides": 0, "figs": 0, "imgs": 0, "tables": 0}} for sid in section_ids}

        df = self.loaded_dataframe
        if df is None or type(df) is not pd.DataFrame or df.empty:
            return out

        ppt_ui_warmup()

        collection_name_for_cache = re.sub(
            r"[^\w\-_.]",
            "_",
            self.loaded_dataframe_source or "unknown_dataset",
        )

        base_cache = {
            "source": str(self.loaded_dataframe_source or ""),
            "collection_name": collection_name_for_cache,
            "rows": int(df.shape[0]) if type(df) is pd.DataFrame else 0,
            "cols": [str(c) for c in df.columns] if type(df) is pd.DataFrame else [],
        }

        controls = {}
        ppt_widget = self._ppt_widget
        if ppt_widget is not None:
            params = ppt_widget._params
            if type(params) is dict:
                controls = dict(params)
        base_cache["controls"] = controls

        plan = self._ppt_section_plan()
        dispatch = self._analysis_dispatch_table()

        SECTION_SPEC = {
            "Data_summary": {"mode": "direct"},
            "Scope_and_shape": {"mode": "direct"},
            "Authors_overview": {"mode": "direct"},
            "Institutions_overview": {"mode": "direct"},
            "Citations_overview": {"mode": "plan"},
            "Words_and_topics": {"mode": "plan"},
            "Affiliations_geo": {"mode": "plan"},
            "Temporal_analysis": {"mode": "plan"},
            "Research_design": {"mode": "plan"},
            "Profiles": {"mode": "plan"},
            "Categorical_keywords": {"mode": "plan"},
        }

        def _planned_jobs_for_section(section_id: str) -> int:
            spec = SECTION_SPEC[section_id]
            if spec["mode"] == "direct":
                return 1
            return len(plan.get(section_id, []))

        def _cache_key_for_section(section_id: str):
            spec = SECTION_SPEC[section_id]
            if spec["mode"] == "direct":
                return _ppt_cache_key({"section": section_id, "base": base_cache})

            jobs = plan.get(section_id, [])
            cache_items = [{"analysis_id": analysis_id, "params": params, "title": title} for analysis_id, params, title
                           in jobs]
            return _ppt_cache_key({"section": section_id, "base": base_cache, "items": cache_items})

        def _run_direct_section(section_id: str) -> dict:
            prs_tmp = Presentation()

            if section_id == "Data_summary":
                res = add_data_summary_slides(
                    prs=prs_tmp,
                    df=df,
                    collection_name="Collection",
                    slide_notes=False,
                    return_payload=True,
                    export=False,
                )
                payload0 = _stamp_section(res, section_id) if type(res) is dict else {"slides": [],
                                                                                      "error": "data_summary_failed"}
                return _normalize_deck(payload0, section_id)

            if section_id == "Scope_and_shape":
                res = shape_scope(
                    prs=prs_tmp,
                    df=df,
                    collection_name="Collection",
                    slide_notes=False,
                    export=False,
                    return_payload=True,
                )
                payload0 = _stamp_section(res, section_id) if type(res) is dict else {"slides": [],
                                                                                      "error": "scope_shape_failed"}
                return _normalize_deck(payload0, section_id)

            if section_id == "Authors_overview":
                res = authors_overview(
                    prs=prs_tmp,
                    df=df,
                    collection_name="Collection",
                    slide_notes=False,
                    export=True,
                    return_payload=True,
                )
                payload0 = _stamp_section(res, section_id) if type(res) is dict else {"slides": [],
                                                                                      "error": "authorship_failed"}
                return _normalize_deck(payload0, section_id)

            if section_id == "Institutions_overview":
                res = add_institution_section(
                    prs=prs_tmp,
                    df=df,
                    collection_name="Collection",
                    slide_notes=False,
                    export=True,
                    return_payload=True,
                )
                payload0 = _stamp_section(res, section_id) if type(res) is dict else {"slides": [],
                                                                                      "error": "institutions_failed"}
                return _normalize_deck(payload0, section_id)

            raise ValueError(f"Unknown direct section '{section_id}'")

        def _run_plan_section(section_id: str) -> tuple[dict, int]:
            jobs = plan.get(section_id, [])
            section_called_local = 0

            slides_acc: list[dict] = []

            for analysis_id, params, title in jobs:
                if analysis_id not in dispatch:
                    logging.warning("PPT preload skipped analysis_id=%r (missing dispatch entry)", analysis_id)
                    nonlocal_skipped[0] += 1
                    continue

                fn_name, df_result, fig_or_custom = self._run_analysis_core(
                    analysis_id,
                    params if type(params) is dict else {},
                    cache_name=collection_name_for_cache,
                    quiet=True,
                )
                nonlocal_called[0] += 1
                section_called_local += 1

                payload0 = self._ppt_payload_from_result(
                    section_id=section_id,
                    title=title,
                    df_result=df_result,
                    fig_or_custom=fig_or_custom,
                )

                built = payload0.get("slides") if type(payload0) is dict else None
                if type(built) is list and built:
                    for s in built:
                        if type(s) is dict:
                            slides_acc.append(s)

                if VERBOSE_PPT_PRELOAD:
                    print(
                        f"[{fn_name}] section={section_id} slides={len(built) if type(built) is list else 0}",
                        flush=True,
                    )

            payload_sec = _stamp_section({"slides": slides_acc}, section_id)
            return _normalize_deck(payload_sec, section_id), section_called_local

        nonlocal_called = [0]
        nonlocal_skipped = [0]

        for section_id in section_ids:
            planned = _planned_jobs_for_section(section_id)
            planned_jobs += planned

            key = _cache_key_for_section(section_id)
            cached = self._ppt_cache_get(key) if use_cache else None

            if use_cache and _cache_ready(cached):
                out[section_id] = cached
                sections_cached += 1
                _print_section_line(section_id, out[section_id], cached=True, planned_jobs=planned, called_jobs=0)
                continue

            spec = SECTION_SPEC[section_id]
            if spec["mode"] == "direct":
                out[section_id] = _run_direct_section(section_id)
                self._ppt_cache_set(key, out[section_id])
                sections_called += 1
                _print_section_line(section_id, out[section_id], cached=False, planned_jobs=planned, called_jobs=1)
                continue

            built_payload, section_called = _run_plan_section(section_id)
            out[section_id] = built_payload
            self._ppt_cache_set(key, out[section_id])
            sections_called += 1
            _print_section_line(section_id, out[section_id], cached=False, planned_jobs=planned,
                                called_jobs=section_called)

        called_jobs = int(nonlocal_called[0])
        skipped_missing_dispatch = int(nonlocal_skipped[0])

        self._save_ppt_cache_to_disk()
        print(
            f"[PPT preload] use_cache={use_cache} sections_called={sections_called} "
            f"sections_cached={sections_cached} planned_jobs={planned_jobs} "
            f"called_jobs={called_jobs} missing_dispatch={skipped_missing_dispatch}",
            flush=True,
        )
        return out
