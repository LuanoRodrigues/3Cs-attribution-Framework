#!/usr/bin/env python3
"""
Visualise host for Electron.
Actions:
  - sections: return schema + section list
  - preview: run selected analyses on a DataFrame table and return slide payload
"""
from __future__ import annotations

import json
import math
import sys
import traceback
from pathlib import Path
from typing import Any, Dict, List

BASE_DIR = Path(__file__).resolve().parent
PY_BACKEND_DIR = BASE_DIR.parent
SHARED_DIR = PY_BACKEND_DIR.parent
REPO_ROOT = next((parent for parent in BASE_DIR.parents if (parent / "package.json").is_file()), BASE_DIR.parents[-1])
for entry in (SHARED_DIR, BASE_DIR, REPO_ROOT):
    entry_str = str(entry)
    if entry_str not in sys.path:
        sys.path.insert(0, entry_str)


def _safe_json(obj: Any) -> str:
    def _clean(v: Any) -> Any:
        try:
            if hasattr(v, "item") and not isinstance(v, (bytes, str)):
                v = v.item()
        except Exception:
            pass
        if v is None or isinstance(v, (str, bool, int)):
            return v
        if isinstance(v, float):
            return v if math.isfinite(v) else None
        if isinstance(v, dict):
            return {str(k): _clean(vv) for k, vv in v.items()}
        if isinstance(v, (list, tuple, set)):
            return [_clean(x) for x in v]
        return str(v)

    cleaned = _clean(obj)
    return json.dumps(cleaned, ensure_ascii=True, allow_nan=False)


def _read_payload() -> Dict[str, Any]:
    raw = sys.stdin.read()
    if not raw:
        return {}
    try:
        return json.loads(raw)
    except json.JSONDecodeError as exc:
        return {"action": "error", "message": f"invalid_json:{exc}"}


def _repo_root() -> Path:
    return REPO_ROOT


def _load_visual_helpers() -> tuple[Any | None, list[str]]:
    try:
        import importlib

        module = importlib.import_module("python_backend.visualise.visual_helpers")
        return importlib.reload(module), []
    except Exception:
        return None, ["[visualise][load][error] failed to import visual_helpers", traceback.format_exc()]


def _coerce_table(table: Dict[str, Any]):
    import pandas as pd

    cols = table.get("columns") if isinstance(table, dict) else None
    rows = table.get("rows") if isinstance(table, dict) else None
    if not isinstance(cols, list) or not isinstance(rows, list):
        return pd.DataFrame()
    return pd.DataFrame(rows, columns=[str(c) for c in cols])

def _is_table_payload(table: Any) -> bool:
    try:
        return isinstance(table, dict) and isinstance(table.get("columns"), list) and isinstance(table.get("rows"), list)
    except Exception:
        return False


def _load_last_table_from_datahub_cache(*, logs: list[str]) -> Dict[str, Any] | None:
    """
    Best-effort fallback when the renderer didn't pass a table.
    Reads Electron's DataHub cache marker and returns the cached `{columns, rows}` table.
    """
    try:
        import os
        from pathlib import Path

        # Prefer explicit env (main-process can set this), then XDG-style default used by Electron on Linux.
        candidates: list[Path] = []
        env = str(os.environ.get("DATAHUB_CACHE_DIR") or "").strip()
        if env:
            candidates.append(Path(env).expanduser())
        xdg_config = Path(os.environ.get("XDG_CONFIG_HOME") or (Path.home() / ".config"))
        candidates.append(xdg_config / "my-electron-app" / "data-hub-cache")

        for cache_dir in candidates:
            last_path = cache_dir / "last.json"
            if not last_path.exists():
                continue
            try:
                last = json.loads(last_path.read_text(encoding="utf-8"))
            except Exception as exc:
                logs.append(f"[visualise][cache][warn] failed to read {last_path}: {exc}")
                continue
            cache_path = str((last or {}).get("cachePath") or "").strip()
            if cache_path:
                p = Path(cache_path).expanduser()
                if p.exists():
                    try:
                        payload = json.loads(p.read_text(encoding="utf-8"))
                        table = payload.get("table") if isinstance(payload, dict) else None
                        if _is_table_payload(table):
                            logs.append(f"[visualise][cache][loaded] {p}")
                            return table  # type: ignore[return-value]
                    except Exception as exc:
                        logs.append(f"[visualise][cache][warn] failed to read cached table {p}: {exc}")
            # Fallback: scan cache dir for any table json.
            try:
                newest: tuple[float, Path] | None = None
                for cand in cache_dir.glob("*.json"):
                    if cand.name.lower() in {
                        "references.json",
                        "references_library.json",
                        "references.used.json",
                        "references_library.used.json",
                        "last.json",
                    }:
                        continue
                    try:
                        m = cand.stat().st_mtime
                        if newest is None or m > newest[0]:
                            newest = (m, cand)
                    except Exception:
                        continue
                if newest is not None:
                    payload = json.loads(newest[1].read_text(encoding="utf-8"))
                    table = payload.get("table") if isinstance(payload, dict) else None
                    if _is_table_payload(table):
                        logs.append(f"[visualise][cache][loaded] {newest[1]}")
                        return table  # type: ignore[return-value]
            except Exception:
                continue
    except Exception as exc:
        logs.append(f"[visualise][cache][warn] datahub cache load failed: {exc}")
    return None


def _schema_and_sections() -> dict:
    schema = [
        {"type": "number", "key": "top_n_authors", "label": "Top N authors", "default": 15, "min": 5, "max": 100},
        {"type": "number", "key": "production_top_n", "label": "Production top N", "default": 10, "min": 3, "max": 50},
        {
            "type": "select",
            "key": "year_bucket",
            "label": "Year bucket",
            "default": 1,
            "options": [
                {"label": "1 year", "value": "1"},
                {"label": "2 years", "value": "2"},
                {"label": "5 years", "value": "5"},
            ],
        },
        {
            "type": "select",
            "key": "data_source",
            "label": "Text source",
            "default": "controlled_vocabulary_terms",
            "options": [
                {"label": "Keywords", "value": "controlled_vocabulary_terms"},
                {"label": "Title", "value": "title"},
                {"label": "Abstract", "value": "abstract"},
                {"label": "Full text", "value": "fulltext"},
            ],
        },
        {
            "type": "select",
            "key": "words_topics_view",
            "label": "Words & topics",
            "default": "both",
            "options": [
                {"label": "Words", "value": "words"},
                {"label": "N-grams", "value": "ngrams"},
                {"label": "Both", "value": "both"},
            ],
        },
        {
            "type": "select",
            "key": "word_plot_type",
            "label": "Word analysis plot",
            "default": "bar_vertical",
            "options": [
                {"label": "Frequency bar (vertical)", "value": "bar_vertical"},
                {"label": "Frequency bar (horizontal)", "value": "bar_horizontal"},
                {"label": "Word cloud", "value": "word_cloud"},
                {"label": "Treemap", "value": "treemap"},
                {"label": "Frequency heatmap", "value": "heatmap"},
                {"label": "Co-occurrence network", "value": "cooccurrence_network"},
                {"label": "Words over time", "value": "words_over_time"},
                {"label": "Term diversity over time", "value": "term_diversity_over_time"},
            ],
        },
        {"type": "number", "key": "top_n_words", "label": "Top N words", "default": 30, "min": 5, "max": 500},
        {"type": "number", "key": "max_words", "label": "Max words (cloud)", "default": 100, "min": 10, "max": 500},
        {"type": "number", "key": "min_frequency", "label": "Min frequency", "default": 2, "min": 1, "max": 100},
        {"type": "number", "key": "min_cooccurrence", "label": "Min co-occurrence", "default": 2, "min": 1, "max": 100},
        {"type": "number", "key": "min_keyword_frequency", "label": "Min keyword freq (network)", "default": 3, "min": 1, "max": 100},
        {"type": "number", "key": "max_nodes_for_plot", "label": "Max nodes (network)", "default": 50, "min": 10, "max": 200},
        {"type": "number", "key": "num_top_words", "label": "Top N words (over time)", "default": 5, "min": 1, "max": 20},
        {
            "type": "textarea",
            "key": "specific_words_to_track",
            "label": "Specific words (optional)",
            "default": "",
            "asList": True,
            "placeholder": "semicolon or newline separated…",
        },
        {
            "type": "textarea",
            "key": "specific_ngrams_to_track",
            "label": "Specific n-grams (evolution, optional)",
            "default": "",
            "asList": True,
            "placeholder": "Use when N-gram plot is 'Evolution over time' (semicolon or newline separated)…",
        },
        {
            "type": "select",
            "key": "wordcloud_colormap",
            "label": "Word cloud colormap",
            "default": "viridis",
            "options": [
                {"label": "viridis", "value": "viridis"},
                {"label": "plasma", "value": "plasma"},
                {"label": "inferno", "value": "inferno"},
                {"label": "magma", "value": "magma"},
                {"label": "cividis", "value": "cividis"},
                {"label": "Pastel1", "value": "Pastel1"},
                {"label": "Paired", "value": "Paired"},
                {"label": "Spectral", "value": "Spectral"},
                {"label": "coolwarm", "value": "coolwarm"},
                {"label": "RdYlGn", "value": "RdYlGn"},
            ],
        },
        {"type": "number", "key": "contour_width", "label": "Cloud contour width", "default": 0, "min": 0, "max": 5},
        {"type": "text", "key": "contour_color", "label": "Cloud contour color", "default": "steelblue"},
        {
            "type": "select",
            "key": "ngram_plot_type",
            "label": "N-gram plot",
            "default": "bar_chart",
            "options": [
                {"label": "Bar chart (frequency)", "value": "bar_chart"},
                {"label": "Evolution over time", "value": "ngram_evolution_time_series"},
                {"label": "Co-occurrence network", "value": "ngram_cooccurrence_network"},
                {"label": "Frequency heatmap", "value": "ngram_frequency_heatmap"},
            ],
        },
        {"type": "number", "key": "ngram_n", "label": "N-gram size (n)", "default": 2, "min": 1, "max": 7},
        {"type": "number", "key": "top_n_ngrams", "label": "Top N n-grams", "default": 20, "min": 5, "max": 200},
        {"type": "number", "key": "num_top_ngrams_for_evolution", "label": "N-grams to track (evolution)", "default": 7, "min": 1, "max": 15},
        {"type": "number", "key": "min_ngram_cooccurrence", "label": "Min n-gram co-occurrence", "default": 2, "min": 1, "max": 20},
        {"type": "number", "key": "max_nodes_for_ngram_network", "label": "Max nodes (ngram net)", "default": 30, "min": 10, "max": 150},
        {"type": "number", "key": "num_ngrams_for_heatmap_cols", "label": "N-grams (heatmap cols)", "default": 25, "min": 5, "max": 50},
        {"type": "number", "key": "num_docs_for_heatmap_rows", "label": "Docs (heatmap rows)", "default": 30, "min": 5, "max": 100},
        # Thematic & method
        {"type": "text", "key": "phase_tag", "label": "Phase column/tag", "default": "phase_focus_value"},
        {"type": "text", "key": "year_col", "label": "Year column", "default": "year"},
        {"type": "number", "key": "topn_cols", "label": "Cross-tab max columns", "default": 14, "min": 5, "max": 50},
        # Back-compat keys (older UI)
        {"type": "number", "key": "top_ngram", "label": "Top n-gram (legacy)", "default": 20, "min": 5, "max": 200},
        {
            "type": "select",
            "key": "slide_notes",
            "label": "slide_notes",
            "default": "false",
            "options": [{"label": "false", "value": "false"}, {"label": "true", "value": "true"}],
        },
    ]
    sections = [
        {"id": "Data_summary", "label": "Data summary", "hint": "Summary cards and key totals."},
        {"id": "Scope_and_shape", "label": "Scope and shape", "hint": "Volume, years, document types."},
        {"id": "Authors_overview", "label": "Authors overview", "hint": "Top authors and collaboration."},
        {"id": "Authorship_institution", "label": "Authorship & institutions", "hint": "Top authors, co-authorship and institutional actors."},
        {"id": "Citations_overview", "label": "Citations overview", "hint": "Citation distribution and leaders."},
        {"id": "Citations_influence", "label": "Citations & influence", "hint": "Stats, distributions and top-cited works."},
        {"id": "Words_and_topics", "label": "Words and topics", "hint": "Keywords, n-grams, topic signals."},
        {"id": "Ngrams", "label": "N-grams", "hint": "N-gram frequency and co-occurrence."},
        {"id": "Affiliations_geo", "label": "Affiliations (geo)", "hint": "Institutions and geography."},
        {"id": "Temporal_analysis", "label": "Temporal analysis", "hint": "Trends over time."},
        {"id": "Research_design", "label": "Research design", "hint": "Design outputs and mix."},
        {"id": "Thematic_method", "label": "Thematic & method", "hint": "Cross-tabs and phase mapping."},
        {"id": "Categorical_keywords", "label": "Categorical keywords", "hint": "Categorical keyword outputs."},
    ]
    return {"schema": schema, "sections": sections}


def _normalize_slide(slide: dict) -> dict:
    out = dict(slide)
    fig_json = out.get("fig_json")
    if isinstance(fig_json, str):
        s = fig_json.strip()
        if s.startswith("{") or s.startswith("["):
            try:
                out["fig_json"] = json.loads(s)
            except json.JSONDecodeError:
                out["fig_json"] = None
    return out


def _flatten_slides(payload: Any, *, section: str) -> list[dict]:
    if isinstance(payload, dict) and isinstance(payload.get("slides"), list):
        out: list[dict] = []
        for s in payload.get("slides"):
            if not isinstance(s, dict):
                continue
            slide = _normalize_slide(s)
            if section and "section" not in slide:
                slide["section"] = section
            out.append(slide)
        return out
    return []

def _export_pptx_from_slides(
    slides: list[dict],
    *,
    output_path: Path,
    slide_notes: bool,
    logs: list[str],
) -> Path:
    import base64
    import io

    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    def _table_text(html: str, *, max_rows: int = 22, max_cols: int = 10) -> str:
        if not isinstance(html, str) or not html.strip():
            return ""
        try:
            from bs4 import BeautifulSoup  # type: ignore

            soup = BeautifulSoup(html, "html.parser")
            rows = []
            for r_i, tr in enumerate(soup.find_all("tr")):
                if r_i >= max_rows:
                    break
                cells = tr.find_all(["th", "td"])
                vals = [c.get_text(" ", strip=True) for c in cells[:max_cols]]
                if vals:
                    rows.append("\t".join(vals))
            return "\n".join(rows).strip()
        except Exception:
            import re

            txt = re.sub(r"<[^>]+>", " ", html)
            txt = re.sub(r"\s+", " ", txt).strip()
            return txt[:4000]

    def _decode_data_url_to_bytes(s: str) -> bytes | None:
        s = str(s or "").strip()
        if not s:
            return None
        if s.startswith("data:image/") and "base64," in s:
            try:
                b64 = s.split("base64,", 1)[1].strip()
                return base64.b64decode(b64)
            except Exception:
                return None
        return None

    def _fig_png_bytes(fig_json: Any) -> bytes | None:
        if not fig_json:
            return None
        try:
            import plotly.graph_objects as go  # type: ignore
            import plotly.io as pio  # type: ignore

            fig = go.Figure(fig_json)
            return pio.to_image(fig, format="png", width=1400, height=800, scale=2)
        except Exception as exc:
            logs.append(f"[visualise][export_pptx][warn] fig->png failed: {exc}")
            logs.append(traceback.format_exc())
            return None

    def _add_picture_fit(slide, png_bytes: bytes, *, title_height_in: float = 0.95) -> None:
        from pptx.util import Inches

        # Default widescreen: 13.33 x 7.5 inches
        slide_w_in = prs.slide_width / 914400.0
        slide_h_in = prs.slide_height / 914400.0
        left = 0.6
        right = 0.6
        top = title_height_in
        bottom = 0.6
        box_w = max(1.0, slide_w_in - left - right)
        box_h = max(1.0, slide_h_in - top - bottom)

        img_w_px = None
        img_h_px = None
        try:
            from PIL import Image  # type: ignore

            im = Image.open(io.BytesIO(png_bytes))
            img_w_px, img_h_px = im.size
        except Exception:
            img_w_px, img_h_px = None, None

        if img_w_px and img_h_px:
            aspect = float(img_w_px) / float(img_h_px)
            box_aspect = float(box_w) / float(box_h)
            if aspect >= box_aspect:
                w = box_w
                h = w / aspect
            else:
                h = box_h
                w = h * aspect
            x = left + (box_w - w) / 2.0
            y = top + (box_h - h) / 2.0
            slide.shapes.add_picture(io.BytesIO(png_bytes), left=Inches(x), top=Inches(y), width=Inches(w), height=Inches(h))
            return

        # Fallback: width only.
        slide.shapes.add_picture(io.BytesIO(png_bytes), left=Inches(left), top=Inches(top), width=Inches(box_w))

    def _extract_table(html: str, *, max_rows: int = 26, max_cols: int = 10) -> tuple[list[list[str]], bool]:
        if not isinstance(html, str) or not html.strip():
            return ([], False)
        try:
            from bs4 import BeautifulSoup  # type: ignore

            soup = BeautifulSoup(html, "html.parser")
            rows_out: list[list[str]] = []
            header = False
            for r_i, tr in enumerate(soup.find_all("tr")):
                if r_i >= max_rows:
                    break
                cells = tr.find_all(["th", "td"])
                vals = [c.get_text(" ", strip=True) for c in cells[:max_cols]]
                if vals:
                    rows_out.append(vals)
                if tr.find("th") is not None:
                    header = True
            return (rows_out, header)
        except Exception:
            return ([], False)

    def _add_table_styled(slide, table_html: str, *, title_height_in: float = 0.95) -> bool:
        rows, has_header = _extract_table(table_html)
        if not rows:
            return False
        # normalize to rectangle
        cols_n = max(len(r) for r in rows)
        rows = [r + [""] * (cols_n - len(r)) for r in rows]

        slide_w_in = prs.slide_width / 914400.0
        slide_h_in = prs.slide_height / 914400.0
        left = 0.6
        right = 0.6
        top = title_height_in
        bottom = 0.6
        box_w = max(1.0, slide_w_in - left - right)
        box_h = max(1.0, slide_h_in - top - bottom)

        tbl_shape = slide.shapes.add_table(len(rows), cols_n, Inches(left), Inches(top), Inches(box_w), Inches(box_h))
        table = tbl_shape.table

        # Basic styling
        header_fill = RGBColor(31, 41, 55)  # slate-800
        header_text = RGBColor(255, 255, 255)
        body_text = RGBColor(17, 24, 39)  # slate-900
        stripe_a = RGBColor(255, 255, 255)
        stripe_b = RGBColor(243, 244, 246)  # gray-100

        for r in range(len(rows)):
            for c in range(cols_n):
                cell = table.cell(r, c)
                cell.text = str(rows[r][c] or "")
                tf = cell.text_frame
                tf.word_wrap = True
                for p in tf.paragraphs:
                    p.alignment = PP_ALIGN.LEFT
                    for run in p.runs:
                        run.font.size = Pt(10)
                        run.font.name = "Calibri"
                        run.font.color.rgb = header_text if (has_header and r == 0) else body_text
                        run.font.bold = bool(has_header and r == 0)
                fill = cell.fill
                fill.solid()
                if has_header and r == 0:
                    fill.fore_color.rgb = header_fill
                else:
                    fill.fore_color.rgb = stripe_a if (r % 2 == 0) else stripe_b
        return True

    for idx, slide_dict in enumerate(slides):
        if not isinstance(slide_dict, dict):
            continue
        title = str(slide_dict.get("title") or f"Slide {idx + 1}")
        notes = str(slide_dict.get("notes") or "")
        table_html = str(slide_dict.get("table_html") or "")
        fig_json = slide_dict.get("fig_json")
        bullets = slide_dict.get("bullets")
        img_value = slide_dict.get("img") or slide_dict.get("image") or slide_dict.get("png") or slide_dict.get("data_url") or ""

        slide = prs.slides.add_slide(blank_layout)

        tx = slide.shapes.add_textbox(left=Inches(0.6), top=Inches(0.35), width=Inches(12.0), height=Inches(0.6))
        tf = tx.text_frame
        tf.clear()
        tf.text = title
        try:
            tf.paragraphs[0].runs[0].font.size = Pt(26)
        except Exception:
            pass

        png = _decode_data_url_to_bytes(str(img_value)) if isinstance(img_value, str) else None
        if png is None:
            png = _fig_png_bytes(fig_json)

        if png:
            _add_picture_fit(slide, png, title_height_in=1.05)
        elif table_html.strip() and _add_table_styled(slide, table_html, title_height_in=1.05):
            pass
        else:
            body = slide.shapes.add_textbox(left=Inches(0.7), top=Inches(1.1), width=Inches(12.0), height=Inches(5.7))
            btf = body.text_frame
            btf.clear()

            table_text = _table_text(table_html)
            if table_text:
                lines = table_text.splitlines()
                btf.text = (lines[0] if lines else "")[:160]
                for line in lines[1:40]:
                    p = btf.add_paragraph()
                    p.text = line[:180]
                    p.level = 0
            elif isinstance(bullets, list) and bullets:
                btf.text = str(bullets[0] or "")
                for b in bullets[1:20]:
                    p = btf.add_paragraph()
                    p.text = str(b or "")
                    p.level = 0
            elif notes.strip():
                btf.text = notes[:2000]
            else:
                btf.text = "No figure available for this slide."

        if slide_notes and notes.strip():
            try:
                slide.notes_slide.notes_text_frame.text = notes
            except Exception:
                pass

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def _describe_slide(slide: dict, *, collection_name: str, params: dict) -> tuple[str, list[str]]:
    logs: list[str] = []
    title = str(slide.get("title") or "Figure").strip()
    section = str(slide.get("section") or "").strip()
    fig_json = slide.get("fig_json")
    table_html = str(slide.get("table_html") or "").strip()

    if isinstance(fig_json, str):
        s = fig_json.strip()
        if s.startswith("{") or s.startswith("["):
            try:
                fig_json = json.loads(s)
            except Exception:
                fig_json = None

    def _strip_html(s: str, max_chars: int = 2400) -> str:
        import re

        txt = re.sub(r"<[^>]+>", " ", s or "")
        txt = re.sub(r"\s+", " ", txt).strip()
        return txt[:max_chars]

    def _top_pairs(cats: list[str], vals: list[float], n: int = 5) -> list[tuple[str, float]]:
        pairs = []
        for c, v in zip(cats, vals):
            if c is None:
                continue
            name = str(c).strip()
            if not name:
                continue
            try:
                fv = float(v)
            except Exception:
                continue
            pairs.append((name, fv))
        pairs.sort(key=lambda x: x[1], reverse=True)
        return pairs[:n]

    desc_bits: list[str] = []
    desc_bits.append(f"{title}".strip())
    if section:
        desc_bits.append(f"Section: {section.replace('_', ' ')}")
    if collection_name:
        desc_bits.append(f"Collection: {collection_name}")

    # --- Heuristic summary from fig_json
    summary = ""
    try:
        if isinstance(fig_json, dict):
            data = fig_json.get("data") if isinstance(fig_json.get("data"), list) else []
            layout = fig_json.get("layout") if isinstance(fig_json.get("layout"), dict) else {}
            types = [str(t.get("type") or "scatter") for t in data if isinstance(t, dict)]
            types_uniq = sorted({t for t in types if t})
            if types_uniq:
                summary += f"Plot type(s): {', '.join(types_uniq)}. "
            # Prefer first trace for stats
            t0 = next((t for t in data if isinstance(t, dict)), None)
            if isinstance(t0, dict):
                ttype = str(t0.get("type") or "scatter")
                orient = str(t0.get("orientation") or "v")
                if ttype == "bar":
                    cats_raw = t0.get("y") if orient == "h" else t0.get("x")
                    vals_raw = t0.get("x") if orient == "h" else t0.get("y")
                    if isinstance(cats_raw, list) and isinstance(vals_raw, list) and len(cats_raw) == len(vals_raw):
                        pairs = _top_pairs([str(c) for c in cats_raw], vals_raw, n=6)
                        if pairs:
                            top = ", ".join([f"{c} ({v:g})" for c, v in pairs])
                            summary += f"Top categories: {top}. "
                elif ttype == "histogram":
                    xs = t0.get("x")
                    if isinstance(xs, list) and xs:
                        import numpy as np

                        vals = []
                        for v in xs:
                            try:
                                vals.append(float(v))
                            except Exception:
                                pass
                        if vals:
                            arr = np.array(vals, dtype=float)
                            summary += f"N={len(vals)} mean={float(arr.mean()):.2f} median={float(np.median(arr)):.2f}. "
                elif ttype in {"scatter", "scattergl"}:
                    xs = t0.get("x")
                    ys = t0.get("y")
                    if isinstance(xs, list) and isinstance(ys, list) and len(xs) == len(ys) and len(xs) >= 3:
                        # attempt numeric trend on y
                        import numpy as np

                        yv = []
                        for v in ys:
                            try:
                                yv.append(float(v))
                            except Exception:
                                yv.append(np.nan)
                        arr = np.array(yv, dtype=float)
                        arr = arr[np.isfinite(arr)]
                        if arr.size >= 3:
                            delta = float(arr[-1] - arr[0])
                            direction = "increasing" if delta > 0 else "decreasing" if delta < 0 else "flat"
                            summary += f"Trend appears {direction} over the plotted range. "
                elif ttype == "pie":
                    labels = t0.get("labels")
                    values = t0.get("values")
                    if isinstance(labels, list) and isinstance(values, list) and len(labels) == len(values):
                        pairs = _top_pairs([str(c) for c in labels], values, n=6)
                        if pairs:
                            top = ", ".join([f"{c} ({v:g})" for c, v in pairs])
                            summary += f"Top slices: {top}. "
            # Include title from layout if present
            try:
                lt = layout.get("title")
                if isinstance(lt, dict) and lt.get("text"):
                    summary = f"{str(lt.get('text')).strip()}. " + summary
                elif isinstance(lt, str) and lt.strip():
                    summary = f"{lt.strip()}. " + summary
            except Exception:
                pass
    except Exception as exc:
        logs.append(f"[visualise][describe][warn] fig summary failed: {exc}")

    if summary:
        desc_bits.append("")
        desc_bits.append("Summary:")
        desc_bits.append(summary.strip())

    # --- Table cue
    if table_html:
        desc_bits.append("")
        desc_bits.append("Table (preview):")
        desc_bits.append(_strip_html(table_html))

    # --- LLM call (optional/stubbed in this repo), fallback to heuristic.
    desc = "\n".join([b for b in desc_bits if b is not None]).strip()
    try:
        from python_backend.core.utils.calling_models import call_models_plots  # type: ignore
        import os

        prompt_text = (
            "[TASK]\n"
            "Write speaker notes to explain the chart for a presentation slide.\n"
            "Be concise, concrete, and mention key patterns + caveats.\n\n"
            "[CONTEXT]\n"
            f"Collection: {collection_name}\n"
            f"Section: {section}\n"
            f"Title: {title}\n\n"
            "[SLIDE DATA]\n"
            f"{desc}\n"
        )
        ai_out = call_models_plots(
            mode="analyze",
            prompt_text=prompt_text,
            model_api_name=os.getenv("OPENAI_VISION_MODEL", "gpt-5.1-mini"),
            images=None,
            image_detail="low",
            max_tokens=500,
            use_cache=True,
            analysis_key_suffix="visualise_slide_description",
            section_title=collection_name or "collection",
            store_only=False,
            read=False,
        )
        ai_text = str((ai_out or {}).get("text") or "").strip()
        if ai_text:
            return ai_text, logs
    except Exception as exc:
        logs.append(f"[visualise][describe][warn] ai describe failed: {exc}")

    return desc, logs


def _run_preview(
    df,
    params: dict,
    include: list[str],
    collection_name: str,
    *,
    mode: str = "",
    selection: dict | None = None,
) -> tuple[dict, list[str]]:
    visual, load_logs = _load_visual_helpers()
    logs: list[str] = []
    logs.extend(load_logs)
    if visual is None:
        return (
            {
                "slides": [
                    {
                        "title": "Visualise",
                        "bullets": ["Visualise backend failed to load. See logs."],
                        "notes": "",
                    }
                ],
                "index": 0,
            },
            logs,
        )

    from pptx import Presentation

    prs = Presentation()
    slide_notes = str(params.get("slide_notes", "false")).strip().lower() in {"1", "true", "yes", "y", "on"}
    is_build = "build" in str(mode or "").lower()
    if is_build:
        try:
            from python_backend.visualise.power_point_export import ppt_ui_warmup

            warm = ppt_ui_warmup()
            logs.append(
                "[visualise][ppt_ui_warmup] "
                + _safe_json(
                    {
                        "res_dir": warm.get("res_dir"),
                        "template_len": warm.get("template_len"),
                        "plotly_js_len": warm.get("plotly_js_len"),
                        "qweb_js_len": warm.get("qweb_js_len"),
                    }
                )
            )
        except Exception as exc:
            logs.append(f"[visualise][ppt_ui_warmup][error] {exc}")
            logs.append(traceback.format_exc())

    slides: list[dict] = []

    def _error_slide(section_id: str, exc: Exception) -> dict:
        return _normalize_slide(
            {
                "title": f"{section_id.replace('_', ' ')} — Error",
                "bullets": [str(exc)],
                "notes": "",
                "section": section_id,
            }
        )

    def _cb(_msg: str):
        return None

    defaults = {
        "word": {"plot_type": "bar_vertical", "data_source": "controlled_vocabulary_terms"},
        "ngrams": {"plot_type": "bar_chart", "data_source": "controlled_vocabulary_terms", "ngram_n": 2, "top_n_ngrams": 20},
        # Prefer offline-safe plots by default (world maps require topojson fetch in the renderer).
        "affiliations": {"plot_type": "country_bar", "top_n": 30},
        "temporal": {"plot_type": "multi_line_trend"},
        "citations": {"plot_type": "top_cited_bar"},
        "categorical": {"selected_category": "#theme", "plot_type": "bar_chart"},
    }

    if not include:
        include = [s["id"] for s in _schema_and_sections().get("sections", []) if isinstance(s, dict) and s.get("id")]

    for sec in include:
        logs.append(f"[visualise][section][start] {sec}")
        if sec == "Data_summary":
            try:
                payload = visual.add_data_summary_slides(
                    prs, df, collection_name=collection_name, slide_notes=slide_notes, return_payload=True, export=False
                )
                slides.extend(_flatten_slides(payload, section=sec))
                logs.append(f"[visualise][section][ok] Data_summary")
            except Exception as exc:
                logs.append(f"[visualise][section][error] Data_summary: {exc}")
                logs.append(traceback.format_exc())
                slides.append(_error_slide(sec, exc))
        elif sec == "Scope_and_shape":
            try:
                payload = visual.shape_scope(
                    prs,
                    df,
                    collection_name=collection_name,
                    slide_notes=slide_notes,
                    return_payload=True,
                    export=False,
                    progress_callback=_cb,
                )
                slides.extend(_flatten_slides(payload, section=sec))
                logs.append(f"[visualise][section][ok] Scope_and_shape")
            except Exception as exc:
                logs.append(f"[visualise][section][error] Scope_and_shape: {exc}")
                logs.append(traceback.format_exc())
                slides.append(_error_slide(sec, exc))
        elif sec == "Authors_overview":
            try:
                payload = visual.authors_overview(
                    prs, df, collection_name=collection_name, slide_notes=slide_notes, return_payload=True, export=False
                )
                slides.extend(_flatten_slides(payload, section=sec))
                logs.append(f"[visualise][section][ok] Authors_overview")
            except Exception as exc:
                logs.append(f"[visualise][section][error] Authors_overview: {exc}")
                logs.append(traceback.format_exc())
                slides.append(_error_slide(sec, exc))
        elif sec == "Authorship_institution":
            try:
                p = dict(params or {})
                top_n_authors = int(p.get("top_n_authors", 20) or 20)
                payload = visual.add_authorship_and_institution_section(
                    prs,
                    df,
                    collection_name=collection_name,
                    slide_notes=slide_notes,
                    top_n_authors=top_n_authors,
                    return_payload=True,
                    export=False,
                    progress_callback=_cb,
                )
                slides.extend(_flatten_slides(payload, section=sec))
                logs.append(f"[visualise][section][ok] Authorship_institution")
            except Exception as exc:
                logs.append(f"[visualise][section][error] Authorship_institution: {exc}")
                logs.append(traceback.format_exc())
                slides.append(_error_slide(sec, exc))
        elif sec == "Citations_overview":
            try:
                p = dict(defaults["citations"], **(params or {}))
                payload = visual.analyze_citations_overview_plotly(df, p, _cb, return_payload=True, export=False)
                slides.extend(_flatten_slides(payload, section=sec))
                logs.append(f"[visualise][section][ok] Citations_overview")
            except Exception as exc:
                logs.append(f"[visualise][section][error] Citations_overview: {exc}")
                logs.append(traceback.format_exc())
                slides.append(_error_slide(sec, exc))
        elif sec == "Citations_influence":
            try:
                p = dict(params or {})
                top_n = int(p.get("production_top_n", 10) or 10)
                citations_col = p.get("citations_col") if isinstance(p.get("citations_col"), str) else None
                payload = visual.add_citations_and_influence_section(
                    prs,
                    df,
                    collection_name=collection_name,
                    slide_notes=slide_notes,
                    citations_col=citations_col,
                    top_n=top_n,
                    return_payload=True,
                    export=False,
                    progress_callback=_cb,
                )
                slides.extend(_flatten_slides(payload, section=sec))
                logs.append(f"[visualise][section][ok] Citations_influence")
            except Exception as exc:
                logs.append(f"[visualise][section][error] Citations_influence: {exc}")
                logs.append(traceback.format_exc())
                slides.append(_error_slide(sec, exc))
        elif sec == "Words_and_topics":
            try:
                p0 = dict(params or {})
                view = str(p0.get("words_topics_view") or "both").strip().lower()
                data_source = str(p0.get("data_source") or defaults["word"].get("data_source") or "controlled_vocabulary_terms")
                word_plot_type = str(p0.get("word_plot_type") or defaults["word"].get("plot_type") or "bar_vertical")
                ngram_plot_type = str(p0.get("ngram_plot_type") or defaults["ngrams"].get("plot_type") or "bar_chart")

                include_has_ngrams = "Ngrams" in set(include or [])

                want_words = view in {"words", "both", ""}  # default to both
                want_ngrams = view in {"ngrams", "both"} and not include_has_ngrams

                if want_words:
                    wp = dict(defaults["word"], **p0)
                    wp["data_source"] = data_source
                    wp["plot_type"] = word_plot_type
                    payload = visual.analyze_wordanalysis_dispatcher(
                        df,
                        wp,
                        _cb,
                        slide_notes=slide_notes,
                        return_payload=True,
                        export=False,
                        zotero_client_for_pdf=None,
                        collection_name_for_cache=collection_name,
                    )
                    slides.extend(_flatten_slides(payload, section=sec))

                if want_ngrams:
                    np = dict(defaults["ngrams"], **p0)
                    np["data_source"] = data_source
                    np["plot_type"] = ngram_plot_type
                    if "top_n_ngrams" not in np and "top_ngram" in p0:
                        np["top_n_ngrams"] = p0.get("top_ngram")
                    payload2 = visual.analyze_ngrams(
                        df,
                        np,
                        _cb,
                        collection_name_for_cache=collection_name,
                        slide_notes=slide_notes,
                        return_payload=True,
                        export=False,
                    )
                    slides.extend(_flatten_slides(payload2, section=sec))

                if not want_words and not want_ngrams:
                    if view in {"ngrams"} and include_has_ngrams:
                        slides.append(
                            _normalize_slide(
                                {
                                    "title": "Words and topics — N-grams",
                                    "bullets": ["N-grams are generated in the N-grams section."],
                                    "notes": "",
                                    "section": sec,
                                }
                            )
                        )
                    elif view not in {"words", "ngrams", "both", ""}:
                        slides.append(
                            _normalize_slide(
                                {
                                    "title": "Words and topics — No selection",
                                    "bullets": [f"Unsupported Words & topics view: '{view}'"],
                                    "notes": "",
                                    "section": sec,
                                }
                            )
                        )
                logs.append(f"[visualise][section][ok] Words_and_topics")
            except Exception as exc:
                logs.append(f"[visualise][section][error] Words_and_topics: {exc}")
                logs.append(traceback.format_exc())
                slides.append(_error_slide(sec, exc))
        elif sec == "Ngrams":
            try:
                p0 = dict(params or {})
                p = dict(defaults["ngrams"], **p0)
                # UI uses `ngram_plot_type`; always honor it over defaults.
                if "ngram_plot_type" in p0:
                    p["plot_type"] = p0.get("ngram_plot_type")
                if "data_source" not in p and "data_source" in p0:
                    p["data_source"] = p0.get("data_source")
                if "top_n_ngrams" not in p and "top_ngram" in p0:
                    p["top_n_ngrams"] = p0.get("top_ngram")
                payload = visual.analyze_ngrams(
                    df,
                    p,
                    _cb,
                    collection_name_for_cache=collection_name,
                    slide_notes=slide_notes,
                    return_payload=True,
                    export=False,
                )
                slides.extend(_flatten_slides(payload, section=sec))
                logs.append(f"[visualise][section][ok] Ngrams")
            except Exception as exc:
                logs.append(f"[visualise][section][error] Ngrams: {exc}")
                logs.append(traceback.format_exc())
                slides.append(_error_slide(sec, exc))
        elif sec == "Affiliations_geo":
            try:
                p = dict(defaults["affiliations"], **(params or {}))
                payload = visual.analyze_affiliations(df, p, _cb, return_payload=True, export=False)
                slides.extend(_flatten_slides(payload, section=sec))
                logs.append(f"[visualise][section][ok] Affiliations_geo")
            except Exception as exc:
                logs.append(f"[visualise][section][error] Affiliations_geo: {exc}")
                logs.append(traceback.format_exc())
                slides.append(_error_slide(sec, exc))
        elif sec == "Geo_sector":
            logs.append("[visualise][section][warn] Geo_sector section disabled (removed from UI).")
        elif sec == "Temporal_analysis":
            try:
                p = dict(defaults["temporal"], **(params or {}))
                payload = visual.analyze_temporal_analysis(
                    df,
                    p,
                    _cb,
                    return_payload=True,
                    export=False,
                    collection_name=collection_name,
                )
                slides.extend(_flatten_slides(payload, section=sec))
                logs.append(f"[visualise][section][ok] Temporal_analysis")
            except Exception as exc:
                logs.append(f"[visualise][section][error] Temporal_analysis: {exc}")
                logs.append(traceback.format_exc())
                slides.append(_error_slide(sec, exc))
        elif sec == "Research_design":
            try:
                payload = visual.analyze_research_design_suite(df, params or {}, _cb, return_payload=True, export=False)
                slides.extend(_flatten_slides(payload, section=sec))
                logs.append(f"[visualise][section][ok] Research_design")
            except Exception as exc:
                logs.append(f"[visualise][section][error] Research_design: {exc}")
                logs.append(traceback.format_exc())
                slides.append(_error_slide(sec, exc))
        elif sec == "Thematic_method":
            try:
                p = dict(params or {})
                phase_tag = str(p.get("phase_tag") or "").strip()
                if not phase_tag:
                    for cand in ("phase_focus_value", "phase", "mission_phase", "phase_focus"):
                        if cand in df.columns:
                            phase_tag = cand
                            break
                year_col = str(p.get("year_col") or "").strip()
                if not year_col:
                    for cand in ("year", "publication_year", "year_numeric"):
                        if cand in df.columns:
                            year_col = cand
                            break
                payload = visual.add_thematic_and_method_section(
                    prs,
                    df,
                    collection_name=collection_name,
                    slide_notes=slide_notes,
                    cross_tabs=p.get("cross_tabs") if isinstance(p.get("cross_tabs"), list) else None,
                    phase_tag=phase_tag or str(p.get("phase_tag") or "phase_focus_value"),
                    year_col=year_col or str(p.get("year_col") or "publication_year"),
                    topn_cols=int(p.get("topn_cols", 14) or 14),
                    return_payload=True,
                    export=False,
                    progress_callback=_cb,
                )
                slides.extend(_flatten_slides(payload, section=sec))
                logs.append(f"[visualise][section][ok] Thematic_method")
            except Exception as exc:
                logs.append(f"[visualise][section][error] Thematic_method: {exc}")
                logs.append(traceback.format_exc())
                slides.append(_error_slide(sec, exc))
        elif sec == "Profiles":
            logs.append("[visualise][section][warn] Profiles section disabled (removed from UI).")
        elif sec == "Categorical_keywords":
            try:
                p = dict(defaults["categorical"], **(params or {}))
                results_df, fig = visual.analyze_categorical_keywords(df, p, _cb)
                table_html = "<div>No data available.</div>"
                try:
                    if getattr(results_df, "empty", True):
                        table_html = "<div>No data available.</div>"
                    else:
                        table_html = results_df.head(50).to_html(index=False, escape=True)
                except Exception:
                    table_html = "<div>Unable to render table.</div>"

                fig_json = None
                if fig is not None:
                    try:
                        import json as _json
                        from plotly.utils import PlotlyJSONEncoder

                        fig_json = _json.loads(_json.dumps(fig.to_plotly_json(), cls=PlotlyJSONEncoder))
                    except Exception as exc2:
                        logs.append(f"[visualise][section][warn] Categorical_keywords fig_json failed: {exc2}")

                slides.append(
                    _normalize_slide(
                        {
                            "title": "Categorical keywords",
                            "table_html": table_html,
                            "fig_json": fig_json,
                            "section": sec,
                            "notes": "",
                        }
                    )
                )
                logs.append(f"[visualise][section][ok] Categorical_keywords")
            except Exception as exc:
                logs.append(f"[visualise][section][error] Categorical_keywords: {exc}")
                logs.append(traceback.format_exc())
                slides.append(_error_slide(sec, exc))
        else:
            slides.append({"title": sec.replace("_", " "), "bullets": ["No handler for section."], "notes": "", "section": sec})
            logs.append(f"[visualise][section][warn] No handler for {sec}")

    # Assign stable ids so the renderer can keep focus between reruns with different params.
    try:
        by_section_count: dict[str, int] = {}
        for s in slides:
            if not isinstance(s, dict):
                continue
            if s.get("slide_id"):
                continue
            sec = str(s.get("section") or "").strip() or "Visualise"
            by_section_count[sec] = int(by_section_count.get(sec, 0)) + 1
            n = by_section_count[sec]
            title = str(s.get("title") or "")
            kind = "fig" if s.get("fig_json") else "table" if str(s.get("table_html") or "").strip() else "text"
            # Keep the ID independent of title so changing top-N params doesn't break selection/focus.
            s["slide_id"] = f"{sec}:{n}:{kind}"
    except Exception:
        # Never fail the deck due to metadata.
        pass

    # Build-mode filtering: keep all slides in selected sections OR explicitly selected slide ids.
    if is_build and isinstance(selection, dict):
        try:
            selected_sections = {
                str(x).strip() for x in (selection.get("sections") or []) if str(x).strip()
            }
            selected_slide_ids = {
                str(x).strip() for x in (selection.get("slideIds") or []) if str(x).strip()
            }
            if selected_sections or selected_slide_ids:
                slides = [
                    s
                    for s in slides
                    if isinstance(s, dict)
                    and (
                        (str(s.get("section") or "").strip() in selected_sections)
                        or (str(s.get("slide_id") or "").strip() in selected_slide_ids)
                    )
                ]
        except Exception as exc:
            logs.append(f"[visualise][selection][error] {exc}")
            logs.append(traceback.format_exc())

    if is_build:
        try:
            from python_backend.visualise.power_point_export import _print_ui_debug_dump

            def _fig_shape(fig_json: Any) -> str:
                if fig_json is None:
                    return "none"
                if isinstance(fig_json, dict):
                    ok = isinstance(fig_json.get("data"), list) and isinstance(fig_json.get("layout"), dict)
                    return "dict_ok" if ok else "dict_bad"
                if isinstance(fig_json, str):
                    return "string"
                return "other"

            dbg_slides: list[dict] = []
            for i, s in enumerate(slides):
                if not isinstance(s, dict):
                    continue
                fig_json = s.get("fig_json")
                fig_shape = _fig_shape(fig_json)
                dbg_slides.append(
                    {
                        "i": int(i),
                        "title": str(s.get("title") or "")[:120],
                        "keys": sorted(list(s.keys())),
                        "has_fig": bool(fig_json),
                        "has_img": bool(str(s.get("img") or "").strip()),
                        "has_table": bool(str(s.get("table_html") or "").strip()),
                        "has_text": bool(str(s.get("notes") or "").strip()),
                        "fig_shape": fig_shape,
                    }
                )

            _print_ui_debug_dump(
                {
                    "slides": dbg_slides,
                    "slides_len": len(dbg_slides),
                    "index": 0,
                    "title": str((slides[0] or {}).get("title") if slides else "Visualise"),
                }
            )
        except Exception as exc:
            logs.append(f"[visualise][ui_debug_dump][error] {exc}")
            logs.append(traceback.format_exc())

    if not slides:
        slides = [{"title": "Visualise", "bullets": ["No slides produced."], "notes": ""}]

    return {"slides": slides, "index": 0}, logs


def main() -> None:
    import os
    import sys

    def _handle(payload: dict) -> dict:
        action = str(payload.get("action") or "")
        if action == "sections":
            return {"status": "ok", **_schema_and_sections()}

        if action == "preview":
            table = payload.get("table") or {}
            params = payload.get("params") or {}
            include = payload.get("include") or []
            selection = payload.get("selection") if isinstance(payload.get("selection"), dict) else None
            collection_name = str(payload.get("collectionName") or "Collection")
            mode = str(payload.get("mode") or "")
            cache_logs: list[str] = []
            if not _is_table_payload(table):
                cached = _load_last_table_from_datahub_cache(logs=cache_logs)
                if cached is not None:
                    table = cached
            df = _coerce_table(table)
            deck, logs = _run_preview(
                df,
                params,
                [str(s) for s in include if str(s).strip()],
                collection_name,
                mode=mode,
                selection=selection,
            )
            return {"status": "ok", "deck": deck, "logs": [*cache_logs, *logs]}

        if action == "export_pptx":
            table = payload.get("table") or {}
            params = payload.get("params") or {}
            include = payload.get("include") or []
            selection = payload.get("selection") if isinstance(payload.get("selection"), dict) else None
            collection_name = str(payload.get("collectionName") or "Collection")
            outp_raw = str(payload.get("outputPath") or "").strip()
            notes_overrides = payload.get("notesOverrides") if isinstance(payload.get("notesOverrides"), dict) else {}
            rendered_images = payload.get("renderedImages") if isinstance(payload.get("renderedImages"), dict) else {}
            if not outp_raw:
                return {"status": "error", "message": "Missing outputPath"}
            cache_logs: list[str] = []
            if not _is_table_payload(table):
                cached = _load_last_table_from_datahub_cache(logs=cache_logs)
                if cached is not None:
                    table = cached
            df = _coerce_table(table)
            mode = "build_pptx"
            deck, logs = _run_preview(
                df,
                params,
                [str(s) for s in include if str(s).strip()],
                collection_name,
                mode=mode,
                selection=selection,
            )
            logs = [*cache_logs, *logs]
            slides = deck.get("slides") if isinstance(deck, dict) else None
            slides_list = slides if isinstance(slides, list) else []
            # Apply per-slide notes overrides (from UI "Describe" opt-in).
            try:
                for s in slides_list:
                    if not isinstance(s, dict):
                        continue
                    sid = str(s.get("slide_id") or "").strip()
                    if sid and sid in notes_overrides:
                        s["notes"] = str(notes_overrides.get(sid) or "").strip()
                    if sid and sid in rendered_images and str(rendered_images.get(sid) or "").strip():
                        s["img"] = str(rendered_images.get(sid) or "").strip()
            except Exception as exc:
                logs.append(f"[visualise][notes_override][warn] {exc}")
                logs.append(traceback.format_exc())

            slide_notes = str(params.get("slide_notes", "false")).strip().lower() in {"1", "true", "yes", "y", "on"} or bool(notes_overrides)
            try:
                saved = _export_pptx_from_slides(
                    [s for s in slides_list if isinstance(s, dict)],
                    output_path=Path(outp_raw).expanduser().resolve(),
                    slide_notes=slide_notes,
                    logs=logs,
                )
                return {"status": "ok", "path": str(saved), "logs": logs}
            except Exception as exc:
                logs.append(f"[visualise][export_pptx][fatal] {exc}")
                logs.append(traceback.format_exc())
                return {"status": "error", "message": str(exc), "logs": logs}

        if action == "describe_slide":
            slide = payload.get("slide") if isinstance(payload.get("slide"), dict) else {}
            params = payload.get("params") if isinstance(payload.get("params"), dict) else {}
            collection_name = str(payload.get("collectionName") or "Collection")
            text, logs = _describe_slide(slide, collection_name=collection_name, params=params)
            return {"status": "ok", "description": text, "logs": logs}

        return {"status": "error", "message": f"unknown_action:{action}"}

    if str(os.environ.get("VISUALISE_SERVER") or "").strip() in {"1", "true", "yes", "on"}:
        for line in sys.stdin:
            raw = (line or "").strip()
            if not raw:
                continue
            try:
                incoming = json.loads(raw)
            except Exception as exc:
                sys.stdout.write(_safe_json({"status": "error", "message": f"invalid_json:{exc}"}) + "\n")
                sys.stdout.flush()
                continue
            try:
                out = _handle(incoming if isinstance(incoming, dict) else {})
            except Exception as exc:
                out = {"status": "error", "message": str(exc), "trace": traceback.format_exc()}
            sys.stdout.write(_safe_json(out) + "\n")
            sys.stdout.flush()
        return

    payload = _read_payload()
    try:
        out = _handle(payload if isinstance(payload, dict) else {})
    except Exception as exc:
        out = {"status": "error", "message": str(exc), "trace": traceback.format_exc()}
    print(_safe_json(out))


if __name__ == "__main__":
    main()
