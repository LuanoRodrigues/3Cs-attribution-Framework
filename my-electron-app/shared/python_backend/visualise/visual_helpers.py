import base64
import functools
import html
import io
import itertools
import json
import logging
import re
import sys
from collections import Counter, defaultdict
from pathlib import Path
from typing import List, Optional, Tuple

import networkx as nx
import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from dotenv import load_dotenv
from plotly.graph_objs import Figure
from plotly.utils import PlotlyJSONEncoder
from pptx.enum.text import PP_ALIGN  # <-- Correct name
from pptx.util import Emu, Inches, Pt
from wordcloud import WordCloud

BASE_DIR = Path(__file__).resolve().parent
REPO_ROOT = next((parent for parent in BASE_DIR.parents if (parent / "package.json").is_file()), BASE_DIR.parents[-1])
for entry in (BASE_DIR, REPO_ROOT):
    entry_str = str(entry)
    if entry_str not in sys.path:
        sys.path.insert(0, entry_str)

from .data_processing import get_text_corpus_from_df, preprocess_text, _country_counts, \
    _make_world_map, _centroid, _split_institutions, _set_notes, _calculate_all_author_stats, get_document_type, \
    _geocode, _get_ngrams_from_text, _choose_corpus_text, _scan_text_for_keywords, _normalise_and_strip_html, \
    tokenize_for_ngram_context_refined, _get_ngram_data_by_year
from python_backend.core.utils.calling_models import call_models_plots

# Zotero Setup
load_dotenv()
ZOTERO_CLASS_AVAILABLE = False
zot = None  # This will be the usable Zotero instance for this module
# Cache Directory Setup for this module

# from bibliometric_analysis_tool.utils.zotero_class import Zotero


APP_CACHE_DIR = Path(".") / ".bibliometric_tool_dp_cache"  # Fallback
try:
    APP_CACHE_DIR.mkdir(parents=True, exist_ok=True)
except Exception as e_mkdir:
    logging.error(f"Failed to create fallback cache dir: {e_mkdir}")

if 'MISTRAL_API_KEY_ENV_VAR' not in globals(): MISTRAL_API_KEY_ENV_VAR = "MISTRAL_API_KEY"

if 'STOP_WORDS' not in globals(): STOP_WORDS = set(["the", "a", "is", "of", "and"])  # Minimal


def analyze_words_over_time(df: pd.DataFrame, params: dict, progress_callback=None,
                            zotero_client_for_pdf=None,
                            collection_name_for_cache=None):
    """
    Analyze frequency of top or specified words over time (keywords, abstract, fulltext, or title) and produce
    a responsive Plotly time-series with both user-specified and top keywords.
    Uses per-document token caching to avoid repeated full-text extraction.
    """
    def _cb(msg: str):
        if progress_callback:
            progress_callback(f"WordsOverTime: {msg}")
        logging.info(f"WordsOverTime: {msg}")

    # 1. Params & validation
    source = params.get("data_source", "controlled_vocabulary_terms")
    top_n = params.get("num_top_words", 5)
    specific = [w.strip().lower() for w in (params.get("specific_words_to_track") or []) if w.strip()]
    caption = params.get("plot_caption", "")
    _cb(f"Analyzing word usage over time from '{source}'.")

    valid_sources = {"controlled_vocabulary_terms", "abstract", "fulltext", "title"}
    if source not in valid_sources:
        fig = go.Figure().add_annotation(text=f"Unsupported source '{source}'", showarrow=False)
        return pd.DataFrame(), fig

    text_col = 'key' if source == 'fulltext' else source
    if 'year' not in df.columns or (text_col not in df.columns and source != 'title'):
        fig = go.Figure().add_annotation(text="Missing 'year' or source column", showarrow=False)
        return pd.DataFrame(), fig

    df2 = df.copy()
    df2['year'] = pd.to_numeric(df2['year'], errors='coerce').astype('Int64')
    df2 = df2.dropna(subset=['year']).astype({'year': int})
    if df2.empty:
        fig = go.Figure().add_annotation(text="No valid year data.", showarrow=False)
        return pd.DataFrame(), fig

    # 2. Token extraction per document (cache tokens in column)
    def extract_for_row(row):
        subdf = pd.DataFrame([row])
        if source == 'fulltext':
            return get_text_corpus_from_df(subdf, source,
                                           collection_name_for_cache, zotero_client_for_pdf, None)
        if source == 'title':
            return preprocess_text(str(row.get('title', '')))
        return get_text_corpus_from_df(subdf, source,
                                       collection_name_for_cache, zotero_client_for_pdf, None)

    _cb("Tokenizing each document once...")
    df2['tokens'] = df2.apply(extract_for_row, axis=1)
    first_count = len(df2.iloc[0]['tokens']) if not df2.empty else 0
    _cb(f"Tokenization complete: first doc has {first_count} tokens.")

    # 3. Determine words to track
    all_tokens = [tok for tokens in df2['tokens'] for tok in tokens]
    counts = Counter(all_tokens)
    top_words = [w for w, _ in counts.most_common(top_n) if w not in specific]
    if specific:
        words = specific + top_words
        _cb(f"Tracking specified + top words: {words}")
    else:
        words = top_words
        _cb(f"Tracking top words: {words}")

    if not words:
        fig = go.Figure().add_annotation(text="No words to track.", showarrow=False)
        return pd.DataFrame(), fig

    # 4. Aggregate yearly frequencies using precomputed tokens
    records = []
    for year_val, grp in df2.groupby('year'):
        year_tokens = [tok for tokens in grp['tokens'] for tok in tokens]
        cnt = Counter(year_tokens)
        for w in words:
            records.append({'Year': year_val, 'Word': w, 'Frequency': cnt.get(w, 0)})
    time_df = pd.DataFrame(records)

    # 5. Plotting with desired style
    fig = go.Figure()
    palette = px.colors.qualitative.Plotly
    for i, w in enumerate(words):
        sub = time_df[time_df['Word'] == w].sort_values('Year')
        fig.add_trace(go.Scatter(
            x=sub['Year'], y=sub['Frequency'], name=w,
            mode='lines+markers',
            line=dict(width=2, shape='linear', color=palette[i % len(palette)]),
            marker=dict(size=6, symbol='circle')
        ))

    # Optional caption
    annotations = []
    if caption:
        annotations.append(dict(
            text=caption,
            showarrow=False,
            xref='paper', yref='paper',
            x=0.5, y=-0.15,
            xanchor='center', yanchor='top',
            font=dict(size=12, color='grey')
        ))

    # 6. Layout: white background, blue grid squares, square ticks
    fig.update_layout(
        autosize=True,
        # template='plotly_white',
        # paper_bgcolor='white',
        # plot_bgcolor='white',
        margin=dict(l=30, r=30, t=60, b=60),
        title=dict(text=f"Word Trends Over Time ({source.title()})", x=0.5,
                   font=dict(size=20, family='Arial', color='#1f77b4')),
        xaxis=dict(title='Year', dtick=1,
                   showgrid=True, gridcolor='#cfe3f3', gridwidth=1,
                   zeroline=False, showline=True, linecolor='#1f77b4', linewidth=1,
                   tickfont=dict(size=11, color='#1f77b4'),
                   ticks='outside', ticklen=5),
        yaxis=dict(title='Frequency',
                   showgrid=True, gridcolor='#cfe3f3', gridwidth=1,
                   zeroline=False, showline=True, linecolor='#1f77b4', linewidth=1,
                   tickfont=dict(size=11, color='#1f77b4'),
                   ticks='outside', ticklen=5),
        legend=dict(orientation='h', y=1.02, x=0.5, xanchor='center',
                    font=dict(size=11, color='#1f77b4')),
        annotations=annotations
    )

    _cb("Finished words over time.")
    return time_df, fig

def add_data_summary_slides(
    prs,
    df: "pd.DataFrame",
    collection_name: str = "Collection",
    *,
    top_n: int = 3,
    slide_notes: bool = True,
    year_min: int = 1500,
    return_payload: bool = False,
    export: bool = False,
):
    """
    ###1. compute summary metrics from an item-level deduplicated dataframe
    ###2. always build a Plotly "card grid" figure that renders cleanly in the UI (no overflow)
    ###3. if export: also render the same card layout into a pptx Presentation (one slide)
    ###4. if return_payload: return {"slides":[...], "figs":[...], "index":0, "payload":{...}}

    Assumptions (brittle by design):
      - df is a pandas.DataFrame (may be empty; empty returns a 1-slide payload if preview enabled).
      - Plotly is available; go.Figure().to_plotly_json() works.
      - If export=True, prs is a python-pptx Presentation.
    """
    import re
    import json

    import numpy as np
    import pandas as pd
    import plotly.graph_objects as go
    from plotly.utils import PlotlyJSONEncoder

    def _explode_semicolon_or_list(series: pd.Series) -> pd.Series:
        if series is None or series.empty:
            return pd.Series([], dtype=str)
        s = series.dropna()
        s = s.apply(lambda x: list(x) if type(x) in (set, tuple) else x)
        is_list = s.apply(lambda x: type(x) is list)
        if bool(is_list.any()):
            s = pd.concat([s[~is_list], s[is_list].explode()])
        s = s.astype(str).str.split(";").explode().str.strip()
        s = s.replace("", pd.NA).dropna()
        return s

    def _top_n_from_col(series: pd.Series, n: int = 3) -> list[tuple[str, int]]:
        if series is None or series.empty:
            return []
        s = _explode_semicolon_or_list(series)
        if s.empty:
            return []
        vc = s.value_counts().head(int(n))
        return [(str(k), int(v)) for k, v in vc.items()]

    def _best_year_bounds(df_: pd.DataFrame, *, ymin: int = 1500) -> tuple[str, str]:
        years = None
        for col in ("year_numeric", "year"):
            if col in df_.columns:
                y = pd.to_numeric(df_[col], errors="coerce")
                years = y if years is None else years.fillna(y)
        if years is None:
            return "N/A", "N/A"
        years = years.dropna().astype(float)
        if years.empty:
            return "N/A", "N/A"
        current_year = pd.Timestamp.today().year + 1
        years = years[(years >= ymin) & (years <= current_year)]
        if years.empty:
            return "N/A", "N/A"
        return str(int(years.min())), str(int(years.max()))

    def _clean_numeric_series(raw: pd.Series) -> pd.Series:
        s = raw.astype(str).str.replace(",", "", regex=False).str.strip()
        s = s.str.extract(r"([0-9]*\.?[0-9]+)", expand=False)
        return pd.to_numeric(s, errors="coerce")

    def _best_citations_series(df_: pd.DataFrame) -> pd.Series | None:
        preferred = ["citations", "cited_by", "cited_by_count", "citations_count"]
        pattern_hits = [c for c in df_.columns if re.search(r"(citat|cited|times_cited)", str(c), re.I)]
        scan = preferred + [c for c in pattern_hits if c not in preferred]
        best, best_n = None, -1
        for col in scan:
            if col in df_.columns:
                s = _clean_numeric_series(df_[col]).dropna()
                if (not s.empty) and (int(s.shape[0]) > int(best_n)):
                    best, best_n = s, int(s.shape[0])
        return best

    def _fmt_number(x: float | int | str) -> str:
        if type(x) in (int, np.integer):
            return f"{int(x):,}"
        if type(x) in (float, np.floating):
            return f"{float(x):,.2f}"
        return str(x)

    def _keywords_series(df_: pd.DataFrame) -> pd.Series:
        candidates = [
            "controlled_vocabulary_terms", "controlled vocabulary terms",
            "keywords", "keyword", "tags", "terms", "topics",
        ]
        possible = [c for c in df_.columns if any(re.search(rf"\b{re.escape(k)}\b", str(c), re.I) for k in candidates)]
        cols = [c for c in candidates if c in df_.columns] + possible

        out = []
        for col in cols:
            if col not in df_.columns:
                continue
            if bool(df_[col].isna().all()):
                continue
            ser = _explode_semicolon_or_list(df_[col]).astype(str)
            theme_vals = ser.str.extract(r"^\s*#\s*theme\s*:\s*(.+?)\s*$", flags=re.I, expand=False)
            theme_vals = theme_vals.dropna().astype(str).str.strip()
            theme_vals = theme_vals.str.replace(r"\s+", " ", regex=True)
            theme_vals = theme_vals.replace("", pd.NA).dropna()
            if not theme_vals.empty:
                out.append(theme_vals)

        if not out:
            return pd.Series([], dtype=str)

        ser_all = pd.concat(out, ignore_index=True)
        ser_all = ser_all.str.strip().str.casefold()
        ser_all = ser_all.replace("", pd.NA).dropna()
        return ser_all

    def _year_series_local(frame: pd.DataFrame) -> pd.Series:
        ycol = "year_numeric" if "year_numeric" in frame.columns else ("year" if "year" in frame.columns else None)
        if ycol is None:
            return pd.Series([np.nan] * len(frame), index=frame.index)
        return pd.to_numeric(frame[ycol], errors="coerce")

    def _norm_title_for_key(s: str) -> str:
        s = "" if s is None or (type(s) is float and pd.isna(s)) else str(s)
        s = s.lower()
        s = re.sub(r"[^a-z0-9]+", " ", s)
        s = re.sub(r"\s+", " ", s).strip()
        return s

    def _wrap_line(s: str, *, max_chars: int) -> list[str]:
        t = "" if s is None else str(s)
        t = " ".join(t.split())
        if not t:
            return [""]
        words = t.split(" ")
        out = []
        cur = ""
        for w in words:
            if not cur:
                cur = w
                continue
            if (len(cur) + 1 + len(w)) <= int(max_chars):
                cur = cur + " " + w
            else:
                out.append(cur)
                cur = w
        out.append(cur)
        return out

    def _fmt_top_lines(pairs: list[tuple[str, int]]) -> list[str]:
        if not pairs:
            return ["N/A"]
        out = []
        for name, count in pairs:
            out.append(f"{name} ({count})")
        return out

    items = df.copy()

    if ("key" in items.columns) and bool(items["key"].notna().any()):
        items = items.sort_values("key").drop_duplicates("key", keep="first").copy()
    else:
        tkey = items["title"].map(_norm_title_for_key) if "title" in items.columns else pd.Series([""] * len(items), index=items.index)
        yser = _year_series_local(items)
        ykey = yser.fillna(-1).astype("int64")
        items = (
            items.assign(__tkey=tkey, __ykey=ykey)
            .sort_values(["__tkey", "__ykey"])
            .drop_duplicates(["__tkey", "__ykey"], keep="first")
            .drop(columns=["__tkey", "__ykey"])
            .copy()
        )

    total_docs = int(len(items))

    authors_col = items["authors"] if "authors" in items.columns else pd.Series(dtype=str)
    exploded_authors = authors_col.dropna().astype(str).str.split(";").explode().str.strip()
    exploded_authors = exploded_authors[exploded_authors != ""]
    unique_auth = int(exploded_authors.str.casefold().nunique()) if not exploded_authors.empty else 0

    if "citations" in items.columns:
        def _find_citations_col_local(frame: pd.DataFrame) -> str | None:
            preferred = ["citations", "cited_by", "cited_by_count", "citations_count"]
            for p in preferred:
                if p in frame.columns:
                    return p
            pats = (r"^citations?$", r"^cites?$", r"times_?cited", r"^n_?cit", r"^num_?cit", r"cited_by")
            for c in frame.columns:
                if any(re.search(p, str(c), flags=re.I) for p in pats):
                    return c
            return None

        _cit_col = _find_citations_col_local(items)
        cit_series = pd.to_numeric(items[_cit_col].copy(), errors="coerce") if _cit_col is not None else pd.Series(dtype=float)
    else:
        _best = _best_citations_series(items)
        cit_series = _best.copy() if _best is not None else pd.Series(dtype=float)

    cit_series = cit_series.replace([np.inf, -np.inf], np.nan).fillna(0.0)

    tot_cit_str = "0"
    avg_cit_str = "0.00"
    if not cit_series.empty:
        tot_cit_str = _fmt_number(float(cit_series[cit_series > 0].sum()))
        avg_cit_str = f"{float(cit_series.mean()):.2f}"

    min_year, max_year = _best_year_bounds(items, ymin=year_min)

    top_types = _top_n_from_col(items["item_type"], n=top_n) if "item_type" in items.columns else []
    top_sources = _top_n_from_col(items["source"], n=top_n) if "source" in items.columns else []
    top_authors = _top_n_from_col(exploded_authors, n=top_n) if not exploded_authors.empty else []

    kw_series = _keywords_series(items)
    top_keywords = []
    if not kw_series.empty:
        vc = kw_series.value_counts().head(int(top_n))
        top_keywords = [(f"{k}", int(v)) for k, v in vc.items()]

    ai_text = ""
    if slide_notes:
        stats_payload = {
            "collection": str(collection_name or ""),
            "total_documents": int(total_docs),
            "unique_authors": int(unique_auth),
            "total_citations": str(tot_cit_str),
            "avg_citations_per_doc": str(avg_cit_str),
            "year_span": [str(min_year), str(max_year)],
            "filters": {"year_min_inclusive": int(year_min)},
            "top_n": int(top_n),
            "top_document_types": top_types,
            "top_authors": top_authors,
            "top_sources": top_sources,
            "top_keywords": top_keywords,
        }

        prompt_text = (
            "[TASK]\n"
            "You are a senior bibliometrics analyst. Draft 1–2 minutes of SPEAKER NOTES for a title-less summary slide.\n"
            "Deliver exactly:\n"
            " • 3–5 executive takeaways (bullet-like sentences)\n"
            " • context & caveats (coverage bias, year filter, name disambiguation, unnormalised citations)\n"
            " • notable concentrations AND notable absences\n"
            " • 3 practical next steps for validation/action\n\n"
            "[CONTEXT]\n"
            f"Collection: {collection_name}\n"
            f"Cards: Total Documents, Unique Authors, Total Citations, Publication Timespan, "
            f"Top {top_n} Document Types, Top {top_n} Authors, Top {top_n} Sources, Top {top_n} Keywords\n\n"
            "[RAW STATS — JSON]\n"
            f"{json.dumps(stats_payload, ensure_ascii=False)}\n"
        )

        import os
        ai_out = call_models_plots(
            mode="analyze",
            prompt_text=prompt_text,
            model_api_name=os.getenv("OPENAI_VISION_MODEL", "gpt-5.1-mini"),
            images=None,
            image_detail="low",
            max_tokens=900,
            use_cache=True,
            analysis_key_suffix="data_summary_notes",
            section_title=collection_name or "collection",
            store_only=False,
            read=False,
        )
        ai_text = (ai_out.get("text") or "").strip()

        if not ai_text:
            def _fmt_pairs(pairs):
                return ", ".join([f"{k} ({v})" for k, v in pairs]) if pairs else "N/A"

            ai_text = (
                f"{collection_name} • Summary. We count {total_docs:,} documents by ~{unique_auth:,} unique authors, "
                f"with {tot_cit_str} total citations (avg {avg_cit_str} per item). "
                f"Timespan {min_year}–{max_year} (values < {year_min} ignored). "
                f"Top types: {_fmt_pairs(top_types)}. Top authors: {_fmt_pairs(top_authors)}. "
                f"Top sources: {_fmt_pairs(top_sources)}. Top keywords: {_fmt_pairs(top_keywords)}. "
                "Caveats: coverage/metadata quality, author name disambiguation, and unnormalised citations may bias results. "
                "Next steps: (1) disambiguate authors; (2) field- and year-normalise citations; "
                "(3) audit a sample of item types/keywords for precision."
            )

    payload = {
        "section": "Data_summary",
        "collection_name": str(collection_name or ""),
        "metrics": {
            "total_docs": int(total_docs),
            "unique_authors": int(unique_auth),
            "total_citations_str": str(tot_cit_str),
            "avg_citations_str": str(avg_cit_str),
            "min_year": str(min_year),
            "max_year": str(max_year),
            "year_min_filter": int(year_min),
        },
        "top": {
            "document_types": top_types,
            "authors": top_authors,
            "sources": top_sources,
            "keywords": top_keywords,
        },
        "notes_text": str(ai_text or ""),
    }

    # -----------------------------
    # Plotly UI figure (clean card grid)
    # -----------------------------
    def _card(fig, x0, y0, x1, y1, heading: str, lines: list[str]):
        fig.add_shape(
            type="rect",
            xref="paper",
            yref="paper",
            x0=float(x0),
            y0=float(y0),
            x1=float(x1),
            y1=float(y1),
            line={"width": 1, "color": "rgba(205,205,205,1)"},
            fillcolor="rgba(247,247,247,1)",
            layer="below",
        )
        fig.add_annotation(
            xref="paper",
            yref="paper",
            x=float(x0 + 0.012),
            y=float(y1 - 0.018),
            xanchor="left",
            yanchor="top",
            text=str(heading).upper(),
            showarrow=False,
            align="left",
            font={"size": 12, "family": "Segoe UI", "color": "rgba(95,95,95,1)"},
        )

        # Typography + spacing tuned to prevent overflow.
        # We render up to 4 wrapped lines; extras are truncated by design.
        max_lines = 4
        inner_top = float(y1 - 0.070)
        inner_bottom = float(y0 + 0.025)
        avail = float(inner_top - inner_bottom)

        # If only a single value line, make it larger.
        base_size = 22 if len(lines) == 1 else 15
        step = float(avail / max(1, max_lines))

        y = float(inner_top)
        for i, ln in enumerate(lines[:max_lines]):
            fig.add_annotation(
                xref="paper",
                yref="paper",
                x=float(x0 + 0.012),
                y=float(y),
                xanchor="left",
                yanchor="top",
                text=str(ln),
                showarrow=False,
                align="left",
                font={"size": int(base_size if i == 0 else max(12, base_size - 2)), "family": "Segoe UI", "color": "rgba(30,30,30,1)"},
            )
            y = float(y - step)

    def _top_block_lines(pairs: list[tuple[str, int]], *, max_chars: int) -> list[str]:
        raw = _fmt_top_lines(pairs)
        # Wrap each entry; keep at most 4 rendered lines total.
        out = []
        for r in raw:
            for w in _wrap_line(r, max_chars=int(max_chars)):
                out.append(w)
                if len(out) >= 4:
                    return out
        return out[:4] if out else ["N/A"]

    fig = go.Figure()
    fig.update_layout(
        margin={"l": 26, "r": 26, "t": 18, "b": 18},
        paper_bgcolor="white",
        plot_bgcolor="white",
        xaxis={"visible": False},
        yaxis={"visible": False},
        width=1280,
        height=720,
        showlegend=False,
    )
    # Ensure the "card grid" is treated as a real Plotly figure everywhere (thumb snapshot/export paths often
    # key off trace types). This trace is invisible and does not affect layout.
    fig.add_trace(
        go.Scatter(
            x=[0],
            y=[0],
            mode="markers",
            marker={"opacity": 0},
            hoverinfo="skip",
            showlegend=False,
        )
    )

    # Grid geometry in paper coords
    left_x0, left_x1 = 0.03, 0.49
    right_x0, right_x1 = 0.51, 0.97
    y_tops = [0.96, 0.73, 0.50, 0.27]
    h = 0.205

    _card(fig, left_x0, y_tops[0] - h, left_x1, y_tops[0], "Total Documents", [_fmt_number(total_docs)])
    _card(fig, left_x0, y_tops[1] - h, left_x1, y_tops[1], "Unique Authors", [_fmt_number(unique_auth)])
    _card(fig, left_x0, y_tops[2] - h, left_x1, y_tops[2], "Total Citations", [str(tot_cit_str)])
    _card(fig, left_x0, y_tops[3] - h, left_x1, y_tops[3], "Publication Timespan", [f"{min_year}–{max_year}"])

    _card(fig, right_x0, y_tops[0] - h, right_x1, y_tops[0], f"Top Document Types (Top {top_n})",
          _top_block_lines(top_types, max_chars=34))
    _card(fig, right_x0, y_tops[1] - h, right_x1, y_tops[1], f"Top Authors (Top {top_n})",
          _top_block_lines(top_authors, max_chars=34))
    _card(fig, right_x0, y_tops[2] - h, right_x1, y_tops[2], f"Top Sources (Top {top_n})",
          _top_block_lines(top_sources, max_chars=34))
    _card(fig, right_x0, y_tops[3] - h, right_x1, y_tops[3], f"Top Keywords (Top {top_n})",
          _top_block_lines(top_keywords, max_chars=34))

    fig_json_str = json.dumps(fig.to_plotly_json(), cls=PlotlyJSONEncoder)

    slide_spec = {
        "title": f"{collection_name} — Data summary",
        "fig_json": fig_json_str,
        "table_html": "",
        "notes": ai_text if slide_notes else "",
    }

    ui_payload = {
        "slides": [slide_spec],
        "figs": [fig_json_str],
        "index": 0,
        "imgs": [],
        "tables": [],
        "payload": payload,
    }

    if (not export) or (prs is None):
        if return_payload:
            return ui_payload
        return [fig]

    # -----------------------------
    # PPTX export (same card grid)
    # -----------------------------
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_AUTO_SIZE

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    def _add_card(
        x_in,
        y_in,
        w_in,
        h_in,
        heading: str,
        *,
        value: str | None = None,
        lines: list[str] | None = None,
    ):
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(float(x_in)),
            Inches(float(y_in)),
            Inches(float(w_in)),
            Inches(float(h_in)),
        )
        fill = shp.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(247, 247, 247)
        shp.line.color.rgb = RGBColor(205, 205, 205)

        tf = shp.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = Pt(10)
        tf.margin_right = Pt(10)
        tf.margin_top = Pt(8)
        tf.margin_bottom = Pt(6)

        p0 = tf.paragraphs[0]
        p0.text = str(heading).upper()
        p0.font.size = Pt(10)
        p0.font.bold = True
        p0.font.color.rgb = RGBColor(95, 95, 95)
        p0.space_after = Pt(2)

        if lines is not None:
            # Reduce font to prevent overflow; each line is already wrapped in the Plotly path
            for line in lines[:4]:
                p = tf.add_paragraph()
                p.text = str(line)
                p.level = 0
                p.font.size = Pt(11)
                p.font.bold = False
                p.font.color.rgb = RGBColor(30, 30, 30)
                p.space_before = Pt(0)
                p.space_after = Pt(0)
        else:
            p1 = tf.add_paragraph()
            p1.text = str(value or "")
            p1.level = 0
            p1.font.size = Pt(22)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(30, 30, 30)
            p1.space_before = Pt(0)
            p1.space_after = Pt(0)

        return shp

    # Slide geometry tuned for default widescreen (13.333x7.5). Works acceptably for other sizes.
    left_x, right_x = 0.50, 5.20
    col_w = 4.55
    top_y = 0.55
    card_h = 1.23
    gutter_y = 0.20

    _add_card(left_x, top_y + 0 * (card_h + gutter_y), col_w, card_h, "Total Documents", value=_fmt_number(total_docs))
    _add_card(left_x, top_y + 1 * (card_h + gutter_y), col_w, card_h, "Unique Authors", value=_fmt_number(unique_auth))
    _add_card(left_x, top_y + 2 * (card_h + gutter_y), col_w, card_h, "Total Citations", value=str(tot_cit_str))
    _add_card(left_x, top_y + 3 * (card_h + gutter_y), col_w, card_h, "Publication Timespan", value=f"{min_year}–{max_year}")

    _add_card(right_x, top_y + 0 * (card_h + gutter_y), col_w, card_h, f"Top Document Types (Top {top_n})",
              lines=_top_block_lines(top_types, max_chars=34))
    _add_card(right_x, top_y + 1 * (card_h + gutter_y), col_w, card_h, f"Top Authors (Top {top_n})",
              lines=_top_block_lines(top_authors, max_chars=34))
    _add_card(right_x, top_y + 2 * (card_h + gutter_y), col_w, card_h, f"Top Sources (Top {top_n})",
              lines=_top_block_lines(top_sources, max_chars=34))
    _add_card(right_x, top_y + 3 * (card_h + gutter_y), col_w, card_h, f"Top Keywords (Top {top_n})",
              lines=_top_block_lines(top_keywords, max_chars=34))

    if slide_notes and ai_text:
        slide.notes_slide.notes_text_frame.text = ai_text

    if return_payload:
        return ui_payload

    return None


def _outside_text_headroom(
    fig: "go.Figure",
    *,
    orientation: str,           # "h" for horizontal bars, "v" for vertical bars
    values,                      # 1-D iterable of numeric bar magnitudes
    labels=None,                 # 1-D iterable of text labels (optional)
    pad_frac: float = 0.28,      # enlarge axis by ~28% to fit text tails
    extra_margin_px: int = 160   # add right/top margin in px to avoid PNG trimming
):
    """
    Expands the plotting range + margins so Plotly won't trim bar-end text like "123 (45.6%)".
    Call AFTER building traces and BEFORE export.
    """
    import numpy as _np
    vals = _np.asarray(values, dtype=float)
    vals = vals[_np.isfinite(vals)]
    if vals.size == 0:
        return

    vmax = float(vals.max())
    if vmax <= 0:
        vmax = 1.0

    if orientation.lower().startswith("h"):
        # Extend x-range and right margin
        fig.update_xaxes(range=[0, vmax * (1.0 + max(0.05, pad_frac))])
        m = fig.layout.margin.to_plotly_json() if fig.layout.margin is not None else {}
        fig.update_layout(margin=dict(
            l=int(m.get("l", 70)),
            r=int(max(m.get("r", 0), extra_margin_px)),
            t=int(m.get("t", 70)),
            b=int(m.get("b", 60)),
        ))
        # Make sure the text is actually placed "outside"
        fig.update_traces(selector=dict(type="bar"), textposition="outside")
    else:
        # Vertical bars: extend y-range and top margin
        fig.update_yaxes(range=[0, vmax * (1.0 + max(0.05, pad_frac))])
        m = fig.layout.margin.to_plotly_json() if fig.layout.margin is not None else {}
        fig.update_layout(margin=dict(
            l=int(m.get("l", 70)),
            r=int(m.get("r", 80)),
            t=int(max(m.get("t", 0), extra_margin_px)),
            b=int(m.get("b", 70)),
        ))
        fig.update_traces(selector=dict(type="bar"), textposition="outside")


def shape_scope(
    prs,
    df: "pd.DataFrame",
    collection_name: str = "Collection",
    *,
    slide_notes: bool = True,
    year_range: tuple[int, int] = (2010, 2025),
    top_n_cats: int = 20,
    include_percentages: bool = True,
    export: bool = False,
    return_payload: bool = False,
    progress_callback=None,
):
    preview_slides = [] if (export or return_payload) else None

    """
    Build 'Scope' slides (2 per metric: a Table slide, then a Figure slide):

      1) Publications by year
      2) Distribution by document type
      3) Split by publisher type

    Tables and figures are centered and constrained to fit the slide.
    """
    # ---- imports (kept local) ----


    # ---------- progress ----------
    def _cb(msg: str):
        if progress_callback is not None:
            progress_callback(f"SCOPE: {msg}")

    # ---------- global style ----------
    FONT_FAMILY = "Segoe UI"
    COLORWAY = ["#2563EB", "#10B981", "#F59E0B", "#EF4444", "#8B5CF6", "#14B8A6", "#F97316", "#3B82F6"]
    PLOT_W, PLOT_H = 1200, 650  # render slightly smaller -> crisper + easier to fit

    INCH = float(Inches(1))  # EMU-per-inch to float
    SLIDE_W_IN = float(prs.slide_width) / INCH
    SLIDE_H_IN = float(prs.slide_height) / INCH

    # margins (inches)
    OUTER_MARGIN_X = 0.6
    OUTER_MARGIN_Y = 0.6
    TITLE_HEIGHT_IN = 0.7

    # Consistent Plotly defaults
    def _apply_fig_style(fig: "go.Figure", title: str):
        fig.update_layout(
            template="plotly_white",
            width=PLOT_W, height=PLOT_H,
            margin=dict(l=60, r=40, t=60, b=50),
            title=dict(text=title, x=0.02, xanchor="left", font=dict(size=20, family=FONT_FAMILY)),
            font=dict(family=FONT_FAMILY, size=13),
            colorway=COLORWAY,
        )
        return fig

    # ---------- safe helpers (reuse if present globally) ----------
    _ensure_numeric_local = globals().get("_ensure_numeric", None)
    _safe_ai_notes_local = globals().get("_safe_ai_notes", None)
    _write_notes_local   = globals().get("_write_notes", None)

    def _ensure_numeric_df(x: "pd.DataFrame | None", fallback=("Value",), min_rows=1):
        if callable(_ensure_numeric_local):
            return _ensure_numeric_local(x, fallback_cols=fallback, min_rows=min_rows)
        if x is None or not isinstance(x, pd.DataFrame) or x.empty:
            return pd.DataFrame({fallback[0]: [0]})
        out = pd.DataFrame({c: pd.to_numeric(x[c], errors="coerce") for c in x.columns}).fillna(0)
        if out.empty:
            out = pd.DataFrame({fallback[0]: [0]})
        return out

    # ---------- parsing helpers ----------
    def _year_series(frame: pd.DataFrame) -> pd.Series:
        for cand in ("publication_year", "year_numeric", "year", "Year", "pub_year", "issued", "date"):
            if cand in frame.columns:
                s = pd.to_numeric(frame[cand], errors="coerce")
                if s.notna().any():
                    return s
                s2 = pd.to_datetime(frame[cand], errors="coerce", utc=True)
                if s2.notna().any():
                    return pd.to_numeric(s2.dt.year, errors="coerce")
        for c in frame.columns:
            if any(k in str(c).lower() for k in ("year", "date", "issued", "created", "publication")):
                s = pd.to_datetime(frame[c], errors="coerce", utc=True)
                if s.notna().any():
                    return pd.to_numeric(s.dt.year, errors="coerce")
        return pd.Series([np.nan] * len(frame), index=frame.index)

    def _doc_type(frame: pd.DataFrame) -> pd.Series:
        """
        Prefer 'document_type_value' if present; fall back to common type columns.
        """
        if "document_type_value" in frame.columns:
            s = frame["document_type_value"].astype(str).str.strip()
            if s.notna().any():
                # Guard against coded/ordinal doc-type columns (e.g. 0..N) that render as numbers.
                # If it mostly looks numeric, prefer the human-readable type columns instead.
                try:
                    non_empty = s.replace("nan", "").replace("None", "").replace("", pd.NA).dropna()
                    if not non_empty.empty:
                        frac_numeric = float(non_empty.str.fullmatch(r"\d+").mean())
                        if frac_numeric < 0.8:
                            return s
                except Exception:
                    return s
        for cand in ("document_type", "item_type", "itemType", "item type", "type"):
            if cand in frame.columns:
                s = frame[cand].astype(str).str.strip()
                if s.notna().any():
                    return s
        return pd.Series(["Unknown"] * len(frame), index=frame.index)

    def _source(frame: pd.DataFrame) -> pd.Series:
        for cand in ("source", "publication_outlet", "publicationTitle", "publisher", "institution"):
            if cand in frame.columns:
                s = frame[cand].astype(str).str.strip()
                if s.notna().any():
                    return s
        return pd.Series([""] * len(frame), index=frame.index)

    def _publisher_bucket(s: str) -> str:
        txt = (s or "").lower()
        think_tank = ["rand","heritage","brookings","chatham","carnegie","csis","rusi","atlantic council",
                      "ifri","dgap","cepa","nupi","sipri","institut","foundation"]
        policy = ["ministry","government","gov","commission","council","parliament","oecd","world bank","imf",
                  "un ","united nations","european commission","nato","white house","cabinet","home office","senate","house of commons"]
        academic = ["journal","revue","review","press","university","universität","universidade","springer","elsevier",
                    "wiley","taylor & francis","sage","ieee","acm","oxford","cambridge"]
        if any(k in txt for k in academic): return "Academic"
        if any(k in txt for k in policy): return "Policy/Government"
        if any(k in txt for k in think_tank): return "Think Tank"
        return "Academic" if "journal" in txt else "Other/Unknown"

    def _extract_tag_values(frame: pd.DataFrame, prefix: str) -> pd.Series:
        pref = prefix.lower().strip().lstrip("#")
        for cand in (f"{pref}", f"{pref}_value", f"{pref}s", f"{pref}_values"):
            if cand in frame.columns:
                return frame[cand].fillna("").astype(str)
        if "tags" in frame.columns: raw = frame["tags"]
        elif "Tags" in frame.columns: raw = frame["Tags"]
        else:
            tag_col = next((c for c in frame.columns if "tag" in str(c).lower()), None)
            raw = frame[tag_col] if tag_col else pd.Series([""]*len(frame), index=frame.index)
        out = []
        for v in raw:
            vals, candidates = [], ([str(x) for x in v] if isinstance(v,(list,tuple))
                                    else re.split(r"[;|,]\s*", str(v or "")) if isinstance(v,str) else [])
            for t in candidates:
                s = str(t).strip(); s2 = s.lstrip("#").strip()
                if s2.lower().startswith(pref + ":"):
                    parts = s2.split(":")
                    if len(parts) >= 2: vals.append(parts[-1].strip())
            seen=set(); uniq=[]
            for val in vals:
                k=val.casefold()
                if k not in seen: seen.add(k); uniq.append(val)
            out.append(", ".join(uniq))
        return pd.Series(out, index=frame.index)

    def _with_other_bin(vc: pd.Series, *, top_n: int) -> pd.DataFrame:
        vc = vc[vc.index.notna()].copy()
        if top_n and len(vc) > top_n:
            head = vc.head(top_n)
            other = int(vc.iloc[top_n:].sum())
            if other > 0: head.loc["Other"] = other
            return head
        return vc

    def _maybe_percent(df_counts: pd.DataFrame, count_col="Count"):
        if not include_percentages or df_counts[count_col].sum() <= 0:
            return df_counts
        total = df_counts[count_col].sum()
        df_counts = df_counts.copy()
        df_counts["%"] = (100.0 * df_counts[count_col] / total).round(1)
        return df_counts

    def _new_blank_slide_with_title(title_text: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(
            Inches(0.0),
            Inches(0.15),
            Inches(SLIDE_W_IN),
            Inches(TITLE_HEIGHT_IN),
        )
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        return slide

    def _add_table_slide(title_text: str, df_small: pd.DataFrame):
        slide = _new_blank_slide_with_title(title_text)

        if df_small is None or df_small.empty:
            tb = slide.shapes.add_textbox(Inches(OUTER_MARGIN_X), Inches(2.7),
                                          Inches(SLIDE_W_IN - 2 * OUTER_MARGIN_X), Inches(1.0))
            p = tb.text_frame.paragraphs[0]
            p.text = "No data available."
            p.font.name = FONT_FAMILY
            p.font.size = Pt(16)
            p.alignment = PP_ALIGN.CENTER

            if preview_slides is not None:
                preview_slides.append({
                    "title": title_text,
                    "table_html": "<div>No data available.</div>",
                })
            return slide

        max_w = SLIDE_W_IN - 2 * OUTER_MARGIN_X
        max_h = SLIDE_H_IN - (TITLE_HEIGHT_IN + 2 * OUTER_MARGIN_Y)

        rows = max(1, len(df_small) + 1)
        cols = max(1, df_small.shape[1])

        def _len_inch_est(s: object) -> float:
            t = "" if s is None or (isinstance(s, float) and pd.isna(s)) else str(s)
            t = " ".join(t.split())
            return 0.085 * max(4, len(t)) + 0.14

        col_needs = []
        for j, col in enumerate(df_small.columns):
            header_need = _len_inch_est(col)
            longest_need = header_need
            for i in range(len(df_small)):
                longest_need = max(longest_need, min(_len_inch_est(df_small.iloc[i, j]), 3.8))
            col_needs.append(longest_need)

        MIN_W = 0.9
        MAX_W = 3.6

        col_needs = [min(MAX_W, max(MIN_W, x)) for x in col_needs]
        total_need = sum(col_needs)

        if total_need > (max_w - 0.2):
            scale = (max_w - 0.2) / total_need
            col_widths_in = [max(MIN_W, x * scale) for x in col_needs]
        else:
            col_widths_in = col_needs[:]

        table_w = sum(col_widths_in)
        est_row_h = min(0.42, max(0.26, max_h / (rows + 1)))
        table_h = min(max_h, rows * est_row_h + 0.35)

        left_in = (SLIDE_W_IN - table_w) / 2.0
        top_in = TITLE_HEIGHT_IN + ((SLIDE_H_IN - TITLE_HEIGHT_IN) - table_h) / 2.0

        table = slide.shapes.add_table(
            rows, cols, Inches(left_in), Inches(top_in), Inches(table_w), Inches(table_h)
        ).table

        for j, w in enumerate(col_widths_in):
            table.columns[j].width = Emu(Inches(w))

        for j, col in enumerate(df_small.columns):
            cell = table.cell(0, j)
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(col)
            p.font.bold = True
            p.font.name = FONT_FAMILY
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER
            cell.margin_left = cell.margin_right = Emu(900)
            cell.margin_top = Emu(600)
            cell.margin_bottom = Emu(600)

        for i, (_, r) in enumerate(df_small.iterrows(), start=1):
            for j, col in enumerate(df_small.columns):
                cell = table.cell(i, j)
                txt = "" if pd.isna(r[col]) else " ".join(str(r[col]).split())
                tf = cell.text_frame
                tf.clear()
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = txt
                p.font.name = FONT_FAMILY
                p.font.size = Pt(11)
                p.alignment = PP_ALIGN.LEFT
                cell.margin_left = Emu(800)
                cell.margin_right = Emu(800)
                cell.margin_top = Emu(500)
                cell.margin_bottom = Emu(500)

        if preview_slides is not None:
            preview_slides.append({
                "title": title_text,
                "table_html": df_small.to_html(index=False, escape=True),
            })

        return slide

    def _coerce_plotly_figure(obj):
        """
        ###1. accept go.Figure, plotly-json dicts, or containers that include a go.Figure
        ###2. return go.Figure or None
        """
        if isinstance(obj, go.Figure):
            return obj

        if isinstance(obj, dict):
            if "data" in obj and "layout" in obj:
                return go.Figure(obj)
            return None

        if isinstance(obj, (list, tuple)):
            for item in obj:
                if isinstance(item, go.Figure):
                    return item
                if isinstance(item, dict) and ("data" in item and "layout" in item):
                    return go.Figure(item)
            return None

        return None

    def _add_figure_slide(title_text: str, fig, *, notes: str = ""):
        """
        ###1. create a figure slide with a title and optional notes
        ###2. when preview/export payload is enabled, emit fig_json (Plotly spec) and optionally fig_html
        ###3. keep the PPT slide body as a placeholder text box (no PNG rendering here)
        """
        slide = _new_blank_slide_with_title(title_text)
        _set_notes(slide, notes)

        fig2 = _coerce_plotly_figure(fig)

        if preview_slides is not None:
            payload = {"title": title_text, "notes": notes}

            if isinstance(fig2, go.Figure):
                # Ensure json-safe payload (avoid numpy arrays becoming strings in the host serializer).
                payload["fig_json"] = json.loads(json.dumps(fig2.to_plotly_json(), cls=PlotlyJSONEncoder))
                payload["fig_html"] = fig2.to_html(
                    include_plotlyjs=False,
                    full_html=False,
                    config={
                        "responsive": True,
                        "displayModeBar": False,
                        "scrollZoom": True,
                        "doubleClick": "reset",
                    },
                )
            else:
                payload["fig_json"] = None
                payload["fig_html"] = None
                payload["bullets"] = ["No figure available."]

            preview_slides.append(payload)

        body_top = TITLE_HEIGHT_IN + 0.35
        body_h = SLIDE_H_IN - body_top - OUTER_MARGIN_Y
        box = slide.shapes.add_textbox(
            Inches(OUTER_MARGIN_X),
            Inches(body_top),
            Inches(SLIDE_W_IN - 2 * OUTER_MARGIN_X),
            Inches(max(1.0, body_h)),
        )
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "Figure is delivered in the preview/export payload."
        p.font.name = FONT_FAMILY
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.CENTER

        return slide

    # ---------- auto-insight notes ----------
    def _auto_insights(df_small: pd.DataFrame, label_col: str | None, value_col: str, context: str) -> str:
        if df_small is None or df_small.empty or value_col not in df_small:
            return f"{context}: No data."
        s = df_small[value_col]
        total = float(s.sum())
        if total <= 0: return f"{context}: No counts."
        parts = [f"{context}: total n={int(total)}."]
        if label_col and label_col in df_small.columns and label_col != value_col:
            top = df_small.sort_values(value_col, ascending=False).head(3)
            tops = [f"{row[label_col]} ({int(row[value_col])})" for _, row in top.iterrows()]
            parts.append("Top: " + ", ".join(tops) + ".")
        if label_col and label_col.lower() == "year":
            tail = df_small.sort_values(label_col).tail(5)
            if len(tail) >= 2:
                delta = int(tail[value_col].iloc[-1]) - int(tail[value_col].iloc[-2])
                direction = "↑" if delta > 0 else "↓" if delta < 0 else "→"
                parts.append(f"Most recent vs prior: {direction} {delta:+d}.")
        top3 = int(df_small.sort_values(value_col, ascending=False).head(3)[value_col].sum())
        parts.append(f"Top-3 concentration: {top3/total:.0%}.")
        return " ".join(parts)

    def _notes(slide, title_for_prompt, purpose, display_df=None, label_col=None, value_col=None):
        text = None
        if _safe_ai_notes_local is not None:
            text = _safe_ai_notes_local(
                title=title_for_prompt,
                purpose=purpose,
                data=display_df.to_dict("records") if display_df is not None else None,
            )
        if text is None:
            text = _auto_insights(display_df, label_col=label_col, value_col=value_col, context=title_for_prompt)
        if slide_notes and text:
            slide.notes_slide.notes_text_frame.text = text

    # ---------- figure builders ----------
    def _barh(df_counts, xcol, ycol, title):
        fig = px.bar(df_counts, x=xcol, y=ycol, orientation="h")
        fig.update_traces(hovertemplate=f"{ycol}: %{{y}}<br>{xcol}: %{{x}}<extra></extra>",
                          marker_line_width=0.5, marker_line_color="rgba(0,0,0,0.15)")
        fig.update_yaxes(automargin=True)
        return _apply_fig_style(fig, title)

    def _bar(df_counts, xcol, ycol, title):
        fig = px.bar(df_counts, x=xcol, y=ycol)
        fig.update_traces(hovertemplate=f"{xcol}: %{{x}}<br>{ycol}: %{{y}}<extra></extra>",
                          marker_line_width=0.5, marker_line_color="rgba(0,0,0,0.15)")
        fig.update_xaxes(tickangle=-25, automargin=True)
        return _apply_fig_style(fig, title)

    # ===================== build data =====================
    work = df.copy()

    # ===================== 1) Publications by year =====================
    _cb("Publications by year")

    work["__year"] = pd.to_numeric(_year_series(work), errors="coerce")

    y0, y1 = sorted([int(year_range[0]), int(year_range[1])])
    mask = (work["__year"] >= y0) & (work["__year"] <= y1)
    ydf = work.loc[mask].copy()

    if ydf.empty:
        counts_year = pd.DataFrame(
            {"Year": list(range(y0, y1 + 1)), "Publications": [0] * (y1 - y0 + 1)}
        )
    else:
        tmp = (
            ydf.dropna(subset=["__year"])
            .assign(__year=lambda d: d["__year"].astype(int))
            .groupby("__year", dropna=True)
            .size()
            .rename("Publications")
            .reset_index()
            .rename(columns={"__year": "Year"})
        )

        base = pd.DataFrame({"Year": list(range(y0, y1 + 1))})
        counts_year = (
            base.merge(tmp, on="Year", how="left")
            .assign(
                Publications=lambda d: (
                    pd.to_numeric(d["Publications"], errors="coerce")
                    .replace([np.inf, -np.inf], np.nan)
                    .infer_objects(copy=False)
                    .fillna(0)
                    .astype(int)
                ),
                Year=lambda d: pd.to_numeric(d["Year"], errors="coerce").fillna(0).astype(int),
            )[["Year", "Publications"]]
        )

    _add_table_slide(f"Publications by Year ({y0}–{y1}) — Table", counts_year.copy())

    plot_df = counts_year.copy()

    fig = px.bar(plot_df, x="Year", y="Publications")
    fig.update_traces(
        cliponaxis=False,
        marker_line_width=0.5,
        marker_line_color="rgba(0,0,0,0.15)",
    )
    _ymax = float(np.nanmax(pd.to_numeric(plot_df["Publications"], errors="coerce")))
    if np.isfinite(_ymax) and _ymax > 0:
        fig.update_yaxes(range=[0, _ymax * 1.25], automargin=True)

    fig.update_layout(margin=dict(l=80, r=120, t=120, b=80))
    fig.update_xaxes(tickangle=-25, automargin=True)
    fig = _apply_fig_style(fig, f"{collection_name} · Publications per Year")

    fig_slide = _add_figure_slide(f"Publications by Year ({y0}–{y1}) — Figure", fig)
    _notes(
        fig_slide,
        f"Publications by Year ({y0}–{y1})",
        "Identify growth phases and recency undercounts.",
        counts_year,
        "Year",
        "Publications",
    )

    # ===================== 2) Document type =====================
    _cb("Document type distribution")

    if "document_type_value" in work.columns:
        s = work["document_type_value"].fillna("").astype(str).str.strip()
    else:
        s = _doc_type(work).fillna("").astype(str).str.strip()
    work["__doctype"] = s.where(s.ne(""), "Unknown")

    vc = work["__doctype"].value_counts()
    vc = _with_other_bin(vc, top_n=top_n_cats)
    dt_counts = vc.rename_axis("Document Type").reset_index(name="Count")
    dt_counts = _maybe_percent(dt_counts, "Count")

    _add_table_slide(
        "Distribution by Document Type — Table",
        dt_counts.sort_values("Count", ascending=False),
    )

    plot_df = dt_counts.sort_values("Count", ascending=True).copy()
    if "%" in plot_df.columns:
        plot_df["Label"] = plot_df.apply(lambda r: f"{int(r['Count'])} ({float(r['%']):.1f}%)", axis=1)
    else:
        plot_df["Label"] = plot_df["Count"].astype(int).astype(str)

    fig = px.bar(plot_df, x="Count", y="Document Type", orientation="h")
    fig = _apply_fig_style(fig, f"{collection_name} · Document Types")
    _outside_text_headroom(fig, orientation="h", values=plot_df["Count"], labels=plot_df["Label"])

    fig_slide = _add_figure_slide("Distribution by Document Type — Figure", fig)
    _notes(
        fig_slide,
        "Document Type Distribution",
        "Call out dominant formats and whether grey lit dominates.",
        dt_counts,
        "Document Type",
        "Count",
    )

    # ===================== 3) Publisher type =====================
    _cb("Publisher type split")

    if "publisher_type_value" in work.columns:
        pub_type = work["publisher_type_value"].fillna("").astype(str).str.strip()
    elif "publisher type" in work.columns:
        pub_type = work["publisher type"].fillna("").astype(str).str.strip()
    elif "publisher_type" in work.columns:
        pub_type = work["publisher_type"].fillna("").astype(str).str.strip()
    else:
        pub_type = _source(work).fillna("").astype(str).apply(_publisher_bucket).astype(str).str.strip()

    pub_type = pub_type.where(pub_type.ne(""), np.nan).fillna("Other/Unknown")

    pub_counts = pub_type.value_counts().rename_axis("Publisher Type").reset_index(name="Count")
    pub_counts = _maybe_percent(pub_counts, "Count")

    _add_table_slide(
        "Publisher Type Split — Table",
        pub_counts.sort_values("Count", ascending=False),
    )

    plot_df = pub_counts.sort_values("Count", ascending=False).copy()
    if "%" in plot_df.columns:
        plot_df["Label"] = plot_df.apply(lambda r: f"{int(r['Count'])} ({float(r['%']):.1f}%)", axis=1)
    else:
        plot_df["Label"] = plot_df["Count"].astype(int).astype(str)

    fig = px.bar(plot_df, x="Publisher Type", y="Count")
    fig.update_xaxes(tickangle=-25, automargin=True)
    fig = _apply_fig_style(fig, f"{collection_name} · Publisher Type")
    _outside_text_headroom(fig, orientation="v", values=plot_df["Count"], labels=plot_df["Label"])

    fig_slide = _add_figure_slide("Publisher Type Split — Figure", fig)
    _notes(
        fig_slide,
        "Publisher Type Split",
        "Assess venue balance; policy/think-tank heavy corpora may imply different evidence standards.",
        pub_counts,
        "Publisher Type",
        "Count",
    )

    # Removed by request: Country focus / Phase focus / Sector focus graphs.
    # Keep the earlier 3 scope plots only.
    _cb("Scope slides completed.")

    if preview_slides is not None:
        fig_json_n = sum(1 for s in preview_slides if isinstance(s, dict) and s.get("fig_json") is not None)
        fig_html_n = sum(1 for s in preview_slides if isinstance(s, dict) and str(s.get("fig_html") or "").strip())
        print(
            f"[shape_scope] export={export} return_payload={return_payload} slides={len(preview_slides)} fig_json={fig_json_n} fig_html={fig_html_n}"
        )
        return {"slides": preview_slides, "index": 0}

    return None

    # ===================== 4) Country focus =====================
    _cb("Country focus distribution — Top 20")

    for cand in ("country_focus_value", "country_focus", "country focus"):
        if cand in work.columns:
            country_series = work[cand].fillna("").astype(str)
            break
    else:
        country_series = _extract_tag_values(work, "country_focus_value")

    if country_series.eq("").all():
        country_series = _extract_tag_values(work, "affiliation:country")
    if country_series.eq("").all():
        country_series = _extract_tag_values(work, "affiliation")

    country_series = country_series.fillna("").astype(str)
    cf = pd.DataFrame({"Country Focus": country_series.where(country_series.ne(""), np.nan)}).dropna()

    # --- 4b) Canonicalize country names / codes to reduce duplicates ---
    def _norm_country(x: str) -> str:
        s = (x or "").strip()
        if not s:
            return s
        s_key = s.lower().replace(".", "").replace("_", " ").strip()

        MAP = {
            # United States
            "us": "United States", "u s": "United States", "usa": "United States",
            "united states of america": "United States",
            # United Kingdom
            "uk": "United Kingdom", "u k": "United Kingdom", "gb": "United Kingdom",
            "gbr": "United Kingdom", "great britain": "United Kingdom",
            "england": "United Kingdom", "scotland": "United Kingdom",
            "wales": "United Kingdom", "northern ireland": "United Kingdom",
            # European Union
            "eu": "European Union", "e u": "European Union", "european union": "European Union",
            # Global / non-specific
            "global": "Global/International", "world": "Global/International",
            "international": "Global/International", "non-country specific": "Global/International",
            "non country specific": "Global/International",
        }
        if s_key in MAP:
            return MAP[s_key]
        # Title-case fallback for regular country names
        return s.title()

    def _uniq_normed_countries(cell: str) -> list[str]:
        """
        Split on commas/semicolons/pipes/slashes, normalize each token,
        and de-duplicate *within the same row* (so 'Brazil, Brazil' counts once for that record).
        """
        tokens = re.split(r"\s*[,;|/]\s*", str(cell or ""))
        seen, out = set(), []
        for tok in tokens:
            tok = tok.strip()
            if not tok:
                continue
            norm = _norm_country(tok)
            key = norm.casefold()
            if key not in seen:
                seen.add(key)
                out.append(norm)
        return out

    TOPN_COUNTRIES = 20

    if cf.empty:
        country_counts = pd.DataFrame({"Country Focus": ["N/A"], "Count": [0]})
    else:
        # explode AFTER per-row de-duplication
        exp = cf["Country Focus"].apply(_uniq_normed_countries).explode()
        vc = exp.value_counts()
        # Top 20 ONLY (no "Other" aggregation to avoid a large 'Other' bucket)
        vc_top = vc.head(TOPN_COUNTRIES)
        country_counts = vc_top.rename_axis("Country Focus").reset_index(name="Count")

    # (Optional) percentage column
    country_counts = _maybe_percent(country_counts, "Count")

    # ---------- Slides ----------
    # Table (sorted by Count desc)
    table_slide = _add_table_slide("Distribution by Country Focus — Top 20 — Table",
                                   country_counts.sort_values("Count", ascending=False))

    # Figure (horizontal bars)
    plot_df = country_counts.sort_values("Count", ascending=True).copy()
    plot_df["Label"] = (
        plot_df.apply(lambda r: f"{int(r['Count'])} ({float(r['%']):.1f}%)", axis=1)
        if "%" in plot_df.columns
        else plot_df["Count"].astype(int).astype(str)
    )

    fig = px.bar(plot_df, x="Count", y="Country Focus", orientation="h")
    fig = _apply_fig_style(fig, f"{collection_name} · Country Focus — Top 20")

    # prevent cropping of the percent tail / value labels
    _outside_text_headroom(
        fig,
        orientation="h",
        values=plot_df["Count"],
        labels=plot_df["Label"],
    )

    fig_slide = _add_figure_slide("Distribution by Country Focus — Top 20 — Figure", fig)

    # ---------- Notes (explicit, reliable) ----------
    table_notes = (
            "Methodology (Country Focus — Top 20):\n"
            "• Source columns (in order): 'country_focus_value', 'country_focus', 'country focus'. If all are empty, "
            "we parse tag-derived columns 'affiliation:country' then 'affiliation'.\n"
            "• Per-record parsing: split country strings on commas/semicolons/pipes/slashes; trim whitespace.\n"
            "• Normalization: common aliases are canonicalized — US/USA → United States; UK/GB/Great Britain/England/"
            "Scotland/Wales/Northern Ireland → United Kingdom; EU → European Union; "
            "Global/World/International → Global/International.\n"
            "• De-duplication: within each record, repeated countries are counted once; after that we aggregate across all records.\n"
            "• Display: top 20 countries only (no 'Other' bucket), to avoid inflating a catch-all category.\n"
            + ("• Percentages are relative to total country mentions after per-record de-duplication."
               if "%" in country_counts.columns else "")
    )
    table_slide.notes_slide.notes_text_frame.text = table_notes

    fig_notes = (
        "Figure shows the same Top 20 country counts as the table, using the normalization and per-record de-duplication rules "
        "described in the table notes."
    )
    fig_slide.notes_slide.notes_text_frame.text = fig_notes

    # ===================== 5) Phase focus =====================
    _cb("Phase focus distribution")
    # Prefer explicit phase columns (handles: 'phase_focus_value', 'phase_focus', 'phase focus')
    for cand in ("phase_focus_value", "phase_focus", "phase focus"):
        if cand in work.columns:
            phase_series = work[cand].fillna("").astype(str)
            break
    else:
        phase_series = _extract_tag_values(work, "phase_focus_value")

    ph = pd.DataFrame({"Phase Focus": phase_series.where(phase_series.ne(""), np.nan)}).dropna()

    if ph.empty:
        phase_counts = pd.DataFrame({"Phase Focus": ["N/A"], "Count": [0]})
    else:
        exp = ph["Phase Focus"].str.split(",").explode().str.strip()
        exp = exp[exp != ""]
        vc = _with_other_bin(exp.value_counts(), top_n=top_n_cats)
        phase_counts = vc.rename_axis("Phase Focus").reset_index(name="Count")

    phase_counts = _maybe_percent(phase_counts, "Count")

    _add_table_slide("Number of Studies by Phase Focus — Table",
                     phase_counts.sort_values("Count", ascending=False))

    plot_df = phase_counts.sort_values("Count", ascending=False).copy()
    if "%" in plot_df.columns:
        plot_df["Label"] = plot_df.apply(lambda r: f"{int(r['Count'])} ({float(r['%']):.1f}%)", axis=1)
    else:
        plot_df["Label"] = plot_df["Count"].astype(int).astype(str)

    fig = px.bar(plot_df, x="Phase Focus", y="Count")
    fig.update_xaxes(tickangle=-25, automargin=True)
    fig = _apply_fig_style(fig, f"{collection_name} · Phase Focus")

    _outside_text_headroom(
        fig,
        orientation="v",
        values=plot_df["Count"],
        labels=plot_df["Label"],
    )

    fig_slide = _add_figure_slide("Number of Studies by Phase Focus — Figure", fig)
    _notes(fig_slide, "Phase Focus",
           "Assess clustering (e.g., Designing/Implementing) and gaps for future sampling.",
           phase_counts, "Phase Focus", "Count")

    # ===================== 6) Sector focus =====================
    _cb("Sector focus distribution")
    # Prefer explicit sector columns (handles: 'sector_focus_value', 'sector_focus', 'sector focus')
    for cand in ("sector_focus_value", "sector_focus", "sector focus"):
        if cand in work.columns:
            sector_series = work[cand].fillna("").astype(str)
            break
    else:
        sector_series = _extract_tag_values(work, "sector_focus_value")

    sf = pd.DataFrame({"Sector Focus": sector_series.where(sector_series.ne(""), np.nan)}).dropna()

    if sf.empty:
        sector_counts = pd.DataFrame({"Sector Focus": ["N/A"], "Count": [0]})
    else:
        exp = sf["Sector Focus"].str.split(",").explode().str.strip()
        exp = exp[exp != ""]
        vc = _with_other_bin(exp.value_counts(), top_n=top_n_cats)
        sector_counts = vc.rename_axis("Sector Focus").reset_index(name="Count")

    sector_counts = _maybe_percent(sector_counts, "Count")

    _add_table_slide("Number of Studies by Sector Focus — Table",
                     sector_counts.sort_values("Count", ascending=False))

    plot_df = sector_counts.sort_values("Count", ascending=True).copy()
    if "%" in plot_df.columns:
        plot_df["Label"] = plot_df.apply(lambda r: f"{int(r['Count'])} ({float(r['%']):.1f}%)", axis=1)
    else:
        plot_df["Label"] = plot_df["Count"].astype(int).astype(str)

    fig = px.bar(plot_df, x="Count", y="Sector Focus", orientation="h")
    fig = _apply_fig_style(fig, f"{collection_name} · Sector Focus")

    _outside_text_headroom(
        fig,
        orientation="h",
        values=plot_df["Count"],
        labels=plot_df["Label"],
    )

    fig_slide = _add_figure_slide("Number of Studies by Sector Focus — Figure", fig)
    _notes(fig_slide, "Sector Focus",
           "Identify dominant sectors and non-sector specific load; check tag consistency.",
           sector_counts, "Sector Focus", "Count")

    _cb("Scope slides completed.")

    if preview_slides is not None:
        fig_json_n = sum(1 for s in preview_slides if isinstance(s, dict) and s.get("fig_json") is not None)
        fig_html_n = sum(1 for s in preview_slides if isinstance(s, dict) and str(s.get("fig_html") or "").strip())
        print(
            f"[shape_scope] export={export} return_payload={return_payload} slides={len(preview_slides)} fig_json={fig_json_n} fig_html={fig_html_n}"
        )
        return {"slides": preview_slides, "index": 0}

    return None

# --- Collaboration Plots ---

def _plot_co_authorship_network(df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    """
    Co-authorship network with labels, community colouring, and robust parsing.

    Params (all optional):
      - min_collaborations: int = 2         # edge threshold by co-auth count
      - top_n_authors: int = 50             # node cap by weighted degree
      - layout: str = "force"               # "force" | "kk" | "spectral" | "circular"
      - iterations: int = 50                # layout iterations for force
      - spring_k: float | None              # spring-layout k; default ~ 1.5/sqrt(n)
      - show_labels: bool = True            # draw author names
      - label_strategy: str = "topk"        # "all" | "topk" | "deg_threshold"
      - label_top_k: int = 20               # labels for top-k nodes (by weighted degree)
      - label_min_degree: int = 2           # min degree for labels (if strategy="deg_threshold")
    """
    import math
    import numpy as np
    import pandas as pd
    import networkx as nx
    import itertools
    from collections import defaultdict
    import plotly.graph_objects as go

    # ---- params & defaults ----
    min_collaborations = int(params.get("min_collaborations", 1))
    top_n_nodes = int(params.get("top_n_authors", 50))
    layout = str(params.get("layout", "force")).lower()
    iterations = int(params.get("iterations", 50))
    spring_k = params.get("spring_k", None)
    show_labels = bool(params.get("show_labels", True))
    label_strategy = str(params.get("label_strategy", "topk")).lower()
    label_top_k = int(params.get("label_top_k", 20))
    label_min_degree = int(params.get("label_min_degree", 2))

    # ---- parse authors from DataFrame ----
    if "authors_list" in df.columns and df["authors_list"].notna().any():
        series = df["authors_list"].dropna().apply(
            lambda xs: [str(a).strip() for a in xs if str(a).strip()]
            if isinstance(xs, (list, tuple, set)) else []
        )
    elif "authors" in df.columns and df["authors"].notna().any():
        series = df["authors"].dropna().astype(str).apply(
            lambda s: [a.strip() for a in s.split(";") if a.strip()]
        )
    else:
        empty = pd.DataFrame(columns=["Source", "Target", "Weight"])
        fig = go.Figure().add_annotation(
            text="Missing 'authors' column.", showarrow=False
        )
        return empty, fig

    # build weighted co-auth edges
    edge_counts = defaultdict(int)
    for authors in series:
        if not authors or len(authors) < 2:
            continue
        unique_sorted = sorted(set(authors))
        for a, b in itertools.combinations(unique_sorted, 2):
            edge_counts[(a, b)] += 1

    edges_df = pd.DataFrame(
        [{"Source": a, "Target": b, "Weight": w} for (a, b), w in edge_counts.items()]
    )
    if edges_df.empty:
        # Fallback: still show something useful for collaboration even when all papers are single-authored.
        counts = series.apply(lambda xs: len(xs) if isinstance(xs, (list, tuple, set)) else 0)
        dist = (
            pd.DataFrame({"Authors per document": counts})
            .value_counts()
            .rename_axis(["Authors per document"])
            .reset_index(name="Count")
            .sort_values("Authors per document")
            .reset_index(drop=True)
        )
        fig = go.Figure()
        fig.add_bar(x=dist["Authors per document"], y=dist["Count"])
        fig.update_layout(
            title="Authorship per document (no co-authorship edges)",
            xaxis_title="Authors per document",
            yaxis_title="Count",
            margin=dict(l=60, r=40, t=40, b=60),
        )
        return edges_df, fig

    # threshold by min_collaborations
    edges_df = edges_df.loc[edges_df["Weight"] >= min_collaborations].copy()
    if edges_df.empty:
        counts = series.apply(lambda xs: len(xs) if isinstance(xs, (list, tuple, set)) else 0)
        dist = (
            pd.DataFrame({"Authors per document": counts})
            .value_counts()
            .rename_axis(["Authors per document"])
            .reset_index(name="Count")
            .sort_values("Authors per document")
            .reset_index(drop=True)
        )
        fig = go.Figure()
        fig.add_bar(x=dist["Authors per document"], y=dist["Count"])
        fig.update_layout(
            title=f"Authorship per document (no edges ≥ {min_collaborations})",
            xaxis_title="Authors per document",
            yaxis_title="Count",
            margin=dict(l=60, r=40, t=40, b=60),
        )
        return edges_df, fig

    # graph
    G = nx.Graph()
    for _, r in edges_df.iterrows():
        G.add_edge(r["Source"], r["Target"], weight=float(r["Weight"]))

    # rank nodes by weighted degree ("strength")
    strength = dict(G.degree(weight="weight"))
    if len(strength) > top_n_nodes:
        top_nodes = [n for n, _ in sorted(strength.items(), key=lambda t: t[1], reverse=True)[:top_n_nodes]]
        G = G.subgraph(top_nodes).copy()
        # re-filter edges_df to induced subgraph
        keep = set(G.nodes())
        edges_df = edges_df[edges_df["Source"].isin(keep) & edges_df["Target"].isin(keep)].reset_index(drop=True)
        strength = dict(G.degree(weight="weight"))

    n_nodes = max(1, G.number_of_nodes())

    # community detection (greedy modularity as a robust default)
    try:
        from networkx.algorithms.community import greedy_modularity_communities
        comms = list(greedy_modularity_communities(G, weight="weight"))
        node_community = {}
        for cid, members in enumerate(comms):
            for m in members:
                node_community[m] = cid
    except Exception:
        node_community = {n: 0 for n in G.nodes()}

    # centralities
    try:
        bet = nx.betweenness_centrality(G, weight="weight", normalized=True)
    except Exception:
        bet = {n: 0.0 for n in G.nodes()}

    # layout
    if layout == "kk":
        pos = nx.kamada_kawai_layout(G, weight="weight")
    elif layout == "spectral":
        pos = nx.spectral_layout(G)
    elif layout == "circular":
        pos = nx.circular_layout(G)
    else:
        k = spring_k if isinstance(spring_k, (int, float)) else 1.5 / math.sqrt(n_nodes)
        pos = nx.spring_layout(G, k=k, iterations=iterations, weight="weight", seed=42)

    # edge traces with weight-coded width and hover
    if edges_df["Weight"].max() > 0:
        w_norm = (edges_df["Weight"] - edges_df["Weight"].min()) / max(1e-9, edges_df["Weight"].max() - edges_df["Weight"].min())
    else:
        w_norm = pd.Series([0.0] * len(edges_df))
    edge_x, edge_y, edge_text, edge_width = [], [], [], []
    for s, t, w in edges_df[["Source", "Target", "Weight"]].itertuples(index=False):
        x0, y0 = pos[s]; x1, y1 = pos[t]
        edge_x += [x0, x1, None]
        edge_y += [y0, y1, None]
        edge_text.append(f"{s} ↔ {t}<br>Co-authored: {int(w)}")
        edge_width.append(0.6 + 2.0 * float(w_norm.loc[(edges_df["Source"] == s) & (edges_df["Target"] == t)].iloc[0]))

    edge_trace = go.Scatter(
        x=edge_x, y=edge_y,
        mode="lines",
        line=dict(width=1, color="rgba(130,130,130,0.6)"),
        hoverinfo="none",
        showlegend=False
    )

    # node attributes
    deg = dict(G.degree())
    s_vals = np.array([strength.get(n, 0.0) for n in G.nodes()], dtype=float)
    s_min, s_max = (float(np.nanmin(s_vals)), float(np.nanmax(s_vals))) if len(s_vals) else (0.0, 0.0)
    def _scale_size(s):
        if s_max <= 0:
            return 10.0
        # sqrt scaling for visual balance
        return 10.0 + 24.0 * math.sqrt(s / s_max)

    node_x, node_y, node_size, node_color, node_text = [], [], [], [], []
    for n in G.nodes():
        x, y = pos[n]
        node_x.append(x); node_y.append(y)
        node_size.append(_scale_size(strength.get(n, 0.0)))
        node_color.append(node_community.get(n, 0))
        node_text.append(
            f"{n}"
            f"<br>Publications (weighted degree): {strength.get(n, 0):.0f}"
            f"<br>Collaborators (degree): {deg.get(n, 0)}"
            f"<br>Betweenness: {bet.get(n, 0.0):.3f}"
            f"<br>Community: {node_community.get(n, 0)}"
        )

    node_trace = go.Scatter(
        x=node_x, y=node_y,
        mode="markers",
        text=node_text,
        hoverinfo="text",
        marker=dict(
            size=node_size,
            color=node_color,
            colorscale="Turbo",
            showscale=True,
            colorbar=dict(title="Community")
        ),
        showlegend=False
    )

    # label selection
    label_nodes = set()
    if show_labels:
        if label_strategy == "all":
            label_nodes = set(G.nodes())
        elif label_strategy == "deg_threshold":
            label_nodes = {n for n in G.nodes() if deg.get(n, 0) >= label_min_degree}
        else:
            # "topk" by weighted degree
            label_nodes = {n for n, _ in sorted(strength.items(), key=lambda t: t[1], reverse=True)[:max(1, label_top_k)]}

    label_x, label_y, label_txt = [], [], []
    for n in label_nodes:
        label_x.append(pos[n][0]); label_y.append(pos[n][1]); label_txt.append(n)

    label_trace = go.Scatter(
        x=label_x, y=label_y,
        mode="text",
        text=label_txt,
        textposition="middle center",
        textfont=dict(size=10),
        hoverinfo="skip",
        showlegend=False
    ) if show_labels and label_nodes else None

    # figure
    data_traces = [edge_trace, node_trace]
    if label_trace is not None:
        data_traces.append(label_trace)

    fig = go.Figure(
        data=data_traces,
        layout=go.Layout(
            title="Co-Authorship Network",
            showlegend=False,
            margin=dict(b=20, l=20, r=20, t=30),
            xaxis=dict(visible=False),
            yaxis=dict(visible=False),
            hovermode="closest"
        )
    )

    return edges_df.reset_index(drop=True), fig

def _plot_inter_entity_collaboration_chord(df: pd.DataFrame, params: dict, entity_type: str) -> tuple[pd.DataFrame, go.Figure]:
    """Creates a chord diagram showing collaborations between institutions or countries."""
    top_n = params.get("top_n_entities", 15)
    if entity_type not in df.columns or df[entity_type].dropna().empty:
        return pd.DataFrame(), go.Figure().add_annotation(text=f"Missing '{entity_type}' column.", showarrow=False)

    # Use explode for list-like columns (like country) or split for string columns
    if isinstance(df[entity_type].dropna().iloc[0], list):
        df_exploded = df.explode(entity_type)
    else:
        df_exploded = df.assign(**{entity_type: df[entity_type].str.split(';')}).explode(entity_type)

    df_exploded[entity_type] = df_exploded[entity_type].str.strip()
    df_exploded.dropna(subset=[entity_type], inplace=True)

    pairs = df_exploded.groupby('title')[entity_type].apply(lambda x: list(itertools.combinations(sorted(x.unique()), 2))).explode().dropna()

    if pairs.empty: return pd.DataFrame(), go.Figure().add_annotation(text=f"No inter-{entity_type} collaborations found.", showarrow=False)

    edge_counts = Counter(pairs); edges_df = pd.DataFrame(edge_counts.items(), columns=['Pair', 'Count']).sort_values('Count', ascending=False)
    all_entities = pd.concat([edges_df['Pair'].apply(lambda x: x[0]), edges_df['Pair'].apply(lambda x: x[1])])
    top_entities = all_entities.value_counts().nlargest(top_n).index
    edges_df = edges_df[edges_df['Pair'].apply(lambda x: x[0] in top_entities and x[1] in top_entities)]

    if edges_df.empty: return pd.DataFrame(), go.Figure().add_annotation(text=f"No links between top {top_n} entities.", showarrow=False)

    all_nodes = sorted(list(top_entities)); node_map = {name: i for i, name in enumerate(all_nodes)}
    sources = edges_df['Pair'].apply(lambda x: node_map[x[0]]).tolist(); targets = edges_df['Pair'].apply(lambda x: node_map[x[1]]).tolist(); values = edges_df['Count'].tolist()
    fig = go.Figure(go.Sankey(arrangement="snap", node=dict(label=all_nodes), link=dict(source=sources, target=targets, value=values)))
    fig.update_layout(title_text=f"Top Inter-{entity_type.capitalize()} Collaboration Flow", font_size=10)
    return edges_df, fig

def _plot_author_collaboration_sunburst(df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    """Creates a sunburst of a focus author's co-authors grouped by institution."""
    author_to_plot = params.get("author_name") or _get_top_author(df)
    if not author_to_plot: return pd.DataFrame(), go.Figure().add_annotation(text="No author available to analyze.", showarrow=False)

    co_authored_papers = df[df['authors'].str.contains(re.escape(author_to_plot), na=False)]
    collaborator_list = [author.strip() for _, row in co_authored_papers.iterrows() for author in row['authors'].split(';') if author.strip() and author.strip() != author_to_plot]

    if not collaborator_list: return pd.DataFrame(), go.Figure().add_annotation(text=f"No collaborators found for {author_to_plot}.", showarrow=False)

    records = []
    for collaborator in set(collaborator_list):
        collab_papers = df[df['authors'].str.contains(re.escape(collaborator), na=False)]
        institution = collab_papers['institution'].dropna().iloc[0] if not collab_papers['institution'].dropna().empty else "Unknown Institution"
        records.append({"Institution": institution, "Collaborator": collaborator, "Count": collaborator_list.count(collaborator)})

    sunburst_df = pd.DataFrame(records).sort_values("Count", ascending=False)
    fig = px.sunburst(sunburst_df, path=['Institution', 'Collaborator'], values='Count', title=f"Collaboration Structure for {author_to_plot}")
    return sunburst_df, fig

def _plot_collaboration_matrix_heatmap(df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    """Creates a heatmap of collaborations between top authors."""
    top_n = params.get("top_n_authors", 20)
    work = df.copy()

    if "authors" not in work.columns:
        top_authors = []
    else:
        a = work["authors"]

        if a.dtype == "object":
            non_null = a.dropna()
            sample = non_null.iloc[0] if len(non_null) else ""

            if isinstance(sample, str):
                work["authors"] = (
                    a.fillna("")
                    .astype(str)
                    .apply(lambda s: [x.strip() for x in s.split(";") if x.strip()])
                )
            else:
                work["authors"] = a

        top_authors = list((_get_top_author(work.explode("authors")) or pd.Series()).index[:top_n])
    if len(top_authors) < 2: return pd.DataFrame(), go.Figure().add_annotation(text="Not enough authors for matrix.", showarrow=False)

    matrix = pd.DataFrame(0, index=top_authors, columns=top_authors)
    for _, row in df.iterrows():
        authors = {a.strip() for a in row.get('authors', '').split(';')}
        for a1, a2 in itertools.combinations(authors, 2):
            if a1 in top_authors and a2 in top_authors:
                matrix.loc[a1, a2] += 1; matrix.loc[a2, a1] += 1

    fig = px.imshow(matrix, text_auto=True, title=f"Collaboration Matrix of Top {top_n} Authors")
    return matrix, fig
def analyze_author_collaboration(df: pd.DataFrame, params: dict, callback: callable) -> tuple[pd.DataFrame, go.Figure]:
    plot_type = params.get("plot_type")
    callback(f"Analyzing Author Collaboration: {plot_type}...")

    plot_map = {
        'co_authorship_network': _plot_co_authorship_network,
        'inter_institution_chord': lambda d, p: _plot_inter_entity_collaboration_chord(d, p, 'institution'),
        'inter_country_chord': lambda d, p: _plot_inter_entity_collaboration_chord(d, p, 'country'),
        'collaboration_matrix': _plot_collaboration_matrix_heatmap,
        'author_collaboration_sunburst': _plot_author_collaboration_sunburst
    }

    plot_func = plot_map.get(plot_type)
    if plot_func:
        return plot_func(df, params)

    msg = f"Unknown collaboration plot type: '{plot_type}'"
    callback(f"Error: {msg}")
    return pd.DataFrame(), go.Figure().add_annotation(text=msg, showarrow=False)



def analyze_lotkas_law(df: pd.DataFrame, params: dict, progress_callback=None):
    """
    Analyzes author productivity based on Lotka's Law.
    params: {} (currently no specific params from UI for this, but can be added)
    Returns: (results_dataframe_with_lotka_distribution, plotly_figure_object)
    """

    def _callback(msg):
        if progress_callback: progress_callback(msg)
        logging.info(msg)

    _callback("Starting Lotka's Law analysis...")
    if df is None or df.empty or 'authors' not in df.columns:
        _callback("Missing 'authors' column or DataFrame is empty for Lotka's Law.")
        return pd.DataFrame(columns=['Num_Publications', 'Num_Authors_Observed', 'Prop_Authors_Observed']), None

    all_authors_flat = [author.strip() for author_list_str in df['authors'].dropna().astype(str)
                        for author in author_list_str.split(';') if author.strip()]
    if not all_authors_flat:
        _callback("No author data found after parsing.")
        return pd.DataFrame(columns=['Num_Publications', 'Num_Authors_Observed', 'Prop_Authors_Observed']), None

    author_publication_counts = Counter(all_authors_flat)

    # Count how many authors published X number of papers
    # productivity_counts: {1: 100, 2: 30, 3:10} means 100 authors published 1 paper, 30 published 2, etc.
    productivity_distribution = Counter(author_publication_counts.values())

    results_df = pd.DataFrame(list(productivity_distribution.items()),
                              columns=['Num_Publications', 'Num_Authors_Observed'])
    results_df = results_df.sort_values(by='Num_Publications').reset_index(drop=True)

    total_unique_authors = results_df['Num_Authors_Observed'].sum()
    results_df['Prop_Authors_Observed'] = results_df['Num_Authors_Observed'] / total_unique_authors

    _callback("Lotka's Law distribution table generated. Creating Plotly figure...")
    fig = go.Figure()
    if not results_df.empty:
        fig.add_trace(go.Scatter(
            x=results_df['Num_Publications'],
            y=results_df['Prop_Authors_Observed'],
            mode='lines+markers',
            name="Observed Proportion of Authors",
            hovertemplate="Publications: %{x}<br>Proportion of Authors: %{y:.3f}<br>(%{customdata[0]} authors)<extra></extra>",
            customdata=results_df[['Num_Authors_Observed']]
        ))

        # Attempt to fit Lotka's Law: F(n) = C / n^alpha
        # This is a simplification. Real fitting involves regression on log-log data.
        # For alpha=2 (common assumption), F(n) is proportional to 1/n^2
        # We can scale C such that the sum of proportions is 1 (or close for the observed range)
        if len(results_df['Num_Publications']) > 1:
            try:
                # Basic estimation of C (constant of proportionality) and alpha
                # Fit log(Prop_Authors) = log(C) - alpha * log(Num_Publications)
                log_n = np.log(results_df['Num_Publications'])
                log_f_n = np.log(results_df['Prop_Authors_Observed'])

                # Remove -inf from log_f_n if Prop_Authors_Observed can be 0 for some n (after reindex or if gaps exist)
                # However, results_df comes from value_counts, so Num_Authors_Observed should be > 0.
                # np.isfinite check is safer
                valid_fit_indices = np.isfinite(log_n) & np.isfinite(log_f_n)
                if np.sum(valid_fit_indices) > 1:  # Need at least 2 points to fit a line
                    log_n_fit = log_n[valid_fit_indices]
                    log_f_n_fit = log_f_n[valid_fit_indices]

                    # Linear regression: log_f_n = m * log_n + b
                    # where m = -alpha, b = log(C)
                    m, b = np.polyfit(log_n_fit, log_f_n_fit, 1)
                    alpha_estimated = -m
                    c_estimated = np.exp(b)

                    results_df['Lotka_Predicted_Prop (alpha~2)'] = (1 / results_df['Num_Publications'] ** 2) / \
                                                                   sum(1 / k ** 2 for k in
                                                                       results_df['Num_Publications'])
                    results_df[f'Lotka_Predicted_Prop (alpha={alpha_estimated:.2f})'] = \
                        (c_estimated / results_df['Num_Publications'] ** alpha_estimated)

                    fig.add_trace(go.Scatter(
                        x=results_df['Num_Publications'],
                        y=results_df['Lotka_Predicted_Prop (alpha~2)'],
                        mode='lines', name="Lotka's Law (alpha=2)",
                        line=dict(dash='dash')
                    ))
                    fig.add_trace(go.Scatter(
                        x=results_df['Num_Publications'],
                        y=results_df[f'Lotka_Predicted_Prop (alpha={alpha_estimated:.2f})'],
                        mode='lines', name=f"Lotka's Law (Fitted alpha={alpha_estimated:.2f})",
                        line=dict(dash='dot')
                    ))
                    _callback(f"Lotka's Law: Estimated alpha = {alpha_estimated:.2f}, C = {c_estimated:.3f}")

            except Exception as e_lotka_fit:
                _callback(f"Could not fit Lotka's Law parameters: {e_lotka_fit}")

        fig.update_layout(
            title="Lotka's Law: Author Productivity Distribution",
            xaxis_title="Number of Publications (n)",
            yaxis_title="Proportion of Authors Publishing n Papers",
            xaxis_type="log",
            yaxis_type="log",
            legend_title_text="Distribution"
        )
    else:
        _callback("No data for Lotka's Law plot.")

    _callback("Lotka's Law analysis complete.")
    return results_df, fig


def analyze_authors_production_over_time(df: pd.DataFrame, params: dict, progress_callback=None):
    """
    Analyzes and visualizes the publication production of top N authors over time.
    params: {"top_n_authors": int}
    Returns: (results_dataframe, plotly_figure_object)
    """

    def _callback(msg):
        if progress_callback: progress_callback(msg)
        logging.info(msg)

    _callback("Starting Authors' Production Over Time analysis...")
    if df is None or df.empty or 'authors' not in df.columns or 'year' not in df.columns:
        _callback("Missing 'authors' or 'year' column, or DataFrame is empty.")
        return pd.DataFrame(), None

    top_n = params.get("top_n_authors", 5)  # Default to top 5 if not specified

    # Identify top N authors based on total publications
    all_authors_flat = [author.strip() for author_list_str in df['authors'].dropna().astype(str)
                        for author in author_list_str.split(';') if author.strip()]
    if not all_authors_flat:
        _callback("No author data found after parsing.")
        return pd.DataFrame(), None

    author_total_counts = Counter(all_authors_flat)
    top_authors_list = [author for author, count in author_total_counts.most_common(top_n)]

    if not top_authors_list:
        _callback(f"Could not determine top {top_n} authors.")
        return pd.DataFrame(), None
    _callback(f"Analyzing production for top {len(top_authors_list)} authors: {', '.join(top_authors_list)}")

    # Prepare data: Ensure 'year' is numeric and clean
    df_analysis = df.copy()
    df_analysis['year_numeric'] = pd.to_numeric(df_analysis['year'], errors='coerce')
    df_analysis = df_analysis.dropna(subset=['year_numeric', 'authors'])
    df_analysis['year_numeric'] = df_analysis['year_numeric'].astype(int)

    min_year, max_year = int(df_analysis['year_numeric'].min()), int(df_analysis['year_numeric'].max())
    all_years_range = list(range(min_year, max_year + 1))

    # Create a pivot table: Authors as index, Years as columns, Counts as values
    author_production_data = []
    for author in top_authors_list:
        # Filter documents by this author
        author_docs = df_analysis[
            df_analysis['authors'].astype(str).str.contains(re.escape(author), case=False, na=False)]
        if not author_docs.empty:
            production_by_year = author_docs.groupby('year_numeric').size().reindex(all_years_range, fill_value=0)
            production_by_year.name = author
            author_production_data.append(production_by_year)

    if not author_production_data:
        _callback("No production data found for the selected top authors.")
        return pd.DataFrame(columns=['Year'] + top_authors_list), None

    results_df = pd.concat(author_production_data, axis=1).fillna(0).astype(int)
    results_df.index.name = 'Year'
    results_df = results_df.reset_index()

    _callback("Production data table generated. Creating Plotly figure...")
    fig = go.Figure()
    for author in top_authors_list:
        if author in results_df.columns:
            fig.add_trace(go.Scatter(
                x=results_df['Year'],
                y=results_df[author],
                mode='lines+markers',
                name=author,
                hovertemplate=f"<b>{author}</b><br>Year: %{{x}}<br>Publications: %{{y}}<extra></extra>"
            ))

    fig.update_layout(
        title=f"Top {len(top_authors_list)} Authors' Publication Production Over Time",
        xaxis_title="Year",
        yaxis_title="Number of Publications",
        legend_title_text="Authors",
        hovermode="x unified"  # Shows all author data for a given year on hover
    )
    # Theming will be applied by the receiving widget (AuthorsAnalysisWidget)

    _callback("Authors' Production Over Time analysis complete.")
    return results_df, fig

def _plot_author_career_trajectory(df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    """
    Creates a dual-axis career trajectory plot for a single author, showing
    their yearly publication count and average citation impact.
    """
    # Use the helper to get the most productive author if none is specified in params
    author_to_plot = params.get("author_name") or _get_top_author(df)

    if not author_to_plot:
        return pd.DataFrame(), go.Figure().add_annotation(text="No author available to analyze trajectory.",
                                                          showarrow=False)

    # Filter the main dataframe to get only the papers by the focus author
    author_df = df[df['authors'].str.contains(re.escape(author_to_plot), na=False)].copy()

    # Ensure 'year' and 'citations' are numeric and clean
    author_df['year'] = pd.to_numeric(author_df['year'], errors='coerce')
    author_df['citations'] = pd.to_numeric(author_df['citations'], errors='coerce').fillna(0)
    author_df.dropna(subset=['year'], inplace=True)
    author_df['year'] = author_df['year'].astype(int)

    if author_df.empty:
        return pd.DataFrame(), go.Figure().add_annotation(
            text=f"No valid publication data found for author: {author_to_plot}", showarrow=False)

    # Aggregate the data by year to get publication counts and average citations
    trajectory_df = author_df.groupby('year').agg(
        Publications=('title', 'count'),
        AvgCitations=('citations', 'mean')
    ).reset_index()

    # Ensure a continuous time series by reindexing to the full year range for this author
    full_year_range = range(trajectory_df['year'].min(), trajectory_df['year'].max() + 1)
    trajectory_df = trajectory_df.set_index('year').reindex(full_year_range, fill_value=0).reset_index()

    trajectory_df['AvgCitations'] = trajectory_df['AvgCitations'].round(2)

    # --- Create the Plot using Plotly Graph Objects for dual-axis control ---
    fig = go.Figure()

    # Add the Bar chart for Publications on the primary y-axis
    fig.add_trace(go.Bar(
        x=trajectory_df['year'],
        y=trajectory_df['Publications'],
        name='Publications',
        marker_color='cornflowerblue',
        hovertemplate="Year: %{x}<br>Publications: %{y}<extra></extra>"
    ))

    # Add the Line chart for Average Citations on the secondary y-axis
    fig.add_trace(go.Scatter(
        x=trajectory_df['year'],
        y=trajectory_df['AvgCitations'],
        name='Avg. Citations',
        mode='lines+markers',
        yaxis='y2',  # Assign this trace to the secondary y-axis
        marker=dict(color='darkorange', size=8),
        line=dict(width=3),
        hovertemplate="Year: %{x}<br>Avg. Citations: %{y:.2f}<extra></extra>"
    ))

    # --- Update the layout for a professional dual-axis chart ---
    fig.update_layout(
        title_text=f"Career Trajectory: <b>{author_to_plot}</b>",
        xaxis_title="Year",
        yaxis=dict(
            title="<b>Publications</b>",
            # titlefont=dict(color="cornflowerblue"),
            tickfont=dict(color="cornflowerblue"),
        ),
        yaxis2=dict(
            title="<b>Average Citations per Paper</b>",
            # titlefont=dict(color="darkorange"),
            tickfont=dict(color="darkorange"),
            anchor="x",
            overlaying="y",
            side="right"
        ),
        legend=dict(x=0.01, y=0.99, xanchor='left', yanchor='top'),
        barmode='group'
    )

    # Return the aggregated data for the table view and the final figure
    return trajectory_df, fig


def _plot_author_thematic_evolution(df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    """Creates a heatmap of an author's keyword usage over time."""
    # <<< FIX: Use helper to get a default author if none is provided >>>
    author_to_plot = params.get("author_name") or _get_top_author(df)
    top_n_keywords = params.get("top_n_keywords", 10)
    if not author_to_plot:
        return pd.DataFrame(), go.Figure().add_annotation(text="No author available to analyze.", showarrow=False)
    # ... (rest of the function is correct)
    author_df = df[df['authors'].str.contains(re.escape(author_to_plot), na=False)].copy()
    author_df['year'] = pd.to_numeric(pd.Series(author_df['year']), errors='coerce').dropna().astype(int)
    author_kws = author_df.explode('controlled_vocabulary_terms')['controlled_vocabulary_terms'].dropna()
    top_author_keywords = [kw for kw, count in Counter(author_kws).most_common(top_n_keywords)]
    if not top_author_keywords: return pd.DataFrame(), go.Figure().add_annotation(text=f"No keywords for {author_to_plot}.", showarrow=False)
    records = [{'Year': row['year'], 'Keyword': kw} for _, row in author_df.iterrows() for kw in set(row.get('controlled_vocabulary_terms', [])) if kw in top_author_keywords]
    if not records: return pd.DataFrame(), go.Figure().add_annotation(text="No keyword data to plot.", showarrow=False)
    pivot_df = pd.DataFrame(records).groupby(['Year', 'Keyword']).size().unstack(fill_value=0)
    fig = px.imshow(pivot_df.T, text_auto=True, aspect="auto", title=f"Thematic Evolution for {author_to_plot}")
    return pivot_df.reset_index(), fig

def _plot_overall_citation_trend(df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    """Plots the average citations per paper across all years."""
    df_analysis = df.dropna(subset=['year', 'citations']).copy()
    df_analysis['year'] = pd.to_numeric(pd.Series(df_analysis['year']), errors='coerce').dropna().astype(int)
    df_analysis['citations'] = pd.to_numeric(pd.Series(df_analysis['citations']), errors='coerce').dropna().astype(int)

    trend_df = df_analysis.groupby('year')['citations'].mean().round(2).reset_index()
    fig = px.line(trend_df, x='year', y='citations', markers=True, title="Overall Citation Trend Per Year")
    fig.update_layout(yaxis_title="Average Citations per Paper")
    return trend_df, fig

def _plot_keyword_emergence_slope(df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    """Creates a slope chart to show emerging and declining keywords over time."""
    top_n = params.get("top_n_keywords", 10)

    df_analysis = df.dropna(subset=['year', 'controlled_vocabulary_terms']).copy()
    df_analysis['year'] = pd.to_numeric(pd.Series(df_analysis['year']), errors='coerce').dropna().astype(int)

    if len(df_analysis['year'].unique()) < 2:
        return pd.DataFrame(), go.Figure().add_annotation(text="Not enough year data for trend analysis.",
                                                          showarrow=False)

    # Divide time into two periods
    median_year = df_analysis['year'].median()
    df_period1 = df_analysis[df_analysis['year'] <= median_year]
    df_period2 = df_analysis[df_analysis['year'] > median_year]

    if df_period1.empty or df_period2.empty:
        return pd.DataFrame(), go.Figure().add_annotation(text="Data missing in one of the time periods.",
                                                          showarrow=False)

    counts_p1 = Counter(df_period1.explode('controlled_vocabulary_terms')['controlled_vocabulary_terms'].dropna())
    counts_p2 = Counter(df_period2.explode('controlled_vocabulary_terms')['controlled_vocabulary_terms'].dropna())

    # Normalize frequencies
    norm_counts_p1 = {kw: count / len(df_period1) for kw, count in counts_p1.items()}
    norm_counts_p2 = {kw: count / len(df_period2) for kw, count in counts_p2.items()}

    all_kws = set(norm_counts_p1.keys()) | set(norm_counts_p2.keys())
    trend_data = [{'Keyword': kw, 'Change': norm_counts_p2.get(kw, 0) - norm_counts_p1.get(kw, 0),
                   'Period1_Freq': norm_counts_p1.get(kw, 0), 'Period2_Freq': norm_counts_p2.get(kw, 0)} for kw in
                  all_kws]

    trends_df = pd.DataFrame(trend_data).sort_values('Change', ascending=False)

    emerging = trends_df.head(top_n)
    fading = trends_df.tail(top_n)
    plot_df_data = pd.concat([emerging, fading]).drop_duplicates()

    fig = go.Figure()
    period1_label = f"Early Period (≤{int(median_year)})"
    period2_label = f"Late Period (>{int(median_year)})"

    for _, row in plot_df_data.iterrows():
        fig.add_trace(go.Scatter(x=[period1_label, period2_label], y=[row['Period1_Freq'], row['Period2_Freq']],
                                 mode='lines+markers', name=row['Keyword'],
                                 line=dict(color='green' if row['Change'] > 0 else 'red', width=1),
                                 marker=dict(size=8)))

    fig.update_layout(title="Keyword Emergence & Decline (Normalized Frequency)", showlegend=False)
    return trends_df, fig

def analyze_author_trends(df: pd.DataFrame, params: dict, callback: callable) -> tuple[pd.DataFrame, go.Figure]:
    plot_type = params.get("plot_type")
    callback(f"Analyzing Author Trends: {plot_type}...")

    plot_map = {
        'top_authors_production_over_time': analyze_authors_production_over_time,
        'author_career_trajectory': _plot_author_career_trajectory,
        'author_thematic_evolution': _plot_author_thematic_evolution,
        'keyword_emergence_slope': _plot_keyword_emergence_slope,
        'overall_citation_trend': _plot_overall_citation_trend
    }

    plot_func = plot_map.get(plot_type)
    if plot_func:
        # Some functions don't expect the 'callback' argument, so we call them differently
        if plot_type == 'top_authors_production_over_time':
            return plot_func(df, params, callback)
        return plot_func(df, params)

    msg = f"Unknown trend plot type: '{plot_type}'"
    callback(f"Error: {msg}")
    return pd.DataFrame(), go.Figure().add_annotation(text=msg, showarrow=False)

def analyze_keyword_heatmap(
    df: pd.DataFrame,
    params: dict,
    progress_callback=None,
    zotero_client=None,
    collection_name_for_cache=None
):
    """
    Build and return a heatmap of keyword frequencies per author.
    """
    def _cb(m):
        if progress_callback:
            progress_callback(f"KeywordHeatmap: {m}")
        logging.info(f"KeywordHeatmap: {m}")

    # 1️⃣ Params
    source       = params.get("data_source", "controlled_vocabulary_terms")
    top_n        = params.get("num_top_words", 20)
    author_field = params.get("author_field", "authors")
    min_freq     = params.get("min_frequency", 1)

    _cb(f"Building heatmap for top {top_n} {source} vs. {author_field}…")

    # 2️⃣ Extract tokens per document
    df2 = df.dropna(subset=[author_field])
    df2['tokens'] = df2.apply(
        lambda row: get_text_corpus_from_df(
            pd.DataFrame([row]), source, collection_name_for_cache, zotero_client, None
        ),
        axis=1
    )

    # 3️⃣ Count per-author
    from collections import defaultdict, Counter
    author_map = defaultdict(Counter)
    for _, row in df2.iterrows():
        authors = [a.strip() for a in str(row[author_field]).split(';') if a.strip()]
        kws = set(row['tokens'])
        for auth in authors:
            author_map[auth].update(kws)

    # 4️⃣ Global top-N keywords
    global_counts = Counter()
    for c in author_map.values():
        global_counts.update(c)
    top_keywords = [kw for kw, _ in global_counts.most_common(top_n)]
    _cb(f"Top keywords: {top_keywords}")

    # 5️⃣ Build author×keyword matrix
    authors = list(author_map.keys())
    matrix = pd.DataFrame(
        [[author_map[a].get(kw, 0) for kw in top_keywords] for a in authors],
        index=authors,
        columns=top_keywords
    )
    # Keep only those keywords whose max frequency ≥ min_freq
    keep = [kw for kw in matrix.columns if matrix[kw].max() >= min_freq]
    matrix = matrix[keep]

    if matrix.empty:
        _cb("No data above threshold for heatmap.")
        fig = go.Figure().add_annotation(text="No data for heatmap", showarrow=False)
        return pd.DataFrame(), fig

    # 6️⃣ Melt for table and plot heatmap
    table_df = (
        matrix
        .reset_index()
        .melt(id_vars=['index'], var_name='Keyword', value_name='Frequency')
        .rename(columns={'index': author_field})
    )

    fig = px.imshow(
        matrix,
        labels=dict(x="Keyword", y=author_field, color="Frequency"),
        x=matrix.columns, y=matrix.index,
        color_continuous_scale="Blues",
        text_auto=True,
        aspect="auto",
        title=f"Keyword vs. {author_field.title()} Heatmap"
    )

    # refine each cell’s hover text
    fig.update_traces(
        hovertemplate="Author: %{y}<br>Keyword: %{x}<br>Count: %{z}<extra></extra>",
        selector=dict(type="heatmap")
    )

    # overall layout polish
    fig.update_layout(
        template="plotly_white",  # crisp white canvas
        paper_bgcolor="white",
        plot_bgcolor="white",
        margin=dict(l=100, r=40, t=80, b=120),

        title=dict(
            text=f"<b>Keyword vs. {author_field.title()} Heatmap</b>",
            x=0.5,  # center title
            font=dict(family="Arial", size=20, color="#1f77b4")
        ),

        xaxis=dict(
            side="bottom",
            tickangle=-45,
            tickfont=dict(family="Arial", size=12, color="#333"),
            title="",  # omit redundant axis title
        ),
        yaxis=dict(
            tickfont=dict(family="Arial", size=12, color="#333"),
            title="",  # omit redundant axis title
        ),

        coloraxis_colorbar=dict(
            title="Frequency",
            tickfont=dict(family="Arial", size=11, color="#333"),
            thickness=15,
            lenmode="fraction", len=0.5,
            yanchor="middle", y=0.5,
            outlinewidth=0,
        ),

        font=dict(family="Arial", color="#333"),  # fall-back font for text outside axes
    )

    _cb("Keyword heatmap ready.")
    return table_df, fig
def analyze_keyword_cooccurrence_network(
        df: pd.DataFrame,
        params: dict,
        progress_callback=None,
        zotero_client_for_pdf=None,
        collection_name_for_cache=None
):
    def _callback(msg):
        if progress_callback: progress_callback(f"KeywordCooccurrence: {msg}")
        logging.info(f"KeywordCooccurrence: {msg}")

    # --- 1. Get Parameters ---
    data_source_param = params.get("data_source", "controlled_vocabulary_terms")
    min_cooccurrence = params.get("min_cooccurrence", 2)
    min_keyword_frequency = params.get("min_keyword_frequency", 3)
    max_nodes_to_visualize = params.get("max_nodes_for_plot", 75)
    plot_type = params.get("plot_type", "network_graph")
    community_algo = params.get("community_algorithm", "louvain")

    _callback(f"Starting Keyword Co-occurrence from '{data_source_param}'...")

    # --- 2. Get Keywords Per Document ---
    docs_keywords_list = []
    for index, row in df.iterrows():
        kw_entry = row.get(data_source_param)
        doc_tokens = []
        if isinstance(kw_entry, list):
            doc_tokens = [str(kw).strip().lower() for kw in kw_entry if isinstance(kw, str) and kw.strip()]
        elif isinstance(kw_entry, str) and data_source_param == "controlled_vocabulary_terms":
            doc_tokens = [kw.strip().lower() for kw in kw_entry.split(';') if kw.strip()]
        elif isinstance(kw_entry, str):  # For abstract/title/fulltext
            doc_tokens = preprocess_text(kw_entry)  # Assumes preprocess_text handles these sources
        docs_keywords_list.append(list(set(doc_tokens)))

    if not any(docs_keywords_list):
        msg = "No keywords/tokens found in the selected data source."
        _callback(msg);
        return pd.DataFrame(), go.Figure().add_annotation(text=msg, showarrow=False)
        # --- 3. Build Graph (remains the same) ---
    all_keywords = [kw for sublist in docs_keywords_list for kw in sublist]
    keyword_overall_freq = Counter(all_keywords)
    valid_keywords = {kw for kw, count in keyword_overall_freq.items() if count >= min_keyword_frequency}
    if not valid_keywords:
        return pd.DataFrame(), go.Figure().add_annotation(text=f"No keywords meet min freq of {min_keyword_frequency}.",
                                                          showarrow=False)

    filtered_docs_keywords = [[kw for kw in doc_kws if kw in valid_keywords] for doc_kws in docs_keywords_list]
    cooccurrence_counts = Counter(pair for doc_kws in filtered_docs_keywords if len(doc_kws) >= 2 for pair in
                                  itertools.combinations(sorted(doc_kws), 2))
    edge_list = [{"Keyword1": p[0], "Keyword2": p[1], "Weight": c} for p, c in cooccurrence_counts.items() if
                 c >= min_cooccurrence]
    if not edge_list:
        return pd.DataFrame(), go.Figure().add_annotation(
            text=f"No pairs meet min co-occurrence of {min_cooccurrence}.", showarrow=False)

    results_df = pd.DataFrame(edge_list).sort_values(by="Weight", ascending=False).reset_index(drop=True)
    G = nx.from_pandas_edgelist(results_df, 'Keyword1', 'Keyword2', ['Weight'])
    nx.set_node_attributes(G, {node: keyword_overall_freq.get(node, 0) for node in G.nodes()}, 'frequency')

    # --- 4. Community Detection (remains the same) ---
    communities = None
    try:
        if community_algo == "louvain":
            import community as community_louvain
            communities = community_louvain.best_partition(G, weight='weight', random_state=42)
        elif community_algo == "label_propagation":
            community_sets = list(nx.community.label_propagation_communities(G))
            communities = {node: i for i, comm_set in enumerate(community_sets) for node in comm_set}
        if communities:
            nx.set_node_attributes(G, communities, 'community')
            _callback(f"Community detection found {len(set(communities.values()))} communities.")
    except Exception as e:
        _callback(f"Community detection failed: {e}")

    # --- 5. Create Plot (NEW SOPHISTICATED VERSION) ---
    fig = go.Figure()

    if plot_type == "network_graph":
        G_viz = G
        if G.number_of_nodes() > max_nodes_to_visualize:
            top_nodes = sorted(G.degree(weight='weight'), key=lambda x: x[1], reverse=True)[:max_nodes_to_visualize]
            G_viz = G.subgraph([n[0] for n in top_nodes]).copy()

        if G_viz.number_of_nodes() == 0:
            return results_df, fig.add_annotation(text="No nodes to display in network.", showarrow=False)

        # --- Layout Algorithm ---
        pos = nx.spring_layout(G_viz, k=0.8 / np.sqrt(G_viz.number_of_nodes()) if G_viz.number_of_nodes() > 1 else 0.5,
                               iterations=70, seed=42, weight='weight')

        # --- Edge Trace ---
        edge_x, edge_y = [], []
        for edge in G_viz.edges():
            x0, y0 = pos[edge[0]];
            x1, y1 = pos[edge[1]]
            edge_x.extend([x0, x1, None]);
            edge_y.extend([y0, y1, None])

        edge_trace = go.Scatter(
            x=edge_x, y=edge_y,
            line=dict(width=0.7, color=THEME.get('BORDER_SECONDARY', '#444')),
            hoverinfo='none',
            mode='lines')
        fig.add_trace(edge_trace)

        # --- Node Trace ---
        node_x, node_y, node_text, node_hover_text, node_sizes, node_colors = [], [], [], [], [], []

        # Determine threshold for showing labels directly on the plot
        degrees = [d for n, d in G_viz.degree(weight='weight')]
        if len(degrees) > 15:
            label_threshold = np.percentile(degrees, 85)  # Show labels for top 15% of nodes by degree
        else:
            label_threshold = 0  # Show all labels if graph is small

        for node in G_viz.nodes():
            x, y = pos[node]
            node_x.append(x);
            node_y.append(y)

            freq = G_viz.nodes[node].get('frequency', 1)
            degree = G_viz.degree(node, weight='weight')

            # Node Size based on log-scaled frequency
            node_sizes.append(10 + 30 * np.log1p(freq) / np.log1p(max(keyword_overall_freq.values())))

            # Node Color based on community
            if communities:
                node_colors.append(G_viz.nodes[node].get('community', 0))
            else:
                node_colors.append(degree)  # Fallback to coloring by degree

            # Hover Text
            hover_info = f"<b>{html.escape(node)}</b><br>Overall Frequency: {freq}<br>Connections Strength: {int(degree)}"
            if communities:
                hover_info += f"<br>Community ID: {G_viz.nodes[node].get('community', 'N/A')}"
            node_hover_text.append(hover_info)

            # On-plot Text Label
            if degree >= label_threshold:
                node_text.append(html.escape(node))
            else:
                node_text.append('')

        node_trace = go.Scatter(
            x=node_x, y=node_y,
            mode='markers+text',
            hoverinfo='text',
            hovertext=node_hover_text,
            text=node_text,
            textposition="bottom center",
            textfont=dict(size=10, color=THEME.get('TEXT_SECONDARY')),
            marker=dict(
                showscale=not bool(communities),  # Show scale only if coloring by a continuous value (degree)
                colorscale='YlGnBu' if not communities else 'Rainbow',
                reversescale=True,
                color=node_colors,
                size=node_sizes,
                line=dict(width=1, color=THEME.get('BORDER_PRIMARY')),
                sizemin=4,
                colorbar=dict(
                    thickness=15,
                    title='Node Strength',
                    xanchor='left',

                ) if not communities else None
            )
        )
        fig.add_trace(node_trace)

        fig.update_layout(
            title_text=f'Thematic Network of Keywords from {data_source_param.title()}',
            showlegend=False,
            hovermode='closest',
            margin=dict(b=10, l=10, r=10, t=50),
            xaxis=dict(showgrid=False, zeroline=False, showticklabels=False),
            yaxis=dict(showgrid=False, zeroline=False, showticklabels=False),
            plot_bgcolor=THEME.get('BACKGROUND_CONTENT_AREA_PLOT', 'rgba(0,0,0,0)'),
            paper_bgcolor='rgba(0,0,0,0)'
        )


    else:
        fig.add_annotation(text=f"Unknown plot type '{plot_type}'.", showarrow=False)

    _callback("Keyword Co-occurrence analysis complete.")
    return results_df, fig


# =========================================================================
# === 2. "PRIVATE" HELPER FUNCTIONS FOR PLOTTING (Using Standardized Columns) ===
# =========================================================================

def _plot_h_index_bar(stats_df: pd.DataFrame, params: dict) -> go.Figure:
    top_n = params.get("top_n_authors", 15)
    plot_df = stats_df.sort_values(by=['H-Index', 'TotalCitations'], ascending=[False, False]).head(top_n)
    fig = px.bar(plot_df, x="Author", y="H-Index", title=f"Top {len(plot_df)} Authors by H-Index",
                 labels={"H-Index": "h-index"}, hover_data=["TotalPublications", "TotalCitations"],
                 color="H-Index", color_continuous_scale=px.colors.sequential.thermal, text_auto=True)
    fig.update_traces(textposition='outside');
    fig.update_xaxes(categoryorder="total descending")
    return fig


def _plot_productivity_bar(stats_df: pd.DataFrame, params: dict) -> go.Figure:
    top_n = params.get("top_n_authors", 15)
    plot_df = stats_df.sort_values(by=['TotalPublications', 'TotalCitations'], ascending=[False, False]).head(top_n)
    fig = px.bar(plot_df, x="Author", y="TotalPublications", title=f"Top {len(plot_df)} Most Productive Authors",
                 labels={"TotalPublications": "Number of Publications"}, hover_data=["TotalCitations", "H-Index"],
                 color="TotalPublications", color_continuous_scale=px.colors.sequential.Teal, text_auto=True)
    fig.update_traces(textposition='outside');
    fig.update_xaxes(categoryorder="total descending")
    return fig


def _plot_total_citations_bar(stats_df: pd.DataFrame, params: dict) -> go.Figure:
    top_n = params.get("top_n_authors", 15)
    plot_df = stats_df.sort_values(by=['TotalCitations', 'TotalPublications'], ascending=[False, False]).head(top_n)
    fig = px.bar(plot_df, x="Author", y="TotalCitations", title=f"Top {len(plot_df)} Authors by Total Citations",
                 labels={"TotalCitations": "Total Citations"}, hover_data=["TotalPublications", "H-Index"],
                 color="TotalCitations", color_continuous_scale=px.colors.sequential.Plasma, text_auto=True)
    fig.update_traces(textposition='outside');
    fig.update_xaxes(categoryorder="total descending")
    return fig


def _plot_avg_citations_bar(stats_df: pd.DataFrame, params: dict) -> go.Figure:
    top_n = params.get("top_n_authors", 15)
    plot_df = stats_df.sort_values(by=['AvgCitationsPerPub', 'TotalCitations'], ascending=[False, False]).head(top_n)
    fig = px.bar(plot_df, x="Author", y="AvgCitationsPerPub",
                 title=f"Top {len(plot_df)} Authors by Avg. Citations per Publication",
                 labels={"AvgCitationsPerPub": "Avg. Citations / Pub"},
                 hover_data=["TotalPublications", "TotalCitations"],
                 color="AvgCitationsPerPub", color_continuous_scale=px.colors.sequential.Viridis, text_auto=".2f")
    fig.update_traces(textposition='outside');
    fig.update_xaxes(categoryorder="total descending")
    return fig


def _plot_productivity_treemap(stats_df: pd.DataFrame, params: dict) -> go.Figure:
    top_n = params.get("top_n_authors", 30)
    plot_df = stats_df.sort_values("TotalPublications", ascending=False).head(top_n)
    fig = px.treemap(plot_df, path=[px.Constant(f"Top {len(plot_df)} Authors"), 'Author'],
                     values='TotalPublications', title=f"Productivity Treemap of Top {len(plot_df)} Authors",
                     color='TotalCitations', color_continuous_scale='Blues',
                     hover_data={'TotalPublications': ':.0f', 'TotalCitations': ':.0f', 'H-Index': ':.0f'})
    fig.update_traces(textinfo='label+value', root_color="lightgrey")
    return fig
# code for replacement
def _plot_productivity_vs_impact_scatter(stats_df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    """
    Returns (plot_df, fig).

    Variants (params['variant']):
      - 'beeswarm' (default): strip/beeswarm by publications bucket (robust on old Plotly).
      - 'lollipop': ranked horizontal lollipop by TotalCitations.
      - 'classic': jittered scatter (for completeness).

    Other params:
      - top_n_authors (int, default 25)
      - label_top_k (int, default min(15, top_n))
      - jitter (float, default 0.08)   # used in beeswarm/classic
      - seed (int, default 42)
    """
    import numpy as np
    import pandas as pd
    import plotly.graph_objects as go

    variant      = str(params.get("variant", "beeswarm")).lower()
    top_n        = int(params.get("top_n_authors", 25))
    label_top_k  = int(params.get("label_top_k", min(15, top_n)))
    jitter       = float(params.get("jitter", 0.08))
    seed         = int(params.get("seed", 42))

    # ---- prepare data ----
    df = stats_df.copy()
    for c in ("TotalPublications", "TotalCitations", "H-Index", "AvgCitationsPerPub"):
        if c in df.columns:
            df[c] = pd.to_numeric(df[c], errors="coerce")
    df = df.dropna(subset=["Author", "TotalPublications", "TotalCitations"])

    if df.empty:
        fig = go.Figure().add_annotation(text="No author data available for scatter.",
                                         showarrow=False, x=0.5, y=0.5, xref="paper", yref="paper")
        return df, fig

    plot_df = df.sort_values("TotalCitations", ascending=False).head(top_n).copy()

    # marker size from H-index (works if missing)
    H = plot_df.get("H-Index")
    if H is not None:
        h = pd.to_numeric(pd.Series(H), errors="coerce").fillna(0).astype(float)
        if h.size > 0 and np.nanmax(h.to_numpy()) > 0:
            sizes = 10 + 30 * np.sqrt((h - h.min()) / (h.max() - h.min() + 1e-12))
        else:
            sizes = pd.Series(16.0, index=plot_df.index)
    else:
        sizes = pd.Series(16.0, index=plot_df.index)

    # label only top-K by citations
    label_set = set(plot_df.sort_values("TotalCitations", ascending=False)["Author"].head(label_top_k))
    labels = plot_df["Author"].where(plot_df["Author"].isin(label_set), "")

    # color by AvgCitationsPerPub as continuous (works on old Plotly GO)
    acp = pd.to_numeric(plot_df.get("AvgCitationsPerPub", pd.Series([np.nan]*len(plot_df))), errors="coerce")

    rng = np.random.RandomState(seed)

    if variant == "lollipop":
        top = plot_df.sort_values("TotalCitations", ascending=False).copy()
        cat = list(top["Author"])[::-1]
        top["AuthorCat"] = pd.Categorical(top["Author"], categories=cat, ordered=True)

        # stems
        x_stem, y_stem = [], []
        for _, r in top.iterrows():
            x_stem += [0, float(r["TotalCitations"]), None]
            y_stem += [r["AuthorCat"], r["AuthorCat"], None]

        fig = go.Figure()
        fig.add_trace(go.Scatter(x=x_stem, y=y_stem, mode="lines",
                                 line=dict(color="rgba(0,0,0,0.25)", width=2),
                                 hoverinfo="skip", showlegend=False))
        fig.add_trace(go.Scatter(
            x=top["TotalCitations"], y=top["AuthorCat"],
            mode="markers+text",
            marker=dict(
                size=sizes.reindex(top.index).fillna(14),
                color=acp.reindex(top.index),
                colorscale="Blues",
                colorbar=dict(title="Avg cites / pub"),
                line=dict(width=1, color="rgba(0,0,0,0.6)")
            ),
            text=top["Author"],
            textposition="middle right",
            textfont=dict(size=10),
            hovertemplate="<b>%{text}</b><br>Citations: %{x}<br>"
                          "Pubs: %{customdata[0]}<br>"
                          "Avg cites/pub: %{marker.color:.2f}<extra></extra>",
            customdata=np.c_[pd.to_numeric(pd.Series(top['TotalPublications']), errors='coerce').fillna(0).to_numpy()],
            showlegend=False
        ))
        fig.update_layout(
            title=f"Top {len(top)} Authors by Total Citations — Lollipop",
            title_x=0.5, template="plotly_white",
            margin=dict(l=100, r=40, t=60, b=50),
            xaxis_title="Total Citations", yaxis_title=""
        )
        fig.update_xaxes(gridcolor="rgba(0,0,0,0.08)")
        fig.update_yaxes(showgrid=False)
        return top, fig

    if variant == "classic":
        crowded = plot_df.groupby("TotalPublications")["Author"].transform("count").clip(lower=1)
        scale = 1 + (crowded - 1) / crowded.max()
        x = plot_df["TotalPublications"] + rng.uniform(-jitter, jitter, size=len(plot_df)) * scale
        fig = go.Figure(go.Scatter(
            x=x, y=plot_df["TotalCitations"],
            mode="markers+text",
            text=labels, textposition="top center", textfont=dict(size=11),
            marker=dict(size=sizes, color="rgba(0, 102, 204, 0.75)",
                        line=dict(width=1, color="rgba(0,0,0,0.6)")),
            hovertemplate="<b>%{text}</b><br>Publications: %{x:.2f}<br>Citations: %{y}<extra></extra>",
            showlegend=False
        ))
        x_min, x_max = float(x.min()), float(x.max())
        y_max = float(plot_df["TotalCitations"].max())
        pad_x = max(0.35, (x_max - x_min) * 0.15 or 0.35)
        pad_y = max(2.0, y_max * 0.10)
        fig.update_layout(
            title=f"Citation Impact vs. Productivity (Top {len(plot_df)} Authors)",
            title_x=0.5, template="plotly_white",
            margin=dict(l=50, r=30, t=60, b=55),
            xaxis_title="Number of Publications", yaxis_title="Total Citations",
            hovermode="closest"
        )
        fig.update_xaxes(range=[max(0.0, x_min - pad_x), x_max + pad_x], gridcolor="rgba(0,0,0,0.08)")
        fig.update_yaxes(range=[-0.5, y_max + pad_y], gridcolor="rgba(0,0,0,0.08)")
        return plot_df, fig

    # ---- beeswarm (default) ----
    plot_df["PubsBucket"] = pd.to_numeric(pd.Series(plot_df["TotalPublications"]), errors="coerce").fillna(0).astype(int)
    # manual jitter per bucket, no px.strip args required
    x_vals = []
    for b, idx in plot_df.groupby("PubsBucket").groups.items():
        n = len(idx)
        # symmetric offsets so the swarm looks balanced on old Plotly too
        if n == 1:
            offsets = np.array([0.0])
        else:
            offsets = np.linspace(-0.35, 0.35, n)
            offsets += rng.uniform(-0.03, 0.03, size=n)  # small randomization
        x_vals.extend(b + offsets)
    x_vals = np.array(x_vals)

    text_vals = plot_df["Author"].where(plot_df["Author"].isin(label_set), "")

    fig = go.Figure(go.Scatter(
        x=x_vals,
        y=plot_df["TotalCitations"],
        mode="markers+text",
        text=text_vals,
        textposition="top center",
        textfont=dict(size=10),
        marker=dict(
            size=sizes,
            color=acp,
            colorscale="Blues",
            colorbar=dict(title="Avg cites / pub"),
            line=dict(width=1, color="rgba(0,0,0,0.55)")
        ),
        hovertemplate="<b>%{text}</b><br>"
                      "Pubs bucket: %{x:.2f}<br>"
                      "Citations: %{y}<br>"
                      "<extra></extra>",
        showlegend=False
    ))
    y_max = float(plot_df["TotalCitations"].max())
    fig.update_layout(
        title=f"Citation Impact vs. Productivity — Beeswarm (Top {len(plot_df)})",
        title_x=0.5, template="plotly_white",
        margin=dict(l=60, r=40, t=60, b=60),
        xaxis_title="Number of Publications (buckets)",
        yaxis_title="Total Citations",
        hovermode="closest"
    )
    # x ticks at integer buckets present
    buckets = sorted(plot_df["PubsBucket"].unique().tolist())
    fig.update_xaxes(tickmode="array", tickvals=buckets, ticktext=[str(b) for b in buckets],
                     gridcolor="rgba(0,0,0,0.08)")
    fig.update_yaxes(range=[-0.5, y_max + max(2.0, y_max * 0.10)],
                     gridcolor="rgba(0,0,0,0.08)")
    return plot_df, fig



def _plot_author_keyword_heatmap(df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    top_n_authors = params.get("top_n_authors", 15)
    top_n_keywords = params.get("top_n_keywords", 20)

    author_counts = Counter(
        author.strip() for authors_str in df['authors'].dropna() for author in authors_str.split(';'))
    top_authors = [author for author, count in author_counts.most_common(top_n_authors)]

    keyword_counts = Counter(kw for kws_list in df['controlled_vocabulary_terms'].dropna() for kw in kws_list)
    top_keywords = [kw for kw, count in keyword_counts.most_common(top_n_keywords)]

    matrix = pd.DataFrame(0, index=top_authors, columns=top_keywords)
    for _, row in df.iterrows():
        authors_in_doc = {author.strip() for author in row.get('authors', '').split(';')}
        keywords_in_doc = set(row.get('controlled_vocabulary_terms', []))
        for author in authors_in_doc:
            if author in matrix.index:
                for keyword in keywords_in_doc:
                    if keyword in matrix.columns:
                        matrix.loc[author, keyword] += 1

    matrix = matrix.loc[(matrix.sum(axis=1) > 0), (matrix.sum(axis=0) > 0)]
    if matrix.empty: return pd.DataFrame(), go.Figure().add_annotation(
        text="No overlapping data for Author-Keyword Heatmap")

    fig = px.imshow(matrix, text_auto=True, aspect="auto",
                    labels=dict(x="Keyword", y="Author", color="Frequency"),
                    title=f"Thematic Focus: Top Authors vs. Top Keywords", color_continuous_scale="Blues")
    return matrix.reset_index(), fig


def _plot_author_keyword_treemap(df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    top_n_authors = params.get("top_n_authors", 10)

    author_counts = Counter(
        author.strip() for authors_str in df['authors'].dropna() for author in authors_str.split(';'))
    top_authors = [author for author, count in author_counts.most_common(top_n_authors)]

    records = []
    for author in top_authors:
        author_df = df[df['authors'].str.contains(re.escape(author), na=False)]
        author_kws = Counter(kw for kws_list in author_df['controlled_vocabulary_terms'].dropna() for kw in kws_list)
        for kw, count in author_kws.most_common(5):  # Top 5 keywords per author
            records.append({'Author': author, 'Keyword': kw, 'Count': count})

    if not records: return pd.DataFrame(), go.Figure().add_annotation(text="No keyword data for top authors.")

    treemap_df = pd.DataFrame(records)
    fig = px.treemap(treemap_df, path=[px.Constant("Top Authors"), 'Author', 'Keyword'], values='Count',
                     title=f"Thematic Focus of Top {top_n_authors} Authors",
                     color='Count', color_continuous_scale='YlGnBu')
    return treemap_df, fig


def _plot_animated_productivity_race(df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    top_n_authors = params.get("top_n_authors", 15)

    author_counts = Counter(
        author.strip() for authors_str in df['authors'].dropna() for author in authors_str.split(';'))
    top_authors = {author for author, count in author_counts.most_common(top_n_authors)}

    df_anim = df[['year', 'authors']].dropna().copy()
    df_anim['year'] = pd.to_numeric(pd.Series(df_anim['year']), errors='coerce').dropna().astype(int)

    # Explode authors string into list before exploding DataFrame rows
    df_anim['authors'] = df_anim['authors'].str.split(';')
    df_anim = df_anim.explode('authors')
    df_anim['authors'] = df_anim['authors'].str.strip()

    df_anim = df_anim[df_anim['authors'].isin(top_authors)]

    yearly_counts = df_anim.groupby(['year', 'authors']).size().reset_index(name='count')

    # Calculate cumulative counts correctly
    yearly_counts.sort_values('year', inplace=True)
    yearly_counts['CumulativePubs'] = yearly_counts.groupby('authors')['count'].cumsum()

    fig = px.bar(yearly_counts, y='authors', x='CumulativePubs', color='authors',
                 animation_frame='year', animation_group='authors',
                 orientation='h', title=f"Animated Productivity Race of Top {len(top_authors)} Authors",
                 labels={'CumulativePubs': 'Cumulative Publications', 'authors': 'Author'})

    fig.update_layout(yaxis={'categoryorder': 'total ascending'})
    fig.layout.updatemenus[0].buttons[0].args[1]['frame']['duration'] = 400

    return yearly_counts, fig


# =========================================================================
# === 3. MASTER DISPATCHER FUNCTION (The one called by MainWindow) ===
# =========================================================================

def analyze_categorical_keywords(df: pd.DataFrame, params: dict, progress_callback=None):
    def _callback(msg):
        if progress_callback: progress_callback(msg)
        logging.info(msg)

    selected_category = params.get("selected_category", None)
    plot_type = params.get("plot_type", "bar_chart")  # e.g., bar_chart, pie_chart

    _callback(f"Analyzing categorical keywords for category: '{selected_category}'.")

    if 'controlled_vocabulary_terms' not in df.columns or df['controlled_vocabulary_terms'].isnull().all():
        _callback("Keywords column missing or empty for categorical analysis.")
        return pd.DataFrame(), None
    if not selected_category:
        _callback("No category selected for analysis.")
        return pd.DataFrame({"CategoryValue": ["No Category Selected"], "Count": [0]}), None

    cat_lc = str(selected_category).strip().lower()

    category_values = []
    for kw_list_entry in df['controlled_vocabulary_terms'].dropna():
        if isinstance(kw_list_entry, list):
            for kw in kw_list_entry:
                if isinstance(kw, str) and ':' in kw:
                    cat, val = kw.split(':', 1)
                    if cat.strip().lower() == cat_lc:
                        category_values.append(val.strip())
        elif isinstance(kw_list_entry, str):  # Fallback if keywords are ; separated strings
            for kw in kw_list_entry.split(';'):
                if isinstance(kw, str) and ':' in kw:
                    cat, val = kw.split(':', 1)
                    if cat.strip().lower() == cat_lc:
                        category_values.append(val.strip())

    if not category_values:
        _callback(f"No values found for category '{selected_category}'.")
        return pd.DataFrame({"CategoryValue": [f"No data for {selected_category}"], "Count": [0]}), None

    value_counts = Counter(category_values)
    results_df = pd.DataFrame(value_counts.items(), columns=["Value", "Count"])
    results_df = results_df.sort_values(by="Count", ascending=False).reset_index(drop=True)

    fig = None
    if not results_df.empty:
        try:
            if plot_type == "bar_chart":
                fig = px.bar(results_df, x="Value", y="Count", text_auto=True,
                             title=f"Distribution of Values for Category: {selected_category}")
                fig.update_traces(textposition='outside')
                fig.update_xaxes(categoryorder="total descending")
            elif plot_type == "pie_chart":
                fig = px.pie(results_df, names="Value", values="Count", hole=0.3,
                             title=f"Proportion of Values for Category: {selected_category}")
                fig.update_traces(textposition='inside', textinfo='percent+label')
        except Exception as e:
            _callback(f"Error creating plot for categorical keywords: {e}")
    else:
        _callback("No data to plot for categorical keywords.")

    _callback("Categorical Keyword analysis complete.")
    return results_df, fig


def _get_top_author(df: pd.DataFrame) -> str | None:
    """Gets the most frequent author from the dataframe."""
    if 'authors' not in df.columns or df['authors'].dropna().empty:
        return None
    # Ensure authors are strings before splitting
    authors_series = df['authors'].dropna().astype(str)
    author_counts = Counter(author.strip() for authors_str in authors_series for author in authors_str.split(';'))
    return author_counts.most_common(1)[0][0] if author_counts else None



def analyze_author_impact(df: pd.DataFrame, params: dict, callback: callable) -> tuple[pd.DataFrame, go.Figure]:
    """
    Main dispatcher for all author performance and impact plots.
    Calculates comprehensive author statistics once and then routes to the
    appropriate plotting function.
    """
    plot_type = params.get("plot_type", "h_index_bar")
    callback(f"Analyzing Author Impact: {plot_type}...")

    # This first part for keyword/animated plots is fine
    if plot_type in ["author_keyword_heatmap", "author_keyword_treemap", "animated_productivity_race"]:
        plot_map = {
            "author_keyword_heatmap": _plot_author_keyword_heatmap,
            "author_keyword_treemap": _plot_author_keyword_treemap,
            "animated_productivity_race": _plot_animated_productivity_race,
        }
        return plot_map[plot_type](df, params)

    # For all other plots, calculate the full stats
    author_stats_df = _calculate_all_author_stats(df)
    if author_stats_df.empty:
        msg = "Author analysis failed: No valid author data to process."
        callback(f"Error: {msg}")
        return pd.DataFrame(), go.Figure().add_annotation(text=msg)

    callback(f"Calculated statistics for {len(author_stats_df)} unique authors. Generating plot...")

    # ===================================================================
    # === THE FIX IS HERE: The map is now complete. ===
    # ===================================================================
    plot_function_map = {
        "h_index_bar": _plot_h_index_bar,
        "productivity_bar": _plot_productivity_bar,
        "productivity_treemap": _plot_productivity_treemap,  # <-- ADDED
        "citations_total_bar": _plot_total_citations_bar,
        "citations_avg_bar": _plot_avg_citations_bar,  # <-- ADDED
        "productivity_vs_impact_scatter": _plot_productivity_vs_impact_scatter,
    }
    # ===================================================================

    plot_func = plot_function_map.get(plot_type)

    if plot_func:
        # Pass the pre-calculated stats DataFrame to the plotting helper.
        fig = plot_func(author_stats_df, params)
        # Return the full stats DataFrame (for the table view) and the figure.
        return author_stats_df, fig
    else:
        msg = f"Unknown impact plot type: '{plot_type}'"
        callback(f"Error: {msg}")
        # Return the full stats DataFrame even if the plot type is unknown
        return author_stats_df, go.Figure().add_annotation(text=msg)



def add_authorship_and_institution_section(
    prs,
    df: "pd.DataFrame",
    collection_name: str = "Collection",
    *,
    slide_notes: bool = True,
    top_n_authors: int = 20,
    export: bool = False,
    return_payload: bool = False,
    progress_callback=None,
):
    """
    ###1. build an Authorship & Institutions slide-deck in the same style as authors_overview
    ###2. create 6 slides:
        - Top authors — Table
        - Top authors — Figure
        - Co-authorship — Strongest pairs (Top 20) — Table
        - Co-authorship network — Figure
        - Major institutional actors (publication_outlet) — Table
        - Major institutional actors (publication_outlet) — Figure
    ###3. if export or return_payload: return {"slides":[...], "figs":[...], "index":0}; else return None

    Assumptions (brittle by design):
      - df is a pandas.DataFrame and non-empty.
      - plotly figures support .to_image().
      - tables are pandas.DataFrames.
    """

    import io, itertools, math, re, json
    import numpy as np
    import pandas as pd
    import plotly.express as px
    import plotly.graph_objects as go
    from plotly.utils import PlotlyJSONEncoder
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN

    preview_slides = [] if (export or return_payload) else None
    preview_figs = [] if (export or return_payload) else None

    FONT_FAMILY = "Segoe UI"
    PLOT_W, PLOT_H = 1200, 650

    INCH = float(Inches(1))
    SLIDE_W_IN = float(prs.slide_width) / INCH
    SLIDE_H_IN = float(prs.slide_height) / INCH

    OUTER_MARGIN_X = 0.6
    OUTER_MARGIN_Y = 0.6
    TITLE_HEIGHT_IN = 0.7

    def _cb(msg: str):
        if progress_callback:
            progress_callback(f"AUTH+INST: {msg}")

    def _new_blank_slide_with_title(title_text: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(
            Inches(0.0),
            Inches(0.15),
            Inches(SLIDE_W_IN),
            Inches(TITLE_HEIGHT_IN),
        )
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        return slide

    def _set_notes(slide, text: str):
        if slide_notes and text:
            slide.notes_slide.notes_text_frame.text = text

    def _add_table_slide(title_text: str, table_df: "pd.DataFrame", notes: str = ""):
        slide = _new_blank_slide_with_title(title_text)
        _set_notes(slide, notes)

        df_small = table_df.head(25).copy()

        max_w = SLIDE_W_IN - 2 * OUTER_MARGIN_X
        max_h = SLIDE_H_IN - (TITLE_HEIGHT_IN + 2 * OUTER_MARGIN_Y)

        rows = len(df_small) + 1
        cols = df_small.shape[1]

        def _len_inch_est(s: object) -> float:
            t = "" if pd.isna(s) else str(s)
            t = " ".join(t.split())
            return 0.085 * max(4, len(t)) + 0.14

        col_needs = []
        for j, col in enumerate(df_small.columns):
            header_need = _len_inch_est(col)
            longest_need = header_need
            for i in range(len(df_small)):
                longest_need = max(longest_need, min(_len_inch_est(df_small.iloc[i, j]), 3.8))
            col_needs.append(longest_need)

        MIN_W = 0.9
        MAX_W = 3.6

        col_needs = [min(MAX_W, max(MIN_W, x)) for x in col_needs]
        total_need = float(np.sum(col_needs))

        if total_need > (max_w - 0.2):
            scale = (max_w - 0.2) / total_need
            col_widths_in = [max(MIN_W, x * scale) for x in col_needs]
        else:
            col_widths_in = col_needs[:]

        table_w = float(np.sum(col_widths_in))
        est_row_h = min(0.42, max(0.26, max_h / (rows + 1)))
        table_h = min(max_h, rows * est_row_h + 0.35)

        left_in = (SLIDE_W_IN - table_w) / 2.0
        top_in = TITLE_HEIGHT_IN + ((SLIDE_H_IN - TITLE_HEIGHT_IN) - table_h) / 2.0

        table = slide.shapes.add_table(
            rows, cols, Inches(left_in), Inches(top_in), Inches(table_w), Inches(table_h)
        ).table

        for j, w in enumerate(col_widths_in):
            table.columns[j].width = Emu(Inches(w))

        for j, col in enumerate(df_small.columns):
            cell = table.cell(0, j)
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(col)
            p.font.bold = True
            p.font.name = FONT_FAMILY
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER
            cell.margin_left = cell.margin_right = Emu(900)
            cell.margin_top = Emu(600)
            cell.margin_bottom = Emu(600)

        for i, (_, r) in enumerate(df_small.iterrows(), start=1):
            for j, col in enumerate(df_small.columns):
                cell = table.cell(i, j)
                txt = "" if pd.isna(r[col]) else " ".join(str(r[col]).split())
                tf = cell.text_frame
                tf.clear()
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = txt
                p.font.name = FONT_FAMILY
                p.font.size = Pt(11)
                p.alignment = PP_ALIGN.LEFT
                cell.margin_left = Emu(800)
                cell.margin_right = Emu(800)
                cell.margin_top = Emu(500)
                cell.margin_bottom = Emu(500)

        if preview_slides is not None:
            preview_slides.append(
                {"title": title_text, "table_html": df_small.to_html(index=False, escape=True), "notes": notes}
            )

        return slide

    def _add_figure_slide(title_text: str, fig: "go.Figure", notes: str = ""):
        slide = _new_blank_slide_with_title(title_text)
        _set_notes(slide, notes)

        content_w = SLIDE_W_IN - 2 * OUTER_MARGIN_X
        content_h = SLIDE_H_IN - (TITLE_HEIGHT_IN + 2 * OUTER_MARGIN_Y)
        if export:
            try:
                png = fig.to_image(format="png", width=PLOT_W, height=PLOT_H, scale=2)
                ar = (PLOT_W / PLOT_H) if PLOT_H else 1.7778

                target_w = content_w
                target_h = target_w / ar
                if target_h > content_h:
                    target_h = content_h
                    target_w = target_h * ar
                target_w *= 0.98
                target_h *= 0.98

                left_in = (SLIDE_W_IN - target_w) / 2.0
                top_in = TITLE_HEIGHT_IN + (SLIDE_H_IN - TITLE_HEIGHT_IN - target_h) / 2.0

                slide.shapes.add_picture(
                    io.BytesIO(png),
                    Inches(left_in),
                    Inches(top_in),
                    width=Inches(target_w),
                    height=Inches(target_h),
                )
            except Exception as exc:
                try:
                    print(f"[visualise][warn] plotly.to_image failed ({title_text}): {exc}")
                except Exception:
                    pass

        if preview_slides is not None:
            fig_json = json.dumps(fig.to_plotly_json(), cls=PlotlyJSONEncoder)
            preview_figs.append({"title": title_text, "fig_json": fig_json, "notes": notes or ""})
            preview_slides.append({"title": title_text, "fig_json": fig_json, "table_html": "", "notes": notes or ""})

        return slide

    if df.empty:
        if preview_slides is not None:
            return {
                "slides": [{"title": "Authorship & institutions", "table_html": "<div>No data loaded.</div>", "notes": ""}],
                "figs": [],
                "index": 0,
            }
        return None

    def sstr(x) -> str:
        return "" if pd.isna(x) else str(x)

    def _authors_from_row(row) -> list[str]:
        if "authors_list" in row.index:
            al = row["authors_list"]
            if type(al) in (list, tuple):
                out = []
                for a in al:
                    if type(a) is dict:
                        ln = sstr(a.get("lastName")).strip()
                        fn = sstr(a.get("firstName")).strip()
                        lab = ", ".join([p for p in [ln, fn] if p])
                        if lab:
                            out.append(lab)
                    else:
                        ss = sstr(a).strip()
                        if ss:
                            out.append(ss)
                if out:
                    return out

        for k in ("authors", "creator_summary"):
            if k in row.index:
                s = sstr(row[k]).strip()
                if s:
                    return [p.strip() for p in re.split(r"[;|,]+", s) if p.strip()]

        return []

    def _modal(series: pd.Series, default: str = ""):
        s = series.dropna().map(sstr).map(str.strip)
        s = s[s != ""]
        if s.empty:
            return default
        return s.value_counts().index[0]

    def _col(work: "pd.DataFrame", *names: str) -> "pd.Series":
        for n in names:
            if n in work.columns:
                return work[n].map(sstr)
        return pd.Series([""] * len(work), index=work.index)

    # ===================== 1) Top authors (Table + Figure)
    _cb("Building top-authors slides")

    work = df.copy().dropna(how="all")
    work["__AuthorsList"] = work.apply(_authors_from_row, axis=1)

    auth_rows = work.explode("__AuthorsList").rename(columns={"__AuthorsList": "Author"})
    auth_rows["Author"] = auth_rows["Author"].map(sstr).map(str.strip)
    auth_rows = auth_rows[auth_rows["Author"] != ""]

    work["__AffText"] = _col(
        work,
        "affiliation_institution",
        "institution",
        "affiliations",
        "affiliation",
        "publisher",
        "publication_outlet",
        "publicationTitle",
    )
    work["__AffCountry"] = _col(work, "affiliation_country", "affiliation:country", "affiliation_country_value", "country_affiliation")
    work["__AffDept"] = _col(work, "affiliation_department", "affiliation:department", "department")

    auth_rows = auth_rows.join(work[["__AffText", "__AffCountry", "__AffDept"]], how="left")

    author_counts = auth_rows.groupby("Author").size().rename("Publications").reset_index()
    modal_aff = auth_rows.groupby("Author").agg(
        {"__AffText": _modal, "__AffCountry": _modal, "__AffDept": _modal}
    ).reset_index().rename(
        columns={
            "__AffText": "Dominant Affiliation",
            "__AffCountry": "Affiliation:Country",
            "__AffDept": "Affiliation:Department",
        }
    )

    toplist = (
        author_counts.merge(modal_aff, on="Author", how="left")
        .sort_values(["Publications", "Author"], ascending=[False, True])
        .head(int(top_n_authors))
        .loc[:, ["Author", "Dominant Affiliation", "Affiliation:Country", "Affiliation:Department", "Publications"]]
    )

    _add_table_slide(
        "(B) Top authors — Table",
        toplist,
        notes=(
            "Method: Authors parsed from structured 'authors_list' when present; otherwise from 'authors'/'creator_summary' "
            "(split on commas/semicolons). Publications counted per record after exploding authors. "
            "Dominant affiliation/country/department are modal strings across that author’s records."
        ),
    )

    fig_auth = px.bar(
        toplist.sort_values("Publications", ascending=True),
        x="Publications",
        y="Author",
        orientation="h",
        text="Publications",
        title=f"{collection_name} · Top authors (by publications)",
    )
    fig_auth.update_traces(texttemplate="%{text}")
    fig_auth.update_layout(
        template="plotly_white",
        width=PLOT_W,
        height=PLOT_H,
        margin=dict(l=68, r=120, t=80, b=80),
        title=dict(x=0.02, xanchor="left", font=dict(size=20, family=FONT_FAMILY)),
        font=dict(family=FONT_FAMILY, size=13),
        legend_title_text="",
    )

    _add_figure_slide(
        "(B) Top authors — Figure",
        fig_auth,
        notes=(
            "Figure: Horizontal bar chart of Top authors by publication count. "
            "Source parsing rules match the Top-authors table slide."
        ),
    )

    # ===================== 2) Co-authorship (Table + Figure)
    _cb("Building co-authorship slides")

    edge_counts = {}
    for _, r in work.iterrows():
        A = _authors_from_row(r)
        A = [a for a in A if a]
        for a, b in itertools.combinations(sorted(set(A)), 2):
            edge_counts[(a, b)] = edge_counts.get((a, b), 0) + 1

    edges = pd.DataFrame(
        [(a, b, w) for (a, b), w in edge_counts.items()],
        columns=["Author A", "Author B", "Shared Publications"],
    ).sort_values("Shared Publications", ascending=False)

    edge_top = edges.head(20).reset_index(drop=True)

    _add_table_slide(
        "(B) Co-authorship — Strongest pairs (Top 20) — Table",
        edge_top,
        notes="Method: For each record, all unique author pairs are generated; counts are shared-publication frequencies.",
    )

    keep = set(edge_top["Author A"]).union(edge_top["Author B"])
    e2 = edges[(edges["Author A"].isin(keep)) & (edges["Author B"].isin(keep))].copy()
    nodes = sorted(set(e2["Author A"]).union(e2["Author B"]))
    n = len(nodes)
    theta = np.linspace(0, 2 * math.pi, n, endpoint=False)
    pos = {nodes[i]: (math.cos(theta[i]), math.sin(theta[i])) for i in range(n)}

    x_edges, y_edges = [], []
    for _, rr in e2.iterrows():
        x0, y0 = pos[rr["Author A"]]
        x1, y1 = pos[rr["Author B"]]
        x_edges += [x0, x1, None]
        y_edges += [y0, y1, None]

    deg = {k: 0 for k in nodes}
    for _, rr in e2.iterrows():
        deg[rr["Author A"]] += 1
        deg[rr["Author B"]] += 1
    sizes = [8 + 6 * math.sqrt(max(1, deg[k])) for k in nodes]

    fig_net = go.Figure()
    fig_net.add_trace(
        go.Scatter(
            x=x_edges,
            y=y_edges,
            mode="lines",
            line=dict(width=1, color="rgba(0,0,0,0.25)"),
            hoverinfo="skip",
            showlegend=False,
        )
    )
    fig_net.add_trace(
        go.Scatter(
            x=[pos[k][0] for k in nodes],
            y=[pos[k][1] for k in nodes],
            mode="markers+text",
            text=nodes,
            textposition="top center",
            marker=dict(size=sizes, line=dict(width=0.5, color="rgba(0,0,0,0.3)")),
            hovertemplate="%{text}<extra></extra>",
            showlegend=False,
        )
    )
    fig_net.update_xaxes(showgrid=False, zeroline=False, visible=False)
    fig_net.update_yaxes(showgrid=False, zeroline=False, visible=False, scaleanchor="x", scaleratio=1)
    fig_net.update_layout(
        template="plotly_white",
        width=PLOT_W,
        height=PLOT_H,
        margin=dict(l=68, r=120, t=80, b=80),
        title=dict(
            text=f"{collection_name} · Co-authorship network (strongest pairs)",
            x=0.02,
            xanchor="left",
            font=dict(size=20, family=FONT_FAMILY),
        ),
        font=dict(family=FONT_FAMILY, size=13),
        legend_title_text="",
    )

    _add_figure_slide(
        "(B) Co-authorship network — Figure",
        fig_net,
        notes=(
            "Network: Nodes are authors; edges indicate co-appearance on records. "
            "Graph includes authors appearing in the Top-20 strongest pairs. "
            "Layout is circular; node size scales with degree within this subgraph."
        ),
    )

    # ===================== 3) Major institutional actors (publication_outlet) — Table + Figure
    _cb("Building institutional-actors slides")

    outlet_col = None
    try:
        want = {
            "publication_outlet",
            "publicationtitle",
            "publication_title",
            "publicationtitle",
            "publication",
            "journal",
            "journal_title",
            "publisher",
        }
        for c in df.columns:
            if str(c).strip().lower() in want:
                outlet_col = c
                break
        if outlet_col is None:
            # tolerate camelCase variants and minor spelling differences
            for c in df.columns:
                lc = str(c).strip().lower()
                if "publication" in lc and ("title" in lc or "outlet" in lc):
                    outlet_col = c
                    break
    except Exception:
        outlet_col = None

    if "publisher_type_value" in df.columns:
        type_col = "publisher_type_value"
    else:
        df["_tmp_pubtype_"] = "Unknown/Other"
        type_col = "_tmp_pubtype_"

    def _split_outlets(s: str) -> list[str]:
        parts = re.split(r"[;|,/]\s*", str(s or ""))
        return [" ".join(p.strip().split()) for p in parts if p and p.strip()]

    if outlet_col is None:
        _cb("No publication outlet column found; skipping institutional-actors slides")
        if type_col == "_tmp_pubtype_" and "_tmp_pubtype_" in df.columns:
            try:
                df.drop(columns=["_tmp_pubtype_"], inplace=True)
            except Exception:
                pass
    else:
        tmp = df[[outlet_col, type_col]].copy().rename(
            columns={outlet_col: "publication_outlet", type_col: "publisher_type_value"}
        )

        tmp["publication_outlet"] = tmp["publication_outlet"].astype(str).map(lambda x: " ".join(x.split()))
        tmp["publisher_type_value"] = tmp["publisher_type_value"].astype(str).replace(
            {"": "Unknown/Other", "nan": "Unknown/Other"}
        )

        tmp["publication_outlet_list"] = tmp["publication_outlet"].apply(_split_outlets)
        tmp = tmp.explode("publication_outlet_list", ignore_index=True)
        tmp["publication_outlet"] = tmp["publication_outlet_list"].fillna("").astype(str).str.strip()
        tmp = tmp.drop(columns=["publication_outlet_list"])
        tmp = tmp[tmp["publication_outlet"] != ""].copy()

        counts = tmp.groupby("publication_outlet").size().rename("Count")
        modal_type = (
            tmp.groupby("publication_outlet")["publisher_type_value"]
            .agg(lambda s: s.value_counts().index[0])
            .rename("Publisher Type")
        )
        agg = (
            pd.concat([counts, modal_type], axis=1)
            .sort_values("Count", ascending=False)
            .head(20)
            .reset_index()
        )

        table_df = agg.rename(columns={"publication_outlet": "Publication Outlet"})[
            ["Publication Outlet", "Publisher Type", "Count"]
        ]
        _add_table_slide(
            "(B) Major institutional actors (publication_outlet) — Table",
            table_df,
            notes=(
                "Method: publication_outlet values split on ';', ',', '/', or '|' and whitespace-normalised. "
                "Counts are outlet mentions (not unique records). Publisher Type is modal publisher_type_value per outlet."
            ),
        )

    fig_inst = px.bar(
        agg.sort_values("Count", ascending=False),
        x="publication_outlet",
        y="Count",
        color="Publisher Type",
        text="Count",
        title=f"{collection_name} · Major institutional actors (publication_outlet)",
    )
    fig_inst.update_traces(textposition="outside", marker_line_width=0.5, marker_line_color="rgba(0,0,0,0.15)")
    fig_inst.update_xaxes(tickangle=-25, automargin=True)
    fig_inst.update_layout(
        template="plotly_white",
        width=PLOT_W,
        height=PLOT_H,
        margin=dict(l=68, r=120, t=80, b=80),
        title=dict(x=0.02, xanchor="left", font=dict(size=20, family=FONT_FAMILY)),
        font=dict(family=FONT_FAMILY, size=13),
        legend_title_text="",
    )

    n_total = int(tmp.shape[0])
    n_used = int(agg["Count"].sum())

    _add_figure_slide(
        "(B) Major institutional actors (publication_outlet) — Figure",
        fig_inst,
        notes=(
            "Rationale: Bars show Top-20 raw values from publication_outlet (no canonical mapping). "
            "Bar colour is modal publisher_type_value for that outlet; blanks labelled 'Unknown/Other'. "
            f"Count basis: {n_used} outlet mentions within {n_total} parsed outlet entries."
        ),
    )

    _cb("Authorship & institutions completed")

    if preview_slides is not None:
        return {"slides": preview_slides, "figs": preview_figs or [], "index": 0}

    return None


def add_citations_and_influence_section(
    prs,
    df: "pd.DataFrame",
    collection_name: str = "Collection",
    *,
    slide_notes: bool = True,
    export: bool = False,
    return_payload: bool = False,
    citations_col: str | None = None,
    top_n: int = 10,
    progress_callback=None,
):
    """
    Produce a citations/influence mini-deck as interactive Plotly payload slides.

    Slides returned (figure slides only; tables are attached as table_html for the Table tab):
      1) Citation distribution (linear, capped at P99) + stats table
      2) Citation distribution (log10(citations+1)) + stats table
      3) Top-N most cited (horizontal bars) + top table
    """
    import json
    import numpy as np
    import pandas as pd
    import plotly.express as px
    import plotly.graph_objects as go
    from plotly.utils import PlotlyJSONEncoder

    def _cb(msg: str):
        if progress_callback:
            try:
                progress_callback(f"CIT: {msg}")
            except Exception:
                pass

    def _fig_json(fig_obj: "go.Figure") -> dict:
        return json.loads(json.dumps(fig_obj.to_plotly_json(), cls=PlotlyJSONEncoder))

    def _table_html(frame: "pd.DataFrame") -> str:
        try:
            if frame is None or frame.empty:
                return "<div>No data available.</div>"
            return frame.to_html(index=False, escape=True)
        except Exception:
            return "<div>Unable to render table.</div>"

    preview_slides = [] if (export or return_payload) else None

    if df is None or not isinstance(df, pd.DataFrame) or df.empty:
        fig = go.Figure()
        fig.add_annotation(text="No data.", showarrow=False, xref="paper", yref="paper")
        if preview_slides is not None:
            preview_slides.append(
                {
                    "title": f"{collection_name} — Citations (no data)",
                    "fig_json": _fig_json(fig),
                    "table_html": "<div>No data.</div>",
                    "notes": "Source: add_citations_and_influence_section()." if slide_notes else "",
                }
            )
            return {"slides": preview_slides, "index": 0}
        return None

    # ---- detect citations column
    cand_cols = [citations_col] if citations_col else []
    cand_cols += ["citations", "citation", "times_cited", "timescited", "n_cit", "num_cit", "cites"]
    cit_col = None
    for c in cand_cols:
        if c and c in df.columns:
            cit_col = c
            break
    if not cit_col:
        # heuristic: first column containing "cit" that is mostly numeric
        for c in df.columns:
            if "cit" not in str(c).lower():
                continue
            s = pd.to_numeric(df[c], errors="coerce")
            if s.notna().mean() >= 0.5:
                cit_col = c
                break

    if not cit_col:
        fig = go.Figure()
        fig.add_annotation(text="No citations column found.", showarrow=False, xref="paper", yref="paper")
        if preview_slides is not None:
            preview_slides.append(
                {
                    "title": f"{collection_name} — Citations (missing column)",
                    "fig_json": _fig_json(fig),
                    "table_html": "<div>No citations column found.</div>",
                    "notes": "Source: add_citations_and_influence_section()." if slide_notes else "",
                }
            )
            return {"slides": preview_slides, "index": 0}
        return None

    s = pd.to_numeric(df[cit_col], errors="coerce")
    n_total = int(len(s))
    n_missing = int(s.isna().sum())
    s = s.fillna(0.0)
    s = s.clip(lower=0)

    _cb(f"citations_col={cit_col} n={n_total}")

    def _h_index(vals: "pd.Series") -> int:
        try:
            v = np.sort(vals.astype(float).to_numpy())[::-1]
            h = 0
            for i, x in enumerate(v, start=1):
                if x >= i:
                    h = i
                else:
                    break
            return int(h)
        except Exception:
            return 0

    stats = {
        "N": n_total,
        "Missing": n_missing,
        "Zero": int((s <= 0).sum()),
        "Sum": float(s.sum()),
        "Mean": float(s.mean()),
        "Median": float(s.median()),
        "P75": float(s.quantile(0.75)),
        "P90": float(s.quantile(0.90)),
        "P95": float(s.quantile(0.95)),
        "P99": float(s.quantile(0.99)),
        "H-index": _h_index(s),
    }
    stats_df = pd.DataFrame([{"Metric": k, "Value": v} for k, v in stats.items()])
    stats_html = _table_html(stats_df)

    # ---- (1) linear distribution capped at P99
    p99 = float(stats["P99"]) if np.isfinite(stats["P99"]) else float(s.max())
    s_lin = s.clip(upper=max(1.0, p99))
    fig_lin = px.histogram(
        x=s_lin,
        nbins=min(40, max(10, int(np.sqrt(max(1, n_total))))),
        title=f"{collection_name} · Citation distribution (capped at P99={p99:.0f})",
        labels={"x": f"{cit_col} (capped)", "y": "Count"},
    )
    fig_lin.update_layout(margin=dict(l=60, r=60, t=70, b=60))

    # ---- (2) log distribution
    fig_log = px.histogram(
        x=np.log10(s + 1.0),
        nbins=min(40, max(10, int(np.sqrt(max(1, n_total))))),
        title=f"{collection_name} · Citation distribution (log10({cit_col}+1))",
        labels={"x": f"log10({cit_col}+1)", "y": "Count"},
    )
    fig_log.update_layout(margin=dict(l=60, r=60, t=70, b=60))

    # ---- (3) top cited
    top = df.copy()
    top["_cit"] = pd.to_numeric(top[cit_col], errors="coerce").fillna(0.0).clip(lower=0)
    top = top.sort_values("_cit", ascending=False).head(int(top_n) if int(top_n) > 0 else 10)
    if "title" in top.columns:
        top["title"] = top["title"].astype(str)
    top_tbl_cols = [c for c in ["title", "authors", "year", cit_col] if c in top.columns]
    top_table = top[top_tbl_cols].copy() if top_tbl_cols else top[["_cit"]].copy()
    if cit_col not in top_table.columns and "_cit" in top_table.columns:
        top_table = top_table.rename(columns={"_cit": cit_col})
    top_html = _table_html(top_table)

    label = None
    if "title" in top.columns:
        label = top["title"].astype(str).str.slice(0, 80)
    elif "key" in top.columns:
        label = top["key"].astype(str)
    else:
        label = pd.Series([f"Item {i+1}" for i in range(len(top))])
    fig_top = px.bar(
        top.assign(_label=label).sort_values("_cit", ascending=True),
        x="_cit",
        y="_label",
        orientation="h",
        text="_cit",
        title=f"{collection_name} · Top {len(top)} by citations",
        labels={"_cit": cit_col, "_label": "Work"},
    )
    fig_top.update_traces(textposition="outside", marker_line_width=0.5, marker_line_color="rgba(0,0,0,0.15)")
    xmax = float(np.nanmax(pd.to_numeric(top["_cit"], errors="coerce"))) if len(top) else 0.0
    if np.isfinite(xmax) and xmax > 0:
        fig_top.update_xaxes(range=[0, xmax * 1.18])
    fig_top.update_layout(margin=dict(l=180, r=140, t=70, b=60))

    if preview_slides is not None:
        preview_slides.extend(
            [
                {
                    "title": "Citations — Distribution (linear)",
                    "fig_json": _fig_json(fig_lin),
                    "table_html": stats_html,
                    "notes": f"Column: {cit_col}. Missing treated as 0. P99 cap={p99:.0f}." if slide_notes else "",
                },
                {
                    "title": "Citations — Distribution (log)",
                    "fig_json": _fig_json(fig_log),
                    "table_html": stats_html,
                    "notes": f"Column: {cit_col}. Log uses log10(citations+1)." if slide_notes else "",
                },
                {
                    "title": "Citations — Top cited",
                    "fig_json": _fig_json(fig_top),
                    "table_html": top_html,
                    "notes": f"Top {len(top)} by {cit_col}." if slide_notes else "",
                },
            ]
        )
        return {"slides": preview_slides, "index": 0}

    return None


def add_geo_and_sector_section(
    prs,
    df: "pd.DataFrame",
    collection_name: str = "Collection",
    *,
    slide_notes: bool = True,
    export: bool = False,
    return_payload: bool = False,
    country_tag: str = "country_focus_value",
    sector_tag: str = "sector_focus_value",
    top_countries: int = 30,
    progress_callback=None,
):
    """
    Geographic & sectoral coverage (payload form).

    Slides returned:
      1) Country focus (Top countries) — Figure + table
      2) Sector coverage (%) — Figure + table
      3) Government missions (optional) — Figure + table (if a mission column exists)
    """
    import json
    import re
    import numpy as np
    import pandas as pd
    import plotly.express as px
    import plotly.graph_objects as go
    from plotly.utils import PlotlyJSONEncoder

    def _cb(msg: str):
        if progress_callback:
            try:
                progress_callback(f"GEO: {msg}")
            except Exception:
                pass

    def _fig_json(fig_obj: "go.Figure") -> dict:
        return json.loads(json.dumps(fig_obj.to_plotly_json(), cls=PlotlyJSONEncoder))

    def _explode_tags(series: "pd.Series") -> "pd.Series":
        s = series.fillna("").astype(str)
        s = s.str.replace(r"\s*[,;|/]\s*", "§", regex=True)
        out = s.str.split("§").explode()
        out = out.astype(str).str.strip()
        return out[out.astype(bool)]

    def _table_html(frame: "pd.DataFrame") -> str:
        try:
            if frame is None or frame.empty:
                return "<div>No data available.</div>"
            return frame.to_html(index=False, escape=True)
        except Exception:
            return "<div>Unable to render table.</div>"

    preview_slides = [] if (export or return_payload) else None
    if df is None or not hasattr(df, "columns") or getattr(df, "empty", True):
        fig = go.Figure()
        fig.add_annotation(text="No data.", showarrow=False, xref="paper", yref="paper")
        if preview_slides is not None:
            preview_slides.append({"title": "Geo & sector (no data)", "fig_json": _fig_json(fig), "table_html": "<div>No data.</div>", "notes": ""})
            return {"slides": preview_slides, "index": 0}
        return None

    # ---- Countries
    if country_tag in df.columns:
        cexp = _explode_tags(df[country_tag])
        ccounts = cexp.value_counts().rename_axis("Country").reset_index(name="Count")
        ccounts = ccounts.sort_values("Count", ascending=False).reset_index(drop=True)
        topN = int(top_countries) if int(top_countries) > 0 else 30
        ctbl = ccounts.head(topN).copy()
        fig_c = px.bar(
            ctbl.sort_values("Count", ascending=True),
            x="Count",
            y="Country",
            orientation="h",
            text="Count",
            title=f"{collection_name} · Country focus (Top {len(ctbl)})",
        )
        fig_c.update_traces(textposition="outside", marker_line_width=0.5, marker_line_color="rgba(0,0,0,0.15)")
        xmax = float(np.nanmax(pd.to_numeric(ctbl["Count"], errors="coerce"))) if len(ctbl) else 0.0
        if np.isfinite(xmax) and xmax > 0:
            fig_c.update_xaxes(range=[0, xmax * 1.18])
        fig_c.update_layout(margin=dict(l=160, r=140, t=70, b=60))

        if preview_slides is not None:
            preview_slides.append(
                {
                    "title": "Geo — Country focus",
                    "fig_json": _fig_json(fig_c),
                    "table_html": _table_html(ctbl),
                    "notes": f"Tag column: {country_tag}." if slide_notes else "",
                }
            )

    # ---- Sectors (%)
    if sector_tag in df.columns:
        sexp = _explode_tags(df[sector_tag])
        if not sexp.empty:
            sc = (
                sexp.value_counts(normalize=True)
                .mul(100.0)
                .round(1)
                .rename_axis("Sector")
                .reset_index(name="%")
            )
            sc["Sector"] = sc["Sector"].astype(str).map(lambda x: re.sub(r"\s+", " ", x).strip())
            sc_plot = sc.sort_values("%", ascending=True)
            sc_plot["Label"] = sc_plot["%"].map(lambda v: f"{float(v):.1f}%")
            fig_s = px.bar(
                sc_plot,
                x="%",
                y="Sector",
                orientation="h",
                text="Label",
                title=f"{collection_name} · Sector coverage (%)",
            )
            fig_s.update_traces(textposition="outside", marker_line_width=0.5, marker_line_color="rgba(0,0,0,0.15)")
            xmax = float(np.nanmax(pd.to_numeric(sc_plot["%"], errors="coerce"))) if len(sc_plot) else 0.0
            if np.isfinite(xmax) and xmax > 0:
                fig_s.update_xaxes(range=[0, xmax * 1.22])
            fig_s.update_layout(margin=dict(l=160, r=170, t=70, b=60))

            if preview_slides is not None:
                preview_slides.append(
                    {
                        "title": "Sector — Coverage (%)",
                        "fig_json": _fig_json(fig_s),
                        "table_html": _table_html(sc),
                        "notes": f"Tag column: {sector_tag}." if slide_notes else "",
                    }
                )

    # ---- Missions (optional)
    mission_col = next((c for c in ("mission", "mission_title", "government_mission", "mission_name") if c in df.columns), None)
    if mission_col:
        mexp = _explode_tags(df[mission_col])
        if not mexp.empty and preview_slides is not None:
            miss = mexp.value_counts().rename_axis("Mission").reset_index(name="Count")
            miss_plot = miss.head(25).sort_values("Count", ascending=True)
            fig_m = px.bar(
                miss_plot,
                x="Count",
                y="Mission",
                orientation="h",
                text="Count",
                title=f"{collection_name} · Government missions (Top {len(miss_plot)})",
            )
            fig_m.update_traces(textposition="outside", marker_line_width=0.5, marker_line_color="rgba(0,0,0,0.15)")
            xmax = float(np.nanmax(pd.to_numeric(miss_plot["Count"], errors="coerce"))) if len(miss_plot) else 0.0
            if np.isfinite(xmax) and xmax > 0:
                fig_m.update_xaxes(range=[0, xmax * 1.18])
            fig_m.update_layout(margin=dict(l=200, r=140, t=70, b=60))

            preview_slides.append(
                {
                    "title": "Missions — Coverage",
                    "fig_json": _fig_json(fig_m),
                    "table_html": _table_html(miss),
                    "notes": f"Column: {mission_col}." if slide_notes else "",
                }
            )

    if preview_slides is not None:
        if not preview_slides:
            fig = go.Figure()
            fig.add_annotation(text="No geo/sector content produced.", showarrow=False, xref="paper", yref="paper")
            preview_slides.append({"title": "Geo & sector", "fig_json": _fig_json(fig), "table_html": "<div>No data available.</div>", "notes": ""})
        return {"slides": preview_slides, "index": 0}
    return None


def add_thematic_and_method_section(
    prs,
    df: "pd.DataFrame",
    collection_name: str = "Collection",
    *,
    slide_notes: bool = True,
    export: bool = False,
    return_payload: bool = False,
    cross_tabs: list[list[str]] | None = None,
    phase_tag: str = "phase_focus_value",
    year_col: str = "publication_year",
    topn_cols: int = 14,
    progress_callback=None,
):
    """
    Thematic & methodological mapping (payload form).

    Default cross-tabs follow create_power_point.py:
      - focus_type_value × publisher_type
      - empirical_theoretical × sector_focus_value
      - mission phase proportions and time trend
    """
    import json
    import numpy as np
    import pandas as pd
    import plotly.express as px
    import plotly.graph_objects as go
    from plotly.utils import PlotlyJSONEncoder

    def _cb(msg: str):
        if progress_callback:
            try:
                progress_callback(f"THEM: {msg}")
            except Exception:
                pass

    def _fig_json(fig_obj: "go.Figure") -> dict:
        return json.loads(json.dumps(fig_obj.to_plotly_json(), cls=PlotlyJSONEncoder))

    def _explode(series: "pd.Series") -> "pd.Series":
        s = series.fillna("").astype(str)
        s = s.str.replace(r"\s*[,;|/]\s*", "§", regex=True)
        out = s.str.split("§").explode()
        out = out.astype(str).str.strip()
        return out[out.astype(bool)]

    def _table_html(frame: "pd.DataFrame") -> str:
        try:
            if frame is None or frame.empty:
                return "<div>No data available.</div>"
            return frame.to_html(index=False, escape=True)
        except Exception:
            return "<div>Unable to render table.</div>"

    def _cap_columns(ct: "pd.DataFrame", n: int) -> "pd.DataFrame":
        if ct is None or ct.empty or ct.shape[1] <= n:
            return ct
        totals = ct.sum(axis=0).sort_values(ascending=False)
        keep = totals.head(n).index.tolist()
        other = ct.drop(columns=keep).sum(axis=1)
        out = ct[keep].copy()
        out["Other"] = other
        return out

    preview_slides = [] if (export or return_payload) else None
    if df is None or not isinstance(df, pd.DataFrame) or df.empty:
        fig = go.Figure()
        fig.add_annotation(text="No data.", showarrow=False, xref="paper", yref="paper")
        if preview_slides is not None:
            preview_slides.append({"title": "Thematic & method (no data)", "fig_json": _fig_json(fig), "table_html": "<div>No data.</div>", "notes": ""})
            return {"slides": preview_slides, "index": 0}
        return None

    if not cross_tabs:
        # Prefer legacy pairs when present, otherwise pick reasonable fallbacks based on available columns.
        legacy = [["focus_type_value", "publisher_type"], ["empirical_theoretical", "sector_focus_value"]]
        present: list[list[str]] = []
        for a, b in legacy:
            if a in df.columns and b in df.columns:
                present.append([a, b])
        if not present:
            fallbacks = [
                ["theme", "item_type"],
                ["theme", "country"],
                ["item_type", "country"],
                ["item_type", "source"],
                ["theme", "source"],
            ]
            for a, b in fallbacks:
                if a in df.columns and b in df.columns:
                    present.append([a, b])
                    break
        cross_tabs = present

    # ---- Cross-tabs
    for (row_col, col_col) in cross_tabs:
        if row_col not in df.columns or col_col not in df.columns:
            continue
        rows = _explode(df[row_col])
        cols = _explode(df[col_col])
        # align by index by building temp DF
        tmp = pd.DataFrame({"__i": np.arange(len(df))})
        tmp[row_col] = df[row_col]
        tmp[col_col] = df[col_col]
        tmp[row_col] = tmp[row_col].fillna("").astype(str).str.replace(r"\s*[,;|/]\s*", "§", regex=True).str.split("§")
        tmp[col_col] = tmp[col_col].fillna("").astype(str).str.replace(r"\s*[,;|/]\s*", "§", regex=True).str.split("§")
        tmp = tmp.explode(row_col).explode(col_col)
        tmp[row_col] = tmp[row_col].astype(str).str.strip()
        tmp[col_col] = tmp[col_col].astype(str).str.strip()
        tmp = tmp[(tmp[row_col].astype(bool)) & (tmp[col_col].astype(bool))]
        if tmp.empty:
            continue

        ct = pd.crosstab(tmp[row_col], tmp[col_col])
        ct = _cap_columns(ct, int(topn_cols) if int(topn_cols) > 0 else 14)
        table = ct.reset_index()
        long = ct.stack().reset_index(name="Count")
        long.columns = [row_col, col_col, "Count"]
        fig = px.bar(long, x=row_col, y="Count", color=col_col, barmode="group", title=f"{collection_name} · {row_col} × {col_col}")
        fig.update_layout(margin=dict(l=60, r=60, t=70, b=140))
        fig.update_xaxes(tickangle=-25, automargin=True)

        if preview_slides is not None:
            preview_slides.append(
                {
                    "title": f"Thematic — {row_col} × {col_col}",
                    "fig_json": _fig_json(fig),
                    "table_html": _table_html(table),
                    "notes": f"Cross-tab of {row_col} vs {col_col}." if slide_notes else "",
                }
            )

    # ---- Phase proportions
    if phase_tag in df.columns:
        phase = _explode(df[phase_tag])
        if not phase.empty and preview_slides is not None:
            ph = phase.value_counts().rename_axis("Phase").reset_index(name="Count")
            figp = px.bar(ph, x="Phase", y="Count", text="Count", title=f"{collection_name} · Mission phase (overall)")
            figp.update_traces(textposition="outside", marker_line_width=0.5, marker_line_color="rgba(0,0,0,0.15)")
            figp.update_layout(margin=dict(l=60, r=160, t=70, b=140))
            figp.update_xaxes(tickangle=-25, automargin=True)
            preview_slides.append(
                {
                    "title": "Phase — Overall",
                    "fig_json": _fig_json(figp),
                    "table_html": _table_html(ph),
                    "notes": f"Tag column: {phase_tag}." if slide_notes else "",
                }
            )

    # ---- Phase trend by year
    if year_col in df.columns and phase_tag in df.columns and preview_slides is not None:
        years = pd.to_numeric(df[year_col], errors="coerce")
        tmp = pd.DataFrame({"Year": years, "Phase": df[phase_tag]})
        tmp = tmp.dropna(subset=["Year"])
        tmp["Year"] = tmp["Year"].astype(int)
        tmp["Phase"] = tmp["Phase"].fillna("").astype(str).str.replace(r"\s*[,;|/]\s*", "§", regex=True).str.split("§")
        tmp = tmp.explode("Phase")
        tmp["Phase"] = tmp["Phase"].astype(str).str.strip()
        tmp = tmp[tmp["Phase"].astype(bool)]
        if not tmp.empty:
            grp = tmp.groupby(["Year", "Phase"]).size().rename("Count").reset_index()
            figt = px.line(grp.sort_values(["Year", "Phase"]), x="Year", y="Count", color="Phase", markers=True, title=f"{collection_name} · Mission phase trend")
            figt.update_layout(margin=dict(l=60, r=60, t=70, b=60))
            preview_slides.append(
                {
                    "title": "Phase — Trend",
                    "fig_json": _fig_json(figt),
                    "table_html": _table_html(grp.head(80)),
                    "notes": f"Year: {year_col}; Phase: {phase_tag}." if slide_notes else "",
                }
            )

    if preview_slides is not None:
        if not preview_slides:
            fig = go.Figure()
            fig.add_annotation(text="No thematic/method content produced.", showarrow=False, xref="paper", yref="paper")
            preview_slides.append({"title": "Thematic & method", "fig_json": _fig_json(fig), "table_html": "<div>No data available.</div>", "notes": ""})
        return {"slides": preview_slides, "index": 0}

    return None


def analyze_most_frequent_words(df: pd.DataFrame, params: dict, progress_callback=None,
                                zotero_client_for_pdf=None,
                                collection_name_for_cache=None):  # This is passed by MainWindow
    def _callback(msg):
        if progress_callback: progress_callback(f"FreqWords: {msg}")
        logging.info(f"FreqWords: {msg}")

    data_source_param = params.get("data_source", "controlled_vocabulary_terms")
    top_n = params.get("top_n_words", 20)
    min_freq = params.get("min_frequency", 1)
    plot_type = params.get("plot_type", "bar_vertical")

    # Use the directly passed collection_name_for_cache.
    # Fallback only if it's explicitly None (which shouldn't happen if MainWindow passes it).
    effective_collection_name = collection_name_for_cache
    if not effective_collection_name:
        effective_collection_name = f"default_collection_for_freq_words_{data_source_param}"
        _callback(f"Warning: collection_name_for_cache not provided, using default: {effective_collection_name}")

    _callback(
        f"Analyzing most frequent words from '{data_source_param}' (Top {top_n}, Min Freq: {min_freq}) using cache key context: '{effective_collection_name}'")

    all_words = get_text_corpus_from_df(
        df,
        source_type=data_source_param,
        collection_name_for_cache=effective_collection_name,  # Pass the consistent name
        zotero_client=zotero_client_for_pdf,
        progress_callback=progress_callback
    )

    if not all_words:
        _callback("No words found after preprocessing for most frequent words.")
        return pd.DataFrame(columns=["Word", "Frequency"]), go.Figure().add_annotation(text="No words found.",
                                                                                       showarrow=False)

    import json

    def _coerce_hashable_token(x):
        if x is None:
            return None
        if isinstance(x, str):
            t = x.strip()
            return t if t else None
        if isinstance(x, dict):
            t = json.dumps(x, ensure_ascii=False, sort_keys=True)
            t = t.strip()
            return t if t else None
        if isinstance(x, (list, tuple, set)):
            parts = []
            for y in x:
                yy = _coerce_hashable_token(y)
                if yy:
                    parts.append(yy)
            t = " ".join(parts).strip()
            return t if t else None
        t = str(x).strip()
        return t if t else None

    all_words = [t for t in (_coerce_hashable_token(x) for x in all_words) if t]
    word_counts = Counter(all_words)
    filtered_word_counts = {w: c for w, c in word_counts.items() if c >= min_freq}
    if not filtered_word_counts:
        _callback(f"No words meet min frequency {min_freq} for most frequent words.")
        return pd.DataFrame(columns=["Word", "Frequency"]), go.Figure().add_annotation(
            text=f"No words meet min freq {min_freq}.", showarrow=False)

    results_df = pd.DataFrame(filtered_word_counts.items(), columns=["Word", "Frequency"]).sort_values(
        by=["Frequency", "Word"], ascending=[False, True]).reset_index(drop=True)

    plot_df = results_df.head(top_n)
    fig = None
    if not plot_df.empty:
        title = f"Top {len(plot_df)} Frequent Words from '{data_source_param.replace('_', ' ').title()}'"
        try:
            if plot_type == "bar_vertical":
                fig = px.bar(plot_df, x="Word", y="Frequency", text_auto=True, title=title,
                             color="Frequency", color_continuous_scale=px.colors.sequential.Teal)
                if fig: fig.update_traces(textposition='outside'); fig.update_xaxes(categoryorder="total descending")
            elif plot_type == "bar_horizontal":
                fig = px.bar(plot_df.sort_values("Frequency", ascending=True),
                             y="Word", x="Frequency", orientation='h', text_auto=True, title=title,
                             color="Frequency", color_continuous_scale=px.colors.sequential.Mint)
            else:
                _callback(f"Plot type '{plot_type}' for FreqWords not specific, defaulting to vertical bar.")
                fig = px.bar(plot_df, x="Word", y="Frequency", text_auto=True, title=title,
                             color="Frequency", color_continuous_scale=px.colors.sequential.Teal)
                if fig: fig.update_traces(textposition='outside'); fig.update_xaxes(categoryorder="total descending")
        except Exception as e:
            _callback(f"Error creating plot for frequent words: {e}")
            fig = go.Figure().add_annotation(text=f"Plot error: {e}", showarrow=False)
    else:
        _callback("No data to plot for frequent words.")
        fig = go.Figure().add_annotation(text="No data to plot.", showarrow=False)

    _callback(f"Finished analyzing most frequent words from '{data_source_param}'.")
    return results_df, fig



def analyze_word_cloud(df: pd.DataFrame, params: dict, progress_callback=None,
                       zotero_client_for_pdf=None,
                       collection_name_for_cache=None):
    def _callback(msg):
        if progress_callback: progress_callback(f"WordCloud: {msg}")
        logging.info(f"WordCloud: {msg}")

    data_source_param = params.get("data_source", "controlled_vocabulary_terms")
    max_words = params.get("max_words", 100)
    colormap = params.get("wordcloud_colormap", "viridis")
    contour_width = params.get("contour_width", 0)
    contour_color = params.get("contour_color", "steelblue")

    _callback(f"Generating Word Cloud from '{data_source_param}' (Max words: {max_words})...")
    corpus_tokens = get_text_corpus_from_df(
        df,
        source_type=data_source_param,
        collection_name_for_cache=collection_name_for_cache,
        zotero_client=zotero_client_for_pdf,
        progress_callback=None,
        cache_mode="read",  # never hit network here
        preserve_item_cache=True,
        return_mode="flat",  # request a flat token list
        include_summary=False,
    ) or []

    # If a caller accidentally switched return_mode, flatten (defensive)
    if corpus_tokens and isinstance(corpus_tokens, list) and isinstance(corpus_tokens[0], tuple):
        # received per_doc [(doc_id, [toks...]), ...] → flatten
        corpus_tokens = [t for _, toks in corpus_tokens for t in (toks or [])]

    # Normalise and filter
    tokens = [str(t).strip().lower() for t in corpus_tokens if isinstance(t, (str, bytes)) and str(t).strip()]

    if not tokens:
        msg = "Word cloud could not be generated: No tokens available from source."
        _callback(msg)
        return pd.DataFrame(columns=["Word", "Frequency"]), {"type": "message", "content": msg}

    word_frequencies = Counter(tokens)
    results_df = pd.DataFrame(word_frequencies.most_common(max_words), columns=["Word", "Frequency"])

    # ------------------ WORD CLOUD IMAGE GENERATION ------------------
    try:
        wc = WordCloud(
            width=800,
            height=550,
            background_color=None,
            mode="RGBA",
            max_words=max_words,
            stopwords=STOP_WORDS,
            colormap=colormap,
            min_font_size=10,
            prefer_horizontal=0.9,
            contour_width=contour_width if contour_width > 0 else 0,
            contour_color=contour_color,
        )
        wc.generate_from_frequencies(dict(word_frequencies))
        buf = io.BytesIO()
        wc.to_image().save(buf, format="PNG")
        buf.seek(0)
        data_uri = "data:image/png;base64," + base64.b64encode(buf.read()).decode("utf-8")
        _callback("Word cloud image generated and encoded to Data URI successfully.")
    except Exception as e:
        msg = f"Error generating image: {e}"
        _callback(msg)
        logging.error("Word cloud image generation error", exc_info=True)
        return results_df, {"type": "message", "content": msg}

    # ------------------ PLOTLY FALLBACK ------------------
    render_style = params.get("wordcloud_render_style", "cloud")  # "cloud" | "bar"
    if render_style == "bar":
        try:
            fig = go.Figure()
            top_words = results_df.head(max_words)
            fig.add_trace(go.Bar(
                x=top_words["Word"],
                y=top_words["Frequency"],
            ))
            fig.update_layout(title="Word Frequencies",
                              xaxis_title="Word", yaxis_title="Frequency")
            fig_result = fig
        except Exception as e_plot:
            _callback(f"Warning: Plotly bar-chart failed: {e_plot}")
            fig_result = {"type": "message",
                          "content": f"Could not build bar chart: {e_plot}"}
    else:  # default → real word-cloud PNG
        fig_result = data_uri

    _callback("Word Cloud analysis complete.")
    return results_df, fig_result

def analyze_word_treemap(df: pd.DataFrame, params: dict, progress_callback=None,
                         zotero_client_for_pdf=None,
                         collection_name_for_cache=None):  # Received from MainWindow
    def _callback(msg):
        if progress_callback: progress_callback(f"WordTreemap: {msg}")
        logging.info(f"WordTreemap: {msg}")

    data_source_param = params.get("data_source", "controlled_vocabulary_terms")
    top_n = params.get("top_n_words", 30)
    min_freq = params.get("min_frequency", 1)

    effective_collection_name = collection_name_for_cache
    if not effective_collection_name:
        effective_collection_name = f"default_collection_for_treemap_{data_source_param}"
        _callback(f"Warning: collection_name_for_cache not provided, using default: {effective_collection_name}")

    _callback(
        f"Generating Word Treemap from '{data_source_param}' (Top {top_n}) using cache key context: '{effective_collection_name}'")

    all_words = get_text_corpus_from_df(
        df, source_type=data_source_param,
        collection_name_for_cache=effective_collection_name,
        zotero_client=zotero_client_for_pdf,
        progress_callback=progress_callback
    )
    if not all_words: _callback("No words for treemap."); return pd.DataFrame(columns=["Word", "Frequency"]), None
    counts = Counter(all_words);
    filtered = {w: c for w, c in counts.items() if c >= min_freq}
    if not filtered: _callback(f"No words meet min freq {min_freq} for treemap."); return pd.DataFrame(
        columns=["Word", "Frequency"]), None
    results_df = pd.DataFrame(filtered.items(), columns=["Word", "Frequency"]).sort_values("Frequency",
                                                                                           ascending=False).reset_index(
        drop=True)
    plot_df = results_df.head(top_n);
    fig = None
    if not plot_df.empty:
        try:
            title = f"Word Treemap from '{data_source_param.replace('_', ' ').title()}' (Top {len(plot_df)})"
            fig = px.treemap(plot_df, path=[px.Constant(title), 'Word'], values='Frequency', color='Frequency',
                             title=title, color_continuous_scale='GnBu')
            if fig.data:
                fig.update_traces(textinfo="label+value+percent parent")
            fig.update_layout(margin=dict(t=50, l=10, r=10, b=10))
        except Exception as e:
            _callback(f"Treemap error: {e}"); fig = go.Figure().add_annotation(text=f"Treemap error: {e}",
                                                                               showarrow=False)
    else:
        _callback("No data for treemap plot."); fig = go.Figure().add_annotation(text="No data to plot.",
                                                                                 showarrow=False)
    return results_df, fig







def analyze_citations_overview_plotly(
    df: "pd.DataFrame",
    params: dict,
    progress_callback=None,
    *,
    export: bool = False,
    return_payload: bool = False,
):
    """
    ###1. compute the Top-N cited table (always)
    ###2. build one plotly figure for the requested plot_type
    ###3. if export or return_payload: return {"slides":[...], "index":0}; else return (results_df, fig)
    """
    import logging
    import numpy as np
    import pandas as pd
    import plotly.express as px
    import plotly.graph_objects as go

    preview_slides = [] if (export or return_payload) else None

    def _callback(msg: str):
        if progress_callback:
            try:
                progress_callback(msg)
            except Exception:
                pass
        logging.info(msg)

    _callback("Starting Citations Overview analysis...")

    if df is None or not isinstance(df, pd.DataFrame) or df.empty or "citations" not in df.columns:
        _callback("Missing 'citations' column or DataFrame is empty for citations analysis.")
        if preview_slides is not None:
            return {"slides": [{"title": "Citations Overview", "table_html": "<div>No data.</div>", "fig_json": None}], "index": 0}
        return pd.DataFrame(), None

    df_analysis = df.copy()

    df_analysis["citations"] = pd.to_numeric(df_analysis["citations"], errors="coerce").fillna(0).astype(int)

    if "year" in df_analysis.columns:
        df_analysis["year"] = pd.to_numeric(df_analysis["year"], errors="coerce")
    else:
        df_analysis["year"] = pd.NA

    if "doc_type_readable" not in df_analysis.columns:
        if "item_type" in df_analysis.columns:
            df_analysis["doc_type_readable"] = df_analysis["item_type"].apply(get_document_type)
        else:
            df_analysis["doc_type_readable"] = "Unknown"
    df_analysis["doc_type_readable"] = df_analysis["doc_type_readable"].fillna("Unknown")

    plot_type = params.get("plot_type", "top_cited_bar")
    top_n_table = int(params.get("top_n_cited_table", 20))
    top_n_plot = int(params.get("top_n_for_plot", 15))

    required_cols_for_table = ["title", "authors_short", "year", "source", "citations", "doi", "key"]
    for col in required_cols_for_table:
        if col not in df_analysis.columns:
            if col == "citations":
                df_analysis[col] = 0
            elif col == "year":
                df_analysis[col] = pd.NA
            else:
                df_analysis[col] = "N/A"

    df_analysis["year"] = df_analysis["year"].astype("Int64")

    results_df = (
        df_analysis.sort_values(by="citations", ascending=False)
        .head(top_n_table)[required_cols_for_table]
        .copy()
    )

    _callback(f"Generating citation plot: {plot_type}")
    fig = None

    if plot_type == "top_cited_bar":
        plot_df = results_df.head(top_n_plot).copy()
        if not plot_df.empty:
            fig = px.bar(
                plot_df,
                x="title",
                y="citations",
                title=f"Top {len(plot_df)} Most Cited Documents",
                hover_data=["authors_short", "year", "source", "doi"],
                labels={"title": "Document Title", "citations": "Citations"},
            )
            fig.update_xaxes(categoryorder="total descending", tickangle=-30, automargin=True)

    elif plot_type == "citations_vs_year_scatter":
        min_cites = int(params.get("min_citations_for_scatter", 0))
        scatter_df = df_analysis[(df_analysis["citations"] > min_cites) & (df_analysis["year"].notna())].copy()
        if not scatter_df.empty:
            scatter_df["year"] = scatter_df["year"].astype(int)
            scatter_df["log_citations"] = np.log1p(scatter_df["citations"])
            fig = px.scatter(
                scatter_df,
                x="year",
                y="citations",
                title="Citations vs. Publication Year",
                size="citations",
                hover_data=["title", "authors_short"],
                color="log_citations",
                trendline="ols",
            )
            fig.update_layout(xaxis_type="category")

    elif plot_type == "citation_distribution_histogram":
        cited_df = df_analysis[df_analysis["citations"] > 0].copy()
        if not cited_df.empty:
            fig = px.histogram(
                cited_df,
                x="citations",
                title="Distribution of Citation Counts (Cited Items > 0)",
                nbins=int(params.get("histogram_bins", 30)),
                labels={"citations": "Number of Citations"},
                marginal="box",
            )
            fig.update_layout(yaxis_title="Number of Documents")

    elif plot_type == "avg_citations_per_year_line":
        yr_df = df_analysis[df_analysis["year"].notna()].copy()
        if not yr_df.empty:
            yr_df["year"] = yr_df["year"].astype(int)
            avg_df = yr_df.groupby("year")["citations"].mean().reset_index().sort_values("year")
            if not avg_df.empty:
                fig = px.line(
                    avg_df,
                    x="year",
                    y="citations",
                    title="Average Citations per Publication by Year",
                    labels={"year": "Publication Year", "citations": "Average Citations"},
                    markers=True,
                )
                fig.update_layout(xaxis_type="category")

    elif plot_type == "citations_by_doc_type_box":
        cited_df = df_analysis[df_analysis["citations"] > 0].copy()
        min_docs = int(params.get("min_docs_per_type_for_box", 5))
        if not cited_df.empty:
            counts = cited_df["doc_type_readable"].value_counts()
            valid = counts[counts >= min_docs].index
            plot_df = cited_df[cited_df["doc_type_readable"].isin(valid)].copy()
            if not plot_df.empty:
                fig = px.box(
                    plot_df,
                    x="doc_type_readable",
                    y="citations",
                    title=f"Citation Distribution by Document Type (Types with ≥{min_docs} docs)",
                    labels={"doc_type_readable": "Document Type", "citations": "Number of Citations"},
                    points="outliers",
                    notched=True,
                )
                fig.update_xaxes(categoryorder="total descending", tickangle=-20, automargin=True)

    elif plot_type == "citations_sunburst_doc_year":
        sun_df = df_analysis[(df_analysis["citations"] > 0) & (df_analysis["year"].notna())].copy()
        if not sun_df.empty:
            sun_df["year_str"] = sun_df["year"].astype(str)
            top_types = sun_df["doc_type_readable"].value_counts().nlargest(7).index
            sun_df = sun_df[sun_df["doc_type_readable"].isin(top_types)].copy()
            if not sun_df.empty:
                fig = px.sunburst(
                    sun_df,
                    path=["doc_type_readable", "year_str"],
                    values="citations",
                    color="citations",
                    hover_data=["citations"],
                    title="Sunburst of Citations by Document Type and Year",
                )
                fig.update_layout(margin=dict(t=50, l=10, r=10, b=10))

    if fig is None:
        _callback(f"Could not generate '{plot_type}' for citations.")
        fig = go.Figure()
        fig.add_annotation(
            text=f"Plot for '{plot_type}' could not be generated.<br>Check data or parameters.",
            xref="paper",
            yref="paper",
            showarrow=False,
        )
        fig.update_layout(xaxis_visible=False, yaxis_visible=False)
    else:
        _callback(f"Generated '{plot_type}' for citations.")

    _callback("Citations Overview analysis complete.")

    if preview_slides is not None:
        preview_slides.append(
            {
                "title": f"Citations Overview — {plot_type}",
                "table_html": results_df.to_html(index=False, escape=True),
                # Ensure json-safe payload (avoid numpy arrays becoming strings in the host serializer).
                "fig_json": json.loads(json.dumps(fig.to_plotly_json(), cls=PlotlyJSONEncoder)),
                "notes": "",
            }
        )
        return {"slides": preview_slides, "index": 0}

    return results_df, fig
def analyze_temporal_analysis(
    df: pd.DataFrame,
    params: dict,
    progress_callback=lambda m: None,
    *,
    prs=None,
    collection_name: str = "Collection",
    slide_notes: bool = True,
    export: bool = False,
    return_payload: bool = False,
):
    """
    ###1. compute temporal table + plotly figure (existing logic)
    ###2. if prs provided, add PPT table + figure slides
    ###3. if export or return_payload, return {"slides":[...], "figs":[...], "index":0}
    """
    import io, json
    import numpy as np
    import pandas as pd
    import plotly.express as px
    import plotly.graph_objects as go
    from plotly.utils import PlotlyJSONEncoder
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN

    chart_type = params.get("plot_type") or "multi_line_trend"
    color_scheme = px.colors.qualitative.Plotly
    _cb = progress_callback or (lambda *_: None)

    preview_slides = [] if (export or return_payload) else None
    preview_figs = [] if (export or return_payload) else None

    FONT_FAMILY = "Segoe UI"
    PLOT_W, PLOT_H = 1200, 650

    if prs is not None:
        INCH = float(Inches(1))
        SLIDE_W_IN = float(prs.slide_width) / INCH
        SLIDE_H_IN = float(prs.slide_height) / INCH
    else:
        SLIDE_W_IN = 13.333
        SLIDE_H_IN = 7.5

    OUTER_MARGIN_X = 0.6
    OUTER_MARGIN_Y = 0.6
    TITLE_HEIGHT_IN = 0.7

    def _new_blank_slide_with_title(title_text: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(
            Inches(0.0),
            Inches(0.15),
            Inches(SLIDE_W_IN),
            Inches(TITLE_HEIGHT_IN),
        )
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        return slide

    def _set_notes(slide, text: str):
        if slide_notes and text:
            slide.notes_slide.notes_text_frame.text = str(text)

    def _add_table_slide(title_text: str, table_df: pd.DataFrame | None, notes: str = ""):
        slide = _new_blank_slide_with_title(title_text)
        _set_notes(slide, notes)

        if table_df is None or not isinstance(table_df, pd.DataFrame) or table_df.empty:
            tb = slide.shapes.add_textbox(
                Inches(OUTER_MARGIN_X),
                Inches(TITLE_HEIGHT_IN + 2.2),
                Inches(SLIDE_W_IN - 2 * OUTER_MARGIN_X),
                Inches(1.0),
            )
            p = tb.text_frame.paragraphs[0]
            p.text = "No data available."
            p.font.name = FONT_FAMILY
            p.font.size = Pt(16)
            p.alignment = PP_ALIGN.CENTER

            if preview_slides is not None:
                preview_slides.append({"title": title_text, "table_html": "<div>No data available.</div>", "notes": notes})
            return slide

        df_small = table_df.head(25).copy()

        max_w = SLIDE_W_IN - 2 * OUTER_MARGIN_X
        max_h = SLIDE_H_IN - (TITLE_HEIGHT_IN + 2 * OUTER_MARGIN_Y)

        rows = max(1, len(df_small) + 1)
        cols = max(1, df_small.shape[1])

        def _len_inch_est(s: object) -> float:
            t = "" if s is None or (isinstance(s, float) and pd.isna(s)) else str(s)
            t = " ".join(t.split())
            return 0.085 * max(4, len(t)) + 0.14

        col_needs = []
        for j, col in enumerate(df_small.columns):
            header_need = _len_inch_est(col)
            longest_need = header_need
            for i in range(len(df_small)):
                longest_need = max(longest_need, min(_len_inch_est(df_small.iloc[i, j]), 3.8))
            col_needs.append(longest_need)

        MIN_W = 0.9
        MAX_W = 3.6

        col_needs = [min(MAX_W, max(MIN_W, x)) for x in col_needs]
        total_need = float(np.sum(col_needs))

        if total_need > (max_w - 0.2):
            scale = (max_w - 0.2) / total_need
            col_widths_in = [max(MIN_W, x * scale) for x in col_needs]
        else:
            col_widths_in = col_needs[:]

        table_w = float(np.sum(col_widths_in))
        est_row_h = min(0.42, max(0.26, max_h / (rows + 1)))
        table_h = min(max_h, rows * est_row_h + 0.35)

        left_in = (SLIDE_W_IN - table_w) / 2.0
        top_in = TITLE_HEIGHT_IN + ((SLIDE_H_IN - TITLE_HEIGHT_IN) - table_h) / 2.0

        table = slide.shapes.add_table(
            rows, cols, Inches(left_in), Inches(top_in), Inches(table_w), Inches(table_h)
        ).table

        for j, w in enumerate(col_widths_in):
            table.columns[j].width = Emu(Inches(w))

        for j, col in enumerate(df_small.columns):
            cell = table.cell(0, j)
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(col)
            p.font.bold = True
            p.font.name = FONT_FAMILY
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER
            cell.margin_left = cell.margin_right = Emu(900)
            cell.margin_top = Emu(600)
            cell.margin_bottom = Emu(600)

        for i, (_, r) in enumerate(df_small.iterrows(), start=1):
            for j, col in enumerate(df_small.columns):
                cell = table.cell(i, j)
                txt = "" if pd.isna(r[col]) else " ".join(str(r[col]).split())
                tf = cell.text_frame
                tf.clear()
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = txt
                p.font.name = FONT_FAMILY
                p.font.size = Pt(11)
                p.alignment = PP_ALIGN.LEFT
                cell.margin_left = Emu(800)
                cell.margin_right = Emu(800)
                cell.margin_top = Emu(500)
                cell.margin_bottom = Emu(500)

        if preview_slides is not None:
            preview_slides.append({"title": title_text, "table_html": df_small.to_html(index=False, escape=True), "notes": notes})

        return slide

    def _add_figure_slide(title_text: str, fig: go.Figure, notes: str = ""):
        slide = _new_blank_slide_with_title(title_text)
        _set_notes(slide, notes)

        content_w = SLIDE_W_IN - 2 * OUTER_MARGIN_X
        content_h = SLIDE_H_IN - (TITLE_HEIGHT_IN + 2 * OUTER_MARGIN_Y)

        png = fig.to_image(format="png", width=PLOT_W, height=PLOT_H, scale=2)
        ar = (PLOT_W / PLOT_H) if PLOT_H else 1.7778

        target_w = content_w
        target_h = target_w / ar
        if target_h > content_h:
            target_h = content_h
            target_w = target_h * ar
        target_w *= 0.98
        target_h *= 0.98

        left_in = (SLIDE_W_IN - target_w) / 2.0
        top_in = TITLE_HEIGHT_IN + (SLIDE_H_IN - TITLE_HEIGHT_IN - target_h) / 2.0

        slide.shapes.add_picture(
            io.BytesIO(png),
            Inches(left_in),
            Inches(top_in),
            width=Inches(target_w),
            height=Inches(target_h),
        )

        if preview_slides is not None:
            fig_json = json.dumps(fig.to_plotly_json(), cls=PlotlyJSONEncoder)
            preview_figs.append({"title": title_text, "fig_json": fig_json, "notes": notes or ""})
            preview_slides.append({"title": title_text, "table_html": "", "notes": notes or ""})

        return slide

    _cb(f"Temporal Analysis: Received request for '{chart_type}'.")

    if chart_type == "crosstab_heatmap":
        feature1 = params.get("feature1")
        feature2 = params.get("feature2")

        _cb(f"Running Crosstab analysis for '{feature1}' vs '{feature2}'...")

        if not all([feature1, feature2]) or feature1 not in df.columns or feature2 not in df.columns:
            table_df = pd.DataFrame()
            fig = go.Figure().add_annotation(text="Invalid features selected for comparison.")
        else:
            df_c = df[[feature1, feature2]].copy().dropna()
            for col in [feature1, feature2]:
                if df_c[col].dtype == "object" and df_c[col].dropna().apply(isinstance, args=(list,)).any():
                    df_c = df_c.explode(col)
            df_c.dropna(inplace=True)

            if df_c.empty:
                table_df = pd.DataFrame()
                fig = go.Figure().add_annotation(text="No overlapping data for selected features.")
            else:
                crosstab_df = pd.crosstab(df_c[feature1], df_c[feature2])
                fig = px.imshow(
                    crosstab_df,
                    text_auto=True,
                    aspect="auto",
                    labels=dict(
                        x=feature2.replace("_", " ").title(),
                        y=feature1.replace("_", " ").title(),
                        color="Count",
                    ),
                    title=f"Comparison: {feature1.replace('_', ' ').title()} vs. {feature2.replace('_', ' ').title()}",
                )
                fig.update_xaxes(side="top")
                table_df = crosstab_df.reset_index()
        return_df = table_df

    else:
        feature_col = params.get("feature", "publications")
        top_n = params.get("top_n", 10)
        year_from = params.get("year_from")
        year_to = params.get("year_to")

        _cb(f"Preparing trend data for '{feature_col}'...")

        if "year" not in df.columns:
            return_df = pd.DataFrame()
            fig = go.Figure().add_annotation(text="Error: 'year' column not found.")
        else:
            d = df.copy()
            d["year"] = pd.to_numeric(d["year"], errors="coerce")
            d.dropna(subset=["year"], inplace=True)
            d["year"] = d["year"].astype(int)

            if year_from:
                d = d[d["year"] >= year_from]
            if year_to:
                d = d[d["year"] <= year_to]

            if d.empty:
                return_df = pd.DataFrame()
                fig = go.Figure().add_annotation(text="No data in selected year range.")
            else:
                agg_df = pd.DataFrame()
                feature_col_for_plot = feature_col if feature_col != "publications" else "feature"

                if feature_col == "publications":
                    agg_df = d.groupby("year").size().reset_index(name="Count")
                    agg_df[feature_col_for_plot] = "Publications"
                else:
                    if feature_col not in d.columns:
                        return pd.DataFrame(), go.Figure().add_annotation(text=f"Error: Feature '{feature_col}' not found.")
                    if d[feature_col].dtype == "object" and d[feature_col].dropna().apply(isinstance, args=(list,)).any():
                        d = d.explode(feature_col)
                    d.dropna(subset=[feature_col], inplace=True)
                    if d.empty:
                        return_df = pd.DataFrame()
                        fig = go.Figure().add_annotation(text=f"No data for '{feature_col}'.")
                    else:
                        top_categories = d[feature_col].value_counts().nlargest(top_n).index
                        d_filtered = d[d[feature_col].isin(top_categories)]
                        agg_df = d_filtered.groupby(["year", feature_col_for_plot]).size().reset_index(name="Count")

                if agg_df.empty:
                    return_df = pd.DataFrame()
                    fig = go.Figure().add_annotation(text=f"No data to plot for '{feature_col}'.")
                else:
                    agg_df.sort_values("year", inplace=True)

                    _cb(f"Creating '{chart_type}' chart for '{feature_col}'...")
                    title_feature = feature_col.replace("_", " ").title()
                    title_chart = chart_type.replace("_", " ").title()
                    full_title = f"{title_feature} Trend: {title_chart}"

                    if chart_type == "multi_line_trend":
                        fig = px.line(
                            agg_df,
                            x="year",
                            y="Count",
                            color=feature_col_for_plot,
                            markers=True,
                            title=full_title,
                            color_discrete_sequence=color_scheme,
                        )
                        fig.update_layout(legend_title_text=title_feature)
                    elif chart_type == "bar":
                        fig = px.bar(
                            agg_df,
                            x="year",
                            y="Count",
                            color=feature_col_for_plot,
                            title=full_title,
                            barmode="stack",
                            color_discrete_sequence=color_scheme,
                        )
                        fig.update_layout(legend_title_text=title_feature)
                    elif chart_type == "area":
                        fig = px.area(
                            agg_df,
                            x="year",
                            y="Count",
                            color=feature_col_for_plot,
                            title=full_title,
                            color_discrete_sequence=color_scheme,
                        )
                        fig.update_layout(legend_title_text=title_feature)
                    elif chart_type == "cumulative_area":
                        agg_df["Cumulative Count"] = agg_df.groupby(feature_col_for_plot)["Count"].cumsum()
                        fig = px.area(
                            agg_df,
                            x="year",
                            y="Cumulative Count",
                            color=feature_col_for_plot,
                            title=f"Cumulative Growth of {title_feature}",
                            color_discrete_sequence=color_scheme,
                        )
                        fig.update_layout(legend_title_text=title_feature, yaxis_title="Cumulative Count")
                    elif chart_type == "bubble":
                        bubble_df = agg_df.groupby(feature_col_for_plot)["Count"].sum().reset_index()
                        fig = px.scatter(
                            bubble_df,
                            x=feature_col_for_plot,
                            y="Count",
                            size="Count",
                            color=feature_col_for_plot,
                            title=f"Overall Contribution of Top {top_n} {title_feature}",
                            hover_name=feature_col_for_plot,
                            color_discrete_sequence=color_scheme,
                        )
                        fig.update_traces(marker=dict(sizemin=5, line_width=1, line_color="black"))
                        fig.update_layout(xaxis_title=title_feature, yaxis_title="Total Count")
                    elif chart_type == "heat":
                        pivot_df = agg_df.pivot_table(index=feature_col_for_plot, columns="year", values="Count", fill_value=0)
                        fig = px.imshow(
                            pivot_df,
                            text_auto=".0f",
                            aspect="auto",
                            title=f"Activity Heatmap: {title_feature} Over Time",
                            color_continuous_scale="YlGnBu",
                        )
                        fig.update_layout(xaxis_title="Year", yaxis_title=title_feature)
                        return_df = pivot_df.reset_index()
                    else:
                        fig = px.line(agg_df, x="year", y="Count", color=feature_col_for_plot, markers=True, title=full_title)

                    if "year" in agg_df.columns and agg_df["year"].nunique() < 25:
                        fig.update_xaxes(dtick=1)

                    if "return_df" not in locals():
                        return_df = agg_df

    if prs is not None:
        base_title = f"{collection_name} · Temporal Analysis — {chart_type}"
        _add_table_slide(f"{base_title} — Table", return_df, notes=f"Source: analyze_temporal_analysis(plot_type='{chart_type}').")
        _add_figure_slide(f"{base_title} — Figure", fig, notes=f"Figure: analyze_temporal_analysis(plot_type='{chart_type}').")

    if preview_slides is not None:
        base_title = f"{collection_name} · Temporal Analysis — {chart_type}"
        try:
            table_html = (
                return_df.head(50).to_html(index=False, escape=True)
                if isinstance(return_df, pd.DataFrame) and not return_df.empty
                else "<div>No data available.</div>"
            )
        except Exception:
            table_html = "<div>Unable to render table.</div>"

        fig_json = None
        try:
            fig_json = fig.to_plotly_json() if fig is not None else None
        except Exception:
            fig_json = None

        preview_slides.append(
            {
                "title": f"{base_title} — Table",
                "table_html": table_html,
                "notes": f"Source: analyze_temporal_analysis(plot_type='{chart_type}').",
            }
        )
        preview_slides.append(
            {
                "title": f"{base_title} — Figure",
                "table_html": "",
                "fig_json": fig_json,
                "notes": f"Figure: analyze_temporal_analysis(plot_type='{chart_type}').",
            }
        )
        if preview_figs is not None and fig_json is not None:
            preview_figs.append(json.dumps(fig_json, cls=PlotlyJSONEncoder))
        return {"slides": preview_slides, "figs": preview_figs or [], "index": 0}

    return return_df, fig




# ──────────────────────────────────────────────────────────────────
# ──────────────────────────────────────────────────────────────────
def _make_department_sunburst(df: pd.DataFrame,
                              *,
                              include_authors: bool = False) -> go.Figure:
    """
    Build  Country → Institution → Department (→ Author)  sunburst.
    Works with the new list-of-countries column.
    """
    cols = ["country", "institution", "department"]

    base = (
        df.explode("country")          # <<<  key line: list → rows
          .dropna(subset=["country", "institution"])
          .assign(department=lambda d: d["department"].fillna("—"))
    )

    if include_authors and "authors" in base.columns:
        base = base.explode("authors")
        cols.append("authors")

    # Plotly requires that all internal nodes have at least one child
    def _row_ok(r):
        # if a parent level is NA while a child is not → reject
        for i in range(len(cols) - 1):
            if pd.isna(r[cols[i]]) and pd.notna(r[cols[i + 1]]):
                return False
        return True

    base = base.loc[base.apply(_row_ok, axis=1)]
    if base.empty:
        return go.Figure().add_annotation(
            text="No rows form a complete hierarchy.", showarrow=False)

    base["value"] = 1
    fig = px.sunburst(base, path=cols, values="value",
                      title=" → ".join(c.title() for c in cols))
    fig.update_layout(margin=dict(l=0, r=0, t=60, b=0))
    return fig
def _make_world_map(df: pd.DataFrame, *, by_authors: bool = False) -> go.Figure:
    """
    Offline-first geo view.

    NOTE: Plotly choropleths typically require loading external topojson (often blocked by CSP/offline).
    Use a scattergeo "bubble map" instead.
    """
    return _make_bubble_map(df, by_authors=by_authors)




def _make_bubble_map(df: pd.DataFrame, *, by_authors: bool = False) -> go.Figure:
    cc = _country_counts(df, by_authors=by_authors)
    if cc.empty:
        return go.Figure().add_annotation(text="No country data.", showarrow=False)

    # attach centroids
    cc[["lat", "lon"]] = cc["country"].apply(
        lambda c: pd.Series(_centroid(c)))
    cc = cc.dropna(subset=["lat", "lon"])
    if cc.empty:
        return go.Figure().add_annotation(text="No coordinates.", showarrow=False)

    scale = max(cc["count"].max() / 60, 1)
    fig = go.Figure(go.Scattergeo(
        lat=cc["lat"], lon=cc["lon"],
        text=cc["country"] + "<br>" + cc["count"].astype(str),
        marker=dict(size=cc["count"] / scale,
                    color=cc["count"],
                    colorscale="Plasma",
                    colorbar_title="Count",
                    line_width=0.4),
        hovertemplate="%{text}<extra></extra>"
    ))
    fig.update_layout(title=("Author" if by_authors else "Publication") +
                             " Bubble-map",
                      margin=dict(l=0, r=0, t=60, b=0),
                      geo=dict(showland=True, landcolor="#E5E5E5",
                               showcountries=True, countrycolor="white"))
    return fig

#  Country collaboration network  –  force-directed graph
# ─────────────────────────────────────────────────────────────────────────────
def _make_collab_network_country(df: pd.DataFrame,
                                 min_edge: int = 2,
                                 top_n:   int = 30) -> go.Figure:
    """
    Undirected co-authorship network between countries.
    • Each publication → all unordered country pairs
    • Keep edges with ≥ `min_edge` occurrences
    • Restrict to `top_n` most-connected nodes
    """
    import itertools, networkx as nx, plotly.graph_objects as go, pandas as pd

    # -------- build edge list ------------------------------------------------
    pairs = (df["country"]
             .explode()                               # list → rows
             .groupby(level=0)                        # back to publication
             .apply(lambda s: list(itertools.combinations(sorted(set(s)), 2)))
             .explode()
             .dropna())

    if pairs.empty:
        return go.Figure().add_annotation(text="No country data.", showarrow=False)

    edges = (pairs.value_counts()
                   .rename_axis("pair")
                   .reset_index(name="weight")
                   .query("weight >= @min_edge"))

    if edges.empty:
        return go.Figure().add_annotation(
            text=f"No links with ≥ {min_edge} collaborations.",
            showarrow=False)

    # -------- keep most-connected nodes -------------------------------------
    node_rank = (pd.Series([n for a, b in edges["pair"] for n in (a, b)])
                   .value_counts()
                   .head(top_n))
    keep = set(node_rank.index)
    edges = edges[edges["pair"].apply(lambda p: p[0] in keep and p[1] in keep)]

    if edges.empty:
        return go.Figure().add_annotation(text="Nothing after top-N filter.", showarrow=False)

    # -------- build NetworkX graph ------------------------------------------
    G = nx.Graph()
    for (a, b), w in zip(edges["pair"], edges["weight"]):
        G.add_edge(a, b, weight=w)

    pos = nx.spring_layout(G, k=0.5, seed=42)           # coordinates
    node_x, node_y, node_text = [], [], []
    for n in G.nodes:
        x, y = pos[n]
        node_x.append(x); node_y.append(y)
        node_text.append(f"{n} ({G.degree[n]})")

    edge_x, edge_y = [], []
    for a, b in G.edges:
        x0, y0 = pos[a]; x1, y1 = pos[b]
        edge_x += [x0, x1, None]
        edge_y += [y0, y1, None]

    # -------- plotly figure --------------------------------------------------
    fig = go.Figure()

    fig.add_trace(go.Scatter(
        x=edge_x, y=edge_y,
        mode="lines",
        line=dict(width=0.5, color="#888"),
        hoverinfo="none"))

    fig.add_trace(go.Scatter(
        x=node_x, y=node_y,
        mode="markers+text",
        text=node_text,
        textposition="top center",
        marker=dict(size=[4+3*G.degree[n] for n in G.nodes],
                    color="#1f77b4",
                    line=dict(width=1, color="#333")),
        hovertemplate="%{text}<extra></extra>"
    ))

    fig.update_layout(
        title="Country collaboration network",
        showlegend=False,
        margin=dict(l=20, r=20, t=60, b=20),
        xaxis=dict(showgrid=False, zeroline=False, visible=False),
        yaxis=dict(showgrid=False, zeroline=False, visible=False))
    return fig


# ─────────────────────────────────────────────────────────────────────────────
#  Institution collaboration network  –  force-directed graph
# ─────────────────────────────────────────────────────────────────────────────
# ──────────────────────────────────────────────────────────────────
# 2)  Institution scatter  (explode country list before groupby)
# ──────────────────────────────────────────────────────────────────
def _make_scatter_institutions(df: pd.DataFrame) -> tuple[pd.DataFrame, go.Figure]:
    # explode once – afterwards every row has exactly one country string
    inst = df.explode("country")[["institution", "country"]].dropna(how="any")
    if inst.empty:
        return pd.DataFrame(), go.Figure().add_annotation(
            text="No institution data.", showarrow=False)

    cnt = (inst.groupby(["institution", "country"])
                .size()
                .reset_index(name="count"))

    fig = px.scatter_geo(
        cnt, locations="country", locationmode="country names",
        size="count", color="count", hover_name="institution",
        projection="natural earth", color_continuous_scale="Viridis",
        title="Institutions by Country"
    )
    fig.update_layout(margin=dict(l=0, r=0, t=50, b=0))
    return cnt, fig



# ──────────────────────────────────────────────────────────────────────────────
#  Institution → (lat, lon) cache
# ──────────────────────────────────────────────────────────────────────────────
@functools.lru_cache(maxsize=1024)
def get_institution_latlon(institution: str) -> Tuple[Optional[float], Optional[float]]:
    """
    Geocode an institution name to (latitude, longitude).
    Returns (None, None) if it can’t be found or on error.
    """
    if not institution or institution.strip().lower() in ("", "none", "none specified"):
        return None, None

    query = institution.strip()
    try:
        location = _geocode(query)
        if location:
            return location.latitude, location.longitude
        else:
            # Try appending “University” if not already present
            if "univ" not in query.lower():
                location = _geocode(query + " University")
                if location:
                    return location.latitude, location.longitude
    except Exception as e:
        logging.warning(f"Geocoding error for '{institution}': {e}")
    return None, None

def _make_country_pie(cc: pd.DataFrame, top_n: int) -> go.Figure:
    """
    cc: DataFrame with columns ['country','count'], already top-N filtered.
    """
    if cc.empty:
        return go.Figure().add_annotation(
            text="No country data available.",
            showarrow=False
        )

    fig = px.pie(
        cc,
        names="country",
        values="count",
        title=f"Country distribution (Top {top_n})",
    )
    fig.update_traces(
        textinfo="label+percent",
        textposition="inside",
        hovertemplate="%{label}: %{value} publications<extra></extra>"
    )
    fig.update_layout(
        margin=dict(l=0, r=0, t=60, b=0),
        legend_title="Country"
    )
    return fig

def _make_collab_network_institution(df: pd.DataFrame,
                                     min_edge: int = 2,
                                     top_n:   int = 30) -> go.Figure:
    """
    Undirected co-authorship network between institutions.
    • Splits ‘institution’ on “;”
    • Edge weight = joint publications
    """
    import itertools, networkx as nx, plotly.graph_objects as go, pandas as pd

    # -------- edge list ------------------------------------------------------
    pairs = (df["institution"]
             .apply(_split_institutions)
             .apply(lambda lst: list(itertools.combinations(sorted(set(lst)), 2)))
             .explode()
             .dropna())

    if pairs.empty:
        return go.Figure().add_annotation(text="No institution data.", showarrow=False)

    edges = (pairs.value_counts()
                   .rename_axis("pair")
                   .reset_index(name="weight")
                   .query("weight >= @min_edge"))

    if edges.empty:
        return go.Figure().add_annotation(
            text=f"No links with ≥ {min_edge} collaborations.",
            showarrow=False)

    # -------- top nodes ------------------------------------------------------
    node_rank = (pd.Series([n for a, b in edges["pair"] for n in (a, b)])
                   .value_counts()
                   .head(top_n))
    keep = set(node_rank.index)
    edges = edges[edges["pair"].apply(lambda p: p[0] in keep and p[1] in keep)]
    if edges.empty:
        return go.Figure().add_annotation(text="Nothing after top-N filter.", showarrow=False)

    # -------- graph & layout -------------------------------------------------
    G = nx.Graph()
    for (a, b), w in zip(edges["pair"], edges["weight"]):
        G.add_edge(a, b, weight=w)

    pos = nx.spring_layout(G, k=0.6, seed=42)
    node_x, node_y, node_text = [], [], []
    for n in G.nodes:
        x, y = pos[n]
        node_x.append(x); node_y.append(y)
        node_text.append(f"{n} ({G.degree[n]})")

    edge_x, edge_y = [], []
    for a, b in G.edges:
        x0, y0 = pos[a]; x1, y1 = pos[b]
        edge_x += [x0, x1, None]
        edge_y += [y0, y1, None]

    # -------- figure ---------------------------------------------------------
    fig = go.Figure()

    fig.add_trace(go.Scatter(
        x=edge_x, y=edge_y,
        mode="lines",
        line=dict(width=0.4, color="#aaa"),
        hoverinfo="none"))

    fig.add_trace(go.Scatter(
        x=node_x, y=node_y,
        mode="markers",
        text=node_text,
        marker=dict(size=[4+2*G.degree[n] for n in G.nodes],
                    color="#e377c2",
                    line=dict(width=1, color="#333")),
        hovertemplate="%{text}<extra></extra>"
    ))

    fig.update_layout(
        title="Institution collaboration network",
        showlegend=False,
        margin=dict(l=20, r=20, t=60, b=20),
        xaxis=dict(showgrid=False, zeroline=False, visible=False),
        yaxis=dict(showgrid=False, zeroline=False, visible=False))
    return fig

def _make_chord_inst_country(df: pd.DataFrame, top_n: int = 20) -> go.Figure:
    """
    Circular Sankey: institution-to-country links.
    """
    tmp = df.explode("country")                        # <<< new
    sub = (tmp[["country", "institution"]]
           .dropna()
           .groupby(["country", "institution"])
           .size()
           .reset_index(name="value")
           .sort_values("value", ascending=False)
           .head(top_n))

    if sub.empty:
        return go.Figure().add_annotation(
            text="No data to build chord.", showarrow=False)

    countries    = sorted(sub["country"].unique())
    institutions = sorted(sub["institution"].unique())
    nodes        = countries + institutions
    idx          = {name: i for i, name in enumerate(nodes)}

    fig = go.Figure(go.Sankey(
        arrangement="fixed",
        node=dict(pad=6, thickness=12,
                  label=nodes,
                  color=["#3182bd"] * len(countries) +
                        ["#6baed6"] * len(institutions)),
        link=dict(
            source=[idx[c] for c in sub["country"]],
            target=[idx[i] for i in sub["institution"]],
            value=sub["value"]
        )
    ))
    fig.update_layout(title=f"Top {top_n} Institution ⇄ Country links",
                      font=dict(size=10),
                      margin=dict(l=10, r=10, t=50, b=10))
    return fig

def analyze_affiliations(
    df: "pd.DataFrame",
    params: dict,
    progress_callback=None,
    *,
    slide_notes: bool = True,
    export: bool = False,
    return_payload: bool = False,
):
    """
    ###1. dispatch plot_type -> (table_df, fig)
    ###2. if export or return_payload: return {"slides":[...], "figs":[...], "imgs":[...], "tables":[...], "payload":{...}, "index":0}
         else: return (table_df, fig)
    """
    import json
    import pandas as pd
    import plotly.graph_objects as go
    from plotly.utils import PlotlyJSONEncoder

    def _cb(msg: str):
        if progress_callback:
            try:
                progress_callback(f"AFFIL: {msg}")
            except Exception:
                pass

    def _fig_json_str(fig_obj: "go.Figure") -> str:
        return json.dumps(fig_obj.to_plotly_json(), cls=PlotlyJSONEncoder)

    preview_slides = [] if (export or return_payload) else None
    preview_figs = [] if (export or return_payload) else None
    preview_tables = [] if (export or return_payload) else None
    preview_imgs = [] if (export or return_payload) else None

    PLOT = str(params.get("plot_type", "world_map_pubs") or "world_map_pubs")
    TOPN = int(params.get("top_n", 20) or 20)
    MIN_EDGE = int(params.get("min_collaborations", 2) or 2)

    if df is None or not isinstance(df, pd.DataFrame) or df.empty:
        _cb("No data.")
        empty_df = pd.DataFrame()
        fig = go.Figure()
        fig.add_annotation(text="No data.", showarrow=False, xref="paper", yref="paper")
        if preview_slides is not None:
            slide = {
                "title": f"Affiliations — {PLOT}",
                "table_html": "<div>No data.</div>",
                "fig_json": _fig_json_str(fig),
                "notes": "Source: analyze_affiliations()." if slide_notes else "",
            }
            preview_slides.append(slide)
            preview_figs.append(slide["fig_json"])
            preview_tables.append(slide["table_html"])
            return {
                "slides": preview_slides,
                "figs": preview_figs,
                "imgs": preview_imgs,
                "tables": preview_tables,
                "payload": {"section": "Affiliations_geo", "plot_type": PLOT, "top_n": TOPN, "min_edge": MIN_EDGE, "n_rows": 0},
                "index": 0,
            }
        return empty_df, fig

    df_clean = df

    table_df = pd.DataFrame()
    fig = None

    _cb(f"plot_type={PLOT}")

    if PLOT == "world_map_pubs":
        table_df = _country_counts(df_clean, by_authors=False)
        fig = _make_world_map(df_clean, by_authors=False)

    elif PLOT == "world_map_authors":
        table_df = _country_counts(df_clean, by_authors=True)
        fig = _make_world_map(df_clean, by_authors=True)

    elif PLOT == "bubble_map_pubs":
        table_df = _country_counts(df_clean, by_authors=False)
        fig = _make_bubble_map(df_clean, by_authors=False)

    elif PLOT == "bubble_map_authors":
        table_df = _country_counts(df_clean, by_authors=True)
        fig = _make_bubble_map(df_clean, by_authors=True)

    elif PLOT == "geo_scatter_institutions":
        table_df, fig = _make_scatter_institutions(df_clean)

    elif PLOT == "country_bar":
        import plotly.express as px

        cc = _country_counts(df_clean)
        if isinstance(cc, pd.DataFrame) and ("count" in cc.columns):
            table_df = cc.nlargest(TOPN, "count").copy()
        else:
            table_df = pd.DataFrame()

        fig = px.bar(
            table_df if not table_df.empty else pd.DataFrame({"country": [], "count": []}),
            x="country",
            y="count",
            title=f"Top {TOPN} Countries",
            color="count",
            color_continuous_scale="Viridis",
        )
        fig.update_layout(xaxis_tickangle=-45, margin=dict(t=60))

    elif PLOT == "country_pie":
        cc = _country_counts(df_clean)
        if isinstance(cc, pd.DataFrame) and ("count" in cc.columns):
            table_df = cc.nlargest(TOPN, "count").copy()
        else:
            table_df = pd.DataFrame()
        fig = _make_country_pie(table_df, TOPN)

    elif PLOT == "institution_bar":
        import plotly.express as px

        if "institution" in df_clean.columns:
            inst = (
                df_clean["institution"]
                .astype(str)
                .replace("nan", "")
                .str.strip()
                .replace("", pd.NA)
                .dropna()
                .value_counts()
                .head(TOPN)
                .reset_index(name="count")
                .rename(columns={"index": "institution"})
            )
            table_df = inst.copy()
        else:
            table_df = pd.DataFrame({"institution": [], "count": []})

        fig = px.bar(
            table_df,
            x="institution",
            y="count",
            title=f"Top {TOPN} Institutions",
            color="count",
            color_continuous_scale="Viridis",
        )
        fig.update_layout(xaxis_tickangle=-45, margin=dict(t=60))

    elif PLOT == "department_sunburst":
        fig = _make_department_sunburst(df_clean, include_authors=bool(params.get("with_authors", False)))
        table_df = pd.DataFrame()

    elif PLOT == "collab_network_country":
        fig = _make_collab_network_country(df_clean, MIN_EDGE, TOPN)
        table_df = pd.DataFrame()

    elif PLOT == "collab_network_institution":
        fig = _make_collab_network_institution(df_clean, MIN_EDGE, TOPN)
        table_df = pd.DataFrame()

    elif PLOT == "chord_inst_country":
        fig = _make_chord_inst_country(df_clean, TOPN)
        table_df = pd.DataFrame()

    else:
        fig = go.Figure()
        fig.add_annotation(text=f"Unknown plot_type '{PLOT}'.", showarrow=False, xref="paper", yref="paper")
        table_df = pd.DataFrame()

    if fig is None:
        fig = go.Figure()
        fig.add_annotation(text=f"Plot '{PLOT}' not generated.", showarrow=False, xref="paper", yref="paper")

    payload = {
        "section": "Affiliations_geo",
        "plot_type": PLOT,
        "top_n": int(TOPN),
        "min_edge": int(MIN_EDGE),
        "n_rows": int(len(df_clean)),
    }

    if preview_slides is not None:
        table_html = table_df.to_html(index=False, escape=True) if isinstance(table_df, pd.DataFrame) and not table_df.empty else "<div>No table.</div>"
        notes = f"Source: analyze_affiliations(plot_type='{PLOT}')." if slide_notes else ""

        slide = {
            "title": f"Affiliations — {PLOT}",
            "table_html": table_html,
            "fig_json": _fig_json_str(fig),
            "notes": notes,
        }

        preview_slides.append(slide)
        preview_figs.append(slide["fig_json"])
        preview_tables.append(slide["table_html"])

        return {
            "slides": preview_slides,
            "figs": preview_figs,
            "imgs": preview_imgs,
            "tables": preview_tables,
            "payload": payload,
            "index": 0,
        }

    return table_df, fig


def _plot_design_crosstab_heatmap(df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    """
    Heatmap for co-occurrence of two design dimensions that truly reflects the UI selections.

    Why your axis “look stuck”:
      • In some datasets the UI text may not exactly match column names (spacing/case).
      • Some columns are duplicated or carry list/; separated values.
    Fixes:
      1) Resolve the requested names to real dataframe columns (slug-aware).
      2) Tokenise multi-coded cells (lists or ';').
      3) Build pairs and crosstab using temporary axis names to avoid duplicate-name issues.
      4) Always label axes with the *resolved* column names so you can see what’s actually used.
    """
    import re, itertools

    # --- 1) Resolve UI names to actual DF columns (slug-aware) ---
    dfu = df.loc[:, ~df.columns.duplicated()].copy()

    def _slug(s: str) -> str:
        s = re.sub(r"[^A-Za-z0-9]+", "_", s.strip().lower())
        return re.sub(r"_+", "_", s).strip("_")

    # build a reverse map from slug -> original column
    slug_map: dict[str, str] = {}
    for c in map(str, dfu.columns):
        slug_map[_slug(c)] = c

    def _resolve(name: str, default: str) -> str:
        if not isinstance(name, str) or not name.strip():
            name = default
        if name in dfu.columns:
            return name
        s = _slug(name)
        if s in slug_map:
            return slug_map[s]
        # last resort: exact default if present, else first col
        return default if default in dfu.columns else list(map(str, dfu.columns))[0]

    x_req = params.get("dimension1", "framework_model")
    y_req = params.get("dimension2", "methodology")

    x_col = _resolve(x_req, "framework_model")
    y_col = _resolve(y_req, "methodology")

    # --- 2) Tokeniser for multi-coded cells ---
    def _tokens(cell):
        if cell is None:
            return []
        try:
            if pd.isna(cell):
                return []
        except Exception:
            pass
        if isinstance(cell, (list, tuple, set)):
            items = [str(v).strip() for v in cell]
        else:
            items = [t.strip() for t in re.split(r"[;,]", str(cell))]
        return [t for t in items if t and t.lower() != "none"]

    # --- 3) Build all (Y, X) pairs row-wise (cartesian product when multi-coded) ---
    pairs = []
    # Use .itertuples to avoid boolean-mask alignment problems; iterate raw values
    for xv, yv in dfu[[x_col, y_col]].itertuples(index=False, name=None):
        xs = _tokens(xv)
        ys = _tokens(yv)
        if not xs or not ys:
            continue
        for a, b in itertools.product(xs, ys):
            pairs.append((b, a))  # rows=Y, cols=X

    if not pairs:
        return pd.DataFrame(), go.Figure().add_annotation(
            text=f"No co-occurrence after tokenising '{x_col}' and '{y_col}'.",
            showarrow=False
        )

    pair_df = pd.DataFrame(pairs, columns=["__y__", "__x__"])

    # --- 4) Crosstab + figure (safe temp names) ---
    ctab = pd.crosstab(pair_df["__y__"], pair_df["__x__"])
    if ctab.empty:
        return pd.DataFrame(), go.Figure().add_annotation(
            text=f"No co-occurrence data for {y_col} × {x_col}.",
            showarrow=False
        )

    fig = px.imshow(
        ctab,
        text_auto=True,
        aspect="auto",
        labels=dict(x=x_col.replace("_", " ").title(),
                    y=y_col.replace("_", " ").title(),
                    color="Count"),
        title=f"Co-occurrence: {y_col.replace('_',' ').title()} × {x_col.replace('_',' ').title()}"
    )

    # Return tidy table with real axis name for rows
    ctab_reset = ctab.reset_index().rename(columns={"__y__": y_col})
    return ctab_reset, fig
def _plot_design_hierarchical_sunburst(df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    """
    Sunburst with two modes:

    • Field→Keywords: use params['sunburst_field'] as the base; leafs are controlled_vocabulary_terms.
    • Path mode: use params['path'] = [col1, col2, ...] (1–3 levels).

    Critical: avoid duplicate column names before calling px.sunburst.
    """
    import re

    # --- helpers ---
    def _dedup_columns(d: pd.DataFrame) -> pd.DataFrame:
        # Narwhals (used by Plotly 6+) requires unique column names
        # https://plotly.com/python-api-reference/generated/plotly.express.sunburst.html
        # https://pandas.pydata.org/docs/user_guide/duplicates.html
        return d.loc[:, ~d.columns.duplicated()].copy()

    def _split_tokens(cell) -> list[str]:
        if cell is None:
            return []
        try:
            if pd.isna(cell):
                return []
        except Exception:
            pass
        if isinstance(cell, (list, tuple, set)):
            toks = [str(x).strip() for x in cell]
        else:
            toks = [t.strip() for t in re.split(r"[;,]", str(cell))]
        return [t for t in toks if t]

    def _ensure_kw_list(d: pd.DataFrame, base_col: str = "controlled_vocabulary_terms") -> str:
        # Alias controlled_vocabulary_terms if needed, then create *_list for keywords
        if base_col not in d.columns and "controlled_vocabulary_terms" in d.columns:
            d[base_col] = d["controlled_vocabulary_terms"]
        list_col = f"{base_col}_list"
        if list_col not in d.columns and base_col in d.columns:
            d[list_col] = d[base_col].apply(_split_tokens)
        return list_col if list_col in d.columns else ""

    dfu = _dedup_columns(df)

    # MODE 1: Field → Keywords
    base_field = (params.get("sunburst_field") or "").strip()
    if base_field:
        kw_list_col = _ensure_kw_list(dfu, "controlled_vocabulary_terms")
        if base_field not in dfu.columns or not kw_list_col:
            txt = f"Missing column(s) for Field→Keywords sunburst: base='{base_field}', keywords='controlled_vocabulary_terms'."
            return pd.DataFrame(), go.Figure().add_annotation(text=txt, showarrow=False)

        # Tokenise the base field too (arrays / ';' are common in your schema)
        base_list_col = f"__{base_field}__list"
        dfu[base_list_col] = dfu[base_field].apply(_split_tokens)

        # Explode both axes
        w = dfu.explode(base_list_col)
        w = w.explode(kw_list_col)

        # Clean
        w = w.dropna(subset=[base_list_col, kw_list_col])
        w[base_list_col] = w[base_list_col].astype(str).str.strip()
        w[kw_list_col] = w[kw_list_col].astype(str).str.strip()
        w = w[(w[base_list_col] != "") & (w[kw_list_col] != "")]
        if w.empty:
            return pd.DataFrame(), go.Figure().add_annotation(text="No (field, keyword) pairs after tokenisation.", showarrow=False)

        # Avoid duplicate column names: drop the original base field before renaming temp → base_field
        if base_field in w.columns:
            w = w.drop(columns=[base_field])
        w = w.rename(columns={base_list_col: base_field, kw_list_col: "keyword"})
        w = _dedup_columns(w)

        # Table + figure
        tbl = (
            w[[base_field, "keyword"]]
            .value_counts()
            .reset_index(name="count")
        )
        fig = px.sunburst(
            w,
            path=[base_field, "keyword"],
            title=f"Sunburst — {base_field.replace('_', ' ').title()} → Keywords"
        )
        return tbl, fig

    # MODE 2: Legacy multi-level path (no keywords)
    path = [c for c in (params.get("path") or []) if isinstance(c, str) and c.strip()]
    if not path:
        path = ["methods", "contribution_type"]

    missing = [c for c in path if c not in dfu.columns]
    if missing:
        txt = f"Missing columns for sunburst path: {missing}"
        return pd.DataFrame(), go.Figure().add_annotation(text=txt, showarrow=False)

    w = dfu.copy()
    for c in path:
        tmp = f"__{c}__list"
        w[tmp] = w[c].apply(_split_tokens).apply(lambda L: L if L else [None])
        # EXPLODE on the tmp list
        w = w.explode(tmp)
        # Drop the original column before renaming tmp → c (prevents duplicate names)
        w = w.drop(columns=[c])
        w = w.rename(columns={tmp: c})

    w = w.dropna(subset=path)
    for c in path:
        w[c] = w[c].astype(str).str.strip()
        w = w[w[c] != ""]
    if w.empty:
        return pd.DataFrame(), go.Figure().add_annotation(text="No data for sunburst path after tokenisation.", showarrow=False)

    w = _dedup_columns(w)  # final guard for Narwhals uniqueness

    tbl = w[path].value_counts().reset_index(name="count")
    fig = px.sunburst(w, path=path, title="Hierarchy")
    return tbl, fig


def _plot_argumentation_funnel(df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    params = dict(params or {})
    params.setdefault("funnel_dimension", "argumentation_logic")
    return _plot_category_funnel(df, params)

def _plot_keyword_network(df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    """
    Co-occurrence network for controlled vocabulary terms.
    Aliases old 'topic_phrase_network' to this function.
    """
    import itertools, math
    from collections import Counter

    # Prefer already-prepared list column; else split the string column.
    feature = "controlled_vocabulary_terms"
    list_col = "controlled_vocabulary_terms_list" if "controlled_vocabulary_terms_list" in df.columns else None

    def _tok(cell):
        if pd.isna(cell): return []
        if isinstance(cell, (list, tuple, set)): return [str(x).strip() for x in cell if str(x).strip()]
        return [t.strip() for t in str(cell).split(";") if t.strip()]

    # Build per-document term lists
    if list_col:
        term_lists = df[list_col].dropna().apply(_tok).tolist()
    elif feature in df.columns:
        term_lists = df[feature].dropna().apply(_tok).tolist()
    else:
        return pd.DataFrame(), go.Figure().add_annotation(text="No controlled vocabulary terms found.", showarrow=False)

    # Edge weights from within-document co-occurrences (unordered pairs)
    edge_counter = Counter()
    node_counter = Counter()
    for terms in term_lists:
        uniq = sorted(set(terms))
        node_counter.update(uniq)
        for a, b in itertools.combinations(uniq, 2):
            edge_counter[(a, b)] += 1

    if not edge_counter:
        return pd.DataFrame(), go.Figure().add_annotation(text="No co-occurrence among terms.", showarrow=False)

    # Keep top edges and corresponding nodes
    max_edges = int(params.get("max_edges", 60))
    top_edges = edge_counter.most_common(max_edges)
    nodes = sorted({n for (a, b), _ in top_edges for n in (a, b)}, key=lambda x: (-node_counter[x], x))

    # Circular layout (no networkx dependency)
    N = len(nodes)
    angles = [2 * math.pi * i / max(N, 1) for i in range(N)]
    coords = {n: (math.cos(ang), math.sin(ang)) for n, ang in zip(nodes, angles)}

    # Build edge traces
    edge_x, edge_y, edge_w = [], [], []
    for (a, b), w in top_edges:
        x0, y0 = coords[a]; x1, y1 = coords[b]
        edge_x += [x0, x1, None]; edge_y += [y0, y1, None]; edge_w.append(w)

    fig = go.Figure()

    # draw edges
    fig.add_trace(go.Scatter(
        x=edge_x, y=edge_y, mode="lines",
        line=dict(width=1), hoverinfo="skip", showlegend=False
    ))

    # draw nodes
    node_x = [coords[n][0] for n in nodes]
    node_y = [coords[n][1] for n in nodes]
    node_sizes = [8 + 2 * math.sqrt(node_counter[n]) for n in nodes]

    fig.add_trace(go.Scatter(
        x=node_x, y=node_y, mode="markers+text",
        text=nodes, textposition="top center",
        marker=dict(size=node_sizes, line=dict(width=1, color="black")),
        hovertext=[f"{n} (freq {node_counter[n]})" for n in nodes],
        hoverinfo="text", showlegend=False
    ))

    fig.update_layout(
        title="Keyword Co-occurrence Network (Controlled Vocabulary)",
        xaxis=dict(visible=False), yaxis=dict(visible=False),
        plot_bgcolor="white"
    )

    # Return edges table
    edge_df = pd.DataFrame(
        [(a, b, w) for (a, b), w in top_edges],
        columns=["term_a", "term_b", "weight"]
    )
    return edge_df, fig
def _plot_upset_multicode(df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    """
    UpSet-style intersections for one multi-code design dimension (e.g., 'methods').
    Graceful fallback if no intersections meet thresholds: progressively relax, then show top singles.
    """
    import itertools
    from collections import Counter

    dim = params.get("multicode_field") or params.get("upset_dimension") or "methods"
    top_k = int(params.get("top_k") or params.get("upset_top_k") or 8)
    min_degree = int(params.get("min_intersection") or params.get("min_degree") or 2)
    min_count = int(params.get("min_count") or 2)
    max_bars = int(params.get("max_bars") or 30)
    max_combo_size = int(params.get("max_combo_size") or 4)

    dfu = df.loc[:, ~df.columns.duplicated()].copy()
    if dim not in dfu.columns:
        return pd.DataFrame(), go.Figure().add_annotation(text=f"Missing column: {dim}", showarrow=False)

    def _tok(cell):
        if isinstance(cell, (list, tuple, set)):
            return [str(x).strip() for x in cell if str(x).strip()]
        if pd.isna(cell):
            return []
        return [t.strip() for t in str(cell).split(";") if t.strip()]

    lists = dfu[dim].apply(_tok)
    if lists.empty:
        return pd.DataFrame(), go.Figure().add_annotation(text=f"No data for {dim}.", showarrow=False)

    # top-k codes overall
    code_counts = Counter([c for L in lists for c in L])
    keep = {c for c, _ in code_counts.most_common(max(top_k, 1))}
    if not keep:
        return pd.DataFrame(), go.Figure().add_annotation(text=f"No frequent codes for {dim}.", showarrow=False)

    def _build_records(min_deg: int) -> list[tuple[tuple[str, ...], int]]:
        cc = Counter()
        for L in lists:
            S = sorted(set([c for c in L if c in keep]))
            if len(S) < min_deg: continue
            for r in range(min_deg, min(len(S), max_combo_size) + 1):
                for combo in itertools.combinations(S, r):
                    cc[combo] += 1
        return [(combo, cnt) for combo, cnt in cc.items() if cnt >= min_count]

    # try with requested min_degree, then relax to 2, then to 1 (singles)
    records = _build_records(min_degree)
    if not records and min_degree > 2:
        records = _build_records(2)
    if not records:
        # fallback to singles chart
        singles = pd.DataFrame(code_counts.most_common(top_k), columns=[dim, "Count"])
        fig = px.bar(
            singles, x=dim, y="Count",
            title=f"Top {len(singles)} {dim.replace('_',' ').title()} (no intersections met thresholds)"
        )
        fig.update_layout(xaxis_tickangle=-30)
        return singles, fig

    records.sort(key=lambda x: (-x[1], len(x[0]), x[0]))
    table = pd.DataFrame(
        [{"intersection": " ∩ ".join(combo), "Count": cnt, "k": len(combo)} for combo, cnt in records]
    )
    fig = px.bar(
        table.head(max_bars),
        x="intersection", y="Count", color="k",
        title=f"UpSet (multi-code): {dim.replace('_',' ').title()}",
        hover_data=["k"],
        labels={"intersection": "Combo (codes)", "k": "Size"}
    )
    fig.update_layout(xaxis_tickangle=-30)
    return table, fig


def _plot_category_funnel(df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    """
    Single generic funnel for ANY categorical design field.
    Makes the column *follow the UI selection* if explicit field missing:
      tries funnel_field → funnel_dimension → dimension1 → 'argumentation_logic'.
    Handles multi-coded cells (semicolon/list) and 'None' noise.
    """
    import re
    dfu = df.loc[:, ~df.columns.duplicated()].copy()

    col = (
        params.get("funnel_field")
        or params.get("funnel_dimension")
        or params.get("dimension1")
        or "argumentation_logic"
    )
    if col not in dfu.columns:
        return pd.DataFrame(), go.Figure().add_annotation(text=f"Missing column: {col}", showarrow=False)

    def _tok(cell):
        if isinstance(cell, (list, tuple, set)):
            return [str(x).strip() for x in cell if str(x).strip()]
        if pd.isna(cell):
            return []
        return [t.strip() for t in re.split(r"[;,]", str(cell)) if t.strip()]

    s = dfu[col].apply(_tok).explode()
    if s is None or s.empty:
        return pd.DataFrame(), go.Figure().add_annotation(text=f"No data for {col}.", showarrow=False)

    s = s.astype(str).str.strip()
    s = s[(s != "") & (~s.str.lower().eq("none"))]
    if s.empty:
        return pd.DataFrame(), go.Figure().add_annotation(text=f"No data for {col}.", showarrow=False)

    cnt = s.value_counts().reset_index()
    cnt.columns = [col, "Count"]
    fig = px.funnel(cnt, x="Count", y=col, title=f"Distribution: {col.replace('_',' ').title()}")
    return cnt, fig


def _evolution_line(df: pd.DataFrame, p: dict) -> Tuple[pd.DataFrame, go.Figure]:
    concept=p.get("concept_to_track","methodology")
    if 'year' not in df.columns or concept not in df.columns:
        return pd.DataFrame(), go.Figure().add_annotation(text="Need year & concept col", showarrow=False)
    sub=df[['year',concept]].dropna(); sub['year']=pd.to_numeric(pd.Series(sub['year']),errors='coerce').dropna().astype(int)
    if isinstance(sub[concept].iloc[0], list): sub=sub.explode(concept)
    sub=sub[sub[concept].astype(str).str.strip()!='']
    cnt=sub.groupby(['year',concept]).size().reset_index(name='Count')
    fig=px.line(cnt,x='year',y='Count',color=concept,markers=True,title=f"{concept} Evolution")
    return cnt, fig

def _bar_dimension(df: pd.DataFrame, p: dict)->Tuple[pd.DataFrame,go.Figure]:
    dim=p.get("dimension1","contribution_type")
    if dim not in df.columns: return pd.DataFrame(),go.Figure().add_annotation(text="Missing column",showarrow=False)
    cnt=df[dim].dropna().value_counts().reset_index()
    cnt.columns=[dim,'Count']
    fig=px.bar(cnt,x=dim,y='Count',title=f"Frequency of {dim.replace('_',' ').title()}",
               text='Count')
    return cnt,fig

def _animated_bar(df: pd.DataFrame, _:dict)->Tuple[pd.DataFrame,Figure]:
    if {'year','framework_model'}-set(df.columns):
        txt="Need year & framework_model"; return pd.DataFrame(), go.Figure().add_annotation(text=txt,showarrow=False)
    sub=df[['year','framework_model']].dropna()
    sub['year']=pd.to_numeric(pd.Series(sub['year']),errors='coerce').dropna().astype(int)
    sub=sub.groupby(['year','framework_model']).size().reset_index(name='Count')
    fig=px.bar(sub,x='framework_model',y='Count',color='framework_model',
               animation_frame='year',range_y=[0,sub['Count'].max()*1.1],
               title="Orientation Annual Counts (animated)")
    return sub, fig

def analyze_research_design_suite(
    df: "pd.DataFrame",
    params: dict,
    callback=(lambda m: None),
    *,
    slide_notes: bool = True,
    export: bool = False,
    return_payload: bool = False,
):
    """
    ###1. dispatch plot_type -> (table_df, fig)
    ###2. if export or return_payload: return {"slides":[...], "figs":[...], "imgs":[...], "tables":[...], "payload":{...}, "index":0}
         else: return (table_df, fig)
    """
    import json
    import pandas as pd
    import plotly.graph_objects as go
    from plotly.utils import PlotlyJSONEncoder

    def _fig_json_str(fig_obj: "go.Figure") -> str:
        return json.dumps(fig_obj.to_plotly_json(), cls=PlotlyJSONEncoder)

    preview_slides = [] if (export or return_payload) else None
    preview_figs = [] if (export or return_payload) else None
    preview_tables = [] if (export or return_payload) else None
    preview_imgs = [] if (export or return_payload) else None

    plot_type = str(params.get("plot_type", "crosstab_heatmap") or "crosstab_heatmap")

    if plot_type == "crosstab_heatmap":
        callback(f"[Research-Design] {plot_type} (X={params.get('dimension1')}, Y={params.get('dimension2')}) …")
    else:
        callback(f"[Research-Design] {plot_type} …")

    plot_map = {
        "crosstab_heatmap": _plot_design_crosstab_heatmap,
        "hierarchical_sunburst": _plot_design_hierarchical_sunburst,
        "argumentation_funnel": _plot_argumentation_funnel,
        "category_funnel": _plot_category_funnel,
        "keyword_network": _plot_keyword_network,
        "topic_phrase_network": _plot_keyword_network,
        "upset_multicode": _plot_upset_multicode,
    }

    if "_evolution_line" in globals():
        plot_map["concept_evolution"] = _evolution_line
    if "_bar_dimension" in globals():
        plot_map["dimension_bar"] = _bar_dimension
    if "_animated_bar" in globals():
        plot_map["animated_barrace"] = _animated_bar

    func = plot_map.get(plot_type)

    if func is None:
        txt = f"Unknown plot_type '{plot_type}'"
        callback(txt)
        fig = go.Figure()
        fig.add_annotation(text=txt, showarrow=False, xref="paper", yref="paper")
        table_df = pd.DataFrame()

        if preview_slides is not None:
            slide = {
                "title": f"Research Design — {plot_type}",
                "table_html": "<div>No table.</div>",
                "fig_json": _fig_json_str(fig),
                "notes": (f"Source: analyze_research_design_suite(plot_type='{plot_type}')." if slide_notes else ""),
            }
            preview_slides.append(slide)
            preview_figs.append(slide["fig_json"])
            preview_tables.append(slide["table_html"])
            return {
                "slides": preview_slides,
                "figs": preview_figs,
                "imgs": preview_imgs,
                "tables": preview_tables,
                "payload": {"section": "Research_design", "plot_type": plot_type, "n_rows": int(len(df)) if hasattr(df, "__len__") else 0},
                "index": 0,
            }

        return table_df, fig

    table_df, fig = func(df, params)

    if not isinstance(table_df, pd.DataFrame):
        table_df = pd.DataFrame()

    if fig is None:
        fig = go.Figure()
        fig.add_annotation(text=f"Plot '{plot_type}' not generated.", showarrow=False, xref="paper", yref="paper")

    payload = {
        "section": "Research_design",
        "plot_type": plot_type,
        "dimension1": params.get("dimension1"),
        "dimension2": params.get("dimension2"),
        "n_rows": int(len(df)) if hasattr(df, "__len__") else 0,
    }

    if preview_slides is not None:
        table_html = table_df.to_html(index=False, escape=True) if not table_df.empty else "<div>No table.</div>"
        notes = f"Source: analyze_research_design_suite(plot_type='{plot_type}')." if slide_notes else ""

        slide = {
            "title": f"Research Design — {plot_type}",
            "table_html": table_html,
            "fig_json": _fig_json_str(fig),
            "notes": notes,
        }

        preview_slides.append(slide)
        preview_figs.append(slide["fig_json"])
        preview_tables.append(slide["table_html"])

        return {
            "slides": preview_slides,
            "figs": preview_figs,
            "imgs": preview_imgs,
            "tables": preview_tables,
            "payload": payload,
            "index": 0,
        }

    return table_df, fig



def _profile_plot_scorecard(profile_df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    """Generates a scorecard of key metrics for the filtered entity data."""
    entity_type = params.get("feature_col", "Entity")
    entity_name = params.get("feature_value", "Profile")

    total_pubs = len(profile_df)
    cit_col = "citations" if "citations" in profile_df.columns else None
    cites = pd.to_numeric(profile_df[cit_col], errors="coerce").fillna(0.0) if cit_col else pd.Series([0.0] * len(profile_df))
    total_cites = int(float(cites.sum())) if not cites.empty else 0
    year_col = "year" if "year" in profile_df.columns else None
    years = pd.to_numeric(profile_df[year_col], errors="coerce") if year_col else pd.Series([pd.NA] * len(profile_df))
    first_year_val = years.min()
    last_year_val = years.max()
    active_years_str = (
        f"{int(first_year_val)} - {int(last_year_val)}" if pd.notna(first_year_val) and pd.notna(last_year_val) else "N/A"
    )

    fig = go.Figure()
    fig.add_trace(go.Indicator(mode="number", value=total_pubs, title={"text": "Total Publications"},
                               domain={'row': 0, 'column': 0}))
    fig.add_trace(go.Indicator(mode="number", value=total_cites, title={"text": "Total Citations"},
                               domain={'row': 0, 'column': 1}))

    table_df = pd.DataFrame(
        [{'Metric': 'Publications', 'Value': total_pubs}, {'Metric': 'Citations', 'Value': total_cites}])

    if entity_type == "authors":
        stats = _calculate_all_author_stats(profile_df)
        author_stats = stats[stats['Author'] == entity_name]
        h_index = author_stats['H-Index'].iloc[0] if not author_stats.empty else 0
        fig.add_trace(
            go.Indicator(mode="number", value=h_index, title={"text": "H-Index"}, domain={'row': 0, 'column': 2}))
        fig.update_layout(grid={'rows': 1, 'columns': 3, 'pattern': "independent"})
        table_df = author_stats[['TotalPublications', 'TotalCitations', 'H-Index']].T.reset_index()
        table_df.columns = ['Metric', 'Value']
    else:
        # Indicator doesn't support arbitrary text values reliably; use an annotation instead.
        fig.update_layout(grid={'rows': 1, 'columns': 2, 'pattern': "independent"})
        fig.add_annotation(
            text=f"Active period: {active_years_str}",
            xref="paper",
            yref="paper",
            x=0.98,
            y=0.5,
            xanchor="right",
            yanchor="middle",
            showarrow=False,
            font=dict(size=16),
        )

    fig.update_layout(title_text=f"Key Metrics for: {entity_name}")
    return table_df, fig


def _profile_plot_radar(profile_df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    """Creates a multi-dimensional 'fingerprint' radar chart."""
    entity_name = params.get("feature_value", "Profile")

    metrics = {
        'Productivity': len(profile_df),
        'Total Impact': profile_df['citations'].sum(),
        'Avg. Impact': profile_df['citations'].mean(),
        'Collaboration': profile_df['authors'].str.split(';').str.len().mean(),
        'Topical Diversity': profile_df['controlled_vocabulary_terms'].explode().nunique()
    }

    # Simple normalization (max value is 100) for visualization
    radar_data = {k: [v] for k, v in metrics.items()}
    radar_df = pd.DataFrame(radar_data)

    fig = go.Figure()
    fig.add_trace(go.Scatterpolar(
        r=list(metrics.values()),
        theta=list(metrics.keys()),
        fill='toself',
        name=entity_name
    ))
    fig.update_layout(polar=dict(radialaxis=dict(visible=True, range=[0, max(metrics.values()) * 1.1])),
                      showlegend=False, title=f"Performance Fingerprint for: {entity_name}")
    return radar_df.T.reset_index(), fig


def _profile_plot_pub_trend_line(profile_df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    """Generates a publication trend line for any entity."""
    entity_name = params.get("feature_value", "Profile")
    trend_df = profile_df.groupby('year').size().reset_index(name='Count')
    trend_df = trend_df.sort_values('year')
    fig = px.line(trend_df, x='year', y='Count', title=f"Publication Trend for: {entity_name}", markers=True)
    fig.update_layout(yaxis_title="Number of Publications", xaxis_title="Year")
    return trend_df, fig


def _profile_plot_thematic_treemap(profile_df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    """Generates a keyword treemap for an entity's publications."""
    entity_name = params.get("feature_value", "Profile")
    if 'controlled_vocabulary_terms' not in profile_df.columns or profile_df['controlled_vocabulary_terms'].dropna().empty:
        return pd.DataFrame(), go.Figure().add_annotation(text="No keyword data available for this entity.")

    keyword_counts = Counter(kw for kws_list in profile_df['controlled_vocabulary_terms'].dropna() for kw in kws_list)
    top_keywords_df = pd.DataFrame(keyword_counts.most_common(40), columns=['Keyword', 'Frequency'])

    fig = px.treemap(top_keywords_df, path=[px.Constant(f"Thematic Focus for {entity_name}"), 'Keyword'],
                     values='Frequency', color='Frequency', color_continuous_scale='YlGnBu',
                     title=f"Thematic Focus for: {entity_name}")
    return top_keywords_df, fig


def _profile_plot_top_collaborators_bar(profile_df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    """Generates a bar chart of top co-authors for the profiled entity."""
    entity_name = params.get("feature_value")

    collaborators = Counter()
    for items in profile_df['authors'].dropna():
        items = [i.strip() for i in items.split(';')]
        # Exclude the entity itself if it's an author profile
        if params.get("feature_col") == 'authors':
            items_without_self = [item for item in items if item != entity_name]
            collaborators.update(items_without_self)
        else:
            collaborators.update(items)

    if not collaborators:
        return pd.DataFrame(), go.Figure().add_annotation(text=f"No collaborators found for {entity_name}.")

    top_collabs_df = pd.DataFrame(collaborators.most_common(15), columns=['Collaborator', 'Count'])
    fig = px.bar(top_collabs_df, x='Collaborator', y='Count', text_auto=True,
                 title=f"Top Co-Authors for: {entity_name}")
    fig.update_xaxes(categoryorder="total descending")
    return top_collabs_df, fig


def _profile_plot_pub_type_donut(profile_df: pd.DataFrame, params: dict) -> tuple[pd.DataFrame, go.Figure]:
    """Generates a donut chart of publication types for an entity."""
    entity_name = params.get("feature_value")
    counts = profile_df['item_type'].value_counts().reset_index()
    counts.columns = ['Publication Type', 'Count']

    fig = px.pie(counts, names='Publication Type', values='Count',
                 title=f"Publication Types for: {entity_name}", hole=0.4)
    fig.update_traces(textinfo='percent+label')
    return counts, fig





# ========================================================
def analyze_feature_profile(
    df: "pd.DataFrame",
    params: dict,
    callback,
    *,
    slide_notes: bool = True,
    export: bool = False,
    return_payload: bool = False,
):
    """
    ###1. filter df to the selected feature/value context
    ###2. dispatch plot_type -> (table_df, fig)
    ###3. if export or return_payload: return {"slides":[...], "figs":[...], "imgs":[...], "tables":[...], "payload":{...}, "index":0}
         else: return (table_df, fig)
    """
    import json
    import logging
    import re
    import pandas as pd
    import plotly.graph_objects as go
    from plotly.utils import PlotlyJSONEncoder

    def _fig_json_str(fig_obj: "go.Figure") -> str:
        return json.dumps(fig_obj.to_plotly_json(), cls=PlotlyJSONEncoder)

    preview_slides = [] if (export or return_payload) else None
    preview_figs = [] if (export or return_payload) else None
    preview_tables = [] if (export or return_payload) else None
    preview_imgs = [] if (export or return_payload) else None

    def _fail(msg: str, *, title_suffix: str = "Error"):
        fig0 = go.Figure()
        fig0.add_annotation(text=msg, showarrow=False, xref="paper", yref="paper")
        fig0.update_layout(xaxis_visible=False, yaxis_visible=False)

        table0 = pd.DataFrame()
        if preview_slides is not None:
            slide = {
                "title": f"Profile — {title_suffix}",
                "table_html": "<div>No table.</div>",
                "fig_json": _fig_json_str(fig0),
                "notes": (msg if slide_notes else ""),
            }
            preview_slides.append(slide)
            preview_figs.append(slide["fig_json"])
            preview_tables.append(slide["table_html"])
            return {
                "slides": preview_slides,
                "figs": preview_figs,
                "imgs": preview_imgs,
                "tables": preview_tables,
                "payload": {"section": "Profiles", "error": msg},
                "index": 0,
            }
        return table0, fig0

    feature_col = params.get("feature_col")
    feature_value = params.get("feature_value")
    plot_type = params.get("plot_type", "scorecard")

    if not feature_col or feature_value is None or not plot_type:
        return _fail("Please select a feature and a value to profile.", title_suffix="Missing selection")

    callback(f"Generating '{plot_type}' profile for '{feature_value}' in '{feature_col}'...")

    if df is None or not isinstance(df, pd.DataFrame) or df.empty:
        return _fail("No data loaded.", title_suffix="No data")

    if feature_col not in df.columns:
        return _fail(f"Feature column '{feature_col}' not found.", title_suffix="Missing column")

    series = df[feature_col].dropna()
    if series.empty:
        return _fail(f"No data in '{feature_col}' column.", title_suffix="Empty column")

    first = series.iloc[0]
    is_list_col = isinstance(first, list)

    if is_list_col:
        want = str(feature_value).strip()
        mask = series.apply(lambda x: want in [str(i).strip() for i in (x or [])])
        profile_df = df.loc[mask.reindex(df.index, fill_value=False)].copy()
    else:
        pat = re.escape(str(feature_value))
        mask = df[feature_col].astype(str).str.contains(pat, na=False, regex=True)
        profile_df = df.loc[mask].copy()

    if profile_df.empty:
        return _fail(f"No data found for '{feature_value}' in '{feature_col}'.", title_suffix="No matches")

    plot_function_map = {
        "scorecard": _profile_plot_scorecard,
        "performance_radar": (lambda d, p: _profile_plot_radar(d, p, df)),
        "pub_trend_line": _profile_plot_pub_trend_line,
        "thematic_treemap": _profile_plot_thematic_treemap,
        "pub_type_donut": _profile_plot_pub_type_donut,
        "co_author_network": None,
        "top_authors_bar": _profile_plot_top_collaborators_bar,
        "career_trajectory": _plot_author_career_trajectory,
        "co_occurrence_network": analyze_keyword_cooccurrence_network,
    }

    plot_func = plot_function_map.get(str(plot_type))

    if not plot_func:
        return _fail(f"Plot type '{plot_type}' is not yet implemented.", title_suffix="Unsupported plot")

    try:
        table_df, fig = plot_func(profile_df, params)
    except Exception as e:
        logging.error(f"Error generating plot '{plot_type}': {e}", exc_info=True)
        return _fail(f"Error generating plot: {e}", title_suffix="Plot error")

    if not isinstance(table_df, pd.DataFrame):
        table_df = pd.DataFrame()

    if fig is None:
        fig = go.Figure()
        fig.add_annotation(text=f"Plot '{plot_type}' not generated.", showarrow=False, xref="paper", yref="paper")
        fig.update_layout(xaxis_visible=False, yaxis_visible=False)

    payload = {
        "section": "Profiles",
        "feature_col": str(feature_col),
        "feature_value": str(feature_value),
        "plot_type": str(plot_type),
        "n_rows_total": int(len(df)),
        "n_rows_profile": int(len(profile_df)),
    }

    if preview_slides is not None:
        table_html = table_df.to_html(index=False, escape=True) if not table_df.empty else "<div>No table.</div>"
        notes = f"Source: analyze_feature_profile(plot_type='{plot_type}', feature_col='{feature_col}')." if slide_notes else ""

        slide = {
            "title": f"Profile — {feature_value} — {plot_type}",
            "table_html": table_html,
            "fig_json": _fig_json_str(fig),
            "notes": notes,
        }

        preview_slides.append(slide)
        preview_figs.append(slide["fig_json"])
        preview_tables.append(slide["table_html"])

        return {
            "slides": preview_slides,
            "figs": preview_figs,
            "imgs": preview_imgs,
            "tables": preview_tables,
            "payload": payload,
            "index": 0,
        }

    return table_df, fig


def _plot_ngram_bar(results_df: pd.DataFrame, ngram_n_size: int, top_n_to_plot: int) -> go.Figure:
    """Reusable bar chart for top N-grams."""
    plot_df = results_df.head(top_n_to_plot)
    title = f"Top {len(plot_df)} {ngram_n_size}-grams"
    fig = px.bar(
        plot_df,
        x="Ngram",
        y="Frequency",
        text_auto=True,
        title=title,
        color="Frequency",
        color_continuous_scale=px.colors.sequential.Purples,
        template='plotly_white'
    )
    fig.update_traces(textposition='outside')
    fig.update_xaxes(categoryorder="total descending", tickangle=-30)
    return fig


def _plot_ngram_cooccurrence_network(
    df: pd.DataFrame,
    source: str,
    n: int,
    min_cooc: int,
    max_nodes: int,
    zotero_client=None,
    cache_name=None
) -> tuple[pd.DataFrame, go.Figure]:
    """
    Builds and returns (edge_table_df, network_fig) for n-gram co-occurrence.
    - df: your document DataFrame
    - source: one of 'abstract','title','controlled_vocabulary_terms','fulltext'
    - n: n-gram size
    - min_cooc: minimum co-occurrence threshold
    - max_nodes: max number of nodes to display
    """
    # --- 1. Extract per-document n-grams ---
    docs_ngrams = []
    for _, row in df.iterrows():
        if source == 'abstract':
            text = str(row.get('abstract','') or '')
        elif source == 'title':
            text = str(row.get('title','') or '')
        elif source == 'controlled_vocabulary_terms':
            kws = row.get('controlled_vocabulary_terms', [])
            text = ' '.join(k for k in kws if isinstance(k, str))
        elif source == 'fulltext':
            tokens = get_text_corpus_from_df(
                pd.DataFrame([row]),
                source,
                cache_name,
                zotero_client,
                None
            )
            text = ' '.join(tokens)
        else:
            text = ''
        # your existing helper to split into n-grams:
        doc_ngrams = _get_ngrams_from_text(text, n)
        docs_ngrams.append(list(set(doc_ngrams)))  # unique per doc

    # --- 2. Global frequency & select top candidates ---
    global_counts = Counter(ng for doc in docs_ngrams for ng in doc)
    # if fewer than max_nodes exist, we'll just use all
    candidates = [ng for ng, _ in global_counts.most_common(max_nodes)]
    if len(candidates) < 2:
        # nothing to co-occur
        empty_fig = go.Figure().add_annotation(text="Not enough distinct n-grams", showarrow=False)
        return pd.DataFrame(columns=["Ngram1","Ngram2","Cooccurrences"]), empty_fig

    # --- 3. Build co-occurrence counts among chosen candidates ---
    cooc = Counter()
    for doc in docs_ngrams:
        present = sorted(set(doc) & set(candidates))
        for a, b in itertools.combinations(present, 2):
            cooc[(a,b)] += 1

    # filter by min_cooc
    edges = [(a, b, c) for (a,b), c in cooc.items() if c >= min_cooc]
    if not edges:
        empty_fig = go.Figure().add_annotation(text="No co-occurrences meet threshold", showarrow=False)
        return pd.DataFrame(columns=["Ngram1","Ngram2","Cooccurrences"]), empty_fig

    edge_df = pd.DataFrame(edges, columns=["Ngram1","Ngram2","Cooccurrences"])

    # --- 4. Build NetworkX graph and prune to top max_nodes by weighted degree ---
    G = nx.Graph()
    for a, b, w in edges:
        G.add_edge(a, b, weight=w)
    # annotate node frequency
    for node in G.nodes():
        G.nodes[node]['freq'] = global_counts[node]

    # if too many nodes, keep only the highest-degree ones
    if G.number_of_nodes() > max_nodes:
        deg = G.degree(weight='weight')
        top_nodes = [n for n,_ in sorted(deg, key=lambda x: x[1], reverse=True)[:max_nodes]]
        G = G.subgraph(top_nodes).copy()

    # --- 5. Layout & Plotly traces ---
    pos = nx.spring_layout(G, k=0.5/np.sqrt(max(1, G.number_of_nodes())), iterations=30, weight='weight', seed=42)

    # edges
    edge_x, edge_y = [], []
    for u, v in G.edges():
        x0, y0 = pos[u]; x1, y1 = pos[v]
        edge_x += [x0, x1, None]
        edge_y += [y0, y1, None]

    edge_trace = go.Scatter(
        x=edge_x, y=edge_y,
        mode='lines',
        line=dict(width=0.8, color='#888'),
        hoverinfo='none'
    )

    # nodes
    node_x, node_y, node_text, node_size, node_color = [], [], [], [], []
    for node in G.nodes():
        x, y = pos[node]
        node_x.append(x)
        node_y.append(y)
        freq = G.nodes[node]['freq']
        deg_w = G.degree(node, weight='weight')
        node_text.append(f"{node}<br>Freq: {freq}<br>Cooc‐deg: {deg_w}")
        # size ∝ log(freq)
        node_size.append(5 + np.log1p(freq) * 4)
        node_color.append(deg_w)

    node_trace = go.Scatter(
        x=node_x, y=node_y,
        mode='markers+text',
        text=[n for n in G.nodes()],
        textposition='top center',
        hovertext=node_text,
        hoverinfo='text',
        marker=dict(
            showscale=True,
            colorscale='YlGnBu',
            color=node_color,
            size=node_size,
            colorbar=dict(title='Cooc-deg'),
            line_width=1
        )
    )

    fig = go.Figure(
        data=[edge_trace, node_trace],
        layout=go.Layout(
            title=f"{n}-gram Co-occurrence Network",
            showlegend=False,
            hovermode='closest',
            margin=dict(l=20, r=20, t=40, b=20),
            xaxis=dict(showgrid=False, zeroline=False, showticklabels=False),
            yaxis=dict(showgrid=False, zeroline=False, showticklabels=False)
        )
    )

    return edge_df, fig


def _plot_ngram_heatmap_by_authors(df: pd.DataFrame, data_source: str, ngram_n: int,
                                   top_n_ngrams: int, zotero_client, cache_name: str) -> tuple[pd.DataFrame, go.Figure]:
    """Reusable heatmap of n-gram frequency by authors."""
    author_map = defaultdict(Counter)
    for _, row in df.iterrows():
        authors = row.get('authors', 'Unknown')
        # extract text
        if data_source == 'abstract':
            text = str(row.get('abstract', '') or '')
        elif data_source == 'title':
            text = str(row.get('title', '') or '')
        elif data_source == 'controlled_vocabulary_terms':
            kws = row.get('controlled_vocabulary_terms', [])
            text = ' '.join(k for k in kws if isinstance(k, str))
        elif data_source == 'fulltext':
            tokens = get_text_corpus_from_df(
                pd.DataFrame([row]), data_source, cache_name, zotero_client, None
            )
            text = ' '.join(tokens)
        else:
            text = ''
        if not text:
            continue
        author_map[authors].update(_get_ngrams_from_text(text, ngram_n))

    # fix overall_counts
    overall_counts = Counter()
    for cnt in author_map.values():
        overall_counts.update(cnt)

    if not overall_counts:
        fig = go.Figure().add_annotation(text="No data for heatmap.", showarrow=False)
        return pd.DataFrame(columns=['Author','Ngram','Frequency']), fig

    top_ngrams = [ng for ng, _ in overall_counts.most_common(top_n_ngrams)]

    heat_rows = []
    for author, cnt in author_map.items():
        row = {'Author': author}
        for ng in top_ngrams:
            row[ng] = cnt.get(ng, 0)
        heat_rows.append(row)
    df_heat = pd.DataFrame(heat_rows).set_index('Author')
    df_heat = df_heat.loc[(df_heat > 0).any(axis=1)]

    if df_heat.empty:
        fig = go.Figure().add_annotation(text="No data for heatmap.", showarrow=False)
        table_df = pd.DataFrame(columns=['Author','Ngram','Frequency'])
        return table_df, fig

    fig = px.imshow(
        df_heat,
        text_auto=True,
        aspect='auto',
        template='plotly_white',
        title=f"N-gram Frequency by Author ({ngram_n}-grams)",
        color_continuous_scale='Blues',
        labels={'x':'N-gram','y':'Author','color':'Frequency'}
    )
    fig.update_xaxes(side='bottom', tickangle=-45)
    table_df = df_heat.reset_index().melt(id_vars=['Author'], var_name='Ngram', value_name='Frequency')
    table_df = table_df[table_df['Frequency']>0]
    return table_df, fig


def analyze_ngrams(
    df: "pd.DataFrame",
    params: dict,
    progress_callback=None,
    zotero_client_for_pdf=None,
    collection_name_for_cache=None,
    *,
    slide_notes: bool = True,
    export: bool = False,
    return_payload: bool = False,
):
    """
    ###1. compute overall n-gram frequency table (when applicable)
    ###2. dispatch plot_type -> (table_df, fig)
    ###3. if export or return_payload: return {"slides":[...], "figs":[...], "imgs":[...], "tables":[...], "payload":{...}, "index":0}
         else: return (table_df, fig)
    """
    import json
    import pandas as pd
    import plotly.graph_objects as go
    from plotly.utils import PlotlyJSONEncoder
    from collections import Counter

    preview_slides = [] if (export or return_payload) else None
    preview_figs = [] if (export or return_payload) else None
    preview_tables = [] if (export or return_payload) else None
    preview_imgs = [] if (export or return_payload) else None

    def _cb(msg: str):
        if progress_callback:
            progress_callback(f"Ngrams: {msg}")

    def _fig_json_str(fig_obj: "go.Figure") -> str:
        return json.dumps(fig_obj.to_plotly_json(), cls=PlotlyJSONEncoder)

    def _empty(msg: str):
        fig0 = go.Figure()
        fig0.add_annotation(text=msg, showarrow=False, xref="paper", yref="paper")
        fig0.update_layout(xaxis_visible=False, yaxis_visible=False)
        return pd.DataFrame(), fig0

    if df is None or not isinstance(df, pd.DataFrame) or df.empty:
        table0, fig0 = _empty("No data.")
        if preview_slides is not None:
            slide = {
                "title": "N-grams — No data",
                "table_html": "<div>No data.</div>",
                "fig_json": _fig_json_str(fig0),
                "notes": ("No data loaded." if slide_notes else ""),
            }
            preview_slides.append(slide)
            preview_figs.append(slide["fig_json"])
            preview_tables.append(slide["table_html"])
            return {
                "slides": preview_slides,
                "figs": preview_figs,
                "imgs": preview_imgs,
                "tables": preview_tables,
                "payload": {"section": "Words_and_topics", "plot_type": None, "data_source": None},
                "index": 0,
            }
        return table0, fig0

    source = params.get("data_source", "abstract")
    ngram_n = int(params.get("ngram_n", 2))
    top_n = int(params.get("top_n_ngrams", 20))
    plot_type = params.get("plot_type", "bar_chart")
    min_freq = int(params.get("min_frequency", 2))

    _cb(f"Computing n={ngram_n} from {source} (plot={plot_type})")

    results_df = pd.DataFrame(columns=["Ngram", "Frequency"])

    if plot_type in ("bar_chart", "ngram_frequency_heatmap"):
        all_ngrams = []
        for _, row in df.iterrows():
            if source == "abstract":
                text = str(row.get("abstract", "") or "")
            elif source == "title":
                text = str(row.get("title", "") or "")
            elif source == "controlled_vocabulary_terms":
                kws = row.get("controlled_vocabulary_terms", [])
                if isinstance(kws, (list, tuple, set)):
                    text = " ".join(k for k in kws if isinstance(k, str))
                else:
                    text = str(kws or "")
            elif source == "fulltext":
                tokens = get_text_corpus_from_df(
                    pd.DataFrame([row]),
                    source,
                    collection_name_for_cache,
                    zotero_client_for_pdf,
                    None,
                )
                text = " ".join(tokens) if isinstance(tokens, list) else str(tokens or "")
            else:
                text = ""

            all_ngrams.extend(_get_ngrams_from_text(text, ngram_n))

        counts = Counter(all_ngrams)
        filtered = {ng: c for ng, c in counts.items() if int(c) >= min_freq}
        if filtered:
            results_df = (
                pd.DataFrame(filtered.items(), columns=["Ngram", "Frequency"])
                .sort_values("Frequency", ascending=False)
                .reset_index(drop=True)
            )

    fig = None
    table_df = results_df

    if plot_type == "bar_chart":
        fig = _plot_ngram_bar(results_df, ngram_n, top_n)

    elif plot_type == "ngram_cooccurrence_network":
        table_df, fig = _plot_ngram_cooccurrence_network(
            df,
            source,
            ngram_n,
            int(params.get("min_ngram_cooccurrence", 2)),
            int(params.get("max_nodes_for_ngram_network", 30)),
            zotero_client_for_pdf,
            collection_name_for_cache,
        )

    elif plot_type == "ngram_evolution_time_series":
        ngrams_by_year = _get_ngram_data_by_year(
            df,
            source,
            ngram_n,
            zotero_client_for_pdf,
            collection_name_for_cache,
            progress_callback,
        )

        if not ngrams_by_year:
            table_df, fig = _empty("No data for evolution.")
        else:
            records = []
            for year, ng_list in ngrams_by_year.items():
                for ng in (ng_list or []):
                    records.append({"year": year, "controlled_vocabulary_terms": ng})
            df_for_trends = pd.DataFrame.from_records(records)

            wot_params = {
                "data_source": "controlled_vocabulary_terms",
                "num_top_words": int(params.get("num_top_ngrams_for_evolution", 7)),
            }

            out = analyze_words_over_time(
                df_for_trends,
                wot_params,
                progress_callback,
                zotero_client_for_pdf,
                collection_name_for_cache,
                export=export,
                return_payload=return_payload,
            )

            if isinstance(out, dict) and "slides" in out:
                return out

            table_df, fig = out

    elif plot_type == "ngram_frequency_heatmap":
        table_df, fig = _plot_ngram_heatmap_by_authors(
            df,
            source,
            ngram_n,
            top_n,
            zotero_client_for_pdf,
            collection_name_for_cache,
        )

    else:
        _cb(f"Unknown plot type '{plot_type}', defaulting to bar chart.")
        fig = _plot_ngram_bar(results_df, ngram_n, top_n)

    if fig is None:
        fig = go.Figure()
        fig.add_annotation(
            text=f"Plot for '{plot_type}' could not be generated.",
            showarrow=False,
            xref="paper",
            yref="paper",
        )
        fig.update_layout(xaxis_visible=False, yaxis_visible=False)

    payload = {
        "section": "Words_and_topics",
        "plot_type": str(plot_type),
        "data_source": str(source),
        "ngram_n": int(ngram_n),
        "top_n": int(top_n),
        "min_frequency": int(min_freq),
        "n_rows": int(len(df)),
    }

    if preview_slides is not None:
        table_html = table_df.to_html(index=False, escape=True) if isinstance(table_df, pd.DataFrame) and not table_df.empty else "<div>No table.</div>"
        notes = f"Source: analyze_ngrams(plot_type='{plot_type}', n={ngram_n}, source='{source}')." if slide_notes else ""

        slide = {
            "title": f"N-grams — {plot_type}",
            "table_html": table_html,
            "fig_json": _fig_json_str(fig),
            "notes": notes,
        }

        preview_slides.append(slide)
        preview_figs.append(slide["fig_json"])
        preview_tables.append(slide["table_html"])

        return {
            "slides": preview_slides,
            "figs": preview_figs,
            "imgs": preview_imgs,
            "tables": preview_tables,
            "payload": payload,
            "index": 0,
        }

    return table_df, fig
def analyze_pdf_keywords_and_trends(
    df: "pd.DataFrame",
    params: dict,
    progress_callback=None,
    zotero_client=None,
    *,
    slide_notes: bool = True,
    export: bool = False,
    return_payload: bool = False,
):
    """
    ###1. scan selectable corpus for keywords (with caching) and build a snippets table
    ###2. dispatch plot_type -> either a plotly figure, or a snippets list / message
    ###3. if export or return_payload: return {"slides":[...], "figs":[...], "imgs":[...], "tables":[...], "payload":{...}, "index":0}
         else: return (table_df, visual_obj)
    """
    import json
    import logging
    import os
    import re
    import hashlib
    import pickle
    import pandas as pd
    import plotly.express as px
    import plotly.graph_objects as go
    from plotly.utils import PlotlyJSONEncoder
    from collections import defaultdict, Counter

    preview_slides = [] if (export or return_payload) else None
    preview_figs = [] if (export or return_payload) else None
    preview_tables = [] if (export or return_payload) else None
    preview_imgs = [] if (export or return_payload) else None

    def _cb(msg: str):
        if progress_callback:
            progress_callback(f"PDFScan: {msg}")
        logging.info(f"PDFScan: {msg}")

    def _fig_json_str(fig_obj: "go.Figure") -> str:
        return json.dumps(fig_obj.to_plotly_json(), cls=PlotlyJSONEncoder)

    def _as_message_fig(text: str) -> "go.Figure":
        fig0 = go.Figure()
        fig0.add_annotation(text=text, showarrow=False, xref="paper", yref="paper")
        fig0.update_layout(xaxis_visible=False, yaxis_visible=False, margin=dict(l=40, r=40, t=40, b=40))
        return fig0

    def _payload_return(slides, payload, index=0):
        return {
            "slides": slides,
            "figs": preview_figs or [],
            "imgs": preview_imgs or [],
            "tables": preview_tables or [],
            "payload": payload,
            "index": index,
        }

    user_keywords = [
        str(k).lower().strip()
        for k in (params.get("keywords_to_scan_in_pdf", []) or [])
        if str(k).strip()
    ]
    plot_type = str(params.get("plot_type", "snippets_only") or "snippets_only")
    corpus_source = str(params.get("corpus_source", "full_text") or "full_text").strip().lower()
    collection_name = str(params.get("collection_name_for_cache", "unknown_coll_pdfscan") or "unknown_coll_pdfscan")
    safe_collection = re.sub(r"[^\w\-.]+", "_", collection_name)

    focus_kw = None
    if plot_type == "single_keyword_focus_ngrams":
        raw_focus = params.get("focus_keyword_for_details", "")
        if isinstance(raw_focus, str) and raw_focus.strip():
            focus_kw = raw_focus.lower().strip()
        elif user_keywords:
            focus_kw = user_keywords[0]
        if not focus_kw:
            _cb("No focus keyword for n-grams; falling back to snippets_only.")
            plot_type = "snippets_only"

    _cb(f"Scan start: keywords={user_keywords[:3]} corpus='{corpus_source}' view='{plot_type}' collection='{safe_collection}'")

    payload = {
        "section": "PDF_keywords_and_trends",
        "collection": safe_collection,
        "corpus_source": corpus_source,
        "plot_type": plot_type,
        "keywords": user_keywords[:200],
        "focus_keyword": focus_kw,
    }

    if not user_keywords:
        table0 = pd.DataFrame()
        visual0 = {"type": "message", "content": "No keywords provided to scan."}
        if preview_slides is not None:
            fig0 = _as_message_fig("No keywords provided to scan.")
            slide = {
                "title": "PDF Keyword Scan — No keywords",
                "table_html": "<div>No keywords provided.</div>",
                "fig_json": _fig_json_str(fig0),
                "notes": ("No keywords provided." if slide_notes else ""),
            }
            preview_slides.append(slide)
            preview_figs.append(slide["fig_json"])
            preview_tables.append(slide["table_html"])
            return _payload_return(preview_slides, payload, 0)
        return table0, visual0

    if df is None or not isinstance(df, pd.DataFrame) or df.empty or "key" not in df.columns:
        table0 = pd.DataFrame()
        visual0 = {"type": "message", "content": "'key' column missing or data empty."}
        if preview_slides is not None:
            fig0 = _as_message_fig("'key' column missing or data empty.")
            slide = {
                "title": "PDF Keyword Scan — No data",
                "table_html": "<div>'key' column missing or data empty.</div>",
                "fig_json": _fig_json_str(fig0),
                "notes": ("No scan performed." if slide_notes else ""),
            }
            preview_slides.append(slide)
            preview_figs.append(slide["fig_json"])
            preview_tables.append(slide["table_html"])
            return _payload_return(preview_slides, payload, 0)
        return table0, visual0

    if corpus_source == "full_text" and zotero_client is None:
        table0 = pd.DataFrame()
        visual0 = {"type": "message", "content": "Zotero client not configured for full-text scan."}
        if preview_slides is not None:
            fig0 = _as_message_fig("Zotero client not configured for full-text scan.")
            slide = {
                "title": "PDF Keyword Scan — Zotero not configured",
                "table_html": "<div>Zotero client not configured for full-text scan.</div>",
                "fig_json": _fig_json_str(fig0),
                "notes": ("Full-text scan requires Zotero client." if slide_notes else ""),
            }
            preview_slides.append(slide)
            preview_figs.append(slide["fig_json"])
            preview_tables.append(slide["table_html"])
            return _payload_return(preview_slides, payload, 0)
        return table0, visual0

    sorted_kw = tuple(sorted(set(user_keywords)))
    kw_hash = hashlib.md5(str(sorted_kw).encode("utf-8")).hexdigest()[:8]
    cache_version = "v1.6_pdfscan_data"
    cache_filename = f"pdfscan_{safe_collection}_kw{kw_hash}_{corpus_source}_{cache_version}.pkl"
    cache_file = MAIN_APP_CACHE_DIR / cache_filename

    all_snippets = []
    hits_per_year = defaultdict(lambda: defaultdict(int))
    docs_mentioning = defaultdict(list)

    if cache_file.exists():
        _cb(f"Loading cache: {cache_file.name}")
        try:
            with open(cache_file, "rb") as f:
                cached = pickle.load(f)
            all_snippets = cached.get("all_snippets", []) or []
            loaded_hits = cached.get("hits_per_year_detailed", {}) or {}
            for yr, kwc in loaded_hits.items():
                hits_per_year[int(yr)] = defaultdict(int, kwc)
            loaded_docs = cached.get("docs_mentioning_keyword", {}) or {}
            for kw, docs in loaded_docs.items():
                docs_mentioning[str(kw).lower().strip()].extend(docs or [])
            if all_snippets or hits_per_year or docs_mentioning:
                _cb(f"Cache loaded: snippets={len(all_snippets)}")
        except Exception as e:
            _cb(f"Cache load failed: {e}; recomputing.")
            all_snippets = []
            hits_per_year = defaultdict(lambda: defaultdict(int))
            docs_mentioning = defaultdict(list)

    if not all_snippets and not hits_per_year and not docs_mentioning:
        _cb("Cache miss; scanning corpus.")
        df_work = df.copy()
        if "year" in df_work.columns:
            df_work["year_numeric"] = pd.to_numeric(df_work["year"], errors="coerce")
        else:
            df_work["year_numeric"] = pd.NA
            _cb("Warning: 'year' column missing.")

        per_item_cache = MAIN_APP_CACHE_DIR / ".pdfscan_flat_text"
        per_item_cache.mkdir(parents=True, exist_ok=True)

        total = len(df_work)
        for i, (_idx, row) in enumerate(df_work.iterrows(), 1):
            key = row.get("key")
            year_val = row.get("year_numeric")
            year_int = int(year_val) if pd.notna(year_val) else None
            title = str(row.get("title", f"Item_{key}" if key else "Unknown Title"))
            authors = str(row.get("authors", "N/A"))

            text = _choose_corpus_text(
                row,
                corpus_source,
                zotero_client=zotero_client,
                mistral_api_key=os.environ.get("MISTRAL_API_KEY"),
                per_item_cache_dir=per_item_cache,
                progress_cb=_cb,
            )

            if (i % 20 == 0) or (i == total):
                _cb(f"Processed {i}/{total} items")

            if not text:
                continue

            matches = _scan_text_for_keywords(text, user_keywords, window=160)
            if not matches:
                continue

            per_doc_kw = defaultdict(lambda: {"hits": 0, "plain_text_snippets": [], "rich_snippets": []})

            for m in matches:
                kw_found = str(m.get("keyword_found", "")).lower().strip()
                snippet_plain = str(m.get("snippet", "") or "").strip()

                snippet_rich = {
                    "paragraph_context": snippet_plain,
                    "original_paragraph": snippet_plain,
                    "keyword_found": m.get("keyword_found", ""),
                    "page_number": None,
                    "source_item_key": key,
                    "source_pdf_path": None,
                    "source_bib_header": f"{authors} ({year_int if year_int is not None else 'N/A'})",
                    "source_title": title,
                }
                all_snippets.append(snippet_rich)

                if year_int is not None:
                    hits_per_year[year_int]["_overall_"] += 1
                    hits_per_year[year_int][kw_found] += 1

                per_doc_kw[kw_found]["hits"] += 1
                if snippet_plain:
                    per_doc_kw[kw_found]["plain_text_snippets"].append(snippet_plain)
                per_doc_kw[kw_found]["rich_snippets"].append(snippet_rich)

            for kw_norm, dval in per_doc_kw.items():
                docs_mentioning[str(kw_norm).lower().strip()].append(
                    {
                        "key": key,
                        "title": title,
                        "authors": authors,
                        "year": year_int,
                        "hits_in_this_doc": dval["hits"],
                        "plain_text_snippets": dval["plain_text_snippets"],
                        "rich_snippets": dval["rich_snippets"],
                    }
                )

        if all_snippets or hits_per_year or docs_mentioning:
            try:
                with open(cache_file, "wb") as f:
                    pickle.dump(
                        {
                            "all_snippets": all_snippets,
                            "hits_per_year_detailed": {yr: dict(kwc) for yr, kwc in hits_per_year.items()},
                            "docs_mentioning_keyword": {k: v for k, v in docs_mentioning.items()},
                        },
                        f,
                    )
                _cb(f"Saved cache: {cache_file.name}")
            except Exception as e:
                _cb(f"Cache save failed: {e}")

    df_result = pd.DataFrame(all_snippets) if all_snippets else pd.DataFrame()
    visual = {"type": "message", "content": "Select display or no visual."}

    if plot_type == "overall_hits_trend":
        oy = {year: d.get("_overall_", 0) for year, d in hits_per_year.items() if int(d.get("_overall_", 0)) > 0}
        if oy and ("year" in df.columns) and df["year"].notna().any():
            pub_trend = pd.to_numeric(pd.Series(df["year"]), errors="coerce").dropna().astype(int).value_counts().sort_index()
            hit_trend = pd.Series(oy).sort_index()

            if (not pub_trend.empty) and (not hit_trend.empty):
                fig = go.Figure()
                fig.add_trace(go.Scatter(x=pub_trend.index.astype(str), y=pub_trend.values, name="Publications", mode="lines+markers"))
                fig.add_trace(go.Scatter(x=hit_trend.index.astype(str), y=hit_trend.values, name="Total Keyword Hits", mode="lines+markers", yaxis="y2"))
                fig.update_layout(
                    title_text=f"Publications vs Total Keyword Hits ({corpus_source})",
                    xaxis_title="Year",
                    yaxis_title="Publications",
                    yaxis2=dict(title="Keyword Hits", overlaying="y", side="right", showgrid=False, zeroline=False),
                    xaxis_type="category",
                    legend_title_text="Metric",
                )
                visual = fig
                df_result = pd.DataFrame(
                    {
                        "Year": pub_trend.index.astype(int),
                        "Publications": pub_trend.values.astype(int),
                    }
                ).merge(
                    pd.DataFrame({"Year": hit_trend.index.astype(int), "Keyword Hits": hit_trend.values.astype(int)}),
                    on="Year",
                    how="outer",
                ).fillna(0)
            else:
                visual = {"type": "message", "content": "Not enough data for overall trend."}
        else:
            visual = {"type": "message", "content": "No yearly keyword hits for trend."}

    elif plot_type == "single_keyword_focus_ngrams" and focus_kw:
        def _contains_phrase_in_ngram(ngram_tuple, phrase_tokens) -> bool:
            L = len(phrase_tokens)
            if L == 0:
                return False
            if L == 1:
                return phrase_tokens[0] in ngram_tuple
            N = len(ngram_tuple)
            for j in range(0, N - L + 1):
                if list(ngram_tuple[j : j + L]) == phrase_tokens:
                    return True
            return False

        docs_data = docs_mentioning.get(focus_kw, [])

        if (not docs_data) and (" " in focus_kw) and all_snippets:
            phrase = focus_kw.lower().strip()
            grouped = defaultdict(
                lambda: {
                    "key": None,
                    "title": "",
                    "authors": "",
                    "year": None,
                    "hits_in_this_doc": 0,
                    "plain_text_snippets": [],
                    "rich_snippets": [],
                }
            )
            for snip in all_snippets:
                para_plain = (
                    snip.get("original_paragraph")
                    or _normalise_and_strip_html(snip.get("paragraph_context") or "")
                    or ""
                )
                if phrase in str(para_plain).lower():
                    dkey = snip.get("source_item_key") or snip.get("key") or snip.get("source_pdf_path")
                    g = grouped[dkey]
                    g["key"] = dkey
                    g["title"] = snip.get("source_title", "")
                    g["authors"] = snip.get("source_bib_header", "")
                    g["year"] = snip.get("year")
                    g["plain_text_snippets"].append(para_plain)
                    g["rich_snippets"].append(snip)
                    g["hits_in_this_doc"] += 1
            docs_data = list(grouped.values())

        if not docs_data:
            visual = {"type": "message", "content": f"Focus keyword '{focus_kw}' not found in snippets."}
        else:
            from nltk import ngrams

            n = int(params.get("ngram_n_for_focus_kw", 2))
            min_freq = int(params.get("min_freq_for_focus_kw_ngrams", 2))
            top_n = int(params.get("top_n_ngrams_for_focus_kw_plot", 15))

            phrase_tokens = tokenize_for_ngram_context_refined(focus_kw)
            if not phrase_tokens:
                phrase_tokens = [t for t in focus_kw.lower().split() if t]

            n_for_gen = max(n, len(phrase_tokens))

            ngram_list = []
            for doc in docs_data:
                for plain in (doc.get("plain_text_snippets", []) or []):
                    tokens = tokenize_for_ngram_context_refined(plain)
                    if len(tokens) < n_for_gen:
                        continue
                    for tup in ngrams(tokens, n_for_gen):
                        if _contains_phrase_in_ngram(tup, phrase_tokens):
                            ngram_list.append(" ".join(tup))

            cnt = Counter(ngram_list)
            df_ngrams = (
                pd.DataFrame([(k, v) for k, v in cnt.items() if int(v) >= min_freq], columns=["Ngram", "Frequency"])
                .sort_values("Frequency", ascending=False)
                .head(top_n)
                .reset_index(drop=True)
            )

            if not df_ngrams.empty:
                fig = px.bar(
                    df_ngrams,
                    x="Ngram",
                    y="Frequency",
                    title=f"Top contextual {n_for_gen}-grams containing '{focus_kw}' ({corpus_source})",
                    text_auto=True,
                )
                fig.update_traces(textposition="outside")
                fig.update_xaxes(categoryorder="total descending", tickangle=-45)
                visual = fig
                df_result = df_ngrams
            else:
                visual = {"type": "message", "content": f"No frequent n-grams containing '{focus_kw}'."}

    elif all_snippets:
        visual = all_snippets
    else:
        visual = {"type": "message", "content": "No keyword snippets found."}

    _cb(f"Return: visual={type(visual).__name__} table_rows={len(df_result)}")

    if preview_slides is not None:
        if isinstance(visual, go.Figure):
            fig = visual
        else:
            msg = visual.get("content") if isinstance(visual, dict) else "Snippets available."
            fig = _as_message_fig(str(msg))

        table_html = df_result.to_html(index=False, escape=True) if isinstance(df_result, pd.DataFrame) and not df_result.empty else "<div>No table.</div>"
        notes = ""
        if slide_notes:
            notes = f"Source: analyze_pdf_keywords_and_trends(plot_type='{plot_type}', corpus='{corpus_source}')."

        slide = {
            "title": f"PDF Keyword Scan — {plot_type}",
            "table_html": table_html,
            "fig_json": _fig_json_str(fig),
            "notes": notes,
        }

        preview_slides.append(slide)
        preview_figs.append(slide["fig_json"])
        preview_tables.append(slide["table_html"])

        return _payload_return(preview_slides, payload, 0)

    return df_result, visual

def analyze_wordanalysis_dispatcher(
    df: "pd.DataFrame",
    params: dict,
    progress_callback=None,
    *,
    zotero_client_for_pdf=None,
    collection_name_for_cache=None,
    slide_notes: bool = True,
    export: bool = False,
    return_payload: bool = False,
):
    """
    ###1. dispatch to one word-analysis function based on params["analysis"] or params["plot_type"]
    ###2. normalize outputs into a preview/export payload schema:
         - slides: [{title, table_html, fig_json?, img?, notes?}, ...]
         - figs: [fig_json_str, ...]
         - imgs: [data_uri_or_file_uri, ...]
         - tables: [table_html, ...]
         - payload: {section, analysis, plot_type, data_source, ...}
         - index: 0
    ###3. never use kaleido; for preview, send plotly JSON; for word-cloud, send data URI image
    """
    import json
    import logging
    import pandas as pd
    import plotly.graph_objects as go
    from plotly.utils import PlotlyJSONEncoder

    preview_slides = [] if (export or return_payload) else None
    preview_figs = [] if (export or return_payload) else None
    preview_imgs = [] if (export or return_payload) else None
    preview_tables = [] if (export or return_payload) else None

    def _cb(msg: str):
        if progress_callback:
            progress_callback(f"WordAnalysis: {msg}")
        logging.info(f"WordAnalysis: {msg}")

    def _fig_json_str(fig_obj: "go.Figure") -> str:
        return json.dumps(fig_obj.to_plotly_json(), cls=PlotlyJSONEncoder)

    def _as_message_fig(text: str) -> "go.Figure":
        fig0 = go.Figure()
        fig0.add_annotation(text=text, showarrow=False, xref="paper", yref="paper")
        fig0.update_layout(xaxis_visible=False, yaxis_visible=False, margin=dict(l=40, r=40, t=40, b=40))
        return fig0

    def _table_html(df0: "pd.DataFrame") -> str:
        if df0 is None or (isinstance(df0, pd.DataFrame) and df0.empty):
            return "<div>No table.</div>"
        if not isinstance(df0, pd.DataFrame):
            return "<div>No table.</div>"
        return df0.to_html(index=False, escape=True)

    def _payload_return(slides, payload, index=0):
        return {
            "slides": slides,
            "figs": preview_figs or [],
            "imgs": preview_imgs or [],
            "tables": preview_tables or [],
            "payload": payload,
            "index": index,
        }

    def _append_slide(title: str, table_df: "pd.DataFrame", fig=None, img=None, notes: str = ""):
        slide = {"title": title, "table_html": _table_html(table_df), "notes": notes}

        if fig is not None and isinstance(fig, go.Figure):
            slide["fig_json"] = _fig_json_str(fig)
            preview_figs.append(slide["fig_json"])

        if img is not None:
            slide["img"] = img
            preview_imgs.append(img)

        preview_tables.append(slide["table_html"])
        preview_slides.append(slide)

    def _dispatch_key() -> str:
        a = params.get("analysis", "")
        if isinstance(a, str) and a.strip():
            return a.strip().lower()

        pt = params.get("plot_type", "")
        if not isinstance(pt, str):
            return "freq_words"
        pt = pt.strip().lower()

        if pt in {"bar_vertical", "bar_horizontal"}:
            return "freq_words"
        if pt in {"treemap"}:
            return "treemap"
        if pt in {"word_cloud", "cloud", "wordcloud"}:
            return "word_cloud"
        if pt in {"heatmap", "frequency_heatmap", "keyword_heatmap"}:
            return "heatmap"
        if pt in {"cooccurrence_network", "co-occurrence network", "cooccurrence", "network_graph", "network"}:
            return "cooccurrence_network"
        if pt in {"words_over_time", "time_series", "trend"}:
            return "words_over_time"
        return "freq_words"

    analysis = _dispatch_key()
    plot_type = params.get("plot_type", "")
    data_source = params.get("data_source", "controlled_vocabulary_terms")

    payload = {
        "section": "WordAnalysis",
        "analysis": analysis,
        "plot_type": plot_type,
        "data_source": data_source,
        "collection_name_for_cache": collection_name_for_cache,
    }

    if df is None or not isinstance(df, pd.DataFrame) or df.empty:
        _cb("No data.")
        if preview_slides is not None:
            fig0 = _as_message_fig("No data.")
            _append_slide(
                "Word Analysis — No data",
                pd.DataFrame(),
                fig=fig0,
                img=None,
                notes=("No data loaded." if slide_notes else ""),
            )
            return _payload_return(preview_slides, payload, 0)
        return pd.DataFrame(), None

    _cb(f"Dispatching analysis='{analysis}' (source='{data_source}').")

    table_df = pd.DataFrame()
    fig_obj = None
    img_obj = None

    if analysis == "freq_words":
        table_df, fig_obj = analyze_most_frequent_words(
            df,
            params,
            progress_callback=progress_callback,
            zotero_client_for_pdf=zotero_client_for_pdf,
            collection_name_for_cache=collection_name_for_cache,
        )
        if fig_obj is None:
            fig_obj = _as_message_fig("No figure returned.")
        title = f"Word Analysis — Frequent Words — {params.get('plot_type', 'bar')}"
        notes = ""
        if slide_notes:
            notes = f"Source: analyze_most_frequent_words(data_source='{data_source}')."
        if preview_slides is not None:
            _append_slide(title, table_df, fig=fig_obj, img=None, notes=notes)
            return _payload_return(preview_slides, payload, 0)
        return table_df, fig_obj

    if analysis == "word_cloud":
        table_df, wc_visual = analyze_word_cloud(
            df,
            params,
            progress_callback=progress_callback,
            zotero_client_for_pdf=zotero_client_for_pdf,
            collection_name_for_cache=collection_name_for_cache,
        )

        if isinstance(wc_visual, go.Figure):
            fig_obj = wc_visual
        elif isinstance(wc_visual, str) and wc_visual.strip():
            if wc_visual.strip().startswith("data:image/"):
                img_obj = wc_visual.strip()
            else:
                img_obj = wc_visual.strip()
        elif isinstance(wc_visual, dict) and wc_visual.get("type") == "message":
            fig_obj = _as_message_fig(str(wc_visual.get("content", "No visual.")))
        else:
            fig_obj = _as_message_fig("Word cloud returned no visual.")

        title = "Word Analysis — Word Cloud"
        notes = ""
        if slide_notes:
            notes = f"Source: analyze_word_cloud(data_source='{data_source}')."

        if preview_slides is not None:
            _append_slide(title, table_df, fig=fig_obj, img=img_obj, notes=notes)
            return _payload_return(preview_slides, payload, 0)

        return table_df, (img_obj if img_obj is not None else fig_obj)

    if analysis == "treemap":
        table_df, fig_obj = analyze_word_treemap(
            df,
            params,
            progress_callback=progress_callback,
            zotero_client_for_pdf=zotero_client_for_pdf,
            collection_name_for_cache=collection_name_for_cache,
        )
        if fig_obj is None:
            fig_obj = _as_message_fig("No treemap figure returned.")
        title = "Word Analysis — Treemap"
        notes = ""
        if slide_notes:
            notes = f"Source: analyze_word_treemap(data_source='{data_source}')."
        if preview_slides is not None:
            _append_slide(title, table_df, fig=fig_obj, img=None, notes=notes)
            return _payload_return(preview_slides, payload, 0)
        return table_df, fig_obj

    if analysis == "words_over_time":
        table_df, fig_obj = analyze_words_over_time(
            df,
            params,
            progress_callback=progress_callback,
            zotero_client_for_pdf=zotero_client_for_pdf,
            collection_name_for_cache=collection_name_for_cache,
        )
        if fig_obj is None:
            fig_obj = _as_message_fig("No time-series figure returned.")
        title = "Word Analysis — Words Over Time"
        notes = ""
        if slide_notes:
            notes = f"Source: analyze_words_over_time(data_source='{data_source}')."
        if preview_slides is not None:
            _append_slide(title, table_df, fig=fig_obj, img=None, notes=notes)
            return _payload_return(preview_slides, payload, 0)
        return table_df, fig_obj

    if analysis == "heatmap":
        table_df, fig_obj = analyze_keyword_heatmap(
            df,
            params,
            progress_callback=progress_callback,
            zotero_client=zotero_client_for_pdf,
            collection_name_for_cache=collection_name_for_cache,
        )
        if fig_obj is None:
            fig_obj = _as_message_fig("No heatmap figure returned.")
        title = "Word Analysis — Frequency Heatmap"
        notes = ""
        if slide_notes:
            notes = f"Source: analyze_keyword_heatmap(data_source='{data_source}')."
        if preview_slides is not None:
            _append_slide(title, table_df, fig=fig_obj, img=None, notes=notes)
            return _payload_return(preview_slides, payload, 0)
        return table_df, fig_obj

    if analysis == "cooccurrence_network":
        p2 = dict(params or {})
        p2["plot_type"] = "network_graph"
        table_df, fig_obj = analyze_keyword_cooccurrence_network(
            df,
            p2,
            progress_callback=progress_callback,
            zotero_client_for_pdf=zotero_client_for_pdf,
            collection_name_for_cache=collection_name_for_cache,
        )
        if fig_obj is None:
            fig_obj = _as_message_fig("No network figure returned.")
        title = "Word Analysis — Co-occurrence Network"
        notes = ""
        if slide_notes:
            notes = f"Source: analyze_keyword_cooccurrence_network(data_source='{data_source}')."
        if preview_slides is not None:
            _append_slide(title, table_df, fig=fig_obj, img=None, notes=notes)
            return _payload_return(preview_slides, payload, 0)
        return table_df, fig_obj

    _cb(f"Unknown analysis '{analysis}'.")
    fig_obj = _as_message_fig(f"Unknown analysis '{analysis}'.")
    if preview_slides is not None:
        _append_slide("Word Analysis — Unknown", pd.DataFrame(), fig=fig_obj, img=None, notes="")
        return _payload_return(preview_slides, payload, 0)

    return pd.DataFrame(), fig_obj
def authors_overview(
    prs,
    df: "pd.DataFrame",
    collection_name: str = "Collection",
    *,
    slide_notes: bool = True,
    impact_plot_types: list[str] | None = None,
    collaboration_plot_types: list[str] | None = None,
    trends_plot_types: list[str] | None = None,
    export: bool = False,
    return_payload: bool = False,
    progress_callback=None,
):
    """
    ###1. build a compact Authors Overview slide-deck (table slide + figure slide per plot)
    ###2. reuse existing analysis dispatchers: analyze_author_impact/collaboration/trends + analyze_lotkas_law
    ###3. if export or return_payload: return {"slides":[...], "figs":[...], "index":0}; else return None

    NOTE (fix for your Plotly 'Scatter: title' crash):
      - Figure slides now store fig_json as a JSON string of fig.to_plotly_json()
      - We do NOT store wrapper dicts in fig_json
      - This ensures go.Figure(json.loads(fig_json)) reconstructs cleanly in _fig_spec_to_thumb_url
    """
    import io
    import json

    import numpy as np
    import pandas as pd
    import plotly.graph_objects as go
    from plotly.utils import PlotlyJSONEncoder
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN

    preview_slides = [] if (export or return_payload) else None
    preview_figs = [] if (export or return_payload) else None

    FONT_FAMILY = "Segoe UI"
    PLOT_W, PLOT_H = 1200, 650

    INCH = float(Inches(1))
    SLIDE_W_IN = float(prs.slide_width) / INCH
    SLIDE_H_IN = float(prs.slide_height) / INCH

    OUTER_MARGIN_X = 0.6
    OUTER_MARGIN_Y = 0.6
    TITLE_HEIGHT_IN = 0.7

    def _cb(msg: str):
        if progress_callback:
            progress_callback(f"AUTHORS: {msg}")

    def _new_blank_slide_with_title(title_text: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(
            Inches(0.0),
            Inches(0.15),
            Inches(SLIDE_W_IN),
            Inches(TITLE_HEIGHT_IN),
        )
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        return slide

    def _set_notes(slide, text: str):
        if slide_notes and text:
            slide.notes_slide.notes_text_frame.text = text

    def _add_table_slide(title_text: str, table_df: "pd.DataFrame | None", notes: str = ""):
        slide = _new_blank_slide_with_title(title_text)
        _set_notes(slide, notes)

        if table_df is None or type(table_df) is not pd.DataFrame or table_df.empty:
            tb = slide.shapes.add_textbox(
                Inches(OUTER_MARGIN_X),
                Inches(TITLE_HEIGHT_IN + 2.2),
                Inches(SLIDE_W_IN - 2 * OUTER_MARGIN_X),
                Inches(1.0),
            )
            p = tb.text_frame.paragraphs[0]
            p.text = "No data available."
            p.font.name = FONT_FAMILY
            p.font.size = Pt(16)
            p.alignment = PP_ALIGN.CENTER

            if preview_slides is not None:
                preview_slides.append({"title": title_text, "table_html": "<div>No data available.</div>", "notes": notes})
            return slide

        df_small = table_df.head(25).copy()

        max_w = SLIDE_W_IN - 2 * OUTER_MARGIN_X
        max_h = SLIDE_H_IN - (TITLE_HEIGHT_IN + 2 * OUTER_MARGIN_Y)

        rows = max(1, len(df_small) + 1)
        cols = max(1, df_small.shape[1])

        def _len_inch_est(s: object) -> float:
            t = "" if pd.isna(s) else str(s)
            t = " ".join(t.split())
            return 0.085 * max(4, len(t)) + 0.14

        col_needs = []
        for j, col in enumerate(df_small.columns):
            header_need = _len_inch_est(col)
            longest_need = header_need
            for i in range(len(df_small)):
                longest_need = max(longest_need, min(_len_inch_est(df_small.iloc[i, j]), 3.8))
            col_needs.append(longest_need)

        MIN_W = 0.9
        MAX_W = 3.6

        col_needs = [min(MAX_W, max(MIN_W, x)) for x in col_needs]
        total_need = float(np.sum(col_needs))

        if total_need > (max_w - 0.2):
            scale = (max_w - 0.2) / total_need
            col_widths_in = [max(MIN_W, x * scale) for x in col_needs]
        else:
            col_widths_in = col_needs[:]

        table_w = float(np.sum(col_widths_in))
        est_row_h = min(0.42, max(0.26, max_h / (rows + 1)))
        table_h = min(max_h, rows * est_row_h + 0.35)

        left_in = (SLIDE_W_IN - table_w) / 2.0
        top_in = TITLE_HEIGHT_IN + ((SLIDE_H_IN - TITLE_HEIGHT_IN) - table_h) / 2.0

        table = slide.shapes.add_table(
            rows, cols, Inches(left_in), Inches(top_in), Inches(table_w), Inches(table_h)
        ).table

        for j, w in enumerate(col_widths_in):
            table.columns[j].width = Emu(Inches(w))

        for j, col in enumerate(df_small.columns):
            cell = table.cell(0, j)
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(col)
            p.font.bold = True
            p.font.name = FONT_FAMILY
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER
            cell.margin_left = cell.margin_right = Emu(900)
            cell.margin_top = Emu(600)
            cell.margin_bottom = Emu(600)

        for i, (_, r) in enumerate(df_small.iterrows(), start=1):
            for j, col in enumerate(df_small.columns):
                cell = table.cell(i, j)
                txt = "" if pd.isna(r[col]) else " ".join(str(r[col]).split())
                tf = cell.text_frame
                tf.clear()
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = txt
                p.font.name = FONT_FAMILY
                p.font.size = Pt(11)
                p.alignment = PP_ALIGN.LEFT
                cell.margin_left = Emu(800)
                cell.margin_right = Emu(800)
                cell.margin_top = Emu(500)
                cell.margin_bottom = Emu(500)

        if preview_slides is not None:
            preview_slides.append({"title": title_text, "table_html": df_small.to_html(index=False, escape=True), "notes": notes})

        return slide

    def _add_figure_slide(title_text: str, fig, notes: str = ""):
        """
        ###1. accept plotly Figure; accept tuple(fig, ...) by taking first element
        ###2. render plotly -> png into slide
        ###3. if preview enabled, store fig_json as a JSON string of the plotly figure spec
        """
        import pandas as pd
        import plotly.graph_objects as go

        if type(fig) is tuple and fig:
            fig = fig[0]

        if type(fig) is pd.DataFrame:
            return _add_table_slide(title_text.replace("— Figure", "— Table"), fig, notes=notes)

        if type(fig) is not go.Figure:
            slide = _new_blank_slide_with_title(title_text)
            _set_notes(slide, notes)

            tb = slide.shapes.add_textbox(
                Inches(OUTER_MARGIN_X),
                Inches(TITLE_HEIGHT_IN + 2.2),
                Inches(SLIDE_W_IN - 2 * OUTER_MARGIN_X),
                Inches(1.2),
            )
            p = tb.text_frame.paragraphs[0]
            p.text = f"Figure object is not Plotly (type={type(fig).__name__})."
            p.font.name = FONT_FAMILY
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER

            if preview_slides is not None:
                preview_slides.append({"title": title_text, "table_html": "<div>Figure unavailable.</div>", "notes": notes or ""})
            return slide

        slide = _new_blank_slide_with_title(title_text)
        _set_notes(slide, notes)

        content_w = SLIDE_W_IN - 2 * OUTER_MARGIN_X
        content_h = SLIDE_H_IN - (TITLE_HEIGHT_IN + 2 * OUTER_MARGIN_Y)
        if export:
            try:
                png = fig.to_image(format="png", width=PLOT_W, height=PLOT_H, scale=2)
                ar = (PLOT_W / PLOT_H) if PLOT_H else 1.7778

                target_w = content_w
                target_h = target_w / ar
                if target_h > content_h:
                    target_h = content_h
                    target_w = target_h * ar
                target_w *= 0.98
                target_h *= 0.98

                left_in = (SLIDE_W_IN - target_w) / 2.0
                top_in = TITLE_HEIGHT_IN + (SLIDE_H_IN - TITLE_HEIGHT_IN - target_h) / 2.0

                slide.shapes.add_picture(
                    io.BytesIO(png),
                    Inches(left_in),
                    Inches(top_in),
                    width=Inches(target_w),
                    height=Inches(target_h),
                )
            except Exception as exc:
                # Kaleido/Chrome is optional for preview, but can block export. Degrade gracefully and log.
                try:
                    print(f"[visualise][warn] plotly.to_image failed ({title_text}): {exc}")
                except Exception:
                    pass
                tb = slide.shapes.add_textbox(
                    Inches(OUTER_MARGIN_X),
                    Inches(TITLE_HEIGHT_IN + 2.2),
                    Inches(SLIDE_W_IN - 2 * OUTER_MARGIN_X),
                    Inches(1.2),
                )
                p = tb.text_frame.paragraphs[0]
                p.text = "Figure export unavailable (Plotly image rendering failed)."
                p.font.name = FONT_FAMILY
                p.font.size = Pt(14)
                p.alignment = PP_ALIGN.CENTER
        else:
            # Preview mode: the Electron UI renders Plotly from fig_json; don't require Kaleido/Chrome.
            tb = slide.shapes.add_textbox(
                Inches(OUTER_MARGIN_X),
                Inches(TITLE_HEIGHT_IN + 2.2),
                Inches(SLIDE_W_IN - 2 * OUTER_MARGIN_X),
                Inches(1.2),
            )
            p = tb.text_frame.paragraphs[0]
            p.text = "Preview-only slide (figure rendered in the app)."
            p.font.name = FONT_FAMILY
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER

        if preview_slides is not None:
            fig_json_str = json.dumps(fig.to_plotly_json(), cls=PlotlyJSONEncoder)
            preview_slides.append({"title": title_text, "fig_json": fig_json_str, "table_html": "", "notes": notes or ""})
            preview_figs.append(fig_json_str)

        return slide

    if impact_plot_types is None:
        impact_plot_types = [
            "h_index_bar",
            "productivity_bar",
            "citations_total_bar",
            "productivity_vs_impact_scatter",
        ]

    if collaboration_plot_types is None:
        collaboration_plot_types = [
            "co_authorship_network",
        ]

    if trends_plot_types is None:
        trends_plot_types = [
            "top_authors_production_over_time",
            "overall_citation_trend",
        ]

    if df is None or type(df) is not pd.DataFrame or df.empty:
        if preview_slides is not None:
            return {
                "slides": [{"title": "Authors overview", "table_html": "<div>No data loaded.</div>", "notes": ""}],
                "figs": [],
                "index": 0,
            }
        return None

    _cb("Building impact slides")
    for pt in impact_plot_types:
        params = {"plot_type": pt}
        stats_df, fig = analyze_author_impact(df, params, progress_callback or (lambda *_: None))

        table_title = f"{collection_name} · Author Impact — {pt} — Table"
        fig_title = f"{collection_name} · Author Impact — {pt} — Figure"

        _add_table_slide(table_title, stats_df, notes=f"Source: analyze_author_impact(plot_type='{pt}').")
        _add_figure_slide(fig_title, fig, notes=f"Figure: analyze_author_impact(plot_type='{pt}').")

    _cb("Building collaboration slides")
    for pt in collaboration_plot_types:
        params = {"plot_type": pt}
        tab_df, fig = analyze_author_collaboration(df, params, progress_callback or (lambda *_: None))

        table_title = f"{collection_name} · Author Collaboration — {pt} — Table"
        fig_title = f"{collection_name} · Author Collaboration — {pt} — Figure"

        _add_table_slide(table_title, tab_df, notes=f"Source: analyze_author_collaboration(plot_type='{pt}').")
        _add_figure_slide(fig_title, fig, notes=f"Source: analyze_author_collaboration(plot_type='{pt}').")

    _cb("Building trend slides")
    for pt in trends_plot_types:
        params = {"plot_type": pt}
        tab_df, fig = analyze_author_trends(df, params, progress_callback or (lambda *_: None))

        table_title = f"{collection_name} · Author Trends — {pt} — Table"
        fig_title = f"{collection_name} · Author Trends — {pt} — Figure"

        _add_table_slide(table_title, tab_df, notes=f"Source: analyze_author_trends(plot_type='{pt}').")
        _add_figure_slide(fig_title, fig, notes=f"Figure: analyze_author_trends(plot_type='{pt}').")

    _cb("Building Lotka slides")
    lotka_df, lotka_fig = analyze_lotkas_law(df, {}, progress_callback)

    _add_table_slide(
        f"{collection_name} · Lotka's Law — Table",
        lotka_df,
        notes="Source: analyze_lotkas_law(). Table is the observed productivity distribution.",
    )
    _add_figure_slide(
        f"{collection_name} · Lotka's Law — Figure",
        lotka_fig,
        notes="Source: analyze_lotkas_law(). Log-log plot of observed proportions and (if fitted) predicted curves.",
    )

    _cb("Authors overview completed")

    if preview_slides is not None:
        return {"slides": preview_slides, "figs": preview_figs or [], "index": 0}

    return None



def add_institution_section(
    prs,
    df: "pd.DataFrame",
    collection_name: str = "Collection",
    *,
    slide_notes: bool = True,
    top_n_outlets: int = 20,
    export: bool = False,
    return_payload: bool = False,
    progress_callback=None,
):
    """
    ###1. build an Institutions Overview slide-deck in the same style as authors_overview
    ###2. create 2 slides:
        - Major institutional actors (publication_outlet) — Table
        - Major institutional actors (publication_outlet) — Figure
    ###3. if export or return_payload: return {"slides":[...], "figs":[...], "index":0}; else return None

    Assumptions (brittle by design):
      - df is a pandas.DataFrame (may be empty; empty returns a 1-slide payload if preview enabled).
      - plotly figures support .to_image().
      - publication outlet column exists: 'publication_outlet' or fallback 'publicationTitle' or 'publisher'.
    """
    import io
    import json
    import re

    import numpy as np
    import pandas as pd
    import plotly.express as px
    from plotly.utils import PlotlyJSONEncoder
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN

    preview_slides = [] if (export or return_payload) else None
    preview_figs = [] if (export or return_payload) else None

    FONT_FAMILY = "Segoe UI"
    PLOT_W, PLOT_H = 1200, 650

    INCH = float(Inches(1))
    SLIDE_W_IN = float(prs.slide_width) / INCH
    SLIDE_H_IN = float(prs.slide_height) / INCH

    OUTER_MARGIN_X = 0.6
    OUTER_MARGIN_Y = 0.6
    TITLE_HEIGHT_IN = 0.7

    def _cb(msg: str):
        if progress_callback:
            progress_callback(f"INST: {msg}")

    def _new_blank_slide_with_title(title_text: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(
            Inches(0.0),
            Inches(0.15),
            Inches(SLIDE_W_IN),
            Inches(TITLE_HEIGHT_IN),
        )
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        return slide

    def _set_notes(slide, text: str):
        if slide_notes and text:
            slide.notes_slide.notes_text_frame.text = text

    def _add_table_slide(title_text: str, table_df: "pd.DataFrame | None", notes: str = ""):
        slide = _new_blank_slide_with_title(title_text)
        _set_notes(slide, notes)

        if table_df is None or type(table_df) is not pd.DataFrame or table_df.empty:
            tb = slide.shapes.add_textbox(
                Inches(OUTER_MARGIN_X),
                Inches(TITLE_HEIGHT_IN + 2.2),
                Inches(SLIDE_W_IN - 2 * OUTER_MARGIN_X),
                Inches(1.0),
            )
            p = tb.text_frame.paragraphs[0]
            p.text = "No data available."
            p.font.name = FONT_FAMILY
            p.font.size = Pt(16)
            p.alignment = PP_ALIGN.CENTER

            if preview_slides is not None:
                preview_slides.append({"title": title_text, "table_html": "<div>No data available.</div>", "notes": notes})
            return slide

        df_small = table_df.head(25).copy()

        max_w = SLIDE_W_IN - 2 * OUTER_MARGIN_X
        max_h = SLIDE_H_IN - (TITLE_HEIGHT_IN + 2 * OUTER_MARGIN_Y)

        rows = max(1, len(df_small) + 1)
        cols = max(1, df_small.shape[1])

        def _len_inch_est(s: object) -> float:
            t = "" if pd.isna(s) else str(s)
            t = " ".join(t.split())
            return 0.085 * max(4, len(t)) + 0.14

        col_needs = []
        for j, col in enumerate(df_small.columns):
            header_need = _len_inch_est(col)
            longest_need = header_need
            for i in range(len(df_small)):
                longest_need = max(longest_need, min(_len_inch_est(df_small.iloc[i, j]), 3.8))
            col_needs.append(longest_need)

        MIN_W = 0.9
        MAX_W = 3.6

        col_needs = [min(MAX_W, max(MIN_W, x)) for x in col_needs]
        total_need = float(np.sum(col_needs))

        if total_need > (max_w - 0.2):
            scale = (max_w - 0.2) / total_need
            col_widths_in = [max(MIN_W, x * scale) for x in col_needs]
        else:
            col_widths_in = col_needs[:]

        table_w = float(np.sum(col_widths_in))
        est_row_h = min(0.42, max(0.26, max_h / (rows + 1)))
        table_h = min(max_h, rows * est_row_h + 0.35)

        left_in = (SLIDE_W_IN - table_w) / 2.0
        top_in = TITLE_HEIGHT_IN + ((SLIDE_H_IN - TITLE_HEIGHT_IN) - table_h) / 2.0

        table = slide.shapes.add_table(
            rows, cols, Inches(left_in), Inches(top_in), Inches(table_w), Inches(table_h)
        ).table

        for j, w in enumerate(col_widths_in):
            table.columns[j].width = Emu(Inches(w))

        for j, col in enumerate(df_small.columns):
            cell = table.cell(0, j)
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(col)
            p.font.bold = True
            p.font.name = FONT_FAMILY
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER
            cell.margin_left = cell.margin_right = Emu(900)
            cell.margin_top = Emu(600)
            cell.margin_bottom = Emu(600)

        for i, (_, r) in enumerate(df_small.iterrows(), start=1):
            for j, col in enumerate(df_small.columns):
                cell = table.cell(i, j)
                txt = "" if pd.isna(r[col]) else " ".join(str(r[col]).split())
                tf = cell.text_frame
                tf.clear()
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = txt
                p.font.name = FONT_FAMILY
                p.font.size = Pt(11)
                p.alignment = PP_ALIGN.LEFT
                cell.margin_left = Emu(800)
                cell.margin_right = Emu(800)
                cell.margin_top = Emu(500)
                cell.margin_bottom = Emu(500)

        if preview_slides is not None:
            preview_slides.append({"title": title_text, "table_html": df_small.to_html(index=False, escape=True), "notes": notes})

        return slide

    def _add_figure_slide(title_text: str, fig, notes: str = ""):
        slide = _new_blank_slide_with_title(title_text)
        _set_notes(slide, notes)

        content_w = SLIDE_W_IN - 2 * OUTER_MARGIN_X
        content_h = SLIDE_H_IN - (TITLE_HEIGHT_IN + 2 * OUTER_MARGIN_Y)
        if export:
            try:
                png = fig.to_image(format="png", width=PLOT_W, height=PLOT_H, scale=2)
                ar = (PLOT_W / PLOT_H) if PLOT_H else 1.7778

                target_w = content_w
                target_h = target_w / ar
                if target_h > content_h:
                    target_h = content_h
                    target_w = target_h * ar
                target_w *= 0.98
                target_h *= 0.98

                left_in = (SLIDE_W_IN - target_w) / 2.0
                top_in = TITLE_HEIGHT_IN + (SLIDE_H_IN - TITLE_HEIGHT_IN - target_h) / 2.0

                slide.shapes.add_picture(
                    io.BytesIO(png),
                    Inches(left_in),
                    Inches(top_in),
                    width=Inches(target_w),
                    height=Inches(target_h),
                )
            except Exception as exc:
                try:
                    print(f"[visualise][warn] plotly.to_image failed ({title_text}): {exc}")
                except Exception:
                    pass

        if preview_slides is not None:
            fig_json_str = json.dumps(fig.to_plotly_json(), cls=PlotlyJSONEncoder)
            preview_slides.append({"title": title_text, "fig_json": fig_json_str, "table_html": "", "notes": notes or ""})
            preview_figs.append(fig_json_str)

        return slide

    if df is None or type(df) is not pd.DataFrame or df.empty:
        if preview_slides is not None:
            return {
                "slides": [{"title": "Institutions overview", "table_html": "<div>No data loaded.</div>", "notes": ""}],
                "figs": [],
                "index": 0,
            }
        return None

    _cb("Building institutional-actors slides")

    outlet_col = None
    for c in ("publication_outlet", "publicationTitle", "publisher"):
        if c in df.columns:
            outlet_col = c
            break
    if outlet_col is None:
        raise KeyError("Missing required column 'publication_outlet' (or fallback 'publicationTitle' / 'publisher').")

    if "publisher_type_value" in df.columns:
        type_col = "publisher_type_value"
    else:
        df["_tmp_pubtype_"] = "Unknown/Other"
        type_col = "_tmp_pubtype_"

    def _split_outlets(s: str) -> list[str]:
        parts = re.split(r"[;|,/]\s*", str(s or ""))
        return [" ".join(p.strip().split()) for p in parts if p and p.strip()]

    tmp = df[[outlet_col, type_col]].copy().rename(
        columns={outlet_col: "publication_outlet", type_col: "publisher_type_value"}
    )

    tmp["publication_outlet"] = tmp["publication_outlet"].astype(str).map(lambda x: " ".join(x.split()))
    tmp["publisher_type_value"] = tmp["publisher_type_value"].astype(str).replace({"": "Unknown/Other", "nan": "Unknown/Other"})

    tmp["publication_outlet_list"] = tmp["publication_outlet"].apply(_split_outlets)
    tmp = tmp.explode("publication_outlet_list", ignore_index=True)
    tmp["publication_outlet"] = tmp["publication_outlet_list"].fillna("").astype(str).str.strip()
    tmp = tmp.drop(columns=["publication_outlet_list"])
    tmp = tmp[tmp["publication_outlet"] != ""].copy()

    if tmp.empty:
        table_df = pd.DataFrame({"Publication Outlet": ["—"], "Publisher Type": ["—"], "Count": [0]})
        _add_table_slide("(B) Major institutional actors (publication_outlet) — Table", table_df)

        fig_inst = px.bar(pd.DataFrame({"publication_outlet": [], "Count": []}), x="publication_outlet", y="Count")
        fig_inst.update_layout(
            template="plotly_white",
            width=PLOT_W,
            height=PLOT_H,
            margin=dict(l=68, r=120, t=80, b=80),
            title=dict(text=f"{collection_name} · Major institutional actors (publication_outlet)", x=0.02, xanchor="left", font=dict(size=20, family=FONT_FAMILY)),
            font=dict(family=FONT_FAMILY, size=13),
            legend_title_text="",
        )
        _add_figure_slide(
            "(B) Major institutional actors (publication_outlet) — Figure",
            fig_inst,
            notes="Method: No usable 'publication_outlet' entries were found in the dataset.",
        )

        if preview_slides is not None:
            return {"slides": preview_slides, "figs": preview_figs or [], "index": 0}
        return None

    counts = tmp.groupby("publication_outlet").size().rename("Count")
    modal_type = (
        tmp.groupby("publication_outlet")["publisher_type_value"]
        .agg(lambda s: s.value_counts().index[0])
        .rename("Publisher Type")
    )
    agg = (
        pd.concat([counts, modal_type], axis=1)
        .sort_values("Count", ascending=False)
        .head(int(top_n_outlets))
        .reset_index()
    )

    table_df = agg.rename(columns={"publication_outlet": "Publication Outlet"})[
        ["Publication Outlet", "Publisher Type", "Count"]
    ]
    _add_table_slide(
        "(B) Major institutional actors (publication_outlet) — Table",
        table_df,
        notes=(
            "Method: publication_outlet values split on ';', ',', '/', or '|' and whitespace-normalised. "
            "Counts are outlet mentions (not unique records). Publisher Type is modal publisher_type_value per outlet."
        ),
    )

    fig_inst = px.bar(
        agg.sort_values("Count", ascending=False),
        x="publication_outlet",
        y="Count",
        color="Publisher Type",
        text="Count",
    )
    fig_inst.update_traces(textposition="outside", marker_line_width=0.5, marker_line_color="rgba(0,0,0,0.15)")
    fig_inst.update_xaxes(tickangle=-25, automargin=True)
    fig_inst.update_layout(
        template="plotly_white",
        width=PLOT_W,
        height=PLOT_H,
        margin=dict(l=68, r=120, t=80, b=80),
        title=dict(text=f"{collection_name} · Major institutional actors (publication_outlet)", x=0.02, xanchor="left", font=dict(size=20, family=FONT_FAMILY)),
        font=dict(family=FONT_FAMILY, size=13),
        legend_title_text="",
    )

    n_total = int(tmp.shape[0])
    n_used = int(agg["Count"].sum())

    _add_figure_slide(
        "(B) Major institutional actors (publication_outlet) — Figure",
        fig_inst,
        notes=(
            "Rationale: Bars show Top outlets from publication_outlet (no canonical mapping). "
            "Bar colour is modal publisher_type_value for that outlet; blanks labelled 'Unknown/Other'. "
            f"Count basis: {n_used} outlet mentions within {n_total} parsed outlet entries."
        ),
    )

    _cb("Institutions overview completed")

    if preview_slides is not None:
        return {"slides": preview_slides, "figs": preview_figs or [], "index": 0}

    return None

