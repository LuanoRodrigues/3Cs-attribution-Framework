from __future__ import annotations

# NOTE: This file is used as a backend PPTX builder. Keep imports resilient:
# - Prefer local `python_backend.*` modules.
# - Avoid legacy `src.*` and `bibliometric_analysis_tool.*` imports (not shipped in the Electron bundle).

try:
    from python_backend.retrieve.Zotero_loader_to_df import load_data_from_source_for_widget  # type: ignore
except Exception:  # pragma: no cover
    load_data_from_source_for_widget = None  # type: ignore[assignment]

try:
    from python_backend.visualise.data_processing import *  # type: ignore  # noqa: F403,F401
except Exception:  # pragma: no cover
    # Allow importing this module even when data_processing deps are missing; callers can still use parts that don't need it.
    pass

# Optional Zotero client instance (legacy code used `src.core.literature_review.zt`).
zt = None
import logging, sys
logging.basicConfig(level=logging.INFO,
                    format="%(asctime)s | %(levelname)s | %(message)s",
                    stream=sys.stdout,
                    force=True)



# Content margins & title height
OUTER_MARGIN_X = 0.6   # left/right margin (inches)
OUTER_MARGIN_Y = 0.6   # top/bottom margin below title (inches)
TITLE_HEIGHT_IN = 0.7  # height reserved for the title row

# ---- plotting render size ----
PLOT_W, PLOT_H = 1100, 600

# ---------- helpers ----------
def _split_authors_semicolon(val) -> list[str]:
    """
    Split authors by ';' and trim; returns [] for NaN/empty.
    Safe with lists that may contain pd.NA/None.
    """
    import pandas as _pd, re as _re
    if isinstance(val, list):
        raw = val
    else:
        if _pd.isna(val):
            raw = []
        else:
            raw = _re.split(r"[;]", str(val))
    out, seen = [], set()
    for p in raw:
        if p is None or _pd.isna(p):
            continue
        s = str(p).strip()
        if not s:
            continue
        if s not in seen:
            seen.add(s)
            out.append(s)
    return out


def _find_year_col(df):
    for cand in ("year_numeric", "year", "Year", "pub_year"):
        if cand in df.columns: return cand
    return None


def _find_citations_col(df):
    import re
    pats = (r"^citations?$", r"^times_?cited$", r"^n_?cit", r"^num_?cit", r"^cites?$")
    for c in df.columns:
        if any(re.search(p, str(c), flags=re.I) for p in pats):
            return c
    return None


def _safe_int(x, default=0):
    import pandas as _pd
    try:
        if x is None or _pd.isna(x) or (isinstance(x, str) and not x.strip()):
            return default
        return int(float(x))
    except Exception:
        return default
# --- drop-in replacements / additions ---

def _canonical_csv(df, max_chars: int = 20000) -> str:
    """
    Make the CSV prompt-stable across runs:
      - consistent column ordering (alphabetical by column name)
      - coerce numerics; round floats; strip strings
      - stable row ordering (mergesort on all columns as strings)
      - Unix newlines
    """
    import pandas as pd

    if df is None or getattr(df, "empty", True):
        return ""

    # Work on a copy
    d = df.copy()

    # Ensure columns exist and are sortable
    cols = list(d.columns)
    cols_sorted = sorted(cols, key=lambda s: str(s))

    # Coerce columns for comparability and stable sort keys
    for c in cols_sorted:
        # Try numeric first
        s = pd.to_numeric(d[c], errors="coerce")
        if s.notna().any():
            # numeric column
            d[c] = s.round(6)  # round floats for stability
        else:
            # treat as text (incl. NaN→"")
            d[c] = d[c].astype(str).str.strip()

    # Reorder columns
    d = d[cols_sorted]

    # Sort rows deterministically by all columns as strings (stable mergesort)
    sort_keys = [c for c in cols_sorted]
    d_for_sort = d.astype(str)
    d = d.iloc[d_for_sort.sort_values(sort_keys, kind="mergesort").index]

    csv = d.to_csv(index=False, lineterminator="\n")
    return csv[:max_chars] + ("\n# …TRUNCATED…" if len(csv) > max_chars else "")


def _stable_cache_dir(base_logs_dir: str | None) -> str:
    """
    Choose a repeatable cache directory that doesn't depend on Path.cwd().
    Prefer a project/logs root if provided; else, anchor to the repo (or home).
    """
    from pathlib import Path
    if base_logs_dir:
        root = Path(base_logs_dir)
    else:
        # fallback: repo-ish anchor (parent of this file) or home
        try:
            root = Path(__file__).resolve().parents[1]
        except Exception:
            root = Path.home()
    p = root / ".annotarium_llm_cache"
    p.mkdir(parents=True, exist_ok=True)
    return str(p)


def _set_notes(slide, slide_notes: bool, text: str) -> None:
    """Write notes only when slide_notes is True; never let failures bubble up."""
    if not slide_notes:
        return
    try:
        slide.notes_slide.notes_text_frame.text = str(text or "")
    except Exception:
        # Last-ditch: never break slide creation because of notes
        try:
            slide.notes_slide.notes_text_frame.text = "(notes unavailable)"
        except Exception:
            pass

def _ensure_numeric(df: "pd.DataFrame | None",
                    fallback_cols: tuple[str, ...] = ("Value",),
                    min_rows: int = 1) -> "pd.DataFrame":
    """
    Return a numeric-only DataFrame with at least one row.
    This prevents downstream handlers from doing .iloc/.iat on empty frames.
    """
    if df is None or not isinstance(df, pd.DataFrame) or df.shape[1] == 0:
        return pd.DataFrame({c: [0] * max(1, min_rows) for c in fallback_cols})
    # Coerce every column to numeric; non-numeric → NaN → 0
    out = pd.DataFrame({c: pd.to_numeric(df[c], errors="coerce") for c in df.columns})
    out = out.fillna(0)
    if out.shape[0] == 0:
        out = pd.DataFrame({c: [0] * max(1, min_rows) for c in out.columns or fallback_cols})
    return out.reset_index(drop=True)

def _safe_ai_notes(
    fig,
    df_numeric: "pd.DataFrame | None",
    *,
    collection_name: str,
    analysis_suffix: str,
    title_for_prompt: str,
    prompt_purpose: str,
    max_tokens: int = 1200,
    effort: str = "low",
    base_logs_dir=None,
    _cb=lambda *_: None,
) -> str:
    """
    Call ai_notes_from_plot with a guaranteed numeric DF.
    Never raises; returns a fallback message on failure.
    """
    # ensure non-empty numeric payload
    df_numeric = _ensure_numeric(df_numeric)
    try:
        # ai_notes_from_plot is assumed to exist in your module
        return ai_notes_from_plot(
            fig,
            df_numeric,
            collection_name=collection_name,
            analysis_suffix=analysis_suffix,
            title_for_prompt=title_for_prompt,
            prompt_purpose=prompt_purpose,
            max_tokens=max_tokens,
            effort=effort,
            base_logs_dir=base_logs_dir,
            _cb=_cb,
        )
    except NameError:
        return "(notes helper unavailable: ai_notes_from_plot not found)"
    except Exception as e:
        return f"(notes generation failed safely: {type(e).__name__})"

def _write_notes(
    slide,
    slide_notes: bool,
    fig,
    df_numeric: "pd.DataFrame | None",
    **kwargs,
) -> None:
    """
    Convenience wrapper: run _safe_ai_notes(...) then write via _set_notes(...).
    """
    text = _safe_ai_notes(fig, df_numeric, **kwargs)
    _set_notes(slide, slide_notes, text)

def ai_notes_from_plot(
    fig: "go.Figure",
    raw_df: "pd.DataFrame",
    *,
    collection_name: str,
    analysis_suffix: str,
    title_for_prompt: str,
    prompt_purpose: str,
    vision_model_env: str = "OPENAI_VISION_MODEL",
    default_model: str = "gpt-5-mini",
    max_tokens: int = 2400,
    effort: str = "low",
    base_logs_dir: str | None = None,
    _cb=lambda _: None,
) -> str:
    import json, tempfile
    from pathlib import Path

    def _extract_text_from_plots(res: dict) -> str:
        if not isinstance(res, dict):
            return ""
        for key in ("text", "raw_text", "output_text"):
            val = res.get(key, "")
            if isinstance(val, str) and val.strip():
                return val.strip()
        try:
            blob = res.get("raw_response") or {}
            s = json.dumps(blob, ensure_ascii=False)
            if s:
                import re
                m = re.search(r'"(output_text|text)"\s*:\s*"([^"]+)"', s)
                if m:
                    return m.group(2).strip()
        except Exception:
            pass
        return ""

    # 1) Export the figure
    with tempfile.TemporaryDirectory() as _td:
        img_path = Path(_td) / f"{analysis_suffix}.png"
        try:
            png_bytes = fig.to_image(format="png", width=1500, height=800, scale=2)
            img_path.write_bytes(png_bytes)
        except Exception as e_img:
            _cb(f"AI analysis skipped (image export failed): {e_img}")
            return ""

        # 2) Canonical raw data for a stable prompt hash
        # 2) Canonical raw data — make it numeric-only and guaranteed ≥1×1
        try:
            df_num = raw_df.select_dtypes(include=["number", "bool"]).copy()
        except Exception:
            df_num = pd.DataFrame()

        if df_num.shape[1] == 0:
            # no numeric columns at all → create one dummy numeric column
            df_num = pd.DataFrame({"_dummy": [0]})
        elif df_num.empty:
            # numeric columns exist but there are no rows → add a zero row
            df_num = pd.DataFrame([{c: 0 for c in df_num.columns}])

        raw_csv = _canonical_csv(df_num, max_chars=20000)

        # 3) Compose prompt
        prompt_text = (
            f"[TASK]\nYou are a senior bibliometrics analyst. {prompt_purpose}\n"
            f"Write a 1–2 minute spoken transcript for slide notes. Return PLAIN TEXT only.\n"
            f"Include:\n"
            f"  • 3–5 executive takeaways\n"
            f"  • trends/turning points\n"
            f"  • outliers & plausible causes\n"
            f"  • methodological caveats\n"
            f"  • 3 concrete next steps\n\n"
            f"[FIGURE]\n{title_for_prompt}\n\n"
            f"[RAW DATA — CSV]\n{raw_csv}\n"
        )


        # 4) Use a STABLE cache dir (not Path.cwd())
        cache_dir = _stable_cache_dir(base_logs_dir)

        try:
            out = call_models_plots(
                mode="analyze",
                prompt_text=prompt_text,
                model_api_name=default_model,
                images=[str(img_path)],  # becomes data: URL internally
                image_detail="high",
                max_tokens=max_tokens,
                use_cache=True,
                cache_dir=cache_dir,            # <<< important
                analysis_key_suffix=analysis_suffix,
                section_title=collection_name or "authors",
                overall_topic=collection_name or "default", store_only=False,
                read=False,
            )
            # optional: bubble up the cache path/key for debugging
            if "cache_path" in out:
                _cb(f"CACHE PATH: {out.get('cache_path')} (key={out.get('cache_key','?')})")
            text = _extract_text_from_plots(out)
            return text if isinstance(text, str) else ""
        except Exception as e:
            _cb(f"AI analyze error: {e}")
            return ""

def add_data_summary_slides(
    prs,
    df: "pd.DataFrame",
    collection_name: str = "Collection",
    *,
    top_n: int = 3,
    slide_notes: bool = True,
    year_min: int = 1500,
    return_payload: bool = False,
):
    """
    ###1. compute deflation-safe summary metrics from an item-level deduplicated dataframe
    ###2. render the one-slide card layout
    ###3. optionally return a payload of computed tables/metrics for preview/testing
    """
    import re
    import pandas as pd
    import numpy as np
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor

    def _explode_semicolon_or_list(series: pd.Series) -> pd.Series:
        if series is None or series.empty:
            return pd.Series([], dtype=str)
        s = series.dropna()
        s = s.apply(lambda x: list(x) if isinstance(x, (set, tuple)) else x)
        is_list = s.apply(lambda x: isinstance(x, list))
        if is_list.any():
            s = pd.concat([s[~is_list], s[is_list].explode()])
        s = s.astype(str).str.split(";").explode().str.strip()
        s = s.replace("", pd.NA).dropna()
        return s

    def _top_n_from_col(series: pd.Series, n: int = 3) -> list[tuple[str, int]]:
        if series is None or series.empty:
            return []
        s = _explode_semicolon_or_list(series)
        if s.empty:
            return []
        vc = s.value_counts().head(int(n))
        return [(str(k), int(v)) for k, v in vc.items()]

    def _best_year_bounds(df_: pd.DataFrame, *, ymin: int = 1500) -> tuple[str, str]:
        years = None
        for col in ("year_numeric", "year"):
            if col in df_.columns:
                y = pd.to_numeric(df_[col], errors="coerce")
                years = y if years is None else years.fillna(y)
        if years is None:
            return "N/A", "N/A"
        years = years.dropna().astype(float)
        if years.empty:
            return "N/A", "N/A"
        current_year = pd.Timestamp.today().year + 1
        years = years[(years >= ymin) & (years <= current_year)]
        if years.empty:
            return "N/A", "N/A"
        return str(int(years.min())), str(int(years.max()))

    def _clean_numeric_series(raw: pd.Series) -> pd.Series:
        s = raw.astype(str).str.replace(",", "", regex=False).str.strip()
        s = s.str.extract(r"([0-9]*\.?[0-9]+)", expand=False)
        return pd.to_numeric(s, errors="coerce")

    def _best_citations_series(df_: pd.DataFrame) -> pd.Series | None:
        preferred = ["citations", "cited_by", "cited_by_count", "citations_count"]
        pattern_hits = [c for c in df_.columns if re.search(r"(citat|cited|times_cited)", str(c), re.I)]
        scan = preferred + [c for c in pattern_hits if c not in preferred]
        best, best_n = None, -1
        for col in scan:
            if col in df_.columns:
                s = _clean_numeric_series(df_[col]).dropna()
                if not s.empty and s.shape[0] > best_n:
                    best, best_n = s, s.shape[0]
        return best

    def _fmt_number(x: float | int | str) -> str:
        if isinstance(x, (int, np.integer)):
            return f"{int(x):,}"
        if isinstance(x, (float, np.floating)):
            return f"{x:,.2f}"
        return str(x)

    def _keywords_series(df_: pd.DataFrame) -> pd.Series:
        candidates = [
            "controlled_vocabulary_terms", "controlled vocabulary terms",
            "keywords", "keyword", "tags", "terms", "topics",
        ]
        possible = [c for c in df_.columns if any(re.search(rf"\b{re.escape(k)}\b", str(c), re.I) for k in candidates)]
        cols = [c for c in candidates if c in df_.columns] + possible

        out = []
        for col in cols:
            if col not in df_.columns or df_[col].isna().all():
                continue
            ser = _explode_semicolon_or_list(df_[col]).astype(str)
            theme_vals = ser.str.extract(r"^\s*#\s*theme\s*:\s*(.+?)\s*$", flags=re.I, expand=False)
            theme_vals = theme_vals.dropna().astype(str).str.strip()
            theme_vals = theme_vals.str.replace(r"\s+", " ", regex=True)
            theme_vals = theme_vals.replace("", pd.NA).dropna()
            if not theme_vals.empty:
                out.append(theme_vals)

        if not out:
            return pd.Series([], dtype=str)

        ser_all = pd.concat(out, ignore_index=True)
        ser_all = ser_all.str.strip().str.casefold()
        ser_all = ser_all.replace("", pd.NA).dropna()
        return ser_all

    def _year_series_local(frame: pd.DataFrame) -> pd.Series:
        ycol = "year_numeric" if "year_numeric" in frame.columns else ("year" if "year" in frame.columns else None)
        if ycol is None:
            return pd.Series([np.nan] * len(frame), index=frame.index)
        return pd.to_numeric(frame[ycol], errors="coerce")

    def _norm_title_for_key(s: str) -> str:
        s = "" if s is None or (isinstance(s, float) and pd.isna(s)) else str(s)
        s = s.lower()
        s = re.sub(r"[^a-z0-9]+", " ", s)
        s = re.sub(r"\s+", " ", s).strip()
        return s

    items = df.copy()

    if "key" in items.columns and items["key"].notna().any():
        items = items.sort_values("key").drop_duplicates("key", keep="first").copy()
    else:
        tkey = items["title"].map(_norm_title_for_key) if "title" in items.columns else pd.Series([""] * len(items), index=items.index)
        yser = _year_series_local(items)
        ykey = yser.fillna(-1).astype("int64")
        items = (
            items.assign(__tkey=tkey, __ykey=ykey)
            .sort_values(["__tkey", "__ykey"])
            .drop_duplicates(["__tkey", "__ykey"], keep="first")
            .drop(columns=["__tkey", "__ykey"])
            .copy()
        )

    total_docs = int(len(items))

    authors_col = items.get("authors", pd.Series(dtype=str))
    exploded_authors = authors_col.dropna().astype(str).str.split(";").explode().str.strip()
    exploded_authors = exploded_authors[exploded_authors != ""]
    unique_auth = int(exploded_authors.str.casefold().nunique()) if not exploded_authors.empty else 0

    if "citations" in items.columns:

        def _find_citations_col_local(frame: pd.DataFrame) -> str | None:
            preferred = ["citations", "cited_by", "cited_by_count", "citations_count"]
            for p in preferred:
                if p in frame.columns:
                    return p
            pats = (r"^citations?$", r"^cites?$", r"times_?cited", r"^n_?cit", r"^num_?cit", r"cited_by")
            for c in frame.columns:
                if any(re.search(p, str(c), flags=re.I) for p in pats):
                    return c
            return None

        _cit_col = _find_citations_col_local(items)
        cit_series = pd.to_numeric(items[_cit_col].copy(), errors="coerce") if _cit_col is not None else pd.Series(dtype=float)
    else:
        _best = _best_citations_series(items)
        cit_series = _best.copy() if _best is not None else pd.Series(dtype=float)

    cit_series = cit_series.replace([np.inf, -np.inf], np.nan).fillna(0.0)

    tot_cit_str = "0"
    avg_cit_str = "0.00"
    if not cit_series.empty:
        tot_cit_str = _fmt_number(float(cit_series[cit_series > 0].sum()))
        avg_cit_str = f"{float(cit_series.mean()):.2f}"

    min_year, max_year = _best_year_bounds(items, ymin=year_min)

    top_types = _top_n_from_col(items["item_type"], n=top_n) if "item_type" in items.columns else []
    top_sources = _top_n_from_col(items["source"], n=top_n) if "source" in items.columns else []
    top_authors = _top_n_from_col(exploded_authors, n=top_n) if not exploded_authors.empty else []

    kw_series = _keywords_series(items)
    top_keywords = []
    if not kw_series.empty:
        vc = kw_series.value_counts().head(int(top_n))
        top_keywords = [(f"{k}", int(v)) for k, v in vc.items()]

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    def _add_card(x_in, y_in, w_in, h_in, heading: str, *, value: str | None = None, lines: list[str] | None = None):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
        fill = shp.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(247, 247, 247)
        shp.line.color.rgb = RGBColor(205, 205, 205)

        tf = shp.text_frame
        tf.clear()

        p0 = tf.paragraphs[0]
        p0.text = heading.upper()
        p0.font.size = Pt(11)
        p0.font.bold = True
        p0.font.color.rgb = RGBColor(95, 95, 95)

        if lines is not None:
            for i, line in enumerate(lines):
                p = tf.add_paragraph()
                p.text = str(line)
                p.level = 1
                p.font.size = Pt(16)
                p.font.bold = (i == 0)
                p.font.color.rgb = RGBColor(30, 30, 30)
        else:
            p1 = tf.add_paragraph()
            p1.text = str(value or "")
            p1.level = 1
            p1.font.size = Pt(24)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(30, 30, 30)
        return shp

    def _fmt_top_lines(pairs: list[tuple[str, int]]) -> list[str]:
        return [f"{name} ({count})" for name, count in pairs] if pairs else ["N/A"]

    left_x, right_x = 0.5, 5.2
    col_w = 4.5
    top_y = 0.6
    card_h = 1.25
    gutter_y = 0.23

    _add_card(left_x, top_y + 0 * (card_h + gutter_y), col_w, card_h, "Total Documents", value=_fmt_number(total_docs))
    _add_card(left_x, top_y + 1 * (card_h + gutter_y), col_w, card_h, "Unique Authors", value=_fmt_number(unique_auth))
    _add_card(left_x, top_y + 2 * (card_h + gutter_y), col_w, card_h, "Total Citations", value=tot_cit_str)
    _add_card(left_x, top_y + 3 * (card_h + gutter_y), col_w, card_h, "Publication Timespan", value=f"{min_year}–{max_year}")

    _add_card(right_x, top_y + 0 * (card_h + gutter_y), col_w, card_h, f"Top Document Types (Top {top_n})", lines=_fmt_top_lines(top_types))
    _add_card(right_x, top_y + 1 * (card_h + gutter_y), col_w, card_h, f"Top Authors (Top {top_n})", lines=_fmt_top_lines(top_authors))
    _add_card(right_x, top_y + 2 * (card_h + gutter_y), col_w, card_h, f"Top Sources (Top {top_n})", lines=_fmt_top_lines(top_sources))
    _add_card(right_x, top_y + 3 * (card_h + gutter_y), col_w, card_h, f"Top Keywords (Top {top_n})", lines=_fmt_top_lines(top_keywords))

    ai_text = ""
    if slide_notes:
        import os
        import json

        stats_payload = {
            "collection": collection_name,
            "total_documents": int(total_docs),
            "unique_authors": int(unique_auth),
            "total_citations": str(tot_cit_str),
            "avg_citations_per_doc": str(avg_cit_str),
            "year_span": [str(min_year), str(max_year)],
            "filters": {"year_min_inclusive": int(year_min)},
            "top_n": int(top_n),
            "top_document_types": top_types,
            "top_authors": top_authors,
            "top_sources": top_sources,
            "top_keywords": top_keywords,
        }

        prompt_text = (
            "[TASK]\n"
            "You are a senior bibliometrics analyst. Draft 1–2 minutes of SPEAKER NOTES for a title-less summary slide.\n"
            "Deliver exactly:\n"
            " • 3–5 executive takeaways (bullet-like sentences)\n"
            " • context & caveats (coverage bias, year filter, name disambiguation, unnormalised citations)\n"
            " • notable concentrations AND notable absences\n"
            " • 3 practical next steps for validation/action\n\n"
            "[CONTEXT]\n"
            f"Collection: {collection_name}\n"
            f"Cards: Total Documents, Unique Authors, Total Citations, Publication Timespan, "
            f"Top {top_n} Document Types, Top {top_n} Authors, Top {top_n} Sources, Top {top_n} Keywords\n\n"
            "[RAW STATS — JSON]\n"
            f"{json.dumps(stats_payload, ensure_ascii=False)}\n"
        )

        ai_out = call_models_plots(
            mode="analyze",
            prompt_text=prompt_text,
            model_api_name=os.getenv("OPENAI_VISION_MODEL", "gpt-5.1-mini"),
            images=None,
            image_detail="low",
            max_tokens=900,
            use_cache=True,
            analysis_key_suffix="data_summary_notes",
            section_title=collection_name or "collection",
            store_only=False,
            read=False,
        )
        ai_text = (ai_out.get("text") or "").strip()

        if not ai_text:

            def _fmt_pairs(pairs):
                return ", ".join([f"{k} ({v})" for k, v in pairs]) if pairs else "N/A"

            ai_text = (
                f"{collection_name} • Summary. We count {total_docs:,} documents by ~{unique_auth:,} unique authors, "
                f"with {tot_cit_str} total citations (avg {avg_cit_str} per item). "
                f"Timespan {min_year}–{max_year} (values < {year_min} ignored). "
                f"Top types: {_fmt_pairs(top_types)}. Top authors: {_fmt_pairs(top_authors)}. "
                f"Top sources: {_fmt_pairs(top_sources)}. Top keywords: {_fmt_pairs(top_keywords)}. "
                "Caveats: coverage/metadata quality, author name disambiguation, and unnormalised citations may bias results. "
                "Next steps: (1) disambiguate authors; (2) field- and year-normalise citations; "
                "(3) audit a sample of item types/keywords for precision."
            )

        slide.notes_slide.notes_text_frame.text = ai_text

    if not return_payload:
        return None

    payload = {
        "collection_name": str(collection_name or ""),
        "metrics": {
            "total_docs": int(total_docs),
            "unique_authors": int(unique_auth),
            "total_citations_str": str(tot_cit_str),
            "avg_citations_str": str(avg_cit_str),
            "min_year": str(min_year),
            "max_year": str(max_year),
            "year_min_filter": int(year_min),
        },
        "top": {
            "document_types": top_types,
            "authors": top_authors,
            "sources": top_sources,
            "keywords": top_keywords,
        },
        "notes_text": str(ai_text or ""),
    }
    return payload


def add_authors_slides(
    prs,
    df: "pd.DataFrame",
    *,
    collection_name: str = "Collection",
    top_n_authors: int = 15,          # for Top lists & labels
    production_top_n: int = 10,       # for author-level time series
    image_width_inches: float = 10.0,
    slide_notes: bool = True,
    progress_callback=None,
) -> None:
    """
    Build THREE slides for the Authors section (updated):

      1) Two-panel time series:
           • Publications over time (count by year)
           • Citations over time (sum by year)

      2) Two stacked bars:
           • Top-15 most productive authors (by publication count)
           • Top-15 most cited authors (sum of per-paper citations)

      3) Scatter of the *Top-10 most collaborative* authors,
         clustered by co-authorship community (color), sized by total interactions.
         (Replaces the prior network graph.)

      Note: The previous Authors×Year heatmap is removed per request.
    """
    # ---- imports kept inside to avoid polluting caller ----
    import io

    import pandas as pd
    import plotly.express as px
    import plotly.graph_objects as go
    from pptx.util import Inches, Emu
    # ---------- global-safe helpers for notes (add once) ----------
    def _ensure_numeric(df, fallback_cols=("Value",), min_rows=1):
        """
        Return a numeric-only DataFrame with at least `min_rows` rows.
        If df is None/empty or ends up with 0 columns, create fallback numeric columns.
        """
        import pandas as _pd
        if df is None:
            df = _pd.DataFrame()

        if not isinstance(df, _pd.DataFrame):
            try:
                df = _pd.DataFrame(df)
            except Exception:
                df = _pd.DataFrame()

        out = _pd.DataFrame()
        for c in df.columns:
            try:
                out[c] = _pd.to_numeric(df[c], errors="coerce").fillna(0)
            except Exception:
                # skip non-coercible columns
                pass

        if out.shape[1] == 0:  # no columns → build a tiny numeric frame
            out = _pd.DataFrame({col: [0] * max(1, int(min_rows)) for col in fallback_cols})

        if out.shape[0] < int(min_rows):
            pad = _pd.DataFrame({c: [0] * (int(min_rows) - out.shape[0]) for c in out.columns})
            out = _pd.concat([out, pad], ignore_index=True)

        return out.reset_index(drop=True)

    def _safe_ai_notes(fig, df_numeric, *, collection_name, analysis_suffix,
                       title_for_prompt, prompt_purpose, max_tokens=1200,
                       effort="low", base_logs_dir=None, _cb=None):
        """
        Wrapper around ai_notes_from_plot that never raises. Returns a string.
        """
        fn = globals().get("ai_notes_from_plot")
        if not callable(fn):
            return "(notes unavailable)"
        try:
            return fn(
                fig, df_numeric,
                collection_name=collection_name,
                analysis_suffix=analysis_suffix,
                title_for_prompt=title_for_prompt,
                prompt_purpose=prompt_purpose,
                max_tokens=max_tokens,
                effort=effort,
                base_logs_dir=base_logs_dir,
                _cb=_cb,
            )
        except Exception as e:
            return f"(notes generation failed safely: {type(e).__name__})"

    def _set_notes(slide, slide_notes, text: str) -> None:
        """Write notes only when slide_notes is True; never let failures bubble up."""
        if not slide_notes:
            return
        try:
            slide.notes_slide.notes_text_frame.text = str(text or "")
        except Exception:
            # Last-ditch: never break slide creation because of notes
            try:
                slide.notes_slide.notes_text_frame.text = "(notes unavailable)"
            except Exception:
                pass

    # Local-safe truncation
    def __trunc(s: str, n: int = 72) -> str:
        s = "" if s is None or (isinstance(s, float) and pd.isna(s)) else str(s).strip()
        return s if len(s) <= n else s[: n - 1].rstrip() + "…"

    def __trunc_abs(s: str, n: int = 520) -> str:
        s = "" if s is None or (isinstance(s, float) and pd.isna(s)) else str(s).strip()
        s = " ".join(s.split())
        return s if len(s) <= n else s[: n - 1].rstrip() + "…"

    # Ensure an authors string exists even if _authors_str is not globally defined
    def __authors_fallback(row) -> str:

        if isinstance(row.get("authors"), str) and row["authors"].strip():
            return "; ".join([p.strip() for p in row["authors"].split(";") if p.strip()])
        if isinstance(row.get("creator_summary"), str) and row["creator_summary"].strip():
            return row["creator_summary"].strip()
        return ""

    _auth_func = globals().get("_authors_str", __authors_fallback)
    def __authors_fallback(row) -> str:
        if "authors_list" in row and isinstance(row["authors_list"], (list, tuple)) and row["authors_list"]:
            parts = []
            for a in row["authors_list"]:
                if isinstance(a, dict):
                    ln, fn = a.get("lastName"), a.get("firstName")
                    if ln and fn:
                        parts.append(f"{ln}, {fn}")
                    elif ln:
                        parts.append(str(ln))
                    elif fn:
                        parts.append(str(fn))
                else:
                    s = str(a).strip()
                    if s:
                        parts.append(s)
            if parts:
                return "; ".join(parts)
        if isinstance(row.get("authors"), str) and row["authors"].strip():
            return "; ".join([p.strip() for p in row["authors"].split(";") if p.strip()])
        if isinstance(row.get("creator_summary"), str) and row["creator_summary"].strip():
            return row["creator_summary"].strip()
    _auth_func = globals().get("_authors_str", __authors_fallback)

    def _cb(msg: str):
        if progress_callback:
            try: progress_callback(f"AUTHORS: {msg}")
            except Exception: pass
        logging.info(f"AUTHORS: {msg}")

    # ---------- helpers ----------
    def _df_to_csv_limited(df: pd.DataFrame, max_chars: int = 15000) -> str:
        if df is None or len(df) == 0:
            return ""
        csv = df.to_csv(index=False)
        return csv[:max_chars] + ("\n# …TRUNCATED…" if len(csv) > max_chars else "")


    def _split_authors_semicolon(val) -> list[str]:
        """Split authors by ';' or accept list; dedupe; safe for pd.NA."""
        import pandas as _pd, re as _re
        if isinstance(val, list):
            parts = [str(p).strip() for p in val
                     if p is not None and not _pd.isna(p) and str(p).strip()]
        else:
            s = "" if _pd.isna(val) else str(val)
            parts = [p.strip() for p in _re.split(r"[;]", s) if p.strip()]
        seen, out = set(), []
        for p in parts:
            if p not in seen:
                seen.add(p); out.append(p)
        return out

    def _year_series(frame: pd.DataFrame) -> pd.Series:
        ycol = "year_numeric" if "year_numeric" in frame.columns else ("year" if "year" in frame.columns else None)
        if ycol is None:
            return pd.Series([np.nan] * len(frame), index=frame.index)
        return pd.to_numeric(frame[ycol], errors="coerce")

    def _find_citations_col(frame: pd.DataFrame) -> str | None:
        import re
        pats = (r"^citations?$", r"^cites?$", r"times_?cited", r"^n_?cit", r"^num_?cit")
        for c in frame.columns:
            if any(re.search(p, str(c), flags=re.I) for p in pats):
                return c
        return None

    def _df_authors_col(frame: pd.DataFrame) -> str | None:
        if "authors_list" in frame.columns: return "authors_list"
        if "authors" in frame.columns: return "authors"
        return None

    def _fig_to_slide(fig: go.Figure, *, width_px=1500, height_px=800):

        try:
            png = fig.to_image(format="png", width=width_px, height=height_px, scale=2)
        except KeyboardInterrupt:
            import logging
            logging.warning("Cancelled by user (Ctrl+C) during image rendering — partial output only.")
            raise

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        left = Emu(int((prs.slide_width - Inches(image_width_inches)) / 2))
        top = Inches(0.6)
        slide.shapes.add_picture(io.BytesIO(png), left, top, width=Inches(image_width_inches))
        return slide

    if df is None or df.empty:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Authors — Overview"
        slide.placeholders[1].text = "No data available."
        return

    # ---------- Common pre-processing ----------
    # Item-level ground truth first: one row per document (prevents post-explode inflation)
    work = df.copy()

    # Prefer 'key' for uniqueness; fall back to (title, year) if needed
    if "key" in work.columns and work["key"].notna().any():
        work = work.sort_values("key").drop_duplicates("key", keep="first")
    else:
        # Conservative fallback (handles datasets without 'key')
        _y = _year_series(work)
        work = (
            work.assign(__tmp_year=_y)
            .sort_values(["title", "__tmp_year"])
            .drop_duplicates(["title", "__tmp_year"], keep="first")
            .drop(columns="__tmp_year")
        )

    work["__year"] = _year_series(work)

    cit_col = _find_citations_col(work)
    if cit_col is None:
        work["__cit"] = 0
    else:
        work["__cit"] = pd.to_numeric(work[cit_col], errors="coerce").fillna(0).astype(int)
    # Quick inflation audit (prints once to logs)
    try:
        import logging
        if "key" in df.columns and df["key"].notna().any():
            true_total = (
                df.sort_values("key").drop_duplicates("key", keep="first")
                .get(cit_col, pd.Series(0, index=df.index))
            )
            true_total = pd.to_numeric(true_total, errors="coerce").fillna(0).astype(float).sum()
            current_total = pd.to_numeric(work["__cit"], errors="coerce").fillna(0).astype(float).sum()
            if true_total > 0 and current_total / true_total > 1.05:
                logging.warning(f"[AUTHORS] Citation inflation detected: x{current_total / true_total:.2f} "
                                f"(current={current_total:.0f}, true={true_total:.0f})")
    except Exception:
        pass

    a_col = _df_authors_col(work)
    if a_col is None:
        work["__authors"] = [[] for _ in range(len(work))]
    else:
        work["__authors"] = work[a_col].apply(_split_authors_semicolon)

    # --- shared helpers (put once near the top of your slide code) ---

    def _authors_str(row) -> str:
        if "authors_list" in row and isinstance(row["authors_list"], (list, tuple)):
            parts = [
                ", ".join([p for p in [a.get("lastName"), a.get("firstName")] if p]) if isinstance(a, dict) else str(a)
                for a in row["authors_list"]
            ]
            parts = [p.strip() for p in parts if p and str(p).strip()]
            return "; ".join(parts)
        if "authors" in row and isinstance(row["authors"], str) and row["authors"].strip():
            return "; ".join([s.strip() for s in row["authors"].split(";") if s.strip()])
        if "creator_summary" in row and isinstance(row["creator_summary"], str) and row["creator_summary"].strip():
            return row["creator_summary"].strip()
        return ""

    def _truncate(s: str, n: int = 90) -> str:
        s = "" if s is None or pd.isna(s) else str(s).strip()
        return s if len(s) <= n else s[: n - 1].rstrip() + "…"

    def _truncate_abs(s: str, n: int = 320) -> str:
        s = "" if s is None or pd.isna(s) else str(s).strip()
        s = " ".join(s.split())
        return s if len(s) <= n else s[: n - 1].rstrip() + "…"


    # --- Slide 0 dataset (top N by citations) ---
    _cb("Building Slide 0: most-cited works/authors (cards)…")

    # Build Slide 0 — treemap “cards” showing Title (full), Authors (full), and Citations.
    # Also generate LLM notes safely by passing a numeric-only table, with rich strings embedded in the prompt.

    import numpy as np
    import textwrap as _tw

    def __trunc_abs(s: str, n: int = 520) -> str:
        s = "" if s is None or (isinstance(s, float) and pd.isna(s)) else str(s).strip()
        s = " ".join(s.split())
        return s if len(s) <= n else s[: n - 1].rstrip() + "…"

    def __authors_fallback(row) -> str:
        if "authors_list" in row and isinstance(row["authors_list"], (list, tuple)) and row["authors_list"]:
            out = []
            for a in row["authors_list"]:
                if isinstance(a, dict):
                    ln, fn = a.get("lastName"), a.get("firstName")
                    if ln and fn:
                        out.append(f"{ln}, {fn}")
                    elif ln:
                        out.append(str(ln))
                    elif fn:
                        out.append(str(fn))
                else:
                    s = str(a).strip()
                    if s: out.append(s)
            if out: return "; ".join(out)
        if isinstance(row.get("authors"), str) and row["authors"].strip():
            return "; ".join([p.strip() for p in row["authors"].split(";") if p.strip()])
        if isinstance(row.get("creator_summary"), str) and row["creator_summary"].strip():
            return row["creator_summary"].strip()
        return ""

    _auth_func = globals().get("_authors_str", __authors_fallback)

    # Build toplist
    work["__cit"] = pd.to_numeric(work.get("__cit", work.get("citations", 0)), errors="coerce").fillna(0).astype(int)
    _top = (
        work.assign(__Authors=work.apply(_auth_func, axis=1))
        .loc[:, ["title", "__Authors", "__cit", "abstract", "url"]]
        .rename(columns={
            "title": "Title",
            "__Authors": "Authors",
            "__cit": "Citations",
            "abstract": "Abstract",
            "url": "URL",
        })
    )

    if _top.empty:
        fig_cards = go.Figure().add_annotation(text="No citation data available.", showarrow=False)
        fig_cards.update_layout(
            template="plotly_white",
            height=520, width=1400,
            title=dict(text=f"{collection_name} · Most-cited works/authors", x=0.02, xanchor="left"),
            margin=dict(l=40, r=40, t=70, b=40),
            font=dict(family="Calibri", size=12, color="#2A2E33"),
            paper_bgcolor="white", plot_bgcolor="white",
        )
    else:
        _top = _top.sort_values(["Citations", "Title"], ascending=[False, True]).head(12).reset_index(drop=True)
        _top.insert(0, "Rank", _top.index + 1)

        # Treemap node labels (required)
        labels = [f"#{int(r.Rank)}" for _, r in _top.iterrows()]

        # Dynamic wrapping WITHOUT truncation for Title/Authors.
        def __wrap_card_full(title, authors, citations, share, fig_w=1400, fig_h=760):
            def _s(x):
                return "" if (x is None or (isinstance(x, float) and pd.isna(x))) else str(x).strip()

            title = _s(title)
            authors = _s(authors)

            # approximate drawable area per tile (account for margins/padding)
            total_area = max(1.0, (fig_w - 120) * (fig_h - 140))
            tile_area = max(1.0, float(share) * float(total_area))

            # width estimate from area (near-square)
            width_est = max(80.0, math.sqrt(tile_area) * 0.95)  # was 120.0 → allow narrower tiles

            # characters-per-line heuristic; allow very small cpl to avoid text suppression
            cpl = max(6, int(width_est / 7.0))  # was max(12, …)

            # soft wrap (no truncation, no ellipsis)
            title_html = _tw.fill(title, width=cpl, break_long_words=False, break_on_hyphens=False).replace("\n",
                                                                                                            "<br>")
            auth_html = _tw.fill(authors, max(6, cpl - 4), break_long_words=False, break_on_hyphens=False).replace("\n",
                                                                                                                   "<br>")
            cit_html = f"{int(citations):,} citations"

            return f"<b>{title_html}</b><br>{auth_html}<br><b>{cit_html}</b>"

        # per-tile shares and wrapped texts
        values = _top["Citations"].astype(int).tolist()
        shares = (np.array(values) / max(1, float(sum(values)))).tolist()

        card_text = [
            __wrap_card_full(r.Title, r.Authors, r.Citations, shares[i])
            for i, (_, r) in enumerate(_top.iterrows())
        ]

        parents = [""] * len(_top)  # flat treemap
        values = _top["Citations"].astype(int).tolist()

        palette = px.colors.qualitative.Set2 + px.colors.qualitative.Safe
        colors = [palette[i % len(palette)] for i in range(len(_top))]

        # Hover: do not truncate Title/Authors; keep Abstract compact
        def __to_s(x):
            return "" if (x is None or (isinstance(x, float) and pd.isna(x))) else str(x).strip()

        custom = list(
            zip(
                _top["Title"].apply(__to_s),  # no truncation
                _top["Authors"].apply(__to_s),  # no truncation
                _top["Citations"].astype(int),
                _top["Abstract"].apply(lambda s: __trunc_abs(s, 520)),  # compact abstract
                _top.get("URL", pd.Series([""] * len(_top))).fillna(""),
            )
        )
        fig_cards = go.Figure(
            go.Treemap(
                labels=labels,
                parents=parents,
                values=values,
                text=card_text,  # three-line “card”: Title / Authors / Citations (wrapped, never truncated)
                textinfo="text",
                texttemplate="%{text}",
                textfont=dict(size=11),  # slightly smaller baseline → fewer hidden tiles
                tiling=dict(pad=10, squarifyratio=1.05),
                marker=dict(colors=colors, line=dict(width=1.6, color="rgba(0,0,0,0.10)")),
                pathbar=dict(visible=False),
                customdata=custom,
                hovertemplate=(
                    "<b>%{customdata[0]}</b><br>"
                    "%{customdata[1]}<br>"
                    "Citations: %{customdata[2]:,}<br><br>"
                    "<b>Abstract:</b> %{customdata[3]}<br>"
                    "%{customdata[4]}<extra></extra>"
                ),
            )
        )
        # Force Plotly to SHOW text even for tiny tiles; scale down uniformly instead of hiding.
        fig_cards.update_layout(
            template="plotly_white",
            height=760, width=1400,
            title=dict(text=f"{collection_name} · Most-cited works/authors", x=0.02, xanchor="left"),
            margin=dict(l=48, r=48, t=78, b=44),
            uniformtext=dict(minsize=8, mode="show"),  # was: mode="hide", minsize=11
            font=dict(family="Calibri", size=12, color="#2A2E33"),
            paper_bgcolor="white",
            plot_bgcolor="white",
        )
        # Extra guard: ensure treemap traces keep the smaller text size even if other defaults override later.
        fig_cards.update_traces(textfont=dict(size=11), selector=dict(type="treemap"))

    # Render slide
    slide0 = _fig_to_slide(fig_cards, height_px=760)

    # LLM notes — pass numeric-only table; embed full Title/Authors/Abstract in the prompt to avoid ', with s' errors
    if slide_notes and not _top.empty:
        safe_df = _top.loc[:, ["Rank", "Citations"]].copy()
        safe_df["Rank"] = pd.to_numeric(safe_df["Rank"], errors="coerce").fillna(0).astype(int)
        safe_df["Citations"] = pd.to_numeric(safe_df["Citations"], errors="coerce").fillna(0).astype(int)

        def _s(x):  # safe string
            return "" if (x is None or (isinstance(x, float) and pd.isna(x))) else str(x).strip()

        bullets = []
        for _, r in _top.iterrows():
            bullets.append(
                f"{int(r['Rank'])}) {_s(r['Title'])} — {_s(r['Authors'])} — "
                f"{int(r['Citations']):,} — {__trunc_abs(_s(r['Abstract']), 200)}"
            )
        context_blob = "\n".join(bullets)

        try:
            _notes_text = _safe_ai_notes(
                fig_cards,
                _ensure_numeric(safe_df if isinstance(safe_df, pd.DataFrame) else pd.DataFrame(),
                                fallback_cols=("Rank", "Citations"), min_rows=1),
                collection_name=collection_name,
                analysis_suffix="authors_slide0_most_cited_cards",
                title_for_prompt=f"{collection_name} · Most-cited works/authors",
                prompt_purpose=(
                        "Analyse the Top-12 most-cited publications using the context list (Title — Authors — Citations — Abstract). "
                        "Identify dominant themes, concentration (few items vs. spread), and anomalies (self-citation, venue effects, "
                        "age bias). Provide 3–5 validation steps (DOI de-duplication, year-normalised rates, author disambiguation, "
                        "source coverage checks). Return crisp slide-notes bullets (≈200–300 words)."
                        "\n\nContext:\n" + context_blob
                ),
                max_tokens=1600,
                effort="low",
                base_logs_dir=None,
                _cb=_cb,
            )
            _set_notes(slide=slide0, slide_notes=slide_notes, text=_notes_text)

        except Exception as e:
            slide0.notes_slide.notes_text_frame.text = (
                f"LLM notes generation failed: {e}\n\nContext:\n{context_blob}"
            )
    # ======================================================================
    # Slide 1: Publications over time & Citations over time
    # ======================================================================



    # -------- Slide 1: Citations over time --------
    _cb("Building Slide 1: citations over time…")

    # try to reuse __year; else derive a minimal year series
    if "__year" not in work.columns:
        # minimal fallback year extraction (robust path likely already ran earlier)
        def __mini_year(x):
            import re, pandas as _pd, numpy as _np
            if isinstance(x, (int, float)) and 1800 <= int(x) <= 2100:
                return float(int(x))
            s = str(x) if x is not None and not _pd.isna(x) else ""
            s = s.strip()
            if not s:
                return _np.nan
            try:
                dt = _pd.to_datetime(s, errors="coerce", utc=True)
                if dt is not _pd.NaT and not _pd.isna(dt):
                    y = int(dt.year)
                    return float(y) if 1800 <= y <= 2100 else _np.nan
            except Exception:
                pass
            m = re.search(r"(18|19|20)\d{2}", s)
            return float(int(m.group(0))) if m else _np.nan

        candidate_year_cols = [
            "year_numeric", "year", "Year", "pub_year", "publication_year",
            "issued", "date", "Date", "publicationDate", "created"
        ]
        _ysrc = None
        for c in candidate_year_cols:
            if c in work.columns:
                _ysrc = c
                break
        if _ysrc is None:
            for c in work.columns:
                if any(k in str(c).lower() for k in ("year", "date", "issued", "created")):
                    _ysrc = c
                    break
        work["__year"] = work[_ysrc].apply(__mini_year) if _ysrc is not None else np.nan

    ts = work.dropna(subset=["__year"]).copy()
    if not ts.empty:
        ts["__year"] = ts["__year"].astype(int)
    else:
        ts = pd.DataFrame(columns=["__year", "__cit"])

    # normalise __cit
    if "__cit" not in ts.columns:
        ts["__cit"] = pd.to_numeric(work.get("__cit", work.get("citations", 0)), errors="coerce").fillna(0).astype(int)
    else:
        ts["__cit"] = pd.to_numeric(ts["__cit"], errors="coerce").fillna(0).astype(int)

    if ts.empty:
        fig_cit_ts = go.Figure().add_annotation(text="No year/citation data to plot.", showarrow=False)
        fig_cit_ts.update_layout(template="plotly_white", height=520, width=1400,
                                 title=f"{collection_name} · Citations Over Time")
        cit_full = pd.DataFrame(columns=["__year", "Citations", "YoY", "IsPeak", "IsSharpIncrease",
                                         "TopAuthors", "TopAuthorsCitations"])
    else:
        # base yearly series
        cit_by_year = (
            ts.groupby("__year", dropna=True)["__cit"].sum()
            .rename("Citations")
            .reset_index()
            .sort_values("__year")
        )
        # fill continuous years for visual continuity
        y0, y1 = int(cit_by_year["__year"].min()), int(cit_by_year["__year"].max())
        idx = pd.DataFrame({"__year": list(range(y0, y1 + 1))})
        cit_full = idx.merge(cit_by_year, on="__year", how="left").fillna({"Citations": 0})
        cit_full["Citations"] = cit_full["Citations"].astype(int)

        # YoY change
        cit_full["YoY"] = cit_full["Citations"].diff().fillna(0).astype(int)

        # simple local-peak flag: strictly greater than neighbours (handles endpoints as non-peaks)
        def _is_peak(i, vals):
            if i <= 0 or i >= len(vals) - 1:
                return False
            return vals[i] > vals[i - 1] and vals[i] > vals[i + 1]

        vals = cit_full["Citations"].tolist()
        cit_full["IsPeak"] = [bool(_is_peak(i, vals)) for i in range(len(vals))]

        # sharp increase flag: YoY above 75th percentile (and positive)
        try:
            q75 = pd.Series(vals).diff().dropna()
            q75 = float(q75.quantile(0.75)) if not q75.empty else 0.0
        except Exception:
            q75 = 0.0
        cit_full["IsSharpIncrease"] = (cit_full["YoY"] > max(q75, 0)).astype(bool)

        # Top contributing author(s) per year:
        #   assign each record's citations to its publication year and author(s); sum by (year, author)
        #   NOTE: this treats __cit as attached to the doc's year (as in your original logic).
        def _extract_authors_any(row) -> list[str]:
            """Return a list of author strings from any available field."""
            import re

            # Prefer already-parsed list (e.g., __authors)
            if isinstance(row.get("__authors"), list):
                return [str(a).strip() for a in row["__authors"] if str(a).strip()]

            # Structured Zotero-like list
            if isinstance(row.get("authors_list"), (list, tuple)) and row["authors_list"]:
                out = []
                for a in row["authors_list"]:
                    if isinstance(a, dict):
                        ln, fn = a.get("lastName"), a.get("firstName")
                        if ln and fn:
                            out.append(f"{ln}, {fn}")
                        elif ln:
                            out.append(str(ln))
                        elif fn:
                            out.append(str(fn))
                    else:
                        s = str(a).strip()
                        if s:
                            out.append(s)
                if out:
                    return out

            # Flat strings
            for key in ("authors", "creator_summary"):
                s = row.get(key)
                if isinstance(s, str) and s.strip():
                    s = s.strip()
                    import re
                    # if no separators, return the single name
                    if not re.search(r"[;,]", s):
                        return [s]
                    # otherwise split on ; or ,
                    return [part.strip() for part in re.split(r"[;|,]+", s) if part.strip()]

            return []

        ts_auth = ts.copy()
        ts_auth["__AuthorsList"] = ts_auth.apply(_extract_authors_any, axis=1)
        auth_year_cits = (
            ts_auth.explode("__AuthorsList")
            .dropna(subset=["__AuthorsList"])
            .groupby(["__year", "__AuthorsList"])["__cit"].sum()
            .rename("AuthorCitations")
            .reset_index()
        )

        # attach top author info to each year row
        top_author_map = {}
        for y in map(int, cit_full["__year"].tolist()):
            sub = auth_year_cits.loc[auth_year_cits["__year"] == y, ["__year", "__AuthorsList", "AuthorCitations"]]
            if sub.empty:
                top_author_map[y] = ("", 0)
                continue

            # normalise type; treat non-numerics as 0
            sub = sub.copy()
            sub["AuthorCitations"] = pd.to_numeric(sub["AuthorCitations"], errors="coerce").fillna(0).astype(int)

            # top citation value for the year
            top_c = int(sub["AuthorCitations"].max())

            # tie-handling: collect up to two authors that share the top value
            ties = sub.loc[sub["AuthorCitations"] == top_c, "__AuthorsList"].astype(str).tolist()
            ties = [t for t in ties if t][:2]

            # stable single-name fallback in case ties list ends up empty (shouldn't happen)
            if not ties:
                sub_sorted = sub.sort_values(["AuthorCitations", "__AuthorsList"], ascending=[False, True])
                names = sub_sorted["__AuthorsList"].astype(str).tolist()
                ties = [names[0]] if names else []

            top_author_map[y] = ("; ".join(ties), top_c)

        cit_full["TopAuthors"] = cit_full["__year"].map(lambda yy: top_author_map.get(yy, ("", 0))[0])
        cit_full["TopAuthorsCitations"] = cit_full["__year"].map(lambda yy: top_author_map.get(yy, ("", 0))[1]).astype(
            int)

        # Build figure
        fig_cit_ts = go.Figure(
            go.Scatter(
                x=cit_full["__year"].astype(int),
                y=cit_full["Citations"].astype(int),
                mode="lines+markers",
                name="Citations",
                hovertemplate="Year %{x}<br>Citations: %{y:,}<extra></extra>",
            )
        )
        fig_cit_ts.update_layout(
            template="plotly_white",
            height=560, width=1400,
            title=dict(text=f"{collection_name} · Citations Over Time", x=0.02, xanchor="left"),
            margin=dict(l=64, r=44, t=76, b=64),
            xaxis_title="Year",
            yaxis_title="Citations",
            font=dict(family="Calibri", size=12, color="#2A2E33"),
        )

        # Annotate up to N notable years (peaks or sharp increases)
        N_ANN = 3
        notable = cit_full[(cit_full["IsPeak"]) | (cit_full["IsSharpIncrease"])].copy()
        # rank by Citations, then YoY magnitude
        if not notable.empty:
            notable["YoYAbs"] = notable["YoY"].abs()
            notable = notable.sort_values(["Citations", "YoYAbs", "__year"], ascending=[False, False, True]).head(N_ANN)
            y_range = (cit_full["Citations"].min(), cit_full["Citations"].max())
            bump = max(1, int((y_range[1] - y_range[0]) * 0.05))  # label offset
            for _, r in notable.iterrows():
                year_i = int(r["__year"])
                val_i = int(r["Citations"])  # use yearly total for y-position on the line plot

                sub_year = ts[ts["__year"] == year_i]
                if sub_year.empty or "__cit" not in sub_year.columns:
                    label = "0 — —"
                else:
                    # robust idxmax that never raises on bad/empty data
                    _series = pd.to_numeric(sub_year["__cit"], errors="coerce").fillna(0)
                    if int(_series.shape[0]) == 0:
                        label = "0 — —"
                    else:
                        idx = _series.idxmax()
                        row_max = sub_year.loc[idx] if idx in sub_year.index else None
                        mdv = pd.to_numeric(row_max["__cit"], errors="coerce") if row_max is not None else 0
                        max_doc = int(0 if pd.isna(mdv) else int(mdv))
                        authors_list = _extract_authors_any(row_max) if row_max is not None else []
                        authors_str = "; ".join([a for a in authors_list if a]) or "—"
                        label = f"{max_doc:,} — {authors_str}"

                fig_cit_ts.add_annotation(
                    x=year_i, y=val_i + bump,
                    text=label,
                    showarrow=True,
                    arrowhead=2,
                    ax=0, ay=-20,
                    bgcolor="rgba(255,255,255,0.8)",
                    bordercolor="rgba(0,0,0,0.15)",
                    borderwidth=1,
                    font=dict(size=11),
                )

    slide1 = _fig_to_slide(fig_cit_ts)

    if slide_notes:
        try:
            # numeric-only table to avoid formatting errors inside ai_notes_from_plot
            raw_ts = cit_full.rename(columns={"__year": "Year"})[["Year", "Citations", "YoY"]].copy()
            raw_ts["Year"] = pd.to_numeric(raw_ts["Year"], errors="coerce").fillna(0).astype(int)
            raw_ts["Citations"] = pd.to_numeric(raw_ts["Citations"], errors="coerce").fillna(0).astype(int)
            raw_ts["YoY"] = pd.to_numeric(raw_ts["YoY"], errors="coerce").fillna(0).astype(int)
        except Exception:
            raw_ts = pd.DataFrame(columns=["Year", "Citations", "YoY"])

        # Put the peak/author context into the prompt text (not the table)
        notable = cit_full[(cit_full.get("IsPeak", False)) | (cit_full.get("IsSharpIncrease", False))].copy()
        notable = notable.sort_values(["Citations", "YoY", "__year"], ascending=[False, False, True]).head(5)
        context_peaks = "\n".join([
            f"{int(r['__year'])}: {int(r['Citations']):,} (Top: {str(r.get('TopAuthors', '')).strip() or '—'} — "
            f"{int(r.get('TopAuthorsCitations', 0)):,})"
            for _, r in notable.iterrows()
        ]) if not notable.empty else "—"

        _notes_text = _safe_ai_notes(
            fig_cit_ts,
            _ensure_numeric(raw_ts, fallback_cols=("Year", "Citations", "YoY"), min_rows=1),
            collection_name=collection_name,
            analysis_suffix="authors_slide1_citations_over_time",
            title_for_prompt=f"{collection_name} · Citations Over Time",
            prompt_purpose=(
                "Use the Year–Citations–YoY table to interpret the trajectory. "
                "Identify growth phases, plateaus, and outliers; explain likely causes (publication batches, venue effects, "
                "indexing lag); and discuss recency undercount risk.\n\n"
                "Peak/Sharp-increase context:\n"
                f"{context_peaks}\n\n"
                "Conclude with validation actions: coverage checks, DOI de-duplication, and age normalisation."
            ),
            max_tokens=1600,
            effort="low",
            base_logs_dir=None,
            _cb=_cb,
        )
        _set_notes(slide=slide1, slide_notes=slide_notes, text=_notes_text)

    # Slide 3 — Top-5 authors: production over time (line chart)
    _cb("Building Slide 3: production of top-5 authors over time…")

    # --- helpers ---
    def _extract_authors_any(row) -> list[str]:
        # Prefer structured 'authors_list' if present
        if "authors_list" in row and isinstance(row["authors_list"], (list, tuple)) and row["authors_list"]:
            out = []
            for a in row["authors_list"]:
                if isinstance(a, dict):
                    ln, fn = a.get("lastName"), a.get("firstName")
                    if ln and fn:
                        out.append(f"{ln}, {fn}")
                    elif ln:
                        out.append(str(ln))
                    elif fn:
                        out.append(str(fn))
                else:
                    s = str(a).strip()
                    if s:
                        out.append(s)
            if out:
                return out
        # Fallback to flat 'authors' string (assume ';' or ',' separators)
        if isinstance(row.get("authors"), str) and row["authors"].strip():
            parts = [p.strip() for p in re.split(r"[;|,]+", row["authors"]) if p.strip()]
            if parts:
                return parts
        # Fallback to 'creator_summary'
        if isinstance(row.get("creator_summary"), str) and row["creator_summary"].strip():
            parts = [p.strip() for p in re.split(r"[;|,]+", row["creator_summary"]) if p.strip()]
            if parts:
                return parts
        return []

    def _safe_title(x) -> str:
        return "" if x is None or (isinstance(x, float) and pd.isna(x)) else str(x).strip()

    # Ensure year available (__year already built earlier; if not, derive minimally)
    if "__year" not in work.columns:
        def __mini_year(x):
            import re as _re, pandas as _pd, numpy as _np
            if isinstance(x, (int, float)) and 1800 <= int(x) <= 2100:
                return float(int(x))
            s = str(x) if x is not None and not _pd.isna(x) else ""
            s = s.strip()
            if not s:
                return _np.nan
            try:
                dt = _pd.to_datetime(s, errors="coerce", utc=True)
                if dt is not _pd.NaT and not _pd.isna(dt):
                    y = int(dt.year)
                    return float(y) if 1800 <= y <= 2100 else _np.nan
            except Exception:
                pass
            m = _re.search(r"(18|19|20)\d{2}", s)
            return float(int(m.group(0))) if m else _np.nan

        _ycands = [
            "year_numeric", "year", "Year", "pub_year", "publication_year",
            "issued", "date", "Date", "publicationDate", "created"
        ]
        _ysrc = next((c for c in _ycands if c in work.columns), None)
        if _ysrc is None:
            _ysrc = next((c for c in work.columns
                          if any(k in str(c).lower() for k in ("year", "date", "issued", "created"))), None)
        work["__year"] = work[_ysrc].apply(__mini_year) if _ysrc is not None else np.nan

    # Explode to author–year–title rows
    w3 = work.copy()
    w3["__year"] = pd.to_numeric(w3["__year"], errors="coerce")
    w3 = w3.dropna(subset=["__year"]).copy()
    if not w3.empty:
        w3["__year"] = w3["__year"].astype(int)
    w3["__AuthorsList"] = w3.apply(_extract_authors_any, axis=1)

    def __dedupe_list(L):
        seen = set()
        out = []
        for x in (L if isinstance(L, list) else []):
            k = re.sub(r"\s+", " ", str(x).strip()).casefold()
            if k and k not in seen:
                seen.add(k);
                out.append(str(x).strip())
        return out

    w3["__AuthorsList"] = w3["__AuthorsList"].apply(__dedupe_list)
    w3["__TitleSafe"] = w3["title"].apply(_safe_title) if "title" in w3.columns else ""

    auth_rows = (
        w3.explode("__AuthorsList")
        .rename(columns={"__AuthorsList": "Author"})
    )
    auth_rows["Author"] = auth_rows["Author"].fillna("").astype(str).str.strip()
    auth_rows = auth_rows[auth_rows["Author"] != ""]

    # Drop institutional aliases (DG R&I, Publications Office, etc.) from ranking entirely.
    def __norm_author_name(s: str) -> str:
        s = re.sub(r"\s+", " ", str(s).strip()).casefold()
        s = re.sub(r"[\u2010-\u2015-]+", "-", s)  # normalise hyphens
        return s

    _ec_aliases = {
        "directorate-general for research and innovation",
        "directorate general for research and innovation",
        "publications office of the european union",
        "publications office of the eu",
        "directorate-general for research & innovation",
        "dg research and innovation",
        "dg research & innovation",
        "directorate general research and innovation",
        "directorate-general research and innovation (european commission)",
    }
    _ec_aliases_norm = {__norm_author_name(x) for x in _ec_aliases}

    auth_rows["__norm"] = auth_rows["Author"].map(__norm_author_name)
    auth_rows = auth_rows[~auth_rows["__norm"].isin(_ec_aliases_norm)].drop(columns="__norm")

    # Identify Top-5 authors by publication count (after filtering aliases)
    author_counts = (
        auth_rows.groupby("Author").size().rename("Count").reset_index()
        .sort_values(["Count", "Author"], ascending=[False, True])
    )
    top5_authors = author_counts.head(5)["Author"].tolist()

    # Time series for Top-5
    auth_top = auth_rows[auth_rows["Author"].isin(top5_authors)].copy()
    ts_auth = (auth_top.groupby(["__year", "Author"]).size()
               .rename("Publications").reset_index())
    years_span = []
    if not ts_auth.empty:
        y0, y1 = int(ts_auth["__year"].min()), int(ts_auth["__year"].max())
        years_span = list(range(y0, y1 + 1))
        # complete grid
        grid = pd.MultiIndex.from_product([years_span, top5_authors], names=["__year", "Author"]).to_frame(index=False)
        ts_auth = (grid.merge(ts_auth, on=["__year", "Author"], how="left")
                   .fillna({"Publications": 0}))
    else:
        ts_auth = pd.DataFrame(columns=["__year", "Author", "Publications"])

    # Plot: one line per author
    if ts_auth.empty:
        fig_auth_ts = go.Figure().add_annotation(text="No author/year data to plot.", showarrow=False)
        fig_auth_ts.update_layout(template="plotly_white", height=560, width=1400,
                                  title=f"{collection_name} · Top-5 Authors — Publications Over Time")
    else:
        fig_auth_ts = go.Figure()
        for a in top5_authors:
            sub = ts_auth[ts_auth["Author"] == a]
            fig_auth_ts.add_trace(
                go.Scatter(
                    x=sub["__year"].astype(int),
                    y=sub["Publications"].astype(int),
                    mode="lines+markers",
                    name=a,
                    hovertemplate="Year %{x}<br>Publications: %{y:,}<extra></extra>",
                )
            )

        fig_auth_ts.update_layout(
            template="plotly_white",
            height=560, width=1400,
            title=dict(text=f"{collection_name} · Top-5 Authors — Publications Over Time", x=0.02, xanchor="left"),
            margin=dict(l=60, r=40, t=70, b=60),
            xaxis_title="Year",
            yaxis_title="Publications",
            legend_title_text="Author",
            font=dict(family="Calibri", size=12, color="#2A2E33"),
        )

    slide3 = _fig_to_slide(fig_auth_ts)



        # 1) Numeric-only pivot for the LLM (safe for formatters)
    try:
        pivot_num = (
            ts_auth.pivot(index="__year", columns="Author", values="Publications")
            .fillna(0).reset_index().rename(columns={"__year": "Year"})
        )
        pivot_num["Year"] = pd.to_numeric(pivot_num["Year"], errors="coerce").fillna(0).astype(int)
        for c in pivot_num.columns:
            if c != "Year":
                pivot_num[c] = pd.to_numeric(pivot_num[c], errors="coerce").fillna(0).astype(int)
    except Exception:
        pivot_num = pd.DataFrame()

    # 2) Totals & shares (stable ordering → deterministic cache)
    totals_df = (
        ts_auth.groupby("Author", as_index=False)["Publications"].sum()
        .sort_values(["Publications", "Author"], ascending=[False, True], kind="mergesort")
    )
    total_all = int(pd.to_numeric(totals_df["Publications"], errors="coerce").fillna(0).sum())

    totals_lines = []
    for _, r in totals_df.iterrows():
        a = str(r["Author"])
        tmp = pd.to_numeric(r["Publications"], errors="coerce")
        v = 0 if pd.isna(tmp) else int(float(tmp))
        share = (100.0 * v / total_all) if total_all else 0.0
        totals_lines.append(f"{a}: {v} pubs ({share:.1f}%)")
    totals_blob = "; ".join(totals_lines) if totals_lines else "—"

    # 3) Consistency (per-author coefficient of variation over years)
    try:
        stats = ts_auth.groupby("Author")["Publications"].agg(mean="mean", std="std").fillna(0.0)
        stats["cv"] = stats.apply(lambda r: (float(r["std"]) / r["mean"]) if r["mean"] > 0 else float("inf"),
                                  axis=1)
        most_consistent = stats.replace([float("inf")], 1e9)["cv"].idxmin() if not stats.empty else "—"
        most_volatile = stats.replace([float("inf")], -1)["cv"].idxmax() if not stats.empty else "—"
    except Exception:
        most_consistent, most_volatile = "—", "—"

    # 4) Momentum (share of output in the last 3 years)
    try:
        last_year = int(ts_auth["__year"].max()) if not ts_auth.empty else None
        recent_years = [y for y in years_span if last_year and y >= last_year - 2]
        recent = ts_auth[ts_auth["__year"].isin(recent_years)].groupby("Author")["Publications"].sum()
        base = totals_df.set_index("Author")["Publications"].replace(0, pd.NA)
        recent_share = (recent / base).fillna(0.0)
        momentum_leader = recent_share.sort_values(ascending=False).index[0] if len(recent_share) else "—"
    except Exception:
        momentum_leader = "—"

    # 5) Influence proxy via citations (if available)
    try:
        _w = work.copy()
        _w["__cit"] = pd.to_numeric(_w.get("__cit", 0), errors="coerce").fillna(0).astype(int)
        _w["__AuthorsList"] = _w.apply(_extract_authors_any, axis=1)
        auth_cit = (_w.explode("__AuthorsList")
                    .rename(columns={"__AuthorsList": "Author"}))
        cit_totals = (auth_cit.groupby("Author")["__cit"].sum()
                      .reindex(totals_df["Author"]).fillna(0).astype(int))
        most_cited_author = cit_totals.idxmax() if len(cit_totals) else "—"
        cit_ctx = "; ".join([f"{a}: {int(cit_totals.get(a, 0))} cites"
                             for a in totals_df["Author"]]) if len(cit_totals) else "n/a"
    except Exception:
        most_cited_author, cit_ctx = "—", "n/a"

    # 6) Compose compact, deterministic context for the LLM
    metrics_blob = (
        f"Totals/Shares: {totals_blob} | "
        f"Most consistent (lowest CV): {most_consistent} | "
        f"Most volatile (highest CV): {most_volatile} | "
        f"Recent momentum leader (last 3y share): {momentum_leader} | "
        f"Citation context (if available): {cit_ctx} | "
        f"Note: ‘influence’ should be inferred primarily from citations, not counts."
    )

    # 7) Ask the LLM for a comparative narrative (evolution, influence, consistency)
    if slide_notes:
        _notes_text = _safe_ai_notes(
            fig_auth_ts,
            _ensure_numeric(pivot_num if isinstance(pivot_num, pd.DataFrame) else pd.DataFrame(),
                            fallback_cols=("Year",), min_rows=1),
            collection_name=collection_name,
            analysis_suffix="authors_slide3_top5_over_time_notes_v2",
            title_for_prompt=f"{collection_name} · Top-5 Authors — Publications Over Time",
            prompt_purpose=(
                    "Explain how each author's production evolved relative to peers using the Year×Author matrix. "
                    "Identify leaders by total output, periods of acceleration/slowdown, crossings between authors, "
                    "and whether output is concentrated in a few bursts or sustained. "
                    "Assess CONSISTENCY using year-to-year variance (low CV = steadier), and MOMENTUM via the last 3 years' share. "
                    "Discuss INFLUENCE carefully: rely on citations if provided (else note the limitation of using counts alone). "
                    "State which author appears most influential and why (e.g., sustained high output + strong citations, or late surge). "
                    "Close with 2–3 validation actions (author disambiguation, venue/field shifts, duplicate removal). "
                    "Helper metrics → " + metrics_blob
            ),
            max_tokens=1600,
            effort="low",
            base_logs_dir=None,
            _cb=_cb,
        )
        _set_notes(slide=slide3, slide_notes=slide_notes, text=_notes_text)

    # Slide 4 — Main themes by Top-5 authors (LLM over top-5×5 abstracts → normalized themes)
    # Slide 4a & 4b — Themes from top-5×5 cited works (LLM via function tools → normalised themes)
    # ────────────────────────────────────────────────────────────────────────────
    _cb("Building Slide 4a/4b: themes from top-5 most cited works per top-5 authors…")

    import re, json
    import numpy as np
    import pandas as pd
    import plotly.graph_objects as go
    try:
        import plotly.express as px
    except Exception:
        px = None  # fallback colors later

    # ---------- tiny helpers ----------
    # ---------- tiny helpers ----------
    def _safe(s):
        s = "" if s is None or (isinstance(s, float) and pd.isna(s)) else str(s)
        return " ".join(s.split()).strip()

    def _extract_authors_any(row) -> list[str]:
        """
        Return a cleaned list of author strings from any available field.

        Normalisation rule:
          • If 'European Commission' appears as a separate author on the SAME record,
            drop institutional aliases such as Directorate-General for Research and Innovation,
            Publications Office of the European Union, DG Research & Innovation, etc.
          • Regardless of co-presence, map those aliases to the canonical label 'European Commission'
            so they never appear as separate 'authors' in Top-N lists.
        """

        # --- helpers ---
        def _norm(s: str) -> str:
            s = re.sub(r"\s+", " ", str(s).strip()).casefold()
            s = re.sub(r"[\u2010-\u2015-]+", "-", s)
            return s

        def _is_ec(s: str) -> bool:
            n = _norm(s)
            return bool(re.search(r"\beuropean commission\b", n))

        def _is_ec_alias(s: str) -> bool:
            n = _norm(s)
            # DG Research & Innovation (many spellings)
            if re.search(r"\bdirectorate[-\s]?general\b.*\bresearch\b.*\binnovation\b", n):
                return True
            if re.search(r"\bdg\b.*\bresearch\b.*\binnovation\b", n):
                return True
            # Publications Office of the EU
            if re.search(r"\bpublications office\b.*\b(european union|eu)\b", n):
                return True
            return False

        EC = "European Commission"

        # --- collect raw authors (prefer structured list) ---
        raw: list[str] = []
        if "authors_list" in row and isinstance(row["authors_list"], (list, tuple)) and row["authors_list"]:
            for a in row["authors_list"]:
                if isinstance(a, dict):
                    ln, fn = a.get("lastName"), a.get("firstName")
                    if ln and fn:
                        raw.append(f"{ln}, {fn}")
                    elif ln:
                        raw.append(str(ln))
                    elif fn:
                        raw.append(str(fn))
                else:
                    s = str(a).strip()
                    if s:
                        raw.append(s)
        elif isinstance(row.get("authors"), str) and row["authors"].strip():
            raw = [p.strip() for p in re.split(r"[;|,]+", row["authors"]) if p.strip()]
        elif isinstance(row.get("creator_summary"), str) and row["creator_summary"].strip():
            raw = [p.strip() for p in re.split(r"[;|,]+", row["creator_summary"]) if p.strip()]

        if not raw:
            return []

        # Is EC explicitly present as a SEPARATE author?
        ec_present = any(_is_ec(a) for a in raw)

        # Build cleaned list with canonicalisation + per-record suppression when EC is present
        cleaned: list[str] = []
        seen_norm: set[str] = set()
        for a in raw:
            if ec_present and _is_ec_alias(a):
                # drop alias when EC is also listed
                continue
            # map any alias (even when EC not listed) to the canonical EC label
            label = EC if _is_ec_alias(a) else a.strip()
            key = _norm(label)
            if key and key not in seen_norm:
                seen_norm.add(key)
                cleaned.append(label)

        return cleaned

    def _extract_theme_tags(cell) -> list[str]:
        """
        Return canonical '#theme:<term>' tokens from a cell (list or string).
        Accepts lists/sets/tuples and ';|,\\n' separators. Ignores non-theme tokens.
        """
        if cell is None or (isinstance(cell, float) and pd.isna(cell)):
            return []
        if isinstance(cell, (list, tuple, set)):
            parts = []
            for v in cell:
                if v is None or (isinstance(v, float) and pd.isna(v)):
                    continue
                parts.extend(re.split(r"[;\|,\n]+", str(v)))
        else:
            parts = re.split(r"[;\|,\n]+", str(cell))
        out: list[str] = []
        for raw in parts:
            s = str(raw or "").strip()
            if not s:
                continue
            if re.match(r"^\s*#?\s*theme\s*:", s, flags=re.I):
                base = re.sub(r"^\s*#?\s*theme\s*:\s*", "", s, flags=re.I).strip()
                base = re.sub(r"\s+", " ", base)
                if base:
                    out.append(f"{base}")
        return out


    def _find_citations_col_local(frame: pd.DataFrame) -> str | None:
        pats = (r"^citations?$", r"^cites?$", r"times_?cited", r"^n_?cit", r"^num_?cit")
        for c in frame.columns:
            if any(re.search(p, str(c), flags=re.I) for p in pats):
                return c
        return None

    # ---------- ensure minimal fields: __year, __cit, __authors ----------
    if "__year" not in work.columns:
        def __mini_year(x):
            import re as _re, pandas as _pd, numpy as _np
            if isinstance(x, (int, float)) and 1800 <= int(x) <= 2100:
                return float(int(x))
            s = str(x) if x is not None and not _pd.isna(x) else ""
            s = s.strip()
            if not s:
                return _np.nan
            try:
                dt = _pd.to_datetime(s, errors="coerce", utc=True)
                if dt is not _pd.NaT and not _pd.isna(dt):
                    y = int(dt.year)
                    return float(y) if 1800 <= y <= 2100 else _np.nan
            except Exception:
                pass
            m = _re.search(r"(18|19|20)\d{2}", s)
            return float(int(m.group(0))) if m else _np.nan

        _ycands = ["year_numeric", "year", "Year", "pub_year", "publication_year",
                   "issued", "date", "Date", "publicationDate", "created"]
        _ysrc = next((c for c in _ycands if c in work.columns), None)
        if _ysrc is None:
            _ysrc = next((c for c in work.columns
                          if any(k in str(c).lower() for k in ("year", "date", "issued", "created"))), None)
        work["__year"] = work[_ysrc].apply(__mini_year) if _ysrc is not None else np.nan

    # normalise citations column
    cit_col = _find_citations_col(work) if "_find_citations_col" in globals() else _find_citations_col_local(work)
    if cit_col is None:
        work["__cit"] = pd.to_numeric(work.get("__cit", 0), errors="coerce").fillna(0).astype(int)
    else:
        work["__cit"] = pd.to_numeric(work[cit_col], errors="coerce").fillna(0).astype(int)

    # authors list
    if "__authors" not in work.columns:
        work["__authors"] = work.apply(_extract_authors_any, axis=1)

    # ---------- Top-5 authors by productivity ----------
    tmp = work.copy()
    tmp["__AuthorsList"] = tmp["__authors"]
    auth_rows = tmp.explode("__AuthorsList").rename(columns={"__AuthorsList": "Author"})
    auth_rows["Author"] = auth_rows["Author"].fillna("").astype(str).str.strip()
    auth_rows = auth_rows[auth_rows["Author"] != ""]
    author_counts = (auth_rows.groupby("Author").size()
                     .rename("Count").reset_index()
                     .sort_values(["Count", "Author"], ascending=[False, True]))
    top5_authors = author_counts.head(5)["Author"].tolist()

    # ---------- Build 25-doc payload: top-5 most cited per author ----------
    # ---------- Build 25-doc payload: top-5 most cited per author (with theme tags) ----------
    docs_payload, rows25 = [], []
    authors_present = []  # authors for whom we actually found docs with abstracts/titles

    # discover potential theme-bearing columns once
    theme_cols = [c for c in work.columns if re.search(r"(theme|keyword|tag|subject)", str(c), flags=re.I)]

    def _norm_theme_tok(s: str) -> str:
        s = re.sub(r"\s+", " ", str(s).strip())
        s = s.casefold()
        s = re.sub(r"[\u2010-\u2015-]+", "-", s)
        return s

    if top5_authors:
        for a in top5_authors:
            # docs for this author
            docs_a = work[work["__authors"].apply(lambda L: a in L if isinstance(L, list) else False)].copy()
            if docs_a.empty:
                continue
            # keep only rows with at least a title or abstract
            docs_a = docs_a[(docs_a.get("title").notna() | docs_a.get("abstract").notna())]
            if docs_a.empty:
                continue
            docs_a["__cit"] = pd.to_numeric(docs_a["__cit"], errors="coerce").fillna(0).astype(int)
            take = (docs_a.sort_values(["__cit", "title"], ascending=[False, True])
                    .head(5)
                    .loc[:, ["title", "abstract", "__cit"] + [c for c in theme_cols if c in docs_a.columns]])
            if take.empty:
                continue
            authors_present.append(a)
            for _, rr in take.iterrows():
                t = _safe(rr.get("title", ""))
                ab = _safe(rr.get("abstract", ""))[:900]  # clip to keep prompt size reasonable
                c = int(rr.get("__cit", 0))

                # collect + normalise + de-duplicate theme tokens across discovered columns
                theme_tokens: list[str] = []
                for col in theme_cols:
                    if col in rr:
                        theme_tokens.extend(_extract_theme_tags(rr.get(col)))

                seen: set[str] = set()
                normed: list[str] = []
                for tok in theme_tokens:
                    key = _norm_theme_tok(tok)
                    if key and key not in seen:
                        seen.add(key)
                        normed.append(tok.strip())

                themes_str = "; ".join(normed)  # semicolon-separated list of canonical theme terms

                docs_payload.append(f"{a}\t{t}\t{ab}\t{c}\t{themes_str}")
                rows25.append({"Author": a, "Title": t, "Citations": c, "Themes": themes_str})

    # If we have no payload, show placeholders (no invented themes)
    if not docs_payload or not authors_present:
        fig_heat_empty = go.Figure().add_annotation(text="No abstracts available for top-5 authors.",
                                                    showarrow=False)
        fig_heat_empty.update_layout(template="plotly_white", height=560, width=1400,
                                     title=f"{collection_name} · Slide 4a — Themes Heatmap")
        slide4a = _fig_to_slide(fig_heat_empty)

        fig_sankey_empty = go.Figure().add_annotation(text="No abstracts available to build author → theme links.",
                                                      showarrow=False)
        fig_sankey_empty.update_layout(template="plotly_white", height=560, width=1400,
                                       title=f"{collection_name} · Slide 4b — Author → Theme Structure")
        slide4b = _fig_to_slide(fig_sankey_empty)
        if slide_notes:
            slide4a.notes_slide.notes_text_frame.text = f"{collection_name} • No abstracts available to infer themes."
            slide4b.notes_slide.notes_text_frame.text = f"{collection_name} • No abstracts available to infer themes."
    else:
        # ---------- LLM function tool definition & router ----------
        # ---------- LLM function tool definition & router ----------
        # We FORCE the model to call this function with the computed results.

        # Authors present (for schema enums + instructions)
        authors_present = (
            top5_authors if isinstance(top5_authors, list) and top5_authors
            else sorted({dp.split("\t", 1)[0] for dp in docs_payload})
        )

        THEMES_TOOL = {
            "type": "function",
            "name": "themes_heatmap",
            "description": (
                "Return a small, unified theme set (8–10 max) derived from the provided 25 documents, "
                "with per-author normalised weights (each author's weights sum to 1.0 over the themes). "
                "Optionally include author→theme flows."
            ),
            "strict": True,
            "parameters": {
                "type": "object",
                "additionalProperties": False,
                "properties": {
                    "themes": {
                        "type": "array",
                        "minItems": 5,
                        "maxItems": 10,
                        "items": {"type": "string", "description": "Concise theme label (≤3 words)"}
                    },
                    "heatmap": {
                        "type": "object",
                        "additionalProperties": False,
                        "properties": {
                            "x": {  # themes order (should match 'themes')
                                "type": "array",
                                "minItems": 5,
                                "maxItems": 10,
                                "items": {"type": "string"}
                            },
                            "y": {  # authors order (restricted to detected authors)
                                "type": "array",
                                "minItems": max(1, len(authors_present)),
                                "items": {"type": "string", "enum": authors_present}
                            },
                            "z": {  # rows align to y; cols to x; values in [0,1]
                                "type": "array",
                                "items": {
                                    "type": "array",
                                    "items": {"type": "number", "minimum": 0.0, "maximum": 1.0}
                                }
                            }
                        },
                        "required": ["x", "y", "z"]
                    },
                    "flows": {
                        "type": ["array", "null"],
                        "items": {
                            "type": "object",
                            "additionalProperties": False,
                            "properties": {
                                "author": {"type": "string", "enum": authors_present},
                                "theme": {"type": "string"},
                                "weight": {"type": "number", "minimum": 0.0, "maximum": 1.0}
                            },
                            "required": ["author", "theme", "weight"]
                        }
                    }
                },
                "required": ["themes", "heatmap", "flows"]
            }
        }

        def _tool_themes_heatmap(**payload):
            """
            Router for the 'themes_heatmap' tool.
            We don't transform payload; returning an ACK string is enough — the model's
            arguments are captured by call_models_plots() and parsed below.
            """
            return json.dumps({"ok": True})

        tools = [THEMES_TOOL]
        router = {"themes_heatmap": _tool_themes_heatmap}

        # ---------- Build TSV payload & instructions ----------
        tsv = "Author\tTitle\tAbstract\tCitations\tThemes\n" + "\n".join(docs_payload)

        authors_clause = "; ".join(authors_present)
        sys_instructions = (
            "You MUST call the 'themes_heatmap' function exactly once with your final results. "
            "Do NOT output prose. "
            "Use concise theme labels (≤3 words), domain-appropriate, non-jargony. "
            "Normalise weights PER AUTHOR so each author's weights sum to 1.0 over the returned themes. "
            f"Restrict authors (heatmap.y) to this set and spellings: {authors_clause}. "
            "Ensure heatmap.x equals the chosen 'themes' (same order). "
            "If a theme is not present for an author, set the corresponding weight to 0."
        )

        contract = """
        You are given a TSV of 25 documents (top-5 most cited per each of the Top-5 productive authors).
        Columns: Author, Title, Abstract, Citations, Themes.

        • The 'Themes' column is optional, semicolon-separated prior tags already normalised (e.g., "mission-oriented; innovation policy").
          Use them as hints, but you must consolidate/merge synonyms into a single, unified theme set.

        Tasks:
        1) Derive a unified set of 8–10 themes across all documents. Keep labels short (≤3 words) and meaningful.
        2) Build a per-author heatmap where rows=authors and columns=themes; values are normalised weights in [0,1].
           Each row must sum to ~1.0 (±1e-3 tolerance). Missing signal → 0.
        3) Optionally provide a list of flows (author→theme with weight) useful for a Sankey diagram.

        Return your results ONLY by calling the 'themes_heatmap' function (no text reply).
        """

        prompt = (
            "From these documents, extract normalised cross-author themes.\n\n"
            f"{contract}\n\nTSV START\n{tsv}\nTSV END\n"
        )

        # ---------- Call the model with tools ----------
        llm_res = None
        llm_error = None
        llm_json = None  # <-- parsed function-call arguments (themes/heatmap/flows)

        try:
            if "call_models_plots" in globals():
                llm_res = call_models_plots(
                    mode="analyze",
                    prompt_text=prompt,
                    model_api_name="gpt-5-mini",
                    max_tokens=10000,
                    use_cache= True,
                    analysis_key_suffix="authors_slide4_top5x5_themes_tools",
                    section_title="authors",
                    overall_topic=collection_name,
                    tools=tools,
                    tool_router=router,
                    tool_choice={"type": "function", "name": "themes_heatmap"},  # FORCE exactly one call
                    allowed_tools=[{"type": "function", "name": "themes_heatmap"}],
                    parallel_tool_calls=False,
                    instructions=sys_instructions,
                    max_function_loops=1,
                )
                print(llm_res)
                # Extract the function-call arguments (the actual results)
                # Extract the function-call arguments (the actual results)
                # Extract the function-call arguments (the actual results)
                tcalls = (llm_res or {}).get("tool_calls", []) or []
                for tc in tcalls:
                    if tc.get("type") == "function_call" and tc.get("name") == "themes_heatmap":
                        args_raw = tc.get("arguments", {})
                        if isinstance(args_raw, dict):
                            llm_json = args_raw
                        elif isinstance(args_raw, str):
                            try:
                                llm_json = json.loads(args_raw)
                            except Exception:
                                llm_json = None
                        else:
                            llm_json = None
                        break

                if llm_json is None:
                    llm_error = "themes_heatmap call missing or arguments could not be parsed."



            else:
                llm_error = "call_models_plots() not available"
        except Exception as e:
            llm_error = f"LLM tool-calling error: {e}"

        # ---------- Extract the function-call payload ----------
        themes, authors_order, Z, flows = [], [], [], []
        if isinstance(llm_res, dict) and isinstance(llm_res.get("tool_calls"), list):
            calls = [tc for tc in llm_res["tool_calls"]
                     if tc.get("type") == "function_call" and tc.get("name") == "themes_heatmap"]
            if calls:
                args_raw = calls[-1].get("arguments", {}) or {}
                if isinstance(args_raw, dict):
                    args = args_raw
                elif isinstance(args_raw, str):
                    try:
                        args = json.loads(args_raw)
                    except Exception:
                        args = {}
                else:
                    args = {}

                hm = args.get("heatmap") or {}
                themes = list((hm.get("x") or args.get("themes") or []))[:10]
                authors_order = [a for a in (hm.get("y") or []) if a in authors_present] or authors_present[:]
                Z_raw = hm.get("z") or []

                R, C = len(authors_order), len(themes)
                if R and C and isinstance(Z_raw, list):
                    try:
                        Z = [[0.0] * C for _ in range(R)]
                        for i in range(min(R, len(Z_raw))):
                            row = Z_raw[i] or []
                            for j in range(min(C, len(row))):
                                v = row[j]
                                if v is None or (isinstance(v, float) and pd.isna(v)):
                                    v = 0.0
                                Z[i][j] = float(v)
                    except Exception:
                        Z = []

                if isinstance(args.get("flows"), list):
                    for it in args["flows"]:
                        try:
                            a = str(it["author"])
                            t = str(it["theme"])
                            w = float(it["weight"])
                            if a in authors_order and t in themes and w >= 0:
                                flows.append((a, t, w))
                        except Exception:
                            pass

        # If the tool-call is missing/invalid, render placeholders without inventing themes
        if not themes or not authors_order or not Z:
            msg = "Theme extraction tool call missing or invalid." + (f" ({llm_error})" if llm_error else "")
            fig_heat_empty = go.Figure().add_annotation(text=msg, showarrow=False)
            fig_heat_empty.update_layout(template="plotly_white", height=560, width=1400,
                                         title=f"{collection_name} · Slide 4a — Themes Heatmap")
            slide4a = _fig_to_slide(fig_heat_empty)

            fig_sankey_empty = go.Figure().add_annotation(text=msg, showarrow=False)
            fig_sankey_empty.update_layout(template="plotly_white", height=560, width=1400,
                                           title=f"{collection_name} · Slide 4b — Author → Theme Structure")
            slide4b = _fig_to_slide(fig_sankey_empty)
            if slide_notes:
                slide4a.notes_slide.notes_text_frame.text = f"{collection_name} • {msg}"
                slide4b.notes_slide.notes_text_frame.text = f"{collection_name} • {msg}"
        else:
            # Sanity: clamp + soft renormalise each author row to sum≈1.0
            R, C = len(authors_order), len(themes)
            for i in range(R):
                Z[i] = [max(0.0, float(v)) for v in Z[i]]
                s = sum(Z[i])
                if s > 0:
                    # if it's far from 1, renormalise
                    if abs(s - 1.0) > 1e-3:
                        Z[i] = [v / s for v in Z[i]]

            # percentage view if data within ~[0..1.5]
            zmax = max((max(r) if r else 0.0) for r in Z) if Z else 1.0
            as_pct = zmax <= 1.5
            Z_plot = [[(v * 100.0 if as_pct else float(v)) for v in row] for row in Z]

            # ---------- Slide 4a: HEATMAP ----------
            fig_heat = go.Figure(
                go.Heatmap(
                    z=Z_plot, x=themes, y=authors_order,
                    colorbar=dict(title="Share (%)" if as_pct else "Weight"),
                    hovertemplate="Author: %{y}<br>Theme: %{x}<br>Value: %{z:.1f}" + (
                        "%" if as_pct else "") + "<extra></extra>",
                )
            )
            fig_heat.update_layout(
                template="plotly_white",
                height=640, width=1400,
                title=dict(text=f"{collection_name} · Slide 4a — Theme intensity across Top-5 authors (normalised)",
                           x=0.02, xanchor="left"),
                margin=dict(l=80, r=40, t=80, b=80),
                xaxis_title="Theme",
                yaxis_title="Author",
                font=dict(family="Calibri", size=12, color="#2A2E33"),
            )
            slide4a = _fig_to_slide(fig_heat, height_px=640)

            if slide_notes:
                df_heat_numeric = pd.DataFrame(Z_plot)
                a_order_txt = "; ".join(authors_order)
                t_order_txt = "; ".join(themes)
                _notes_text = _safe_ai_notes(
                    fig_heat,
                    _ensure_numeric(df_heat_numeric if isinstance(df_heat_numeric, pd.DataFrame) else pd.DataFrame(),
                                    fallback_cols=("Value",), min_rows=1),
                    collection_name=collection_name,
                    analysis_suffix="authors_slide4a_heatmap_tools",
                    title_for_prompt=f"{collection_name} · Slide 4a — Themes heatmap",
                    prompt_purpose=(
                            "Rows correspond to authors (order): " + a_order_txt + ". "
                                                                                   "Columns correspond to themes (order): " + t_order_txt + ". "
                                                                                                                                            "Values are normalised shares per author. Analyse concentration vs. spread, standout specialisations, "
                                                                                                                                            "and overlaps; add 2–3 validation steps (theme label sanity, leakage across labels, author disambiguation)."
                    ),
                    max_tokens=900,
                    effort="low",
                    base_logs_dir=None,
                    _cb=_cb,
                )
                _set_notes(slide=slide4a, slide_notes=slide_notes, text=_notes_text)

            # ---------- Slide 4b: STRUCTURE (Sankey from tool payload) ----------
            # Use flows from the tool-call if provided; otherwise derive top-3 per author from Z (still tool-based values)
            flows_plot = []
            if flows:
                flows_plot = [(a, t, (w * 100.0 if as_pct else w)) for (a, t, w) in flows if
                              a in authors_order and t in themes]
            else:
                for i, a in enumerate(authors_order):
                    if i >= len(Z_plot): break
                    pairs = sorted(list(zip(themes, Z_plot[i])), key=lambda t: t[1], reverse=True)[:3]
                    for tlabel, w in pairs:
                        if w > 0:
                            flows_plot.append((a, tlabel, float(w)))

            author_nodes = authors_order
            theme_nodes = themes
            node_labels = author_nodes + theme_nodes
            if px:
                node_colors = (
                        px.colors.qualitative.Set2[:len(author_nodes)]
                        + px.colors.qualitative.Pastel1[:len(theme_nodes)]
                )
            else:
                node_colors = ["#4C78A8", "#F58518", "#E45756", "#72B7B2", "#54A24B"][:len(author_nodes)] + \
                              ["#EECA3B", "#B279A2", "#FF9DA6", "#9D755D", "#BAB0AC"][:len(theme_nodes)]

            index_map = {lbl: i for i, lbl in enumerate(node_labels)}
            S, T, V = [], [], []
            for a, t, w in flows_plot:
                if a in index_map and t in index_map and w > 0:
                    S.append(index_map[a]);
                    T.append(index_map[t]);
                    V.append(float(w))
            if not V:
                S, T, V = [0], [len(author_nodes)], [1.0]

            fig_sankey = go.Figure(
                go.Sankey(
                    node=dict(label=node_labels, color=node_colors, pad=12, thickness=14,
                              line=dict(width=0.5, color="rgba(0,0,0,0.2)")),
                    link=dict(source=S, target=T, value=V,
                              hovertemplate="%{source.label} → %{target.label}<br>Value: %{value:.1f}" + (
                                  "%" if as_pct else "") + "<extra></extra>"),
                )
            )
            fig_sankey.update_layout(
                template="plotly_white",
                height=640, width=1400,
                title=dict(text=f"{collection_name} · Slide 4b — Author → Theme structure", x=0.02, xanchor="left"),
                margin=dict(l=40, r=40, t=80, b=40),
                font=dict(family="Calibri", size=12, color="#2A2E33"),
            )
            slide4b = _fig_to_slide(fig_sankey, height_px=640)

            if slide_notes:
                try:
                    flows_df = pd.DataFrame(flows_plot, columns=["Author", "Theme", "Weight"])
                    flows_df["Weight"] = pd.to_numeric(flows_df["Weight"], errors="coerce").fillna(0.0)
                    flows_df = flows_df.sort_values(["Author", "Weight"], ascending=[True, False]).groupby(
                        "Author").head(5)

                # ---- Slide 4a notes (LLM analysis over the produced heatmap data) ----

                    try:
                        # Use the *normalised* 0–1 matrix for determinism; round to stabilise cache keys.
                        heat_payload = {
                            "authors": list(authors_order),
                            "themes": list(themes),
                            "matrix": [[round(float(v), 4) for v in row] for row in Z],  # 0–1 weights
                        }
                        prompt4a = (
                                f"{collection_name} · Slide 4a — Themes heatmap\n"
                                "You are given JSON with authors, themes, and a row-normalised matrix of weights (0–1). "
                                "Write 4–6 crisp bullets (≤25 words each) that:\n"
                                "• Assess concentration vs spread per author\n"
                                "• Identify standout specialisations and cross-author overlaps\n"
                                "• Note any caveats (label leakage, small-N)\n"
                                "Return ONLY the bullets (no title, no preamble).\n\n"
                                "JSON:\n" + json.dumps(heat_payload, ensure_ascii=False, sort_keys=True)
                        )
                        res4a = call_models_plots(
                            mode="analyze",
                            prompt_text=prompt4a,
                            model_api_name="gpt-5-mini",
                            max_tokens=700,
                            use_cache=True,
                            analysis_key_suffix="authors_slide4a_heatmap_notes",
                            section_title="authors",
                            overall_topic=collection_name,
                        )
                        slide4a.notes_slide.notes_text_frame.text = (res4a or {}).get("text", "").strip() or "—"
                    except Exception as _e:
                        slide4a.notes_slide.notes_text_frame.text = f"(notes generation failed: {type(_e).__name__})"

                    # ---- Slide 4b notes (LLM analysis over the produced flows structure) ----
                    if slide_notes:
                        try:
                            # Convert flows to a compact, stable JSON (weights back to 0–1 if needed).
                            flows_json = [
                                {"author": a, "theme": t,
                                 "weight": round(float(w) / 100.0 if as_pct else float(w), 4)}
                                for (a, t, w) in flows_plot
                            ]
                            struct_payload = {
                                "authors": list(authors_order),
                                "themes": list(themes),
                                "flows": flows_json,  # author→theme weights in 0–1
                            }
                            prompt4b = (
                                    f"{collection_name} · Slide 4b — Author → Theme structure\n"
                                    "You are given JSON with authors, themes, and author→theme weights (0–1). "
                                    "Write 4–6 bullets (≤25 words each) that:\n"
                                    "• Summarise dominant themes per author and diversification vs concentration\n"
                                    "• Highlight overlaps/bridges between authors\n"
                                    "• Suggest 2–3 validation checks (label leakage, disambiguation, age/citation bias)\n"
                                    "Return ONLY the bullets (no title, no preamble).\n\n"
                                    "JSON:\n" + json.dumps(struct_payload, ensure_ascii=False, sort_keys=True)
                            )
                            res4b = call_models_plots(
                                mode="analyze",
                                prompt_text=prompt4b,
                                model_api_name="gpt-4.1-mini",
                                max_tokens=700,
                                use_cache=True,
                                analysis_key_suffix="authors_slide4b_structure_notes",
                                section_title="authors",
                                overall_topic=collection_name,
                            )
                            slide4b.notes_slide.notes_text_frame.text = (res4b or {}).get("text", "").strip() or "—"
                        except Exception as _e:
                            slide4b.notes_slide.notes_text_frame.text = f"(notes generation failed: {type(_e).__name__})"


                except Exception:

                    flows_df = pd.DataFrame()

                if slide_notes:

                    # numeric-only, never-empty payload to keep the handler safe

                    try:

                        flows_num = pd.DataFrame({

                            "Weight": pd.to_numeric(

                                flows_df["Weight"] if "Weight" in flows_df.columns else pd.Series([], dtype=float),

                                errors="coerce"

                            ).fillna(0.0)

                        })

                    except Exception:

                        flows_num = pd.DataFrame({"Weight": [0.0]})

                    if flows_num.shape[0] == 0:
                        flows_num = pd.DataFrame({"Weight": [0.0]})

                    _notes_text = _safe_ai_notes(
                        fig_sankey,
                        _ensure_numeric(flows_num if isinstance(flows_num, pd.DataFrame) else pd.DataFrame(),
                                        fallback_cols=("Weight",), min_rows=1),
                        collection_name=collection_name,
                        analysis_suffix="authors_slide4b_structure_tools",
                        title_for_prompt=f"{collection_name} · Slide 4b — Author → Theme structure",
                        prompt_purpose=(
                            "Describe which themes dominate for each author and where overlaps occur. "
                            "Call out any authors with highly concentrated portfolios (one dominant theme) versus diversified ones. "
                            "Suggest 2–3 next checks (theme leakage, author disambiguation, citation-age normalisation)."
                        ),
                        max_tokens=900,
                        effort="low",
                        base_logs_dir=None,
                        _cb=_cb,
                    )
                    _set_notes(slide=slide4b, slide_notes=slide_notes, text=_notes_text)

        # ======================================================================
    # Slide 5 (UPDATED): Scatter of Top-10 most collaborative authors
    # ======================================================================
        # ======================================================================
        # Slide 5 (UPDATED): Scatter of Top-10 most collaborative authors
        # ======================================================================

        import re
        import math  # ← needed for sqrt/triangles below
        import plotly.graph_objects as go
        import plotly.express as px
        import networkx as nx
        import pandas as pd

        # ---- guards & fallbacks so this block works regardless of slide order ----

        # 1) Ensure citations column exists
        if "__cit" not in work.columns:
            try:
                cit_col = _find_citations_col(work)  # defined earlier in this file (if present)
            except Exception:
                cit_col = None
            if cit_col is None:
                work["__cit"] = 0
            else:
                work["__cit"] = pd.to_numeric(work[cit_col], errors="coerce").fillna(0).astype(int)

        # 2) Ensure authors list exists and is ALWAYS list[str]
        def _authors_list_row(row) -> list[str]:
            # structured list of dicts (common in Zotero exports)
            if isinstance(row.get("authors_list"), (list, tuple)) and row["authors_list"]:
                out = []
                for a in row["authors_list"]:
                    if isinstance(a, dict):
                        ln, fn = a.get("lastName"), a.get("firstName")
                        if ln and fn:
                            out.append(f"{ln}, {fn}")
                        elif ln:
                            out.append(str(ln))
                        elif fn:
                            out.append(str(fn))
                    else:
                        s = str(a).strip()
                        if s: out.append(s)
                return [s for s in out if s]

            # flat strings (split on ; or ,)
            for key in ("authors", "creator_summary"):
                s = row.get(key)
                if isinstance(s, str) and s.strip():
                    return [p.strip() for p in re.split(r"[;|,]+", s) if p.strip()]

            # already a list?
            if isinstance(row.get("__authors"), (list, tuple)):
                return [str(p).strip() for p in row["__authors"] if str(p).strip()]

            return []

        if "__authors" in work.columns:
            # Hard-normalize whatever is there into list[str] (prevents scalar NaN/float surprises)
            work["__authors"] = work["__authors"].apply(
                lambda v:
                ([str(p).strip() for p in v if str(p).strip()]  # list/tuple
                 if isinstance(v, (list, tuple)) else
                 ([p.strip() for p in re.split(r"[;|,]+", v) if p.strip()]  # split string
                  if isinstance(v, str) and v.strip() else
                  []))
            )
        else:
            work["__authors"] = work.apply(_authors_list_row, axis=1)

        # 3) Build a local publication count dict (fixes earlier 'pubs' reference)
        _pubs: dict[str, int] = {}
        for _, _row in work.iterrows():
            A = [a for a in (_row.get("__authors") or []) if a]
            for a in A:
                _pubs[a] = _pubs.get(a, 0) + 1

        # -------------------- original styling/logic preserved below --------------------

        # Tunables for clarity
        EDGE_DRAW_MAX = 4
        MAX_GROUPS = 20
        TOP_LABELS = 15
        K_CORE_MIN = 2
        SEED = 42
        FONT_FAMILY = "Calibri"
        TITLE_FONT_FAMILY = "Calibri Light"

        # Defaults for title/notes
        top_node_name = "—"
        top_interactions = 0

        # 1) Weighted edge list (only items with >= 2 authors)
        edges: dict[tuple[str, str], int] = {}
        for _, row in work.iterrows():
            A = row.get("__authors") or []
            if not isinstance(A, list):  # extra safety
                A = []
            if len(A) < 2:
                continue
            for i in range(len(A) - 1):
                for j in range(i + 1, len(A)):
                    u, v = sorted((A[i], A[j]))
                    edges[(u, v)] = edges.get((u, v), 0) + 1

        # 2) Graph and attributes
        G = nx.Graph()
        for (u, v), w in edges.items():
            G.add_edge(u, v, weight=int(w))
        for a, p in _pubs.items():
            if a not in G:
                G.add_node(a)
            G.nodes[a]["pubs"] = int(p)

        if G.number_of_nodes() == 0:
            fig_net = go.Figure().add_annotation(text="No co-authorship data.", showarrow=False)
            fig_net.update_layout(template="plotly_white")
        else:
            # Weighted degree = interactions
            interactions_full = {
                n: int(sum(d.get("weight", 1) for _, _, d in G.edges(n, data=True)))
                for n in G.nodes()
            }
            nx.set_node_attributes(G, interactions_full, "interactions")

            # Community detection
            try:
                from networkx.algorithms.community import greedy_modularity_communities
                comms_full = list(greedy_modularity_communities(G)) if G.number_of_edges() > 0 else []
            except Exception:
                comms_full = []
            if not comms_full:
                comms_full = [set(G.nodes())]

            # Keep largest communities
            comms_sorted = sorted(comms_full, key=lambda s: len(s), reverse=True)
            comms_keep = comms_sorted[:min(MAX_GROUPS, len(comms_sorted))]
            keep_nodes = set().union(*comms_keep) if comms_keep else set()

            # Induced subgraph
            H = G.subgraph(keep_nodes).copy() if keep_nodes else nx.Graph()

            # Optional pruning: k-core (+ keep global hubs)
            if H.number_of_nodes() > 0:
                try:
                    Hk = nx.k_core(H, k=K_CORE_MIN)
                    top_global_hubs = sorted(
                        H.nodes(),
                        key=lambda n: (interactions_full.get(n, 0), H.nodes[n].get("pubs", 0), str(n)),
                        reverse=True
                    )[:10]
                    H_nodes_final = set(Hk.nodes()) | set(top_global_hubs)
                    H = H.subgraph(H_nodes_final).copy()
                except Exception:
                    pass

            if H.number_of_nodes() == 0:
                fig_net = go.Figure().add_annotation(text="No co-authorship data (after filtering).", showarrow=False)
                fig_net.update_layout(template="plotly_white")
            else:
                # Recompute interactions within H
                interactions = {
                    n: int(sum(d.get("weight", 1) for _, _, d in H.edges(n, data=True)))
                    for n in H.nodes()
                }
                nx.set_node_attributes(H, interactions, "interactions")

                # Community IDs
                community_id: dict[str, int] = {}
                for ci, com in enumerate(comms_keep):
                    for n in com:
                        if n in H:
                            community_id[n] = ci
                nx.set_node_attributes(H, community_id, "community")

                # Main hub for title/notes
                def _hub_key(n: str) -> tuple[int, int, str]:
                    return (interactions.get(n, 0), H.nodes[n].get("pubs", 0), str(n))

                top_node_name = max(H.nodes(), key=_hub_key)
                top_interactions = interactions.get(top_node_name, 0)

                # 3) Community-aware layout
                num_groups = max(1, len(comms_keep))
                R = 10.0 + 2.0 * math.sqrt(num_groups)
                centroids: dict[int, tuple[float, float]] = {}
                for i in range(num_groups):
                    ang = 2 * math.pi * i / num_groups
                    centroids[i] = (R * math.cos(ang), R * math.sin(ang))
                pos: dict[str, tuple[float, float]] = {}
                for ci, com in enumerate(comms_keep):
                    nodes_ci = [n for n in com if n in H]
                    if not nodes_ci:
                        continue
                    sub = H.subgraph(nodes_ci)
                    scale = 1.5 + 0.1 * math.sqrt(len(nodes_ci))
                    sub_pos = nx.spring_layout(
                        sub, seed=SEED + ci, k=0.8 / max(1, math.sqrt(len(nodes_ci))),
                        iterations=300, weight="weight"
                    )
                    cx, cy = centroids.get(ci, (0.0, 0.0))
                    for n, (x, y) in sub_pos.items():
                        pos[n] = (cx + scale * x, cy + scale * y)

                # 4) Edge trace
                xe, ye = [], []
                for u, v, d in H.edges(data=True):
                    w = int(d.get("weight", 1))
                    if w > EDGE_DRAW_MAX:
                        continue
                    if u in pos and v in pos:
                        xe += [pos[u][0], pos[v][0], None]
                        ye += [pos[u][1], pos[v][1], None]
                edge_trace = go.Scatter(
                    x=xe, y=ye, mode="lines",
                    line=dict(width=0.9, color="rgba(120,124,130,0.28)"),
                    hoverinfo="none", showlegend=False
                )

                # 5) Nodes per community (original style)
                palette = (
                        px.colors.qualitative.Set2
                        + px.colors.qualitative.Pastel1
                        + px.colors.qualitative.Bold
                        + px.colors.qualitative.Safe
                )
                community_nodes: dict[int, list[str]] = {}
                for n in H.nodes():
                    community_nodes.setdefault(community_id.get(n, 0), []).append(n)

                ints = list(interactions.values())
                vmin = min(ints) if ints else 0
                vmax = max(ints) if ints else 1

                def _scale_size(v: int) -> float:
                    if vmax == vmin:
                        return 18.0
                    sv = math.sqrt(v)
                    smin, smax = math.sqrt(vmin), math.sqrt(vmax)
                    t = (sv - smin) / (smax - smin + 1e-9)
                    return 14.0 + 32.0 * max(0.0, min(1.0, t))

                node_traces = []
                for ci, nodes_c in sorted(community_nodes.items(), key=lambda t: (len(t[1]), t[0]), reverse=True):
                    nodes_sorted = sorted(nodes_c, key=_hub_key, reverse=True)
                    x = [pos[n][0] for n in nodes_sorted if n in pos]
                    y = [pos[n][1] for n in nodes_sorted if n in pos]
                    sizes = [_scale_size(interactions.get(n, 0)) for n in nodes_sorted if n in pos]
                    hover = [
                        f"<b>{n}</b>"
                        f"<br>Interactions: {interactions.get(n, 0)}"
                        f"<br>Publications: {H.nodes[n].get('pubs', 0)}"
                        f"<br>Group: {ci + 1}"
                        for n in nodes_sorted if n in pos
                    ]
                    if not x:  # skip empty groups defensively
                        continue
                    node_traces.append(
                        go.Scatter(
                            x=x, y=y, mode="markers",
                            name=f"Group {ci + 1}",
                            hovertext=hover, hoverinfo="text",
                            marker=dict(
                                size=sizes,
                                color=palette[ci % len(palette)],
                                opacity=0.94,
                                line=dict(
                                    width=[2.2 if n == top_node_name else 1.1 for n in nodes_sorted if n in pos],
                                    color="white"
                                ),
                            ),
                            showlegend=True,
                        )
                    )

                # 6) Labels
                group_leaders = {max(nodes, key=_hub_key) for _, nodes in community_nodes.items() if nodes}
                remaining = [n for n in sorted(H.nodes(), key=_hub_key, reverse=True) if n not in group_leaders]
                label_nodes = list(group_leaders | {top_node_name})
                for n in remaining:
                    if len(label_nodes) >= TOP_LABELS:
                        break
                    label_nodes.append(n)

                lx = [pos[n][0] for n in label_nodes if n in pos]
                ly = [pos[n][1] for n in label_nodes if n in pos]
                ltext = [n for n in label_nodes if n in pos]
                label_trace = go.Scatter(
                    x=lx,
                    y=ly,
                    mode="text",
                    text=ltext,
                    textposition="middle right",
                    textfont=dict(size=12),
                    hoverinfo="skip",
                    showlegend=False,
                    cliponaxis=False,
                )

                # 7) Compose figure
                fig_net = go.Figure(data=[edge_trace, *node_traces, label_trace])
                fig_net.update_layout(
                    template="plotly_white",
                    height=780,
                    width=1400,
                    margin=dict(l=20, r=20, t=98, b=20),
                    xaxis=dict(visible=False),
                    yaxis=dict(visible=False),
                    font=dict(family=FONT_FAMILY, size=12, color="#303235"),
                    title_font=dict(family=TITLE_FONT_FAMILY, size=20, color="#181A1B"),
                    legend=dict(
                        orientation="h",
                        yanchor="bottom",
                        y=1.02,
                        xanchor="right",
                        x=1.0,
                        bgcolor="rgba(255,255,255,0.7)",
                        bordercolor="rgba(0,0,0,0.08)",
                        borderwidth=1,
                        font=dict(size=11),
                    ),
                )

        # Insert slide 5
        slide5 = _fig_to_slide(fig_net, height_px=780)

        # Notes for slide 5 — use global helpers to guarantee numeric payload & safe write
        if slide_notes:
            # 1) Build the full (string+numeric) table as before, but tolerate missing locals
            try:
                nodes_df = pd.DataFrame([
                    {
                        "Author": n,
                        "Group": (int(community_id.get(n, -1)) + 1)
                        if ('community_id' in locals() and n in community_id) else -1,
                        "Interactions": int(interactions.get(n, 0)) if 'interactions' in locals() else 0,
                        "Publications": int(H.nodes[n].get("pubs", 0))
                        if ('H' in locals() and isinstance(H, nx.Graph) and n in H) else 0,
                    }
                    for n in (H.nodes() if ('H' in locals() and isinstance(H, nx.Graph)) else [])
                ])
            except Exception:
                nodes_df = pd.DataFrame(columns=["Author", "Group", "Interactions", "Publications"])

            if not nodes_df.empty and "Interactions" in nodes_df.columns:
                nodes_df = nodes_df.sort_values("Interactions", ascending=False, kind="mergesort")

            # 2) Guaranteed numeric-only frame (≥1 row) using helper
            try:
                nodes_num = pd.DataFrame({
                    "Interactions": pd.to_numeric(nodes_df.get("Interactions", pd.Series([], dtype=float)),
                                                  errors="coerce"),
                    "Publications": pd.to_numeric(nodes_df.get("Publications", pd.Series([], dtype=float)),
                                                  errors="coerce"),
                })
            except Exception:
                nodes_num = pd.DataFrame({"Interactions": [], "Publications": []})
            nodes_num = _ensure_numeric(nodes_num, fallback_cols=("Interactions", "Publications"), min_rows=1)

            # 3) Compact text context (kept out of numeric table)
            top_ctx = []
            try:
                for _, r in (nodes_df.head(12) if not nodes_df.empty else pd.DataFrame()).iterrows():
                    a = str(r.get("Author", "—"))
                    g_val = r.get("Group", -1)
                    g = int(-1 if pd.isna(g_val) else int(g_val))
                    I = int(0 if pd.isna(r.get("Interactions", 0)) else int(r.get("Interactions", 0)))
                    P = int(0 if pd.isna(r.get("Publications", 0)) else int(r.get("Publications", 0)))
                    top_ctx.append(f"{a} (G{g}): {I} interactions, {P} pubs")
            except Exception:
                pass
            ctx_text = "; ".join(top_ctx) if top_ctx else "—"

            # 4) Ask notes generator safely with the SAME prompt text
            _notes_text = _safe_ai_notes(
                fig_net,
                nodes_num,
                collection_name=collection_name,
                analysis_suffix="authors_slide5_network",
                title_for_prompt=(
                    f"{collection_name} · Co-authorship Network (largest "
                    f"{len(comms_keep) if 'comms_keep' in locals() else 'N/A'} groups)"
                ),
                prompt_purpose=(
                        "Explain how the network is built: nodes=authors; edges=co-authorship counts; weighted degree=interactions "
                        "(sum of tie weights). Clarify filters (k-core pruning, heavy-edge suppression) and greedy-modularity "
                        "communities. Interpret the biggest hubs and potential bridge nodes, and what higher interactions imply "
                        "(durable, repeated teaming vs. sheer headcount). Critically assess risks of insularity and propose "
                        "3–4 concrete collaboration strategies and validation checks. "
                        "Top hubs context (name/group/interactions/pubs): " + ctx_text
                ),
                max_tokens=2200,
                effort="low",
                base_logs_dir=None,
                _cb=_cb,
            )
            _set_notes(slide=slide5, slide_notes=slide_notes, text=_notes_text)

            # ======================================================================
            # Slide 6 (NEW): Community detail — sizes, densities, and top authors
            # ======================================================================
            _cb("Building community detail slide (sizes, density, top authors)…")

            # guard: if H isn't available (e.g., empty), skip gracefully
            if 'H' in locals() and isinstance(H, nx.Graph) and H.number_of_nodes() > 0:
                # Build per-community summary for the communities shown in Slide 3 (comms_keep)
                group_rows = []
                bar_groups = []
                bar_sizes = []
                bar_colors = []

                # palette consistent with Slide 3
                palette = (
                        px.colors.qualitative.Set2
                        + px.colors.qualitative.Pastel1
                        + px.colors.qualitative.Bold
                        + px.colors.qualitative.Safe
                )

                # helper for ranking within a community
                def _hub_key_local(n: str) -> tuple[int, int, str]:
                    return (H.nodes[n].get("interactions", 0), H.nodes[n].get("pubs", 0), str(n))

                for ci, com in enumerate(comms_keep):
                    nodes_ci = [n for n in com if n in H]
                    if not nodes_ci:
                        continue

                    sub = H.subgraph(nodes_ci)
                    n_nodes = sub.number_of_nodes()
                    n_edges = sub.number_of_edges()
                    # unweighted density (simple graph)
                    try:
                        density = nx.density(sub)
                    except Exception:
                        density = 0.0
                    # total collaboration weight within the community
                    total_w = int(sum(d.get("weight", 1) for _, _, d in sub.edges(data=True)))

                    # top authors by interactions within the community
                    top_authors = sorted(nodes_ci, key=_hub_key_local, reverse=True)[:5]
                    top_authors_str = "; ".join(
                        f"{a} ({H.nodes[a].get('interactions', 0)})" for a in top_authors
                    )

                    # leading author
                    leader = top_authors[0] if top_authors else "—"

                    group_label = f"Group {ci + 1}"
                    group_rows.append([
                        group_label,
                        str(n_nodes),
                        str(n_edges),
                        f"{density:.2f}",
                        str(total_w),
                        leader,
                        top_authors_str,
                    ])

                    bar_groups.append(group_label)
                    bar_sizes.append(n_nodes)
                    bar_colors.append(palette[ci % len(palette)])

                # If nothing to show, emit a simple note
                if not group_rows:
                    fig_groups = go.Figure().add_annotation(
                        text="No community details available after filtering.",
                        showarrow=False
                    )
                    fig_groups.update_layout(template="plotly_white", height=600, width=1400)
                else:
                    # Build a two-panel figure: left = group sizes bar, right = detail table
                    from plotly.subplots import make_subplots
                    panel4 = make_subplots(
                        rows=1, cols=2,
                        column_widths=[0.46, 0.54],
                        specs=[[{"type": "xy"}, {"type": "table"}]],
                        subplot_titles=("Largest Communities (by nodes)", "Community Details")
                    )

                    # Left: horizontal bar of community sizes
                    panel4.add_trace(
                        go.Bar(
                            x=bar_sizes,
                            y=bar_groups,
                            orientation="h",
                            marker=dict(color=bar_colors, line=dict(width=0.5, color="rgba(0,0,0,0.15)")),
                            text=[str(v) for v in bar_sizes],
                            textposition="auto",
                            hovertemplate=" %{y}<br>Nodes: %{x}<extra></extra>",
                            showlegend=False,
                        ),
                        row=1, col=1
                    )
                    panel4.update_xaxes(title_text="Nodes", row=1, col=1)
                    panel4.update_yaxes(title_text=None, row=1, col=1)

                    # Right: detail table
                    header_vals = ["Group", "Nodes", "Edges", "Density", "Total weight", "Leader",
                                   "Top 5 authors (interactions)"]
                    panel4.add_trace(
                        go.Table(
                            header=dict(
                                values=header_vals,
                                fill_color="rgba(0,0,0,0.05)",
                                align="left",
                                font=dict(size=12, family=FONT_FAMILY, color="#111"),
                                line_color="rgba(0,0,0,0.1)",
                            ),
                            cells=dict(
                                values=list(zip(*group_rows))),
                        ),
                        row=1, col=2
                    )

                    panel4.update_layout(
                        template="plotly_white",
                        height=900,
                        width=1400,
                        title_text=f"{collection_name} · Community Overview (shown subset)",
                        margin=dict(l=40, r=40, t=80, b=40),
                        font=dict(family=FONT_FAMILY, size=12, color="#303235"),
                        title_font=dict(family=TITLE_FONT_FAMILY, size=20, color="#181A1B"),
                    )
                    fig_groups = panel4

                slide4 = _fig_to_slide(fig_groups, height_px=900,
                                       )

                if slide_notes:
                    # Build full dataframe (strings allowed) for context
                    try:
                        groups_df = pd.DataFrame(
                            group_rows,
                            columns=["Group", "Nodes", "Edges", "Density", "Total weight", "Leader",
                                     "Top 5 authors (interactions)"]
                        ) if group_rows else pd.DataFrame()
                    except Exception:
                        groups_df = pd.DataFrame()

                    # NUMERIC-ONLY frame via helper (ALWAYS ≥1 row)
                    try:
                        tmp = groups_df.reindex(columns=["Nodes", "Edges", "Density", "Total weight"]).copy()
                        tmp["Nodes"] = pd.to_numeric(tmp["Nodes"], errors="coerce")
                        tmp["Edges"] = pd.to_numeric(tmp["Edges"], errors="coerce")
                        tmp["Density"] = pd.to_numeric(tmp["Density"], errors="coerce")
                        tmp["Total weight"] = pd.to_numeric(tmp["Total weight"], errors="coerce")
                    except Exception:
                        tmp = pd.DataFrame({"Nodes": [], "Edges": [], "Density": [], "Total weight": []})
                    groups_num = _ensure_numeric(tmp, fallback_cols=("Nodes", "Edges", "Density", "Total weight"),
                                                 min_rows=1)

                    # Compact text context kept OUT of the numeric table
                    ctx = ""
                    if not groups_df.empty:
                        try:
                            ctx = "; ".join(
                                f"{row['Group']}: leader {row['Leader']}, nodes {row['Nodes']}, edges {row['Edges']}"
                                for _, row in groups_df.iterrows()
                            )
                        except Exception:
                            ctx = ""

                    _grp_notes = _safe_ai_notes(
                        fig_groups,
                        groups_num,
                        collection_name=collection_name,
                        analysis_suffix="authors_slide4_groups",
                        title_for_prompt=f"{collection_name} · Community overview (sizes, density, leaders)",
                        prompt_purpose=(
                                "Interpret communities by size, density, and total collaboration weight. "
                                "Explain clearly what higher interactions imply (durable, repeated teaming vs. sheer headcount), "
                                "identify hubs and potential bridges, critically assess insularity risks, and give actionable next steps. "
                                "Context: " + (ctx or "—")
                        ),
                        max_tokens=2200,
                        effort="low",
                        base_logs_dir=None,
                        _cb=_cb,
                    )
                    _set_notes(slide=slide4, slide_notes=slide_notes, text=_grp_notes)


# --- REPLACE ENTIRE FUNCTION: add_temporal_trends_slides ---------------------
def add_Metadata_temporal_trends_slides(
    prs,
    df: pd.DataFrame,
    *,
    top_k_sources: int = 20,
    top_k_keywords: int = 20,
    slide_notes: bool = True,
    image_width_inches: float = 10.0,  # unused here but kept for signature compat
    progress_callback=None,
):
    """
    Build ONE slide with a 1×2 grid of plots:

        ┌───────────────────────────────┬───────────────────────────────┐
        │ Top Sources (overall count)   │ Top Controlled Vocabulary     │
        └───────────────────────────────┴───────────────────────────────┘

    (Replaces the former 2×2 temporal layout.)
    """
    import io, re
    import pandas as pd
    import plotly.express as px
    from pptx.util import Inches

    def _cb(msg: str):
        if progress_callback:
            try:
                progress_callback(msg)
            except Exception:
                pass

    if df is None or df.empty:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title & Content
        slide.shapes.title.text = "Top Sources & Vocabulary"
        slide.placeholders[1].text_frame.text = "No data available."
        return

    # ---------- helpers ----------
    def _plotly_to_png(fig, width=1200, height=700, scale=2) -> bytes:
        try:
            return fig.to_image(format="png", width=width, height=height, scale=scale)
        except Exception:
            import plotly.io as pio
            return pio.to_image(fig, format="png", width=width, height=height, scale=scale, engine="kaleido")

    def _place_in_cell(slide, png: bytes, cell_left, cell_top, cell_w, cell_h):
        """Place an image so it fits within a cell (centered, preserving AR)."""
        pic = slide.shapes.add_picture(io.BytesIO(png), left=cell_left, top=cell_top, width=cell_w)
        if pic.height > cell_h:  # too tall → fit by height
            pic.height = cell_h
            pic.left = cell_left + int((cell_w - pic.width) / 2)  # center horizontally
        else:
            pic.top = cell_top + int((cell_h - pic.height) / 2)   # center vertically
        return pic

    # ======================================================================
    # Figure A: Top Sources (overall)
    # ======================================================================
    _cb("Building Top Sources bar chart")

    def _source_is_placeholder(v):
        s = str(v or "").strip().upper()
        return s in {"", "N/A", "NA", "NONE", "FALSE"}

    if "source" in df.columns:
        src = df[["source"]].copy()
        src["source"] = src["source"].astype(str).str.strip()
        src = src[~src["source"].apply(_source_is_placeholder)]
        counts_sources = (
            src["source"]
            .value_counts(dropna=False)
            .head(int(top_k_sources))
            .rename_axis("source")
            .reset_index(name="Count")
        )
        fig_sources = px.bar(
            counts_sources.sort_values("Count"),
            x="Count", y="source", orientation="h",
            title=f"Top Sources (Top {min(top_k_sources, len(counts_sources))})",
            color="source",
            color_discrete_sequence=px.colors.qualitative.Vivid,
        )
        fig_sources.update_layout(showlegend=False, yaxis_title="Source", xaxis_title="Count")
    else:
        fig_sources = px.bar(x=[], y=[], title="No 'source' column found.")

    # ======================================================================
    # Figure B: Top Controlled Vocabulary (overall)
    # ======================================================================
    _cb("Building Top Controlled Vocabulary bar chart")

    if "controlled_vocabulary_terms" in df.columns:
        def _cv_to_list(v):
            if isinstance(v, (list, tuple, set)):
                toks = [str(x).strip() for x in v if str(x).strip()]
            else:
                s = "" if pd.isna(v) else str(v)
                s = re.sub(r"[\r\n]+", ";", s)
                s = re.sub(r"[,\|/]", ";", s)
                toks = [t.strip() for t in s.split(";") if t.strip()]
            out = []
            for t in toks:
                low = t.lower()
                m = re.match(r"^\s*#\s*topic\s*:\s*(.+)$", low, flags=re.I)
                if m:
                    val = m.group(1).strip()
                    if val:
                        out.append(val)
                    continue
                if "#" in t:  # drop hash-taggy junk
                    continue
                out.append(t)
            return out

        cv = df[["controlled_vocabulary_terms"]].copy()
        cv["__terms"] = cv["controlled_vocabulary_terms"].apply(_cv_to_list)
        exploded = cv.explode("__terms")
        exploded = exploded[exploded["__terms"].astype(str).str.strip() != ""]
        counts_terms = (
            exploded["__terms"].value_counts().head(int(top_k_keywords)).rename_axis("term").reset_index(name="Count")
        )
        fig_cv = px.bar(
            counts_terms.sort_values("Count"),
            x="Count", y="term", orientation="h",
            title=f"Top Controlled Vocabulary (Top {min(top_k_keywords, len(counts_terms))})",
            color="term",
            color_discrete_sequence=px.colors.qualitative.Vivid,
        )
        fig_cv.update_layout(showlegend=False, yaxis_title="Term", xaxis_title="Count")
    else:
        fig_cv = px.bar(x=[], y=[], title="No 'controlled_vocabulary_terms' column found.")

    # Convert to PNGs
    png_sources = _plotly_to_png(fig_sources)
    png_cv      = _plotly_to_png(fig_cv)

    # ======================================================================
    # Lay out 1×2 grid on ONE slide (left=Sources, right=Vocabulary)
    # ======================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    margin = Inches(0.5)
    gap = Inches(0.35)
    top = Inches(0.4)

    # Available area (minus margins)
    avail_w = prs.slide_width - 2 * margin
    avail_h = prs.slide_height - 2 * margin

    # Two columns, one row
    cell_w = int((avail_w - gap) / 2)
    cell_h = int(avail_h)

    left1 = margin
    left2 = margin + cell_w + gap
    top_row = top

    _place_in_cell(slide, png_sources, left1, top_row, cell_w, cell_h)
    _place_in_cell(slide, png_cv,      left2, top_row, cell_w, cell_h)

    if slide_notes:
        try:
            slide.notes_slide.notes_text_frame.text = (
                f"Left: Top sources (top {top_k_sources}).  "
                f"Right: Top controlled vocabulary terms (top {top_k_keywords})."
            )
        except Exception:
            pass

# --- END REPLACEMENT ---------------------------------------------------------
def add_metadata_overview_slide(
    prs,
    df: "pd.DataFrame",
    *,
    slide_title: str = "Metadata Overview",   # kept in signature but NOT rendered on slide
    top_n_sources: int = 15,
    top_n_publication_titles: int = 15,
    slide_notes: bool = True,
    image_width_inches: float = 10.0,
    progress_callback=None,
) -> None:
    """
    One slide with a 3-panel overview:
      (Top, spans full width)   Item type distribution over time (stacked bar by year)
      (Bottom-Left)             Top sources (horizontal bar)
      (Bottom-Right)            Top publication titles (horizontal bar)

    Notes:
      • If year or item_type are missing, the top panel falls back to overall item_type counts.
      • The slide uses the Blank layout (no slide title). The figure is centered.
      • Bottom panels have extra horizontal spacing and row height > top panel.
    """
    import io
    import numpy as np
    import pandas as pd
    import plotly.graph_objects as go
    from plotly.subplots import make_subplots
    from pptx.util import Inches

    def _cb(msg: str):
        if progress_callback:
            try:
                progress_callback(msg)
            except Exception:
                pass

    def _plotly_to_png_bytes(fig, width=1500, height=900, scale=2):
        try:
            return fig.to_image(format="png", width=width, height=height, scale=scale)
        except Exception:
            import plotly.io as pio
            return pio.to_image(fig, format="png", width=width, height=height, scale=scale, engine="kaleido")

    def _numeric_year(series: pd.Series) -> pd.Series:
        y = pd.to_numeric(series, errors="coerce")
        return y.dropna().astype(int)

    def _get_year_col(df_: pd.DataFrame) -> str | None:
        if "year_numeric" in df_.columns and df_["year_numeric"].notna().any():
            return "year_numeric"
        if "year" in df_.columns and df_["year"].notna().any():
            return "year"
        return None

    if df is None or df.empty:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        return

    # ---------- Top sources ----------
    if "analyze_top_list" in globals() and callable(globals()["analyze_top_list"]):
        try:
            s_sources = analyze_top_list(df, "source", top_n=top_n_sources, is_list_column=False)
        except Exception:
            s_sources = pd.Series(dtype=int)
    else:
        s_sources = (
            df["source"].dropna().astype(str).str.strip().replace("", np.nan).dropna().value_counts().head(top_n_sources)
            if "source" in df.columns else pd.Series(dtype=int)
        )

    # ---------- Top publication titles ----------
    pub_col = "publicationTitle" if "publicationTitle" in df.columns else (
        "publication_title" if "publication_title" in df.columns else None
    )
    if pub_col:
        s_pubtitles = (
            df[pub_col].dropna().astype(str).str.strip().replace("", np.nan).dropna()
              .value_counts().head(top_n_publication_titles)
        )
    else:
        s_pubtitles = pd.Series(dtype=int)

    # ---------- Item type by year (stacked) ----------
    year_col = _get_year_col(df)
    have_stacked = bool(year_col and "item_type" in df.columns and df["item_type"].notna().any())
    pivot = pd.DataFrame()
    top_types: list[str] = []
    years_sorted: list[int] = []

    if have_stacked:
        d = df[[year_col, "item_type"]].copy()
        d[year_col] = _numeric_year(d[year_col])
        d["item_type"] = d["item_type"].astype(str).str.strip().replace("", np.nan)
        d = d.dropna(subset=[year_col, "item_type"])
        if not d.empty:
            # limit categories to keep the legend tidy
            keep_types = list(d["item_type"].value_counts().head(8).index)
            d["item_type"] = d["item_type"].where(d["item_type"].isin(keep_types), other="Other")
            pivot = pd.crosstab(d[year_col], d["item_type"]).sort_index()
            years_sorted = pivot.index.tolist()
            top_types = pivot.sum(axis=0).sort_values(ascending=False).index.tolist()

    # ---------- Compose a 2×2 subplots figure (top spans full width) ----------
    fig = make_subplots(
        rows=2,
        cols=2,
        specs=[[{"colspan": 2}, None], [{}, {}]],
        # Bottom row slightly taller; create generous gap between bottom panels
        row_heights=[0.42, 0.58],
        vertical_spacing=0.12,
        horizontal_spacing=0.14,
        subplot_titles=(
            "Item Type Distribution over Time" if have_stacked else "Item Type Distribution",
            "Top Sources",
            "Top Publication Titles",
        ),
    )

    # Top panel
    if have_stacked and not pivot.empty:
        for itype in top_types:
            yvals = pivot.get(itype, pd.Series(0, index=pivot.index)).tolist()
            fig.add_trace(
                go.Bar(
                    name=str(itype),
                    x=years_sorted,
                    y=yvals,
                    hovertemplate=f"{itype}<br>Year=%{{x}}<br>Count=%{{y}}<extra></extra>",
                ),
                row=1, col=1,
            )
        fig.update_xaxes(title_text="Year", row=1, col=1)
        fig.update_yaxes(title_text="Count", row=1, col=1)
        # hide phantom (1,2) axes to keep the top centered visually
        fig.update_xaxes(visible=False, row=1, col=2)
        fig.update_yaxes(visible=False, row=1, col=2)
    else:
        # Fallback: overall item_type counts
        if "item_type" in df.columns and df["item_type"].notna().any():
            s_it = (
                df["item_type"].astype(str).str.strip().replace("", np.nan).dropna()
                  .value_counts().head(15)
            )
            fig.add_trace(
                go.Bar(
                    y=s_it.index.tolist()[::-1],
                    x=s_it.values.tolist()[::-1],
                    orientation="h",
                ),
                row=1, col=1,
            )
            fig.update_xaxes(title_text="Count", row=1, col=1)
            fig.update_yaxes(title_text="Item Type", row=1, col=1)
        fig.update_xaxes(visible=False, row=1, col=2)
        fig.update_yaxes(visible=False, row=1, col=2)

    # Bottom-Left: sources
    if not s_sources.empty:
        fig.add_trace(
            go.Bar(
                y=s_sources.index.tolist()[::-1],
                x=s_sources.values.tolist()[::-1],
                orientation="h",
            ),
            row=2, col=1,
        )
        fig.update_xaxes(title_text="Count", row=2, col=1)
        fig.update_yaxes(title_text="Source", row=2, col=1)
    else:
        fig.add_annotation(
            row=2, col=1,
            text="No 'source' data.",
            showarrow=False,
            xref="x2 domain", yref="y2 domain",
        )

    # Bottom-Right: publication titles
    if not s_pubtitles.empty:
        fig.add_trace(
            go.Bar(
                y=s_pubtitles.index.tolist()[::-1],
                x=s_pubtitles.values.tolist()[::-1],
                orientation="h",
            ),
            row=2, col=2,
        )
        fig.update_xaxes(title_text="Count", row=2, col=2)
        fig.update_yaxes(title_text="Publication", row=2, col=2)
    else:
        fig.add_annotation(
            row=2, col=2,
            text="No 'publicationTitle' data.",
            showarrow=False,
            xref="x3 domain", yref="y3 domain",
        )

    # Style: centered look, subtle grid, tidy margins
    fig.update_layout(
        barmode="stack" if have_stacked else "relative",
        template="plotly_white",
        showlegend=True if have_stacked else False,
        legend=dict(orientation="h", yanchor="bottom", y=1.04, xanchor="right", x=1),
        margin=dict(l=60, r=60, t=40, b=60),
        font=dict(family="Calibri", size=12),
    )
    for r in (1, 2):
        for c in (1, 2):
            fig.update_xaxes(gridcolor="rgba(0,0,0,0.08)", row=r, col=c)
            fig.update_yaxes(gridcolor="rgba(0,0,0,0.08)", row=r, col=c)

    # ---------- Add centered image on a Blank slide ----------
    png = _plotly_to_png_bytes(fig, width=1500, height=900, scale=2)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank (NO slide title)
    desired_width = Inches(image_width_inches)
    top = Inches(0.6)  # keep a comfortable top margin
    left = int((prs.slide_width - desired_width) / 2)

    pic = slide.shapes.add_picture(io.BytesIO(png), left=left, top=top, width=desired_width)

    # Ensure the image fits vertically; keep it centered horizontally
    max_h = prs.slide_height - top - Inches(0.4)
    if pic.height > max_h:
        scale = max_h / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)
        pic.left = int((prs.slide_width - pic.width) / 2)

    if slide_notes:
        slide.notes_slide.notes_text_frame.text = (
            "Top: item types over time (stacked; overall fallback if no years). "
            "Bottom: top sources and publication titles. Extra gap added between "
            "bottom panels; bottom row slightly taller than the top."
        )
def add_keyword_kpi_slide(
    prs,
    df: "pd.DataFrame",
    *,
    data_source: str = "controlled_vocabulary_terms",
    image_width_inches: float = 10.0,
    slide_notes: bool = True,
    progress_callback=None,
):
    """
    Slide 1 — Vocabulary KPI cards only (no title on slide).
    """
    import io, math, re, logging, numpy as np, pandas as pd
    from collections import Counter
    from pptx.util import Inches
    from PIL import Image, ImageDraw, ImageFont

    def _cb(msg: str):
        if progress_callback:
            try: progress_callback(f"KPI: {msg}")
            except Exception: pass
        logging.info(f"KPI: {msg}")

    def _cv_to_list(v) -> list[str]:
        if isinstance(v, (list, tuple, set)):
            toks = [str(x).strip() for x in v if str(x).strip()]
        else:
            s = "" if pd.isna(v) else str(v)
            s = re.sub(r"[\r\n]+", ";", s)
            s = re.sub(r"[,\|/]", ";", s)
            toks = [t.strip() for t in s.split(";") if t.strip()]
        out = []
        for t in toks:
            m = re.match(r"^\s*#\s*topic\s*:\s*(.+)$", t, flags=re.I)
            if m:
                val = m.group(1).strip()
                if val:
                    out.append(val)
                continue
            # keep hashtags: strip a leading '#' but don't drop the term
            tt = t.strip()
            if tt.startswith("#"):
                tt = tt.lstrip("#").strip()
                if tt:
                    out.append(tt)
                continue
            out.append(tt)

        return out

    # normalize source and allow fulltext (which is not a column)
    _src = str(data_source).lower().replace(" ", "").replace("-", "_")
    is_fulltext = _src in {"fulltext", "full_text", "full", "pdf"}

    if df is None or df.empty or (not is_fulltext and data_source not in df.columns):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        if slide_notes:
            slide.notes_slide.notes_text_frame.text = (
                "No data to build keyword KPIs." if is_fulltext else f"No '{data_source}' data."
            )
        return slide

    work = df.copy()
    work["__terms"] = work[data_source].apply(_cv_to_list)
    total_docs = len(work)
    docs_with_kw = int(sum(1 for L in work["__terms"] if isinstance(L, list) and len(L) > 0))
    pct_cov = (docs_with_kw / total_docs * 100.0) if total_docs else 0.0
    vocab = Counter([t for L in work["__terms"] for t in (L if isinstance(L, list) else [])])
    vocab_size = len(vocab)
    avg_kw_per_doc = (sum(len(L) for L in work["__terms"]) / total_docs) if total_docs else 0.0
    if vocab:
        top1_share = max(vocab.values()) / total_docs * 100.0 if total_docs else 0.0
        probs = np.array(list(vocab.values()), dtype=float)
        probs = probs / probs.sum() if probs.sum() > 0 else probs
        shannon_H = float(-(probs * np.log2(np.clip(probs, 1e-12, 1.0))).sum())
        hhi = float(np.sum(np.square(probs)))
    else:
        top1_share = 0.0; shannon_H = 0.0; hhi = 0.0

    # --- draw KPI cards on a single canvas ---
    W, H = 1920, 1080
    canvas = Image.new("RGBA", (W, H), (255, 255, 255, 255))
    draw = ImageDraw.Draw(canvas)
    try:
        f_big = ImageFont.truetype("arial.ttf", 36)
        f_mid = ImageFont.truetype("arial.ttf", 22)
    except Exception:
        f_big = ImageFont.load_default()
        f_mid = ImageFont.load_default()

    cards = [
        ("Docs", f"{total_docs}"),
        ("Docs w/ Keywords", f"{docs_with_kw} ({pct_cov:.1f}%)"),
        ("Avg KW / Doc", f"{avg_kw_per_doc:.2f}"),
        ("Vocabulary Size", f"{vocab_size}"),
        ("Top-1 Share", f"{top1_share:.1f}%"),
        ("Entropy (H)", f"{shannon_H:.2f}"),
        ("HHI", f"{hhi:.3f}"),
    ]
    cols = 3
    rows = math.ceil(len(cards) / cols)
    pad = 24
    grid_w = W - pad * (cols + 1)
    grid_h = H - pad * (rows + 1)
    card_w = grid_w // cols
    card_h = grid_h // rows

    k = 0
    for r in range(rows):
        for c in range(cols):
            if k >= len(cards): break
            x = pad + c * (card_w + pad)
            y = pad + r * (card_h + pad)
            draw.rounded_rectangle([x, y, x + card_w, y + card_h], radius=24,
                                   fill=(245, 247, 250, 255), outline=(220, 225, 230, 255), width=2)
            label, value = cards[k]
            draw.text((x + 18, y + 16), label, fill=(98, 104, 109, 255), font=f_mid)
            vw, vh = draw.textbbox((0, 0), value, font=f_big)[2:]
            draw.text((x + 18, y + card_h // 2 - vh // 2 + 6), value, fill=(17, 24, 39, 255), font=f_big)
            k += 1

    # insert into PPT (no title)
    out = io.BytesIO(); canvas.convert("RGB").save(out, format="PNG"); out.seek(0)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    pic = slide.shapes.add_picture(out, left=Inches(0.5), top=Inches(0.5), width=Inches(image_width_inches))
    # vertical fit
    max_h = prs.slide_height - pic.top - Inches(0.5)
    if pic.height > max_h:
        sc = max_h / pic.height
        pic.height = int(pic.height * sc); pic.width = int(pic.width * sc)
    # center horizontally
    pic.left = int((prs.slide_width - pic.width) / 2)
    if slide_notes:
        slide.notes_slide.notes_text_frame.text = "Vocabulary KPIs (coverage, size, concentration, diversity)."
    return slide
def add_keyword_panels_2x2_slide(  # now: FOUR slides, one per plot
    prs,
    df: "pd.DataFrame",
    *,
    slide_title: str | None = None,  # used as a group label in notes
    data_source: str = "full_text",
    # treemap
    top_n_words: int = 30,
    min_frequency: int = 1,
    # heatmaps
    heatmap_top_terms: int = 20,
    heatmap_top_groups: int = 12,
    # wordcloud
    max_words: int = 250,
    wordcloud_colormap: str = "viridis",
    image_width_inches: float = 10.0,
    slide_notes: bool = True,
    progress_callback=None,
):
    """
    Builds FOUR slides (one per visualization), each with LLM-written notes:

      1) Keyword Treemap
      2) Word Cloud (or Top-Terms bar fallback)
      3) Keyword × Author Heatmap
      4) Keyword × ItemType Heatmap

    Parsing is #theme-aware: if any '#theme:<term>' tokens are present, only those are used (prefix stripped).
    """
    import io, re, logging, base64
    import numpy as np, pandas as pd
    from pptx.util import Inches, Emu
    import plotly.express as px
    import plotly.graph_objects as go

    try:
        from PIL import Image
    except Exception:
        Image = None

    # ---------- callbacks ----------
    def _cb(msg: str):
        if progress_callback:
            try:
                progress_callback(f"KW: {msg}")
            except Exception:
                pass
        logging.info(f"KW: {msg}")

    # ---------- guards ----------
    _src_norm = str(data_source).lower().replace(" ", "").replace("-", "_")
    is_fulltext = _src_norm in {"fulltext", "full_text", "full", "pdf"}

    if df is None or getattr(df, "empty", True) or (not is_fulltext and data_source not in df.columns):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        if slide_notes:
            slide.notes_slide.notes_text_frame.text = (
                "No data to build keyword panels."
                if is_fulltext else f"No '{data_source}' data."
            )
        return

    # ---------- helpers ----------
    def _plotly_to_png_bytes(fig: go.Figure, *, width=1500, height=900, scale=2) -> bytes:
        try:
            return fig.to_image(format="png", width=width, height=height, scale=scale)
        except Exception:
            import plotly.io as pio
            return pio.to_image(fig, format="png", width=width, height=height, scale=scale, engine="kaleido")

    def _fig_to_slide(fig: go.Figure, *, width=1500, height=900, scale=2):
        png = _plotly_to_png_bytes(fig, width=width, height=height, scale=scale)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        desired_width = Inches(image_width_inches)
        top = Inches(0.6)
        left = Emu(int((prs.slide_width - desired_width) / 2))
        pic = slide.shapes.add_picture(io.BytesIO(png), left=left, top=top, width=desired_width)
        # vertical fit
        max_h = prs.slide_height - top - Inches(0.4)
        if pic.height > max_h:
            sc = max_h / pic.height
            pic.height = int(pic.height * sc)
            pic.width = int(pic.width * sc)
            pic.left = Emu(int((prs.slide_width - pic.width) / 2))
        return slide

    def _render_plot_slide(fig: go.Figure, *, title_for_prompt: str, analysis_suffix: str,
                           data_df: "pd.DataFrame", prompt_purpose: str):
        slide = _fig_to_slide(fig, width=1500, height=900, scale=2)
        # numeric-only table for notes (guarantee ≥1×1)
        try:
            df_num = data_df.select_dtypes(include=["number", "bool"]).copy()
        except Exception:
            df_num = pd.DataFrame()
        if df_num.shape[1] == 0:
            df_num = pd.DataFrame({"_dummy": [0]})
        elif df_num.empty:
            df_num = pd.DataFrame([{c: 0 for c in df_num.columns}])
        if slide_notes:
            try:
                notes = ai_notes_from_plot(
                    fig,
                    df_num,
                    collection_name=(slide_title or "Keywords"),
                    analysis_suffix=analysis_suffix,
                    title_for_prompt=title_for_prompt,
                    prompt_purpose=prompt_purpose,
                    max_tokens=1200,
                    effort="low",
                    base_logs_dir=None,
                    _cb=_cb,
                )
            except Exception:
                notes = ""
            slide.notes_slide.notes_text_frame.text = (notes or title_for_prompt).strip()
        return slide

    def _extract_theme_tags(cell) -> list[str]:
        """
        Return canonical '#theme:<term>' tokens from a cell (list or string).
        Accepts lists/sets/tuples and ';|,\\n' separators. Ignores non-theme tokens.
        """
        if cell is None or (isinstance(cell, float) and pd.isna(cell)):
            return []
        if isinstance(cell, (list, tuple, set)):
            parts = []
            for v in cell:
                if v is None or (isinstance(v, float) and pd.isna(v)):
                    continue
                parts.extend(re.split(r"[;\|,\n]+", str(v)))
        else:
            parts = re.split(r"[;\|,\n]+", str(cell))
        out: list[str] = []
        for raw in parts:
            s = str(raw or "").strip()
            if not s:
                continue
            if re.match(r"^\s*#?\s*theme\s*:", s, flags=re.I):
                base = re.sub(r"^\s*#?\s*theme\s*:\s*", "", s, flags=re.I).strip()
                base = re.sub(r"\s+", " ", base)
                if base:
                    out.append(f"{base}")
        return out

    def _tokens_from_cell(cell) -> list[str]:
        """
        Theme-aware tokenizer:
          - If '#theme:' tokens exist in the cell → use ONLY those (without prefix).
          - Else: split on ; , | / and newlines; strip leading '#' and whitespace.
        """
        # check for theme tags
        ttags = _extract_theme_tags(cell)
        if ttags:
            return [t.replace("#theme:", "").strip() for t in ttags if t.strip()]
        # generic split
        if cell is None or (isinstance(cell, float) and pd.isna(cell)):
            return []
        if isinstance(cell, (list, tuple, set)):
            parts = []
            for v in cell:
                if v is None or (isinstance(v, float) and pd.isna(v)):
                    continue
                parts.extend(re.split(r"[;\|,\n/]+", str(v)))
        else:
            parts = re.split(r"[;\|,\n/]+", str(cell))
        out = []
        for p in parts:
            s = str(p).strip()
            if not s:
                continue
            if s.startswith("#"):
                s = s.lstrip("#").strip()
            out.append(s)
        return out

    def _authors_to_list(v) -> list[str]:
        if isinstance(v, list):
            return [str(x).strip() for x in v if str(x).strip()]
        s = "" if pd.isna(v) else str(v)
        s = re.sub(r"[\r\n]+", ";", s)
        s = re.sub(r"[,\|/]", ";", s)
        return [t.strip() for t in s.split(";") if t.strip()]

    # ---------- normalize terms once (#theme-aware) ----------
    work = df.copy()
    if is_fulltext:
        key = "controlled_vocabulary_terms" if "controlled_vocabulary_terms" in work.columns else data_source
        work["__terms"] = work.get(key, pd.Series([[]]*len(work))).apply(_tokens_from_cell)
    else:
        work["__terms"] = work[data_source].apply(_tokens_from_cell)

    # ============================================================
    # Slide 1 — Keyword Treemap
    # ============================================================
    _cb("Building slide 1: Keyword Treemap…")
    counts_tm: dict[str, int] = {}
    for L in work["__terms"]:
        for t in (L or []):
            counts_tm[t] = counts_tm.get(t, 0) + 1
    tm_df = (
        pd.Series(counts_tm, name="Count")
        .sort_values(ascending=False)
        .reset_index().rename(columns={"index": "Term"})
    )
    if min_frequency > 1:
        tm_df = tm_df[tm_df["Count"] >= int(min_frequency)]
    tm_df = tm_df.sort_values(["Count", "Term"], ascending=[False, True], kind="mergesort").head(int(top_n_words))

    if tm_df.empty:
        fig_tm = go.Figure().add_annotation(text="No keyword data for treemap.", showarrow=False)
        fig_tm.update_layout(template="plotly_white", title="Keyword Treemap", margin=dict(l=20, r=20, t=60, b=20))
    else:
        fig_tm = px.treemap(tm_df, path=["Term"], values="Count", title="Keyword Treemap")
        fig_tm.update_layout(margin=dict(l=20, r=20, t=60, b=20), template="plotly_white")

    _render_plot_slide(
        fig_tm,
        data_df=tm_df.copy(),
        analysis_suffix="kw_treemap",
        title_for_prompt=f"{slide_title or 'Keywords'} · Keyword Treemap",
        prompt_purpose=(
            "Explain which keywords dominate and whether the distribution is concentrated or long-tailed. "
            "Call out obvious synonyms/variants and the risk of double counting. "
            "Note the min-frequency filter and top-N cutoff; propose normalization & alias merges."
        ),
    )

    # ============================================================
    # Slide 2 — Word Cloud (or Top Terms bar)
    # ============================================================
    _cb("Building slide 2: Word Cloud / Top Terms…")
    flat_terms = [t for L in work["__terms"] for t in (L or [])]
    counts_wc = (
        pd.Series(flat_terms, dtype="object").value_counts().reset_index().rename(columns={"index": "Term", 0: "Count"})
    )
    counts_wc.columns = ["Term", "Count"]
    counts_wc["Count"] = pd.to_numeric(counts_wc["Count"], errors="coerce").fillna(0).astype(int)
    counts_wc = counts_wc.sort_values(["Count", "Term"], ascending=[False, True], kind="mergesort")
    top_counts_wc = counts_wc.head(40).reset_index(drop=True)

    fig_wc = None
    wc_image_source = None
    try:
        _zot = globals().get("zotero_client") or globals().get("zt")
        _coll = (
                globals().get("collection_name_for_cache")
                or globals().get("CURRENT_COLLECTION_NAME")
                or globals().get("collection_name")
        )
        _cb("WordCloud helper: generating…")
        _, wc_data_uri = analyze_word_cloud(
            df,
            {
                "data_source":"full_text",
                "max_words": max_words,
                "wordcloud_colormap": wordcloud_colormap,
                "wordcloud_render_style": "cloud",
            },
            zotero_client_for_pdf=_zot,
            collection_name_for_cache=_coll,
            progress_callback=_cb,
        )
        if isinstance(wc_data_uri, str) and wc_data_uri.startswith("data:image/"):
            wc_image_source = wc_data_uri
    except Exception as e:
        _cb(f"WordCloud helper failed: {e}")

    if wc_image_source:
        fig_wc = go.Figure()
        fig_wc.add_layout_image(
            dict(
                source=wc_image_source,  # data URI works without PIL
                xref="paper", yref="paper",
                x=0, y=1, sizex=1, sizey=1,
                xanchor="left", yanchor="top",
                layer="below"
            )
        )
        fig_wc.update_xaxes(visible=False)
        fig_wc.update_yaxes(visible=False)
        fig_wc.update_layout(title="Word Cloud", template="plotly_white", margin=dict(l=0, r=0, t=60, b=0))
    else:
        if top_counts_wc.empty:
            fig_wc = go.Figure().add_annotation(text="No terms for word cloud.", showarrow=False)
            fig_wc.update_layout(template="plotly_white", title="Top Keywords (fallback)")
        else:
            fig_wc = px.bar(
                top_counts_wc.sort_values("Count"),
                x="Count", y="Term", orientation="h", title="Top Keywords (fallback)",
            )
            fig_wc.update_layout(template="plotly_white", margin=dict(l=120, r=40, t=60, b=40))

    # numeric-only payload for notes (Count + Rank)
    if top_counts_wc.empty:
        wc_notes_df = pd.DataFrame({"Count": [0], "Rank": [1]})
    else:
        wc_notes_df = pd.DataFrame({
            "Count": pd.to_numeric(top_counts_wc["Count"], errors="coerce").fillna(0).astype(int).values,
            "Rank": np.arange(1, len(top_counts_wc) + 1, dtype=int),
        })

    _render_plot_slide(
        fig_wc,
        data_df=wc_notes_df,
        analysis_suffix="kw_wordcloud",
        title_for_prompt=f"{slide_title or 'Keywords'} · Word Cloud / Top Terms",
        prompt_purpose=(
            "Summarize the dominant terms and how quickly frequencies drop. "
            "Emphasize that word-cloud size reflects frequency only (not importance). "
            "List caveats (stopwords, lemmatization, hashtag normalization) and propose clean-ups & deduping."
        ),
    )

    # ============================================================
    # Slide 3 — Keyword × Author Heatmap
    # ============================================================
    _cb("Building slide 3: Keyword × Author Heatmap…")
    if "authors_list" in df.columns:
        work["__authors"] = df["authors_list"].apply(_authors_to_list)
    else:
        work["__authors"] = df.get("authors", pd.Series([""]*len(work))).apply(_authors_to_list)

    la = work.loc[:, ["__terms", "__authors"]].explode("__authors").explode("__terms")
    la = la.dropna(subset=["__authors", "__terms"])
    la["__authors"] = la["__authors"].astype(str).str.strip()
    la["__terms"] = la["__terms"].astype(str).str.strip()
    la = la[(la["__authors"] != "") & (la["__terms"] != "")]

    pv_auth = pd.DataFrame()
    if not la.empty:
        top_terms = la["__terms"].value_counts().head(int(heatmap_top_terms)).index
        la_t = la[la["__terms"].isin(top_terms)]
        top_auth = la_t["__authors"].value_counts().head(int(heatmap_top_groups)).index
        sub = la_t[la_t["__authors"].isin(top_auth)].copy()
        sub["cnt"] = 1
        pv_auth = sub.pivot_table(index="__terms", columns="__authors", values="cnt", aggfunc="sum", fill_value=0)

    if pv_auth.size == 0:
        fig_ha = go.Figure().add_annotation(text="Not enough author/term data for heatmap.", showarrow=False)
        fig_ha.update_layout(template="plotly_white", title="Keyword × Author")
        data_ha = pd.DataFrame({"_dummy": [0]})
    else:
        fig_ha = px.imshow(pv_auth, aspect="auto", color_continuous_scale="Blues", title="Keyword × Author")
        fig_ha.update_layout(template="plotly_white", margin=dict(l=80, r=20, t=60, b=60))
        data_ha = pv_auth.reset_index()

    _render_plot_slide(
        fig_ha,
        data_df=data_ha.select_dtypes(include=["number", "bool"]).copy() if not data_ha.empty else data_ha,
        analysis_suffix="kw_heatmap_authors",
        title_for_prompt=f"{slide_title or 'Keywords'} · Keyword × Author",
        prompt_purpose=(
            "Identify which authors cluster around which keywords. "
            "Call out a few author-keyword hubs and any surprising gaps; recommend validation via sample inspection."
        ),
    )

    # ============================================================
    # Slide 4 — Keyword × ItemType Heatmap
    # ============================================================
    _cb("Building slide 4: Keyword × ItemType Heatmap…")
    it_col = next((c for c in ("item_type", "doc_type_readable", "document_type", "type") if c in df.columns), None)
    pv_type = pd.DataFrame()
    if it_col:
        li = pd.DataFrame({it_col: df[it_col].astype(str).str.strip(), "__terms": work["__terms"]})
        li = li.replace({it_col: {"": np.nan}}).dropna(subset=[it_col])
        li = li.explode("__terms").dropna(subset=["__terms"])
        li["__terms"] = li["__terms"].astype(str).str.strip()
        li = li[li["__terms"] != ""]
        if not li.empty:
            top_terms = li["__terms"].value_counts().head(int(heatmap_top_terms)).index
            li_t = li[li["__terms"].isin(top_terms)]
            top_types = li_t[it_col].value_counts().head(int(heatmap_top_groups)).index
            sub = li_t[li_t[it_col].isin(top_types)].copy()
            sub["cnt"] = 1
            pv_type = sub.pivot_table(index="__terms", columns=it_col, values="cnt", aggfunc="sum", fill_value=0)

    if pv_type.size == 0:
        fig_hi = go.Figure().add_annotation(text="Not enough item-type/term data for heatmap.", showarrow=False)
        fig_hi.update_layout(template="plotly_white", title=f"Keyword × {it_col or 'ItemType'}")
        data_hi = pd.DataFrame({"_dummy": [0]})
    else:
        fig_hi = px.imshow(pv_type, aspect="auto", color_continuous_scale="YlOrRd", title=f"Keyword × {it_col}")
        fig_hi.update_layout(template="plotly_white", margin=dict(l=80, r=20, t=60, b=60))
        data_hi = pv_type.reset_index()

    _render_plot_slide(
        fig_hi,
        data_df=data_hi.select_dtypes(include=["number", "bool"]).copy() if not data_hi.empty else data_hi,
        analysis_suffix="kw_heatmap_itemtype",
        title_for_prompt=f"{slide_title or 'Keywords'} · Keyword × {it_col or 'ItemType'}",
        prompt_purpose=(
            "Show which keywords are concentrated in particular item/document types. "
            "Point out skew and potential coding artefacts; suggest harmonizing item-type labels if fragmented."
        ),
    )
# --- safe numeric table for notes (≥1 row, ≥2 numeric cols) ---
def _safe_notes_2cols(df: pd.DataFrame, cands: tuple[str,...] = ("Count","count","value","values")) -> pd.DataFrame:
    if not isinstance(df, pd.DataFrame) or df.empty:
        return pd.DataFrame({"c1":[0], "c2":[0]})
    num = df.select_dtypes(include=["number"]).copy()
    if num.empty:
        for name in cands:
            if name in df.columns:
                num = pd.DataFrame({"c1": pd.to_numeric(df[name], errors="coerce")})
                break
        if num.empty:
            num = pd.DataFrame({"c1":[0]})
    num = num.apply(lambda s: pd.to_numeric(s, errors="coerce")).fillna(0)
    if num.shape[1] == 1:
        num["c2"] = 0
    return num.reset_index(drop=True)

def add_Metadata_trends(
    prs,
    df: "pd.DataFrame",
    *,
    slide_title: str = "Metadata Trends",   # used for notes grouping
    top_n_sources: int = 15,
    top_n_publication_titles: int = 15,
    top_n_vocab: int = 20,
    top_n_themes: int = 12,
    slide_notes: bool = True,
    image_width_inches: float = 10.0,
    progress_callback=None,
) -> None:
    """
    FOUR slides (one per plot), each with LLM-written notes:
      1) Top Sources
      2) Top Publication Titles
      3) Top Controlled Vocabulary (cleans '#…' except '#topic: …')
      4) Topic Themes (deterministic dummy themes when missing)
    """
    import io, re, hashlib
    import numpy as np
    import pandas as pd
    import plotly.graph_objects as go
    from pptx.util import Inches

    def _cb(msg: str):
        if progress_callback:
            try:
                progress_callback(msg)
            except Exception:
                pass

    if df is None or getattr(df, "empty", True):
        prs.slides.add_slide(prs.slide_layouts[6])  # Blank placeholder
        return

    # ---------- helpers ----------
    def _plotly_to_png_bytes(fig, width=1500, height=900, scale=2):
        try:
            return fig.to_image(format="png", width=width, height=height, scale=scale)
        except Exception:
            import plotly.io as pio
            return pio.to_image(fig, format="png", width=width, height=height, scale=scale, engine="kaleido")

    def _series_to_df(series: pd.Series, label_col: str = "Item") -> pd.DataFrame:
        """
        Convert a value_counts-style Series to a 2-col DataFrame.
        Fix: operate on the Series (not its .values ndarray) so .fillna() works.
        """
        if not isinstance(series, pd.Series) or series.empty:
            return pd.DataFrame({label_col: ["—"], "Count": [0]})

        counts = pd.to_numeric(series, errors="coerce").fillna(0).astype(int)
        out = pd.DataFrame({
            label_col: [str(x) for x in counts.index.tolist()],
            "Count": counts.to_list(),
        })
        # deterministic sort for cache stability
        out = out.sort_values(["Count", label_col], ascending=[False, True], kind="mergesort").reset_index(drop=True)
        return out

    def _hbar_fig(df_xy: pd.DataFrame, label_col: str, title_sub: str, empty_msg: str) -> go.Figure:
        fig = go.Figure()
        if df_xy["Count"].sum() == 0 and (df_xy[label_col] == "—").all():
            fig.add_annotation(text=empty_msg, showarrow=False)
        else:
            fig.add_trace(
                go.Bar(
                    y=df_xy[label_col].tolist()[::-1],
                    x=df_xy["Count"].tolist()[::-1],
                    orientation="h",
                    marker=dict(line=dict(width=0.5, color="rgba(0,0,0,0.6)")),
                    hovertemplate="%{y}<br>Count: %{x:,}<extra></extra>",
                )
            )
        fig.update_layout(
            template="plotly_white",
            title=dict(text=title_sub, x=0.02, xanchor="left"),
            margin=dict(l=80, r=60, t=60, b=60),
            font=dict(family="Calibri", size=12),
            showlegend=False,
        )
        fig.update_xaxes(title_text="Count", gridcolor="rgba(0,0,0,0.08)")
        fig.update_yaxes(title_text=None, gridcolor="rgba(0,0,0,0.08)")
        return fig

    def _render_slide_with_notes(fig: go.Figure, data_df: pd.DataFrame, analysis_suffix: str, title_for_prompt: str, prompt_purpose: str):
        # render onto a blank slide, centered
        png = _plotly_to_png_bytes(fig, width=1500, height=900, scale=2)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        desired_width = Inches(image_width_inches)
        top = Inches(0.6)
        left = int((prs.slide_width - desired_width) / 2)
        pic = slide.shapes.add_picture(io.BytesIO(png), left=left, top=top, width=desired_width)
        # fit vertically while keeping it centered
        max_h = prs.slide_height - top - Inches(0.4)
        if pic.height > max_h:
            scale = max_h / pic.height
            pic.height = int(pic.height * scale)
            pic.width  = int(pic.width  * scale)
            pic.left   = int((prs.slide_width - pic.width) / 2)
        # notes via LLM
        if slide_notes:
            notes_text = ai_notes_from_plot(
                fig,
                data_df,
                collection_name=slide_title,
                analysis_suffix=analysis_suffix,
                title_for_prompt=title_for_prompt,
                prompt_purpose=prompt_purpose,
                max_tokens=1400,
                effort="low",
                base_logs_dir=None,
                _cb=_cb,
            )
            if isinstance(notes_text, str) and notes_text.strip():
                slide.notes_slide.notes_text_frame.text = notes_text
            else:
                slide.notes_slide.notes_text_frame.text = f"{slide_title} • {title_for_prompt}"

    # ---------- TOP SOURCES ----------
    _cb("Metadata: building slide — Top Sources…")
    if "source" in df.columns:
        s_sources = (
            df["source"].astype(str).str.strip().replace({"": np.nan}).dropna()
              .value_counts().head(int(top_n_sources))
        )
    else:
        s_sources = pd.Series(dtype=int)
    df_sources = _series_to_df(s_sources, label_col="Source")
    fig_sources = _hbar_fig(df_sources, "Source", "Top Sources", "No 'source' data.")
    _render_slide_with_notes(
        fig_sources,
        df_sources,
        analysis_suffix="metadata_top_sources",
        title_for_prompt=f"{slide_title} · Top Sources",
        prompt_purpose=(
            "Describe which sources dominate and whether the distribution is long-tailed or concentrated. "
            "Call out coverage bias (e.g., one database/publisher vs. diversified intake) and any naming variants. "
            "Note data limits (deduplication, source normalization) and suggest next steps (harmonize names, add missing feeds, "
            "stratify by year)."
        ),
    )

    # ---------- TOP PUBLICATION TITLES ----------
    _cb("Metadata: building slide — Top Publication Titles…")
    pub_col = "publicationTitle" if "publicationTitle" in df.columns else (
        "publication_title" if "publication_title" in df.columns else None
    )
    if pub_col:
        s_pub = (
            df[pub_col].astype(str).str.strip().replace({"": np.nan}).dropna()
              .value_counts().head(int(top_n_publication_titles))
        )
    else:
        s_pub = pd.Series(dtype=int)
    df_pub = _series_to_df(s_pub, label_col="Publication")
    fig_pub = _hbar_fig(df_pub, "Publication", "Top Publication Titles", "No publication title data.")
    _render_slide_with_notes(
        fig_pub,
        df_pub,
        analysis_suffix="metadata_top_publications",
        title_for_prompt=f"{slide_title} · Top Publication Titles",
        prompt_purpose=(
            "Summarize where the work appears (journals/venues) and what that implies about audience and field. "
            "Highlight any single-venue dominance vs. spread across venues. Flag potential synonym issues "
            "(abbreviations, series names) and propose venue normalization & field-of-study mapping."
        ),
    )

    # ---------- TOP CONTROLLED VOCABULARY ----------
    _cb("Metadata: building slide — Top Controlled Vocabulary…")
    def _cv_to_list(v) -> list[str]:
        if isinstance(v, (list, tuple, set)):
            toks = [str(x).strip() for x in v if str(x).strip()]
        else:
            s = "" if pd.isna(v) else str(v)
            s = re.sub(r"[\r\n]+", ";", s)
            s = re.sub(r"[,\|/]", ";", s)
            toks = [t.strip() for t in s.split(";") if t.strip()]
        out: list[str] = []
        for t in toks:
            m = re.match(r"^\s*#\s*topic\s*:\s*(.+)$", t, flags=re.I)
            if m:
                val = m.group(1).strip()
                if val:
                    out.append(val)
                continue
            tt = t.strip()
            if tt.startswith("#"):
                tt = tt.lstrip("#").strip()
                if tt:
                    out.append(tt)
                continue
            out.append(tt)
        return out

    if "controlled_vocabulary_terms" in df.columns:
        cv_series = df["controlled_vocabulary_terms"].apply(_cv_to_list)
        flat = [kw for lst in cv_series for kw in (lst if isinstance(lst, list) else [])]
        s_vocab = pd.Series(flat).value_counts().head(int(top_n_vocab)) if flat else pd.Series(dtype=int)
    else:
        s_vocab = pd.Series(dtype=int)
    df_vocab = _series_to_df(s_vocab, label_col="Term")
    fig_vocab = _hbar_fig(df_vocab, "Term", "Top Controlled Vocabulary", "No controlled vocabulary.")
    _render_slide_with_notes(
        fig_vocab,
        df_vocab,
        analysis_suffix="metadata_top_vocab",
        title_for_prompt=f"{slide_title} · Top Controlled Vocabulary",
        prompt_purpose=(
            "Explain the most frequent controlled terms, grouping near-synonyms where obvious. "
            "Comment on topical focus vs. breadth, and any taxonomy artifacts (plural/singular, capitalization). "
            "Recommend cleanup (merge aliases, map to ontology IDs) and checks for drift over time."
        ),
    )

    # ---------- THEMES (with deterministic dummy when missing) ----------
    _cb("Metadata: building slide — Topic Themes…")
    used_dummy = False
    if "theme" in df.columns and df["theme"].astype(str).str.strip().replace({"": np.nan}).notna().any():
        s_theme = (
            df["theme"].astype(str).str.strip().replace({"": np.nan}).dropna()
              .value_counts().head(int(top_n_themes))
        )
    else:
        used_dummy = True
        themes = ["Theme A", "Theme B", "Theme C", "Theme D", "Theme E", "Theme F"]
        titles = df["title"] if "title" in df.columns else pd.Series([""] * len(df))
        def _pick_theme(s):
            try:
                h = int(hashlib.md5(str(s).encode("utf-8", "ignore")).hexdigest(), 16)
            except Exception:
                h = 0
            return themes[h % len(themes)]
        tmp = titles.apply(_pick_theme)
        s_theme = tmp.value_counts().head(int(top_n_themes))

    df_theme = _series_to_df(s_theme, label_col="Theme")
    fig_theme = _hbar_fig(df_theme, "Theme", "Topic Themes", "No themes (dummy used).")
    _render_slide_with_notes(
        fig_theme,
        df_theme,
        analysis_suffix="metadata_top_themes",
        title_for_prompt=f"{slide_title} · Topic Themes",
        prompt_purpose=(
            ("These themes are inferred deterministically from titles (dummy fallback); treat as illustrative, not definitive. "
             if used_dummy else
             "Summarize dominant themes and how sharply concentrated they are. ")
            + "Discuss overlap/ambiguity between labels and the implications for downstream analysis. "
              "Propose validation (human review, LLM clustering with guidelines, and mapping to a stable taxonomy)."
        ),
    )

def add_basic_ngram_analysis_slides(
    prs,
    df: "pd.DataFrame",
    keywords: list[str],
    *,
    top_n: int = 15,                        # label only; trends always use provided `keywords`
    n: int = 1,                             # fixed to 1 for this “basic” set
    collection_name: str = "Collection",
    collection_name_for_cache: str | None = None,
    zotero_client=None,
    slide_notes: bool = True,
    image_width_inches: float = 10.0,
    include_abstract_title: bool = True,    # NEW: allow skipping abstract/title slides
    transition_seconds: float | None = None, # NEW: simple auto-advance (seconds). None => no auto-advance
    progress_callback=None,
        rebuild=False
):
    """
    Create *one-big-plot-per-slide* “Basic N-gram (N=1) Analysis” using the provided keywords (Plotly-only).

    Slides (large, single Plotly chart each; exported to PNG via Kaleido):
      1) N1 — Full text • Trend (lines)
      2) N1 — Full text • Stacked bar (year × keyword)
      3) N1 — Abstract • Trend                [optional via include_abstract_title]
      4) N1 — Abstract • Stacked bar          [optional via include_abstract_title]
      5) N1 — Title • Trend                   [optional via include_abstract_title]
      6) N1 — Title • Stacked bar             [optional via include_abstract_title]
      7+) N1 — Author × Keyword (Full text) heatmap, chunked to 3 keywords/slide

    Notes
    -----
    • PowerPoint cannot embed interactive HTML; charts are Plotly figures rendered to PNG (needs kaleido).
    • Phrase matching is contiguous over token streams (multi-word keywords respected).
    • Year is taken from 'year' / 'year_numeric' coerced to int.
    """
    import io
    import logging
    from collections import defaultdict, Counter

    import numpy as np
    import pandas as pd
    from pptx.util import Inches, Emu

    # Plotly (required)
    try:
        import plotly.express as px
        import plotly.graph_objects as go
        import plotly.io as pio
    except Exception as e:
        raise RuntimeError(
            "Plotly is required for n-gram slides. Install with: pip install plotly kaleido"
        ) from e

    # add near the top of add_basic_ngram_analysis_slides (after imports)
    # --- helpers: source normalization ------------------------------------------
    def _src(s: str) -> str:
        s = str(s or "").strip().lower()
        if s in {"full_text", "fulltext", "pdf", "body"}:
            return "fulltext"
        if s in {"title+abstract", "title_abstract", "abstract_title"}:
            return "title_abstract"
        return s

    # --- ONE-TIME preloader per source ------------------------------------------
    _preloaded: dict[str, dict[str, list[str]]] = {}  # source -> {doc_id -> tokens}

    def _ensure_preloaded(source: str):
        src = _src(source)
        if src in _preloaded:
            return

        mapping: dict[str, list[str]] = {}
        try:
            if src == "fulltext":
                # Use the iterator mode so we get per-doc tokens from the *collection cache*
                per_doc_iter = get_text_corpus_from_df(
                    df,
                    source_type="fulltext",
                    collection_name_for_cache=collection_name_for_cache or collection_name,
                    zotero_client=zotero_client,
                    progress_callback=None,
                    cache_mode=("refresh" if rebuild else "read"),
                    preserve_item_cache=True,
                    return_mode="iterator",  # yields (doc_id, tokens)
                    include_summary=False,
                )
                for doc_id, toks in per_doc_iter:
                    mapping[str(doc_id)] = [str(t).lower() for t in (toks or [])]
                _cb(f"Preloaded {len(mapping)} docs (fulltext) from cache_mode={'refresh' if rebuild else 'read'}")
            elif src == "title":
                for i, row0 in df.reset_index(drop=True).iterrows():
                    doc_id = str(row0.get("key") or f"row_{i}")
                    mapping[doc_id] = _tok(str(row0.get("title", "") or ""))
            elif src == "abstract":
                for i, row0 in df.reset_index(drop=True).iterrows():
                    doc_id = str(row0.get("key") or f"row_{i}")
                    mapping[doc_id] = _tok(str(row0.get("abstract", "") or ""))
            elif src == "title_abstract":
                for i, row0 in df.reset_index(drop=True).iterrows():
                    doc_id = str(row0.get("key") or f"row_{i}")
                    mapping[doc_id] = _tok(f"{row0.get('title', '') or ''} {row0.get('abstract', '') or ''}")
            elif src == "controlled_vocabulary_terms":
                for i, row0 in df.reset_index(drop=True).iterrows():
                    doc_id = str(row0.get("key") or f"row_{i}")
                    kws = row0.get("controlled_vocabulary_terms", []) or []
                    mapping[doc_id] = [str(k).lower() for k in kws if isinstance(k, str)]
            else:
                # Fallback = title + abstract
                for i, row0 in df.reset_index(drop=True).iterrows():
                    doc_id = str(row0.get("key") or f"row_{i}")
                    mapping[doc_id] = _tok(f"{row0.get('title', '') or ''} {row0.get('abstract', '') or ''}")
        except Exception as e:
            _cb(f"Preload failed for source={src}: {e}; falling back to per-row tokenization.")
            mapping = {}

        _preloaded[src] = mapping

    def _row_tokens(row: pd.Series, source: str) -> list[str]:
        """
        Return cached tokens for this row+source.
        Ensures we never call Zotero during 'read' runs and we reuse the collection cache.
        """
        src = _src(source)
        _ensure_preloaded(src)
        doc_id = str(row.get("key") or f"row_{getattr(row, 'name', '0')}")
        return _preloaded.get(src, {}).get(doc_id, [])

    # --- callbacks ------------------------------------------------------------
    def _cb(msg: str):
        if progress_callback:
            try:
                progress_callback(f"N1 N-grams: {msg}")
            except Exception:
                pass
        logging.info(f"N1 N-grams: {msg}")

    # --- helpers: export / placement / transitions ---------------------------
    def _plotly_to_png_bytes(fig) -> bytes:
        """Render Plotly figure to PNG using Kaleido (raises on failure)."""
        return pio.to_image(fig, format="png", scale=2)

    def _add_image_slide_centered(png_bytes: bytes, notes: str | None = None):
        """Add a BLANK slide and center the image (no slide title)."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        gutter = Inches(0.5)
        top = Inches(0.8)

        # Fit to width within gutters, then downscale if too tall
        max_w = Emu(int(prs.slide_width - 2 * gutter))
        left = Emu(int((prs.slide_width - max_w) / 2))
        pic = slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, width=max_w)

        max_h = prs.slide_height - top - gutter
        if pic.height > max_h:
            scale = max_h / pic.height
            pic.height = int(pic.height * scale)
            pic.width = int(pic.width * scale)
            pic.left = int((prs.slide_width - pic.width) / 2)

        if slide_notes and notes:
            slide.notes_slide.notes_text_frame.text = notes

        _apply_transition(slide)
        return slide

    def _apply_transition(slide):
        """Apply a subtle transition and optional auto-advance; guard for API presence."""
        try:
            tr = slide.slide_show_transition
            if hasattr(tr, "advance_on_click"):
                tr.advance_on_click = True
            if transition_seconds is not None:
                if hasattr(tr, "advance_on_time"):
                    tr.advance_on_time = True
                if hasattr(tr, "advance_after_time"):
                    tr.advance_after_time = int(max(0.0, float(transition_seconds)) * 1000)
        except Exception:
            pass

    def _add_plotly_slide(fig: "go.Figure", notes: str | None = None):
        png_bytes = _plotly_to_png_bytes(fig)
        _add_image_slide_centered(png_bytes, notes=notes)

    # --- data_processing hooks / tokenization --------------------------------

    from data_processing import get_text_corpus_from_df, tokenize_for_ngram_context_refined


    def _tok(s: str) -> list[str]:
        return [t.lower() for t in tokenize_for_ngram_context_refined(s or "")]

    # normalize provided keywords → tokens (phrase-aware matching)
    kw_phrases = []
    for kw in (keywords or []):
        if not isinstance(kw, str) or not kw.strip():
            continue
        toks = _tok(kw)
        if toks:
            kw_phrases.append((kw.strip(), toks))
    if not kw_phrases:
        _cb("No valid keywords provided; aborting N1 slides.")
        return

    # choose year column
    year_col = "year_numeric" if "year_numeric" in df.columns else ("year" if "year" in df.columns else None)
    if year_col is None:
        _cb("No 'year' column found; temporal plots may be empty.")

    # contiguous phrase counter
    def _count_phrase_in_tokens(tokens: list[str], phrase: list[str]) -> int:
        L = len(phrase)
        if not tokens or L == 0 or L > len(tokens):
            return 0
        cnt = 0
        for i in range(len(tokens) - L + 1):
            if tokens[i:i + L] == phrase:
                cnt += 1
        return cnt

    # token extractor per row+source (prefers cached full text via data_processing)
    # replace _row_tokens with this version
    def _row_tokens(row: pd.Series, source: str) -> list[str]:
        """
        One-shot, collection-level token lookup.
        - Builds per-source token maps once (for the whole `df`) using get_text_corpus_from_df(return_mode="per_doc")
        - Then serves tokens by Zotero key / id / row index without re-hitting Zotero/PDF.
        - Falls back to title/abstract when no full text is available for a given row.
        """
        src = _src(source).lower()

        def _fallback_tokens() -> list[str]:
            if src == "title":
                return _tok(str(row.get("title", "") or ""))
            if src == "abstract":
                return _tok(str(row.get("abstract", "") or ""))
            if src == "controlled_vocabulary_terms":
                kws = row.get("controlled_vocabulary_terms", []) or []
                return [str(k).lower() for k in kws if isinstance(k, str)]
            # default fallback (when we don’t actually have full text)
            txt = f"{row.get('title', '') or ''} {row.get('abstract', '') or ''}"
            return _tok(txt)

        # static, in-function cache: {source -> {doc_id -> [tokens...]}}
        cache = getattr(_row_tokens, "_cache", None)
        if cache is None:
            _row_tokens._cache = {}
            cache = _row_tokens._cache

        # normalize source aliases used by the corpus builder
        if src in {"fulltext", "full_text"}:
            src_norm = "full_text"
        elif src in {"abstract"}:
            src_norm = "abstract"
        elif src in {"title"}:
            src_norm = "title"
        else:
            src_norm = src  # leave others as-is (e.g., controlled_vocabulary_terms)

        # Rebuild-on-demand per source (respect `rebuild`)
        if rebuild:
            cache.pop(src_norm, None)

        # Build once per source using the WHOLE dataframe
        if src_norm not in cache and src_norm in {"full_text", "abstract", "title"}:
            try:
                per_doc = get_text_corpus_from_df(
                    df,
                    source_type=src_norm,
                    collection_name_for_cache=collection_name_for_cache or collection_name,
                    zotero_client=zotero_client,
                    progress_callback=None,
                    cache_mode= "read",
                    preserve_item_cache=True,
                    return_mode="per_doc",
                    include_summary=False,
                ) or []
                tok_map: dict[str, list[str]] = {}
                for item in per_doc:
                    # accept either (doc_id, tokens) tuples or dicts
                    if isinstance(item, tuple) and len(item) == 2:
                        doc_id, toks = item
                    elif isinstance(item, dict):
                        doc_id = item.get("doc_id") or item.get("key") or item.get("id")
                        toks = item.get("tokens") or item.get("flat") or item.get("text") or []
                    else:
                        continue
                    if doc_id is None:
                        continue
                    tok_map[str(doc_id)] = [str(t).lower() for t in (toks or [])]
                cache[src_norm] = tok_map
            except Exception:
                # if the collection-level build fails, fall back per row
                cache[src_norm] = {}

        # look up tokens for this row
        doc_id = str(row.get("key") or row.get("id") or row.get("zotero_key") or f"row_{row.name}")
        toks = cache.get(src_norm, {}).get(doc_id)

        # If nothing cached for this doc, return a sensible fallback for the requested source
        return toks if toks else _fallback_tokens()

    def _has_series_data(years, series_by_kw) -> bool:
        if not years or not series_by_kw:
            return False
        for vals in series_by_kw.values():
            if any(int(v) > 0 for v in vals):
                return True
        return False

    def _add_plotly_slide(fig, notes: str | None = None):
        if fig is None:
            _cb("Skip: figure has no data.")
            return  # do not create a slide
        png_bytes = _plotly_to_png_bytes(fig)
        _add_image_slide_centered(png_bytes, notes=notes)

    # per-source aggregations
    def _build_yearly_and_total_counts(source: str, *, strict: bool = True):
        """
        Compute per-year and total keyword counts for the given `source`.

        Returns
        -------
        years_sorted : list[int]
        series_by_kw : dict[str, list[int]]
        totals       : dict[str, int]

        Raises
        ------
        ValueError/RuntimeError when `strict=True` and we detect empty corpus or no data.
        With strict=False, we just log and return empties so callers can render placeholders.
        """
        import time
        import logging
        from collections import defaultdict, Counter
        import numpy as np
        import pandas as pd

        t0 = time.time()
        src = str(source)

        def _fail(msg, *, hard: bool):
            _cb(msg)
            if hard:
                raise RuntimeError(msg)

        if df is None or len(df) == 0:
            msg = f"_build_yearly_and_total_counts: DataFrame is empty for source='{src}'."
            _cb(msg)
            if strict:
                raise ValueError(msg)
            return [], {}, {}

        if not kw_phrases:
            msg = "_build_yearly_and_total_counts: No valid keywords/phrases to count."
            _cb(msg)
            if strict:
                raise ValueError(msg)
            return [], {}, {}

        if year_col is None:
            _cb(f"_build_yearly_and_total_counts[{src}]: No 'year'/'year_numeric' column; yearly plots may be empty.")

        per_year_kw = defaultdict(lambda: defaultdict(int))
        totals = Counter()

        yrs = (
            pd.to_numeric(df[year_col], errors="coerce")
            if year_col else pd.Series([np.nan] * len(df), index=df.index)
        )
        if year_col:
            _cb(f"_build_yearly_and_total_counts[{src}]: Using year column '{year_col}', "
                f"valid years={int(yrs.notna().sum())}/{len(yrs)}.")

        rows_with_tokens = 0
        total_token_len = 0

        for idx, row in df.iterrows():
            try:
                tokens = _row_tokens(row, src)
            except Exception as e:
                logging.exception("N1: _row_tokens failed for row=%s source=%s", idx, src)
                _cb(f"_row_tokens failed for row={idx} source='{src}': {e}")
                continue

            if not tokens:
                continue

            if not isinstance(tokens, list):
                try:
                    tokens = list(tokens)
                except Exception:
                    tokens = [str(tokens)]
            tokens = [str(t).lower() for t in tokens if t is not None]
            if not tokens:
                continue

            rows_with_tokens += 1
            total_token_len += len(tokens)

            for kw, phrase in kw_phrases:
                c = _count_phrase_in_tokens(tokens, phrase)
                if c:
                    totals[kw] += c

            y = yrs.loc[idx]
            if not pd.isna(y):
                y = int(y)
                for kw, phrase in kw_phrases:
                    c = _count_phrase_in_tokens(tokens, phrase)
                    if c:
                        per_year_kw[y][kw] += c

        elapsed_ms = int((time.time() - t0) * 1000)
        nonzero_kw_totals = sum(1 for v in totals.values() if v > 0)
        _cb(
            f"_build_yearly_and_total_counts[{src}]: "
            f"rows_with_tokens={rows_with_tokens}/{len(df)}, "
            f"avg_tokens_per_row={(total_token_len / rows_with_tokens) if rows_with_tokens else 0:.1f}, "
            f"nonzero_kw_totals={nonzero_kw_totals}, elapsed={elapsed_ms} ms."
        )

        years_sorted = sorted(per_year_kw.keys())
        series_by_kw = (
            {kw: [per_year_kw[y].get(kw, 0) for y in years_sorted] for kw, _ in kw_phrases}
            if years_sorted else {}
        )

        # Hard failure only when strict
        if rows_with_tokens == 0:
            _fail(
                f"_build_yearly_and_total_counts[{src}]: No tokens produced for ANY row. "
                f"Check corpus retrieval / cache.",
                hard=strict,
            )
            return [], {}, dict(totals)

        if nonzero_kw_totals == 0 and not years_sorted:
            _fail(
                f"_build_yearly_and_total_counts[{src}]: All counts are zero and no yearly data. "
                f"Possible causes: keywords never occur, tokenization mismatch, or empty source texts. "
                f"Example first keyword tokens: {kw_phrases[0][1] if kw_phrases else 'n/a'}",
                hard=strict,
            )
            return [], {}, dict(totals)

        try:
            top_preview = ", ".join(f"{k}={totals.get(k, 0)}" for k, _ in kw_phrases[:5])
            _cb(f"_build_yearly_and_total_counts[{src}]: totals preview → {top_preview}")
        except Exception:
            pass

        return years_sorted, series_by_kw, dict(totals)

    # --- Plot factories (Plotly only) ----------------------------------------
    BASE_FONT = dict(family="Calibri", size=12, color="#404446")

    # code to be replaced
    def _fig_trend_plotly(years, series_by_kw, title: str, *, legend_right: bool = False):
        # Remove keywords whose total across the whole timeframe is 0
        series_nonzero = {
            kw: (vals or [])
            for kw, vals in (series_by_kw or {}).items()
            if any(int(v) > 0 for v in (vals or []))
        }
        if not _has_series_data(years, series_nonzero):
            _cb(f"Skip slide '{title}': no temporal data after filtering all-zero keywords.")
            return None

        long_rows = []
        for kw, vals in series_nonzero.items():
            for y, v in zip(years, vals):
                long_rows.append({"Year": int(y), "Keyword": kw, "Count": int(v)})

        import plotly.express as px
        df_long = pd.DataFrame(long_rows)
        if df_long.empty:
            _cb(f"Skip slide '{title}': empty dataframe after melt.")
            return None

        fig = px.line(
            df_long, x="Year", y="Count", color="Keyword",
            markers=True, template="plotly_white", title=title, line_shape="spline"
        )
        fig.update_traces(mode="lines+markers", line=dict(width=2))
        fig.update_layout(
            hovermode="x unified",
            legend=(dict(orientation="v", y=0.5, yanchor="middle", x=1.02, xanchor="left", title=None)
                    if legend_right else dict(orientation="h", yanchor="bottom", y=1.02, x=0, title=None)),
            margin=dict(l=70, r=(150 if legend_right else 30), t=70, b=60),
            xaxis=dict(showgrid=True, gridcolor="rgba(0,0,0,.08)", tickmode="linear"),
            yaxis=dict(showgrid=True, gridcolor="rgba(0,0,0,.08)"),
            font=BASE_FONT,
            title=dict(font=dict(family="Calibri Light", size=22, color="#181A1B")),
        )
        return fig

    # code for replacement
    def _fig_trend_plotly(years, series_by_kw, title: str, *, legend_right: bool = True):
        series_nonzero = {
            kw: (vals or [])
            for kw, vals in (series_by_kw or {}).items()
            if any(int(v) > 0 for v in (vals or []))
        }
        if not _has_series_data(years, series_nonzero):
            _cb(f"Skip slide '{title}': no temporal data after filtering all-zero keywords.")
            return None

        long_rows = [{"Year": int(y), "Keyword": kw, "Count": int(v)}
                     for kw, vals in series_nonzero.items()
                     for y, v in zip(years, vals)]

        import plotly.express as px
        df_long = pd.DataFrame(long_rows)
        if df_long.empty:
            _cb(f"Skip slide '{title}': empty dataframe after melt.")
            return None

        fig = px.line(
            df_long, x="Year", y="Count", color="Keyword",
            markers=True, template="plotly_white", title=title, line_shape="spline"
        )
        fig.update_traces(mode="lines+markers", line=dict(width=2))
        fig.update_layout(
            hovermode="x unified",
            legend=dict(orientation="v", y=0.5, yanchor="middle", x=1.02, xanchor="left", title=None),
            margin=dict(l=70, r=160, t=70, b=60),
            xaxis=dict(showgrid=True, gridcolor="rgba(0,0,0,.08)", tickmode="linear"),
            yaxis=dict(showgrid=True, gridcolor="rgba(0,0,0,.08)"),
            font=BASE_FONT,
            title=dict(font=dict(family="Calibri Light", size=22, color="#181A1B")),
        )
        return fig

    def _fig_stackbar_plotly(years, series_by_kw, title: str):
        series_nonzero = {
            kw: (vals or [])
            for kw, vals in (series_by_kw or {}).items()
            if any(int(v) > 0 for v in (vals or []))
        }
        if not _has_series_data(years, series_nonzero):
            _cb(f"Skip slide '{title}': no per-year counts after filtering all-zero keywords.")
            return None

        long_rows = [{"Year": int(y), "Keyword": kw, "Count": int(v)}
                     for kw, vals in series_nonzero.items()
                     for y, v in zip(years, vals)]

        import plotly.express as px
        df_long = pd.DataFrame(long_rows)
        if df_long.empty:
            _cb(f"Skip slide '{title}': empty dataframe after melt.")
            return None

        fig = px.bar(
            df_long, x="Year", y="Count", color="Keyword",
            barmode="stack", template="plotly_white", title=title
        )
        fig.update_layout(
            legend=dict(orientation="v", y=0.5, yanchor="middle", x=1.02, xanchor="left", title=None),
            margin=dict(l=70, r=160, t=70, b=60),
            xaxis=dict(type="category", showgrid=False),
            yaxis=dict(showgrid=True, gridcolor="rgba(0,0,0,.08)"),
            font=BASE_FONT,
            title=dict(font=dict(family="Calibri Light", size=22, color="#181A1B")),
        )
        return fig

    def _fig_author_keyword_heatmap_plotly(top_authors, kw_subset, matrix, title: str):
        df_heat = pd.DataFrame(matrix, index=top_authors, columns=kw_subset)
        fig = px.imshow(
            df_heat,
            text_auto=True,
            aspect="auto",
            template="plotly_white",
            title=title,
            color_continuous_scale="Blues",
            labels={"x": "Keyword", "y": "Author", "color": "Frequency"},
        )
        fig.update_xaxes(side="bottom", tickangle=-30)
        fig.update_layout(
            margin=dict(l=120, r=40, t=70, b=80),
            font=BASE_FONT,
            title=dict(font=dict(family="Calibri Light", size=22, color="#181A1B")),
        )
        return fig



    # --- 1) FULL TEXT: trend + stacked bar -----------------------------------
    # FULL TEXT — keep strict to catch real issues
    _cb("Computing counts for FULL TEXT…")
    years_ft, series_ft, totals_ft = _build_yearly_and_total_counts("full_text", strict=True)
    if include_abstract_title:
        try:
            nonempty_abs = int(df.get("abstract", pd.Series(dtype=object)).astype(str).str.strip().ne("").sum())
            _cb(f"Abstract sanity: non-empty abstracts={nonempty_abs}/{len(df)}")
        except Exception:
            _cb("Abstract sanity: could not compute non-empty abstract count.")

    # ABSTRACT — lenient; render placeholders if empty
    if include_abstract_title:
        _cb("Computing counts for ABSTRACT…")
        years_abs, series_abs, totals_abs = _build_yearly_and_total_counts("abstract", strict=False)

        title3 = "N1 — Abstract • Trend"
        fig3 = _fig_trend_plotly(years_abs, series_abs, title3, legend_right=True)
        _add_plotly_slide(fig=fig3, notes=f"{collection_name} · Source=abstract")

        title4 = "N1 — Abstract • Stacked Bar by Year"
        fig4 = _fig_stackbar_plotly(years_abs, series_abs, title4)
        _add_plotly_slide(fig=fig4, notes=f"{collection_name} · Source=abstract")

        # TITLE — lenient as well
        _cb("Computing counts for TITLE…")
        years_ti, series_ti, totals_ti = _build_yearly_and_total_counts("title", strict=False)

        title5 = "N1 — Title • Trend"
        fig5 = _fig_trend_plotly(years_ti, series_ti, title5)
        _add_plotly_slide(fig=fig5, notes=f"{collection_name} · Source=title")

        title6 = "N1 — Title • Stacked Bar by Year"
        fig6 = _fig_stackbar_plotly(years_ti, series_ti, title6)
        _add_plotly_slide(fig=fig6, notes=f"{collection_name} · Source=title")

    # --- 3) AUTHORS × KEYWORD heatmaps (full text), 3 keywords/slide ----------
    _cb("Building Author × Keyword heatmap(s) on FULL TEXT…")
    # Build author × keyword counts
    from collections import defaultdict as _ddict
    auth_kw = _ddict(lambda: Counter())
    for _, row in df.iterrows():
        authors_str = str(row.get("authors", "") or "")
        if not authors_str.strip():
            continue
        authors = [a.strip() for a in authors_str.split(";") if a.strip()]
        tokens = _row_tokens(row, "fulltext")
        if not tokens:
            continue
        for kw, phrase in kw_phrases:
            c = _count_phrase_in_tokens(tokens, phrase)
            if c:
                for a in authors:
                    auth_kw[a][kw] += c

    # Choose top authors by total across *all provided keywords*
    totals_by_author = Counter({a: sum(cnts.get(kw, 0) for kw, _ in kw_phrases) for a, cnts in auth_kw.items()})
    top_authors = [a for a, _ in totals_by_author.most_common(25) if totals_by_author[a] > 0]

    # chunk keywords to 3 per slide
    def _chunks(lst, size):
        for i in range(0, len(lst), size):
            yield lst[i:i + size]

    # Precompute non-zero keywords across the chosen top authors,
    # then pack slides with up to 10 actual (non-zero) keywords.
    kw_list = [kw for kw, _ in kw_phrases]
    kw_totals = {kw: sum(auth_kw.get(a, {}).get(kw, 0) for a in top_authors) for kw in kw_list}
    kw_nonzero = [kw for kw in kw_list if kw_totals.get(kw, 0) > 0]

    # If no authors or no non-zero keywords, add a single placeholder and stop.
    if not top_authors or not kw_nonzero:
        _cb("Skip Author×Keyword heatmap: no authors or all keyword totals are zero on full text.")
    else:
        for chunk_idx, kw_chunk in enumerate(_chunks(kw_nonzero, 10), start=1):
            mat = np.zeros((len(top_authors), len(kw_chunk)), dtype=int)
            for r, auth in enumerate(top_authors):
                row_counts = auth_kw.get(auth, {})
                for c, kw in enumerate(kw_chunk):
                    mat[r, c] = int(row_counts.get(kw, 0))

            if mat.sum() == 0:
                _cb(f"Skip Author×Keyword heatmap chunk {chunk_idx}: matrix all zeros.")
                continue

            fig_h = _fig_author_keyword_heatmap_plotly(
                top_authors, kw_chunk, mat, title="N1 — Author × Keyword (Full text)"
            )
            _add_plotly_slide(fig_h, notes=f"Top {len(top_authors)} authors; Keywords: {', '.join(kw_chunk)}")


def add_ngram_higher_order_slides(
    prs,
    df: "pd.DataFrame",
    keywords: list[str],
    *,
    n_values: tuple[int, ...] = (2, 3, 4, 5),       # n > 1
    data_sources: tuple[str, ...] = ("full_text",),  # e.g., ("fulltext",) or ("title","abstract","fulltext")
    top_n: int = 20,                                 # top-N n-grams per subplot
    collection_name: str = "Collection",
    collection_name_for_cache: str | None = None,
    zotero_client=None,
    slide_notes: bool = True,
    image_width_inches: float = 10.0,
    transition_seconds: float | None = None,         # auto-advance seconds; None => click-only
    progress_callback=None,
        rebuild=False
):
    """
    For each keyword, create ONE slide with a 2×2 Plotly subplot panel showing
    the Top-N n-grams that contain that keyword for n = 2,3,4,5 (or provided n_values).
    Figures include their own titles; the slide itself has *no* title.

    - Plotly-only (exported to PNG via Kaleido).
    - Uses cached full text via data_processing.get_text_corpus_from_df when available.
    - Center images horizontally and scale to fit the slide width minus gutters.

    Notes:
      • If a keyword has fewer than 4 configured n-values, empty panels will state “No data”.
      • If multiple data_sources are provided, a *separate slide per source* is produced for each keyword.
    """
    import io
    import logging
    from collections import Counter

    import pandas as pd
    from pptx.util import Inches, Emu

    # Plotly (required)
    try:
        import plotly.express as px
        import plotly.graph_objects as go
        import plotly.io as pio
        from plotly.subplots import make_subplots
    except Exception as e:
        raise RuntimeError(
            "Plotly is required. Install with: pip install plotly kaleido"
        ) from e

    # ---------- callbacks ----------
    def _cb(msg: str):
        if progress_callback:
            try:
                progress_callback(f"n-gram (n>1): {msg}")
            except Exception:
                pass
        logging.info(f"n-gram (n>1): {msg}")

    # ---------- helpers: pptx placement & transitions ----------
    def _apply_transition(slide):
        try:
            tr = slide.slide_show_transition
            if hasattr(tr, "advance_on_click"):
                tr.advance_on_click = True
            if transition_seconds is not None:
                if hasattr(tr, "advance_on_time"):
                    tr.advance_on_time = True
                if hasattr(tr, "advance_after_time"):
                    tr.advance_after_time = int(max(0.0, float(transition_seconds)) * 1000)  # ms
        except Exception:
            pass

    def _add_image_slide_centered(png_bytes: bytes, notes: str | None = None):
        """Add a BLANK (if available) or Title-Only slide, center image, no title text."""
        # Prefer a blank layout (index 6 in most templates). Fallback to title-only (index 5).
        try:
            blank_idx = 6 if len(prs.slide_layouts) > 6 else None
        except Exception:
            blank_idx = None
        if blank_idx is not None:
            slide = prs.slides.add_slide(prs.slide_layouts[blank_idx])
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only; leave title empty

        # center horizontally; keep soft gutters
        gutter = Inches(0.5)
        desired_width = Inches(image_width_inches)
        max_width = Emu(int(prs.slide_width) - int(2 * gutter))
        width = min(desired_width, max_width)
        left = Emu(int((prs.slide_width - width) / 2))
        top = Inches(0.8)
        slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, width=width)

        if slide_notes and notes:
            slide.notes_slide.notes_text_frame.text = notes

        _apply_transition(slide)
        return slide

    def _plotly_to_png_bytes(fig) -> bytes:
        return pio.to_image(fig, format="png", scale=2)

    # ---------- tokenization / corpus ----------
    from data_processing import get_text_corpus_from_df, tokenize_for_ngram_context_refined

    def _tok(s: str) -> list[str]:
        return [t.lower() for t in tokenize_for_ngram_context_refined(s or "")]

    # normalize keywords → token phrases
    kw_phrases: list[tuple[str, list[str]]] = []
    for kw in (keywords or []):
        if not isinstance(kw, str) or not kw.strip():
            continue
        toks = _tok(kw)
        if toks:
            kw_phrases.append((kw.strip(), toks))
    if not kw_phrases:
        _cb("No valid keywords provided; nothing to plot.")
        return

    # ---------- precompute tokens per source (once) ----------
    # ---------- precompute tokens per source (once, via collection cache) ----------
    def _canon_source(s: str) -> str:
        s = (s or "").lower().strip()
        if s in {"full_text", "full text", "fulltext", "pdf", "body"}:
            return "fulltext"
        if s in {"title+abstract", "title_abstract", "abstract_title"}:
            return "title_abstract"
        return s

    data_sources = tuple(_canon_source(s) for s in data_sources)

    def _doc_id_from_row(row: pd.Series) -> str:
        return str(row.get("key") or row.get("id") or row.get("zotero_key") or f"row_{getattr(row, 'name', 0)}")

    def _fallback_tokens_for_source(row: pd.Series, src: str) -> list[str]:
        if src == "title":
            return _tok(str(row.get("title", "") or ""))
        if src == "abstract":
            return _tok(str(row.get("abstract", "") or ""))
        if src == "title_abstract":
            return _tok(f"{row.get('title', '') or ''} {row.get('abstract', '') or ''}")
        if src in {"controlled_vocabulary_terms", "cv"}:
            entry = row.get("controlled_vocabulary_terms", []) or []
            if isinstance(entry, str):
                return [k.strip().lower() for k in entry.split(";") if k.strip()]
            if isinstance(entry, list):
                return [str(k).strip().lower() for k in entry if isinstance(k, str) and k.strip()]
            return []
        # default fallback
        return _tok(f"{row.get('title', '') or ''} {row.get('abstract', '') or ''}")

    tokens_by_source: dict[str, list[list[str]]] = {}

    _cb(f"Extracting token streams for sources (cached): {', '.join(data_sources)} …")
    for src in data_sources:
        # Build a token map for the whole DF using the collection cache when possible
        tok_map: dict[str, list[str]] = {}

        if src in {"fulltext", "title", "abstract", "title_abstract"}:
            try:
                per_doc = get_text_corpus_from_df(
                    df,
                    source_type=src,
                    collection_name_for_cache=collection_name_for_cache or collection_name,
                    zotero_client=zotero_client,
                    progress_callback=None,
                    cache_mode=("refresh" if rebuild else "read"),
                    preserve_item_cache=True,
                    return_mode="per_doc",
                    include_summary=False,
                ) or []
                for item in per_doc:
                    if isinstance(item, tuple) and len(item) == 2:
                        did, toks = item
                    elif isinstance(item, dict):
                        did = item.get("doc_id") or item.get("key") or item.get("id")
                        toks = item.get("tokens") or item.get("flat") or []
                    else:
                        continue
                    if did is None:
                        continue
                    tok_map[str(did)] = [str(t).lower() for t in (toks or [])]
            except Exception as e:
                _cb(f"Cache preload failed for source={src}: {e}; will use per-row fallbacks.")
                tok_map = {}

        # Align to DF order; use sensible per-row fallback when missing
        aligned: list[list[str]] = []
        for _, row in df.iterrows():
            did = _doc_id_from_row(row)
            toks = tok_map.get(did)
            if not toks:
                toks = _fallback_tokens_for_source(row, src)
            aligned.append(toks)

        tokens_by_source[src] = aligned

    # ---------- n-gram mining (contains-phrase) ----------
    def _contains_phrase(ngram_tuple: tuple[str, ...], phrase_tokens: list[str]) -> bool:
        L = len(phrase_tokens)
        N = len(ngram_tuple)
        if L == 0 or L > N:
            return False
        if L == 1:
            return phrase_tokens[0] in ngram_tuple
        for i in range(0, N - L + 1):
            if list(ngram_tuple[i:i + L]) == phrase_tokens:
                return True
        hyp = "-".join(phrase_tokens)
        if hyp in ngram_tuple:
            return True
        return False

    def _top_ngrams_for_keyword(
        tokens_per_doc: list[list[str]],
        phrase_tokens: list[str],
        n: int,
        top_n: int,
        min_freq: int = 1,
    ) -> pd.DataFrame:
        """Return DataFrame with columns=['Ngram','Frequency'] for top n-grams containing the phrase."""
        if n < max(2, len(phrase_tokens)):
            return pd.DataFrame(columns=["Ngram", "Frequency"])
        from nltk import ngrams
        cnt = Counter()
        for toks in tokens_per_doc:
            if not toks or len(toks) < n:
                continue
            for tup in ngrams(toks, n):
                if _contains_phrase(tup, phrase_tokens):
                    cnt[" ".join(tup)] += 1
        if not cnt:
            return pd.DataFrame(columns=["Ngram", "Frequency"])
        items = [(k, v) for k, v in cnt.items() if v >= min_freq]
        df_top = (
            pd.DataFrame(items, columns=["Ngram", "Frequency"])
            .sort_values("Frequency", ascending=False)
            .head(top_n)
            .reset_index(drop=True)
        )
        return df_top

    # ---------- per-keyword slides ----------
    BASE_FONT = dict(family="Calibri", size=12, color="#404446")

    def _subplot_panel_for_keyword(
        keyword_label: str,
        phrase_tokens: list[str],
        source_name: str,
        n_values: tuple[int, ...],
    ) -> "go.Figure":
        """Build a 2×2 subplot panel with horizontal bars for n in n_values (up to 4)."""
        # Prepare up to 4 panels; if fewer than 4 n-values, remaining are placeholders
        n_vals = list(n_values)[:4]
        while len(n_vals) < 4:
            n_vals.append(None)

        # --- build a 2×2 panel with a wide center gutter ---------------------
        desired_n = [n for n in n_vals if isinstance(n, int) and n > 1]
        panels = []
        for n_val in desired_n:
            df_top = _top_ngrams_for_keyword(
                tokens_by_source[source_name], phrase_tokens, n_val, top_n
            )
            if not df_top.empty:
                panels.append((n_val, df_top))

        CENTER_GAP = 0.20  # fraction of figure width reserved as middle spacer
        # Prepare titles for available panels; blank out the rest
        titles = [f"{keyword_label} — {source_name} — {n}-grams" for n, _ in panels]
        while len(titles) < 4:
            titles.append("")  # no title where there is no data to show

        fig = make_subplots(
            rows=2, cols=2,
            horizontal_spacing=CENTER_GAP,
            vertical_spacing=0.12,
            column_widths=[(1.0 - CENTER_GAP) / 2.0, (1.0 - CENTER_GAP) / 2.0],
            subplot_titles=titles[:4],
        )

        # Place only panels that have data
        rowcol = [(1, 1), (1, 2), (2, 1), (2, 2)]
        for (n_val, df_top), (r, c) in zip(panels[:4], rowcol):
            df_top = df_top.iloc[::-1]  # largest on top for h-bars
            fig.add_trace(
                go.Bar(
                    x=df_top["Frequency"],
                    y=df_top["Ngram"],
                    orientation="h",
                    text=df_top["Frequency"],
                    textposition="outside",
                    textfont=dict(size=10),
                    hovertemplate="<b>%{y}</b><br>Count: %{x}<extra></extra>",
                    marker=dict(line=dict(width=0.4, color="rgba(0,0,0,.28)")),
                    name=f"{n_val}-grams",
                    showlegend=False,
                ),
                row=r, col=c
            )

        # If no n-value had data at all, show a single unobtrusive message
        if len(panels) == 0:
            fig.add_annotation(
                x=0.5, y=0.5, xref="paper", yref="paper",
                text="No n-gram data for this keyword/source.",
                showarrow=False, font=dict(size=14, color="rgba(0,0,0,.6)")
            )

        # global layout & spacing polish
        fig.update_layout(
            template="plotly_white",
            height=760, width=1200,
            margin=dict(l=70, r=32, t=64, b=52),
            font=dict(family=BASE_FONT["family"], size=11, color=BASE_FONT["color"]),
            title=dict(
                text=f"<b>{keyword_label}</b> — {source_name.title()} · Top {top_n} n-grams containing the keyword",
                x=0.5, xanchor="center",
                font=dict(family="Calibri Light", size=20, color="#181A1B"),
            ),
            uniformtext_minsize=8,
            uniformtext_mode="hide",
            bargap=0.15,
            bargroupgap=0.04,
        )

        # smaller subplot-title fonts
        fig.update_annotations(selector=dict(type="annotation"), font=dict(size=13))

        # axes: compact ticks + automargin for long n-grams
        for ax in ("xaxis", "xaxis2", "xaxis3", "xaxis4", "yaxis", "yaxis2", "yaxis3", "yaxis4"):
            fig.layout[ax].tickfont = dict(size=10)
            if ax.startswith("yaxis"):
                fig.layout[ax].automargin = True
            else:
                fig.layout[ax].showgrid = True
                fig.layout[ax].gridcolor = "rgba(0,0,0,.08)"

        # draw a subtle vertical "frame" around the center gutter
        x0 = 0.5 - CENTER_GAP / 2.0
        x1 = 0.5 + CENTER_GAP / 2.0
        fig.add_shape(  # transparent spacer (below traces)
            type="rect", xref="paper", yref="paper",
            x0=x0, x1=x1, y0=0.0, y1=1.0,
            fillcolor="rgba(0,0,0,0)", line=dict(width=0), layer="below"
        )
        fig.add_shape(  # left edge
            type="line", xref="paper", yref="paper",
            x0=x0, x1=x0, y0=0.06, y1=0.96,
            line=dict(color="rgba(0,0,0,.07)", width=1), layer="below"
        )
        fig.add_shape(  # right edge
            type="line", xref="paper", yref="paper",
            x0=x1, x1=x1, y0=0.06, y1=0.96,
            line=dict(color="rgba(0,0,0,.07)", width=1), layer="below"
        )

        return fig

    def _has_any_ngram_data(tokens_per_doc: list[list[str]], phrase_tokens: list[str], n_vals: tuple[int, ...]) -> bool:
        """Quick existence check to avoid generating empty panels/slides."""
        from nltk import ngrams
        for n in n_vals:
            if not isinstance(n, int) or n <= 1:
                continue
            for toks in tokens_per_doc:
                if not toks or len(toks) < n:
                    continue
                for tup in ngrams(toks, n):
                    if _contains_phrase(tup, phrase_tokens):
                        return True
        return False

    # Build slides
    for src in data_sources:
        pretty_src = "Full text" if src == "fulltext" else src.title()
        n_list = tuple(n for n in n_values if isinstance(n, int) and n > 1)

        for kw_label, phrase_tokens in kw_phrases:
            # Skip if no n-gram exists for this (keyword, source)
            if not _has_any_ngram_data(tokens_by_source[src], phrase_tokens, n_list):
                _cb(f"Skipping '{kw_label}' · source={pretty_src} (no n-gram data).")
                continue

            _cb(f"Building panel for keyword '{kw_label}' · source={pretty_src}…")
            fig_panel = _subplot_panel_for_keyword(kw_label, phrase_tokens, src, n_list)
            png = _plotly_to_png_bytes(fig_panel)
            notes = f"{collection_name} · Source={pretty_src} · Keyword='{kw_label}' · n={list(n_values)} · TopN={top_n}"
            _add_image_slide_centered(png, notes=notes)

    _cb("n-gram (n>1) slides completed.")

from typing import Any, Iterable
from pathlib import Path
from pptx import Presentation
def add_keyword_thematic(
    prs,
    df,
    *,
    collection_name: str = "Collection",
    theme_col_candidates: tuple[str, ...] = ("controlled_vocabulary", "controlled_vocabulary_terms", "keywords"),
    year_col_candidates: tuple[str, ...] = ("year_numeric", "year", "Year", "pub_year", "publication_year", "issued", "date", "created"),
    target_groups: int = 5,            # kept for backwards-compat; unused in this version
    min_tag_freq: int = 1,             # kept for backwards-compat; filtering is handled upstream
    top_k_terms: int = 25,             # how many keywords per overarching to list in notes
    slide_notes: bool = True,
    progress_callback=None,
    image_width_inches: float = 10.0,
    precomputed_schema: dict | None = None,  # ignored in this version (kept for API stability)
):
    """
    Replacement implementation that uses `thematic_analysis(...)` instead of transformer clustering.

    Creates THREE slides:
      1) Overview — bar chart of Overarching theme totals (from thematic_analysis) + notes listing top keywords per theme.
      2) Theme tree — mind-map sunburst (Overarching → Subtheme → Keyword), values = exact counts from corpus.
      3) Evolution — stacked area chart of Overarching prevalence over time (per-document keyword hits).

    Returns a dict with the thematic results and convenience mappings:
      {
        "theme_col": <column used>,
        "year_col": <column used or None>,
        "thematic": <result from thematic_analysis>,
        "keyword_to_overarching": { "<normalized keyword>": "<Overarching title>", ... },
      }
    """
    # -------------------------- imports & helpers --------------------------
    import io
    import re
    import math
    import numpy as np
    import pandas as pd
    import plotly.graph_objects as go
    from pptx.util import Inches, Emu

    def _cb(msg: str):
        if progress_callback:
            try:
                progress_callback(f"THEMES: {msg}")
            except Exception:
                pass

    def _first_present(cols: tuple[str, ...], frame: pd.DataFrame) -> str | None:
        for c in cols:
            if c in frame.columns:
                return c
        return None

    def _extract_year_series(frame: pd.DataFrame) -> tuple[str | None, pd.Series]:
        """
        Best-effort year extraction: prefer numeric (1800..2100); else parse datetimes.
        Returns (col_used, float Series with NaNs).
        """
        col = _first_present(year_col_candidates, frame)
        if col is None:
            return None, pd.Series([np.nan] * len(frame), index=frame.index)
        s = frame[col]
        # numeric 4-digit years if present
        num = pd.to_numeric(s, errors="coerce")
        out = pd.Series(np.nan, index=frame.index, dtype="float")
        out = np.where((num >= 1800) & (num <= 2100), num, np.nan)
        # fill with parsed datetime years
        dt = pd.to_datetime(s, errors="coerce", utc=True)
        out = pd.Series(out, index=frame.index)
        out = out.fillna(dt.dt.year.astype("float"))
        return col, out

    def _fig_to_slide(fig: go.Figure, *, width_px=1600, height_px=900):
        """Render Plotly fig to a centered PNG on a blank slide."""
        png = fig.to_image(format="png", width=width_px, height=height_px, scale=2)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        left = Emu(int((prs.slide_width - Inches(image_width_inches)) / 2))
        top = Inches(0.6)
        slide.shapes.add_picture(io.BytesIO(png), left, top, width=Inches(image_width_inches))
        return slide

    def _norm_token(s: str) -> str:
        return re.sub(r"\s+", " ", str(s or "").strip()).lower()

    def _tokenize_cell(cell, separators: str = r"[;|,\n]+") -> list[str]:
        if cell is None or (isinstance(cell, float) and math.isnan(cell)):
            return []
        if isinstance(cell, (list, tuple, set)):
            parts = []
            for v in cell:
                if v is None or (isinstance(v, float) and math.isnan(v)):
                    continue
                parts.extend(re.split(separators, str(v)))
        else:
            parts = re.split(separators, str(cell))
        out: list[str] = []
        for raw in parts:
            tok = _norm_token(raw)
            if len(tok) >= 2:
                # Accept either plain keywords or "#theme:<kw>"-style tags
                tok = re.sub(r"^\s*#?\s*theme\s*:\s*", "", tok).strip()
                if tok:
                    out.append(tok)
        return out

    # -------------------------- sanity: data & columns --------------------------
    if df is None or getattr(df, "empty", True):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Themes — Overview"
        slide.placeholders[1].text = "No data provided."
        return {
            "theme_col": None,
            "year_col": None,
            "thematic": None,
            "keyword_to_overarching": {},
        }

    theme_col = _first_present(theme_col_candidates, df)
    if theme_col is None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Themes — Overview"
        slide.placeholders[1].text = "No theme column found."
        return {
            "theme_col": None,
            "year_col": None,
            "thematic": None,
            "keyword_to_overarching": {},
        }

    # -------------------------- run thematic analysis --------------------------
    _cb("Running thematic_analysis…")
    # We pass theme_col so it counts from the right column.
    thematic_result = thematic_analysis(df=df, theme_col=theme_col, verbose=False,section_title=collection_name)

    # Expecting:
    # thematic_result["data"] -> {"flat_themes": {...}, "structural_themes": {...}, "regrouped": {...}?}
    # thematic_result["plots"] -> {"bar_fig": go.Figure, "mind_map_fig": go.Figure}
    data_blk = (thematic_result or {}).get("data", {})
    plots_blk = (thematic_result or {}).get("plots", {})

    flat_themes = data_blk.get("flat_themes", {})          # {Overarching -> [ {term,count,subtheme,synonyms}, ...]}
    structural = data_blk.get("structural_themes", {"overarching": {}, "subthemes": {}, "keywords": {}})
    regrouped = data_blk.get("regrouped") or thematic_result.get("regrouped") or {}

    bar_fig = plots_blk.get("bar_fig")
    mind_map_fig = plots_blk.get("mind_map_fig")

    # -------------------------- Slide 1: overview (bar) --------------------------
    _cb("Slide 1: overview bar…")
    if bar_fig is None:
        # Build a simple bar if not supplied
        og_names = list(structural.get("overarching", {}).keys())
        og_vals = [structural["overarching"][k] for k in og_names]
        show_vals = [v if v > 0 else 1 for v in og_vals]
        bar_fig = go.Figure(go.Bar(x=og_names, y=show_vals, customdata=og_vals,
                                   hovertemplate="<b>%{x}</b><br>Total keyword hits: %{customdata}<extra></extra>"))
        bar_fig.update_layout(
            title={'text': f"{collection_name} · Overarching Themes (keyword hits)", 'x': 0.5, 'xanchor': 'center'},
            xaxis_title="Overarching theme",
            yaxis_title="Total occurrences in corpus",
            margin=dict(t=60, l=40, r=20, b=80),
            template="plotly_white"
        )
    slide1 = _fig_to_slide(bar_fig)

    if slide_notes:
        # Compose notes with top keywords per overarching
        note_lines = []
        for og_title, rows in flat_themes.items():
            if not rows:
                continue
            top_rows = rows[:max(1, min(top_k_terms, len(rows)))]
            note_lines.append(
                f"• {og_title}: " +
                ", ".join(f"{r['term']} ({int(r['count'])})" for r in top_rows if int(r.get("count", 0)) > 0)
                or "—"
            )
        slide1.notes_slide.notes_text_frame.text = (
            "\n".join(note_lines) if note_lines else "Top keywords per overarching theme not available."
        )

    # -------------------------- Slide 2: mind-map sunburst --------------------------
    _cb("Slide 2: mind-map sunburst…")
    if mind_map_fig is None:
        # Build a FULL hierarchy with robust normalisation so EVERY node appears.
        # Key upgrades:
        #  1) Normalise parent lookups (case/spacing) so subthemes/keywords always attach correctly.
        #  2) Give each node a unique id (avoid label collisions).
        #  3) Ensure parent values equal the sum of children (branchvalues='total').
        #  4) Force-include tiny nodes (value >= 1) and aggressively show text.
        import re as _re

        def _slug(s: str) -> str:
            s = "" if s is None else str(s)
            s = s.strip().lower()
            s = _re.sub(r"[^a-z0-9]+", "-", s)
            return s.strip("-") or "x"

        def _norm_key(s: str) -> str:
            s = "" if s is None else str(s)
            return _re.sub(r"[^a-z0-9]+", "", s.strip().lower())

        # Pull structures (tolerate any being missing)
        og_dict_raw = structural.get("overarching", {}) or {}
        sub_dict_raw = structural.get("subthemes", {}) or {}  # {overarching: {sub: count}}
        kw_dict_raw = structural.get("keywords", {}) or {}  # {sub: {kw: count}}

        # Normalise mapping for robust matching
        # Overarching list keeps original labels but tracks a normalised key
        og_items = [(str(og), _norm_key(og), int(og_dict_raw.get(og, 0) or 0)) for og in og_dict_raw.keys()]
        # Subthemes: map by normalised parent key
        sub_by_og_norm = {}
        for og_label, subs in sub_dict_raw.items():
            ogk = _norm_key(og_label)
            bucket = sub_by_og_norm.setdefault(ogk, {})
            for sub_label, cnt in (subs or {}).items():
                bucket[str(sub_label)] = int(cnt or 0)
        # Keywords: map by normalised subtheme key
        kw_by_sub_norm = {}
        for sub_label, kws in kw_dict_raw.items():
            sk = _norm_key(sub_label)
            bucket = kw_by_sub_norm.setdefault(sk, {})
            for kw_label, cnt in (kws or {}).items():
                bucket[str(kw_label)] = int(cnt or 0)

        labels: list[str] = []
        ids: list[str] = []
        parents: list[str] = []
        values: list[int] = []

        root_label = "Thematic Map"
        root_id = "root"
        labels.append(root_label);
        ids.append(root_id);
        parents.append("");
        values.append(1)

        total_first_ring = 0
        for og_label, og_norm, og_total in og_items:
            og_id = f"og::{_slug(og_label)}"
            # placeholder; we'll overwrite with proper sum below
            labels.append(og_label);
            ids.append(og_id);
            parents.append(root_id);
            values.append(1)

            # Subthemes under this OG (normalised lookup)
            sub_for_og = sub_by_og_norm.get(og_norm, {}) or {}
            og_sum = 0
            for sub_label, sub_cnt in sub_for_og.items():
                sub_id = f"sub::{_slug(og_label)}::{_slug(sub_label)}"

                # Keywords under this subtheme (normalised lookup)
                kw_for_sub = kw_by_sub_norm.get(_norm_key(sub_label), {}) or {}
                sub_sum = 0
                for kw_label, kw_cnt in kw_for_sub.items():
                    kw_id = f"kw::{_slug(og_label)}::{_slug(sub_label)}::{_slug(kw_label)}"
                    v = max(1, int(kw_cnt))
                    labels.append(kw_label);
                    ids.append(kw_id);
                    parents.append(sub_id);
                    values.append(v)
                    sub_sum += v

                # Subtheme node value: sum(keywords) OR provided count, but at least 1
                sv = max(sub_sum, max(1, int(sub_cnt)))
                labels.append(sub_label);
                ids.append(sub_id);
                parents.append(og_id);
                values.append(sv)
                og_sum += sv

            # Overarching node value: sum(subthemes) OR provided og_total, at least 1
            ov = max(og_sum, max(1, int(og_total)))
            values[ids.index(og_id)] = ov  # replace placeholder
            total_first_ring += ov

        # Root equals the sum of first ring (required by branchvalues='total' for full fill)
        values[ids.index(root_id)] = max(1, int(total_first_ring))

        mind_map_fig = go.Figure(go.Sunburst(
            labels=labels,
            ids=ids,  # unique ids prevent unwanted merges
            parents=parents,
            values=values,
            branchvalues="total",
            maxdepth=4,
            sort=False,  # keep input order; avoids tiny slices being squeezed unpredictably
            insidetextorientation="radial",
            textinfo="label+value",
            hovertemplate="<b>%{label}</b><br>Occurrences: %{value}<extra></extra>",
            domain=dict(x=[0.0, 1.0], y=[0.0, 1.0]),
        ))
        # Aggressively show text and keep tiny slices visible
        mind_map_fig.update_traces(
            selector=dict(type="sunburst"),
            textfont=dict(size=9),
            leaf=dict(opacity=0.99)
        )
        mind_map_fig.update_layout(
            title=dict(text=f"{collection_name} · Thematic Sunburst", x=0.5, xanchor="center"),
            margin=dict(t=4, l=0, r=0, b=0),  # maximise drawable area
            template="plotly_white",
            uniformtext=dict(mode="show", minsize=6),  # allow very small text before hiding
        )
        # Render at an even larger resolution to reduce label suppression; image is then scaled to slide width.
    slide2 = _fig_to_slide(mind_map_fig, width_px=2800, height_px=1700)

    if slide_notes:
        slide2.notes_slide.notes_text_frame.text = (
            "Sunburst mind-map of themes. Inner ring = Overarching; middle = Subthemes; outer = Keywords."
        )

    # -------------------------- Slide 3: evolution over time --------------------------
    _cb("Slide 3: evolution over time…")
    year_col_used, years_series = _extract_year_series(df)
    work = df.copy().assign(__year=years_series)

    # Build keyword -> Overarching map (normalize keys)
    # Use flat_themes (includes synonyms)
    kw_to_og: dict[str, str] = {}
    for og_title, rows in flat_themes.items():
        for r in (rows or []):
            canon = _norm_token(r.get("term", ""))
            if canon:
                kw_to_og[canon] = og_title
            syns = r.get("synonyms") or []
            for s in syns:
                s_norm = _norm_token(s)
                if s_norm:
                    kw_to_og[s_norm] = og_title

    # Per-row events: count each (Overarching, keyword) at most once per row-year
    ev_rows = []
    if "__year" in work.columns:
        for _, row in work.iterrows():
            y = row["__year"]
            if pd.isna(y):
                continue
            y = int(y)
            toks = _tokenize_cell(row.get(theme_col))
            seen = set()
            for t in toks:
                og = kw_to_og.get(t)
                if not og:
                    continue
                key = (og, t)
                if key in seen:
                    continue
                seen.add(key)
                ev_rows.append({"Year": y, "Overarching": og, "Count": 1})

    import pandas as pd
    ts_df = pd.DataFrame(ev_rows) if ev_rows else pd.DataFrame(columns=["Year", "Overarching", "Count"])

    if ts_df.empty:
        fig_ts = go.Figure().add_annotation(text="No dated theme data to plot.", showarrow=False)
        fig_ts.update_layout(
            template="plotly_white", height=560, width=1400, title=f"{collection_name} · Theme Evolution"
        )
    else:
        # complete grid for nice stacked area continuity
        y0, y1 = int(ts_df["Year"].min()), int(ts_df["Year"].max())
        all_years = pd.DataFrame({"Year": list(range(y0, y1 + 1))})
        all_og = pd.DataFrame({"Overarching": sorted(ts_df["Overarching"].unique().tolist())})
        grid = all_years.assign(key=1).merge(all_og.assign(key=1), on="key").drop(columns="key")
        ts_full = grid.merge(
            ts_df.groupby(["Year", "Overarching"], as_index=False)["Count"].sum(),
            on=["Year", "Overarching"], how="left"
        ).fillna({"Count": 0}).sort_values(["Year", "Overarching"])

        fig_ts = go.Figure()
        for og in sorted(ts_full["Overarching"].unique()):
            sub = ts_full[ts_full["Overarching"] == og]
            fig_ts.add_trace(
                go.Scatter(
                    x=sub["Year"].astype(int),
                    y=sub["Count"].astype(int),
                    mode="lines",
                    stackgroup="one",
                    name=og,
                    hovertemplate="Year %{x}<br>%{y:,} keyword hits<extra></extra>",
                )
            )
        fig_ts.update_layout(
            template="plotly_white",
            height=560, width=1400,
            title=dict(text=f"{collection_name} · Overarching Theme Evolution Over Time", x=0.02, xanchor="left"),
            margin=dict(l=64, r=44, t=76, b=64),
            xaxis_title="Year",
            yaxis_title="Keyword hits (per document-year)",
            font=dict(family="Calibri", size=12, color="#2A2E33"),
        )
    slide3 = _fig_to_slide(fig_ts)

    if slide_notes:
        if ts_df.empty:
            slide3.notes_slide.notes_text_frame.text = "No date information found; evolution chart is empty."
        else:
            # lightweight notes: recent leaders & variability
            latest_year = int(ts_df["Year"].max())
            recent = (
                ts_df[ts_df["Year"] == latest_year]
                .groupby("Overarching", as_index=False)["Count"].sum()
                .sort_values("Count", ascending=False)
            )
            top_line = ", ".join(f"{r.Overarching} ({int(r.Count)})" for r in recent.itertuples(index=False)) or "—"
            slide3.notes_slide.notes_text_frame.text = (
                f"Latest year ({latest_year}) leaders: {top_line}.\n"
                "Stacked area shows per-document keyword hits per overarching theme."
            )

    # -------------------------- Return summary --------------------------
    return {
        "theme_col": theme_col,
        "year_col": year_col_used,
        "thematic": thematic_result,
        "keyword_to_overarching": kw_to_og,
    }

def add_keywords_slides(
    prs,
    df: "pd.DataFrame",
    *,
    data_source: str = "controlled_vocabulary_terms",
    image_width_inches: float = 10.0,
    slide_notes: bool = True,
    progress_callback=None,
    collection_name: str = "Collection",
    theme_target_groups: int = 6,
    theme_min_tag_freq: int = 1,
    **_
) -> None:
    """
    Emits FIVE slides total:
      1) Vocabulary KPI cards (no title)
      2) 2×2 panels: Treemap | Word Cloud / Top Terms
                     Keyword×Author | Keyword×ItemType
      3) Theme Groups — Weights (bar)
      4) Theme Tree — Sunburst (Group → Tag)
      5) Theme Evolution — Stacked Area (by year)
    """
    add_keyword_kpi_slide(
        prs, df,
        data_source=data_source,
        image_width_inches=image_width_inches,
        slide_notes=slide_notes,
        progress_callback=progress_callback,
    )

    add_keyword_panels_2x2_slide(
        prs, df,
        slide_title=None,
        data_source=data_source,
        image_width_inches=image_width_inches,
        slide_notes=slide_notes,
        progress_callback=progress_callback,
    )

    add_keyword_thematic(
        prs, df,
        collection_name=collection_name,
        theme_col_candidates=(data_source, "controlled_vocabulary", "controlled_vocabulary_terms", "keywords"),
        target_groups=theme_target_groups,
        min_tag_freq=theme_min_tag_freq,
        slide_notes=slide_notes,
        progress_callback=progress_callback,
        image_width_inches=image_width_inches,
    )
def add_affiliations_slides(
    prs,
    df: "pd.DataFrame",
    *,
    collection_name: str = "Collection",
    plots: list[str] | None = None,   # ignored (kept for API compat)
    top_n: int = 20,                  # will clamp to 5 where relevant
    min_collaborations: int = 2,      # unused here
    image_width_inches: float = 10.0,
    slide_notes: bool = True,
    progress_callback=None,
) -> None:
    """
    Build FOUR slides:

      1) Publications by country (map)                 -> analyze_affiliations('world_map_pubs')
      2) Institutions by country (geo-scatter map)     -> analyze_affiliations('geo_scatter_institutions')
      3) One slide, two panels: Top-5 Countries & Top-5 Institutions (bars)
      4) One slide, three panels: Top-5 Cities, Top-5 Departments, All Continents (bars)

    Dataframe columns used: institution, department, country, city, continent.
    """
    import io
    import numpy as np
    import pandas as pd
    import plotly.graph_objects as go
    from plotly.subplots import make_subplots
    from pptx.util import Inches, Emu

    k5 = min(int(top_n), 5)

    def _cb(msg: str):
        if progress_callback:
            try:
                progress_callback(f"AFFIL: {msg}")
            except Exception:
                pass

    def _fig_to_slide(fig: "go.Figure", *, width_px=1700, height_px=900):
        """
        Robust static export with fallback away from mapbox traces.
        """
        import copy
        import plotly.graph_objects as go

        def _is_mapbox(f: "go.Figure") -> bool:
            try:
                if getattr(f.layout, "mapbox", None) is not None:
                    return True
                for tr in f.data:
                    if getattr(tr, "type", "").endswith("mapbox"):
                        return True
            except Exception:
                pass
            return False

        def _convert_mapbox_to_geo(f: "go.Figure") -> "go.Figure":
            g = go.Figure()
            for tr in f.data:
                t = getattr(tr, "type", "")
                if t == "scattermapbox":
                    lon = getattr(tr, "lon", None) or getattr(tr, "longitude", None)
                    lat = getattr(tr, "lat", None) or getattr(tr, "latitude", None)
                    g.add_trace(go.Scattergeo(lon=lon, lat=lat, mode="markers",
                                             marker=dict(size=6, opacity=0.9)))
                elif t == "choroplethmapbox":
                    g.add_trace(go.Choropleth(
                        locations=getattr(tr, "locations", None),
                        z=getattr(tr, "z", None),
                        text=getattr(tr, "text", None),
                        colorscale=getattr(tr, "colorscale", None),
                    ))
                else:
                    g.add_trace(copy.deepcopy(tr))
            g.update_layout(
                template="plotly_white",
                geo=dict(scope="world", projection_type="natural earth",
                         showland=True, showcountries=True, landcolor="rgb(240,240,240)"),
                margin=dict(l=60, r=40, t=80, b=60)
            )
            return g

        try:
            png = fig.to_image(format="png", width=width_px, height=height_px, scale=2)
        except Exception:
            try:
                if _is_mapbox(fig):
                    fig2 = _convert_mapbox_to_geo(fig)
                    png = fig2.to_image(format="png", width=width_px, height=height_px, scale=2)
                else:
                    raise
            except Exception as e1:
                fb = go.Figure().add_annotation(text=f"Render fallback: {type(e1).__name__}", showarrow=False)
                fb.update_layout(template="plotly_white", width=width_px, height=height_px,
                                 title="Affiliations — Fallback Visualization")
                png = fb.to_image(format="png", width=width_px, height=height_px, scale=2)

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        left = Emu(int((prs.slide_width - Inches(image_width_inches)) / 2))
        top = Inches(0.6)
        slide.shapes.add_picture(io.BytesIO(png), left, top, width=Inches(image_width_inches))
        return slide

    def _num_only(df_like: "pd.DataFrame") -> "pd.DataFrame":
        try:
            out = df_like.select_dtypes(include=["number", "bool"]).copy()
        except Exception:
            out = pd.DataFrame()
        if out.shape[1] == 0:
            return pd.DataFrame({"_dummy": [0]})
        if out.empty:
            return pd.DataFrame([{c: 0 for c in out.columns}])
        return out

    # ---------------- Slide 1: Publications by country (map)
    _cb("Slide 1: world map — publications by country")
    try:
        cc1, fig1 = analyze_affiliations(
            df, {"plot_type": "world_map_pubs", "top_n": k5}, progress_callback=progress_callback
        )
    except Exception as e:
        fig1 = go.Figure().add_annotation(text=f"Error: {e}", showarrow=False)
        fig1.update_layout(template="plotly_white", title=f"{collection_name} · Publications by Country")
        cc1 = pd.DataFrame()
    slide1 = _fig_to_slide(fig1, width_px=1700, height_px=900)
    if slide_notes:
        try:
            notes1 = ai_notes_from_plot(
                fig1, raw_df=_num_only(cc1), collection_name=collection_name,
                analysis_suffix="affil_world_map_pubs",
                title_for_prompt=f"{collection_name} · Publications by Country",
                prompt_purpose="Summarise geographical concentration and outliers. Flag likely geocoding gaps.",
                max_tokens=600, effort="low", base_logs_dir=None, _cb=_cb
            )
        except Exception:
            notes1 = ""
        slide1.notes_slide.notes_text_frame.text = notes1 or "Publications by country."

    # ---------------- Slide 2: Institutions by country (geo-scatter)
    _cb("Slide 2: geo-scatter — institutions by country")
    try:
        cc2, fig2 = analyze_affiliations(
            df, {"plot_type": "geo_scatter_institutions", "top_n": k5}, progress_callback=progress_callback
        )
        # analyze_affiliations returns (df, fig) here; cc2 may be a DataFrame of plotted points
    except Exception as e:
        fig2 = go.Figure().add_annotation(text=f"Error: {e}", showarrow=False)
        fig2.update_layout(template="plotly_white", title=f"{collection_name} · Institutions by Country")
        cc2 = pd.DataFrame()
    slide2 = _fig_to_slide(fig2, width_px=1700, height_px=900)
    if slide_notes:
        try:
            notes2 = ai_notes_from_plot(
                fig2, raw_df=_num_only(cc2), collection_name=collection_name,
                analysis_suffix="affil_geo_scatter_institutions",
                title_for_prompt=f"{collection_name} · Institutions by Country",
                prompt_purpose="Identify institutional hubs and any country-level gaps. Note normalisation needs.",
                max_tokens=600, effort="low", base_logs_dir=None, _cb=_cb
            )
        except Exception:
            notes2 = ""
        slide2.notes_slide.notes_text_frame.text = notes2 or "Institutions plotted over countries."

    # ---------------- Slide 3: Top-5 Countries & Top-5 Institutions (1×2 bars)
    _cb("Slide 3: top-5 countries & institutions bars")
    def _clean_series(s: "pd.Series") -> "pd.Series":
        if s is None:
            return pd.Series([], dtype="object")
        ss = s.astype(str).str.strip().replace({"": np.nan, "nan": np.nan, "None": np.nan})
        return ss.dropna()

    countries = _clean_series(df.get("country"))
    institutions = _clean_series(df.get("institution"))

    top_countries = (
        countries.value_counts().head(k5).reset_index(name="count").rename(columns={"index": "country"})
    )
    top_institutions = (
        institutions.value_counts().head(k5).reset_index(name="count").rename(columns={"index": "institution"})
    )

    if top_countries.empty:
        top_countries = pd.DataFrame({"country": ["—"], "count": [0]})
    if top_institutions.empty:
        top_institutions = pd.DataFrame({"institution": ["—"], "count": [0]})

    top_countries = top_countries.sort_values(["count", "country"], ascending=[True, True], kind="mergesort")
    top_institutions = top_institutions.sort_values(["count", "institution"], ascending=[True, True], kind="mergesort")

    fig3 = make_subplots(
        rows=1, cols=2,
        subplot_titles=(f"Top {k5} Countries", f"Top {k5} Institutions"),
        horizontal_spacing=0.12
    )
    fig3.add_trace(
        go.Bar(x=top_countries["count"], y=top_countries["country"], orientation="h",
               text=top_countries["count"], textposition="outside",
               hovertemplate="%{y}: %{x}<extra></extra>"),
        row=1, col=1
    )
    fig3.add_trace(
        go.Bar(x=top_institutions["count"], y=top_institutions["institution"], orientation="h",
               text=top_institutions["count"], textposition="outside",
               hovertemplate="%{y}: %{x}<extra></extra>"),
        row=1, col=2
    )
    fig3.update_yaxes(categoryorder="array", categoryarray=top_countries["country"].tolist(), row=1, col=1)
    fig3.update_yaxes(categoryorder="array", categoryarray=top_institutions["institution"].tolist(), row=1, col=2)
    fig3.update_xaxes(title_text="Count", row=1, col=1)
    fig3.update_xaxes(title_text="Count", row=1, col=2)
    fig3.update_layout(
        template="plotly_white",
        height=900, width=1700,
        title=dict(text=f"{collection_name} · Top {k5} Countries & Institutions", x=0.02, xanchor="left"),
        margin=dict(l=60, r=40, t=80, b=60),
        uniformtext=dict(minsize=10, mode="hide"),
        font=dict(family="Calibri", size=12, color="#2A2E33"),
        showlegend=False,
    )
    slide3 = _fig_to_slide(fig3, width_px=1700, height_px=900)
    if slide_notes:
        try:
            raw3 = pd.DataFrame({
                "country_count": [int(x) for x in top_countries["count"]],
                "institution_count": [int(x) for x in top_institutions["count"]],
            })
            notes3 = ai_notes_from_plot(
                fig3, raw_df=_num_only(raw3), collection_name=collection_name,
                analysis_suffix="affil_top5_countries_institutions",
                title_for_prompt=f"{collection_name} · Top {k5} Countries & Institutions",
                prompt_purpose="Compare concentration patterns across the two top-5 lists; note divergences.",
                max_tokens=600, effort="low", base_logs_dir=None, _cb=_cb
            )
        except Exception:
            notes3 = ""
        slide3.notes_slide.notes_text_frame.text = notes3 or "Top-5 comparison."

    # ---------------- Slide 4: Cities/Departments/Continents (top-5, top-5, all)
    _cb("Slide 4: top-5 cities, top-5 departments, all continents")
    add_affiliations_tripanel_top5(
        prs,
        df,
        collection_name=collection_name,
        top_n=5,
        image_width_inches=image_width_inches,
        slide_notes=slide_notes,
        progress_callback=progress_callback,
    )

    _cb("all affiliation slides built.")

# code to be replaced
# (no 'add_affiliations_tripanel_top5' function exists)

def add_affiliations_tripanel_top5(
    prs,
    df: "pd.DataFrame",
    *,
    collection_name: str = "Collection",
    top_n: int = 5,                      # fixed to 5 as requested
    image_width_inches: float = 10.0,
    slide_notes: bool = True,
    progress_callback=None,
) -> None:
    """
    One slide with THREE panels:
      • Top-5 Cities
      • Top-5 Departments
      • All Continents
    """
    import io
    import numpy as np
    import pandas as pd
    import plotly.graph_objects as go
    from plotly.subplots import make_subplots
    from pptx.util import Inches, Emu

    def _cb(msg: str):
        if progress_callback:
            try:
                progress_callback(f"AFFIL: {msg}")
            except Exception:
                pass

    def _fig_to_slide(fig: "go.Figure", *, width_px=1700, height_px=900):
        try:
            png = fig.to_image(format="png", width=width_px, height=height_px, scale=2)
        except Exception as e1:
            fb = go.Figure().add_annotation(text=f"Render fallback: {type(e1).__name__}", showarrow=False)
            fb.update_layout(template="plotly_white", width=width_px, height=height_px,
                             title="Affiliations — Fallback Visualization")
            png = fb.to_image(format="png", width=width_px, height=height_px, scale=2)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        left = Emu(int((prs.slide_width - Inches(image_width_inches)) / 2))
        top = Inches(0.6)
        slide.shapes.add_picture(io.BytesIO(png), left, top, width=Inches(image_width_inches))
        return slide

    def _top_counts(series: "pd.Series", label: str, n: int) -> "pd.DataFrame":
        if series is None:
            return pd.DataFrame({label: ["—"], "count": [0]})
        s = series.astype(str).str.strip()
        s = s.replace({"": np.nan, "nan": np.nan, "None": np.nan}).dropna()
        out = (
            s.value_counts()
             .head(int(n))
             .reset_index(name="count")
             .rename(columns={"index": label})
        )
        if out.empty:
            out = pd.DataFrame({label: ["—"], "count": [0]})
        return out.sort_values(["count", label], ascending=[True, True], kind="mergesort")

    # top-5 for cities and departments; all for continents
    cities = _top_counts(df.get("city"), "city", top_n)
    depts  = _top_counts(df.get("department"), "department", top_n)

    # Continents: take all categories
    cont_series = df.get("continent")
    if cont_series is None:
        continents = pd.DataFrame({"continent": ["—"], "count": [0]})
    else:
        s = cont_series.astype(str).str.strip()

        def _split_cont_tokens(x):
            x = "" if x is None else str(x).strip()
            if x.lower() in {"", "nan", "none", "null"}:
                return []
            # values can be like "Europe;none" → count only "Europe"
            parts = [t.strip() for t in x.replace("\n", ";").split(";")]
            parts = [t for t in parts if t and t.lower() not in {"none", "nan", "null", "-"}]
            return parts

        tokens = s.apply(_split_cont_tokens)
        flat = [t for lst in tokens for t in (lst if isinstance(lst, list) else [])]

        if not flat:
            continents = pd.DataFrame({"continent": ["—"], "count": [0]})
        else:
            continents = (
                pd.Series(flat, dtype=str)
                .value_counts()
                .reset_index(name="count")
                .rename(columns={"index": "continent"})
                .sort_values(["count", "continent"], ascending=[True, True], kind="mergesort")
            )
            if continents.empty:
                continents = pd.DataFrame({"continent": ["—"], "count": [0]})

    fig = make_subplots(
        rows=1, cols=3,
        subplot_titles=(f"Top {top_n} Cities", f"Top {top_n} Departments", "Continents"),
        horizontal_spacing=0.08
    )

    fig.add_trace(go.Bar(x=cities["count"], y=cities["city"], orientation="h",
                         text=cities["count"], textposition="outside",
                         hovertemplate="%{y}: %{x}<extra></extra>"), row=1, col=1)
    fig.add_trace(go.Bar(x=depts["count"], y=depts["department"], orientation="h",
                         text=depts["count"], textposition="outside",
                         hovertemplate="%{y}: %{x}<extra></extra>"), row=1, col=2)
    fig.add_trace(go.Bar(x=continents["count"], y=continents["continent"], orientation="h",
                         text=continents["count"], textposition="outside",
                         hovertemplate="%{y}: %{x}<extra></extra>"), row=1, col=3)

    fig.update_yaxes(categoryorder="array", categoryarray=cities["city"].tolist(), row=1, col=1)
    fig.update_yaxes(categoryorder="array", categoryarray=depts["department"].tolist(), row=1, col=2)
    fig.update_yaxes(categoryorder="array", categoryarray=continents["continent"].tolist(), row=1, col=3)

    fig.update_xaxes(title_text="Count", row=1, col=1)
    fig.update_xaxes(title_text="Count", row=1, col=2)
    fig.update_xaxes(title_text="Count", row=1, col=3)

    fig.update_layout(
        template="plotly_white",
        height=900, width=1700,
        title=dict(text=f"{collection_name} · Cities, Departments, Continents", x=0.02, xanchor="left"),
        margin=dict(l=60, r=40, t=80, b=60),
        uniformtext=dict(minsize=10, mode="hide"),
        font=dict(family="Calibri", size=12, color="#2A2E33"),
        showlegend=False,
    )

    slide = _fig_to_slide(fig, width_px=1700, height_px=900)

    if slide_notes:
        try:
            raw = pd.DataFrame({
                "cities": [int(x) for x in cities["count"]],
                "departments": [int(x) for x in depts["count"]],
                "continents": [int(x) for x in continents["count"]],
            })
            notes = ai_notes_from_plot(
                fig, raw_df=raw, collection_name=collection_name,
                analysis_suffix="affil_tripanel_top5_cities_depts_allcontinents",
                title_for_prompt=f"{collection_name} · Cities, Departments, Continents",
                prompt_purpose="Contrast city and department concentration; relate to continent totals. Suggest normalisation fixes.",
                max_tokens=700, effort="low", base_logs_dir=None, _cb=_cb
            )
        except Exception:
            notes = ""
        slide.notes_slide.notes_text_frame.text = notes or "Top-5 cities and departments; all continents."




# ---------- helpers: registry / dispatcher (core) ----------
def _add_transition_slide(prs: "Presentation", title_text: str) -> None:
    from pptx.util import Pt
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    if slide.shapes.title:
        slide.shapes.title.text = title_text
        try:
            slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(40)
        except Exception:
            pass

def _normalise_section_token(s: str) -> tuple[str, str | None]:
    if not isinstance(s, str):
        return str(s), None
    s = s.strip().lower()
    if "." in s:
        a, b = s.split(".", 1)
        return a.strip(), (b.strip() or None)
    if ":" in s:
        a, b = s.split(":", 1)
        return a.strip(), (b.strip() or None)
    if s.startswith("ngrams_"):
        return "ngrams", s.replace("ngrams_", "", 1)
    return s, None
def _get_section_registry() -> dict[str, Any]:
    def _maybe(fn_name: str):
        fn = globals().get(fn_name)
        return fn if callable(fn) else None

    return {
        "data_summary": {
            "title": "Data Summary",
            "handler": _maybe("add_data_summary_slides"),
            "defaults": {},
            "aliases": [],
        },
        "affiliations": {
            "title": "Affiliations",
            "handler": _maybe("add_affiliations_slides"),
            "defaults": {
                "plots": [
                    "world_map_pubs","bubble_map_pubs","geo_scatter_institutions","country_bar",
                    "country_pie","institution_bar","department_sunburst",
                    "collab_network_country","collab_network_institution","chord_inst_country",
                ],
                "top_n": 20,
                "min_collaborations": 2,
            },
            "aliases": [],
        },
        "authors": {
            "title": "Authors",
            "handler": _maybe("add_authors_slides"),
            "defaults": {},
            "aliases": [],
        },
        "keywords": {
            "title": "Keyword and Thematic analysis",
            "handler": _maybe("add_keywords_slides"),
            "defaults": {"data_source": "controlled_vocabulary_terms"},
            "aliases": [],
        },
        "ngrams": {
            "title": "N-grams",
            "subsections": {
                "basic": {
                    "title": "Basic N-gram Analysis",
                    "handler": _maybe("add_basic_ngram_analysis_slides"),
                    "defaults": {"keywords": [], "include_abstract_title": True, "transition_seconds": None},
                    "aliases": [],
                },
                "higher": {
                    "title": "Higher-order N-gram Analysis",
                    "handler": _maybe("add_ngram_higher_order_slides"),
                    "defaults": {"keywords": ["leadership","equity","diversity","inclusion"], "n_values": (2,3,4,5), "data_sources": ("fulltext",), "top_n": 20},
                    "aliases": [],
                },
            },
        },

        # -------- NEW: non-lettered main sections + aliases ----------
        "Scope_and_shape": {
            "title": "Corpus shape and scope",
            "handler": _maybe("add_scope_and_shape_section") or _maybe("shape_scope"),
            "defaults": {"year_range": (2010, 2025)},
            "aliases": ["scope_and_shape","scope and shape","Scope_And_Shape","scope","Scope and Shape"],
        },
        "Authorship_institution": {
            "title": "Authorship and institutional influence",
            "handler": _maybe("add_authorship_and_institution_section"),
            "defaults": {"top_n_authors": 20, "include_coauth_network": True, "include_institution_bar": True},
            "aliases": ["authorship_institution","Authorship_Institution","authorship"],
        },
        "Thematic_method": {
            "title": "Thematic and methodological mapping",
            "handler": _maybe("add_thematic_and_method_section"),
            "defaults": {
                "cross_tabs": [
                    ["focus_type_value", "publisher_type"],
                    ["empirical_theoretical", "sector_focus_value"],
                ],
                "phase_tag": "phase_focus_value",
                "year_col": "publication_year",
            },
            "aliases": ["thematic_method","Thematic_Method","thematic"],
        },
        "Geo_sector": {
            "title": "Geographic and sectoral coverage",
            "handler": _maybe("add_geo_and_sector_section"),
            "defaults": {"country_tag": "country_focus_value", "sector_tag": "sector_focus_value"},
            "aliases": ["geo_sector","Geo_Sector","geography","sector"],
        },
        "Citations_influence": {
            "title": "Citation and influence indicators",
            "handler": _maybe("add_citations_and_influence_section"),
            "defaults": {"citations_col": "citations", "top_n": 10},
            "aliases": ["citations_influence","Citations_Influence","citations"],
        },
        "Theory_empirical": {
            "title": "Theoretical–empirical balance",
            "handler": _maybe("add_theory_empirical_section"),
            "defaults": {"focus_type_col": "focus_type_value", "method_col": "method_type"},
            "aliases": ["theory_empirical","Theory_Empirical","theory"],
        },
    }



def _resolve_section_key(section_key: str, registry: dict[str, Any]) -> str | None:
    import re
    def norm(s: str) -> str:
        s = (s or "").strip()
        s = s.replace("&", "and")
        s = s.replace("-", "_").replace(" ", "_")
        s = re.sub(r"_+", "_", s)
        return s.lower()
    if section_key in registry:
        return section_key
    nk = norm(section_key)
    for k, meta in registry.items():
        if nk == norm(k):
            return k
        for a in meta.get("aliases", []):
            if nk == norm(a):
                return k
    return None



def _coerce_df(df):
    """
    Accept DataFrame, (DataFrame, ...), [DataFrame, ...], or {'df': DataFrame}.
    Return a pandas DataFrame or raise TypeError.
    """
    import pandas as pd
    if isinstance(df, pd.DataFrame):
        return df
    if isinstance(df, (list, tuple)):
        for x in df:
            if isinstance(x, pd.DataFrame):
                return x
    if isinstance(df, dict) and isinstance(df.get("df"), pd.DataFrame):
        return df["df"]
    raise TypeError(f"_coerce_df: expected DataFrame-like, got {type(df)}")

def _filter_kwargs_for_callable(fn, extra_kwargs: dict) -> dict:
    """
    Keep only kwargs that the callable 'fn' will accept (by signature).
    """
    import inspect
    try:
        sig = inspect.signature(fn)
    except (TypeError, ValueError):
        # Builtins / no signature → pass nothing
        return {}
    params = sig.parameters
    allowed = set()
    has_var_kw = False
    for name, p in params.items():
        if p.kind == inspect.Parameter.VAR_KEYWORD:
            has_var_kw = True
        elif p.kind in (inspect.Parameter.KEYWORD_ONLY, inspect.Parameter.POSITIONAL_OR_KEYWORD):
            allowed.add(name)
    if has_var_kw:
        return dict(extra_kwargs or {})
    return {k: v for k, v in (extra_kwargs or {}).items() if k in allowed}
def _invoke_handler(prs, df, collection_name, section_title, handler, kwargs, progress_callback):
    """
    Robust invocation:
      1) Coerce df → DataFrame (handles tuples from load_data_from_source_for_widget)
      2) Merge kwargs without duplicating 'collection_name'
      3) Try call; on TypeError, strip unknown kwargs and retry
      4) On failure, add a fallback text slide with a clear message
    """
    # helper to merge kwargs while avoiding duplicate 'collection_name'
    def _merge_kwargs(base_kwargs: dict) -> dict:
        merged = dict(base_kwargs or {})
        # inject progress_callback unconditionally (it'll be filtered if not accepted)
        merged["progress_callback"] = progress_callback
        # only inject collection_name if caller didn't already supply it
        merged.setdefault("collection_name", collection_name)
        return merged

    if not callable(handler):
        msg = f"No handler available for '{section_title}'."
    else:
        # Step 1: coerce df if possible
        try:
            df_ok = _coerce_df(df)
        except Exception:
            df_ok = df

        # Step 2/3: attempt call, then retry with filtered kwargs if needed
        try:
            handler(prs, df_ok, **_merge_kwargs(kwargs))
            return
        except TypeError:
            try:
                clean_kwargs = _filter_kwargs_for_callable(handler, _merge_kwargs(kwargs))
                handler(prs, df_ok, **clean_kwargs)
                return
            except Exception as e2:
                msg = f"Handler error in '{section_title}': {e2}"
        except Exception as e:
            msg = f"Handler error in '{section_title}': {e}"

    # Step 4: fallback slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title & Content
    if slide.shapes.title:
        slide.shapes.title.text = section_title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.paragraphs[0].text = msg



def _coerce_sections(sections: Iterable[Any]) -> list[tuple[str, str | None, dict]]:
    out: list[tuple[str, str | None, dict]] = []
    sections = sections or []
    for item in sections:
        if isinstance(item, str):
            sec, sub = _normalise_section_token(item)
            out.append((sec, sub, {}))
        elif isinstance(item, dict):
            for k, v in item.items():
                sec, sub = _normalise_section_token(k)
                if isinstance(v, dict):
                    out.append((sec, sub, v))
                elif isinstance(v, list) and sub is None:
                    for subname in v:
                        out.append((sec, str(subname).strip().lower(), {}))
                else:
                    out.append((sec, sub, {}))
        else:
            out.append((str(item), None, {}))
    return out

def build_bibliometrics_pptx_core(
    df: "pd.DataFrame",
    sections: list[Any],
    *,
    output_path: str | Path,
    collection_name: str = "Collection",
    options: dict[str, Any] | None = None,
    progress_callback=None,
) -> Path:
    try:
        df = _coerce_df(df)
    except Exception:
        # let handlers render a clear fallback slide via _invoke_handler
        pass

    prs = Presentation()
    registry = _get_section_registry()
    options = options or {}
    todo = _coerce_sections(sections)

    # local resolver (robust to case/underscores/spaces); uses registry aliases if defined
    def _resolve_section_key(section_key: str, registry_map: dict[str, Any]) -> str | None:
        import re
        def norm(s: str) -> str:
            s = (s or "").strip()
            s = s.replace("&", "and").replace("-", "_").replace(" ", "_")
            return re.sub(r"_+", "_", s).lower()

        if section_key in registry_map:
            return section_key
        nk = norm(section_key)
        for k, meta in registry_map.items():
            if nk == norm(k):
                return k
            for a in meta.get("aliases", []):
                if nk == norm(a):
                    return k
        return None

    # normalise subsection keys too (simple normaliser)
    def _norm_subkey(subkey: str, subreg: dict[str, Any]) -> str | None:
        import re
        def n(s: str) -> str:
            return re.sub(r"_+", "_", s.replace(" ", "_")).lower()

        if subkey in subreg:
            return subkey
        target = n(subkey)
        for k in subreg.keys():
            if n(k) == target:
                return k
        return None

    def _opt_for(sec: str, sub: str | None) -> dict:
        k_sub = f"{sec}.{sub}" if sub else None
        return ((options.get(k_sub, {}) if k_sub else {}) | options.get(sec, {}))

    for sec, sub, per_item_opts in todo:
        # resolve section key via aliases/normalisation
        resolved = _resolve_section_key(sec, registry)
        if not resolved:
            _add_transition_slide(prs, sec.title())
            _invoke_handler(prs, df, collection_name, sec.title(), handler=None, kwargs={},
                            progress_callback=progress_callback)
            continue
        entry = registry[resolved]

        if sub:
            subsecs = (entry.get("subsections") or {})
            subkey = _norm_subkey(sub, subsecs)
            subentry = subsecs.get(subkey) if subkey else None
            trans_title = (subentry or {}).get("title") or f"{entry['title']} — {(subkey or sub).title()}"
            _add_transition_slide(prs, trans_title)

            if subentry and subentry.get("handler"):
                kwargs = {**subentry.get("defaults", {}), **_opt_for(resolved, subkey or sub), **per_item_opts}
                _invoke_handler(prs, df, collection_name, trans_title, subentry["handler"], kwargs, progress_callback)
            else:
                _invoke_handler(prs, df, collection_name, trans_title, handler=None, kwargs={},
                                progress_callback=progress_callback)
            continue

        trans_title = entry.get("title", resolved.title())
        _add_transition_slide(prs, trans_title)
        handler = entry.get("handler")
        kwargs = {**entry.get("defaults", {}), **_opt_for(resolved, None), **per_item_opts}
        _invoke_handler(prs, df, collection_name, trans_title, handler, kwargs, progress_callback)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path

def _default_sections_all() -> list[Any]:
    """
    All available sections including n-gram subsections.
    Uses *non-lettered* keys per spec.
    """
    return [
        "data_summary",
        "affiliations",
        "authors",
        "keywords",
        {"ngrams": ["basic", "higher"]},
        # main sections (no letter prefixes)
        "Scope_and_shape",
        "Authorship_institution",
        "Thematic_method",
        "Geo_sector",
        "Citations_influence",
        "Theory_empirical",
    ]


def _map_corpus_source_to_data_source(corpus_source: str) -> str:
    """
    Normalise friendly corpus_source to the data_source expected by keyword slides.
    """
    s = (corpus_source or "").strip().lower()
    if s in {"fulltext", "full_text", "pdf", "body"}:
        return "full_text"
    if s in {"title", "titles"}:
        return "title"
    if s in {"abstract", "abstracts"}:
        return "abstract"
    # default → controlled vocabulary / keywords column
    return "controlled_vocabulary_terms"


def _build_df_for_collection(self, collection_name: str, cache: bool = True):
        """
        Build a DataFrame for items in `collection_name` with fields used by extract_content_for_keywords:
          - key, title, year, authors, abstract, pdf_path
        Prefers direct local PDF path via `self.get_pdf_path_for_item(key)`.
        """

        rows = []
        items = self.get_all_items(collection_name=collection_name, cache=cache) or []
        for it in items:
            data = (it or {}).get("data", {})
            key = (it or {}).get("key")
            title = data.get("title", "")
            date_val = data.get("date") or data.get("year") or ""
            year = str(date_val)[:4] if date_val else ""
            creators = data.get("creators", []) or []
            authors = []
            for c in creators:
                if isinstance(c, dict):
                    if c.get("lastName") or c.get("firstName"):
                        last = (c.get("lastName") or "").strip()
                        first = (c.get("firstName") or "").strip()
                        if last and first:
                            authors.append(f"{last}, {first}")
                        else:
                            authors.append((last or first))
                    elif c.get("name"):
                        authors.append(str(c.get("name")))
            try:
                pdf_path = self.get_pdf_path_for_item(key)
            except Exception:
                pdf_path = None

            rows.append({
                "key": key,
                "title": title,
                "year": year,
                "authors": "; ".join(a for a in authors if a),
                "abstract": (data.get("abstractNote") or "").strip(),
                "pdf_path": pdf_path or "",
            })

        import pandas as pd
        return pd.DataFrame(rows, columns=["key", "title", "year", "authors", "abstract", "pdf_path"])


def build_bibliometrics_pptx(
    collection_name: str,

    *,
    keywords: list[str] | None = None,
    top_ngram: int = 20,
    corpus_source: str = "controlled_vocabulary_terms",  # or "full_text"|"title"|"abstract"
    collection_name_for_cache: str | None = None,
    zotero_client=zt,
    title: str | None = "Bibliometric Snapshot",
    subtitle: str | None = None,
    include: list | dict | None = None,
    progress_callback=None,
    slide_notes: bool = True,
        csv_path=None,
    image_width_inches: float = 10.0,
        rebuild=False
):
    """
    High-level convenience wrapper that:
      1) Loads df from Zotero.
      2) Builds a section list (all sections if 'include' is None).
      3) Maps user parameters to per-section options.
      4) Calls the core dispatcher to create the deck.

    Returns: Path to the written PPTX.
    """
    # 1) Load data
    if zotero_client is None and csv_path is None:
        raise ValueError("zotero_client is required to build the deck from a collection name.")
    df =None
    pptx_path=""
    if csv_path:
        csvp = Path(csv_path)
        df = pd.read_csv(csvp)
        pptx_path = csvp.with_name(f"{csvp.stem}_bibliometrics.pptx")
    else:
        df,items,_ =load_data_from_source_for_widget(collection_name=collection_name,source_type="zotero",file_path=None
)
        pptx_path = Path(f"{collection_name.replace(' ', '_')}_bibliometrics.pptx")

    print("df columns:", df.columns.tolist())
    # 2) Build section plan
    if include is None:
        sections = _default_sections_all()
    elif isinstance(include, dict) or isinstance(include, list):
        sections = include
    else:
        # allow comma-separated string
        sections = [s.strip() for s in str(include).split(",") if s.strip()]

    # 3) Global → per-section options
    #    - keywords for n-grams slides
    #    - corpus_source → data_source for keyword slides
    #    - cache and ui params passed where supported
    kws = (keywords or [])
    data_source = _map_corpus_source_to_data_source(corpus_source)

    options: dict[str, Any] = {
        "keywords": {
            "data_source": (
                "controlled_vocabulary_terms"
                if str(data_source).lower().replace(" ", "").replace("-", "_") in {"fulltext", "full_text", "full",
                                                                                   "pdf"}
                else data_source
            ),
            "collection_name_for_cache": collection_name_for_cache or collection_name,
            "zotero_client": zotero_client,
            "slide_notes": slide_notes,
            "image_width_inches": image_width_inches,
        },
        "ngrams.basic": {
            "keywords": kws,
            "collection_name": collection_name,
            "collection_name_for_cache": collection_name_for_cache or collection_name,
            "zotero_client": zotero_client,
            "slide_notes": slide_notes,
            "image_width_inches": image_width_inches,
            "rebuild": rebuild,
        },
        "ngrams.higher": {
            "keywords": kws,
            "top_n": top_ngram,
            "collection_name": collection_name,
            "collection_name_for_cache": collection_name_for_cache or collection_name,
            "zotero_client": zotero_client,
            "slide_notes": slide_notes,
            "image_width_inches": image_width_inches,
            "rebuild": rebuild,
        },
        "authors": {
            "slide_notes": slide_notes,
            "image_width_inches": image_width_inches,
            "include": None,
        },
        "affiliations": {
            "plots": [
                "world_map_pubs",
                "bubble_map_pubs",
                "geo_scatter_institutions",
                "country_bar",
                "country_pie",
                "institution_bar",
                "department_sunburst",
                "collab_network_country",
                "collab_network_institution",
                "chord_inst_country",
            ],
            "top_n": 20,
            "min_collaborations": 2,
            "collection_name": collection_name,
            "slide_notes": slide_notes,
            "image_width_inches": image_width_inches,
        },
        "temporal_trends": {
            "slide_notes": slide_notes,
            "image_width_inches": image_width_inches,
        },
        "metadata": {
            "slide_notes": slide_notes,
            "image_width_inches": image_width_inches,
        },
        "data_summary": {
            "slide_notes": slide_notes,
        },

        # ---- NEW MAIN LETTER SECTIONS ----
        "Scope_and_shape": {
            "title": "Corpus shape and scope",
            "collection_name": collection_name,
            "slide_notes": slide_notes,
            "image_width_inches": image_width_inches,
            "year_range": (2010, 2025),
        },
        "Authorship_institution": {
            "title": "Authorship and institutional influence",
            "top_n_authors": 20,
            "include_coauth_network": True,
            "include_institution_bar": True,
            "slide_notes": slide_notes,
            "image_width_inches": image_width_inches,
        },
        "Thematic_method": {
            "title": "Thematic and methodological mapping",
            "cross_tabs": [
                ["focus_type_value", "publisher_type"],
                ["empirical_theoretical", "sector_focus_value"],
            ],
            "phase_tag": "phase_focus_value",
            "year_col": "publication_year",
            "slide_notes": slide_notes,
            "image_width_inches": image_width_inches,
        },
        "Geo_sector": {
            "title": "Geographic and sectoral coverage",
            "country_tag": "country_focus_value",
            "sector_tag": "sector_focus_value",
            "iso_country_col": None,
            "slide_notes": slide_notes,
            "image_width_inches": image_width_inches,
        },
        "Citations_influence": {
            "title": "Citation and influence indicators",
            "citations_col": "citations",
            "top_n": 10,
            "slide_notes": slide_notes,
            "image_width_inches": image_width_inches,
        },
        "Theory_empirical": {
            "title": "Theoretical–empirical balance",
            "focus_type_col": "focus_type_value",
            "method_col": "method_type",
            "slide_notes": slide_notes,
            "image_width_inches": image_width_inches,
        },
    }

    # 4) Build deck
    out = build_bibliometrics_pptx_core(
        df,
        sections=sections,
        output_path=pptx_path,
        collection_name=collection_name,
        options=options,
        progress_callback=progress_callback,
    )
    return out
# df = _build_df_for_collection(zt,collection_name="b.included", cache=True)




def update_excel_counts_single_scan(

    excel_path: str | Path,
    *,
    sheet_name: str | int | None = None,
    skills_col_candidates: tuple[str, ...] = ("skills", "skill", "Skills", "Skill"),

) -> List:
    """
    One-shot scan for ALL skills:
      1) df = zt._build_df_for_collection(collection_name)
      2) read Excel (first sheet if not specified)
      3) collect DISTINCT skill tokens across the sheet
      4) run extract_content_for_keywords once with all skills as globals
      5) aggregate per-skill counts
      6) compute per-row totals (sum across that row's tokens)
      7) write count column immediately to the right of 'skills', save in place, print path
    """
    excel_path = Path(excel_path)


    # 2) Read Excel, tolerate multi-sheet dict
    xdf = pd.read_excel(excel_path, sheet_name=sheet_name)
    if isinstance(xdf, dict):
        first_sheet = next(iter(xdf))
        xdf = xdf[first_sheet]

    # Locate the skills column
    cols_norm = {c.lower(): c for c in xdf.columns}
    skills_col = None
    for cand in skills_col_candidates:
        if cand.lower() in cols_norm:
            skills_col = cols_norm[cand.lower()]
            break
    if not skills_col:
        raise ValueError(f"Skills column not found. Tried: {skills_col_candidates}. Columns: {list(xdf.columns)}")

    def _norm_text(s: str) -> str:
        return (s or "").strip().lower()

    def _split_tokens(cell: str) -> list[str]:
        if not isinstance(cell, str):
            return []
        tmp = cell.replace(";", ",").replace("|", ",")
        return [t.strip() for t in tmp.split(",") if t.strip()]

    def _norm(s: str) -> str:
        return (s or "").strip().lower()
    # 3) Collect DISTINCT tokens (preserve first-seen casing for output mapping)
    distinct_tokens, seen = [], set()
    for val in xdf[skills_col].astype(str).fillna(""):
        for tok in _split_tokens(val):
            n = _norm(tok)
            if n and n not in seen:
                seen.add(n)
                distinct_tokens.append(tok)

    return distinct_tokens
def build_theme_groups(
    df,
    theme_col: str = "controlled_vocabulary_terms",
    key_col: str = "key",
    *,
    target_groups: int = 5,
    min_tag_freq: int = 2,                   # ↑ lift frequency floor
    synonyms: dict[str, str] | None = None,
    stopwords: set[str] | None = None,
    separators: str = r"[;|,\n]+",
    # --- hardened defaults ---
    weight_mode: str = "ppmi",               # "ppmi" | "assoc" | "jaccard" | "count"
    min_edge_weight: float = 0.50,           # ↑ prune harder (PPMI ~ log2-ish; assoc/jaccard in [0,1])
    auto_balance: bool = True,
    balance_max_share: float = 0.50,         # ↑ no group > ~50% of nodes
    balance_max_iter: int = 5,
    use_louvain: bool = True,                # tries Leiden if available; else Louvain; else greedy
    random_seed: int = 42,
) -> dict:
    """
    Build exactly `target_groups` clusters from #theme:* tags using a weighted co-occurrence graph.

    Patches applied (per your request):
      • PPMI (default) or assoc/jaccard, with harder edge pruning.
      • Frequency floor min_tag_freq>=2.
      • Balance tightened (max 50% share; more iterations).
      • Glue-stopwords & degree P95 de-hubbing (drop weakest half edges of ultra-hubs).
      • Resolution sweep (Leiden/Louvain) + coherence proxy (c-TF-IDF) to pick best split.
      • Two-pass: recluster the largest group with stricter thresholds, then re-coerce to K.
      • Optional semantic assist: w_total = 0.8*PPMI_norm + 0.2*cosine(embedding) when available.
      • Bridges: soft overlaps where 2nd-best affinity ≈ best.

    Output (superset of original schema):
      {
        "params": {...},
        "tag_counts": {"#theme:xyz": 12, ...},
        "groups": [
          {"id": "G1", "tags": [...], "label": "…", "top_terms": [["term", score], ...]},
          ...
        ],
        "tag_to_group": {"#theme:...": "G1", ...},
        "bridges": {"#theme:...": ["G2","G1"], ...}
      }
    """
    import re, math, random
    from collections import Counter, defaultdict

    try:
        import pandas as pd  # noqa: F401
    except Exception as e:
        raise RuntimeError("pandas is required") from e

    rng = random.Random(random_seed)

    # ----------------------------- helpers -----------------------------
    def _is_nan(x) -> bool:
        try:
            import pandas as _pd
            return _pd.isna(x)
        except Exception:
            return x is None

    def _canon_term(s: str) -> str:
        s = str(s or "").strip()
        if not s:
            return ""
        s = re.sub(r"^\s*#?\s*theme\s*:\s*", "", s, flags=re.I)
        s = re.sub(r"\s+", " ", s).strip()
        base = s.casefold()
        if synonyms:
            repl = synonyms.get(base)
            if repl:
                return repl.strip()
        return s

    def _prefixed(s: str) -> str:
        s = str(s or "").strip()
        if not s:
            return ""
        if re.match(r"^\s*#?\s*theme\s*:", s, flags=re.I):
            s = _canon_term(s)
        return f"#theme:{s}"

    # default glue stopwords (additive with user stopwords)
    _glue = {
        "case study","multiple-case","single case","interview","semi-structured","content analysis","documentary",
        "analysis","evaluation","evidence","framework","approach","method","methodology","policy","policies",
        "governance","governance model","programme","program","project","system","systems","strategy","strategies",
        "monitoring","indicator","indicators","metrics","measurement","review","systematic review","scoping review"
    }
    _glue_tokens = {g.lower() for g in _glue}

    def _is_glue(term: str) -> bool:
        # term without prefix
        t = re.sub(r"^\s*#?\s*theme\s*:\s*", "", term, flags=re.I).lower()
        # cheap test: if any glue token appears as a word substring
        for g in _glue_tokens:
            if g in t:
                return True
        return False

    def _extract_theme_tags(cell) -> list[str]:
        out: list[str] = []
        if _is_nan(cell):
            return out
        parts: list[str] = []
        if isinstance(cell, list):
            for v in cell:
                if _is_nan(v):
                    continue
                parts.extend(re.split(separators, str(v)))
        else:
            parts = re.split(separators, str(cell))
        for raw in parts:
            raw = str(raw or "").strip()
            if not raw:
                continue
            if re.match(r"^\s*#?\s*theme\s*:", raw, flags=re.I):
                term = _canon_term(raw)
                if not term:
                    continue
                base = term.casefold()
                if stopwords and base in stopwords:
                    continue
                out.append(_prefixed(term))
        return out

    def _tokenize(text: str) -> list[str]:
        t = text.lower()
        t = re.sub(r"[^a-z0-9\-\s/]+", " ", t)
        t = re.sub(r"\s+", " ", t).strip()
        toks = [w for w in t.split(" ") if len(w) >= 3]
        return toks

    # --------------------- snapshot & normalization ---------------------
    if theme_col not in df.columns:
        raise ValueError(f"DataFrame lacks column '{theme_col}'")

    per_item_tags_raw: list[list[str]] = []
    for _, row in df.iterrows():
        per_item_tags_raw.append(_extract_theme_tags(row.get(theme_col)))

    per_item_sets_raw = [sorted(set(L)) for L in per_item_tags_raw]
    N_items = len(per_item_sets_raw)

    # ----------------------- graph construction ------------------------
    _EMB_SIM_MIN = 0.62
    _EMB_ALPHA = 0.20          # per your α≈0.8 on PPMI: so 1-α = 0.2 on cosine
    _KNN_CAP = 12
    _HIGH_DF_RATIO = 0.35      # drop universal tags (appear in >35% of items)

    def _build_graph(min_tf: int, w_min: float, p95_drop_half: bool = True):
        tag_df: Counter = Counter()
        for s in per_item_sets_raw:
            tag_df.update(s)

        # frequency + high-df pruning
        kept = set()
        for t, c in tag_df.items():
            if c >= max(1, min_tf):
                if N_items == 0 or (c / N_items) <= _HIGH_DF_RATIO:
                    kept.add(t)

        # drop glue stopwords before graphing (quick win)
        kept2 = set()
        for t in kept:
            if not _is_glue(t):
                kept2.add(t)
        kept = kept2 if kept2 else kept

        per_item_sets = [[t for t in s if t in kept] for s in per_item_sets_raw]
        tag_df = Counter()
        for s in per_item_sets:
            tag_df.update(s)

        if not tag_df:
            return tag_df, None, {"edges": 0, "semantic": False}

        co_count = defaultdict(int)
        for s in per_item_sets:
            n = len(s)
            for i in range(n):
                for j in range(i + 1, n):
                    a, b = s[i], s[j]
                    if a > b:
                        a, b = b, a
                    co_count[(a, b)] += 1

        try:
            import networkx as nx
        except Exception:
            return tag_df, None, {"edges": 0, "semantic": False}

        G = nx.Graph()
        for t, dfreq in tag_df.items():
            G.add_node(t, weight=int(dfreq))

        def _struct_weight(a: str, b: str, co_ab: int) -> float:
            da, db = tag_df[a], tag_df[b]
            if weight_mode == "count":
                return float(co_ab)
            elif weight_mode == "jaccard":
                denom = (da + db - co_ab)
                return co_ab / denom if denom > 0 else 0.0
            elif weight_mode == "ppmi":
                if N_items <= 0 or da == 0 or db == 0 or co_ab == 0:
                    return 0.0
                p_ab = co_ab / N_items
                p_a = da / N_items
                p_b = db / N_items
                pmi = math.log2(max(p_ab / (p_a * p_b), 1e-12))
                return max(pmi, 0.0)
            # assoc (cosine-like)
            denom = (da * db) ** 0.5
            return co_ab / denom if denom > 0 else 0.0

        edge_added = 0
        for (a, b), co_ab in co_count.items():
            w = _struct_weight(a, b, co_ab)
            if w >= (w_min or 0.0):
                G.add_edge(a, b, w_struct=float(w), weight=float(w))
                edge_added += 1

        info = {"edges": edge_added, "semantic": False}

        # optional semantic augmentation (MiniLM)
        try:
            from sentence_transformers import SentenceTransformer
            import numpy as np
            tags_sorted = sorted(tag_df.keys())
            raw_terms = [re.sub(r"^\s*#?theme\s*:\s*", "", t, flags=re.I) for t in tags_sorted]
            model = SentenceTransformer("all-MiniLM-L6-v2")
            emb = model.encode(raw_terms, normalize_embeddings=True, show_progress_bar=False)
            emb = np.asarray(emb, dtype="float32")
            for i, ti in enumerate(tags_sorted):
                vec = emb[i]
                sims = emb @ vec
                sims[i] = -1.0
                idx = np.argpartition(-sims, min(_KNN_CAP, len(sims)-1))[:_KNN_CAP]
                for j in idx:
                    sj = float(sims[j])
                    if sj < _EMB_SIM_MIN:
                        continue
                    tj = tags_sorted[j]
                    a, b = (ti, tj) if ti < tj else (tj, ti)
                    # normalize structural weight for mixing
                    if G.has_edge(a, b):
                        ws = float(G[a][b].get("w_struct", G[a][b].get("weight", 0.0)))
                    else:
                        ws = 0.0
                    if weight_mode == "ppmi":
                        ws_norm = 1.0 - math.exp(-ws)          # smooth 0..∞ -> 0..1
                    else:
                        ws_norm = max(0.0, min(1.0, ws))       # assoc/jaccard already 0..1
                    w_total = (1.0 - _EMB_ALPHA) * ws_norm + _EMB_ALPHA * sj
                    # keep numeric "weight" in 0..1 scale post-mix
                    if w_total >= w_min:
                        if G.has_edge(a, b):
                            if w_total > G[a][b]["weight"]:
                                G[a][b]["weight"] = float(w_total)
                                G[a][b]["w_sem"] = float(sj)
                        else:
                            G.add_edge(a, b, w_struct=ws, w_sem=float(sj), weight=float(w_total))
            info["semantic"] = True
        except Exception:
            pass

        # degree P95 de-hubbing: drop weakest half of edges for ultra-hubs
        if p95_drop_half and G.number_of_nodes() > 0 and G.number_of_edges() > 0:
            try:
                import numpy as _np
                degw = {n: sum(d.get("weight", 0.0) for _, _, d in G.edges(n, data=True)) for n in G.nodes()}
                vals = _np.array(list(degw.values()), dtype="float32")
                if len(vals) >= 10:
                    p95 = float(_np.percentile(vals, 95))
                    for n, s in degw.items():
                        if s > p95:
                            neigh = []
                            for u, v, d in G.edges(n, data=True):
                                m = v if u == n else u
                                neigh.append((float(d.get("weight", 0.0)), n, m))
                            neigh.sort(key=lambda x: x[0])  # weakest first
                            cut = len(neigh) // 2
                            for w, u, v in neigh[:cut]:
                                # remove only if below current w_min threshold or among the weakest half regardless?
                                # we follow "skip weakest half" regardless, but never remove strong ties above 2*w_min
                                if w < max(w_min * 2.0, 0.75):  # keep very strong ties
                                    if G.has_edge(u, v):
                                        G.remove_edge(u, v)
            except Exception:
                pass

        # final prune edges that fell below threshold after de-hub
        for u, v, d in list(G.edges(data=True)):
            if float(d.get("weight", 0.0)) < w_min:
                G.remove_edge(u, v)

        return tag_df, G, info

    # ----------------------- partitioning & scoring ---------------------
    def _coherence_score(parts, tag_df, G):
        # c-TF-IDF over tag strings, weighted by doc freq
        if not parts:
            return 0.0
        from collections import Counter as _C
        group_tf = []
        vocab_df = _C()
        def _term_tokens(t):
            term = re.sub(r"^\s*#?theme\s*:\s*", "", t, flags=re.I)
            return _tokenize(term)
        for S in parts:
            tf = _C()
            for t in S:
                w = max(1, int(tag_df.get(t, 1)))
                for tok in _term_tokens(t):
                    tf[tok] += w
            group_tf.append(tf)
            for tok in tf.keys():
                vocab_df[tok] += 1
        import math as _m
        n = max(1, len(group_tf))
        scores = []
        for tf in group_tf:
            total = sum(tf.values()) or 1
            cand = []
            for tok, f in tf.items():
                tf_n = f / total
                idf = _m.log(n / (1e-9 + vocab_df[tok]))
                cand.append(tf_n * idf)
            cand.sort(reverse=True)
            scores.append(sum(cand[:8]) / 8.0 if cand else 0.0)
        return sum(scores) / len(scores)

    def _density_score(parts, G):
        if G is None or G.number_of_edges() == 0:
            return 0.0
        intra_w = 0.0
        edge_cnt = 0
        for S in parts:
            nodes = list(S)
            for i in range(len(nodes)):
                for j in range(i + 1, len(nodes)):
                    a, b = nodes[i], nodes[j]
                    if G.has_edge(a, b):
                        intra_w += float(G[a][b].get("weight", 0.0))
                        edge_cnt += 1
        return intra_w / max(1, edge_cnt)

    def _entropy_score(parts):
        sizes = [len(S) for S in parts if len(S) > 0]
        tot = sum(sizes) or 1
        import math as _m
        probs = [s / tot for s in sizes]
        return -sum(p * _m.log(p + 1e-12) for p in probs)

    def _quality(parts, tag_df, G):
        # combined quality proxy
        return (
            0.55 * _coherence_score(parts, tag_df, G) +
            0.35 * _density_score(parts, G) +
            0.10 * _entropy_score(parts) -
            0.25 * abs(len(parts) - target_groups)
        )

    def _partition_with_resolution(G, tag_df, seed):
        if G is None or G.number_of_nodes() == 0:
            return [set()], "none", 1.0, 0.0
        resolutions = [0.8, 1.0, 1.2, 1.4]
        best = None

        # Try Leiden first
        try:
            import igraph as ig
            import leidenalg
            nodes = list(G.nodes())
            idx = {n: i for i, n in enumerate(nodes)}
            edges = [(idx[u], idx[v], float(d.get("weight", 1.0))) for u, v, d in G.edges(data=True)]
            g = ig.Graph()
            g.add_vertices(len(nodes))
            if edges:
                g.add_edges([(i, j) for i, j, _ in edges])
                g.es["weight"] = [w for _, _, w in edges]
            else:
                parts = [{n} for n in nodes]
                return parts, "singleton", 1.0, _quality(parts, tag_df, G)

            for res in resolutions:
                part = leidenalg.find_partition(
                    g,
                    leidenalg.RBConfigurationVertexPartition,
                    weights="weight",
                    seed=int(seed),
                    resolution_parameter=float(res),
                )
                parts = [set(nodes[i] for i in comm) for comm in part]
                parts = [S for S in parts if len(S) > 0]
                score = _quality(parts, tag_df, G)
                if best is None or score > best[0]:
                    best = (score, parts, f"leiden(res={res})", res)
            return best[1], best[2], best[3], best[0]
        except Exception:
            pass

        # Louvain sweep
        try:
            import networkx as nx
            from networkx.algorithms.community import louvain_communities
            best = None
            for res in resolutions:
                parts = list(louvain_communities(G, weight="weight", seed=int(seed), resolution=float(res)))
                parts = [set(S) for S in parts if len(S) > 0]
                score = _quality(parts, tag_df, G)
                if best is None or score > best[0]:
                    best = (score, parts, f"louvain(res={res})", res)
            return best[1], best[2], best[3], best[0]
        except Exception:
            pass

        # Greedy fallback
        try:
            from networkx.algorithms.community import greedy_modularity_communities
            parts = list(greedy_modularity_communities(G, weight="weight"))
            parts = [set(S) for S in parts if len(S) > 0]
            return parts, "greedy_modularity", 1.0, _quality(parts, tag_df, G)
        except Exception:
            parts = [{n} for n in G.nodes()]
            return parts, "singleton", 1.0, _quality(parts, tag_df, G)

    # ---------------- merge/split to exactly target_groups --------------
    def _merge_split_to_target(communities, G, tag_df):
        def _comm_weight(S): return sum(tag_df.get(t, 0) for t in S)
        def _between_weight(A, B):
            w = 0.0
            if G is None:
                return w
            for a in A:
                for b in B:
                    if G.has_edge(a, b):
                        w += float(G[a][b].get("weight", 0.0))
            return w

        comms = [set(S) for S in communities]
        while len(comms) > target_groups:
            comms.sort(key=_comm_weight)
            best = (-1.0, None, None)
            for i in range(min(5, len(comms))):
                for j in range(i + 1, len(comms)):
                    bw = _between_weight(comms[i], comms[j])
                    if bw > best[0]:
                        best = (bw, i, j)
            if best[1] is None:
                B = comms.pop()
                A = comms.pop()
                comms.append(A | B)
            else:
                _, i, j = best
                B = comms.pop(j); A = comms.pop(i)
                comms.append(A | B)

        try:
            from networkx.algorithms.community import greedy_modularity_communities as _gmc
        except Exception:
            _gmc = None

        while len(comms) < target_groups:
            comms.sort(key=lambda s: -_comm_weight(s))
            largest = comms.pop(0)
            if len(largest) <= 2:
                comms.extend([{n} for n in largest])
                continue
            if _gmc is not None and G is not None:
                subG = G.subgraph(largest).copy()
                parts = list(_gmc(subG, weight="weight"))
                if len(parts) <= 1:
                    # heuristic split by degree
                    deg_sorted = sorted(list(largest), key=lambda n: subG.degree(n, weight="weight"), reverse=True)
                    mid = max(1, len(deg_sorted) // 2)
                    comms.extend([set(deg_sorted[:mid]), set(deg_sorted[mid:])])
                else:
                    comms.extend([set(p) for p in parts])
            else:
                sorted_tags = sorted(list(largest), key=lambda t: (-tag_df.get(t, 0), t))
                mid = max(1, len(sorted_tags) // 2)
                comms.extend([set(sorted_tags[:mid]), set(sorted_tags[mid:])])

        return comms

    # ----------------------- fine rebalance -----------------------------
    def _rebalance_capacity(communities, G, cap_ratio: float, min_aff: float, margin_max: float, max_moves: int):
        if G is None or G.number_of_nodes() == 0:
            return communities
        total_nodes = sum(len(c) for c in communities)
        if total_nodes == 0:
            return communities
        import math as _m
        cap = max(1, int(_m.ceil((total_nodes / max(1, target_groups)) * cap_ratio)))
        comms = [set(S) for S in communities]
        group_ids = list(range(len(comms)))

        def _aff(n, S):
            if not S or G is None:
                return 0.0
            s = 0.0; c = 0
            for m in S:
                if m == n:
                    continue
                if G.has_edge(n, m):
                    s += float(G[n][m].get("weight", 0.0)); c += 1
            return s / c if c > 0 else 0.0

        moves = 0
        changed = True
        while changed and moves < max_moves:
            changed = False
            oversized = [i for i, S in enumerate(comms) if len(S) > cap]
            if not oversized:
                break
            for i in oversized:
                S = comms[i]
                candidates = []
                for n in list(S):
                    scores = []
                    for j in group_ids:
                        if i == j:
                            continue
                        scores.append((j, _aff(n, comms[j])))
                    scores.sort(key=lambda x: x[1], reverse=True)
                    best_other, best_aff = (scores[0] if scores else (-1, 0.0))
                    here_aff = _aff(n, S - {n})
                    margin = max(0.0, here_aff - best_aff)
                    if best_other != -1 and best_aff >= min_aff and margin <= margin_max:
                        candidates.append((margin, n, best_other))
                candidates.sort(key=lambda x: x[0])
                for _, n, j in candidates:
                    if len(comms[i]) <= cap:
                        break
                    if len(comms[j]) >= cap:
                        continue
                    comms[i].remove(n); comms[j].add(n)
                    moves += 1; changed = True
                    if moves >= max_moves:
                        break
        return comms

    # ----------------------- consensus + balance ------------------------
    best_overall = None
    last_info = {}
    used_tf = int(min_tag_freq)
    used_wmin = float(min_edge_weight)

    seeds = [random_seed + s for s in range(6)]
    edge_jitters = [0.0, 0.05, -0.03, 0.08, -0.06, 0.10]

    for s, ej in zip(seeds, edge_jitters):
        tag_df, G, info = _build_graph(used_tf, max(0.0, used_wmin + ej))
        last_info = info
        communities, algo, res_used, score = _partition_with_resolution(G, tag_df, s)

        # autobalance by merge/split when oversized
        total_nodes = sum(len(c) for c in communities)
        largest = max((len(c) for c in communities), default=0)
        share = (largest / total_nodes) if total_nodes else 1.0
        if auto_balance and share > balance_max_share:
            for _ in range(balance_max_iter):
                communities = _merge_split_to_target(communities, G, tag_df)
                total_nodes = sum(len(c) for c in communities)
                largest = max((len(c) for c in communities), default=0)
                share = (largest / total_nodes) if total_nodes else 1.0
                if share <= balance_max_share:
                    break

        communities = _rebalance_capacity(
            communities=communities, G=G,
            cap_ratio=min(balance_max_share, 0.6),
            min_aff=0.10, margin_max=0.07, max_moves=250,
        )

        # score after balancing with our composite quality
        sc = _quality(communities, tag_df, G)
        if (best_overall is None) or (sc > best_overall[0]):
            best_overall = (sc, communities, algo, res_used, tag_df, G)

    _, communities, algo, res_used, tag_df, G = best_overall

    # ------------------- two-pass: split largest group ------------------
    # If largest group still too big or low coherence, recluster its subgraph harder
    sizes = [len(S) for S in communities]
    total_nodes = sum(sizes) or 1
    largest_idx = max(range(len(communities)), key=lambda i: sizes[i]) if communities else 0
    largest_share = (sizes[largest_idx] / total_nodes) if communities else 1.0

    need_split = (largest_share > balance_max_share)
    # also check coherence of largest
    if not need_split and communities:
        largest_part = [communities[largest_idx]]
        if _coherence_score(largest_part, tag_df, G) < 0.08:
            need_split = True

    if need_split and communities:
        # stricter thresholds for the subgraph
        stricter_w = used_wmin + (0.5 if weight_mode == "ppmi" else 0.15)
        sub_nodes = sorted(list(communities[largest_idx]))
        try:
            import networkx as nx
            subG = (G.subgraph(sub_nodes).copy() if G is not None else None)
            # If G missing (shouldn't), skip two-pass
            if subG is not None and subG.number_of_nodes() > 2:
                # Build a temporary "df" over the subgraph only by filtering tags
                # We reuse the global tag_df, so just partition the existing subgraph
                # Reweight: drop edges below stricter_w
                for u, v, d in list(subG.edges(data=True)):
                    if float(d.get("weight", 0.0)) < stricter_w:
                        subG.remove_edge(u, v)
                # Partition the subgraph
                sub_parts, sub_algo, sub_res, _ = _partition_with_resolution(subG, tag_df, random_seed + 99)
                # Replace largest group with its sub-parts
                new_comms = [communities[i] for i in range(len(communities)) if i != largest_idx] + sub_parts
                # Re-coerce to exactly K with our merge/split
                new_comms = _merge_split_to_target(new_comms, G, tag_df)
                # Keep if better quality
                if _quality(new_comms, tag_df, G) >= _quality(communities, tag_df, G):
                    communities = new_comms
                    algo = f"{algo}+split({sub_algo})"
        except Exception:
            pass

    # ------------------------ final coercion to K -----------------------
    communities = _merge_split_to_target(communities, G, tag_df)

    # deterministic order: by in-group frequency then name
    def _group_key(S: set[str]) -> tuple[int, str]:
        return (-sum(tag_df.get(t, 0) for t in S), ",".join(sorted(S)))

    communities = sorted(communities, key=_group_key)
    if len(communities) > target_groups:
        communities = communities[:target_groups]
    elif len(communities) < target_groups:
        communities.extend([set() for _ in range(target_groups - len(communities))])

    group_ids = [f"G{i+1}" for i in range(target_groups)]

    # ----------------------------- bridges ------------------------------
    bridges: dict[str, list[str]] = {}
    if G is not None:
        for t in G.nodes():
            scores = []
            for idx, S in enumerate(communities):
                if len(S) == 0:
                    scores.append((idx, 0.0)); continue
                wsum = 0.0; cnt = 0
                for s in S:
                    if t == s:
                        continue
                    if G.has_edge(t, s):
                        wsum += float(G[t][s].get("weight", 0.0)); cnt += 1
                scores.append((idx, (wsum / cnt) if cnt > 0 else 0.0))
            scores.sort(key=lambda x: x[1], reverse=True)
            if len(scores) >= 2:
                best_idx, best_w = scores[0]
                second_idx, second_w = scores[1]
                if best_w > 0 and (second_w / (best_w + 1e-9)) >= 0.85 and second_w >= 0.10:
                    bridges[t] = [group_ids[best_idx], group_ids[second_idx]]

    # ------------------------- map + label output -----------------------
    tag_to_group: dict[str, str] = {}
    groups = []
    for gid, S in zip(group_ids, communities):
        for t in S:
            tag_to_group[t] = gid
        groups.append({"id": gid, "tags": sorted(list(S), key=lambda t: (-tag_df.get(t, 0), t))})

    # c-TF-IDF labels
    try:
        from collections import Counter as _C
        group_tf = []
        vocab_df = _C()
        for S in communities:
            tf = _C()
            for t in S:
                w = max(1, int(tag_df.get(t, 1)))
                term = re.sub(r"^\s*#?theme\s*:\s*", "", t, flags=re.I)
                for tok in _tokenize(term):
                    tf[tok] += w
            group_tf.append(tf)
            for tok in tf.keys():
                vocab_df[tok] += 1
        import math as _m
        n_groups = max(1, len(group_tf))
        labels = []
        top_k = 8
        for tf in group_tf:
            total = sum(tf.values()) or 1
            scored = []
            for tok, f in tf.items():
                tf_norm = f / total
                idf = _m.log(n_groups / (1e-9 + vocab_df[tok]))
                scored.append((tok, tf_norm * idf))
            scored.sort(key=lambda x: x[1], reverse=True)
            labels.append(scored[:top_k])
        for g, top_terms in zip(groups, labels):
            g["top_terms"] = top_terms
            g["label"] = ", ".join([w for (w, _) in top_terms[:4]])
    except Exception:
        pass

    # ------------------------------ return ------------------------------
    return {
        "params": {
            "target_groups": target_groups,
            "min_tag_freq": min_tag_freq,
            "used_min_tag_freq": min_tag_freq,
            "weight_mode": weight_mode,
            "min_edge_weight": min_edge_weight,
            "used_min_edge_weight": min_edge_weight,
            "auto_balance": bool(auto_balance),
            "balance_max_share": balance_max_share,
            "balance_max_iter": balance_max_iter,
            "algo": algo,
            "resolution": res_used,
            "used_networkx": bool(G is not None),
            "semantic_edges": bool(last_info.get("semantic", False)),
            "high_df_prune": _HIGH_DF_RATIO,
            "consensus_runs": 6,
            "capacity_ratio": min(balance_max_share, 0.6),
            "rebalance_min_aff": 0.10,
            "rebalance_margin_max": 0.07,
            "rebalance_max_moves": 250,
            "two_pass_split": True,
            "glue_stopwords_enabled": True,
            "p95_dehub": True,
        },
        "tag_counts": dict(tag_df),
        "groups": groups,
        "tag_to_group": tag_to_group,
        "bridges": bridges,
    }


def add_scope_and_shape_section(prs, df, **kwargs):
    """
    Section A: corpus shape/scope (delegates to shape_scope with safe kwargs).
    """
    collection_name   = kwargs.get("collection_name") or kwargs.get("title") or "Collection"
    slide_notes       = bool(kwargs.get("slide_notes", False))
    year_range        = tuple(kwargs.get("year_range", (2010, 2025)))
    top_n_cats        = int(kwargs.get("top_n_cats", 20))
    include_percent   = bool(kwargs.get("include_percentages", True))
    progress_callback = kwargs.get("progress_callback")

    # call your shape_scope without unknown kwargs
    shape_scope(
        prs,
        df,
        collection_name=collection_name,
        slide_notes=slide_notes,
        year_range=year_range,
        top_n_cats=top_n_cats,
        include_percentages=include_percent,
        progress_callback=progress_callback,
    )


def add_authorship_and_institution_section(prs, df, **kwargs):
    """
    (B) Authorship and institutional influence

    Creates SIX slides:
      1) Top authors — Table
      2) Top authors — Figure
      3) Co-authorship — Strongest pairs (Top 20) — Table
      4) Co-authorship network — Figure
      5) Major institutional actors (publication_outlet) — Table
      6) Major institutional actors (publication_outlet) — Figure
    """
    # ---- imports
    import io, itertools, math, re
    import numpy as np
    import pandas as pd
    import plotly.express as px
    import plotly.graph_objects as go
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN

    # ---- knobs
    collection_name = kwargs.get("collection_name") or kwargs.get("title") or "Collection"
    top_n_authors   = int(kwargs.get("top_n_authors", 20))
    slide_notes = bool(kwargs.get("slide_notes", True))

    # ---- slide metrics
    SLIDE_W_IN = float(prs.slide_width)  / float(Inches(1))
    SLIDE_H_IN = float(prs.slide_height) / float(Inches(1))
    OUTER_X, OUTER_Y, TITLE_H = 0.7, 0.7, 0.7
    FONT = "Segoe UI"
    PLOT_W, PLOT_H = 1200, 650

    # ---------- small utils (robust to floats/NaN)
    def sstr(x) -> str:
        return "" if (x is None or (isinstance(x, float) and np.isnan(x))) else str(x)

    # ---------- slide primitives
    def _new_slide(title: str):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tb = s.shapes.add_textbox(Inches(OUTER_X), Inches(0.2), Inches(SLIDE_W_IN - 2*OUTER_X), Inches(TITLE_H))
        tf = tb.text_frame; tf.clear()
        p = tf.paragraphs[0]; p.text = title; p.font.name = FONT; p.font.size = Pt(24); p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        return s

    def _add_notes(slide, text: str):
        slide.notes_slide.notes_text_frame.text = text

    def _table_slide(title: str, frame: pd.DataFrame):
        s = _new_slide(title)
        if frame is None or frame.empty:
            tb = s.shapes.add_textbox(Inches(OUTER_X), Inches(3.0),
                                      Inches(SLIDE_W_IN-2*OUTER_X), Inches(1.0))
            p = tb.text_frame.paragraphs[0]; p.text = "No data available."
            p.font.name = FONT; p.font.size = Pt(16); p.alignment = PP_ALIGN.CENTER
            return s

        # layout box
        max_w = SLIDE_W_IN - 2*OUTER_X
        max_h = SLIDE_H_IN - (TITLE_H + 2*OUTER_Y)
        rows, cols = len(frame)+1, frame.shape[1]

        def _score(x):
            ss = " ".join(sstr(x).split())
            return min(42, max(4, len(ss)))

        col_scores = []
        for j, c in enumerate(frame.columns):
            longest = max([_score(c)] + [_score(frame.iloc[i, j]) for i in range(len(frame))])
            col_scores.append(longest)
        total = sum(col_scores) or 1
        col_w = [(sc/total)*(max_w-0.5) for sc in col_scores]
        col_w = [max(0.95, w) for w in col_w]
        scale = min(1.0, (max_w-0.25)/sum(col_w))
        col_w = [w*scale for w in col_w]

        table_w = sum(col_w)
        est_row_h = min(0.42, max(0.26, max_h/(rows+1)))
        table_h = min(max_h, rows*est_row_h + 0.35)
        left = (SLIDE_W_IN - table_w)/2.0
        top  = TITLE_H + ((SLIDE_H_IN - TITLE_H) - table_h)/2.0

        tbl = s.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(table_w), Inches(table_h)).table
        for j, w in enumerate(col_w):
            tbl.columns[j].width = Emu(Inches(w))

        for j, c in enumerate(frame.columns):
            cell = tbl.cell(0, j); cell.text = str(c)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.font.bold = True; p.font.name = FONT; p.font.size = Pt(12); p.alignment = PP_ALIGN.CENTER

        for i, (_, r) in enumerate(frame.iterrows(), start=1):
            for j, c in enumerate(frame.columns):
                txt = " ".join(sstr(r[c]).split())
                cell = tbl.cell(i, j); tf = cell.text_frame; tf.clear(); tf.word_wrap = True
                p = tf.paragraphs[0]; p.text = txt; p.font.name = FONT; p.font.size = Pt(11); p.alignment = PP_ALIGN.LEFT
        return s

    def _fig_slide(title: str, fig, notes: str | None = None):
        s = _new_slide(title)
        png = fig.to_image(format="png", width=PLOT_W, height=PLOT_H, scale=2)
        ratio = PLOT_W/PLOT_H
        content_w = SLIDE_W_IN - 2*OUTER_X
        content_h = SLIDE_H_IN - (TITLE_H + 2*OUTER_Y)
        disp_w = min(content_w, content_h*ratio)*0.96
        left = (SLIDE_W_IN - disp_w)/2.0
        top  = TITLE_H + (SLIDE_H_IN - TITLE_H - (disp_w/ratio))/2.0
        s.shapes.add_picture(io.BytesIO(png), Inches(left), Inches(top), width=Inches(disp_w))
        if notes:
            _add_notes(s, notes)
        return s

    def _style(fig, title):
        fig.update_layout(template="plotly_white", width=PLOT_W, height=PLOT_H,
                          margin=dict(l=68, r=120, t=80, b=80),
                          title=dict(text=title, x=0.02, xanchor="left", font=dict(size=20, family=FONT)),
                          font=dict(family=FONT, size=13), legend_title_text="")
        return fig

    # ---------- helpers
    def _authors_from_row(row) -> list[str]:
        # structured
        if "authors_list" in row and isinstance(row["authors_list"], (list, tuple)):
            out = []
            for a in row["authors_list"]:
                if isinstance(a, dict):
                    ln, fn = sstr(a.get("lastName")).strip(), sstr(a.get("firstName")).strip()
                    lab = ", ".join([p for p in [ln, fn] if p])
                    if lab: out.append(lab)
                else:
                    ss = sstr(a).strip()
                    if ss: out.append(ss)
            if out: return out
        # flat strings
        for k in ("authors", "creator_summary"):
            s = sstr(row.get(k)).strip()
            if s:
                return [p.strip() for p in re.split(r"[;|,]+", s) if p.strip()]
        return []

    def _modal(series: pd.Series, default=""):
        s = series.dropna().map(sstr).map(str.strip)
        s = s[s != ""]
        if s.empty: return default
        return s.value_counts().index[0]

    # ===================== 1) Top authors (Table + Figure)
    work = df.copy().dropna(how="all")
    work["__AuthorsList"] = work.apply(_authors_from_row, axis=1)
    auth_rows = work.explode("__AuthorsList").rename(columns={"__AuthorsList": "Author"})
    auth_rows["Author"] = auth_rows["Author"].map(sstr).map(str.strip)
    auth_rows = auth_rows[auth_rows["Author"] != ""]

    # affiliation text/country/department (best-effort)
    def _col(*names):
        for n in names:
            if n in work.columns:
                return work[n].map(sstr)
        return pd.Series([""]*len(work), index=work.index)

    work["__AffText"]    = _col("affiliation_institution","institution","affiliations","affiliation","publisher","publication_outlet","publicationTitle")
    work["__AffCountry"] = _col("affiliation_country","affiliation:country","affiliation_country_value","country_affiliation")
    work["__AffDept"]    = _col("affiliation_department","affiliation:department","department")

    auth_rows = auth_rows.join(work[["__AffText","__AffCountry","__AffDept"]], how="left")

    author_counts = auth_rows.groupby("Author").size().rename("Publications").reset_index()
    modal_aff = auth_rows.groupby("Author").agg({
        "__AffText": _modal, "__AffCountry": _modal, "__AffDept": _modal
    }).reset_index().rename(columns={
        "__AffText":"Dominant Affiliation","__AffCountry":"Affiliation:Country","__AffDept":"Affiliation:Department"
    })

    toplist = (
        author_counts.merge(modal_aff, on="Author", how="left")
        .sort_values(["Publications","Author"], ascending=[False, True])
        .head(top_n_authors)
        .loc[:, ["Author","Dominant Affiliation","Affiliation:Country","Affiliation:Department","Publications"]]
    )

    _table_slide("(B) Top authors — Table", toplist)

    fig_auth = px.bar(toplist.sort_values("Publications", ascending=True),
                      x="Publications", y="Author", orientation="h", text="Publications")

    fig_auth.update_traces(texttemplate="%{text}")
    _fig_slide("(B) Top authors — Figure",
               _style(fig_auth, f"{collection_name} · Top authors (by publications)"),
               notes=("Authors parsed from structured 'authors_list' when present; otherwise from 'authors' / "
                      "'creator_summary' (split on commas/semicolons). Publications counted per record after "
                      "exploding authors. Dominant affiliation & country/department are modal strings across that "
                      "author’s records."))

    # ===================== 2) Co-authorship (Table + Figure)
    edge_counts = {}
    for _, r in work.iterrows():
        A = _authors_from_row(r)
        A = [a for a in A if a]
        for a, b in itertools.combinations(sorted(set(A)), 2):
            edge_counts[(a, b)] = edge_counts.get((a, b), 0) + 1

    if edge_counts:
        edges = pd.DataFrame([(a, b, w) for (a, b), w in edge_counts.items()],
                             columns=["Author A","Author B","Shared Publications"])
        edge_top = edges.sort_values("Shared Publications", ascending=False).head(20).reset_index(drop=True)
    else:
        edge_top = pd.DataFrame(columns=["Author A","Author B","Shared Publications"])

    _table_slide("(B) Co-authorship — Strongest pairs (Top 20) — Table",
                 edge_top if not edge_top.empty else pd.DataFrame({"Message":["No co-authorship pairs found."]}))

    if not edge_counts:
        fig_net = go.Figure().add_annotation(text="No co-authorship pairs found.", showarrow=False)
        _fig_slide("(B) Co-authorship network — Figure",
                   _style(fig_net, f"{collection_name} · Co-authorship network"),
                   notes="No shared-publication pairs were detected; network not constructed.")
    else:
        keep = set(edge_top["Author A"]).union(edge_top["Author B"])
        e2 = edges[(edges["Author A"].isin(keep)) & (edges["Author B"].isin(keep))].copy()
        nodes = sorted(set(e2["Author A"]).union(e2["Author B"]))
        n = len(nodes) if nodes else 1
        theta = np.linspace(0, 2*math.pi, n, endpoint=False)
        pos = {nodes[i]: (math.cos(theta[i]), math.sin(theta[i])) for i in range(n)}

        x_edges, y_edges = [], []
        for _, r in e2.iterrows():
            x0, y0 = pos[r["Author A"]]; x1, y1 = pos[r["Author B"]]
            x_edges += [x0, x1, None]; y_edges += [y0, y1, None]

        deg = {k: 0 for k in nodes}
        for _, r in e2.iterrows():
            deg[r["Author A"]] += 1; deg[r["Author B"]] += 1
        sizes = [8 + 6*math.sqrt(max(1, deg[k])) for k in nodes]

        fig_net = go.Figure()
        fig_net.add_trace(go.Scatter(x=x_edges, y=y_edges, mode="lines",
                                     line=dict(width=1, color="rgba(0,0,0,0.25)"),
                                     hoverinfo="skip", showlegend=False))
        fig_net.add_trace(go.Scatter(x=[pos[k][0] for k in nodes], y=[pos[k][1] for k in nodes],
                                     mode="markers+text", text=nodes, textposition="top center",
                                     marker=dict(size=sizes, line=dict(width=0.5, color="rgba(0,0,0,0.3)")),
                                     hovertemplate="%{text}<extra></extra>", showlegend=False))
        fig_net.update_xaxes(showgrid=False, zeroline=False, visible=False)
        fig_net.update_yaxes(showgrid=False, zeroline=False, visible=False, scaleanchor="x", scaleratio=1)
        _fig_slide("(B) Co-authorship network — Figure",
                   _style(fig_net, f"{collection_name} · Co-authorship network (strongest pairs)"),
                   notes=("Network nodes are authors; an edge exists when two authors co-appear on a record. "
                          "Layout is circular; node size scales with degree among the strongest 20 pairs."))

    # ===================== 3) Major institutional actors (publication_outlet) — Table + Figure =====================
    # Build directly from raw `publication_outlet`; color by modal `publisher_type_value` (no name mapping).

    # 1) locate columns
    outlet_col = next((c for c in ("publication_outlet", "publicationTitle", "publisher") if c in df.columns), None)
    if outlet_col is None:
        raise KeyError("Missing required column 'publication_outlet' (or fallback 'publicationTitle' / 'publisher').")

    type_col = "publisher_type_value" if "publisher_type_value" in df.columns else None
    if type_col is None:
        df["_tmp_pubtype_"] = "Unknown/Other"
        type_col = "_tmp_pubtype_"

    # 2) normalize + explode (properly)
    def _split_outlets(s: str) -> list[str]:
        parts = re.split(r"[;|,/]\s*", str(s or ""))
        return [" ".join(p.strip().split()) for p in parts if p and p.strip()]

    tmp = (
        df[[outlet_col, type_col]].copy()
        .rename(columns={outlet_col: "publication_outlet", type_col: "publisher_type_value"})
    )

    # normalize strings
    tmp["publication_outlet"] = tmp["publication_outlet"].astype(str).map(lambda x: " ".join(x.split()))
    tmp["publisher_type_value"] = tmp["publisher_type_value"].astype(str).replace(
        {"": "Unknown/Other", "nan": "Unknown/Other"}
    )

    # split into list, then explode by column name (correct usage)
    tmp["publication_outlet_list"] = tmp["publication_outlet"].apply(_split_outlets)
    tmp = tmp.explode("publication_outlet_list", ignore_index=True)
    tmp["publication_outlet"] = tmp["publication_outlet_list"].fillna("").astype(str).str.strip()
    tmp = tmp.drop(columns=["publication_outlet_list"])
    tmp = tmp[tmp["publication_outlet"] != ""].copy()

    if tmp.empty:
        table_df = pd.DataFrame({"Publication Outlet": ["—"], "Publisher Type": ["—"], "Count": [0]})
        _table_slide("(B) Major institutional actors (publication_outlet) — Table", table_df)

        fig_inst = go.Figure().add_annotation(text="No publication_outlet values found.", showarrow=False)
        fig_inst = _style(fig_inst, f"{collection_name} · Major institutional actors (publication_outlet)")
        slide_inst_fig = _fig_slide("(B) Major institutional actors (publication_outlet) — Figure", fig_inst)
        if slide_notes:
            _add_notes(slide_inst_fig, "Method: No usable 'publication_outlet' entries were found in the dataset.")
    else:
        # 3) counts per outlet + modal publisher type (for bar color)
        counts = tmp.groupby("publication_outlet").size().rename("Count")
        modal_type = (
            tmp.groupby("publication_outlet")["publisher_type_value"]
            .agg(lambda s: s.value_counts().index[0])
            .rename("Publisher Type")
        )
        agg = (
            pd.concat([counts, modal_type], axis=1)
            .sort_values("Count", ascending=False)
            .head(20)
            .reset_index()
        )

        # 4) table
        table_df = agg.rename(columns={"publication_outlet": "Publication Outlet"})[
            ["Publication Outlet", "Publisher Type", "Count"]
        ]
        _table_slide("(B) Major institutional actors (publication_outlet) — Table", table_df)

        # 5) figure
        fig_inst = px.bar(
            agg.sort_values("Count", ascending=False),
            x="publication_outlet",
            y="Count",
            color="Publisher Type",
            text="Count",
            title=f"{collection_name} · Major institutional actors (publication_outlet)",
        )
        fig_inst.update_traces(textposition="outside",
                               marker_line_width=0.5, marker_line_color="rgba(0,0,0,0.15)")
        fig_inst.update_xaxes(tickangle=-25, automargin=True)
        fig_inst = _style(fig_inst, f"{collection_name} · Major institutional actors (publication_outlet)")

        slide_inst_fig = _fig_slide("(B) Major institutional actors (publication_outlet) — Figure", fig_inst)

        # 6) notes
        if slide_notes:
            n_total = int(tmp.shape[0])
            n_used = int(agg["Count"].sum())
            _add_notes(
                slide_inst_fig,
                "Rationale: Bars show the Top-20 raw values from 'publication_outlet' (no canonical mapping). "
                "Each bar’s color is the modal 'publisher_type_value' for that outlet; blanks labeled 'Unknown/Other'. "
                "Cells with multiple outlets were split on ';', ',', '/', or '|', whitespace normalized, then counted. "
                f"Count basis: {n_used} outlet mentions within {n_total} parsed outlet entries."
            )


def add_thematic_and_method_section(prs, df, **kwargs):
    """
    Section (C): Thematic & methodological mapping
    """
    import io, re
    import numpy as np, pandas as pd
    import plotly.express as px
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN

    # ---------- centered layout helpers (use REAL slide size) ----------
    FONT = "Segoe UI"
    SLIDE_W = float(prs.slide_width) / Inches(1)   # in inches
    SLIDE_H = float(prs.slide_height) / Inches(1)  # in inches
    OUTER_X, OUTER_Y = 0.7, 0.7
    TITLE_H = 0.7

    def _content_box():
        left = OUTER_X
        top = TITLE_H + OUTER_Y
        width = SLIDE_W - 2 * OUTER_X
        height = SLIDE_H - TITLE_H - 2 * OUTER_Y
        return left, top, width, height

    def _new_slide(title: str):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        title_w = SLIDE_W - 2 * OUTER_X
        tb = s.shapes.add_textbox(Inches(OUTER_X), Inches(0.2), Inches(title_w), Inches(TITLE_H))
        tf = tb.text_frame; tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        return s

    def _len_score(v) -> int:
        return min(28, max(4, len(str(v if v is not None else "").strip())))

    def _auto_table(slide, df_small: pd.DataFrame):
        """Create a table centered within the content box; column widths match text."""
        if df_small is None or df_small.empty:
            L, T, W, H = _content_box()
            tb = slide.shapes.add_textbox(Inches(L), Inches(T + H/2 - 0.5), Inches(W), Inches(1.0))
            p = tb.text_frame.paragraphs[0]
            p.text = "No data available."
            p.font.name = FONT; p.font.size = Pt(16); p.alignment = PP_ALIGN.CENTER
            return

        df_small = df_small.copy()
        for c in df_small.columns:
            if df_small[c].dtype == object:
                df_small[c] = df_small[c].astype(str).map(lambda x: " ".join(x.split()))

        L, T, W, H = _content_box()
        rows, cols = len(df_small) + 1, df_small.shape[1]

        scores = []
        for j, col in enumerate(df_small.columns):
            longest = max([_len_score(col)] + [_len_score(df_small.iloc[i, j]) for i in range(len(df_small))])
            scores.append(longest)
        total = sum(scores) or 1

        col_w = [min(3.0, max(1.0, (s / total) * (W - 0.5))) for s in scores]
        scale = min(1.0, (W - 0.25) / sum(col_w))
        col_w = [w * scale for w in col_w]
        table_w = sum(col_w)

        est_row_h = min(0.42, max(0.23, H / (rows + 1)))
        table_h = min(H, rows * est_row_h + 0.4)

        # CENTER inside content box
        left = L + (W - table_w) / 2.0
        top  = T + (H - table_h) / 2.0

        tbl = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                     Inches(table_w), Inches(table_h)).table

        # IMPORTANT: set widths with Inches(...) (no Emu())
        for j, w in enumerate(col_w):
            tbl.columns[j].width = Inches(w)

        # header
        for j, col in enumerate(df_small.columns):
            cell = tbl.cell(0, j)
            cell.text = str(col)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.font.bold = True; p.font.name = FONT; p.font.size = Pt(10.5)
            p.alignment = PP_ALIGN.CENTER
            cell.margin_left = cell.margin_right = Emu(900)
            cell.margin_top = cell.margin_bottom = Emu(420)

        # body
        for i, (_, r) in enumerate(df_small.iterrows(), start=1):
            for j, col in enumerate(df_small.columns):
                cell = tbl.cell(i, j)
                txt = "" if pd.isna(r[col]) else str(r[col])
                tf = cell.text_frame; tf.clear(); tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = txt; p.font.name = FONT; p.font.size = Pt(10)
                p.alignment = PP_ALIGN.LEFT
                cell.margin_left = cell.margin_right = Emu(820)
                cell.margin_top = cell.margin_bottom = Emu(400)

    def _fig_to_slide(slide, fig, wpx=1300, hpx=720):
        """Export Plotly fig and center it inside the content box (keeps aspect)."""
        png = fig.to_image(format="png", width=wpx, height=hpx, scale=2)
        L, T, W, H = _content_box()
        ratio = wpx / hpx
        fit_w = min(W, H * ratio) * 0.96
        fit_h = fit_w / ratio
        left = L + (W - fit_w) / 2.0
        top  = T + (H - fit_h) / 2.0
        slide.shapes.add_picture(io.BytesIO(png), Inches(left), Inches(top), width=Inches(fit_w))


    # ---------- label helpers ----------
    def _nice_label(s: str) -> str:
        if s is None: return ""
        s = str(s).strip()
        s = re.sub(r"^sector[_:\s-]+", "", s, flags=re.I)
        s = s.replace("_", " ")
        s = re.sub(r"\s+", " ", s).strip().title()
        s = re.sub(r"\bEu\b", "EU", s)
        return s

    SHORT_MAP = {
        "Electricity Gas Supply": "Energy",
        "Information And Communication": "ICT",
        "Public Administration And Defence": "Public Admin",
        "Transportation And Storage": "Transport",
        "Professional Scientific And Technical Activities": "Professional/Tech",
        "Water Supply Sewerage Waste Management And Remediation Activities": "Water & Waste",
        "Administrative And Support Service Activities": "Admin & Support",
        "Human Health And Social Work Activities": "Health & Social",
        "Arts Entertainment And Recreation": "Arts & Recreation",
        "Financial And Insurance Activities": "Financial",
        "Wholesale And Retail Trade": "Wholesale/Retail",
        "Accommodation And Food Service Activities": "Hospitality",
        "Education": "Education",
        "Manufacturing": "Manufacturing",
        "Agriculture Forestry And Fishing": "Agriculture",
        "Household": "Household",
    }

    def _shorten(s: str) -> str:
        t = _nice_label(s)
        return SHORT_MAP.get(t, t)

    # ---------- derive fields ----------
    work = df.copy()

    # focus & empirical/theoretical
    if "focus_type_value" not in work.columns:
        work["focus_type_value"] = work.get("focus_type", "")
    et_raw = work.get("focus_type_value", "").astype(str).str.lower()
    work["empirical_theoretical"] = et_raw.map(
        lambda x: "Theoretical" if re.search(r"propositional|concept|framework|theor", x or "")
        else ("Empirical" if (x or "").strip() else "Uncoded")
    )

    # publisher type (robust)
    def _publisher_type_from_row(row) -> str:
        txt = ""
        for c in ("publication_outlet", "publicationTitle", "source", "publisher", "institution"):
            if c in row and isinstance(row[c], str) and row[c].strip():
                txt = row[c].lower()
                break
        txt = (txt or "").strip()
        if not txt:
            return "Other/Unknown"
        if re.search(r"journal|press|springer|elsevier|wiley|taylor|sage|ieee|acm|oxford|cambridge", txt):
            return "Academic"
        if re.search(r"ministry|government|gov|commission|council|parliament|oecd|world bank|imf|un\b|european commission|eu commission|nato", txt):
            return "Policy/Government"
        if re.search(r"rand|brookings|chatham|carnegie|csis|rusi|atlantic council|ifri|sipri|nupi|foundation|institute", txt):
            return "Think Tank"
        return "Academic"
    work["publisher_type"] = work.apply(_publisher_type_from_row, axis=1)

    # sector focus explode
    sector_col = "sector_focus_value" if "sector_focus_value" in work.columns else ("sector_focus" if "sector_focus" in work.columns else None)
    def _explode_series(s: pd.Series) -> pd.Series:
        if s is None or s.empty: return pd.Series([], dtype=str)
        out = s.fillna("").astype(str)
        out = out.str.replace(r"\s*[,;]\s*", "§", regex=True).str.split("§").explode().str.strip()
        return out[out != ""]
    sector_exp = _explode_series(work[sector_col]) if sector_col else pd.Series([], dtype=str)
    sector_exp = sector_exp.map(_shorten)

    # phase focus explode
    phase_col = "phase_focus_value" if "phase_focus_value" in work.columns else ("phase_focus" if "phase_focus" in work.columns else None)
    phase_exp = _explode_series(work[phase_col]).map(_nice_label) if phase_col else pd.Series([], dtype=str)

    # year
    year_col = kwargs.get("year_col") or ("publication_year" if "publication_year" in work.columns else ("year_numeric" if "year_numeric" in work.columns else "year"))
    years = pd.to_numeric(work.get(year_col, np.nan), errors="coerce").astype("Int64")

    # ---------- (1) focus_type_value × publisher_type ----------
    f_focus = work["focus_type_value"].astype(str).map(lambda x: " ".join(x.split()))
    f_pub   = work["publisher_type"].astype(str)
    ct1 = pd.crosstab(f_focus.replace("", pd.NA).dropna(), f_pub.replace("", pd.NA).dropna())

    s = _new_slide("(C) Cross-tab: focus_type_value × publisher_type — Table")
    _auto_table(s, ct1.reset_index() if not ct1.empty else pd.DataFrame({"Message": ["No data to cross-tabulate."]}))

    s = _new_slide("(C) Cross-tab: focus_type_value × publisher_type — Figure")
    if ct1.empty:
        _auto_table(s, pd.DataFrame({"Message": ["No data to plot."]}))
    else:
        long1 = ct1.stack().reset_index(name="Count")
        long1.columns = ["focus_type_value", "publisher_type", "Count"]
        fig1 = px.bar(long1, x="focus_type_value", y="Count", color="publisher_type",
                      barmode="group", title="Focus type × Publisher type")
        fig1.update_layout(margin=dict(l=60, r=60, t=60, b=120))
        fig1.update_xaxes(tickangle=-35, automargin=True)
        _fig_to_slide(s, fig1)

    # ---------- (2) empirical_theoretical × sector_focus_value ----------
    s = _new_slide("(C) Cross-tab: empirical_theoretical × sector_focus_value — Table")
    if sector_exp.empty:
        _auto_table(s, pd.DataFrame({"Message": ["No sector tags to cross-tabulate."]}))
    else:
        tmp = work.loc[:, ["empirical_theoretical"]].copy()
        tmp["__i"] = np.arange(len(tmp))
        sec = work.loc[:, [sector_col]].copy() if sector_col else pd.DataFrame({sector_col: []})
        sec["__i"] = np.arange(len(sec))
        tmp = tmp.merge(sec, on="__i", how="left").drop(columns="__i")
        tmp = tmp.rename(columns={sector_col: "sector_raw"})
        tmp["sector_raw"] = tmp["sector_raw"].fillna("").astype(str)
        tmp["sector_raw"] = tmp["sector_raw"].str.replace(r"\s*[,;]\s*", "§", regex=True)
        tmp = tmp.assign(sector_raw=tmp["sector_raw"].str.split("§")).explode("sector_raw")
        tmp["sector"] = tmp["sector_raw"].map(_shorten)
        tmp = tmp.replace({"sector": {"": np.nan}}).dropna(subset=["sector"])

        ct2 = pd.crosstab(tmp["empirical_theoretical"], tmp["sector"]).sort_index(axis=1)

        TOPN = 14
        if ct2.shape[1] > TOPN:
            totals = ct2.sum(axis=0).sort_values(ascending=False)
            keep = totals.head(TOPN).index.tolist()
            other = ct2.drop(columns=keep).sum(axis=1)
            ct2 = ct2[keep].copy()
            ct2["Other"] = other

        _auto_table(s, ct2.reset_index())

        s2 = _new_slide("(C) Cross-tab: empirical_theoretical × sector_focus_value — Figure")
        long2 = ct2.stack().reset_index(name="Count")
        long2.columns = ["empirical_theoretical", "sector", "Count"]
        long2 = long2.sort_values(["sector", "empirical_theoretical"])
        fig2 = px.bar(long2, x="Count", y="sector", color="empirical_theoretical",
                      orientation="h", title="Empirical/Theoretical × Sector")
        fig2.update_traces(marker_line_width=0.5, marker_line_color="rgba(0,0,0,0.15)")
        fig2.update_layout(margin=dict(l=140, r=120, t=60, b=60), legend_title_text="")
        _fig_to_slide(s2, fig2)

    # ---------- (3) Mission Phase – Overall Proportions ----------
    s = _new_slide("(C) Mission Phase – Overall Proportions — Table")
    if phase_exp.empty:
        _auto_table(s, pd.DataFrame({"Message": ["No phase tags available."]}))
    else:
        ph_counts = phase_exp.value_counts().rename_axis("Phase").reset_index(name="Count")
        ph_counts = ph_counts.sort_values("Count", ascending=False).reset_index(drop=True)
        _auto_table(s, ph_counts)

        s2 = _new_slide("(C) Mission Phase – Overall Proportions — Figure")
        figp = px.bar(ph_counts, x="Phase", y="Count", text="Count", title="Mission Phase – Overall")
        figp.update_traces(textposition="outside", marker_line_width=0.5, marker_line_color="rgba(0,0,0,0.15)")
        figp.update_layout(margin=dict(l=60, r=140, t=60, b=120))
        figp.update_xaxes(tickangle=-25, automargin=True)
        _fig_to_slide(s2, figp)

    # ---------- (4) Mission Phase – Time Trend ----------
    s = _new_slide("(C) Mission Phase – Time Trend — Table")
    if phase_col is None or years.isna().all():
        _auto_table(s, pd.DataFrame({"Message": ["Missing year or phase data."]}))
    else:
        tmp = pd.DataFrame({"Year": years, "Phase": work.get(phase_col, "")})
        tmp["Phase"] = tmp["Phase"].astype(str)
        tmp = tmp.dropna(subset=["Year"])
        tmp = tmp.assign(Phase=tmp["Phase"].str.replace(r"\s*[,;]\s*", "§", regex=True).str.split("§")).explode("Phase")
        tmp["Phase"] = tmp["Phase"].map(_nice_label)
        tmp = tmp.replace({"Phase": {"": np.nan}}).dropna(subset=["Phase"])
        tmp["Year"] = tmp["Year"].astype("Int64")
        grp = tmp.groupby(["Year", "Phase"]).size().rename("Count").reset_index()
        grp["Year"] = grp["Year"].astype(int)
        grp = grp.sort_values(["Year", "Phase"])

        _auto_table(s, grp.head(40))

        s2 = _new_slide("(C) Mission Phase – Time Trend — Figure")
        figt = px.line(grp, x="Year", y="Count", color="Phase", markers=True,
                       title="Mission Phase – Counts over time")
        figt.update_layout(margin=dict(l=60, r=60, t=60, b=60))
        _fig_to_slide(s2, figt)

def add_geo_and_sector_section(prs, df, **kwargs):
    """
    Section (D): Geographic and sectoral coverage

    - Map literature focus by geographical coverage (ISO codes or country names)
    - Calculate sector coverage percentages (ISIC-style tidy labels)
    - Note dominant sectors and countries
    - Optional: summarize coverage for government's missions if present
    """
    import io, re
    import numpy as np
    import pandas as pd
    import plotly.express as px
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN

    # -------------- layout helpers (use REAL slide size; center content) --------------
    FONT = "Segoe UI"
    SLIDE_W = float(prs.slide_width) / Inches(1)
    SLIDE_H = float(prs.slide_height) / Inches(1)
    OUTER_X, OUTER_Y, TITLE_H = 0.7, 0.7, 0.7

    def _content_box():
        left = OUTER_X
        top = TITLE_H + OUTER_Y
        width = SLIDE_W - 2 * OUTER_X
        height = SLIDE_H - TITLE_H - 2 * OUTER_Y
        return left, top, width, height

    def _new_slide(title: str):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        title_w = SLIDE_W - 2 * OUTER_X
        tb = s.shapes.add_textbox(Inches(OUTER_X), Inches(0.2), Inches(title_w), Inches(TITLE_H))
        tf = tb.text_frame; tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        return s

    def _len_score(v) -> int:
        return min(28, max(4, len(str(v if v is not None else "").strip())))

    def _auto_table(slide, df_small: "pd.DataFrame"):
        """Create a table centered within the content box, sized to text."""
        if df_small is None or df_small.empty:
            L, T, W, H = _content_box()
            tb = slide.shapes.add_textbox(Inches(L), Inches(T + H/2 - 0.5), Inches(W), Inches(1.0))
            p = tb.text_frame.paragraphs[0]
            p.text = "No data available."
            p.font.name = FONT; p.font.size = Pt(16); p.alignment = PP_ALIGN.CENTER
            return

        df_small = df_small.copy()
        for c in df_small.columns:
            if df_small[c].dtype == object:
                df_small[c] = df_small[c].astype(str).map(lambda x: " ".join(x.split()))

        L, T, W, H = _content_box()
        rows, cols = len(df_small) + 1, df_small.shape[1]

        scores = []
        for j, col in enumerate(df_small.columns):
            longest = max([_len_score(col)] + [_len_score(df_small.iloc[i, j]) for i in range(len(df_small))])
            scores.append(longest)
        total = sum(scores) or 1

        col_w = [min(3.2, max(1.0, (s / total) * (W - 0.5))) for s in scores]
        scale = min(1.0, (W - 0.25) / sum(col_w))
        col_w = [w * scale for w in col_w]
        table_w = sum(col_w)

        est_row_h = min(0.42, max(0.23, H / (rows + 1)))
        table_h = min(H, rows * est_row_h + 0.4)

        left = L + (W - table_w) / 2.0
        top  = T + (H - table_h) / 2.0

        tbl = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                     Inches(table_w), Inches(table_h)).table
        # Set with Inches(...) (no double conversion)
        for j, w in enumerate(col_w):
            tbl.columns[j].width = Inches(w)

        # header
        for j, col in enumerate(df_small.columns):
            cell = tbl.cell(0, j)
            cell.text = str(col)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.font.bold = True; p.font.name = FONT; p.font.size = Pt(10.5)
            p.alignment = PP_ALIGN.CENTER
            cell.margin_left = cell.margin_right = Emu(900)
            cell.margin_top = cell.margin_bottom = Emu(420)

        # body
        for i, (_, r) in enumerate(df_small.iterrows(), start=1):
            for j, col in enumerate(df_small.columns):
                cell = tbl.cell(i, j)
                txt = "" if pd.isna(r[col]) else str(r[col])
                tf = cell.text_frame; tf.clear(); tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = txt; p.font.name = FONT; p.font.size = Pt(10)
                p.alignment = PP_ALIGN.LEFT
                cell.margin_left = cell.margin_right = Emu(820)
                cell.margin_top = cell.margin_bottom = Emu(400)

    def _fig_to_slide(slide, fig, wpx=1300, hpx=720):
        """Center the exported figure inside the content box (keeps aspect)."""
        png = fig.to_image(format="png", width=wpx, height=hpx, scale=2)
        L, T, W, H = _content_box()
        ratio = wpx / hpx
        fit_w = min(W, H * ratio) * 0.96
        fit_h = fit_w / ratio
        left = L + (W - fit_w) / 2.0
        top  = T + (H - fit_h) / 2.0
        slide.shapes.add_picture(io.BytesIO(png), Inches(left), Inches(top), width=Inches(fit_w))

    # -------------- parsing utilities --------------
    def _explode_tags(s: pd.Series) -> pd.Series:
        if s is None or s.empty: return pd.Series([], dtype=str)
        out = s.fillna("").astype(str)
        out = out.str.replace(r"\s*[,;]\s*", "§", regex=True)
        out = out.str.split("§").explode().str.strip()
        out = out[out != ""]
        return out

    def _nice_sector(s: str) -> str:
        """snake_case or 'sector_x' -> readable short label (ISIC-style)."""
        if s is None: return ""
        s = str(s)
        s = re.sub(r"^sector[_:\s-]+", "", s, flags=re.I)
        s = s.replace("_", " ")
        s = re.sub(r"\s+", " ", s).strip().title()
        # Shorten long ISIC names
        map_short = {
            "Electricity Gas Supply": "Energy",
            "Information And Communication": "ICT",
            "Public Administration And Defence": "Public Admin",
            "Transportation And Storage": "Transport",
            "Professional Scientific And Technical Activities": "Professional/Tech",
            "Water Supply Sewerage Waste Management And Remediation Activities": "Water & Waste",
            "Administrative And Support Service Activities": "Admin & Support",
            "Human Health And Social Work Activities": "Health & Social",
            "Arts Entertainment And Recreation": "Arts & Recreation",
            "Financial And Insurance Activities": "Financial",
            "Wholesale And Retail Trade": "Wholesale/Retail",
            "Accommodation And Food Service Activities": "Hospitality",
            "Agriculture Forestry And Fishing": "Agriculture",
        }
        return map_short.get(s, s)

    def _detect_locationmode(values: pd.Series) -> str:
        """Return 'ISO-3', 'ISO-2', or 'country names' based on the tags."""
        vals = values.dropna().astype(str).str.strip()
        if vals.empty: return "country names"
        if vals.str.fullmatch(r"[A-Z]{3}").all():
            return "ISO-3"
        if vals.str.fullmatch(r"[A-Z]{2}").all():
            return "ISO-2"
        return "country names"

    # -------------- (1) Geographic coverage --------------
    country_col = kwargs.get("country_tag", "country_focus_value")
    if country_col in df.columns:
        exp = _explode_tags(df[country_col])
        if not exp.empty:
            geo_counts = exp.value_counts().rename_axis("Country").reset_index(name="Count")
            # Figure (map)
            slide_g_map = _new_slide("(D) Geographic Coverage — Map")
            locmode = _detect_locationmode(geo_counts["Country"])
            if locmode.startswith("ISO"):
                loc_field = "Country"
                fig_geo = px.choropleth(
                    geo_counts.rename(columns={"Country": loc_field, "Count": "Value"}),
                    locations=loc_field, locationmode=locmode,
                    color="Value", title="Country focus (by count)"
                )
            else:
                fig_geo = px.choropleth(
                    geo_counts.rename(columns={"Country": "location", "Count": "Value"}),
                    locations="location", locationmode="country names",
                    color="Value", title="Country focus (by count)"
                )
            fig_geo.update_layout(margin=dict(l=40, r=40, t=60, b=40), coloraxis_colorbar_title="Count")
            _fig_to_slide(slide_g_map, fig_geo)

            # Table (Top N)
            TOP_COUNTRIES = int(kwargs.get("top_countries", 30))
            slide_g_tbl = _new_slide("(D) Geographic Coverage — Top Countries (Table)")
            _auto_table(slide_g_tbl, geo_counts.head(TOP_COUNTRIES))

            # Add quick insight to notes
            try:
                top3 = ", ".join([f"{r.Country} ({int(r.Count)})" for _, r in geo_counts.head(3).iterrows()])
                slide_g_map.notes_slide.notes_text_frame.text = f"Top countries: {top3}."
            except Exception:
                pass

    # -------------- (2) Sectoral coverage --------------
    sector_col = kwargs.get("sector_tag", "sector_focus_value")
    if sector_col in df.columns:
        sexp = _explode_tags(df[sector_col]).map(_nice_sector)
        if not sexp.empty:
            sc = (
                sexp.value_counts(normalize=True)
                .mul(100).round(1)
                .rename_axis("Sector").reset_index(name="%")
            )
            sc["Sector"] = sc["Sector"].astype(str).map(lambda x: re.sub(r"\s+", " ", x).strip())

            # Table
            slide_s_tbl = _new_slide("(D) Sectoral Coverage — Percentages (Table)")
            _auto_table(slide_s_tbl, sc)

            # Horizontal % bar with text outside (no cropping)
            slide_s_fig = _new_slide("(D) Sectoral Coverage — Percentages (Figure)")
            sc_plot = sc.sort_values("%", ascending=True)
            sc_plot["Label"] = sc_plot["%"].map(lambda v: f"{v:.1f}%")
            fig_sc = px.bar(sc_plot, x="%", y="Sector", orientation="h", text="Label",
                            title="Sector coverage (%)")
            fig_sc.update_traces(textposition="outside",
                                 marker_line_width=0.5,
                                 marker_line_color="rgba(0,0,0,0.15)")
            # Some headroom for outside labels
            xmax = float(np.nanmax(pd.to_numeric(sc_plot["%"], errors="coerce")))
            if np.isfinite(xmax) and xmax > 0:
                fig_sc.update_xaxes(range=[0, xmax * 1.22])
            fig_sc.update_layout(margin=dict(l=140, r=160, t=60, b=60))
            _fig_to_slide(slide_s_fig, fig_sc)

            # Quick insight in notes
            try:
                topS = ", ".join([f"{r.Sector} ({r['%']:.1f}%)" for _, r in sc.head(3).iterrows()])
                slide_s_fig.notes_slide.notes_text_frame.text = f"Dominant sectors: {topS}."
            except Exception:
                pass

    # -------------- (3) Government missions coverage (optional) --------------
    # Look for a plausible mission column.
    mission_col = None
    for cand in ("mission", "mission_title", "government_mission", "mission_name"):
        if cand in df.columns:
            mission_col = cand; break

    if mission_col:
        mexp = _explode_tags(df[mission_col])
        if not mexp.empty:
            miss = mexp.value_counts().rename_axis("Mission").reset_index(name="Count")
            slide_m_tbl = _new_slide("(D) Government Missions — Coverage (Table)")
            _auto_table(slide_m_tbl, miss)

            slide_m_fig = _new_slide("(D) Government Missions — Coverage (Figure)")
            miss_plot = miss.sort_values("Count", ascending=True)
            miss_plot["Label"] = miss_plot["Count"].astype(int).astype(str)
            fig_m = px.bar(miss_plot, x="Count", y="Mission", orientation="h", text="Label",
                           title="Government missions (counts)")
            fig_m.update_traces(textposition="outside",
                                marker_line_width=0.5,
                                marker_line_color="rgba(0,0,0,0.15)")
            xmax = float(np.nanmax(pd.to_numeric(miss_plot["Count"], errors="coerce")))
            if np.isfinite(xmax) and xmax > 0:
                fig_m.update_xaxes(range=[0, xmax * 1.18])
            fig_m.update_layout(margin=dict(l=160, r=140, t=60, b=60))
            _fig_to_slide(slide_m_fig, fig_m)

            try:
                slide_m_fig.notes_slide.notes_text_frame.text = (
                    f"Top missions: {', '.join(miss.Mission.head(3).tolist())}."
                )
            except Exception:
                pass
def add_citations_and_influence_section(prs, df, **kwargs):
    """
    Section (E): Citation and influence indicators

    Outputs (5 slides):
      1) Citation statistics — Table (rich summary: N, zeros, missing, sum, mean/median, stdev, quantiles, h-index)
      2) Citation distribution (linear scale, right-tail capped at P99) — Figure
      3) Citation distribution (log10(citations+1)) — Figure
      4) Top-N most-cited works — Table (with first affiliation)
      5) Top-N most-cited — Figure (horizontal bars)

    Notes are added to each slide with clear data-handling rationale.
    """
    # ---------------- layout + centering helpers ----------------
    import io, re, math
    import numpy as np
    import pandas as pd
    import plotly.express as px
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    FONT = "Segoe UI"
    SLIDE_W = float(prs.slide_width) / Inches(1)
    SLIDE_H = float(prs.slide_height) / Inches(1)
    OUTER_X, OUTER_Y, TITLE_H = 0.7, 0.7, 0.7

    def _content_box():
        left = OUTER_X
        top = TITLE_H + OUTER_Y
        width = SLIDE_W - 2 * OUTER_X
        height = SLIDE_H - TITLE_H - 2 * OUTER_Y
        return left, top, width, height

    def _new_slide(title: str):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        title_w = SLIDE_W - 2 * OUTER_X
        tb = s.shapes.add_textbox(Inches(OUTER_X), Inches(0.2), Inches(title_w), Inches(TITLE_H))
        tf = tb.text_frame; tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        return s

    def _add_note(slide, text: str):
        try:
            slide.notes_slide.notes_text_frame.text = text
        except Exception:
            pass

    def _len_score(v) -> int:
        return min(28, max(4, len(str(v if v is not None else "").strip())))

    def _auto_table(slide, df_small: "pd.DataFrame"):
        """Center a table in the content area, sized to text."""
        if df_small is None or df_small.empty:
            L, T, W, H = _content_box()
            tb = slide.shapes.add_textbox(Inches(L), Inches(T + H/2 - 0.5), Inches(W), Inches(1.0))
            p = tb.text_frame.paragraphs[0]
            p.text = "No data available."
            p.font.name = FONT; p.font.size = Pt(16); p.alignment = PP_ALIGN.CENTER
            return

        df_small = df_small.copy()
        for c in df_small.columns:
            if df_small[c].dtype == object:
                df_small[c] = df_small[c].astype(str).map(lambda x: " ".join(x.split()))

        L, T, W, H = _content_box()
        rows, cols = len(df_small) + 1, df_small.shape[1]

        scores = []
        for j, col in enumerate(df_small.columns):
            longest = max([_len_score(col)] + [_len_score(df_small.iloc[i, j]) for i in range(len(df_small))])
            scores.append(longest)
        total = sum(scores) or 1

        col_w = [min(3.2, max(1.0, (s / total) * (W - 0.5))) for s in scores]
        scale = min(1.0, (W - 0.25) / sum(col_w))
        col_w = [w * scale for w in col_w]
        table_w = sum(col_w)

        est_row_h = min(0.42, max(0.23, H / (rows + 1)))
        table_h = min(H, rows * est_row_h + 0.4)

        left = L + (W - table_w) / 2.0
        top  = T + (H - table_h) / 2.0

        tbl = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                     Inches(table_w), Inches(table_h)).table
        for j, w in enumerate(col_w):
            tbl.columns[j].width = Inches(w)

        # header
        for j, col in enumerate(df_small.columns):
            cell = tbl.cell(0, j)
            cell.text = str(col)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.font.bold = True; p.font.name = FONT; p.font.size = Pt(10.5)
            p.alignment = PP_ALIGN.CENTER

        # body
        for i, (_, r) in enumerate(df_small.iterrows(), start=1):
            for j, col in enumerate(df_small.columns):
                cell = tbl.cell(i, j)
                txt = "" if pd.isna(r[col]) else str(r[col])
                tf = cell.text_frame; tf.clear(); tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = txt; p.font.name = FONT; p.font.size = Pt(10)
                p.alignment = PP_ALIGN.LEFT

    def _fig_to_slide(slide, fig, wpx=1300, hpx=720):
        """Center the image of a Plotly fig inside the content box."""
        png = fig.to_image(format="png", width=wpx, height=hpx, scale=2)
        L, T, W, H = _content_box()
        ratio = wpx / hpx
        fit_w = min(W, H * ratio) * 0.96
        fit_h = fit_w / ratio
        left = L + (W - fit_w) / 2.0
        top  = T + (H - fit_h) / 2.0
        slide.shapes.add_picture(io.BytesIO(png), Inches(left), Inches(top), width=Inches(fit_w))

    # ---------------- helpers ----------------
    def _find_citations_col(frame: pd.DataFrame) -> str | None:
        preferred = ["citations", "cited_by", "cited_by_count", "citations_count",
                     "times_cited", "n_cit", "num_cit"]
        for p in preferred:
            if p in frame.columns:
                return p
        pats = (r"^citations?$", r"^cited[_ ]?by", r"times[_ ]?cited", r"^n[_ ]?cit", r"^num[_ ]?cit")
        for c in frame.columns:
            if any(re.search(p, str(c), flags=re.I) for p in pats):
                return c
        return None

    def _first_affiliation(row) -> str:
        # Try a few common columns; return the first non-empty string
        for c in (
            "affiliation", "affiliations", "institution", "publisher", "publication_outlet",
            "affiliation:department", "affiliation_department",
            "affiliation:country", "affiliation_country"
        ):
            v = row.get(c, "")
            if isinstance(v, str) and v.strip():
                return " ".join(v.split())
        return ""

    # ---------------- working copy + cleaned citations ----------------
    work = df.copy()
    cit_col = kwargs.get("citations_col") or _find_citations_col(work) or "citations"
    work[cit_col] = (
        pd.to_numeric(work.get(cit_col, np.nan), errors="coerce")
        .fillna(0)  # empty / non-numeric -> 0
        .clip(lower=0)  # no negatives
        .astype(int)  # optional: ints for display
    )
    s_cit = work[cit_col].astype(float)  # use this for all downstream stats/plots
    34
    # Create a numeric citations series; leave original df untouched except where displayed
    s_cit = pd.to_numeric(work.get(cit_col, np.nan), errors="coerce")

    # Basic sets
    n_total = int(len(work))
    n_missing = int(s_cit.isna().sum())
    s_valid = s_cit.dropna()
    n_with = int(len(s_valid))
    # Treat negative as 0 (defensive)
    s_valid = s_valid.clip(lower=0)
    n_zeros = int((s_valid == 0).sum())
    s_pos = s_valid[s_valid > 0]

    # Summary stats (defensive for empty)
    total_cites = float(s_valid.sum()) if n_with else 0.0
    mean = float(s_valid.mean()) if n_with else 0.0
    median = float(s_valid.median()) if n_with else 0.0
    stdev = float(s_valid.std(ddof=1)) if n_with > 1 else 0.0
    p10 = float(s_valid.quantile(0.10)) if n_with else 0.0
    p25 = float(s_valid.quantile(0.25)) if n_with else 0.0
    p75 = float(s_valid.quantile(0.75)) if n_with else 0.0
    p90 = float(s_valid.quantile(0.90)) if n_with else 0.0
    p95 = float(s_valid.quantile(0.95)) if n_with else 0.0
    p99 = float(s_valid.quantile(0.99)) if n_with else 0.0
    cmin = int(s_valid.min()) if n_with else 0
    cmax = int(s_valid.max()) if n_with else 0

    # h-index (simple definition on current corpus)
    def _h_index(vals: pd.Series) -> int:
        if vals.empty: return 0
        arr = np.sort(np.array(vals, dtype=float))[::-1]
        h = 0
        for i, v in enumerate(arr, start=1):
            if v >= i: h = i
            else: break
        return int(h)

    h_index = _h_index(s_valid)

    stats = pd.DataFrame({
        "Metric": [
            "Records (total)", "Records with citations", "Records with 0 citations", "Missing citations",
            "Total citations (sum)", "Mean citations", "Median citations", "Std. deviation",
            "P10", "P25 (Q1)", "P75 (Q3)", "P90", "P95", "P99", "Min", "Max", "h-index"
        ],
        "Value": [
            n_total, n_with, n_zeros, n_missing,
            int(total_cites), round(mean, 2), round(median, 2), round(stdev, 2),
            round(p10, 2), round(p25, 2), round(p75, 2), round(p90, 2), round(p95, 2), round(p99, 2),
            cmin, cmax, h_index
        ]
    })

    # -------- Slide 1: stats table --------
    slide1 = _new_slide("(E) Citation Statistics — Table")
    _auto_table(slide1, stats)
    _add_note(
        slide1,
        "Rationale: Citations parsed from a best-match column (e.g., 'citations', 'cited_by', 'times_cited'); "
        "values converted to numeric with non-numeric treated as missing. Negative values clipped to 0. "
        "We report N totals, zeros, missing, sum, mean/median, stdev, key quantiles, min/max and h-index "
        "(max h such that ≥h items have ≥h citations)."
    )

    # -------- Slide 2: distribution (linear, P99 cap) --------
    slide2 = _new_slide("(E) Citation Distribution — Linear (cap at P99) — Figure")
    if s_valid.empty:
        _auto_table(slide2, pd.DataFrame({"Message": ["No citation data available."]}))
    else:
        # Cap at P99 to prevent a few outliers from flattening the rest
        cap = p99 if np.isfinite(p99) and p99 > 0 else cmax
        s_plot = s_valid.clip(upper=cap)
        # Freedman–Diaconis rule for bins (robust to heavy tails); fallback to 30
        try:
            iqr = float(np.subtract(*np.percentile(s_plot, [75, 25])))
            bw = 2 * iqr * (len(s_plot) ** (-1/3)) if iqr > 0 else None
            bins = max(10, min(80, int(math.ceil((s_plot.max() - s_plot.min()) / bw)))) if bw and bw > 0 else 30
        except Exception:
            bins = 30

        fig_hist = px.histogram(
            s_plot, nbins=bins, title="Citation Distribution (capped at 99th percentile)"
        )
        fig_hist.update_layout(
            margin=dict(l=70, r=70, t=70, b=70),
            xaxis_title="Citations (capped at P99)",
            yaxis_title="Records",
            bargap=0.03
        )
        _fig_to_slide(slide2, fig_hist)
    _add_note(
        slide2,
        "Figure caps values at the 99th percentile so the body of the distribution is readable. "
        "Bins chosen via Freedman–Diaconis rule (fallback to 30)."
    )

    # -------- Slide 3: distribution (log10(citations+1)) --------
    slide3 = _new_slide("(E) Citation Distribution — log10(citations+1) — Figure")
    if s_valid.empty:
        _auto_table(slide3, pd.DataFrame({"Message": ["No citation data available."]}))
    else:
        s_log = np.log10(s_valid + 1.0)
        fig_log = px.histogram(
            s_log, nbins=30, title="Citation Distribution on log10 scale"
        )
        fig_log.update_layout(
            margin=dict(l=70, r=70, t=70, b=70),
            xaxis_title="log10(Citations + 1)",
            yaxis_title="Records",
            bargap=0.03
        )
        _fig_to_slide(slide3, fig_log)
    _add_note(
        slide3,
        "Log transform (log10(citations+1)) reveals structure obscured by heavy tails. "
        "‘+1’ allows zero-citation items to be included."
    )

    # -------- Slide 4: Top-N most-cited — Table --------
    top_n = int(kwargs.get("top_n", 10))
    # Ensure display columns exist
    for c in ("title", "authors"):
        if c not in work.columns:
            work[c] = ""

    work[cit_col] = pd.to_numeric(work.get(cit_col, np.nan), errors="coerce").fillna(0).clip(lower=0).astype(int)

    disp = work.copy()
    disp["Affiliation (first)"] = disp.apply(_first_affiliation, axis=1)

    def _short(s, n=120):
        s = "" if s is None or (isinstance(s, float) and np.isnan(s)) else str(s).strip()
        s = " ".join(s.split())
        return s if len(s) <= n else s[: n - 1].rstrip() + "…"

    table_cols = ["title", "authors", "Affiliation (first)", cit_col]
    for c in table_cols:
        if c not in disp.columns:
            disp[c] = ""

    top_tbl = (
        disp.sort_values(cit_col, ascending=False)
            .loc[:, table_cols]
            .rename(columns={"title": "Title", "authors": "Authors", cit_col: "Citations"})
            .head(top_n)
            .reset_index(drop=True)
    )
    top_tbl["Title"] = top_tbl["Title"].map(lambda x: _short(x, 140))
    top_tbl["Authors"] = top_tbl["Authors"].map(lambda x: _short(x, 110))
    top_tbl["Affiliation (first)"] = top_tbl["Affiliation (first)"].map(lambda x: _short(x, 80))

    slide4 = _new_slide(f"(E) Top {top_n} Most-Cited Works — Table")
    _auto_table(slide4, top_tbl)
    _add_note(
        slide4,
        "Top list sorted by numeric citations (negatives clipped to 0). "
        "‘Affiliation (first)’ is the first non-empty value found across common affiliation/outlet fields. "
        "Long strings are truncated for readability; full text remains in the source data."
    )

    # -------- Slide 5: Top-N most-cited — Figure --------
    if not top_tbl.empty:
        plot_df = top_tbl.copy()
        plot_df.insert(0, "Rank", plot_df.index + 1)
        # label author(s) prominently; title in tooltip via hovertemplate
        plot_df["Label"] = plot_df.apply(lambda r: f"#{int(r['Rank'])} — {r['Authors']}", axis=1)

        fig_top = px.bar(
            plot_df.sort_values("Citations", ascending=True),
            x="Citations", y="Label", orientation="h",
            text="Citations", title=f"Top {top_n} Most-Cited — By Authors"
        )
        fig_top.update_traces(
            textposition="outside",
            marker_line_width=0.5, marker_line_color="rgba(0,0,0,0.15)",
            hovertemplate="<b>%{y}</b><br>Citations: %{x}<extra></extra>"
        )
        try:
            xmax = float(np.nanmax(pd.to_numeric(plot_df["Citations"], errors="coerce")))
            if np.isfinite(xmax) and xmax > 0:
                fig_top.update_xaxes(range=[0, xmax * 1.18])
        except Exception:
            pass
        fig_top.update_layout(margin=dict(l=190, r=140, t=70, b=60))
        _fig_to_slide(_new_slide(f"(E) Top {top_n} Most-Cited — Figure"), fig_top)

def add_theory_empirical_section(prs, df, **kwargs):
    """
    Theoretical – empirical balance

    1) Count papers by Theoretical vs Empirical from 'empirical_theoretical_value'.
    2) For the Empirical subset only, classify method types from 'methodology_value'
       using explicit qualitative / quantitative / mixed term lists (substring match).
       Tokens that don't match are counted as 'Unmapped' (reported in notes, not in the table/figure).
    3) Identify papers that explicitly test/evaluate using ONLY 'evaluates_tests_value'
       (truthy values → flagged). For the flagged subset, show the phase distribution
       using 'phase_focus_value' if available. Also show an alternative phase-based ratio
       that counts items with 'Implementing' or 'Evaluating' in 'phase_focus_value'.

    Every slide is centered and includes methodology notes tied to the data actually used.
    """
    # --------------------------- imports ---------------------------
    import io, re
    import numpy as np
    import pandas as pd
    import plotly.express as px
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    # --------------------------- required columns (fail fast) ---------------------------
    TE_COL    = kwargs.get("te_col",    "empirical_theoretical_value")
    METH_COL  = kwargs.get("meth_col",  "methodology_value")
    EVAL_COL  = kwargs.get("eval_col",  "evaluates_tests_value")
    PHASE_COL = kwargs.get("phase_col", "phase_focus_value")  # optional for phase charts

    missing = [c for c in (TE_COL, METH_COL, EVAL_COL) if c not in df.columns]
    if missing:
        raise KeyError(f"Missing required column(s): {', '.join(missing)}")

    # --------------------------- layout helpers ---------------------------
    FONT = "Segoe UI"
    SLIDE_W = float(prs.slide_width) / Inches(1)
    SLIDE_H = float(prs.slide_height) / Inches(1)
    OUTER_X, OUTER_Y, TITLE_H = 0.7, 0.7, 0.7

    def _content_box():
        left = OUTER_X
        top = TITLE_H + OUTER_Y
        width = SLIDE_W - 2 * OUTER_X
        height = SLIDE_H - TITLE_H - 2 * OUTER_Y
        return left, top, width, height

    def _new_slide(title: str):
        s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        title_w = SLIDE_W - 2 * OUTER_X
        tb = s.shapes.add_textbox(Inches(OUTER_X), Inches(0.2), Inches(title_w), Inches(TITLE_H))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        return s

    def _len_score(v) -> int:
        return min(28, max(4, len(str("" if v is None else v).strip())))

    def _auto_table(slide, df_small: "pd.DataFrame"):
        """Center a table in the content area, sized to text; compact whitespace."""
        if df_small is None or df_small.empty:
            L, T, W, H = _content_box()
            tb = slide.shapes.add_textbox(Inches(L), Inches(T + H / 2 - 0.5), Inches(W), Inches(1.0))
            p = tb.text_frame.paragraphs[0]
            p.text = "No data available."
            p.font.name = FONT
            p.font.size = Pt(16)
            p.alignment = PP_ALIGN.CENTER
            return

        df_small = df_small.copy()
        for c in df_small.columns:
            if df_small[c].dtype == object:
                df_small[c] = df_small[c].astype(str).map(lambda x: " ".join(x.split()))

        L, T, W, H = _content_box()
        rows, cols = len(df_small) + 1, df_small.shape[1]

        scores = []
        for j, col in enumerate(df_small.columns):
            longest = max([_len_score(col)] + [_len_score(df_small.iloc[i, j]) for i in range(len(df_small))])
            scores.append(longest)
        total = sum(scores) or 1

        col_w = [min(3.0, max(1.0, (s / total) * (W - 0.5))) for s in scores]
        scale = min(1.0, (W - 0.25) / sum(col_w))
        col_w = [w * scale for w in col_w]
        table_w = sum(col_w)

        est_row_h = min(0.42, max(0.23, H / (rows + 1)))
        table_h = min(H, rows * est_row_h + 0.4)

        left = L + (W - table_w) / 2.0
        top = T + (H - table_h) / 2.0

        tbl = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                     Inches(table_w), Inches(table_h)).table
        for j, w in enumerate(col_w):
            tbl.columns[j].width = Inches(w)

        # header
        for j, col in enumerate(df_small.columns):
            cell = tbl.cell(0, j)
            cell.text = str(col)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.font.bold = True
            p.font.name = FONT
            p.font.size = Pt(10.5)
            p.alignment = PP_ALIGN.CENTER

        # body
        for i, (_, r) in enumerate(df_small.iterrows(), start=1):
            for j, col in enumerate(df_small.columns):
                cell = tbl.cell(i, j)
                txt = "" if pd.isna(r[col]) else str(r[col])
                tf = cell.text_frame
                tf.clear()
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = txt
                p.font.name = FONT
                p.font.size = Pt(10)
                p.alignment = PP_ALIGN.LEFT

    def _fig_to_slide(slide, fig, wpx=1300, hpx=720):
        """Center a Plotly figure image in the content box."""
        png = fig.to_image(format="png", width=wpx, height=hpx, scale=2)
        L, T, W, H = _content_box()
        ratio = wpx / hpx
        fit_w = min(W, H * ratio) * 0.96
        fit_h = fit_w / ratio
        left = L + (W - fit_w) / 2.0
        top = T + (H - fit_h) / 2.0
        slide.shapes.add_picture(io.BytesIO(png), Inches(left), Inches(top), width=Inches(fit_w))

    def _add_note(slide, text: str):
        slide.notes_slide.notes_text_frame.text = text  # no try/except

    # =============================================================================
    # (1) Theoretical vs Empirical — from 'empirical_theoretical_value'
    # =============================================================================
    te = df[TE_COL].astype(str).str.strip()
    te_counts = (
        te.replace("", np.nan)
          .fillna("Uncoded")
          .value_counts()
          .rename_axis("Type")
          .reset_index(name="Count")
    )

    s1 = _new_slide("Theoretical vs Empirical — Table")
    _auto_table(s1, te_counts)
    _add_note(
        s1,
        "Source: 'empirical_theoretical_value'. Exact value counts; blanks labelled 'Uncoded'. "
        "No heuristics or other columns were used. One label per record."
    )

    s1b = _new_slide("Theoretical vs Empirical — Figure")
    _te_plot = te_counts.sort_values("Count", ascending=False).copy()
    _te_plot["Label"] = _te_plot["Count"].astype(int).astype(str)
    fig1 = px.bar(_te_plot, x="Type", y="Count", text="Label", title="Theoretical–Empirical Balance")
    fig1.update_traces(textposition="outside")
    fig1.update_layout(margin=dict(l=60, r=120, t=60, b=60))
    _fig_to_slide(s1b, fig1)
    _add_note(s1b, "Figure mirrors the table counts from 'empirical_theoretical_value' (sorted by Count).")

    # =============================================================================
    # (2) Methodological type — Empirical subset ONLY; from 'methodology_value'
    # =============================================================================
    emp_mask = te.str.lower() == "empirical"
    emp = df.loc[emp_mask, [METH_COL]].copy()

    # Explicit term lists (substring, case-insensitive)
    QUAL_TERMS = [
        "case study", "interview", "focus group", "ethnograph", "qualitative",
        "process tracing", "document analysis", "content analysis", "discourse",
        "comparative case", "thematic", "fieldwork"
    ]
    QUANT_TERMS = [
        "regression", "econometric", "experiment", "survey", "quantitative",
        "randomized", "randomised", "rct", "statistic", "panel", "time series",
        "difference-in-differences", "diff-in-diff", "did", "propensity",
        "matching", "bayesian", "simulation", "network analysis", "machine learning", "ml",
        "quasi-experiment", "quasi experimental"
    ]
    MIXED_TERMS = ["mixed-method", "mixed method", "mixed methods", "mixed"]

    def _bucket_from_methodology(value: str, unused_counter: dict) -> str:
        """
        Classify a single 'methodology_value' cell into Mixed / Quantitative / Qualitative.
        Priority per record: Mixed > Quantitative > Qualitative.
        If nothing matches, record the raw token(s) to 'unused_counter' and return 'Unmapped'.
        """
        if value is None:
            return "Unmapped"
        text = str(value).lower().replace("-", " ")
        if any(term in text for term in MIXED_TERMS):
            return "Mixed"
        if any(term in text for term in QUANT_TERMS):
            return "Quantitative"
        if any(term in text for term in QUAL_TERMS):
            return "Qualitative"
        # record unmapped fragments for transparency
        parts = re.split(r"\s*[,;|/]\s*|\s{2,}", str(value))
        for p in parts:
            p = p.strip()
            if p:
                unused_counter[p] = unused_counter.get(p, 0) + 1
        return "Unmapped"

    if emp.empty:
        s2 = _new_slide("Methodological Types — Empirical subset (Table)")
        _auto_table(s2, pd.DataFrame({"Message": ["No empirical papers identified."]}))
        _add_note(
            s2,
            "Empirical subset is empty (no rows where 'empirical_theoretical_value' == 'Empirical')."
        )
    else:
        unmapped = {}
        method_by_row = emp[METH_COL].map(lambda v: _bucket_from_methodology(v, unmapped))

        wanted = ["Qualitative", "Quantitative", "Mixed"]
        mcounts = (
            method_by_row[method_by_row.isin(wanted)]
            .value_counts()
            .reindex(wanted, fill_value=0)
            .rename_axis("Method Type")
            .reset_index(name="Count")
        )

        s2 = _new_slide("Methodological Types — Empirical subset (Table)")
        _auto_table(s2, mcounts)

        if unmapped:
            unm_df = pd.Series(unmapped).sort_values(ascending=False).head(20)
            unm_txt = ", ".join([f"{k} ({int(v)})" for k, v in unm_df.items()])
        else:
            unm_txt = "None"

        def _list_terms(lst): return ", ".join(sorted(lst))

        _add_note(
            s2,
            "Empirical subset is defined strictly by 'empirical_theoretical_value' == 'Empirical'. "
            "Method categories come ONLY from 'methodology_value' via substring matching:\n"
            f"• Qualitative terms: {_list_terms(QUAL_TERMS)}\n"
            f"• Quantitative terms: {_list_terms(QUANT_TERMS)}\n"
            f"• Mixed terms: {_list_terms(MIXED_TERMS)}\n"
            "Priority per record: Mixed > Quantitative > Qualitative. "
            "If none of these terms appear, the item is counted as Unmapped (not shown in the table). "
            f"Top Unmapped tokens: {unm_txt}."
        )

        s2b = _new_slide("Methodological Types — Empirical subset (Figure)")
        _m_plot = mcounts.sort_values("Count", ascending=False).copy()
        _m_plot["Label"] = _m_plot["Count"].astype(int).astype(str)
        fig2 = px.bar(_m_plot, x="Method Type", y="Count", text="Label", title="Methods (Empirical papers)")
        fig2.update_traces(textposition="outside")
        fig2.update_layout(margin=dict(l=60, r=120, t=60, b=60))
        _fig_to_slide(s2b, fig2)
        _add_note(s2b, "Bar chart of Qualitative / Quantitative / Mixed derived from 'methodology_value' for empirical papers.")

    # =============================================================================
    # (3) Testing/Evaluating — from 'evaluates_tests_value' (+ alternative phase view)
    # =============================================================================
    eval_series = df[EVAL_COL].astype(str).str.strip().str.lower()
    truthy = {"yes", "true", "y", "1", "t"}
    flag_yes = eval_series.isin(truthy)

    n_yes = int(flag_yes.sum())
    n_no = int((~flag_yes).sum())
    total = max(1, n_yes + n_no)

    eval_counts = pd.DataFrame(
        {"Status": ["Testing/Evaluating: Yes", "Testing/Evaluating: No"],
         "Count": [n_yes, n_no]}
    )
    eval_counts["%"] = (100.0 * eval_counts["Count"] / total).round(1)

    # Table
    s3a = _new_slide("Testing/Evaluating — Summary (from 'evaluates_tests_value') — Table")
    _auto_table(s3a, eval_counts.sort_values("Count", ascending=False))
    _add_note(
        s3a,
        "A record is 'Yes' when evaluates_tests_value ∈ {Yes, True, Y, 1, T} (case-insensitive). "
        "All other values (including blanks) are 'No'. Percentages are out of all records."
    )

    # Figure (ratio) — compute text after sorting to avoid misalignment
    s3a_fig = _new_slide("Testing/Evaluating — Summary (from 'evaluates_tests_value') — Figure")
    _eval_plot = eval_counts.sort_values("Count", ascending=False).copy()
    _eval_plot["Label"] = _eval_plot["%"].map(lambda v: f"{v:.1f}%")
    fig_eval = px.bar(
        _eval_plot,
        x="Status",
        y="Count",
        text="Label",
        title="Share of papers that explicitly test/evaluate missions (Yes/No)"
    )
    fig_eval.update_traces(textposition="outside")
    fig_eval.update_layout(margin=dict(l=60, r=120, t=60, b=60))
    _fig_to_slide(s3a_fig, fig_eval)
    _add_note(
        s3a_fig,
        "Labels show percent of total records; bar heights show absolute counts. Labels are computed from the sorted data to keep them aligned."
    )

    # Examples (Top 12 'Yes'), if cols exist
    example_cols = [c for c in ("title", "authors", PHASE_COL) if c in df.columns]
    if example_cols:
        ex = df.loc[flag_yes, example_cols].copy().head(12)

        def _short(s, n):
            s = "" if s is None else str(s).strip()
            s = " ".join(s.split())
            return s if len(s) <= n else s[: n - 1].rstrip() + "…"

        if "title" in ex.columns:
            ex["title"] = ex["title"].map(lambda x: _short(x, 120))
        if "authors" in ex.columns:
            ex["authors"] = ex["authors"].map(lambda x: _short(x, 90))

        s3a_ex = _new_slide("Testing/Evaluating — Examples (Top 12)")
        _auto_table(
            s3a_ex,
            ex.rename(columns={"title": "Title", "authors": "Authors", PHASE_COL: "Phase (if any)"})
        )
        _add_note(
            s3a_ex,
            "Illustrative subset of records flagged 'Yes' by evaluates_tests_value. Text truncated for readability."
        )

    # Phase distribution among flagged items (if phase column exists)
    if PHASE_COL in df.columns and flag_yes.any():
        phase_series = df.loc[flag_yes, PHASE_COL].astype(str)
        phase_tokens = (
            phase_series.str.replace(r"\s*[,;|/]\s*", "§", regex=True)
                        .str.split("§")
                        .explode()
                        .str.strip()
        )
        phase_tokens = phase_tokens[phase_tokens != ""]

        phase_counts = (
            phase_tokens.value_counts()
            .rename_axis("Phase")
            .reset_index(name="Count")
            .sort_values("Count", ascending=False)
        )

        # # =============================================================================
        # # (3) Testing/Evaluating — Phase × Yes/No split (table + figure)
        # # =============================================================================
        # if PHASE_COL in df.columns:
        #     # Build a per-row Yes/No flag from evaluates_tests_value (already defined earlier: flag_yes)
        #     flag_series = pd.Series(np.where(flag_yes, "Yes", "No"), index=df.index)
        #
        #     # Tokenize phases per row, keep non-empty tokens, and align with the Yes/No flag
        #     phase_raw = df[PHASE_COL].astype(str)
        #     phase_tokens = (
        #         phase_raw.str.replace(r"\s*[,;|/]\s*", "§", regex=True)
        #         .str.split("§")
        #         .explode()
        #         .astype(str)
        #         .str.strip()
        #     )
        #     mask_nonempty = phase_tokens.ne("")
        #     if mask_nonempty.any():
        #         phase_flag = pd.DataFrame({
        #             "Phase": phase_tokens[mask_nonempty].values,
        #             "Flag": flag_series.loc[phase_tokens[mask_nonempty].index].values
        #         })
        #         # Tidy casing for display; aggregate counts
        #         phase_flag["Phase"] = phase_flag["Phase"].str.title()
        #         counts_long = (
        #             phase_flag.groupby(["Phase", "Flag"])
        #             .size().rename("Count")
        #             .reset_index()
        #         )
        #         if counts_long.empty:
        #             # No usable phase labels after cleaning
        #             s3b_tbl = _new_slide("Testing/Evaluating — Phase × Yes/No — Table")
        #             _auto_table(s3b_tbl, pd.DataFrame({"Message": ["No phase tags found across records."]}))
        #             s3b_fig = _new_slide("Testing/Evaluating — Phase × Yes/No — Figure")
        #             _auto_table(s3b_fig, pd.DataFrame({"Message": ["No phase tags found across records."]}))
        #             _add_note(s3b_fig, "No phase labels were available to split by Testing/Evaluating (Yes/No).")
        #         else:
        #             # Pivot to Yes/No columns, compute totals & Yes%
        #             wide = counts_long.pivot(index="Phase", columns="Flag", values="Count").fillna(0)
        #             for col in ("Yes", "No"):
        #                 if col not in wide.columns:
        #                     wide[col] = 0
        #             wide = wide[["Yes", "No"]].astype(int)
        #             wide["Total"] = wide["Yes"] + wide["No"]
        #             wide["Yes %"] = (100.0 * np.where(wide["Total"] > 0, wide["Yes"] / wide["Total"], 0)).round(1)
        #
        #             # Order by Total (desc) for both table and figure
        #             wide_sorted = wide.sort_values("Total", ascending=False).reset_index()
        #
        #             # ----- TABLE -----
        #             s3b_tbl = _new_slide("Testing/Evaluating — Phase × Yes/No — Table")
        #             _auto_table(s3b_tbl, wide_sorted.rename(
        #                 columns={"Phase": "Phase", "Yes": "Yes (Count)", "No": "No (Count)"}))
        #             _add_note(
        #                 s3b_tbl,
        #                 "Counts split phases by the explicit Testing/Evaluating flag from 'evaluates_tests_value' "
        #                 "(Yes ∈ {Yes, True, Y, 1, T}; otherwise No). Phase tags are taken from 'phase_focus_value' "
        #                 "by splitting on commas/semicolons/pipes/slashes. 'Yes %' is Yes/(Yes+No) per phase."
        #             )
        #
        #             # ----- FIGURE (stacked bars: Yes in blue, No in red) -----
        #             # Build a long DF with percentages for labels (count + share within phase)
        #             long_plot = counts_long.merge(
        #                 wide_sorted[["Phase", "Total"]], on="Phase", how="left"
        #             )
        #             long_plot["Pct"] = np.where(long_plot["Total"] > 0,
        #                                         (100.0 * long_plot["Count"] / long_plot["Total"]).round(1), 0.0)
        #
        #             # Keep phase order consistent with table (Total desc)
        #             phase_order = wide_sorted["Phase"].tolist()
        #             long_plot["Phase"] = pd.Categorical(long_plot["Phase"], categories=phase_order, ordered=True)
        #             long_plot = long_plot.sort_values(["Phase", "Flag"])
        #
        #             # Text label shows both count and phase-share to convey "ratio"
        #             long_plot["Label"] = long_plot.apply(lambda r: f"{int(r['Count'])} ({r['Pct']:.1f}%)", axis=1)
        #
        #             s3b_fig = _new_slide("Testing/Evaluating — Phase × Yes/No — Figure")
        #             fig_phase = px.bar(
        #                 long_plot,
        #                 x="Phase", y="Count", color="Flag", text="Label",
        #                 category_orders={"Phase": phase_order, "Flag": ["Yes", "No"]},
        #                 title="Per-phase split of Testing/Evaluating (Yes vs No)"
        #             )
        #             fig_phase.update_layout(
        #                 barmode="stack",
        #                 margin=dict(l=60, r=120, t=60, b=80),
        #                 legend_title_text="Flag",
        #                 coloraxis_showscale=False
        #             )
        #             # Fix colors explicitly: Yes = blue, No = red
        #             fig_phase.for_each_trace(
        #                 lambda tr: tr.update(marker_color=("#2563EB" if tr.name == "Yes" else "#EF4444"))
        #             )
        #             fig_phase.update_traces(textposition="outside", cliponaxis=False)
        #             _fig_to_slide(s3b_fig, fig_phase)
        #
        #             _add_note(
        #                 s3b_fig,
        #                 "Stacked bars show absolute counts by phase, colored by the explicit flag in 'evaluates_tests_value' "
        #                 "(Yes in blue, No in red). Labels combine counts with within-phase percentages, so each phase’s ratio "
        #                 "of Yes/No is immediately visible. Ordering follows Total (Yes+No) descending."
        #             )
        #     else:
        #         s3b_tbl = _new_slide("Testing/Evaluating — Phase × Yes/No — Table")
        #         _auto_table(s3b_tbl, pd.DataFrame({"Message": ["No phase tags found across records."]}))
        #         s3b_fig = _new_slide("Testing/Evaluating — Phase × Yes/No — Figure")
        #         _auto_table(s3b_fig, pd.DataFrame({"Message": ["No phase tags found across records."]}))
        #         _add_note(s3b_fig, "No phase labels were available to split by Testing/Evaluating (Yes/No).")
        # else:
        #     s3b = _new_slide("Testing/Evaluating — Phase × Yes/No")
        #     _auto_table(s3b, pd.DataFrame({"Message": ["Column 'phase_focus_value' is not available."]}))
        #     _add_note(s3b, "Cannot compute phase split without 'phase_focus_value'.")

    # --- ADD: keep full outside text on bar charts (no cropping) ---
def _outside_text_headroom(
    fig: "go.Figure",
    *,
    orientation: str,           # "h" for horizontal bars, "v" for vertical bars
    values,                      # 1-D iterable of numeric bar magnitudes
    labels=None,                 # 1-D iterable of text labels (optional)
    pad_frac: float = 0.28,      # enlarge axis by ~28% to fit text tails
    extra_margin_px: int = 160   # add right/top margin in px to avoid PNG trimming
):
    """
    Expands the plotting range + margins so Plotly won't trim bar-end text like "123 (45.6%)".
    Call AFTER building traces and BEFORE export.
    """
    import numpy as _np
    vals = _np.asarray(values, dtype=float)
    vals = vals[_np.isfinite(vals)]
    if vals.size == 0:
        return

    vmax = float(vals.max())
    if vmax <= 0:
        vmax = 1.0

    if orientation.lower().startswith("h"):
        # Extend x-range and right margin
        fig.update_xaxes(range=[0, vmax * (1.0 + max(0.05, pad_frac))])
        m = fig.layout.margin.to_plotly_json() if fig.layout.margin is not None else {}
        fig.update_layout(margin=dict(
            l=int(m.get("l", 70)),
            r=int(max(m.get("r", 0), extra_margin_px)),
            t=int(m.get("t", 70)),
            b=int(m.get("b", 60)),
        ))
        # Make sure the text is actually placed "outside"
        fig.update_traces(selector=dict(type="bar"), textposition="outside")
    else:
        # Vertical bars: extend y-range and top margin
        fig.update_yaxes(range=[0, vmax * (1.0 + max(0.05, pad_frac))])
        m = fig.layout.margin.to_plotly_json() if fig.layout.margin is not None else {}
        fig.update_layout(margin=dict(
            l=int(m.get("l", 70)),
            r=int(m.get("r", 80)),
            t=int(max(m.get("t", 0), extra_margin_px)),
            b=int(m.get("b", 70)),
        ))
        fig.update_traces(selector=dict(type="bar"), textposition="outside")

def shape_scope(
    prs,
    df: "pd.DataFrame",
    collection_name: str = "Collection",
    *,
    slide_notes: bool = True,
    year_range: tuple[int, int] = (2010, 2025),
    top_n_cats: int = 20,
    include_percentages: bool = True,
    progress_callback=None,
) -> None:
    """
    Build 'Scope' slides (2 per metric: a Table slide, then a Figure slide):

      1) Publications by year
      2) Distribution by document type
      3) Split by publisher type
      4) Distribution by country focus
      5) Number of studies by phase_focus
      6) Number of studies by sector_focus

    Tables and figures are centered and constrained to fit the slide.
    """
    # ---- imports (kept local) ----
    import io, re
    import numpy as np, pandas as pd, plotly.express as px, plotly.graph_objects as go
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN  # <-- Correct name
    from pptx.dml.color import RGBColor

    # ---------- progress ----------
    def _cb(msg: str):
        if progress_callback:
            try: progress_callback(f"SCOPE: {msg}")
            except Exception: pass

    # ---------- global style ----------
    FONT_FAMILY = "Segoe UI"
    COLORWAY = ["#2563EB", "#10B981", "#F59E0B", "#EF4444", "#8B5CF6", "#14B8A6", "#F97316", "#3B82F6"]
    PLOT_W, PLOT_H = 1200, 650  # render slightly smaller -> crisper + easier to fit

    INCH = float(Inches(1))  # EMU-per-inch to float
    SLIDE_W_IN = float(prs.slide_width) / INCH
    SLIDE_H_IN = float(prs.slide_height) / INCH

    # margins (inches)
    OUTER_MARGIN_X = 0.6
    OUTER_MARGIN_Y = 0.6
    TITLE_HEIGHT_IN = 0.7

    # Consistent Plotly defaults
    def _apply_fig_style(fig: "go.Figure", title: str):
        fig.update_layout(
            template="plotly_white",
            width=PLOT_W, height=PLOT_H,
            margin=dict(l=60, r=40, t=60, b=50),
            title=dict(text=title, x=0.02, xanchor="left", font=dict(size=20, family=FONT_FAMILY)),
            font=dict(family=FONT_FAMILY, size=13),
            colorway=COLORWAY,
        )
        return fig

    # ---------- safe helpers (reuse if present globally) ----------
    _ensure_numeric_local = globals().get("_ensure_numeric", None)
    _safe_ai_notes_local = globals().get("_safe_ai_notes", None)
    _write_notes_local   = globals().get("_write_notes", None)

    def _ensure_numeric_df(x: "pd.DataFrame | None", fallback=("Value",), min_rows=1):
        if callable(_ensure_numeric_local):
            return _ensure_numeric_local(x, fallback_cols=fallback, min_rows=min_rows)
        if x is None or not isinstance(x, pd.DataFrame) or x.empty:
            return pd.DataFrame({fallback[0]: [0]})
        out = pd.DataFrame({c: pd.to_numeric(x[c], errors="coerce") for c in x.columns}).fillna(0)
        if out.empty:
            out = pd.DataFrame({fallback[0]: [0]})
        return out

    # ---------- parsing helpers ----------
    def _year_series(frame: pd.DataFrame) -> pd.Series:
        for cand in ("publication_year", "year_numeric", "year", "Year", "pub_year", "issued", "date"):
            if cand in frame.columns:
                s = pd.to_numeric(frame[cand], errors="coerce")
                if s.notna().any():
                    return s
                s2 = pd.to_datetime(frame[cand], errors="coerce", utc=True)
                if s2.notna().any():
                    return pd.to_numeric(s2.dt.year, errors="coerce")
        for c in frame.columns:
            if any(k in str(c).lower() for k in ("year", "date", "issued", "created", "publication")):
                s = pd.to_datetime(frame[c], errors="coerce", utc=True)
                if s.notna().any():
                    return pd.to_numeric(s.dt.year, errors="coerce")
        return pd.Series([np.nan] * len(frame), index=frame.index)

    def _doc_type(frame: pd.DataFrame) -> pd.Series:
        """
        Prefer 'document_type_value' if present; fall back to common type columns.
        """
        # 1) exact preferred column
        if "document_type_value" in frame.columns:
            s = frame["document_type_value"].astype(str).str.strip()
            if s.notna().any():
                return s

        # 2) other likely columns
        for cand in ("document_type", "item_type", "itemType", "item type", "type"):
            if cand in frame.columns:
                s = frame[cand].astype(str).str.strip()
                if s.notna().any():
                    return s

        # 3) fallback
        return pd.Series(["Unknown"] * len(frame), index=frame.index)

    def _source(frame: pd.DataFrame) -> pd.Series:
        for cand in ("source", "publication_outlet", "publicationTitle", "publisher", "institution"):
            if cand in frame.columns:
                s = frame[cand].astype(str).str.strip()
                if s.notna().any():
                    return s
        return pd.Series([""] * len(frame), index=frame.index)

    def _publisher_bucket(s: str) -> str:
        txt = (s or "").lower()
        think_tank = ["rand","heritage","brookings","chatham","carnegie","csis","rusi","atlantic council",
                      "ifri","dgap","cepa","nupi","sipri","institut","foundation"]
        policy = ["ministry","government","gov","commission","council","parliament","oecd","world bank","imf",
                  "un ","united nations","european commission","nato","white house","cabinet","home office","senate","house of commons"]
        academic = ["journal","revue","review","press","university","universität","universidade","springer","elsevier",
                    "wiley","taylor & francis","sage","ieee","acm","oxford","cambridge"]
        if any(k in txt for k in academic): return "Academic"
        if any(k in txt for k in policy): return "Policy/Government"
        if any(k in txt for k in think_tank): return "Think Tank"
        return "Academic" if "journal" in txt else "Other/Unknown"

    def _extract_tag_values(frame: pd.DataFrame, prefix: str) -> pd.Series:
        pref = prefix.lower().strip().lstrip("#")
        for cand in (f"{pref}", f"{pref}_value", f"{pref}s", f"{pref}_values"):
            if cand in frame.columns:
                return frame[cand].fillna("").astype(str)
        if "tags" in frame.columns: raw = frame["tags"]
        elif "Tags" in frame.columns: raw = frame["Tags"]
        else:
            tag_col = next((c for c in frame.columns if "tag" in str(c).lower()), None)
            raw = frame[tag_col] if tag_col else pd.Series([""]*len(frame), index=frame.index)
        out = []
        for v in raw:
            vals, candidates = [], ([str(x) for x in v] if isinstance(v,(list,tuple))
                                    else re.split(r"[;|,]\s*", str(v or "")) if isinstance(v,str) else [])
            for t in candidates:
                s = str(t).strip(); s2 = s.lstrip("#").strip()
                if s2.lower().startswith(pref + ":"):
                    parts = s2.split(":")
                    if len(parts) >= 2: vals.append(parts[-1].strip())
            seen=set(); uniq=[]
            for val in vals:
                k=val.casefold()
                if k not in seen: seen.add(k); uniq.append(val)
            out.append(", ".join(uniq))
        return pd.Series(out, index=frame.index)

    # ---------- small utilities ----------
    def _with_other_bin(vc: pd.Series, *, top_n: int) -> pd.DataFrame:
        vc = vc[vc.index.notna()].copy()
        if top_n and len(vc) > top_n:
            head = vc.head(top_n)
            other = int(vc.iloc[top_n:].sum())
            if other > 0: head.loc["Other"] = other
            return head
        return vc

    def _maybe_percent(df_counts: pd.DataFrame, count_col="Count"):
        if not include_percentages or df_counts[count_col].sum() <= 0:
            return df_counts
        total = df_counts[count_col].sum()
        df_counts = df_counts.copy()
        df_counts["%"] = (100.0 * df_counts[count_col] / total).round(1)
        return df_counts

    # ---------- slide primitives (centered table / centered figure) ----------
    def _new_blank_slide_with_title(title_text: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Title spans the entire slide width, so centered text is truly centered
        box = slide.shapes.add_textbox(
            Inches(0.0),
            Inches(0.15),
            Inches(SLIDE_W_IN),
            Inches(TITLE_HEIGHT_IN),
        )
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        return slide

    def _add_table_slide(title_text: str, df_small: pd.DataFrame):
        slide = _new_blank_slide_with_title(title_text)

        # Graceful empty
        if df_small is None or df_small.empty:
            tb = slide.shapes.add_textbox(Inches(OUTER_MARGIN_X), Inches(2.7),
                                          Inches(SLIDE_W_IN - 2 * OUTER_MARGIN_X), Inches(1.0))
            p = tb.text_frame.paragraphs[0]
            p.text = "No data available."
            p.font.name = FONT_FAMILY;
            p.font.size = Pt(16)
            p.alignment = PP_ALIGN.CENTER
            return slide

        # Available canvas for the table
        max_w = SLIDE_W_IN - 2 * OUTER_MARGIN_X
        max_h = SLIDE_H_IN - (TITLE_HEIGHT_IN + 2 * OUTER_MARGIN_Y)

        rows = max(1, len(df_small) + 1)
        cols = max(1, df_small.shape[1])

        # --- 1) Measure content need per column (approximate inches) ---
        # Heuristic: 11pt Segoe ≈ 0.085" per character for A–Z digits, plus padding.
        # Also consider header + longest cell in each column (after stripping).
        def _len_inch_est(s: object) -> float:
            t = "" if s is None or (isinstance(s, float) and pd.isna(s)) else str(s)
            t = " ".join(t.split())  # collapse spaces/newlines
            # 0.085 inch/char + small base
            return 0.085 * max(4, len(t)) + 0.14

        col_needs = []
        for j, col in enumerate(df_small.columns):
            header_need = _len_inch_est(col)
            longest_need = header_need
            # sample all rows; limit extreme outliers by a soft cap per cell
            for i in range(len(df_small)):
                longest_need = max(longest_need, min(_len_inch_est(df_small.iloc[i, j]), 3.8))
            col_needs.append(longest_need)

        # --- 2) Map needs → actual widths under total width cap ---
        # Bounds to keep columns readable but tight
        MIN_W = 0.9
        MAX_W = 3.6

        # clamp individual needs
        col_needs = [min(MAX_W, max(MIN_W, x)) for x in col_needs]
        total_need = sum(col_needs)

        # if total exceeds max_w, scale everything down proportionally
        if total_need > (max_w - 0.2):
            scale = (max_w - 0.2) / total_need
            col_widths_in = [max(MIN_W, x * scale) for x in col_needs]
        else:
            # otherwise, use the natural needs (no forced extra whitespace)
            col_widths_in = col_needs[:]

        table_w = sum(col_widths_in)
        # Row height heuristic: keep compact but never squash
        est_row_h = min(0.42, max(0.26, max_h / (rows + 1)))
        table_h = min(max_h, rows * est_row_h + 0.35)

        # Center the table box on the slide
        left_in = (SLIDE_W_IN - table_w) / 2.0
        top_in = TITLE_HEIGHT_IN + ((SLIDE_H_IN - TITLE_HEIGHT_IN) - table_h) / 2.0

        table = slide.shapes.add_table(
            rows, cols, Inches(left_in), Inches(top_in), Inches(table_w), Inches(table_h)
        ).table

        # Apply precise column widths
        for j, w in enumerate(col_widths_in):
            table.columns[j].width = Emu(Inches(w))

        # Header row
        for j, col in enumerate(df_small.columns):
            cell = table.cell(0, j)
            tf = cell.text_frame
            tf.clear();
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(col)
            p.font.bold = True
            p.font.name = FONT_FAMILY
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER
            # tighter margins; keep text from touching borders
            cell.margin_left = cell.margin_right = Emu(900)
            cell.margin_top = Emu(600);
            cell.margin_bottom = Emu(600)

        # Body cells
        for i, (_, r) in enumerate(df_small.iterrows(), start=1):
            for j, col in enumerate(df_small.columns):
                cell = table.cell(i, j)
                txt = "" if pd.isna(r[col]) else " ".join(str(r[col]).split())
                tf = cell.text_frame
                tf.clear();
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = txt
                p.font.name = FONT_FAMILY
                p.font.size = Pt(11)
                p.alignment = PP_ALIGN.LEFT
                cell.margin_left = Emu(800);
                cell.margin_right = Emu(800)
                cell.margin_top = Emu(500);
                cell.margin_bottom = Emu(500)

        return slide

    def _add_figure_slide(title_text: str, fig: "go.Figure"):
        slide = _new_blank_slide_with_title(title_text)

        content_w = SLIDE_W_IN - 2 * OUTER_MARGIN_X
        content_h = SLIDE_H_IN - (TITLE_HEIGHT_IN + 2 * OUTER_MARGIN_Y)
        content_left = OUTER_MARGIN_X
        content_top = TITLE_HEIGHT_IN + OUTER_MARGIN_Y

        try:
            png = fig.to_image(format="png", width=PLOT_W, height=PLOT_H, scale=2)
            ar = (PLOT_W / PLOT_H) if PLOT_H else 1.7778

            # Fit into content box while preserving aspect ratio
            target_w = content_w
            target_h = target_w / ar
            if target_h > content_h:
                target_h = content_h
                target_w = target_h * ar

            target_w *= 0.98
            target_h *= 0.98

            # Center within the full slide (not just content box)
            left_in = (SLIDE_W_IN - target_w) / 2.0
            top_in = TITLE_HEIGHT_IN + (SLIDE_H_IN - TITLE_HEIGHT_IN - target_h) / 2.0

            slide.shapes.add_picture(
                io.BytesIO(png),
                Inches(left_in),
                Inches(top_in),
                width=Inches(target_w),
                height=Inches(target_h),
            )
        except Exception as e:
            tb = slide.shapes.add_textbox(
                Inches(content_left),
                Inches(content_top + content_h / 2.0 - 0.5),
                Inches(content_w),
                Inches(1.0),
            )
            p = tb.text_frame.paragraphs[0]
            p.text = f"[Figure rendering failed: {type(e).__name__}]"
            p.font.name = FONT_FAMILY
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER

        return slide

    # ---------- auto-insight notes ----------
    def _auto_insights(df_small: pd.DataFrame, label_col: str | None, value_col: str, context: str) -> str:
        if df_small is None or df_small.empty or value_col not in df_small:
            return f"{context}: No data."
        s = df_small[value_col]
        total = float(s.sum())
        if total <= 0: return f"{context}: No counts."
        parts = [f"{context}: total n={int(total)}."]
        if label_col and label_col in df_small.columns and label_col != value_col:
            top = df_small.sort_values(value_col, ascending=False).head(3)
            tops = [f"{row[label_col]} ({int(row[value_col])})" for _, row in top.iterrows()]
            parts.append("Top: " + ", ".join(tops) + ".")
        if label_col and label_col.lower() == "year":
            tail = df_small.sort_values(label_col).tail(5)
            if len(tail) >= 2:
                delta = int(tail[value_col].iloc[-1]) - int(tail[value_col].iloc[-2])
                direction = "↑" if delta > 0 else "↓" if delta < 0 else "→"
                parts.append(f"Most recent vs prior: {direction} {delta:+d}.")
        top3 = int(df_small.sort_values(value_col, ascending=False).head(3)[value_col].sum())
        parts.append(f"Top-3 concentration: {top3/total:.0%}.")
        return " ".join(parts)

    def _notes(slide, title_for_prompt, purpose, display_df=None, label_col=None, value_col=None):
        text = None
        if callable(_safe_ai_notes_local):
            try:
                text = _safe_ai_notes_local(
                    title=title_for_prompt,
                    purpose=purpose,
                    data=display_df.to_dict("records") if display_df is not None else None,
                )
            except Exception:
                text = None
        if text is None:
            text = _auto_insights(display_df, label_col=label_col, value_col=value_col, context=title_for_prompt)
        if slide_notes and text:
            try: slide.notes_slide.notes_text_frame.text = text
            except Exception: pass

    # ---------- figure builders ----------
    def _barh(df_counts, xcol, ycol, title):
        fig = px.bar(df_counts, x=xcol, y=ycol, orientation="h")
        fig.update_traces(hovertemplate=f"{ycol}: %{{y}}<br>{xcol}: %{{x}}<extra></extra>",
                          marker_line_width=0.5, marker_line_color="rgba(0,0,0,0.15)")
        fig.update_yaxes(automargin=True)
        return _apply_fig_style(fig, title)

    def _bar(df_counts, xcol, ycol, title):
        fig = px.bar(df_counts, x=xcol, y=ycol)
        fig.update_traces(hovertemplate=f"{xcol}: %{{x}}<br>{ycol}: %{{y}}<extra></extra>",
                          marker_line_width=0.5, marker_line_color="rgba(0,0,0,0.15)")
        fig.update_xaxes(tickangle=-25, automargin=True)
        return _apply_fig_style(fig, title)

    # ===================== build data =====================
    work = df.copy()

    # 1) Publications by year  — REPLACEMENT
    _cb("Publications by year")

    # robust year series
    work["__year"] = pd.to_numeric(_year_series(work), errors="coerce")

    # guard: ensure valid range (swap if passed in reverse)
    y0, y1 = sorted([int(year_range[0]), int(year_range[1])])
    mask = (work["__year"] >= y0) & (work["__year"] <= y1)
    ydf = work.loc[mask].copy()

    if ydf.empty:
        counts_year = pd.DataFrame({
            "Year": list(range(y0, y1 + 1)),
            "Publications": [0] * (y1 - y0 + 1),
        })
    else:
        # count per year
        tmp = (
            ydf.dropna(subset=["__year"])
            .assign(__year=lambda d: d["__year"].astype(int))
            .groupby("__year", dropna=True).size()
            .rename("Publications").reset_index()
            .rename(columns={"__year": "Year"})
        )

        # complete the year index and coerce safely
        base = pd.DataFrame({"Year": list(range(y0, y1 + 1))})
        counts_year = (
            base.merge(tmp, on="Year", how="left")
            .assign(
                Publications=lambda d: pd.to_numeric(d["Publications"], errors="coerce")
                .replace([np.inf, -np.inf], np.nan)
                .fillna(0)
                .astype(int),
                Year=lambda d: pd.to_numeric(d["Year"], errors="coerce").fillna(0).astype(int),
            )[["Year", "Publications"]]
        )

    # slides
    _add_table_slide(f"Publications by Year ({y0}–{y1}) — Table", counts_year.copy())

    plot_df = counts_year.copy()

    fig = px.bar(plot_df, x="Year", y="Publications")
    fig.update_traces(
        cliponaxis=False,  # prevent cropping of labels above bars
        marker_line_width=0.5,
        marker_line_color="rgba(0,0,0,0.15)",
    )

    # add top headroom to keep labels clear
    try:
        _ymax = float(np.nanmax(pd.to_numeric(plot_df["Publications"], errors="coerce")))
        if np.isfinite(_ymax) and _ymax > 0:
            fig.update_yaxes(range=[0, _ymax * 1.25], automargin=True)
    except Exception:
        pass

    # extra right/top margins to avoid PNG cropping
    fig.update_layout(margin=dict(l=80, r=120, t=120, b=80))
    fig.update_xaxes(tickangle=-25, automargin=True)

    fig = _apply_fig_style(fig, f"{collection_name} · Publications per Year")
    fig_slide = _add_figure_slide(f"Publications by Year ({y0}–{y1}) — Figure", fig)

    _notes(
        fig_slide,
        f"Publications by Year ({y0}–{y1})",
        "Identify growth phases and recency undercounts.",
        counts_year,
        "Year",
        "Publications",
    )

    # ===================== 2) Document type =====================
    _cb("Document type distribution")

    # prefer 'document_type_value' if available
    if "document_type_value" in work.columns:
        work["__doctype"] = work["document_type_value"].fillna("").astype(str).str.strip()
    else:
        work["__doctype"] = _doc_type(work).replace("", "Unknown")

    vc = work["__doctype"].replace("", "Unknown").value_counts()
    vc = _with_other_bin(vc, top_n=top_n_cats)
    dt_counts = vc.rename_axis("Document Type").reset_index(name="Count")
    dt_counts = _maybe_percent(dt_counts, "Count")

    # Table slide (sorted by Count desc)
    _add_table_slide("Distribution by Document Type — Table",
                     dt_counts.sort_values("Count", ascending=False))

    # Figure
    plot_df = dt_counts.sort_values("Count", ascending=True).copy()
    if "%" in plot_df.columns:
        plot_df["Label"] = plot_df.apply(lambda r: f"{int(r['Count'])} ({float(r['%']):.1f}%)", axis=1)
    else:
        plot_df["Label"] = plot_df["Count"].astype(int).astype(str)

    fig = px.bar(plot_df, x="Count", y="Document Type", orientation="h")
    fig = _apply_fig_style(fig, f"{collection_name} · Document Types")

    # << stop truncation of ' (xx.x%)' >>
    _outside_text_headroom(
        fig,
        orientation="h",
        values=plot_df["Count"],
        labels=plot_df["Label"],
    )

    fig_slide = _add_figure_slide("Distribution by Document Type — Figure", fig)
    _notes(fig_slide, "Document Type Distribution",
           "Call out dominant formats and whether grey lit dominates.",
           dt_counts, "Document Type", "Count")

    # 3) Publisher type
    # ===================== 3) Publisher type =====================
    _cb("Publisher type split")

    # Prefer explicit publisher type columns as exported (with/without '_value', spaces, or underscores)
    if "publisher_type_value" in work.columns:
        pub_type = work["publisher_type_value"].astype(str).str.strip().replace("", np.nan)
    elif "publisher type" in work.columns:
        pub_type = work["publisher type"].astype(str).str.strip().replace("", np.nan)
    elif "publisher_type" in work.columns:
        pub_type = work["publisher_type"].astype(str).str.strip().replace("", np.nan)
    else:
        # Fallback to inference from outlet text
        pub_type = _source(work).apply(_publisher_bucket).replace("", np.nan)

    pub_type = pub_type.fillna("Other/Unknown")

    pub_counts = pub_type.value_counts().rename_axis("Publisher Type").reset_index(name="Count")
    pub_counts = _maybe_percent(pub_counts, "Count")

    _add_table_slide("Publisher Type Split — Table",
                     pub_counts.sort_values("Count", ascending=False))

    plot_df = pub_counts.sort_values("Count", ascending=False).copy()
    if "%" in plot_df.columns:
        plot_df["Label"] = plot_df.apply(lambda r: f"{int(r['Count'])} ({float(r['%']):.1f}%)", axis=1)
    else:
        plot_df["Label"] = plot_df["Count"].astype(int).astype(str)

    fig = px.bar(plot_df, x="Publisher Type", y="Count")
    fig.update_xaxes(tickangle=-25, automargin=True)
    fig = _apply_fig_style(fig, f"{collection_name} · Publisher Type")

    # vertical bars → extend Y so text above bars isn't cut
    _outside_text_headroom(
        fig,
        orientation="v",
        values=plot_df["Count"],
        labels=plot_df["Label"],
    )

    fig_slide = _add_figure_slide("Publisher Type Split — Figure", fig)
    _notes(fig_slide, "Publisher Type Split",
           "Assess venue balance; policy/think-tank heavy corpora may imply different evidence standards.",
           pub_counts, "Publisher Type", "Count")
    # ===================== 4) Country focus =====================
    _cb("Country focus distribution — Top 20")

    import re

    # --- 4a) Prefer explicit country-focus columns from the CSV ---
    for cand in ("country_focus_value", "country_focus", "country focus"):
        if cand in work.columns:
            country_series = work[cand].fillna("").astype(str)
            break
    else:
        country_series = _extract_tag_values(work, "country_focus_value")

    # If still empty, try affiliation:country tags explicitly (then generic affiliation)
    if country_series.eq("").all():
        country_series = _extract_tag_values(work, "affiliation:country")
    if country_series.eq("").all():
        country_series = _extract_tag_values(work, "affiliation")

    cf = pd.DataFrame({"Country Focus": country_series.replace("", np.nan)}).dropna()

    # --- 4b) Canonicalize country names / codes to reduce duplicates ---
    def _norm_country(x: str) -> str:
        s = (x or "").strip()
        if not s:
            return s
        s_key = s.lower().replace(".", "").replace("_", " ").strip()

        MAP = {
            # United States
            "us": "United States", "u s": "United States", "usa": "United States",
            "united states of america": "United States",
            # United Kingdom
            "uk": "United Kingdom", "u k": "United Kingdom", "gb": "United Kingdom",
            "gbr": "United Kingdom", "great britain": "United Kingdom",
            "england": "United Kingdom", "scotland": "United Kingdom",
            "wales": "United Kingdom", "northern ireland": "United Kingdom",
            # European Union
            "eu": "European Union", "e u": "European Union", "european union": "European Union",
            # Global / non-specific
            "global": "Global/International", "world": "Global/International",
            "international": "Global/International", "non-country specific": "Global/International",
            "non country specific": "Global/International",
        }
        if s_key in MAP:
            return MAP[s_key]
        # Title-case fallback for regular country names
        return s.title()

    def _uniq_normed_countries(cell: str) -> list[str]:
        """
        Split on commas/semicolons/pipes/slashes, normalize each token,
        and de-duplicate *within the same row* (so 'Brazil, Brazil' counts once for that record).
        """
        tokens = re.split(r"\s*[,;|/]\s*", str(cell or ""))
        seen, out = set(), []
        for tok in tokens:
            tok = tok.strip()
            if not tok:
                continue
            norm = _norm_country(tok)
            key = norm.casefold()
            if key not in seen:
                seen.add(key)
                out.append(norm)
        return out

    TOPN_COUNTRIES = 20

    if cf.empty:
        country_counts = pd.DataFrame({"Country Focus": ["N/A"], "Count": [0]})
    else:
        # explode AFTER per-row de-duplication
        exp = cf["Country Focus"].apply(_uniq_normed_countries).explode()
        vc = exp.value_counts()
        # Top 20 ONLY (no "Other" aggregation to avoid a large 'Other' bucket)
        vc_top = vc.head(TOPN_COUNTRIES)
        country_counts = vc_top.rename_axis("Country Focus").reset_index(name="Count")

    # (Optional) percentage column
    country_counts = _maybe_percent(country_counts, "Count")

    # ---------- Slides ----------
    # Table (sorted by Count desc)
    table_slide = _add_table_slide("Distribution by Country Focus — Top 20 — Table",
                                   country_counts.sort_values("Count", ascending=False))

    # Figure (horizontal bars)
    plot_df = country_counts.sort_values("Count", ascending=True).copy()
    plot_df["Label"] = (
        plot_df.apply(lambda r: f"{int(r['Count'])} ({float(r['%']):.1f}%)", axis=1)
        if "%" in plot_df.columns
        else plot_df["Count"].astype(int).astype(str)
    )

    fig = px.bar(plot_df, x="Count", y="Country Focus", orientation="h")
    fig = _apply_fig_style(fig, f"{collection_name} · Country Focus — Top 20")

    # prevent cropping of the percent tail / value labels
    _outside_text_headroom(
        fig,
        orientation="h",
        values=plot_df["Count"],
        labels=plot_df["Label"],
    )

    fig_slide = _add_figure_slide("Distribution by Country Focus — Top 20 — Figure", fig)

    # ---------- Notes (explicit, reliable) ----------
    table_notes = (
            "Methodology (Country Focus — Top 20):\n"
            "• Source columns (in order): 'country_focus_value', 'country_focus', 'country focus'. If all are empty, "
            "we parse tag-derived columns 'affiliation:country' then 'affiliation'.\n"
            "• Per-record parsing: split country strings on commas/semicolons/pipes/slashes; trim whitespace.\n"
            "• Normalization: common aliases are canonicalized — US/USA → United States; UK/GB/Great Britain/England/"
            "Scotland/Wales/Northern Ireland → United Kingdom; EU → European Union; "
            "Global/World/International → Global/International.\n"
            "• De-duplication: within each record, repeated countries are counted once; after that we aggregate across all records.\n"
            "• Display: top 20 countries only (no 'Other' bucket), to avoid inflating a catch-all category.\n"
            + ("• Percentages are relative to total country mentions after per-record de-duplication."
               if "%" in country_counts.columns else "")
    )
    table_slide.notes_slide.notes_text_frame.text = table_notes

    fig_notes = (
        "Figure shows the same Top 20 country counts as the table, using the normalization and per-record de-duplication rules "
        "described in the table notes."
    )
    fig_slide.notes_slide.notes_text_frame.text = fig_notes

    # ===================== 5) Phase focus =====================
    _cb("Phase focus distribution")
    # Prefer explicit phase columns (handles: 'phase_focus_value', 'phase_focus', 'phase focus')
    for cand in ("phase_focus_value", "phase_focus", "phase focus"):
        if cand in work.columns:
            phase_series = work[cand].fillna("").astype(str)
            break
    else:
        phase_series = _extract_tag_values(work, "phase_focus_value")

    ph = pd.DataFrame({"Phase Focus": phase_series.where(phase_series.ne(""), np.nan)}).dropna()

    if ph.empty:
        phase_counts = pd.DataFrame({"Phase Focus": ["N/A"], "Count": [0]})
    else:
        exp = ph["Phase Focus"].str.split(",").explode().str.strip()
        exp = exp[exp != ""]
        vc = _with_other_bin(exp.value_counts(), top_n=top_n_cats)
        phase_counts = vc.rename_axis("Phase Focus").reset_index(name="Count")

    phase_counts = _maybe_percent(phase_counts, "Count")

    _add_table_slide("Number of Studies by Phase Focus — Table",
                     phase_counts.sort_values("Count", ascending=False))

    plot_df = phase_counts.sort_values("Count", ascending=False).copy()
    if "%" in plot_df.columns:
        plot_df["Label"] = plot_df.apply(lambda r: f"{int(r['Count'])} ({float(r['%']):.1f}%)", axis=1)
    else:
        plot_df["Label"] = plot_df["Count"].astype(int).astype(str)

    fig = px.bar(plot_df, x="Phase Focus", y="Count")
    fig.update_xaxes(tickangle=-25, automargin=True)
    fig = _apply_fig_style(fig, f"{collection_name} · Phase Focus")

    _outside_text_headroom(
        fig,
        orientation="v",
        values=plot_df["Count"],
        labels=plot_df["Label"],
    )

    fig_slide = _add_figure_slide("Number of Studies by Phase Focus — Figure", fig)
    _notes(fig_slide, "Phase Focus",
           "Assess clustering (e.g., Designing/Implementing) and gaps for future sampling.",
           phase_counts, "Phase Focus", "Count")

    # ===================== 6) Sector focus =====================
    _cb("Sector focus distribution")
    # Prefer explicit sector columns (handles: 'sector_focus_value', 'sector_focus', 'sector focus')
    for cand in ("sector_focus_value", "sector_focus", "sector focus"):
        if cand in work.columns:
            sector_series = work[cand].fillna("").astype(str)
            break
    else:
        sector_series = _extract_tag_values(work, "sector_focus_value")

    sf = pd.DataFrame({"Sector Focus": sector_series.replace("", np.nan)}).dropna()

    if sf.empty:
        sector_counts = pd.DataFrame({"Sector Focus": ["N/A"], "Count": [0]})
    else:
        exp = sf["Sector Focus"].str.split(",").explode().str.strip()
        exp = exp[exp != ""]
        vc = _with_other_bin(exp.value_counts(), top_n=top_n_cats)
        sector_counts = vc.rename_axis("Sector Focus").reset_index(name="Count")

    sector_counts = _maybe_percent(sector_counts, "Count")

    _add_table_slide("Number of Studies by Sector Focus — Table",
                     sector_counts.sort_values("Count", ascending=False))

    plot_df = sector_counts.sort_values("Count", ascending=True).copy()
    if "%" in plot_df.columns:
        plot_df["Label"] = plot_df.apply(lambda r: f"{int(r['Count'])} ({float(r['%']):.1f}%)", axis=1)
    else:
        plot_df["Label"] = plot_df["Count"].astype(int).astype(str)

    fig = px.bar(plot_df, x="Count", y="Sector Focus", orientation="h")
    fig = _apply_fig_style(fig, f"{collection_name} · Sector Focus")

    _outside_text_headroom(
        fig,
        orientation="h",
        values=plot_df["Count"],
        labels=plot_df["Label"],
    )

    fig_slide = _add_figure_slide("Number of Studies by Sector Focus — Figure", fig)
    _notes(fig_slide, "Sector Focus",
           "Identify dominant sectors and non-sector specific load; check tag consistency.",
           sector_counts, "Sector Focus", "Count")

    _cb("Scope slides completed.")


# keywords =update_excel_counts_single_scan(excel_path=r"C:\Users\luano\Downloads\Skills_paragraph_cout.xlsx")

col1="0.13_cyber_attribution_corpus_records_total_included"
#
# print(keywords)
#
# input("aaaaaa")
#
# df, items, _ = load_data_from_source_for_widget(collection_name=col1, source_type="zotero", file_path=None)
#
# d=build_theme_groups(df=df)

keywords= [
"indicators of compromise",
"command and control",
"malware",
"critical infrastructure",
"ip address",
"ip spoofing",
"proxy",
"tor",
"vpn ",
"traffic",
"tampering",
"timestamp",
"signature",
"exfiltration",
"packet ",
"network telemetry",
"sensor bias",
"tactics techniques and procedures",
    "TTPs"
"chain of custody",
"burden of proof",
"standard of proof",
"rules of evidence",
"admissible evidence",
"circumstantial evidence",
    "probative"
"classified information",
"national security",
"state responsibility",
"effective control",
"overall control",
"due process",
"procedural fairness",
"jurisdiction conflicts",
"document authenticity",
"witness reliability",
"record preservation",
"control test",
"proportionality",
"sovereignty",
"countermeasures",
"threshold",
"standard of",
"international court of justice",
    "ICJ",
"articles on responsibility",
"naming and shaming",
"strategic ambiguity",
"public attribution",
"private attribution",
"deterrence ",
"escalation",
"collective attribution",
"sanctions policy",
"information sharing",
"intelligence sharing",
"capability protection",
"intelligence",
"trust deficits",
"burden of",
"reputational costs",
"public private partnership",
"confidence building measures",
"norm development",
"attributing",
    "attribution",
    "attributed",
    "misattribution",
    "sources and methods"
    
"deterrence by punishment"
]

col1= "skills and mission"
# pptx_out = build_bibliometrics_pptx(
#     collection_name=col1,
#     csv_path=r"C:\Users\luano\Downloads\Results_skills_total_included_official.csv",
#
#     zotero_client=zt,
#     # optional knobs:
#     keywords=keywords[:3],
#     top_ngram=20,
#     corpus_source="full_text",  # or "full_text" | "title" | "abstract"
#     collection_name_for_cache=col1,
#     # include can be omitted to render ALL sections + n-gram subsections
#     include=[
#         "Scope_and_shape",
#         "Authorship_institution",
#         "Thematic_method",
#         "Theory_empirical",
#         "Citations_influence",
#         "Geo_sector",
#
#         # "data_summary",
#         # "affiliations",
#         # "authors",
#         # "keywords",
#     #     #      {"ngrams":["basic"]}
#              ],
#     slide_notes=False,
#     image_width_inches=10.0,
#     rebuild=False
# )
#
# """
# def describe_trend_data(trend"""
# # print(f"Saved PowerPoint to: {pptx_out}")
#


import os
import json
from pathlib import Path

from src.core.utils.calling_models import call_models_plots


def test_call_models_plots_once() -> None:
    """
    ###1. set explicit env for model + cache dir
    ###2. build a prompt that must return visible text
    ###3. call call_models_plots(mode="analyze") with stable suffix/topic
    ###4. print text + key debug fields + cache path contents hint
    """
    os.environ["OPENAI_VISION_MODEL"] = os.environ.get("OPENAI_VISION_MODEL", "gpt-5-mini")
    cache_dir = str(Path.cwd() / ".annotarium_llm_cache_test")
    os.environ["ANNOTARIUM_CACHE_DIR"] = cache_dir

    model_name = os.getenv("OPENAI_VISION_MODEL", "gpt-5-mini")

    analysis_id = "author_impact_analysis"
    plot_type = "citations_avg_bar"
    params = {"analysis_type": analysis_id, "plot_type": plot_type, "top_n_authors": 15}

    prompt_text = (
        "[TASK]\n"
        "Write a short description for a chart in an author bibliometrics dashboard.\n"
        "You MUST output visible plain text.\n\n"
        "[OUTPUT]\n"
        "Return exactly 2 paragraphs of plain text.\n"
        "Paragraph 1: what the chart shows (2–4 sentences).\n"
        "Paragraph 2: how to read it + caveats (2–4 sentences).\n\n"
        f"[CONTEXT]\nanalysis_id={analysis_id}\nplot_type={plot_type}\n\n"
        "[PARAMS_JSON]\n"
        f"{json.dumps(params, ensure_ascii=False, sort_keys=True)}\n"
    )

    out = call_models_plots(
        mode="analyze",
        prompt_text=prompt_text,
        model_api_name=model_name,
        images=None,
        image_detail="low",
        max_tokens=5000,
        use_cache=True,
        cache_dir=cache_dir,
        analysis_key_suffix=f"describe_{analysis_id}__{plot_type}",
        section_title=analysis_id,
        overall_topic=analysis_id,
        store_only=False,
        read=False,
        tools=None,
        tool_router=None,
        tool_choice="auto",
        allowed_tools=None,
        parallel_tool_calls=None,
        instructions=(
            "Return plain text only. "
            "Do not return JSON. "
            "Do not omit the final answer."
        ),
        max_function_loops=2,
    )

    text = (out.get("text") or "").strip()
    print("\n=== call_models_plots TEST RESULT ===", flush=True)
    print(f"model_used={out.get('model_used')!r}", flush=True)
    print(f"text_len={len(text)}", flush=True)
    print(f"text_preview={text[:300]!r}", flush=True)
    print(f"usage={out.get('usage')}", flush=True)
    print(f"cache_path={out.get('cache_path')!r}", flush=True)
    print(f"cache_key={out.get('cache_key')!r}", flush=True)
    print(f"raw_response_present={out.get('raw_response') is not None}", flush=True)


# if __name__ == "__main__":
#     test_call_models_plots_once()
